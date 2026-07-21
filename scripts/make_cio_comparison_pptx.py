#!/usr/bin/env python3
"""
make_cio_comparison_pptx.py
============================
PAX channel histograms: raw  vs  CIO (no-RB)  vs  CIO-RB, per dataset.

Slides
------
  1. Title
  2-5. Per-dataset: 2-row × 3-col grid
         row 0 — full range:  raw | CIO | CIO-RB
         row 1 — zoomed:      raw [0,0.1] | CIO [-0.05,1] | CIO-RB [-0.05,1]
         each curve = one source frame (all its mr10 patches pooled)  bold = mean
  6. Summary: all 4 dataset mean curves overlaid (3 panels)
  7. Summary zoomed (same x-limits)

Usage:
  cd /net/projects/CLS/lding/gitcode/SubCellAE
  PYTHONPATH=... python3 scripts/make_cio_comparison_pptx.py
"""

from __future__ import annotations

import collections
import io
import re
from pathlib import Path

import czifile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT     = Path("/net/projects/CLS/lding/data/fa_data_analysis")
BASE_RB  = ROOT / "ae_results/patches/cio_rb"
BASE_CIO = ROOT / "ae_results/patches/cio"
OUT_PPTX = Path("cio_comparison_pax.pptx")

PAX_CH = 1
N_BINS = 200
PATCH_RE = re.compile(r"_f(\d{4})x(\d{4})y(\d{4})ps(\d+)\.tif$")

DATASETS = [
    {
        "key":   "vinc",
        "label": "ds1 - vinc",
        "color": "#2E86C1",
        "conditions": [
            {"cond": "control",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Control"},
            {"cond": "ycomp",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Ycomp"},
        ],
    },
    {
        "key":   "pfak",
        "label": "ds2 - pfak",
        "color": "#C44E52",
        "conditions": [
            {"cond": "control",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Control"},
            {"cond": "ycomp",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Ycomp"},
        ],
    },
    {
        "key":   "ppax",
        "label": "ds3 - ppax",
        "color": "#55A868",
        "conditions": [
            {"cond": "control",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Control"},
            {"cond": "ycomp",
             "czi_dir": ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Y-comp"},
        ],
    },
    {
        "key":   "nih3t3",
        "label": "ds4 - nih3t3",
        "color": "#DD8452",
        "conditions": [
            {"cond": "control",
             "czi_dir": ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/Control"},
            {"cond": "ycomp",
             "czi_dir": ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/YCompound"},
        ],
    },
]


# ── data loading ─────────────────────────────────────────────────────────────

def _parse_patch(name: str):
    m = PATCH_RE.search(name)
    return (int(m.group(1)), int(m.group(2)), int(m.group(3)), int(m.group(4))) if m else None


def _load_czi_pax(czi_path: Path) -> np.ndarray:
    arr = czifile.imread(str(czi_path)).squeeze()
    ch = arr[PAX_CH] if arr.ndim == 3 else arr
    return ch.astype(np.float32) / 65535.0


def _pool_patches_by_frame(patch_dir: Path) -> dict[int, np.ndarray]:
    """Return {frame_idx: pooled_pixel_array} from all tif patches in patch_dir."""
    frame_pixels: dict[int, list] = collections.defaultdict(list)
    if not patch_dir.exists():
        return {}
    for pf in patch_dir.glob("*.tif"):
        parsed = _parse_patch(pf.name)
        if parsed:
            fidx = parsed[0]
            frame_pixels[fidx].append(tifffile.imread(str(pf)).astype(np.float32).ravel())
    return {f: np.concatenate(arrs) for f, arrs in frame_pixels.items()}


def gather_frames(ds: dict) -> list[tuple[np.ndarray, np.ndarray, np.ndarray, str]]:
    """
    Returns list of (raw_pixels, cio_pixels, cio_rb_pixels, label) per frame.
    Only frames present in BOTH cio and cio_rb are included.
    raw_pixels: CZI PAX channel at the same patch locations (uint16/65535).
    """
    results = []
    for cond in ds["conditions"]:
        cond_name = cond["cond"]
        czi_dir   = Path(cond["czi_dir"])
        cio_dir   = BASE_CIO / ds["key"] / cond_name / "tiff_patches32_mr10"
        rb_dir    = BASE_RB  / ds["key"] / cond_name / "tiff_patches32_mr10"

        cio_frames = _pool_patches_by_frame(cio_dir)
        rb_frames  = _pool_patches_by_frame(rb_dir)
        common     = sorted(set(cio_frames) & set(rb_frames))

        if not common:
            print(f"    [{cond_name}] no matched frames — cio:{len(cio_frames)} rb:{len(rb_frames)}")
            continue

        czi_files = sorted(czi_dir.glob("*.czi")) if czi_dir.exists() else []
        print(f"    [{cond_name}] {len(common)} matched frames, {len(czi_files)} CZI files")

        # Build raw pixels at same patch locations
        raw_by_frame: dict[int, list] = collections.defaultdict(list)
        for pf in sorted((BASE_RB / ds["key"] / cond_name / "tiff_patches32_mr10").glob("*.tif")):
            parsed = _parse_patch(pf.name)
            if not parsed:
                continue
            fidx, x, y, ps = parsed
            if fidx not in common:
                continue
            if fidx >= len(czi_files):
                continue
            raw_by_frame[fidx].append((x, y, ps, czi_files[fidx]))

        for fidx in common:
            label = f"{cond_name}_f{fidx:04d}"

            # raw pixels: load CZI once per frame
            raw_parts = []
            czi_loaded = None
            for x, y, ps, czi_path in raw_by_frame.get(fidx, []):
                if czi_loaded is None:
                    try:
                        czi_loaded = _load_czi_pax(czi_path)
                    except Exception as e:
                        print(f"      WARN CZI load failed {czi_path.name}: {e}")
                        break
                H, W = czi_loaded.shape
                raw_parts.append(czi_loaded[y: min(y+ps, H), x: min(x+ps, W)].ravel())

            if not raw_parts:
                continue

            results.append((
                np.concatenate(raw_parts),
                cio_frames[fidx],
                rb_frames[fidx],
                label,
            ))

    return results


# ── plotting helpers ──────────────────────────────────────────────────────────

def _plot_curves(ax, arrays: list[np.ndarray], bold_color: str,
                 xlim=None, alpha=0.3, lw=0.8):
    if not arrays:
        return None, None
    n = len(arrays)
    cmap = plt.get_cmap("gist_rainbow")
    all_vals = np.concatenate(arrays)
    if xlim is not None:
        lo, hi = xlim
        all_vals = all_vals[(all_vals >= lo) & (all_vals <= hi)]
        if len(all_vals) == 0:
            return None, None
        bins = np.linspace(lo, hi, N_BINS + 1)
    else:
        bins = np.linspace(float(all_vals.min()), float(all_vals.max()), N_BINS + 1)

    mean_hist = np.zeros(N_BINS, dtype=np.float64)
    for i, arr in enumerate(arrays):
        counts, _ = np.histogram(arr, bins=bins, density=True)
        ax.plot(bins[:-1], counts, color=cmap(i / max(n-1, 1)), alpha=alpha, lw=lw)
        mean_hist += counts
    mean_hist /= n
    ax.plot(bins[:-1], mean_hist, color=bold_color, lw=2.2, alpha=0.95, label="mean")
    return bins, mean_hist


def make_dataset_figure(ds: dict, frames: list) -> plt.Figure:
    raw_arrs  = [r  for r, c, rb, _ in frames]
    cio_arrs  = [c  for r, c, rb, _ in frames]
    rb_arrs   = [rb for r, c, rb, _ in frames]
    n = len(frames)

    fig, axes = plt.subplots(2, 3, figsize=(15, 8))
    fig.suptitle(f"{ds['label']}  —  PAX channel: raw vs CIO vs CIO-RB  ({n} frames)",
                 fontsize=13, fontweight="bold")

    col_cfg = [
        ("Raw (uint16/65535)",    raw_arrs,  None,            (0.0, 0.1)),
        ("After CIO (no RB)",     cio_arrs,  None,            (-0.05, 1.0)),
        ("After CIO-RB",          rb_arrs,   None,            (-0.05, 1.0)),
    ]

    for col, (title, arrs, _, zoom) in enumerate(col_cfg):
        # full range
        ax = axes[0, col]
        _plot_curves(ax, arrs, bold_color=ds["color"])
        ax.set_title(title, fontsize=10)
        ax.set_xlabel("pixel value"); ax.set_ylabel("density")
        ax.legend(fontsize=7)

        # zoomed
        ax = axes[1, col]
        _plot_curves(ax, arrs, bold_color=ds["color"], xlim=zoom)
        ax.set_title(f"{title}  [zoom: {zoom[0]} – {zoom[1]}]", fontsize=9)
        ax.set_xlabel("pixel value"); ax.set_ylabel("density")
        ax.set_xlim(zoom)
        ax.legend(fontsize=7)
        if col > 0:
            ax.axvline(0.0, color="gray", ls=":", lw=1.0, alpha=0.7)
            ax.axvline(1.0, color="red",  ls="--", lw=1.2, alpha=0.8, label="cell mean =1")
            ax.legend(fontsize=7)

    fig.text(0.5, 0.01,
             f"{n} frames  |  each curve = one frame (all mr10 patches pooled)  |  bold = mean",
             ha="center", fontsize=8, color="#555555")
    plt.tight_layout(rect=[0, 0.03, 1, 0.96])
    return fig


def _violin(ax, data_list: list[np.ndarray], labels: list[str],
            colors: list[str], ylim=None):
    """Draw violin + median line for each group; clip data to ylim if given."""
    clipped = []
    for d in data_list:
        if ylim is not None:
            d = d[(d >= ylim[0]) & (d <= ylim[1])]
        clipped.append(d if d.size else np.array([0.0]))

    parts = ax.violinplot(clipped, positions=range(len(clipped)),
                          showmedians=True, showextrema=False)
    for i, (pc, col) in enumerate(zip(parts["bodies"], colors)):
        pc.set_facecolor(col)
        pc.set_alpha(0.65)
    parts["cmedians"].set_color("black")
    parts["cmedians"].set_linewidth(1.5)

    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=9)
    ax.axhline(0.0, color="gray", ls="--", lw=0.9, alpha=0.8)
    if ylim is not None:
        ax.set_ylim(ylim)
    ax.set_ylabel("CIO − CIO_RB")


def make_diff_figure(ds: dict, frames: list) -> plt.Figure:
    """Violin plots of CIO − CIO_RB per condition: full range + zoomed ±0.05."""
    n = len(frames)
    fig, axes = plt.subplots(1, 2, figsize=(10, 6))
    fig.suptitle(
        f"{ds['label']}  —  CIO − CIO_RB pixel difference  ({n} frames)",
        fontsize=13, fontweight="bold",
    )

    # group diffs by condition
    cond_diffs: dict[str, list] = collections.defaultdict(list)
    for _, cio, rb, label in frames:
        cond = label.split("_f")[0]
        cond_diffs[cond].append(cio - rb)

    cond_names  = sorted(cond_diffs.keys())
    data_list   = [np.concatenate(cond_diffs[c]) for c in cond_names]
    colors      = [ds["color"]] * len(cond_names)

    all_diff = np.concatenate(data_list)
    std_d = all_diff.std()
    pct   = 100.0 * np.mean(np.abs(all_diff) > 0.01)

    stats_lines = []
    for cname, d in zip(cond_names, data_list):
        p = 100.0 * np.mean(np.abs(d) > 0.01)
        stats_lines.append(f"{cname}: std={d.std():.5f}  |Δ|>0.01={p:.1f}%")

    _violin(axes[0], data_list, cond_names, colors)
    axes[0].set_title("Full range", fontsize=10)

    _violin(axes[1], data_list, cond_names, colors, ylim=(-0.05, 0.05))
    axes[1].set_title("Zoomed  ±0.05", fontsize=10)

    fig.text(
        0.5, 0.01,
        f"Overall: std={std_d:.5f}  |Δ|>0.01={pct:.1f}%     "
        + "     ".join(stats_lines),
        ha="center", fontsize=8, color="#555555",
    )
    plt.tight_layout(rect=[0, 0.05, 1, 0.96])
    return fig


def make_diff_summary_figure(ds_list: list[dict], ds_frames: list[list]) -> plt.Figure:
    """All 4 datasets as violins: full range + zoomed ±0.05."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 6))
    fig.suptitle("All datasets — CIO − CIO_RB difference (PAX, all conditions pooled)",
                 fontsize=13, fontweight="bold")

    data_list, colors, labels = [], [], []
    legend_lines = []
    for ds, frames in zip(ds_list, ds_frames):
        d = np.concatenate([cio - rb for _, cio, rb, _ in frames])
        pct = 100.0 * np.mean(np.abs(d) > 0.01)
        data_list.append(d)
        colors.append(ds["color"])
        labels.append(ds["label"])
        legend_lines.append(f"{ds['label']}: std={d.std():.5f}  |Δ|>0.01={pct:.1f}%")

    _violin(axes[0], data_list, labels, colors)
    axes[0].set_title("Full range", fontsize=10)
    axes[0].tick_params(axis="x", labelrotation=15)

    _violin(axes[1], data_list, labels, colors, ylim=(-0.05, 0.05))
    axes[1].set_title("Zoomed  ±0.05", fontsize=10)
    axes[1].tick_params(axis="x", labelrotation=15)

    fig.text(
        0.5, 0.01,
        "     ".join(legend_lines),
        ha="center", fontsize=8, color="#555555",
    )
    plt.tight_layout(rect=[0, 0.06, 1, 0.96])
    return fig


def make_summary_figure(ds_list: list[dict],
                        ds_frames: list[list],
                        zoomed: bool = False) -> plt.Figure:
    fig, axes = plt.subplots(1, 3, figsize=(15, 5))
    zoom_tag = " (zoomed)" if zoomed else " (full range)"
    fig.suptitle(f"All datasets — PAX mean distributions{zoom_tag}",
                 fontsize=13, fontweight="bold")

    xlims = [None, None, None] if not zoomed else [(0.0, 0.1), (-0.05, 1.0), (-0.05, 1.0)]
    titles = ["Raw (uint16/65535)", "After CIO (no RB)", "After CIO-RB"]

    for ds, frames in zip(ds_list, ds_frames):
        raw_arrs = [r  for r, c, rb, _ in frames]
        cio_arrs = [c  for r, c, rb, _ in frames]
        rb_arrs  = [rb for r, c, rb, _ in frames]

        for col, arrs in enumerate([raw_arrs, cio_arrs, rb_arrs]):
            # get mean curve via silent plot
            fig_tmp, ax_tmp = plt.subplots()
            bins, mean = _plot_curves(ax_tmp, arrs, bold_color="k",
                                      xlim=xlims[col], alpha=0, lw=0)
            plt.close(fig_tmp)
            if mean is not None:
                axes[col].plot(bins[:-1], mean, color=ds["color"],
                               lw=1.5, label=ds["label"])

    for col, (ax, title, xlim) in enumerate(zip(axes, titles, xlims)):
        ax.set_title(title, fontsize=10)
        ax.set_xlabel("pixel value"); ax.set_ylabel("density")
        ax.legend(fontsize=9)
        if xlim:
            ax.set_xlim(xlim)
        if col > 0:
            ax.axvline(0.0, color="gray", ls=":", lw=1.0, alpha=0.7)
            ax.axvline(1.0, color="red",  ls="--", lw=1.2, alpha=0.8)

    fig.text(0.5, 0.01, "Mean over all frames per dataset",
             ha="center", fontsize=9, color="#555555")
    plt.tight_layout(rect=[0, 0.04, 1, 0.96])
    return fig


# ── PPT helpers ───────────────────────────────────────────────────────────────

def _new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def _fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "PAX Channel: Raw vs CIO vs CIO-RB"
    p.font.size = Pt(34); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.33), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = ("mr10 patches  |  each curve = one frame  |  "
               "zoomed: raw [0,0.1]  /  normed [-0.05,1]  |  bold = mean")
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC)
    p2.alignment = PP_ALIGN.CENTER


def _image_slide(prs, buf):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(buf, Inches(0.1), Inches(0.1),
                             width=Inches(13.13), height=Inches(7.3))


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = _new_prs()
    _title_slide(prs)

    ds_list, ds_frames = [], []
    for ds in DATASETS:
        print(f"\n=== {ds['label']} ===")
        frames = gather_frames(ds)
        print(f"  {len(frames)} frames with matched cio+cio_rb+raw")
        if not frames:
            continue
        ds_list.append(ds)
        ds_frames.append(frames)
        fig = make_dataset_figure(ds, frames)
        _image_slide(prs, _fig_to_buf(fig))

    if ds_list:
        print("\n=== Summary slides ===")
        _image_slide(prs, _fig_to_buf(make_summary_figure(ds_list, ds_frames, zoomed=False)))
        _image_slide(prs, _fig_to_buf(make_summary_figure(ds_list, ds_frames, zoomed=True)))

        print("\n=== CIO − CIO_RB diff slides ===")
        for ds, frames in zip(ds_list, ds_frames):
            _image_slide(prs, _fig_to_buf(make_diff_figure(ds, frames)))
        _image_slide(prs, _fig_to_buf(make_diff_summary_figure(ds_list, ds_frames)))

    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
