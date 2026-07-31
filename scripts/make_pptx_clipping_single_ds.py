#!/usr/bin/env python3
"""
make_pptx_clipping_single_ds.py
================================
Short comparison PPT: single-dataset training only, three clipping conditions.

  Group A: ds_combo_enlcrop_sc2          no clip  | ÷2 | nL1  | λ=0.25
  Group B: ds_combo_enlcrop_clip01_l1    clip[0,1]|    | L1   | λ=0.10
  Group C: ds_combo_enlcrop_sc2_clip02_l1 clip[0,2]|÷2 | L1   | λ=0.10

Datasets: vinc, nih3t3, ppax, pfak  (single-dataset training only)

Per dataset: UMAP z_proj (FA annotation) | UMAP K-means (10 clusters)
             + one slide per cluster (4×4 patch grid, A/B/C side-by-side)

Run from the SubCellAE repo root:
  python scripts/make_pptx_clipping_single_ds.py
"""
from __future__ import annotations

import io
from pathlib import Path

import numpy as np
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
OUT  = Path("clipping_single_ds.pptx")

GROUPS = [
    ("A: no clip (÷2, nL1, λ=0.25)",  "ds_combo_enlcrop_sc2"),
    ("B: clip[0,1] (L1, λ=0.10)",     "ds_combo_enlcrop_clip01_l1"),
    ("C: clip[0,2]÷2 (L1, λ=0.10)",   "ds_combo_enlcrop_sc2_clip02_l1"),
]

SINGLE_DS = ["vinc", "nih3t3", "ppax", "pfak"]

# ── slide geometry ─────────────────────────────────────────────────────────────

SW = Inches(13.33)
SH = Inches(7.5)
TITLE_H = Inches(0.65)
PAD     = Inches(0.12)

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_ALT   = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)

# ── helpers ────────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=14, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _img(slide, src, l, t, max_w, max_h):
    if src is None:
        return
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return
        data = p.read_bytes()
    else:
        data = bytes(src)
    try:
        im = Image.open(io.BytesIO(data))
        iw, ih = im.size
    except Exception:
        return
    scale = min(max_w / iw, max_h / ih)
    w, h  = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(io.BytesIO(data),
                             l + (max_w - w) // 2,
                             t + (max_h - h) // 2, w, h)

def _title_bar(slide, text, size=18, bg=C_DARK):
    _rect(slide, 0, 0, SW, TITLE_H, fill=bg)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         text, size=size, bold=True, color=C_WHITE)

def _norm_u8(arr: np.ndarray) -> np.ndarray:
    arr = arr.astype(np.float32)
    vmin, vmax = arr.min(), arr.max()
    if vmax > vmin:
        arr = ((arr - vmin) / (vmax - vmin) * 255).clip(0, 255)
    return arr.astype(np.uint8)


def _strip_to_4x4(strip: np.ndarray) -> np.ndarray:
    """Convert a (32, 512) 1×16 patch strip → 4×4 grid image with 2-px gaps."""
    H, W = strip.shape
    pw   = H                        # patch width/height = 32
    n    = W // pw                  # number of patches (16)
    cols = 4
    rows = (n + cols - 1) // cols   # = 4
    gap  = 2
    ch   = rows * pw + (rows + 1) * gap
    cw   = cols * pw + (cols + 1) * gap
    canvas = np.zeros((ch, cw), dtype=np.uint8)
    for idx in range(min(n, rows * cols)):
        r, c = divmod(idx, cols)
        y = gap + r * (pw + gap)
        x = gap + c * (pw + gap)
        canvas[y:y+pw, x:x+pw] = strip[:, idx*pw:(idx+1)*pw]
    return canvas


def _load_cluster_strip(cp_dir: Path, k: int) -> np.ndarray | None:
    """Load and normalise a single cluster strip TIF; return uint8 (32,512) or None."""
    tif = None
    lcsv = cp_dir / "cluster_labels.csv"
    if lcsv.exists():
        try:
            import pandas as pd
            df = pd.read_csv(lcsv)
            match = df[df["cluster"] == k]
            if len(match):
                tif = cp_dir / match.iloc[0]["file"]
        except Exception:
            pass
    if tif is None or not tif.exists():
        candidates = sorted(cp_dir.glob(f"cluster_{k:02d}_*.tif"))
        tif = candidates[0] if candidates else None
    if tif is None or not tif.exists():
        return None
    try:
        return _norm_u8(tifffile.imread(str(tif)))
    except Exception:
        return None


def _all_clusters_mosaic(combo: str, dn: str, n_clusters: int = 10) -> bytes | None:
    """Build a 2×5 mosaic of 4×4 cluster panels for one group/combo, with cluster ID labels."""
    from PIL import ImageDraw, ImageFont
    cp_dir = RUNS / dn / combo / "eval" / "cluster_panels_combo"
    if not cp_dir.is_dir():
        return None

    pw   = 32          # patch size
    gap  = 2           # gap between patches inside a cluster panel
    cell = 4*pw + 5*gap  # 138 px: size of one 4×4 cluster panel
    lbl  = 14          # height reserved for cluster ID label below each panel
    mg   = 4           # margin between cells in the mosaic

    ncols, nrows = 5, 2
    W = ncols * cell + (ncols - 1) * mg
    H = nrows * (cell + lbl) + (nrows - 1) * mg
    canvas = Image.new("L", (W, H), color=30)   # dark background
    draw   = ImageDraw.Draw(canvas)

    for k in range(n_clusters):
        strip = _load_cluster_strip(cp_dir, k)
        if strip is None:
            continue
        grid_arr = _strip_to_4x4(strip)
        grid_img = Image.fromarray(grid_arr)

        col = k % ncols
        row = k // ncols
        x = col * (cell + mg)
        y = row * (cell + lbl + mg)
        canvas.paste(grid_img, (x, y))

        # cluster ID label centred below the panel
        draw.text((x + cell // 2, y + cell + 1), f"C{k}",
                  fill=200, anchor="mt")

    buf = io.BytesIO()
    canvas.save(buf, format="PNG")
    return buf.getvalue()

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_DARK)
    _txt(slide, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.8),
         title, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(2),
             subtitle, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

def slide_section(prs, title, detail=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_MID)
    _txt(slide, Inches(1), Inches(2.4), SW - Inches(2), Inches(1.6),
         title, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if detail:
        _txt(slide, Inches(1.5), Inches(4.2), SW - Inches(3), Inches(2),
             detail, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

def slide_3img(prs, title, imgs, caps=None):
    slide = _blank(prs)
    _title_bar(slide, title)
    caps  = (list(caps) + [""]*3)[:3] if caps else [""]*3
    imgs  = (list(imgs) + [None]*3)[:3]
    top   = TITLE_H + PAD
    cap_h = Inches(0.28) if any(caps) else PAD
    h     = SH - top - cap_h - PAD
    w     = (SW - 4*PAD) / 3
    for i, (im, cap) in enumerate(zip(imgs, caps)):
        l = PAD + i * (w + PAD)
        if im:
            _img(slide, im, l, top, w, h)
        else:
            _txt(slide, l, top + h/2 - Inches(0.25), w, Inches(0.5),
                 "[not available]", size=10, color=C_GREY, align=PP_ALIGN.CENTER)
        if cap:
            _txt(slide, l, SH - cap_h, w, cap_h,
                 cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def _style_cell(cell, text, size=8, bold=False, align=PP_ALIGN.CENTER,
                bg=None, fg=C_BLACK):
    cell.text = str(text)
    tf = cell.text_frame
    tf.paragraphs[0].alignment = align
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

def slide_table(prs, title, headers, rows, font_size=10):
    slide = _blank(prs)
    _title_bar(slide, title)
    nr, nc = len(rows) + 1, len(headers)
    row_h  = min(Inches(0.40), (SH - TITLE_H - 2*PAD) / nr)
    tbl_h  = row_h * nr
    tbl_w  = SW - 2*PAD
    top    = TITLE_H + PAD + max(0, (SH - TITLE_H - 2*PAD - tbl_h) / 2)
    tbl    = slide.shapes.add_table(nr, nc, PAD, top, tbl_w, tbl_h).table
    for j, h in enumerate(headers):
        _style_cell(tbl.cell(0, j), h, size=font_size, bold=True, bg=C_DARK, fg=C_WHITE)
    for i, row in enumerate(rows):
        bg = C_ALT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            _style_cell(tbl.cell(i+1, j), str(val), size=font_size, bg=bg)

# ── per-dataset comparison block ───────────────────────────────────────────────

SHORT_CAPS = ["A: no clip", "B: clip[0,1]", "C: clip[0,2]÷2"]

def _get(combo, *rel):
    paths = [RUNS / dn / combo / Path(*rel) for _, dn in GROUPS]
    return [p if p.exists() else None for p in paths]

def combo_slides(prs, combo: str):
    slide_section(prs, combo,
                  "A: no clip (÷2, nL1, λ=0.25)   |   "
                  "B: clip[0,1] (L1, λ=0.10)   |   "
                  "C: clip[0,2]÷2 (L1, λ=0.10)")

    # UMAP z_proj FA annotation
    slide_3img(prs, f"{combo} — UMAP z_proj (FA annotation)",
               _get(combo, "eval", "umap_proj_4ds_annotation.png"), SHORT_CAPS)

    # UMAP K-means cluster IDs (10 clusters)
    slide_3img(prs, f"{combo} — UMAP z_proj (K-means 10 clusters)",
               _get(combo, "eval", "cluster_panels_combo", "umap_combo_kmeans_k10.png"),
               SHORT_CAPS)

    # All 10 clusters on one slide (2×5 mosaic of 4×4 patch grids)
    print(f"    building cluster mosaic for {combo} ...")
    mosaics = [_all_clusters_mosaic(combo, dn) for _, dn in GROUPS]
    slide_3img(prs, f"{combo} — All 10 Cluster Panels (4×4 patches each)",
               mosaics, SHORT_CAPS)

# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs,
        "ConAE Input Normalisation — Single-Dataset Training",
        "no clip (÷2, nL1)  vs  clip[0,1] (L1)  vs  clip[0,2]÷2 (L1)\n"
        "datasets: vinc  |  nih3t3  |  ppax  |  pfak\n\n"
        "2026-07-27")

    slide_table(prs, "Group Configurations",
                ["Group", "Parent dir", "Clip", "Scale", "Recon loss", "λ_contrast"],
                [["A", "ds_combo_enlcrop_sc2",          "none",    "÷2 (sc2)", "nL1", "0.25"],
                 ["B", "ds_combo_enlcrop_clip01_l1",    "[0,1]",   "—",        "L1",  "0.10"],
                 ["C", "ds_combo_enlcrop_sc2_clip02_l1","[0,2]÷2", "÷2 (sc2)", "L1",  "0.10"]])

    for ds in SINGLE_DS:
        print(f"  {ds} ...")
        combo_slides(prs, ds)

    prs.save(str(OUT))
    n = len(prs.slides)
    print(f"\nSaved: {OUT}  ({n} slides)")
    print(f"Rsync to laptop:  rsync -avh --progress {OUT} lding@lding-Precision-3680:/home/lding/Desktop/")

if __name__ == "__main__":
    main()
