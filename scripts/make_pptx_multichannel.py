#!/usr/bin/env python3
"""
make_pptx_multichannel.py
=========================
PPT covering all multichannel ConAE models, organised by channel combo then dim/lambda.

  2ch  pax+act        rb   lat12        λ=0.5 / 0.25
  3ch  pax+zyx+act    CIO  lat18/24     λ=0.5 / 100 / 0.0001
  3ch  pax+zyx+act    rb   lat12/24     λ=0.5 / 0.25
  4ch  vinc+pax+zyx+act  CIO  lat18/24  λ=0.5 / 100 / 0.0001
  4ch  vinc+pax+zyx+act  rb   lat12/32  λ=0.5 / 0.25

Run from the SubCellAE repo root:
  python scripts/make_pptx_multichannel.py
"""
from __future__ import annotations

import io
from pathlib import Path

import numpy as np
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
OUT  = Path("multichannel_results.pptx")

# ── model list ─────────────────────────────────────────────────────────────────
# (label, dirname, ch, norm, lat, proj, lambda, note)

MODELS = [
    # ── 2ch pax+act ────────────────────────────────────────────────────────────
    ("2ch  lat12  λ=0.5",
     "contrastive_cio_rb_vinc_lat12proj8_sc2_nl1_2ch_pax_act",
     2, "rb", 12, 8, 0.5, "pax+act  no enlcrop"),
    ("2ch  lat12  λ=0.25",
     "contrastive_cio_rb_vinc_lat12proj8_sc2_nl1_lc025_2ch_pax_act",
     2, "rb", 12, 8, 0.25, "pax+act  no enlcrop"),

    # ── 3ch pax+zyx+act  CIO ───────────────────────────────────────────────────
    ("3ch (pza)  CIO  lat18  λ=0.5",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza",
     3, "CIO", 18, 12, 0.5, "pax+zyx+act"),
    ("3ch (pza)  CIO  lat18  λ=100",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza_lc100",
     3, "CIO", 18, 12, 100, "pax+zyx+act"),
    ("3ch (pza)  CIO  lat18  λ=0.0001",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza_lc1e4",
     3, "CIO", 18, 12, 0.0001, "pax+zyx+act"),
    ("3ch (pza)  CIO  lat24  λ=0.5",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza",
     3, "CIO", 24, 16, 0.5, "pax+zyx+act"),
    ("3ch (pza)  CIO  lat24  λ=100",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza_lc100",
     3, "CIO", 24, 16, 100, "pax+zyx+act"),
    ("3ch (pza)  CIO  lat24  λ=0.0001",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza_lc1e4",
     3, "CIO", 24, 16, 0.0001, "pax+zyx+act"),

    # ── 3ch pax+zyx+act  rb ────────────────────────────────────────────────────
    ("3ch (pza)  rb  lat12  λ=0.5",
     "contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_3ch_pza",
     3, "rb", 12, 8, 0.5, "pax+zyx+act"),
    ("3ch (pza)  rb  lat12  λ=0.25",
     "contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_lc025_3ch_pza",
     3, "rb", 12, 8, 0.25, "pax+zyx+act"),
    ("3ch (pza)  rb  lat24  λ=0.5",
     "contrastive_cio_rb_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza",
     3, "rb", 24, 16, 0.5, "pax+zyx+act"),
    ("3ch (pza)  rb  lat24  λ=0.25",
     "contrastive_cio_rb_vinc_lat24proj16_enlcrop_sc2_nl1_lc025_3ch_pza",
     3, "rb", 24, 16, 0.25, "pax+zyx+act"),

    # ── 4ch vinc+pax+zyx+act  CIO ─────────────────────────────────────────────
    ("4ch (vinc)  CIO  lat18  λ=0.5",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc",
     4, "CIO", 18, 12, 0.5, "vinc+pax+zyx+act"),
    ("4ch (vinc)  CIO  lat18  λ=100",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc_lc100",
     4, "CIO", 18, 12, 100, "vinc+pax+zyx+act"),
    ("4ch (vinc)  CIO  lat18  λ=0.0001",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc_lc1e4",
     4, "CIO", 18, 12, 0.0001, "vinc+pax+zyx+act"),
    ("4ch (vinc)  CIO  lat24  λ=0.5",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc",
     4, "CIO", 24, 16, 0.5, "vinc+pax+zyx+act"),
    ("4ch (vinc)  CIO  lat24  λ=100",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc_lc100",
     4, "CIO", 24, 16, 100, "vinc+pax+zyx+act"),
    ("4ch (vinc)  CIO  lat24  λ=0.0001",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc_lc1e4",
     4, "CIO", 24, 16, 0.0001, "vinc+pax+zyx+act"),

    # ── 4ch vinc+pax+zyx+act  rb ──────────────────────────────────────────────
    ("4ch (vinc)  rb  lat12  λ=0.5",
     "contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_4ch_vinc",
     4, "rb", 12, 8, 0.5, "vinc+pax+zyx+act"),
    ("4ch (vinc)  rb  lat12  λ=0.25",
     "contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_lc025_4ch_vinc",
     4, "rb", 12, 8, 0.25, "vinc+pax+zyx+act"),
    ("4ch (vinc)  rb  lat32  λ=0.5",
     "contrastive_cio_rb_vinc_lat32proj16_enlcrop_sc2_nl1_4ch_vinc",
     4, "rb", 32, 16, 0.5, "vinc+pax+zyx+act"),
    ("4ch (vinc)  rb  lat32  λ=0.25",
     "contrastive_cio_rb_vinc_lat32proj16_enlcrop_sc2_nl1_lc025_4ch_vinc",
     4, "rb", 32, 16, 0.25, "vinc+pax+zyx+act"),
]

SECTIONS = [
    ("2ch  pax+act  (rb, lat12)",
     "2-channel: paxillin + actin  |  rb normalisation  |  nL1 recon",
     lambda m: m[2] == 2),
    ("3ch  pax+zyx+act  (CIO, lat18/24, λ sweep)",
     "3-channel: paxillin + z-yx + actin  |  CIO normalisation  |  nL1 recon",
     lambda m: m[2] == 3 and m[3] == "CIO"),
    ("3ch  pax+zyx+act  (rb, lat12/24)",
     "3-channel: paxillin + z-yx + actin  |  rb normalisation  |  nL1 recon",
     lambda m: m[2] == 3 and m[3] == "rb"),
    ("4ch  vinc+pax+zyx+act  (CIO, lat18/24, λ sweep)",
     "4-channel: vinculin + paxillin + z-yx + actin  |  CIO normalisation  |  nL1 recon",
     lambda m: m[2] == 4 and m[3] == "CIO"),
    ("4ch  vinc+pax+zyx+act  (rb, lat12/32)",
     "4-channel: vinculin + paxillin + z-yx + actin  |  rb normalisation  |  nL1 recon",
     lambda m: m[2] == 4 and m[3] == "rb"),
]

# ── slide geometry ─────────────────────────────────────────────────────────────

SW = Inches(13.33)
SH = Inches(7.5)
TITLE_H = Inches(0.65)
PAD     = Inches(0.12)

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_ALT   = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)

# ── low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=14, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _img(slide, src, l, t, max_w, max_h):
    if src is None:
        return
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return
        data = p.read_bytes()
    else:
        data = bytes(src)
    try:
        im = Image.open(io.BytesIO(data))
        iw, ih = im.size
    except Exception:
        return
    scale = min(max_w / iw, max_h / ih)
    w, h  = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(io.BytesIO(data),
                             l + (max_w - w) // 2,
                             t + (max_h - h) // 2, w, h)

def _title_bar(slide, text, size=16, bg=C_DARK):
    _rect(slide, 0, 0, SW, TITLE_H, fill=bg)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         text, size=size, bold=True, color=C_WHITE)

def _tif_to_png(path: Path) -> bytes | None:
    try:
        arr = tifffile.imread(str(path)).astype(np.float32)
    except Exception:
        return None
    if arr.ndim == 2:
        img2d = arr
    elif arr.ndim == 3 and arr.shape[2] in (3, 4):
        img2d = arr
    elif arr.ndim == 3:
        K, H, W = arr.shape
        ncols = min(K, 5)
        nrows = (K + ncols - 1) // ncols
        canvas = np.zeros((nrows * H, ncols * W), dtype=np.float32)
        for i in range(K):
            r, c = divmod(i, ncols)
            canvas[r*H:(r+1)*H, c*W:(c+1)*W] = arr[i]
        img2d = canvas
    else:
        img2d = arr[0]
    vmin, vmax = img2d.min(), img2d.max()
    if vmax > vmin:
        img2d = ((img2d - vmin) / (vmax - vmin) * 255).clip(0, 255).astype(np.uint8)
    else:
        img2d = np.zeros_like(img2d, dtype=np.uint8)
    buf = io.BytesIO()
    Image.fromarray(img2d).save(buf, format='PNG')
    return buf.getvalue()

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_DARK)
    _txt(slide, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.8),
         title, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(2),
             subtitle, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

def slide_section(prs, title, detail=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_MID)
    _txt(slide, Inches(1), Inches(2.4), SW - Inches(2), Inches(1.6),
         title, size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if detail:
        _txt(slide, Inches(1.5), Inches(4.2), SW - Inches(3), Inches(2),
             detail, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

def slide_1img(prs, title, img, caption=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top = TITLE_H + PAD
    avail_h = SH - top - (Inches(0.32) if caption else PAD)
    _img(slide, img, PAD, top, SW - 2*PAD, avail_h)
    if caption:
        _txt(slide, PAD, SH - Inches(0.32), SW - 2*PAD, Inches(0.28),
             caption, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def slide_2img(prs, title, img1, img2, cap1="", cap2=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top   = TITLE_H + PAD
    cap_h = Inches(0.28) if (cap1 or cap2) else PAD
    h = SH - top - cap_h - PAD
    w = (SW - 3*PAD) / 2
    _img(slide, img1, PAD,           top, w, h)
    _img(slide, img2, PAD + w + PAD, top, w, h)
    if cap1:
        _txt(slide, PAD, SH - cap_h, w, cap_h,
             cap1, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
    if cap2:
        _txt(slide, PAD + w + PAD, SH - cap_h, w, cap_h,
             cap2, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def slide_3img(prs, title, imgs, caps=None):
    slide = _blank(prs)
    _title_bar(slide, title)
    caps  = (list(caps) + [""]*3)[:3] if caps else [""]*3
    imgs  = (list(imgs) + [None]*3)[:3]
    top   = TITLE_H + PAD
    cap_h = Inches(0.28) if any(caps) else PAD
    h     = SH - top - cap_h - PAD
    w     = (SW - 4*PAD) / 3
    for i, (im, cap) in enumerate(zip(imgs, caps)):
        l = PAD + i * (w + PAD)
        if im:
            _img(slide, im, l, top, w, h)
        else:
            _txt(slide, l, top + h/2 - Inches(0.25), w, Inches(0.5),
                 "[not available]", size=10, color=C_GREY, align=PP_ALIGN.CENTER)
        if cap:
            _txt(slide, l, SH - cap_h, w, cap_h,
                 cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def _style_cell(cell, text, size=8, bold=False, align=PP_ALIGN.CENTER,
                bg=None, fg=C_BLACK):
    cell.text = str(text)
    tf = cell.text_frame
    tf.paragraphs[0].alignment = align
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

def slide_table(prs, title, headers, rows, font_size=9):
    slide = _blank(prs)
    _title_bar(slide, title)
    if not rows:
        return
    nr, nc = len(rows) + 1, len(headers)
    row_h  = min(Inches(0.33), (SH - TITLE_H - 2*PAD) / nr)
    tbl_h  = row_h * nr
    tbl_w  = SW - 2*PAD
    top    = TITLE_H + PAD + max(0, (SH - TITLE_H - 2*PAD - tbl_h) / 2)
    tbl    = slide.shapes.add_table(nr, nc, PAD, top, tbl_w, tbl_h).table
    for j, h in enumerate(headers):
        _style_cell(tbl.cell(0, j), h, size=font_size, bold=True, bg=C_DARK, fg=C_WHITE)
    for i, row in enumerate(rows):
        bg = C_ALT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            _style_cell(tbl.cell(i+1, j), str(val), size=font_size, bg=bg)

# ── cluster panel helper ───────────────────────────────────────────────────────

def _cluster_panel_png(d: Path) -> bytes | None:
    """Return PNG of all_clusters.tif from any cluster_panels* subdir."""
    for cp_name in ("cluster_panels", "cluster_panels_proj"):
        cp_dir = d / "eval" / cp_name
        if not cp_dir.is_dir():
            cp_dir = d / cp_name          # old-style path
        if not cp_dir.is_dir():
            continue
        tif = cp_dir / "all_clusters.tif"
        if not tif.exists():
            tif = next(iter(sorted(cp_dir.glob("all*.tif"))), None)
        if tif:
            return _tif_to_png(tif)
    return None

# ── per-model slides ───────────────────────────────────────────────────────────

def _model_slides(prs, label: str, d: Path):
    if not d.exists():
        slide_section(prs, f"{label}  [MISSING]",
                      f"Directory not found: {d.name}")
        return

    ev    = d / "eval"
    pfx   = "contrastive"
    has_h5 = (d / "model.h5").exists()

    # 1. Training loss + component losses
    loss = d / f"{pfx}_train_val_loss.png"
    comp = d / f"{pfx}_component_losses.png"
    if loss.exists() and comp.exists():
        slide_2img(prs, f"{label} — Training Loss",
                   loss, comp, "Train/Val total loss", "Component losses")
    elif loss.exists():
        slide_1img(prs, f"{label} — Training Loss", loss)

    # 2. ep500 views + reconstruction
    views = d / f"{pfx}_views_ep500.png"
    recon = d / f"{pfx}_recon_ep500.png"
    if views.exists() and recon.exists():
        slide_2img(prs, f"{label} — Epoch 500",
                   views, recon, "Contrastive views", "Reconstruction")
    elif recon.exists():
        slide_1img(prs, f"{label} — Epoch 500 Recon", recon)

    if not ev.is_dir():
        return   # no eval — training slides only

    # ── new-style eval (3ch/4ch CIO) ──────────────────────────────────────────

    # 3. Cross-dataset recon (overall + per-channel)
    cross_nl1 = d / "cross_dataset_recon_nl1.png"
    cross_l1  = d / "cross_dataset_recon_l1.png"
    cross_mse = d / "cross_dataset_recon_mse.png"
    if cross_nl1.exists() or cross_l1.exists():
        slide_3img(prs, f"{label} — Cross-Dataset Recon Metrics",
                   [cross_nl1 if cross_nl1.exists() else None,
                    cross_l1  if cross_l1.exists()  else None,
                    cross_mse if cross_mse.exists() else None],
                   ["nL1 (normalised)", "L1 (absolute)", "MSE"])
    # per-channel (2 shown)
    nl1_ch = [d / f"cross_dataset_recon_nl1_ch{c}.png" for c in range(4)]
    l1_ch  = [d / f"cross_dataset_recon_l1_ch{c}.png"  for c in range(4)]
    for c in range(4):
        if nl1_ch[c].exists() and l1_ch[c].exists():
            slide_2img(prs, f"{label} — Cross-Dataset Recon ch{c}",
                       nl1_ch[c], l1_ch[c], f"nL1 ch{c}", f"L1 ch{c}")

    # 4. UMAP z_proj 4ds (new-style)
    p4_ann  = ev / "umap_proj_4ds_annotation.png"
    p4_cond = ev / "umap_proj_4ds_condition.png"
    p4_ds   = ev / "umap_proj_4ds_dataset.png"
    p4_km   = ev / "umap_proj_4ds_kmeans.png"
    if p4_ann.exists():
        slide_3img(prs, f"{label} — UMAP z_proj (cross-dataset)",
                   [p4_ann,
                    p4_cond if p4_cond.exists() else None,
                    p4_ds   if p4_ds.exists()   else None],
                   ["FA annotation", "Condition", "Dataset"])
    if p4_km.exists():
        slide_2img(prs, f"{label} — UMAP z_proj: K-means & condition",
                   p4_km,
                   p4_cond if p4_cond.exists() else None,
                   "K-means cluster ID", "Condition")

    # 5. PHATE z_proj 4ds (new-style)
    pp_ann  = ev / "phate_proj_4ds_annotation.png"
    pp_cond = ev / "phate_proj_4ds_condition.png"
    if pp_ann.exists():
        slide_2img(prs, f"{label} — PHATE z_proj (cross-dataset)",
                   pp_ann,
                   pp_cond if pp_cond.exists() else None,
                   "FA annotation", "Condition")

    # 6. Old-style eval (2ch): local UMAP/PHATE + classification
    u_ann  = ev / "umap_annotation.png"
    u_cond = ev / "umap_condition.png"
    u_km   = ev / "umap_kmeans.png"
    if u_ann.exists() and not p4_ann.exists():  # only if new-style absent
        slide_3img(prs, f"{label} — UMAP (training data)",
                   [u_ann,
                    u_cond if u_cond.exists() else None,
                    u_km   if u_km.exists()   else None],
                   ["FA annotation", "Condition", "K-means"])

    up_ann  = ev / "umap_proj_annotation.png"
    up_cond = ev / "umap_proj_condition.png"
    if up_ann.exists():
        slide_2img(prs, f"{label} — UMAP z_proj (training data)",
                   up_ann,
                   up_cond if up_cond.exists() else None,
                   "FA annotation", "Condition")

    ph_ann  = ev / "phate_annotation.png"
    ph_cond = ev / "phate_condition.png"
    if ph_ann.exists():
        slide_2img(prs, f"{label} — PHATE (training data)",
                   ph_ann,
                   ph_cond if ph_cond.exists() else None,
                   "FA annotation", "Condition")

    # 7. Confusion matrix (both eval styles)
    conf_v  = ev / "confusion_vinc_val_norm.png"
    conf_pp = ev / "confusion_ppax_norm.png"
    if conf_v.exists() and conf_pp.exists():
        slide_2img(prs, f"{label} — FA Classification",
                   conf_v, conf_pp, "vinc val", "ppax test")
    elif conf_v.exists():
        slide_1img(prs, f"{label} — FA Classification (vinc val)", conf_v)

    # 8. Cluster panels
    cp_png = _cluster_panel_png(d)
    if cp_png:
        slide_1img(prs, f"{label} — K-means Cluster Panels", cp_png)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs,
        "ConAE Multichannel Models",
        "2ch (pax+act)  |  3ch (pax+zyx+act)  |  4ch (vinc+pax+zyx+act)\n"
        "dims: lat12 / lat18 / lat24 / lat32   "
        "norm: CIO / rb   "
        "λ sweep where available\n\n"
        "2026-07-27")

    # Overview table
    hdrs = ["Label", "Ch", "Channels", "Norm", "Lat", "Proj", "λ", "model.h5", "eval"]
    rows = []
    for lbl, dn, ch, norm, lat, proj, lc, note in MODELS:
        d  = RUNS / dn
        h5 = "✓" if (d / "model.h5").exists() else "✗"
        ev = "✓" if (d / "eval").is_dir() else "✗"
        rows.append([lbl, str(ch), note, norm, str(lat), str(proj), str(lc), h5, ev])
    slide_table(prs, "Model Overview", hdrs, rows, font_size=8)

    # Per-section deep-dive
    for sec_title, sec_desc, predicate in SECTIONS:
        sec_models = [(lbl, dn, ch, norm, lat, proj, lc, note)
                      for lbl, dn, ch, norm, lat, proj, lc, note in MODELS
                      if predicate((lbl, dn, ch, norm, lat, proj, lc, note))]
        if not sec_models:
            continue

        slide_section(prs, sec_title, sec_desc)

        # Section table
        hdrs2 = ["Label", "Lat", "Proj", "λ", "model.h5", "eval"]
        rows2 = []
        for lbl, dn, ch, norm, lat, proj, lc, note in sec_models:
            d  = RUNS / dn
            h5 = "✓" if (d / "model.h5").exists() else "✗"
            ev = "✓" if (d / "eval").is_dir() else "✗"
            rows2.append([lbl, str(lat), str(proj), str(lc), h5, ev])
        slide_table(prs, f"{sec_title} — Configurations", hdrs2, rows2, font_size=10)

        for lbl, dn, ch, norm, lat, proj, lc, note in sec_models:
            d = RUNS / dn
            h5 = "✓" if (d / "model.h5").exists() else "✗ no h5"
            ev = "✓ eval" if (d / "eval").is_dir() else "no eval"
            print(f"  {lbl} ({h5}, {ev})")
            slide_section(prs, lbl,
                          f"{dn}\n{note}  |  lat{lat} proj{proj}  |  λ={lc}  |  {norm}  "
                          f"|  {'model.h5 ✓' if (d/'model.h5').exists() else 'model.h5 missing'}  "
                          f"|  {'eval ✓' if (d/'eval').is_dir() else 'no eval'}")
            _model_slides(prs, lbl, d)

    prs.save(str(OUT))
    n = len(prs.slides)
    print(f"\nSaved: {OUT}  ({n} slides)")
    print(f"Rsync to laptop:  rsync -avh --progress {OUT} lding@lding-Precision-3680:/home/lding/Desktop/")


if __name__ == "__main__":
    main()
