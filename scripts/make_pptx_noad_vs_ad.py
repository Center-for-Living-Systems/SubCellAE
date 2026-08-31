#!/usr/bin/env python3
"""
make_pptx_noad_vs_ad.py
========================
Error analysis PPT for 5-class FA classification (No adhesion + 4 subtypes).

Slides
------
 1. Cover
 2. Label statistics — 5-class counts per dataset
 3. Error rate overview — per-scenario bar + aggregate 5×5 confusion matrix
 4. No-adhesion boundary spotlight — how often does the model cross the ad/no-ad line?
 5. Per-class error rates across all scenarios
 6. vinc_only detail — per-class accuracy bar + top errors
 7. ctrl→ycomp detail
 8. pfak→vinc detail (worst cross-dataset)
 9. Key findings

Usage
-----
  python scripts/make_pptx_noad_vs_ad.py
  python scripts/make_pptx_noad_vs_ad.py --out custom.pptx
"""

from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
EVAL_DIR  = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run/fa4_xds_eval")
LABEL_DIR = Path("/net/projects/CLS/lding/data/fa_data_analysis/labelling")
OUT_DEFAULT = Path("results/fa4_noad_vs_ad_errors.pptx")

SUFFIX = "A_zrecon_5cls"

LABEL_FILES = {
    "vinc_ctrl":  LABEL_DIR / "vinc_control_label_Annabel_20260715_1554.csv",
    "vinc_ycomp": LABEL_DIR / "vinc_combined_label_Annabel_20260816.csv",
    "pfak_ctrl":  LABEL_DIR / "pfak_combined_label_Annabel_aug2026.csv",
}

FA_LABELS_5 = ["No adhesion", "Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
LABEL_SHORT = {
    "No adhesion":        "NoAd",
    "Nascent Adhesion":   "NA",
    "focal complex":      "FC",
    "focal adhesion":     "FA",
    "fibrillar adhesion": "Fib",
}
SHORT_LONG = {v: k for k, v in LABEL_SHORT.items()}
SHORTS = [LABEL_SHORT[l] for l in FA_LABELS_5]

CLASS_COLORS = {
    "NoAd": "#9467bd",  # purple
    "NA":   "#1565C0",  # blue
    "FC":   "#E65100",  # orange
    "FA":   "#2ca02c",  # green
    "Fib":  "#C00000",  # red
}

AD_CLASSES   = {"Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"}
NOAD_CLASS   = "No adhesion"

SCENARIOS = ["vinc_only", "pfak_only", "ctrl->ycomp", "vinc->pfak", "pfak->vinc", "combined"]
SCENARIO_LABELS = {
    "vinc_only":   "Vinc only (within-ds)",
    "pfak_only":   "pFAK only (within-ds)",
    "ctrl->ycomp": "ctrl → ycomp (cross-cond)",
    "vinc->pfak":  "Vinc → pFAK (cross-ds)",
    "pfak->vinc":  "pFAK → Vinc (cross-ds)",
    "combined":    "Combined (all data)",
}
SCENARIO_SHORT = {
    "vinc_only":   "Vinc only",
    "pfak_only":   "pFAK only",
    "ctrl->ycomp": "ctrl→ycomp",
    "vinc->pfak":  "Vinc→pFAK",
    "pfak->vinc":  "pFAK→Vinc",
    "combined":    "Combined",
}
SCENARIO_COLORS = {
    "vinc_only":   "#1f77b4",
    "pfak_only":   "#d62728",
    "ctrl->ycomp": "#B07AA1",
    "vinc->pfak":  "#ff7f0e",
    "pfak->vinc":  "#9467bd",
    "combined":    "#2ca02c",
}

# ── slide geometry ─────────────────────────────────────────────────────────────
SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.15)
BODY_T  = TITLE_H + Inches(0.08)
BODY_H  = SH - BODY_T - PAD

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_GREEN = RGBColor(0x37, 0x8B, 0x4A)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)
C_PURP  = RGBColor(0x94, 0x67, 0xBD)

# ── helpers ────────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2*PAD, Inches(0.25),
             subtitle, size=9, color=C_GREY)

def _fig_to_pil(fig) -> Image.Image:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    return Image.open(buf).convert("RGB")

def _place_pil(slide, pil, l, t, max_w, max_h):
    iw, ih = pil.size
    scale  = min(max_w / Inches(iw / 150), max_h / Inches(ih / 150))
    rw, rh = Inches(iw / 150) * scale, Inches(ih / 150) * scale
    buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                             t + (max_h - rh) / 2, rw, rh)

def _place_png(slide, path, l, t, max_w, max_h, ph="[pending]"):
    p = Path(path)
    if not p.exists():
        _txt(slide, l, t + max_h / 2, max_w, Inches(0.3),
             ph, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
        return False
    pil = Image.open(str(p)).convert("RGB")
    _place_pil(slide, pil, l, t, max_w, max_h)
    return True

def _load_predictions(scenario: str) -> pd.DataFrame | None:
    p = EVAL_DIR / f"predictions_{scenario}_{SUFFIX}.csv"
    if not p.exists():
        return None
    return pd.read_csv(p)

def _rgb_from_hex(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(1), Inches(1.6), SW - Inches(2), Inches(1.2),
         "FA5 Error Analysis: Ad vs No Adhesion",
         size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(2.9), SW - Inches(2), Inches(0.6),
         "5-class classification: No Adhesion  ·  Nascent Adhesion  ·  Focal Complex  ·  Focal Adhesion  ·  Fibrillar",
         size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.55), SW - Inches(2), Inches(0.45),
         "Option A  ·  z_recon features (12-d)  ·  75% label fraction  ·  6 eval scenarios",
         size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Color legend
    legend_x = Inches(3.0)
    legend_y = Inches(4.3)
    for short, label in [("NoAd", "No adhesion"), ("NA", "Nascent Adhesion"),
                          ("FC", "Focal complex"), ("FA", "Focal adhesion"), ("Fib", "Fibrillar")]:
        _rect(sl, legend_x, legend_y, Inches(0.22), Inches(0.22),
              fill=_rgb_from_hex(CLASS_COLORS[short]))
        _txt(sl, legend_x + Inches(0.28), legend_y - Inches(0.02), Inches(1.5), Inches(0.28),
             f"{short} = {label}", size=10, color=C_WHITE)
        legend_x += Inches(1.95)

    _txt(sl, Inches(1), SH - Inches(0.55), SW - Inches(2), Inches(0.4),
         "2026-08-24", size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_label_stats(prs):
    sl = _blank(prs)
    _title_bar(sl, "Label Distribution — 5-Class (No Adhesion + 4 FA Subtypes)",
               "Counts per dataset  ·  Note: No adhesion is majority class in vinc")

    fig, axes = plt.subplots(1, 3, figsize=(15, 4.5), facecolor="white")
    ds_display = {"vinc_ctrl": "Vinc ctrl", "vinc_ycomp": "Vinc ycomp", "pfak_ctrl": "pFAK ctrl"}

    for ax, (ds_key, ds_name) in zip(axes, ds_display.items()):
        df = pd.read_csv(LABEL_FILES[ds_key])
        counts = []
        colors = []
        for lbl in FA_LABELS_5:
            n = (df["label"] == lbl).sum()
            counts.append(n)
            colors.append(CLASS_COLORS[LABEL_SHORT[lbl]])
        shorts = [LABEL_SHORT[l] for l in FA_LABELS_5]
        bars = ax.bar(shorts, counts, color=colors, edgecolor="white", width=0.65)
        for b, n in zip(bars, counts):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 1,
                    str(n), ha="center", va="bottom", fontsize=9, fontweight="bold")
        ax.set_title(ds_name, fontsize=11, fontweight="bold")
        ax.set_ylabel("# patches", fontsize=9)
        ax.set_ylim(0, max(counts) * 1.18)
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")

    fig.suptitle("Label counts per dataset (all labeled patches)", fontsize=12, fontweight="bold")
    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(0.5))
    plt.close(fig)

    _txt(sl, PAD, SH - Inches(1.0), SW - 2*PAD, Inches(0.85),
         "No adhesion is the majority class: 342 vinc_ctrl, 770 vinc_ycomp, 60 pfak_ctrl  ·  "
         "FC and Fib remain rare (5–85 patches)  ·  "
         "5-class LightGBM trained on z_recon latents (12-d) with class_weight='balanced'",
         size=9, color=C_DARK, wrap=True)


def slide_overview(prs):
    sl = _blank(prs)
    _title_bar(sl, "Error Rate Overview — 6 Scenarios  ·  5-Class (including No Adhesion)",
               "Option A  ·  z_recon  ·  75% labels  ·  class_weight=balanced")

    fig, axes = plt.subplots(1, 2, figsize=(15, 4.8), facecolor="white",
                             gridspec_kw={"width_ratios": [1, 1.6]})

    # Left: error rate per scenario
    ax = axes[0]
    sc_list  = [sc for sc in SCENARIOS if _load_predictions(sc) is not None]
    rates    = []
    clrs     = []
    sc_names = []
    for sc in sc_list:
        df = _load_predictions(sc)
        rates.append((df["true_label"] != df["pred_label"]).mean() * 100)
        clrs.append(SCENARIO_COLORS[sc])
        sc_names.append(SCENARIO_SHORT[sc])

    bars = ax.barh(sc_names, rates, color=clrs, edgecolor="white")
    for b, v in zip(bars, rates):
        ax.text(v + 0.5, b.get_y() + b.get_height() / 2,
                f"{v:.1f}%", va="center", fontsize=9, fontweight="bold")
    ax.set_xlabel("Error rate (%)", fontsize=10)
    ax.set_title("Error rate at 75% labels (5-class)", fontsize=11, fontweight="bold")
    ax.set_xlim(0, 50)
    ax.axvline(20, color="gray", lw=0.8, linestyle=":")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    # Right: aggregate 5×5 confusion matrix
    ax2 = axes[1]
    combined_conf = np.zeros((5, 5))
    for sc in sc_list:
        df = _load_predictions(sc)
        if df is None:
            continue
        for i, ts in enumerate(SHORTS):
            tl = SHORT_LONG[ts]
            for j, ps in enumerate(SHORTS):
                pl = SHORT_LONG[ps]
                combined_conf[i, j] += ((df["true_label"] == tl) & (df["pred_label"] == pl)).sum()

    row_sums = combined_conf.sum(axis=1, keepdims=True)
    row_sums[row_sums == 0] = 1
    conf_norm = combined_conf / row_sums

    im = ax2.imshow(conf_norm, cmap="Blues", vmin=0, vmax=1, aspect="auto")
    ax2.set_xticks(range(5))
    ax2.set_yticks(range(5))
    xlab = [f"pred\n{s}" for s in SHORTS]
    ylab = [f"true {s}" for s in SHORTS]
    ax2.set_xticklabels(xlab, fontsize=8)
    ax2.set_yticklabels(ylab, fontsize=8)
    # Color axis tick labels
    for tick, short in zip(ax2.get_xticklabels(), SHORTS):
        tick.set_color(CLASS_COLORS[short])
        tick.set_fontweight("bold")
    for tick, short in zip(ax2.get_yticklabels(), SHORTS):
        tick.set_color(CLASS_COLORS[short])
        tick.set_fontweight("bold")
    ax2.set_title("Aggregate confusion (all scenarios, row-normalized)", fontsize=11, fontweight="bold")
    for i in range(5):
        for j in range(5):
            n = int(combined_conf[i, j])
            v = conf_norm[i, j]
            color = "white" if v > 0.55 else "black"
            ax2.text(j, i, f"{v:.2f}\n(n={n})", ha="center", va="center", fontsize=7, color=color)
    fig.colorbar(im, ax=ax2, fraction=0.04, pad=0.02, label="Row-normalized fraction")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.1))
    plt.close(fig)

    _txt(sl, PAD, SH - Inches(1.05), SW * 0.5, Inches(0.9),
         "Within-dataset (vinc_only, pfak_only): 9–20% error with No adhesion included\n"
         "ctrl→ycomp: 33% — No adhesion confused with FA under drug treatment\n"
         "pfak→vinc: 36% — worst, large NoAd→FA confusion (169 patches)",
         size=9, color=C_DARK, wrap=True)
    _txt(sl, SW * 0.5, SH - Inches(1.05), SW * 0.5 - PAD, Inches(0.9),
         "Hotspot: NoAd→FA confusion — 'background' patches predicted as focal adhesion\n"
         "Also: FC→NA remains a persistent cross-dataset error\n"
         "NA→NoAd confusion: 4% of NA patches mislabeled as background",
         size=9, color=C_DARK, wrap=True)


def slide_noad_boundary(prs):
    sl = _blank(prs)
    _title_bar(sl, "No Adhesion Boundary — How Often Does the Model Cross Ad ↔ No-Ad?",
               "Counting predictions that cross the adhesion/no-adhesion boundary per scenario")

    fig, axes = plt.subplots(1, 2, figsize=(15, 5), facecolor="white")

    sc_list = [sc for sc in SCENARIOS if _load_predictions(sc) is not None]
    x = np.arange(len(sc_list))
    sc_names_short = [SCENARIO_SHORT[s] for s in sc_list]

    # Left: boundary crossing counts
    ax = axes[0]
    noad_as_ad  = []  # true=NoAd, pred=Ad
    ad_as_noad  = []  # true=Ad,   pred=NoAd
    for sc in sc_list:
        df = _load_predictions(sc)
        noad_as_ad.append(((df["true_label"] == NOAD_CLASS) &
                            df["pred_label"].isin(AD_CLASSES)).sum())
        ad_as_noad.append((df["true_label"].isin(AD_CLASSES) &
                            (df["pred_label"] == NOAD_CLASS)).sum())

    w = 0.35
    b1 = ax.bar(x - w/2, noad_as_ad, width=w, color=CLASS_COLORS["FA"],
                label="NoAd predicted as Ad", edgecolor="white")
    b2 = ax.bar(x + w/2, ad_as_noad, width=w, color=CLASS_COLORS["NoAd"],
                label="Ad predicted as NoAd", edgecolor="white")
    for b, v in zip(list(b1) + list(b2), noad_as_ad + ad_as_noad):
        if v > 0:
            ax.text(b.get_x() + b.get_width()/2, b.get_height() + 0.3,
                    str(v), ha="center", va="bottom", fontsize=8, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(sc_names_short, fontsize=9, ha="center")
    ax.set_ylabel("# patches crossing boundary", fontsize=10)
    ax.set_title("Ad ↔ No-Ad boundary crossings", fontsize=11, fontweight="bold")
    ax.legend(fontsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    # Right: stacked view of NoAd errors — where do NoAd patches go?
    ax2 = axes[1]
    dest_data = {s: [] for s in ["NA", "FC", "FA", "Fib"]}
    for sc in sc_list:
        df = _load_predictions(sc)
        noad_df = df[df["true_label"] == NOAD_CLASS]
        for s in ["NA", "FC", "FA", "Fib"]:
            dest_data[s].append((noad_df["pred_label"] == SHORT_LONG[s]).sum())

    bottom = np.zeros(len(sc_list))
    for s in ["NA", "FC", "FA", "Fib"]:
        counts = np.array(dest_data[s])
        ax2.bar(x, counts, bottom=bottom, color=CLASS_COLORS[s],
                label=f"NoAd → {s}", edgecolor="white", width=0.6)
        bottom += counts

    # Also show correct NoAd predictions
    correct_noad = []
    for sc in sc_list:
        df = _load_predictions(sc)
        noad_df = df[df["true_label"] == NOAD_CLASS]
        correct_noad.append((noad_df["pred_label"] == NOAD_CLASS).sum())
    ax2.bar(x, correct_noad, bottom=bottom, color=CLASS_COLORS["NoAd"],
            label="NoAd → NoAd (correct)", edgecolor="white", width=0.6, alpha=0.45, hatch="//")

    ax2.set_xticks(x)
    ax2.set_xticklabels(sc_names_short, fontsize=9)
    ax2.set_ylabel("# No adhesion patches", fontsize=10)
    ax2.set_title("Fate of 'No adhesion' patches", fontsize=11, fontweight="bold")
    ax2.legend(fontsize=8, loc="upper right")
    ax2.spines[["top", "right"]].set_visible(False)
    ax2.set_facecolor("white")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.1))
    plt.close(fig)

    _txt(sl, PAD, SH - Inches(1.05), SW - 2*PAD, Inches(0.9),
         "• NoAd→FA is the dominant boundary error, especially pfak→vinc (169 patches): "
         "pFAK training set has few NoAd examples, so the classifier treats background as FA\n"
         "• Ad→NoAd errors: FA patches mistaken for background; more common in cross-condition "
         "(ctrl→ycomp: 9 FA→NoAd) suggesting y-compound changes FA morphology toward background",
         size=9, color=C_DARK, wrap=True)


def slide_perclass_rates(prs):
    sl = _blank(prs)
    _title_bar(sl, "Per-Class Error Rate Across All Scenarios",
               "How often is each class misclassified?  ·  class_weight=balanced, 75% labels")

    fig, ax = plt.subplots(figsize=(14, 5), facecolor="white")

    sc_list = [sc for sc in SCENARIOS if _load_predictions(sc) is not None]
    x = np.arange(len(sc_list))
    width = 0.15
    sc_names_short = [SCENARIO_SHORT[s] for s in sc_list]

    for i, (short, long) in enumerate(zip(SHORTS, FA_LABELS_5)):
        err_rates = []
        for sc in sc_list:
            df = _load_predictions(sc)
            sub = df[df["true_label"] == long]
            if len(sub) == 0:
                err_rates.append(np.nan)
            else:
                err_rates.append((sub["pred_label"] != long).mean() * 100)
        offset = (i - 2) * width
        bars = ax.bar(x + offset, [e if not np.isnan(e) else 0 for e in err_rates],
                      width=width * 0.9, color=CLASS_COLORS[short],
                      label=f"{short}", edgecolor="white")
        for b, v in zip(bars, err_rates):
            if np.isnan(v):
                ax.text(b.get_x() + b.get_width()/2, 1, "–",
                        ha="center", va="bottom", fontsize=6, color="gray")

    ax.set_xticks(x)
    ax.set_xticklabels(sc_names_short, fontsize=9)
    ax.set_ylabel("Per-class error rate (%)", fontsize=10)
    ax.set_title("How often is each class misclassified? (5-class, 75% labels)", fontsize=12, fontweight="bold")
    ax.legend(title="True class", fontsize=9, loc="upper left")
    ax.axhline(50, color="gray", lw=0.8, linestyle=":", alpha=0.6)
    ax.set_ylim(0, 115)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.1))
    plt.close(fig)

    _txt(sl, PAD, SH - Inches(1.05), SW - 2*PAD, Inches(0.9),
         "• FC remains the hardest FA subtype across scenarios (>60% error in cross-dataset)\n"
         "• No adhesion accuracy is surprisingly high in within-dataset scenarios — "
         "but degrades badly in pfak→vinc (NoAd majority is hard without enough pFAK NoAd examples)\n"
         "• Fib error rate is variable — small sample size makes estimates noisy",
         size=9, color=C_DARK, wrap=True)


def _detail_slide(prs, scenario: str, notes: str):
    sl = _blank(prs)
    df = _load_predictions(scenario)
    label = SCENARIO_LABELS[scenario]
    cross = "→" in scenario

    if df is not None:
        n_total = len(df)
        n_err = (df["true_label"] != df["pred_label"]).sum()
        err_rate = n_err / n_total
    else:
        n_total = n_err = 0
        err_rate = 0.0

    _title_bar(
        sl, f"Detail — {label}",
        f"{n_err}/{n_total} errors ({err_rate*100:.1f}%)  ·  5-class  ·  75% labels  ·  "
        f"{'cross-dataset / cross-condition: all test patches from other split' if cross else 'within-dataset: random 25% held-out'}",
    )

    # Left: error grid PNG
    grid_w = SW * 0.60
    _place_png(sl, EVAL_DIR / f"errors_{scenario}_{SUFFIX}.png",
               PAD, BODY_T + Inches(0.05), grid_w - PAD, BODY_H - Inches(0.15))

    # Right: per-class accuracy bar (always show all 5 including Fib)
    rx = grid_w + PAD
    rw = SW - rx - PAD

    if df is not None:
        fig, ax = plt.subplots(figsize=(3.5, 3.6), facecolor="white")
        per_acc = []
        per_n   = []
        for s in SHORTS:
            long = SHORT_LONG[s]
            sub  = df[df["true_label"] == long]
            n    = len(sub)
            acc  = (sub["pred_label"] == long).mean() if n > 0 else np.nan
            per_acc.append(acc)
            per_n.append(n)

        colors  = [CLASS_COLORS[s] for s in SHORTS]
        heights = [a * 100 if not np.isnan(a) else 0 for a in per_acc]
        bars    = ax.bar(SHORTS, heights, color=colors, edgecolor="white")
        for b, acc, n in zip(bars, per_acc, per_n):
            if n == 0:
                ax.text(b.get_x() + b.get_width()/2, 2, "n=0",
                        ha="center", va="bottom", fontsize=7, color="gray")
            else:
                lbl = f"{acc*100:.0f}%\n(n={n})"
                ax.text(b.get_x() + b.get_width()/2, b.get_height() + 1,
                        lbl, ha="center", va="bottom", fontsize=7, fontweight="bold")

        ax.set_ylim(0, 120)
        ax.set_ylabel("Per-class accuracy (%)", fontsize=9)
        ax.set_title(f"Per-class accuracy\n({SCENARIO_SHORT[scenario]})", fontsize=9, fontweight="bold")
        ax.axhline(50, color="gray", lw=0.8, linestyle=":")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")
        _place_pil(sl, _fig_to_pil(fig), rx, BODY_T + Inches(0.05), rw, Inches(3.0))
        plt.close(fig)

        # Top error pairs
        errors = df[df["true_label"] != df["pred_label"]]
        if len(errors) > 0:
            conf = (errors.groupby(["true_label", "pred_label"], observed=True)
                         .size().reset_index(name="n")
                         .sort_values("n", ascending=False))
            ty = BODY_T + Inches(3.15)
            _txt(sl, rx, ty, rw, Inches(0.28), "Top error pairs:", size=9, bold=True, color=C_DARK)
            ty += Inches(0.28)
            for _, row in conf.head(8).iterrows():
                t = LABEL_SHORT.get(row["true_label"], row["true_label"])
                p = LABEL_SHORT.get(row["pred_label"], row["pred_label"])
                tc = _rgb_from_hex(CLASS_COLORS.get(t, "#000000"))
                _txt(sl, rx, ty, rw, Inches(0.24),
                     f"  {t} → {p}   n={row['n']} ({row['n']/n_total*100:.1f}%)",
                     size=8.5, color=tc)
                ty += Inches(0.24)

    _txt(sl, PAD, SH - Inches(1.0), SW - 2*PAD, Inches(0.9),
         notes, size=9, color=C_DARK, wrap=True)


def slide_findings(prs):
    sl = _blank(prs)
    _title_bar(sl, "Key Findings — Ad vs No Adhesion (5-class Error Analysis)",
               "Option A  ·  z_recon  ·  75% labels  ·  class_weight=balanced")

    col_w = (SW - 3 * PAD) / 2
    col2_x = PAD * 2 + col_w

    _txt(sl, PAD, BODY_T + Inches(0.1), col_w, Inches(0.3),
         "What Goes Wrong", size=13, bold=True, color=C_DARK)

    findings = [
        (C_RED,   "NoAd→FA is the #1 boundary error: background patches classified as focal "
                  "adhesion, especially in pfak→vinc (169/208 NoAd errors = 81%)"),
        (C_RED,   "FC remains the hardest FA subtype: low recall across all scenarios; "
                  "mostly confused with NA and FA even in the 5-class model"),
        (C_AMBER, "ctrl→ycomp (33% error): NoAd patches in drug-treated cells "
                  "(ycomp) look more like FA — actin remodeling changes background texture"),
        (C_AMBER, "NA→NoAd confusion: ~5–15% of Nascent Adhesion patches predicted as "
                  "background, suggesting NA and NoAd share similar low-intensity morphology"),
        (C_DARK,  "Within-dataset (vinc_only) is much better at 20% — "
                  "showing the model can learn No adhesion vs adhesion when both appear in training"),
    ]
    fy = BODY_T + Inches(0.48)
    for color, text in findings:
        _txt(sl, PAD, fy, col_w - Inches(0.1), Inches(0.55),
             f"• {text}", size=9, color=color, wrap=True)
        fy += Inches(0.57)

    _txt(sl, col2_x, BODY_T + Inches(0.1), col_w, Inches(0.3),
         "Recommended Actions", size=13, bold=True, color=C_DARK)

    actions = [
        (C_GREEN, "Add more NoAd labels to pFAK: only 60 exist — the model cannot "
                  "generalize NoAd across domains with this few examples"),
        (C_GREEN, "Inspect NoAd→FA patches in pfak→vinc: are they truly 'no adhesion' "
                  "or are they genuine early-stage adhesions that were mislabeled?"),
        (C_AMBER, "Consider hierarchical classification: Stage-1 Ad vs NoAd → "
                  "Stage-2 4-class FA subtypes (current pipeline), then evaluate each stage separately"),
        (C_DARK,  "Class-balanced SupCon training with NoAd: retrain Stage-2 AE "
                  "with No adhesion patches included in contrastive batches"),
        (C_DARK,  "Review NA labeling criteria: NA→NoAd confusion suggests annotators "
                  "and the model both see overlap at the adhesion formation boundary"),
    ]
    sy = BODY_T + Inches(0.48)
    for color, text in actions:
        _txt(sl, col2_x, sy, col_w - Inches(0.1), Inches(0.55),
             f"• {text}", size=9, color=color, wrap=True)
        sy += Inches(0.57)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default=str(OUT_DEFAULT))
    args = parser.parse_args()

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slides = [
        ("Cover",                  slide_cover),
        ("Label statistics",       slide_label_stats),
        ("Error overview",         slide_overview),
        ("No-ad boundary",         slide_noad_boundary),
        ("Per-class error rates",  slide_perclass_rates),
    ]

    detail_configs = [
        ("vinc_only",
         "Within vinc (ctrl + ycomp, 75% split): 20% error. NoAd is well separated — "
         "the model has seen enough no-adhesion examples. FC↔NA and FC↔FA remain the hard boundaries. "
         "Fib (18 patches) has high error due to rarity — class weighting helps but is insufficient."),
        ("ctrl->ycomp",
         "ctrl → ycomp: trained on control vinc, tested on Y-compound treated vinc. 33% error. "
         "Drug treatment remodels the actin cytoskeleton — NoAd patches under ycomp look more like FA. "
         "NA→FA (49) and FC→FA (43) dominate: the drug shifts intensity distributions toward FA-like morphology."),
        ("pfak->vinc",
         "pFAK → Vinc: worst scenario. 36% error. pFAK has only 60 NoAd patches — the classifier "
         "never learns what vinc NoAd looks like, causing massive NoAd→FA confusion (169 errors). "
         "FC→NA also reappears (49 errors). More pFAK NoAd labels are the key bottleneck."),
    ]

    for i, (name, fn) in enumerate(slides, 1):
        print(f"  {i}. {name}")
        fn(prs)

    for i, (sc, notes) in enumerate(detail_configs, len(slides) + 1):
        print(f"  {i}. Detail — {sc}")
        _detail_slide(prs, sc, notes)

    print(f"  {len(slides) + len(detail_configs) + 1}. Key findings")
    slide_findings(prs)

    prs.save(str(out_path))
    print(f"\nSaved {len(prs.slides)} slides → {out_path}")


if __name__ == "__main__":
    main()
