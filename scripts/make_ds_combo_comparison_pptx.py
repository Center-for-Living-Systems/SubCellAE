#!/usr/bin/env python3
"""
make_ds_combo_comparison_pptx.py
Short summary PPT comparing ds_combo v1 (lc025) and v2 (lc010 balanced).
"""
from __future__ import annotations
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W_IN, H_IN = 13.33, 7.5
MARGIN = 0.35

C_DARK  = RGBColor(0x1F, 0x2D, 0x3D)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
C_LGRAY2 = RGBColor(0xDD, 0xDD, 0xDD)
C_BLUE  = RGBColor(0x2E, 0x86, 0xC1)
C_GREEN = RGBColor(0x1A, 0x7A, 0x4A)
C_V1    = RGBColor(0xC0, 0x60, 0x20)   # orange-ish for v1
C_V2    = RGBColor(0x1A, 0x6A, 0xB0)   # blue for v2


def _px(v): return Inches(v)

def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _tb(slide, text, left, top, w, h,
        size=11, bold=False, italic=False,
        color=C_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(w), _px(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def _hline(slide, y, x0=MARGIN, x1=None, color=C_LGRAY2, width_pt=0.75):
    if x1 is None: x1 = W_IN - MARGIN
    from pptx.util import Pt as Pt2
    line = slide.shapes.add_shape(1, _px(x0), _px(y), _px(x1 - x0), _px(0.01))
    line.fill.background()
    line.line.color.rgb = color
    line.line.width = Pt2(width_pt)

def _cell(slide, text, cx, cy, cw, ch,
          size=9, bold=False, color=C_DARK, align=PP_ALIGN.CENTER,
          border=True, fill_color=None):
    box = slide.shapes.add_shape(1, _px(cx), _px(cy), _px(cw), _px(ch))
    if fill_color:
        box.fill.solid(); box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    if border:
        box.line.color.rgb = C_LGRAY2
    else:
        box.line.fill.background()
    _tb(slide, text, cx + 0.04, cy + 0.04, cw - 0.08, ch - 0.06,
        size=size, bold=bold, color=color, align=align)


# ── Slide 1: Title ────────────────────────────────────────────────────────────
def slide_title(prs):
    sl = _add_slide(prs)
    _tb(sl, "Dataset Combination ConAE Sweep",
        1.0, 1.4, W_IN - 2.0, 1.1,
        size=36, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _tb(sl, "v1 (λ=0.25, equal oversampling)  vs  v2 (λ=0.10, balanced split)",
        1.0, 2.65, W_IN - 2.0, 0.5,
        size=17, color=C_GRAY, align=PP_ALIGN.CENTER)
    _hline(sl, 3.3, x0=2.0, x1=W_IN - 2.0)
    items = [
        "15 combinations of {ds1, ds2, ds3, ds4}  (singles, pairs, triples, all-4)",
        "ConAE · enlcrop · sc2 · nL1 loss · latent 12 · proj 8 · 500 epochs · Adam cosine LR",
        "ds1 = vinc   ds2 = pfak   ds3 = ppax   ds4 = nih3t3",
    ]
    y = 3.55
    for item in items:
        _tb(sl, f"• {item}", 2.0, y, W_IN - 4.0, 0.38, size=13, color=C_GRAY)
        y += 0.42


# ── Slide 2: Side-by-side comparison ─────────────────────────────────────────
def slide_comparison(prs):
    sl = _add_slide(prs)
    _tb(sl, "Sweep Comparison: v1 vs v2", MARGIN, 0.12, W_IN - 2*MARGIN, 0.42,
        size=22, bold=True, color=C_DARK)
    _hline(sl, 0.60)

    # column layout: label | v1 | v2
    cx = [MARGIN, 3.0, 8.3]
    cw = [2.6,    5.2, 4.7]
    rh = 0.40
    y0 = 0.72

    headers = ["Setting", "v1  (lc025)", "v2  (lc010 balanced)"]
    hcols   = [C_DARK, C_V1, C_V2]
    for i, (hdr, hc, x, w) in enumerate(zip(headers, hcols, cx, cw)):
        _cell(sl, hdr, x, y0, w, rh, size=10, bold=True, color=hc,
              align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)

    rows = [
        ("λ_contrast",
         "0.25",
         "0.10  (reduced — contrastive loss ~5 at end of v1)"),
        ("ds1 (vinc) split",
         "80% train / 20% val",
         "40% train / 60% val  (per-entry val_split)"),
        ("ds1 repeats",   "×1",  "×1"),
        ("ds2 (pfak) repeats", "×8", "×3"),
        ("ds3 (ppax) repeats", "×4", "×2"),
        ("ds4 (nih3t3) repeats", "×4", "×2"),
        ("ds1 train patches", "~22,110", "~11,055"),
        ("ds2 train patches", "~17,400", "~8,100"),
        ("ds3 train patches", "~21,300", "~10,400"),
        ("ds4 train patches", "~23,200", "~11,600"),
        ("Total train patches (all-4)", "~84,000", "~41,000  (≈½ → ~2× faster/epoch)"),
        ("Reconstruction loss", "nL1  ✓  (plot label was wrong → fixed)", "nL1  ✓"),
        ("Output dir", "ds_combo_enlcrop_sc2/", "ds_combo_enlcrop_sc2_lc010_bal/"),
    ]

    for ri, (label, v1, v2) in enumerate(rows):
        y = y0 + rh * (ri + 1)
        vals = [label, v1, v2]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
        colors = [C_DARK, C_V1, C_V2]
        for i, (val, x, w, aln, col) in enumerate(zip(vals, cx, cw, aligns, colors)):
            _cell(sl, val, x, y, w, rh, size=8.5,
                  color=col, align=aln)


# ── Slide 3: Status ───────────────────────────────────────────────────────────
def slide_status(prs):
    sl = _add_slide(prs)
    _tb(sl, "Training Status", MARGIN, 0.12, W_IN - 2*MARGIN, 0.42,
        size=22, bold=True, color=C_DARK)
    _hline(sl, 0.60)

    combos = [
        "vinc", "nih3t3", "ppax", "pfak",
        "vinc_nih3t3", "vinc_ppax", "vinc_pfak",
        "nih3t3_ppax", "nih3t3_pfak", "ppax_pfak",
        "vinc_nih3t3_ppax", "vinc_nih3t3_pfak", "vinc_ppax_pfak",
        "nih3t3_ppax_pfak", "vinc_nih3t3_ppax_pfak",
    ]
    run_root = Path("/net/projects/CLS/lding/data/fa_data_analysis/"
                    "ae_results/contrastive_run/ds_combo_enlcrop_sc2_lc010_bal")

    # check checkpoint count
    def status(combo):
        d = run_root / combo
        n = len(list(d.glob("*.pt"))) if d.exists() else 0
        if n >= 12: return "Done ✓", C_GREEN
        if n > 0:   return f"Running ({n}/12 ckpt)", C_V1
        return "Pending", C_LGRAY

    # two columns of combos
    col_split = 8
    for col_idx, col_combos in enumerate([combos[:col_split], combos[col_split:]]):
        cx = MARGIN + col_idx * 6.5
        rh = 0.40
        y0 = 0.72
        cw_lbl, cw_st = 3.5, 2.8
        _cell(sl, "Combo", cx, y0, cw_lbl, rh, size=9, bold=True, color=C_DARK)
        _cell(sl, "v1",    cx + cw_lbl,           y0, 0.8, rh, size=9, bold=True, color=C_V1, align=PP_ALIGN.CENTER)
        _cell(sl, "v2",    cx + cw_lbl + 0.8,     y0, cw_st, rh, size=9, bold=True, color=C_V2, align=PP_ALIGN.CENTER)

        for ri, combo in enumerate(col_combos):
            y = y0 + rh * (ri + 1)
            label = combo.replace("_", "+").replace("nih3t3","ds4").replace("ppax","ds3").replace("pfak","ds2").replace("vinc","ds1")
            st_v2, col_v2 = status(combo)

            # v1 always done (from prior session)
            _cell(sl, label, cx, y, cw_lbl, rh, size=8.5, color=C_DARK, align=PP_ALIGN.LEFT)
            _cell(sl, "Done ✓", cx + cw_lbl, y, 0.8, rh, size=8, color=C_GREEN, align=PP_ALIGN.CENTER)
            _cell(sl, st_v2, cx + cw_lbl + 0.8, y, cw_st, rh, size=8, color=col_v2, align=PP_ALIGN.CENTER)


# ── build PPT ─────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    slide_title(prs)
    slide_comparison(prs)
    slide_status(prs)

    out = Path("ds_combo_sweep_comparison.pptx")
    prs.save(str(out))
    print(f"Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
