#!/usr/bin/env python3
"""Generate a PowerPoint explaining the EnlargedJitterCrop augmentation."""

import io
import math
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, Rectangle
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── colour palette ─────────────────────────────────────────────────────────────
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE   = RGBColor(0x1A, 0x3A, 0x6B)   # dark navy
C_ACCENT  = RGBColor(0x2E, 0x86, 0xC1)   # medium blue
C_ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # orange highlight
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_GRAY    = RGBColor(0x66, 0x66, 0x66)
C_LGRAY   = RGBColor(0xEE, 0xEE, 0xEE)

W = Inches(13.33)
H = Inches(7.5)


def rgb_hex(c: RGBColor):
    return (c[0]/255, c[1]/255, c[2]/255)


def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def clear_placeholders(slide):
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)


def text_box(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             bg_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def slide_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, title_text, subtitle=None):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()

    text_box(slide, title_text,
             Inches(0.3), Inches(0.1), Inches(12), Inches(0.65),
             font_size=28, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    if subtitle:
        text_box(slide, subtitle,
                 Inches(0.3), Inches(0.7), Inches(12), Inches(0.35),
                 font_size=14, color=RGBColor(0xCC,0xDD,0xFF))


def add_image(slide, stream, left, top, width=None, height=None):
    return slide.shapes.add_picture(stream, left, top, width=width, height=height)


# =============================================================================
# Slide 1 — Title
# =============================================================================
def slide_title(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)

    # full-slide navy background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_TITLE

    # big title
    text_box(slide, "Enlarged Jitter Crop Augmentation",
             Inches(1), Inches(1.8), Inches(11), Inches(1.5),
             font_size=40, bold=True,
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    text_box(slide, "View generation for contrastive learning on FA patches",
             Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
             font_size=22, color=RGBColor(0xAA,0xCC,0xFF),
             align=PP_ALIGN.CENTER)

    text_box(slide, "SubCellAE  ·  Contrastive AE training",
             Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             font_size=16, color=RGBColor(0x88,0xAA,0xDD),
             align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 2 — Motivation
# =============================================================================
def slide_motivation(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "Motivation: Why not augment the 32×32 patch directly?")

    # ── diagram: naive augmentation ──
    fig, axes = plt.subplots(1, 3, figsize=(10, 3.2), facecolor="white")
    ax_patch, ax_rot, ax_crop = axes

    rng = np.random.default_rng(7)
    # simulate an FA patch — bright blob in center
    x = np.linspace(-1, 1, 58)
    XX, YY = np.meshgrid(x, x)
    blob = np.exp(-(XX**2 + YY**2) / 0.15) * 0.8
    blob += rng.normal(0, 0.04, blob.shape)
    blob = np.clip(blob, 0, 1)

    # 32x32 center crop
    cx = 58 // 2
    patch32 = blob[cx-16:cx+16, cx-16:cx+16]

    # rotate patch32 by 30 deg (naive) — corners go black/gray
    from scipy.ndimage import rotate as ndrotate
    rot32 = ndrotate(patch32, 30, reshape=False, cval=0.0)

    ax_patch.imshow(patch32, cmap="gray", vmin=0, vmax=1)
    ax_patch.set_title("32×32 patch\n(stored on disk)", fontsize=11)
    ax_patch.axis("off")
    rect = Rectangle((0,0), 31, 31, lw=2, edgecolor="steelblue", facecolor="none")
    ax_patch.add_patch(rect)

    ax_rot.imshow(rot32, cmap="gray", vmin=0, vmax=1)
    ax_rot.set_title("Rotate 30°\n→ corners empty / reflect-padded", fontsize=11)
    ax_rot.axis("off")
    for corner_x, corner_y in [(0,0),(28,0),(0,28),(28,28)]:
        rect = Rectangle((corner_x,corner_y), 3, 3, lw=0,
                          facecolor=(1,0.3,0.3,0.4))
        ax_rot.add_patch(rect)
    ax_rot.text(16, 30.5, "⚠ boundary artifacts", ha="center",
                color="red", fontsize=9)

    # double interpolation illustration
    rot58 = ndrotate(blob, 30, reshape=False, cval=0.0)
    crop_from_rot = rot58[cx-16:cx+16, cx-16:cx+16]
    ax_crop.imshow(crop_from_rot, cmap="gray", vmin=0, vmax=1)
    ax_crop.set_title("Rotate full frame → crop\n(2 interpolations, slow per-aug)", fontsize=11)
    ax_crop.axis("off")

    for ax in axes:
        ax.set_facecolor("white")

    fig.tight_layout(pad=1.0)
    stream = fig_to_stream(fig)
    plt.close(fig)
    add_image(slide, stream, Inches(0.8), Inches(1.25), width=Inches(11.5))

    # bullet points
    bullets = [
        "• Rotating a 32×32 patch: corners fall outside → must pad with zeros or reflect → border artifacts",
        "• Rotating the full frame first, then cropping: double interpolation → blurring; slow at training time",
        "• Storing pre-rotated patches: explosion of disk space; fixed augmentation set",
    ]
    y = Inches(4.55)
    for b in bullets:
        text_box(slide, b, Inches(0.5), y, Inches(12.3), Inches(0.45),
                 font_size=14, color=C_GRAY)
        y += Inches(0.44)

    text_box(slide,
             "Solution: extract a larger context window once; apply a single GPU affine "
             "transform per training step.",
             Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.55),
             font_size=15, bold=True, color=C_TITLE)


# =============================================================================
# Slide 3 — Overview of EnlargedCrop pipeline
# =============================================================================
def slide_overview(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "EnlargedJitterCrop — Two-Stage Pipeline")

    fig, axes = plt.subplots(1, 4, figsize=(13, 3.5), facecolor="white",
                              gridspec_kw={"width_ratios": [3,2,2,2]})

    rng = np.random.default_rng(12)
    x = np.linspace(-1, 1, 200)
    XX, YY = np.meshgrid(x, x)
    frame = np.exp(-(XX**2 + YY**2) / 0.08) * 0.7
    frame += rng.normal(0, 0.03, frame.shape)
    frame = np.clip(frame, 0, 1)

    # Add some texture
    frame += 0.15 * np.exp(-((XX-0.5)**2 + (YY-0.3)**2)/0.05)
    frame += 0.1  * np.exp(-((XX+0.4)**2 + (YY+0.6)**2)/0.04)
    frame = np.clip(frame, 0, 1)

    ax0, ax1, ax2, ax3 = axes
    ax0.imshow(frame, cmap="gray", vmin=0, vmax=1)
    ax0.set_title("Full frame\n(CIO-RB normalised)", fontsize=10)
    # show context window
    cx = 100
    half = 29  # 58/2
    rect = Rectangle((cx-half, cx-half), 58, 58, lw=2.5,
                      edgecolor="orange", facecolor="none")
    ax0.add_patch(rect)
    ax0.text(cx, cx-half-3, "58×58 context\n(extracted at init)",
             ha="center", va="bottom", fontsize=8,
             color="orange", fontweight="bold")
    ax0.axis("off")

    # context patch
    ctx = frame[cx-half:cx+half, cx-half:cx+half]
    ax1.imshow(ctx, cmap="gray", vmin=0, vmax=1)
    ax1.set_title("58×58 context patch\n(stored in RAM)", fontsize=10)
    # show inner 32x32
    inner_off = (58-32)//2
    rect2 = Rectangle((inner_off, inner_off), 32, 32, lw=2,
                       edgecolor="steelblue", facecolor="none", linestyle="--")
    ax1.add_patch(rect2)
    ax1.text(29, 57, "32×32 inner\n(target size)",
             ha="center", va="bottom", fontsize=7.5,
             color="steelblue")
    ax1.axis("off")

    # view 1 — rotated + shifted
    from scipy.ndimage import rotate as ndrotate, shift as ndshift
    v1 = ndrotate(ctx, 12, reshape=False, cval=0)
    v1 = ndshift(v1, (3, -2))
    v1_crop = v1[inner_off:inner_off+32, inner_off:inner_off+32]
    ax2.imshow(v1_crop, cmap="gray", vmin=0, vmax=1)
    ax2.set_title("View 1\n(θ=+12°, dx=−2, dy=+3)", fontsize=10)
    ax2.axis("off")
    ax2.set_facecolor("#E8F4FD")

    v2 = ndrotate(ctx, -8, reshape=False, cval=0)
    v2 = ndshift(v2, (-2, 3))
    v2_crop = v2[inner_off:inner_off+32, inner_off:inner_off+32]
    ax3.imshow(v2_crop, cmap="gray", vmin=0, vmax=1)
    ax3.set_title("View 2\n(θ=−8°, dx=+3, dy=−2)", fontsize=10)
    ax3.axis("off")
    ax3.set_facecolor("#FEF9E7")

    # arrows between panels
    fig.text(0.285, 0.55, "① extract\nat init", ha="center",
             fontsize=9, color="darkorange", fontweight="bold")
    fig.text(0.285, 0.46, "→", ha="center", fontsize=18, color="darkorange")

    fig.text(0.515, 0.60, "② affine\nview 1", ha="center",
             fontsize=9, color="steelblue", fontweight="bold")
    fig.text(0.515, 0.51, "↗", ha="center", fontsize=18, color="steelblue")

    fig.text(0.515, 0.40, "② affine\nview 2", ha="center",
             fontsize=9, color="goldenrod", fontweight="bold")
    fig.text(0.515, 0.31, "↘", ha="center", fontsize=18, color="goldenrod")

    fig.tight_layout()
    stream = fig_to_stream(fig)
    plt.close(fig)
    add_image(slide, stream, Inches(0.1), Inches(1.2), width=Inches(13.0))

    bullets = [
        ("① Init-time", "context patch (58×58) extracted from full normalised frame — stored in RAM, no disk overhead"),
        ("② Train-time", "two independent random affines (GPU) → two augmented 32×32 views; single bilinear interpolation"),
    ]
    y = Inches(5.0)
    for label, body in bullets:
        text_box(slide, label, Inches(0.4), y, Inches(1.8), Inches(0.42),
                 font_size=14, bold=True, color=C_ACCENT)
        text_box(slide, body,  Inches(2.1), y, Inches(10.7), Inches(0.42),
                 font_size=14, color=C_GRAY)
        y += Inches(0.48)

    text_box(slide,
             "Same context patch → two structurally-similar but geometrically-distinct views "
             "→ positive pair for contrastive loss",
             Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.5),
             font_size=14, bold=True, color=C_TITLE)


# =============================================================================
# Slide 4 — Context size formula
# =============================================================================
def slide_formula(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "Context Size Formula")

    # diagram
    fig, ax = plt.subplots(1, 1, figsize=(5.5, 5.5), facecolor="white")

    ax.set_xlim(-40, 40); ax.set_ylim(-40, 40)
    ax.set_aspect("equal")
    ax.axis("off")

    # outer context
    ctx_half = 29  # 58/2
    outer = Rectangle((-ctx_half, -ctx_half), 2*ctx_half, 2*ctx_half,
                       lw=2.5, edgecolor="darkorange", facecolor="#FFF3E0")
    ax.add_patch(outer)
    ax.text(0, ctx_half+1.5, "58×58 context", ha="center", fontsize=10,
            color="darkorange", fontweight="bold")

    # inner target 32×32
    inner = Rectangle((-16, -16), 32, 32, lw=2, edgecolor="steelblue",
                       facecolor="#E8F4FD", linestyle="--", zorder=3)
    ax.add_patch(inner)
    ax.text(0, -18.5, "32×32 target", ha="center", fontsize=10,
            color="steelblue", fontweight="bold")

    # max shift indicator
    for dx, dy, col in [(4, 0, "green"), (0, 4, "green"),
                         (-4, 0, "green"), (0, -4, "green")]:
        ax.annotate("", xy=(dx, dy), xytext=(0, 0),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.5))
    ax.text(5, 2, "max_shift_px = 4", fontsize=9, color="green")

    # worst-case diagonal (rotation + shift)
    worst = math.sqrt(2) * (16 + 4)  # sqrt(2)*(ps/2 + shift)
    circle = plt.Circle((0, 0), worst, color="red", fill=False,
                          linestyle=":", lw=2, zorder=4)
    ax.add_patch(circle)
    ax.annotate("", xy=(worst/math.sqrt(2), worst/math.sqrt(2)), xytext=(0, 0),
                arrowprops=dict(arrowstyle="->", color="red", lw=1.5))
    ax.text(worst/math.sqrt(2)+0.5, worst/math.sqrt(2)+0.5,
            f"√2 × 20 ≈ {worst:.1f}", fontsize=9, color="red")

    ax.set_facecolor("white")
    fig.tight_layout()
    stream = fig_to_stream(fig)
    plt.close(fig)
    add_image(slide, stream, Inches(7.5), Inches(1.2), width=Inches(5.5))

    # formula text
    lines = [
        ("Formula:", True, 20, C_TITLE),
        ("context_size = 2 × ⌈ √2 × (patch_size/2 + max_shift_px) ⌉", True, 18, C_ACCENT),
        ("", False, 10, C_GRAY),
        ("Derivation:", True, 16, C_TITLE),
        ("• Worst case: 45° rotation + full shift at every corner", False, 14, C_GRAY),
        ("• Corner of 32×32 sits at radius = patch_size/2 = 16 px from center", False, 14, C_GRAY),
        ("• After shift, effective radius = 16 + 4 = 20 px", False, 14, C_GRAY),
        ("• After 45° rotation, farthest reach = √2 × 20 ≈ 28.3 px", False, 14, C_GRAY),
        ("• Round up: ⌈28.3⌉ = 29  →  context = 2 × 29 = 58", False, 14, C_GRAY),
        ("", False, 10, C_GRAY),
        ("For defaults  ps = 32,  max_shift = 4:", True, 15, C_TITLE),
        ("    context_size  =  2 × ⌈ √2 × 20 ⌉  =  2 × 29  =  58 px", True, 15, C_ORANGE),
    ]
    y = Inches(1.25)
    for txt, bold, fsize, color in lines:
        text_box(slide, txt, Inches(0.4), y, Inches(7.0), Inches(0.42),
                 font_size=fsize, bold=bold, color=color)
        y += Inches(0.42)


# =============================================================================
# Slide 5 — Affine transform details
# =============================================================================
def slide_affine(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "GPU Affine Transform — Single Interpolation")

    # code box
    code = (
        "angles = Uniform(−max_angle_deg, +max_angle_deg)   # per image\n"
        "dx, dy = Uniform(−max_shift_px,  +max_shift_px)    # per image\n\n"
        "s  = out_size / context_size   # scale: 32/58 ≈ 0.552\n"
        "tx = 2 · dx / context_size     # normalised translation\n"
        "ty = 2 · dy / context_size\n\n"
        "θ = [[ cos·s,  sin·s,  tx ],\n"
        "      [−sin·s,  cos·s,  ty ]]\n\n"
        "grid   = affine_grid(θ, output=(B, C, 32, 32))\n"
        "output = grid_sample(context_batch, grid, mode='bilinear')"
    )
    tb = text_box(slide, code,
                  Inches(0.3), Inches(1.25), Inches(8.3), Inches(4.5),
                  font_size=12, color=RGBColor(0xDD, 0xEE, 0xFF),
                  bg_color=RGBColor(0x1A, 0x1A, 0x2E))
    tb.text_frame.paragraphs[0].runs[0].font.name = "Courier New"

    key_points = [
        ("Independent per image", "Each image in the batch gets its own random angle and translation — no batch-level correlation"),
        ("Single interpolation", "Rotation + translation + scale fused into one affine_grid + grid_sample call — avoids double blurring"),
        ("GPU-native", "No CPU↔GPU transfer; the 58×58 context tensors live on GPU, transform applied in the forward pass"),
        ("Contrastive views", "Called twice independently on the same context batch → view1 and view2 share FA content but differ in pose"),
        ("Validation", "At eval time: zero shift, zero angle → deterministic center crop; same code path"),
    ]

    y = Inches(1.25)
    for title, body in key_points:
        # coloured dot
        dot = slide.shapes.add_shape(1,
            Inches(8.9), y + Inches(0.08), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()

        text_box(slide, title, Inches(9.15), y, Inches(3.9), Inches(0.3),
                 font_size=13, bold=True, color=C_TITLE)
        text_box(slide, body,  Inches(9.15), y+Inches(0.3), Inches(3.9), Inches(0.55),
                 font_size=12, color=C_GRAY)
        y += Inches(1.0)


# =============================================================================
# Slide 6 — Parameters in use
# =============================================================================
def slide_params(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "Parameters Used in Current Models")

    headers  = ["Parameter", "Value", "Notes"]
    rows = [
        ("patch_size",       "32 px",    "Target output size"),
        ("context_size",     "58 px",    "Enlarged context extracted from full frame at init"),
        ("max_shift_px",     "4 px",     "Max translation per view (±4 px, independent x/y)"),
        ("max_angle_deg",    "15°",      "Max rotation per view (uniform ±15°)"),
        ("mode",             "bilinear", "Interpolation mode in grid_sample"),
        ("padding_mode",     "border",   "Pads context edges with border values (avoids black)"),
        ("Applies to",       "actin-only models", "2ch models: not yet (MultiChannelEnlargedCropDataset pending)"),
    ]

    col_w = [Inches(3.2), Inches(2.2), Inches(7.0)]
    col_x = [Inches(0.25), Inches(3.5), Inches(5.75)]
    row_h = Inches(0.52)
    y0    = Inches(1.3)

    # header row
    for j, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        box = slide.shapes.add_shape(1, cx, y0, cw, row_h)
        box.fill.solid(); box.fill.fore_color.rgb = C_TITLE
        box.line.fill.background()
        text_box(slide, hdr, cx+Inches(0.08), y0+Inches(0.1),
                 cw-Inches(0.1), row_h-Inches(0.1),
                 font_size=14, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    for i, (p, v, n) in enumerate(rows):
        y = y0 + row_h * (i + 1)
        bg = C_LGRAY if i % 2 == 0 else C_BG
        for j, (val, cw, cx) in enumerate(zip([p, v, n], col_w, col_x)):
            box = slide.shapes.add_shape(1, cx, y, cw, row_h)
            box.fill.solid(); box.fill.fore_color.rgb = bg
            box.line.fill.background()
            fc = C_ORANGE if j == 1 else C_GRAY
            bold = j == 1
            text_box(slide, val, cx+Inches(0.08), y+Inches(0.08),
                     cw-Inches(0.1), row_h-Inches(0.1),
                     font_size=13, bold=bold, color=fc)

    text_box(slide,
             "EnlargedJitterCrop is active during contrastive / supervised-contrastive training only. "
             "Baseline AE uses standard 32×32 patches with no geometric augmentation.",
             Inches(0.25), Inches(6.05), Inches(12.8), Inches(0.55),
             font_size=13, color=C_GRAY)


# =============================================================================
# Slide 7 — Summary
# =============================================================================
def slide_summary(prs):
    slide = add_slide(prs)
    clear_placeholders(slide)
    slide_bg(slide)
    title_bar(slide, "Summary")

    points = [
        (C_ACCENT,  "What",
         "Enlarged context extraction + GPU affine transform to generate two "
         "geometrically diverse views of the same focal adhesion patch"),
        (C_GREEN,   "Why",
         "Avoids boundary artifacts (no zero-padding at patch edges) and double "
         "blurring (single bilinear interpolation). Context extracted once at init; "
         "no per-epoch disk I/O overhead."),
        (C_ORANGE,  "How — contrastive training",
         "Each mini-batch: context_batch (B × 1 × 58 × 58) → _jitter_rot_crop called "
         "twice independently → view1, view2 (B × 1 × 32 × 32). Positive pairs share "
         "FA identity; negative pairs come from different FAs."),
        (C_TITLE,   "Current settings",
         "context = 58, shift ≤ 4 px, angle ≤ 15°. Applies to actin-only contrastive "
         "models (ch3). 2-channel models use fixed 32×32 patches (no enlcrop yet)."),
    ]

    y = Inches(1.3)
    for color, label, body in points:
        bar = slide.shapes.add_shape(1, Inches(0.25), y, Inches(0.08), Inches(0.9))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        text_box(slide, label, Inches(0.5), y, Inches(12.3), Inches(0.35),
                 font_size=16, bold=True, color=color)
        text_box(slide, body,  Inches(0.5), y+Inches(0.35), Inches(12.3), Inches(0.65),
                 font_size=14, color=C_GRAY)
        y += Inches(1.25)


# =============================================================================
# Build presentation
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_motivation(prs)
    slide_overview(prs)
    slide_formula(prs)
    slide_affine(prs)
    slide_params(prs)
    slide_summary(prs)

    out = "slides_enlcrop.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
