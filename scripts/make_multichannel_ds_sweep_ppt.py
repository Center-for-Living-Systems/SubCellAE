#!/usr/bin/env python3
"""
make_multichannel_ds_sweep_ppt.py

Dataset-combination sweep for multi-channel ConAE models:
  Part 1 — 3ch pza (pax+zyx+act): 15 ds combos × 3 losses
  Part 2 — 4ch vinc (vinc+pax+zyx+act): v / nv / n training variants
  Part 3 — Observation: nL1 scale is dataset-intrinsic, not training-dependent

Shows whatever outputs exist; missing files appear as "pending" placeholders.

Usage:
    python scripts/make_multichannel_ds_sweep_ppt.py [--out slides_multichannel_ds_sweep.pptx]
"""

from __future__ import annotations
import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ─────────────────────────────────────────────────────────────────────

SWEEP = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/protein_sweep")

# ── combo definitions ─────────────────────────────────────────────────────────

# letter codes: f=pfak(ds2), n=nih3t3(ds4), p=ppax(ds3), v=vinc(ds1)
DS_LETTER = {"v": "ds1", "f": "ds2", "p": "ds3", "n": "ds4"}
DS_FULL   = {"v": "vinc", "f": "pfak", "p": "ppax", "n": "nih3t3"}

SINGLES   = ["v", "f", "p", "n"]
PAIRS_V   = ["fv", "pv", "nv"]           # pairs containing ds1
PAIRS_NOV = ["fn", "fp", "np"]           # pairs without ds1
TRIPLES   = ["fnp", "fnv", "fpv", "npv"]
ALL_DS    = ["fnpv"]

GROUPS_3CH = [SINGLES, PAIRS_V, PAIRS_NOV, TRIPLES, ALL_DS]
GROUP_LABELS = [
    "Singles  (ds1 · ds2 · ds3 · ds4)",
    "Pairs with ds1  (ds1+ds2 · ds1+ds3 · ds1+ds4)",
    "Pairs without ds1  (ds2+ds4 · ds2+ds3 · ds3+ds4)",
    "Triples  (all 4 triples)",
    "All 4 datasets",
]

LOSSES = ["nl1", "mse", "l1"]
LOSS_LABEL = {"nl1": "nL1", "mse": "MSE", "l1": "L1"}

# 4ch vinc training variants
VARIANTS_4CH = ["v", "nv", "n"]
VARIANT_4CH_LABEL = {"v": "ds1 only", "nv": "ds1+ds4", "n": "ds4 only"}

# ── slide geometry & colours ──────────────────────────────────────────────────

W_IN, H_IN = 13.33, 7.5
MARGIN   = 0.25
HEADER_H = 0.68

C_TITLE = RGBColor(0x1F, 0x2D, 0x3D)
C_GRAY  = RGBColor(0x66, 0x66, 0x66)
C_LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_3CH   = RGBColor(0x62, 0x39, 0xA8)   # purple
C_4CH   = RGBColor(0xD4, 0x6C, 0x08)   # orange
C_NOTE  = RGBColor(0x1A, 0x7A, 0x40)   # green for observation slide

# ── pptx helpers ──────────────────────────────────────────────────────────────

def _px(v): return Inches(v)
def _add_slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def _textbox(slide, text, left, top, width, height,
             font_size=11, bold=False, color=C_TITLE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def _slide_header(slide, title, subtitle=None, color=C_TITLE):
    _textbox(slide, title, MARGIN, 0.05, W_IN - 0.5, 0.45,
             font_size=20, bold=True, color=color)
    rule = slide.shapes.add_shape(1, _px(MARGIN), _px(0.55), _px(W_IN - 0.5), _px(0.025))
    rule.fill.solid(); rule.fill.fore_color.rgb = color; rule.line.fill.background()
    if subtitle:
        _textbox(slide, subtitle, MARGIN, 0.57, W_IN - 0.5, 0.22,
                 font_size=10, color=C_GRAY, italic=True)

def _section_slide(prs, title, subtitle, color):
    slide = _add_slide(prs)
    bar = slide.shapes.add_shape(1, _px(0), _px(2.6), _px(W_IN), _px(1.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    _textbox(slide, title, 0.5, 2.75, W_IN - 1.0, 1.2,
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, subtitle, 0.5, 4.25, W_IN - 1.0, 0.5,
             font_size=15, color=C_WHITE, align=PP_ALIGN.CENTER)

def _add_image_bytes(slide, img_bytes, left, top, width=None, height=None):
    return slide.shapes.add_picture(
        io.BytesIO(img_bytes), _px(left), _px(top),
        width=_px(width) if width is not None else None,
        height=_px(height) if height is not None else None,
    )

def _fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()

def _load_png(path: Path) -> np.ndarray | None:
    if path.exists():
        return np.array(Image.open(str(path)).convert("RGB"))
    return None


# ── combo helpers ─────────────────────────────────────────────────────────────

def _combo_label(combo: str) -> str:
    return "+".join(DS_LETTER[c] for c in sorted(combo))

def _model_dir_3ch(combo: str, loss: str) -> Path:
    return SWEEP / f"conae_3ch_pza_{combo}_{loss}"

def _model_dir_4ch(variant: str, loss: str) -> Path:
    return SWEEP / f"conae_4ch_vinc_{variant}_{loss}"


# ── figure builders ───────────────────────────────────────────────────────────

def _violin_grid(combos: list[str], loss: str,
                 metric: str = "recon_nl1",
                 model_type: str = "3ch") -> bytes:
    """One violin per combo, arranged in a row."""
    n = len(combos)
    n_cols = min(n, 4)
    n_rows = (n + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                             figsize=(n_cols * 5.0, n_rows * 3.5),
                             facecolor="white")
    axes = np.array(axes).reshape(n_rows, n_cols)

    for i, combo in enumerate(combos):
        r, c = i // n_cols, i % n_cols
        ax = axes[r, c]
        d = _model_dir_3ch(combo, loss) if model_type == "3ch" \
            else _model_dir_4ch(combo, loss)
        img = _load_png(d / f"cross_dataset_{metric}.png")
        if img is not None:
            ax.imshow(img)
        else:
            ax.text(0.5, 0.5, "pending", ha="center", va="center",
                    fontsize=12, color="gray", transform=ax.transAxes)
        ax.set_title(_combo_label(combo) if model_type == "3ch"
                     else f"{VARIANT_4CH_LABEL[combo]}",
                     fontsize=11, fontweight="bold", pad=3)
        ax.axis("off")

    for i in range(len(combos), n_rows * n_cols):
        axes[i // n_cols, i % n_cols].axis("off")

    fig.tight_layout(pad=0.5)
    return _fig_to_bytes(fig)


def _umap_grid(combos: list[str], loss: str,
               umap_key: str = "umap_annotation",
               model_type: str = "3ch") -> bytes:
    n = len(combos)
    n_cols = min(n, 4)
    n_rows = (n + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                             figsize=(n_cols * 3.8, n_rows * 3.5),
                             facecolor="white")
    axes = np.array(axes).reshape(n_rows, n_cols)

    for i, combo in enumerate(combos):
        r, c = i // n_cols, i % n_cols
        ax = axes[r, c]
        d = _model_dir_3ch(combo, loss) if model_type == "3ch" \
            else _model_dir_4ch(combo, loss)
        img = _load_png(d / "eval" / f"{umap_key}.png")
        if img is not None:
            ax.imshow(img)
        else:
            ax.text(0.5, 0.5, "pending", ha="center", va="center",
                    fontsize=12, color="gray", transform=ax.transAxes)
        ax.set_title(_combo_label(combo) if model_type == "3ch"
                     else f"{VARIANT_4CH_LABEL[combo]}",
                     fontsize=10, fontweight="bold", pad=3)
        ax.axis("off")

    for i in range(len(combos), n_rows * n_cols):
        axes[i // n_cols, i % n_cols].axis("off")

    fig.tight_layout(pad=0.5)
    return _fig_to_bytes(fig)


def _mean_nl1_table(prs, loss: str, color: RGBColor) -> None:
    """Heat-mapped mean nL1 table: combos × datasets."""
    rows_data = []
    for group in GROUPS_3CH:
        for combo in group:
            d = _model_dir_3ch(combo, loss)
            csv = d / "cross_dataset_recon_metrics.csv"
            row = {"combo": combo}
            if csv.exists():
                df = pd.read_csv(csv)
                if "recon_nl1" in df.columns and "dataset" in df.columns:
                    for ds_key, ds_label in [("vinc","ds1"),("pfak","ds2"),
                                              ("ppax","ds3"),("nih3t3","ds4")]:
                        sub = df[df["dataset"] == ds_key]["recon_nl1"]
                        row[ds_label] = float(sub.mean()) if len(sub) else None
            rows_data.append(row)

    slide = _add_slide(prs)
    _slide_header(slide,
                  f"3ch pza — Mean nL1  ({LOSS_LABEL[loss]} training)",
                  subtitle="Lower = better  ·  bold = training dataset  ·  — = not evaluated yet",
                  color=color)

    ds_cols = ["ds1", "ds2", "ds3", "ds4"]
    col_w   = [2.8] + [2.1] * 4
    col_x: list[float] = []
    x = MARGIN
    for cw in col_w:
        col_x.append(x); x += cw
    headers = ["Combo"] + ds_cols
    row_h = 0.33; y0 = HEADER_H

    for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        box = slide.shapes.add_shape(1, _px(cx), _px(y0), _px(cw), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, hdr, cx+0.04, y0+0.04, cw-0.08, row_h-0.06,
                 font_size=9, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)

    all_vals = [r[d] for r in rows_data for d in ds_cols if r.get(d) is not None]
    vmin = min(all_vals) if all_vals else 0.0
    vmax = max(all_vals) if all_vals else 1.0

    def _heat(val):
        t = (val - vmin) / max(vmax - vmin, 1e-6)
        return RGBColor(int(255*(0.2+0.6*t)), int(255*(0.9-0.5*t)), int(255*(0.9-0.7*t)))

    for i, row in enumerate(rows_data):
        combo = row["combo"]
        parts = set(combo)
        y = y0 + row_h * (i + 1)
        box = slide.shapes.add_shape(1, _px(col_x[0]), _px(y), _px(col_w[0]), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, _combo_label(combo),
                 col_x[0]+0.04, y+0.04, col_w[0]-0.08, row_h-0.06,
                 font_size=8, color=C_TITLE)
        for ci, ds in enumerate(ds_cols, start=1):
            val = row.get(ds)
            cx, cw = col_x[ci], col_w[ci]
            box = slide.shapes.add_shape(1, _px(cx), _px(y), _px(cw), _px(row_h))
            box.fill.background(); box.line.color.rgb = C_LGRAY
            txt = f"{val:.3f}" if val is not None else "—"
            ds_key = {"ds1":"v","ds2":"f","ds3":"p","ds4":"n"}[ds]
            is_train = ds_key in parts
            _textbox(slide, txt, cx+0.04, y+0.04, cw-0.08, row_h-0.06,
                     font_size=8, bold=is_train, color=C_TITLE, align=PP_ALIGN.CENTER)


# ── slide builders ─────────────────────────────────────────────────────────────

def _slide_title(prs):
    slide = _add_slide(prs)
    _textbox(slide, "Multi-Channel ConAE — Dataset Combination Sweep",
             1.0, 1.2, W_IN-2.0, 1.0,
             font_size=34, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    _textbox(slide, "3ch (pax+zyx+act)  ·  15 dataset combinations  ·  3 training losses",
             1.0, 2.3, W_IN-2.0, 0.5,
             font_size=18, color=C_3CH, align=PP_ALIGN.CENTER)
    _textbox(slide, "4ch (vinc+pax+zyx+act)  ·  ds1 vs ds1+ds4 vs ds4 training",
             1.0, 2.9, W_IN-2.0, 0.5,
             font_size=18, color=C_4CH, align=PP_ALIGN.CENTER)
    items = [
        "ds1=vinc  ·  ds2=pfak  ·  ds3=ppax  ·  ds4=nih3t3",
        "All models: ConAE nL1/MSE/L1 · enlcrop sc2 · lat 12/proj 8 · 500 epochs",
        "Trained on vinc (ds1) or multi-dataset combination, evaluated on all 4",
        "Observation: nL1 scale is dataset-intrinsic (intensity characteristics), "
        "independent of training set",
    ]
    y = 3.7
    for item in items:
        _textbox(slide, f"• {item}", 1.5, y, W_IN-3.0, 0.42, font_size=13, color=C_GRAY)
        y += 0.44


def _slide_overview(prs):
    slide = _add_slide(prs)
    _slide_header(slide, "Experiment Overview — 3ch pza: 15 × 3 = 45 models",
                  subtitle="pax+zyx+act  ·  enlcrop sc2  ·  lat 12/proj 8  ·  λ_c=0.03")

    all_combos = SINGLES + PAIRS_V + PAIRS_NOV + TRIPLES + ALL_DS
    col_x = [MARGIN, 1.5, 3.0, 4.5, 6.0]
    col_w = [1.3,    1.4, 1.4, 1.4, 1.4]
    headers = ["#", "Combo", "ds1", "ds2/ds3", "ds4"]
    row_h = 0.36; y0 = HEADER_H

    for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        box = slide.shapes.add_shape(1, _px(cx), _px(y0), _px(cw), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, hdr, cx+0.04, y0+0.05, cw-0.08, row_h-0.08,
                 font_size=9, bold=True, color=C_TITLE)

    for i, combo in enumerate(all_combos):
        y = y0 + row_h * (i + 1)
        parts = set(combo)
        row_vals = [
            str(i+1), _combo_label(combo),
            "✓" if "v" in parts else "—",
            "+".join(DS_LETTER[d] for d in ["f","p"] if d in parts) or "—",
            "✓" if "n" in parts else "—",
        ]
        for ci, (val, cw, cx) in enumerate(zip(row_vals, col_w, col_x)):
            box = slide.shapes.add_shape(1, _px(cx), _px(y), _px(cw), _px(row_h))
            box.fill.background(); box.line.color.rgb = C_LGRAY
            col = C_3CH if ci == 1 else C_GRAY
            _textbox(slide, val, cx+0.04, y+0.04, cw-0.08, row_h-0.06,
                     font_size=8, color=col)

    # 4ch panel on the right
    rx = 7.5
    _textbox(slide, "4ch vinc variants", rx, HEADER_H, 5.5, 0.38,
             font_size=13, bold=True, color=C_4CH)
    rows_4ch = [
        ("v",  "vinc (ds1) only",    "ds1"),
        ("nv", "vinc+nih3t3 (ds1+ds4)", "ds1+ds4"),
        ("n",  "nih3t3 (ds4) only",  "ds4"),
    ]
    y = HEADER_H + 0.45
    for key, desc, train in rows_4ch:
        _textbox(slide, f"{VARIANT_4CH_LABEL[key]}:  {desc}",
                 rx, y, 5.5, 0.38, font_size=10, color=C_TITLE)
        y += 0.42

    _textbox(slide, "All 4ch: 3 variants × 3 losses = 9 models\n"
             "Channel order: vinc · pax · zyx · act",
             rx, y+0.1, 5.5, 0.6, font_size=10, italic=True, color=C_LGRAY)


def _slide_nl1_observation(prs):
    """Text slide explaining the dataset-intrinsic nL1 scale phenomenon."""
    slide = _add_slide(prs)
    _slide_header(slide,
                  "Observation: nL1 scale is dataset-intrinsic, not training-dependent",
                  color=C_NOTE)

    _textbox(slide,
             "What we see in the violin plots:",
             MARGIN, 0.85, W_IN-0.5, 0.38,
             font_size=14, bold=True, color=C_TITLE)

    observations = [
        "ds1 (vinc) consistently shows LOW nL1  regardless of which datasets were used for training",
        "ds4 (nih3t3) consistently shows HIGH nL1  regardless of training set",
        "Adding ds4 to training reduces its nL1 somewhat, but the ranking ds1 < ds2 ≈ ds3 < ds4 persists",
    ]
    y = 1.3
    for obs in observations:
        _textbox(slide, f"• {obs}", MARGIN+0.2, y, W_IN-0.7, 0.45,
                 font_size=12, color=C_TITLE)
        y += 0.48

    _textbox(slide, "Why:", MARGIN, y+0.1, W_IN-0.5, 0.35,
             font_size=14, bold=True, color=C_NOTE)
    y += 0.52

    reasons = [
        "nL1 = L1 / mean|raw|  — the denominator (mean patch intensity) varies across datasets",
        "ds1 (vinc): bright, high-contrast paxillin signal → large mean|raw| → small nL1 even if absolute L1 is similar",
        "ds4 (nih3t3): different cell type, lower paxillin signal intensity → smaller mean|raw| → inflated nL1",
        "This is an intensity normalisation artefact, not a true difference in reconstruction quality",
    ]
    for reason in reasons:
        _textbox(slide, f"• {reason}", MARGIN+0.2, y, W_IN-0.7, 0.45,
                 font_size=12, color=C_GRAY)
        y += 0.48

    _textbox(slide,
             "→  Compare models within the same dataset column, not across datasets.\n"
             "→  For cross-dataset comparison use raw L1 (MAE) or MSE instead of nL1.",
             MARGIN, y+0.1, W_IN-0.5, 0.7,
             font_size=12, bold=True, color=C_NOTE)


def _slide_violin_group(prs, combos, group_label, loss, color, model_type="3ch"):
    slide = _add_slide(prs)
    _slide_header(slide,
                  f"{'3ch pza' if model_type=='3ch' else '4ch vinc'} — {LOSS_LABEL[loss]} training  ·  {group_label}  ·  nL1 eval",
                  subtitle="Each panel: cross-dataset violin (pfak/ppax/nih3t3 test, vinc train if trained on it)",
                  color=color)
    img = _violin_grid(combos, loss, model_type=model_type)
    _add_image_bytes(slide, img, MARGIN, HEADER_H+0.05, width=W_IN-2*MARGIN)


def _slide_umap_group(prs, combos, group_label, loss, color,
                      umap_key="umap_annotation", model_type="3ch"):
    slide = _add_slide(prs)
    key_label = {"umap_annotation": "FA annotation",
                 "umap_condition":  "condition",
                 "umap_kmeans":     "KMeans k=10"}.get(umap_key, umap_key)
    _slide_header(slide,
                  f"{'3ch pza' if model_type=='3ch' else '4ch vinc'} — UMAP ({key_label})  ·  {group_label}",
                  subtitle=f"{LOSS_LABEL[loss]} training",
                  color=color)
    img = _umap_grid(combos, loss, umap_key=umap_key, model_type=model_type)
    _add_image_bytes(slide, img, MARGIN, HEADER_H+0.05, width=W_IN-2*MARGIN)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path,
                        default=Path("slides_multichannel_ds_sweep.pptx"))
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    print("  Title …"); _slide_title(prs)
    print("  Overview …"); _slide_overview(prs)

    # ── 3ch pza section ───────────────────────────────────────────────────────
    _section_slide(prs,
                   "3ch — paxillin + zyxin + actin",
                   "15 dataset combinations  ·  3 training losses  ·  lat 12/proj 8",
                   C_3CH)

    # nL1 mean table
    print("  Mean nL1 table (nL1 training) …")
    _mean_nl1_table(prs, "nl1", C_3CH)

    # Violin slides — nL1 training only (most interpretable)
    for combos, glbl in zip(GROUPS_3CH, GROUP_LABELS):
        print(f"  Violin [nL1] {glbl} …")
        _slide_violin_group(prs, combos, glbl, "nl1", C_3CH)

    # UMAP slides — nL1 training
    for combos, glbl in zip(GROUPS_3CH, GROUP_LABELS):
        print(f"  UMAP {glbl} …")
        _slide_umap_group(prs, combos, glbl, "nl1", C_3CH,
                          umap_key="umap_annotation")

    # ── 4ch vinc section ──────────────────────────────────────────────────────
    _section_slide(prs,
                   "4ch — vinculin + paxillin + zyxin + actin",
                   "ds1 only  ·  ds1+ds4  ·  ds4 only  ·  nL1 training",
                   C_4CH)

    print("  4ch violin (all variants, nL1) …")
    _slide_violin_group(prs, VARIANTS_4CH, "ds1 vs ds1+ds4 vs ds4", "nl1",
                        C_4CH, model_type="4ch")

    print("  4ch UMAP (all variants, nL1) …")
    _slide_umap_group(prs, VARIANTS_4CH, "ds1 vs ds1+ds4 vs ds4", "nl1",
                      C_4CH, umap_key="umap_annotation", model_type="4ch")

    # ── observation slide ─────────────────────────────────────────────────────
    _section_slide(prs,
                   "Why nL1 is always small on ds1 and large on ds4",
                   "Dataset-intrinsic intensity characteristics",
                   C_NOTE)
    print("  nL1 observation slide …")
    _slide_nl1_observation(prs)

    prs.save(str(args.out))
    print(f"\nSaved → {args.out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
