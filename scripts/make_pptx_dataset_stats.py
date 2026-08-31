#!/usr/bin/env python3
"""
make_pptx_dataset_stats.py

Generates a PPT summarising DS1 B2 / DS1 B12 / DS2 B12 label-set statistics:
  - Overview comparison table
  - Per-fold breakdown (each dataset)
  - Per-condition breakdown (each dataset)
  - Budget lists
"""
import argparse
import io
import re
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
DATA = Path("/net/projects/CLS/lding/data/fa_data_analysis")
REPO = Path(__file__).resolve().parents[1]
OUT  = REPO / "results" / "dataset_stats.pptx"

DATASETS = [
    {
        "name":       "DS1 B2",
        "ann_dir":    DATA / "labelling" / "le_b2_supcon",
        "fs_file":    "fold_splits_ds1.csv",
        "job_re_pat": r"le_b2_ds1_fv(\d)_nb(\w+)_r(\d)",
        "channel":    "PAX (ch1)",
        "drug":       "Vincristine",
        "conditions": "ctrl + ycomp",
        "annotator":  "B2 (Annabel)",
        "patch_px":   "32×32",
        "n_folds":    5,
        "n_repeats":  5,
        "classifier": "LGBM / logreg",
        "note":       "cold-start SupCon-AE per budget",
    },
    {
        "name":       "DS1 B12",
        "ann_dir":    DATA / "labelling" / "le_b12_supcon",
        "fs_file":    "fold_splits_ds1.csv",
        "job_re_pat": r"le_b12_ds1_fv(\d)_nb(\w+)_r(\d)",
        "channel":    "PAX (ch1)",
        "drug":       "Vincristine",
        "conditions": "ctrl + ycomp",
        "annotator":  "B1 (Margaret) + B2 (Annabel)",
        "patch_px":   "32×32",
        "n_folds":    5,
        "n_repeats":  5,
        "classifier": "LGBM / logreg",
        "note":       "53 conflicts dropped, B2 priority",
    },
    {
        "name":       "DS2 B12",
        "ann_dir":    DATA / "labelling" / "le_b12_supcon",
        "fs_file":    "fold_splits_ds2.csv",
        "job_re_pat": r"le_b12_ds2_fv(\d)_nb(\w+)_r(\d)",
        "channel":    "FAK (ch2)",
        "drug":       "—",
        "conditions": "ctrl only",
        "annotator":  "B1 (Margaret) + B2 (Annabel)",
        "patch_px":   "32×32",
        "n_folds":    5,
        "n_repeats":  5,
        "classifier": "LGBM / logreg",
        "note":       "pfak channel; heavily imbalanced",
    },
]

# ---------------------------------------------------------------------------
# Slide geometry
# ---------------------------------------------------------------------------
SW = Inches(13.33)
SH = Inches(7.5)
BG = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_H    = Inches(0.75)
SUBTITLE_H = Inches(0.35)
BODY_TOP   = Inches(1.2)

COL_HEADER = RGBColor(0x1f, 0x77, 0xb4)   # blue
COL_AD     = RGBColor(0x22, 0x8B, 0x22)   # green
COL_NOAD   = RGBColor(0x7b, 0x52, 0xab)   # purple
COL_TOTAL  = RGBColor(0x33, 0x33, 0x33)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def _txbox(slide, text, left, top, width, height,
           fontsize=11, bold=False, color=RGBColor(0,0,0),
           align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _header(slide, title, subtitle=""):
    _txbox(slide, title,
           Inches(0.3), Inches(0.08), SW - Inches(0.6), TITLE_H,
           fontsize=20, bold=True, color=RGBColor(0x1f, 0x77, 0xb4))
    if subtitle:
        _txbox(slide, subtitle,
               Inches(0.3), Inches(0.72), SW - Inches(0.6), SUBTITLE_H,
               fontsize=9, color=RGBColor(0x55, 0x55, 0x55))


def _fig_to_slide(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    slide.shapes.add_picture(buf, left, top, width, height)


# ---------------------------------------------------------------------------
# Load all dataset info
# ---------------------------------------------------------------------------
def _load(cfg):
    fs = pd.read_csv(cfg["ann_dir"] / cfg["fs_file"])
    fs["condition"] = fs["unique_ID"].str.extract(r"^([^-]+)-")[0]

    jobs = sorted(cfg["ann_dir"].glob("*_fv0_nb*_r0.csv"))
    budgets = sorted(
        {p.stem.split("_nb")[1].split("_r")[0] for p in jobs if "nball" not in p.stem},
        key=lambda x: int(x) if x.isdigit() else 9999,
    )

    # --- per-fold test counts ---
    fold_test = fs.groupby("fold")["label"].value_counts().unstack(fill_value=0)
    for col in ("adhesion", "No adhesion"):
        if col not in fold_test.columns:
            fold_test[col] = 0

    # --- per-fold train counts (all patches NOT in that fold) ---
    train_rows = []
    for fld in sorted(fs["fold"].unique()):
        train = fs[fs["fold"] != fld]
        train_rows.append({
            "fold": fld,
            "train_ad":   int((train["label"] == "adhesion").sum()),
            "train_noad": int((train["label"] == "No adhesion").sum()),
            "test_ad":    int(fold_test.loc[fld, "adhesion"]),
            "test_noad":  int(fold_test.loc[fld, "No adhesion"]),
        })
    fold_split_tab = pd.DataFrame(train_rows).set_index("fold")
    fold_split_tab["train_total"] = fold_split_tab["train_ad"] + fold_split_tab["train_noad"]
    fold_split_tab["test_total"]  = fold_split_tab["test_ad"]  + fold_split_tab["test_noad"]
    fold_split_tab["train_adh%"]  = (fold_split_tab["train_ad"] / fold_split_tab["train_total"] * 100).round(1)
    fold_split_tab["test_adh%"]   = (fold_split_tab["test_ad"]  / fold_split_tab["test_total"]  * 100).round(1)

    # --- per-budget train label counts (mean over folds × repeats) ---
    job_re_str = cfg.get("job_re_pat", "")
    JOB_RE = re.compile(job_re_str) if job_re_str else None
    budget_rows = []
    if JOB_RE:
        for jp in sorted(cfg["ann_dir"].glob("*.csv")):
            m = JOB_RE.search(jp.stem)
            if not m:
                continue
            fold, budget, repeat = int(m.group(1)), m.group(2), int(m.group(3))
            if budget == "all":
                continue
            df = pd.read_csv(jp)
            budget_rows.append(dict(
                fold=fold, budget=int(budget), repeat=repeat,
                train_ad=int((df["label"] == "adhesion").sum()),
                train_noad=int((df["label"] == "No adhesion").sum()),
            ))
    budget_df = pd.DataFrame(budget_rows)

    # Per-budget summary: mean train counts across folds+repeats; test counts per fold
    avg_test_ad   = fold_split_tab["test_ad"].mean()
    avg_test_noad = fold_split_tab["test_noad"].mean()
    budget_tab = None
    if len(budget_df):
        grp = budget_df.groupby("budget")[["train_ad", "train_noad"]].mean().round(1)
        grp["train_total"] = (grp["train_ad"] + grp["train_noad"]).round(1)
        grp["train_adh%"]  = (grp["train_ad"] / grp["train_total"] * 100).round(1)
        grp["test_ad"]     = round(avg_test_ad, 1)
        grp["test_noad"]   = round(avg_test_noad, 1)
        grp["test_total"]  = round(avg_test_ad + avg_test_noad, 1)
        grp["test_adh%"]   = round(avg_test_ad / (avg_test_ad + avg_test_noad) * 100, 1)
        budget_tab = grp

    # --- per-condition ---
    cond_tab = fs.groupby(["condition", "label"]).size().unstack(fill_value=0)
    for col in ("adhesion", "No adhesion"):
        if col not in cond_tab.columns:
            cond_tab[col] = 0
    cond_tab = cond_tab[["adhesion", "No adhesion"]]
    cond_tab["total"]  = cond_tab.sum(axis=1)
    cond_tab["adh %"]  = (cond_tab["adhesion"] / cond_tab["total"] * 100).round(1)

    n_adh  = int((fs["label"] == "adhesion").sum())
    n_noad = int((fs["label"] == "No adhesion").sum())
    total  = len(fs)

    return dict(
        fs=fs, fold_split_tab=fold_split_tab, cond_tab=cond_tab,
        budget_tab=budget_tab,
        n_adh=n_adh, n_noad=n_noad, total=total,
        budgets=budgets, **cfg,
    )


# ---------------------------------------------------------------------------
# Slide: overview comparison table
# ---------------------------------------------------------------------------
def _slide_overview(prs, all_info):
    slide = _blank(prs)
    _header(slide, "Dataset Overview — All Label Sets",
            "DS1 B2 · DS1 B12 · DS2 B12  |  all use 5-fold stratified CV × 5 repeats")

    cols = ["Dataset", "Conditions", "Annotator(s)",
            "Total patches", "adhesion", "No adhesion", "adh %",
            "Test / fold", "Patch size", "Budgets (n)"]

    rows = []
    for info in all_info:
        adh_pct = round(info["n_adh"] / info["total"] * 100, 1)
        test_fold = round(info["total"] / info["n_folds"])
        rows.append([
            info["name"],
            info["conditions"],
            info["annotator"],
            info["total"],
            info["n_adh"],
            info["n_noad"],
            f"{adh_pct}%",
            f"~{test_fold}",
            info["patch_px"],
            ", ".join(info["budgets"]),
        ])

    fig, ax = plt.subplots(figsize=(12.5, 2.5), facecolor="white")
    ax.axis("off")

    col_widths = [0.07, 0.09, 0.18,
                  0.08, 0.07, 0.08, 0.05, 0.07, 0.06, 0.15]

    tbl = ax.table(
        cellText=rows,
        colLabels=cols,
        cellLoc="center",
        loc="center",
        colWidths=col_widths,
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(8.5)

    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#cccccc")
        if r == 0:
            cell.set_facecolor("#1f77b4")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif r % 2 == 0:
            cell.set_facecolor("#f0f4fa")
        else:
            cell.set_facecolor("white")
        # highlight class count columns
        if r > 0:
            if c == 4:   # adhesion
                cell.get_text().set_color("#228B22")
                cell.get_text().set_fontweight("bold")
            elif c == 5:  # No adhesion
                cell.get_text().set_color("#7b52ab")
                cell.get_text().set_fontweight("bold")
        cell.set_height(0.28)

    fig.tight_layout(pad=0.2)
    _fig_to_slide(slide, fig,
                  Inches(0.15), Inches(1.15), Inches(13.0), Inches(2.8))

    # Notes below
    y = Inches(4.1)
    for info in all_info:
        _txbox(slide, f"  {info['name']}: {info['note']}",
               Inches(0.4), y, Inches(12.5), Inches(0.3),
               fontsize=8.5, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.28)


# ---------------------------------------------------------------------------
# Slide: per-fold + per-condition breakdown for one dataset
# ---------------------------------------------------------------------------
def _make_table(ax, data, col_labels, col_widths, fontsize=9.5,
                header_bg="#1f77b4", alt_bg="#f5f5f5", summary_rows=0,
                ad_cols=(), noad_cols=()):
    tbl = ax.table(cellText=data, colLabels=col_labels,
                   cellLoc="center", loc="center", colWidths=col_widths)
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(fontsize)
    n_data = len(data)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#cccccc")
        if r == 0:
            cell.set_facecolor(header_bg)
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif r > n_data - summary_rows:
            cell.set_facecolor("#dde8f8")
            cell.get_text().set_fontweight("bold")
        elif r % 2 == 0:
            cell.set_facecolor(alt_bg)
        else:
            cell.set_facecolor("white")
        if r > 0:
            if c in ad_cols:
                cell.get_text().set_color("#228B22")
            elif c in noad_cols:
                cell.get_text().set_color("#7b52ab")
        cell.set_height(0.13)
    return tbl


def _slide_detail(prs, info):
    slide = _blank(prs)
    _header(slide,
            f"{info['name']} — Fold Breakdown: Train / Test by Class",
            f"{info['annotator']} · {info['conditions']} · {info['total']} patches · "
            f"{info['n_adh']} adhesion / {info['n_noad']} No adhesion")

    fig, axes = plt.subplots(1, 2, figsize=(12.5, 4.8), facecolor="white")

    # ---- Left: per-fold train/test split table ----
    ax = axes[0]
    ax.axis("off")

    fst = info["fold_split_tab"].reset_index()
    # summary row (totals / means)
    table_data = []
    for _, r in fst.iterrows():
        table_data.append([
            str(int(r["fold"])),
            str(r["train_ad"]), str(r["train_noad"]), str(r["train_total"]), f"{r['train_adh%']}%",
            str(r["test_ad"]),  str(r["test_noad"]),  str(r["test_total"]),  f"{r['test_adh%']}%",
        ])
    # totals row
    table_data.append([
        "ALL",
        str(info["n_adh"]), str(info["n_noad"]), str(info["total"]), f"{info['n_adh']/info['total']*100:.1f}%",
        "—", "—", "—", "—",
    ])

    _make_table(
        ax, table_data,
        col_labels=["Fold",
                    "tr-adh", "tr-no-adh", "tr-total", "tr-adh%",
                    "te-adh",  "te-no-adh",  "te-total",  "te-adh%"],
        col_widths=[0.09, 0.12, 0.14, 0.12, 0.10, 0.12, 0.14, 0.12, 0.10],
        fontsize=9,
        summary_rows=1,
        ad_cols=(1, 5), noad_cols=(2, 6),
    )
    ax.set_title("Per-Fold Train / Test Split", fontsize=11, fontweight="bold", pad=8)

    # ---- Right: per-condition table ----
    ax2 = axes[1]
    ax2.axis("off")

    ct = info["cond_tab"].reset_index()
    ct.columns = ["Condition", "adhesion", "No adhesion", "Total", "adh %"]
    cond_data = [[str(r["Condition"]), str(r["adhesion"]), str(r["No adhesion"]),
                  str(r["Total"]), f"{r['adh %']}%"] for _, r in ct.iterrows()]

    _make_table(
        ax2, cond_data,
        col_labels=["Condition", "adhesion", "No adhesion", "Total", "adh %"],
        col_widths=[0.22, 0.22, 0.26, 0.16, 0.14],
        fontsize=10,
        ad_cols=(1,), noad_cols=(2,),
    )
    ax2.set_title("Per-Condition Breakdown", fontsize=11, fontweight="bold", pad=8)

    fig.tight_layout(pad=1.0)
    _fig_to_slide(slide, fig,
                  Inches(0.1), Inches(1.1), Inches(13.1), Inches(5.2))


def _slide_budget_tab(prs, info):
    """Per-budget train-ad / train-noad / test-ad / test-noad table."""
    if info["budget_tab"] is None:
        return
    slide = _blank(prs)
    _header(slide,
            f"{info['name']} — Per-Budget Label Counts",
            "Mean over 5 folds × 5 repeats · train counts = sampled budget · "
            "test counts = avg held-out fold (fixed)")

    bt = info["budget_tab"].reset_index()
    table_data = []
    for _, r in bt.iterrows():
        table_data.append([
            str(int(r["budget"])),
            f"{r['train_ad']:.1f}",  f"{r['train_noad']:.1f}",
            f"{r['train_total']:.1f}", f"{r['train_adh%']:.1f}%",
            f"{r['test_ad']:.1f}",   f"{r['test_noad']:.1f}",
            f"{r['test_total']:.1f}", f"{r['test_adh%']:.1f}%",
        ])

    fig, ax = plt.subplots(figsize=(11.5, len(table_data) * 0.42 + 1.2), facecolor="white")
    ax.axis("off")

    _make_table(
        ax, table_data,
        col_labels=["nb (budget)",
                    "train-adh", "train-no-adh", "train-total", "train-adh%",
                    "test-adh",  "test-no-adh",  "test-total",  "test-adh%"],
        col_widths=[0.11, 0.11, 0.13, 0.11, 0.10, 0.11, 0.13, 0.11, 0.09],
        fontsize=9.5,
        ad_cols=(1, 5), noad_cols=(2, 6),
    )

    fig.tight_layout(pad=0.5)
    h = min(Inches(6.0), Inches(len(table_data) * 0.42 + 1.4))
    _fig_to_slide(slide, fig, Inches(0.6), Inches(1.1), Inches(12.1), h)


# ---------------------------------------------------------------------------
# Slide: budget lists + class bar chart
# ---------------------------------------------------------------------------
def _slide_budgets(prs, all_info):
    slide = _blank(prs)
    _header(slide, "Label Budgets & Class Distribution",
            "Training budgets (n patches drawn per fold × repeat) · bar = class balance per dataset")

    fig, axes = plt.subplots(1, 2, figsize=(12.0, 4.8), facecolor="white",
                             gridspec_kw={"width_ratios": [1.6, 1]})

    # ---- Left: budget table ----
    ax = axes[0]
    ax.axis("off")

    max_len = max(len(info["budgets"]) for info in all_info)
    rows = []
    for i, info in enumerate(all_info):
        row = [info["name"]] + info["budgets"] + [""] * (max_len - len(info["budgets"]))
        rows.append(row)

    col_labels = ["Dataset"] + [f"#{i+1}" for i in range(max_len)]
    col_w = [0.14] + [0.07] * max_len

    tbl = ax.table(
        cellText=rows,
        colLabels=col_labels,
        cellLoc="center",
        loc="center",
        colWidths=col_w,
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#cccccc")
        if r == 0:
            cell.set_facecolor("#1f77b4")
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
        elif r % 2 == 0:
            cell.set_facecolor("#f0f4fa")
        else:
            cell.set_facecolor("white")
        cell.set_height(0.18)
    ax.set_title("Training Budgets (n labeled patches)", fontsize=11,
                 fontweight="bold", pad=8)

    # ---- Right: stacked bar class balance ----
    ax2 = axes[1]
    names  = [info["name"]  for info in all_info]
    n_ad   = [info["n_adh"] for info in all_info]
    n_noad = [info["n_noad"] for info in all_info]
    totals = [info["total"] for info in all_info]

    x = np.arange(len(names))
    w = 0.45
    bars_noad = ax2.bar(x, n_noad, w, label="No adhesion", color="#7b52ab", alpha=0.85)
    bars_ad   = ax2.bar(x, n_ad,   w, bottom=n_noad, label="adhesion", color="#228B22", alpha=0.85)

    for xi, (ad, noad, tot) in enumerate(zip(n_ad, n_noad, totals)):
        ax2.text(xi, noad / 2, str(noad), ha="center", va="center",
                 fontsize=9, color="white", fontweight="bold")
        ax2.text(xi, noad + ad / 2, str(ad), ha="center", va="center",
                 fontsize=9, color="white", fontweight="bold")
        ax2.text(xi, tot + tot * 0.02, f"n={tot}", ha="center", va="bottom",
                 fontsize=9, fontweight="bold", color="#333333")

    ax2.set_xticks(x)
    ax2.set_xticklabels(names, fontsize=10)
    ax2.set_ylabel("Patch count", fontsize=10)
    ax2.set_title("Class Balance per Dataset", fontsize=11, fontweight="bold")
    ax2.legend(fontsize=9, loc="upper right")
    ax2.spines[["top", "right"]].set_visible(False)
    ax2.set_facecolor("white")
    ax2.grid(axis="y", color="#EEEEEE", linewidth=0.7)

    fig.tight_layout(pad=1.0)
    _fig_to_slide(slide, fig,
                  Inches(0.15), Inches(1.1), Inches(13.0), Inches(5.3))


# ---------------------------------------------------------------------------
# Slide: combined per-fold table (all datasets side by side)
# ---------------------------------------------------------------------------
def _slide_combined_folds(prs, all_info):
    slide = _blank(prs)
    _header(slide, "Per-Fold Patch Counts — All Datasets",
            "5-fold stratified CV · counts show test-fold size (≈ 1/5 total)")

    fig, axes = plt.subplots(1, 3, figsize=(12.5, 3.8), facecolor="white")

    for ax, info in zip(axes, all_info):
        ax.axis("off")
        fst = info["fold_split_tab"].reset_index()
        totals_row = ["ALL", info["n_adh"], info["n_noad"], info["total"],
                      f"{info['n_adh']/info['total']*100:.1f}%"]
        table_data = [[str(int(r["fold"])), str(r["test_ad"]),
                       str(r["test_noad"]), str(r["test_total"]),
                       f"{r['test_adh%']}%"] for _, r in fst.iterrows()]
        table_data.append(totals_row)

        tbl = ax.table(
            cellText=table_data,
            colLabels=["Fold", "adh", "no-adh", "Total", "adh %"],
            cellLoc="center",
            loc="center",
            colWidths=[0.15, 0.22, 0.25, 0.20, 0.18],
        )
        tbl.auto_set_font_size(False)
        tbl.set_fontsize(9.5)
        for (r, c), cell in tbl.get_celld().items():
            cell.set_edgecolor("#cccccc")
            if r == 0:
                cell.set_facecolor("#1f77b4")
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
            elif r == len(table_data):
                cell.set_facecolor("#dde8f8")
                cell.get_text().set_fontweight("bold")
            elif r % 2 == 0:
                cell.set_facecolor("#f5f5f5")
            else:
                cell.set_facecolor("white")
            if r > 0 and r <= len(fst):
                if c == 1:
                    cell.get_text().set_color("#228B22")
                elif c == 2:
                    cell.get_text().set_color("#7b52ab")
            cell.set_height(0.14)
        ax.set_title(info["name"], fontsize=12, fontweight="bold", pad=6)

    fig.tight_layout(pad=1.0)
    _fig_to_slide(slide, fig,
                  Inches(0.15), Inches(1.1), Inches(13.0), Inches(4.5))


# ---------------------------------------------------------------------------
# Title slide
# ---------------------------------------------------------------------------
def _slide_title(prs):
    slide = _blank(prs)
    _txbox(slide, "Label Efficiency Benchmark — Dataset Statistics",
           Inches(0.5), Inches(2.5), SW - Inches(1.0), Inches(1.0),
           fontsize=28, bold=True, color=RGBColor(0x1f, 0x77, 0xb4),
           align=PP_ALIGN.CENTER)
    _txbox(slide,
           "DS1 B2 · DS1 B12 · DS2 B12\n"
           "Patch counts · Fold splits · Condition breakdown · Label budgets",
           Inches(0.5), Inches(3.6), SW - Inches(1.0), Inches(0.9),
           fontsize=13, color=RGBColor(0x55, 0x55, 0x55),
           align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build(out_path: Path):
    print("Loading dataset info …")
    all_info = [_load(cfg) for cfg in DATASETS]

    prs = _prs()

    print("  Slide 1 — Title")
    _slide_title(prs)

    print("  Slide 2 — Overview comparison table")
    _slide_overview(prs, all_info)

    print("  Slide 3 — Combined per-fold counts")
    _slide_combined_folds(prs, all_info)

    for info in all_info:
        print(f"  Slide — {info['name']} detail (fold + condition)")
        _slide_detail(prs, info)
        print(f"  Slide — {info['name']} per-budget counts")
        _slide_budget_tab(prs, info)

    print("  Slide — Budgets & class balance")
    _slide_budgets(prs, all_info)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"\nSaved → {out_path}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(OUT))
    args = ap.parse_args()
    build(Path(args.out))
