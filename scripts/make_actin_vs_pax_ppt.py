#!/usr/bin/env python3
"""
make_actin_vs_pax_ppt.py

PPT comparing ConAE results on actin (ch3) vs paxillin (ch1).
Shows: reconstruction samples, UMAP (FA annotation + KMeans), cluster panels,
       cross-dataset violin plots.

Output: slides_actin_vs_pax.pptx  (in the current working directory)
"""
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── constants ──────────────────────────────────────────────────────────────────
RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")

ACTIN_MODELS = [
    ("baseline_vinc_only_ch3",
     "AE Baseline\n(actin ch3)"),
    ("contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_ch3",
     "ConAE nL1\n(actin ch3)"),
    ("contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_lc025_ch3",
     "ConAE nL1 λ=¼\n(actin ch3)"),
]

PAX_MODELS = [
    ("baseline_vinc_only_pax",
     "Baseline\nPax"),
    ("contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1",
     "ConAE nL1\n(pax ch1)"),
    ("contrastive_cio_rb_vinc_lat12proj8_enlcrop_sc2_nl1_lc025",
     "ConAE nL1 λ=¼\n(pax ch1)"),
]

CH2_MODELS = [
    ("baseline_vinc_2ch_pax_act",
     "AE Baseline\n(pax+actin 2ch)"),
    ("contrastive_cio_rb_vinc_lat12proj8_sc2_nl1_2ch_pax_act",
     "ConAE nL1\n(pax+actin 2ch)"),
    ("contrastive_cio_rb_vinc_lat12proj8_sc2_nl1_lc025_2ch_pax_act",
     "ConAE nL1 λ=¼\n(pax+actin 2ch)"),
]

W_IN, H_IN = 13.33, 7.5
MARGIN     = 0.25
IMG_W      = 9.6    # plot area width; text fills the rest
TEXT_X     = IMG_W + MARGIN + 0.1
TEXT_W     = W_IN - TEXT_X - MARGIN

C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE = RGBColor(0x1F, 0x2D, 0x3D)
C_ACT   = RGBColor(0x1A, 0x73, 0xE8)
C_PAX   = RGBColor(0xE8, 0x43, 0x1A)
C_2CH   = RGBColor(0x1B, 0x8A, 0x5A)
C_GRAY  = RGBColor(0x66, 0x66, 0x66)
C_LGRAY = RGBColor(0xAA, 0xAA, 0xAA)


# ── pptx helpers ──────────────────────────────────────────────────────────────

def _px(inches): return Inches(inches)

def _add_textbox(slide, text, left, top, width, height,
                 font_size=11, bold=False, color=C_TITLE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, bullets, left, top, width, height,
                 font_size=10, color=C_TITLE, head=None, head_color=None):
    """Add a text box with an optional bold heading and bullet list."""
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = True

    if head:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = head
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = head_color or color

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if (head or i > 0) else tf.paragraphs[0]
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def _add_image_bytes(slide, img_bytes, left, top, width=None, height=None):
    return slide.shapes.add_picture(
        io.BytesIO(img_bytes), _px(left), _px(top),
        width=_px(width) if width else None,
        height=_px(height) if height else None,
    )


def _fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_header(slide, title, color=C_TITLE, font_size=18):
    """Title bar: colored text + thin rule, all on white background."""
    _add_textbox(slide, title, MARGIN, 0.08, W_IN - 2*MARGIN, 0.5,
                 font_size=font_size, bold=True, color=color)
    rule = slide.shapes.add_shape(
        1, _px(MARGIN), _px(0.60), _px(W_IN - 2*MARGIN), _px(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def _note_panel(slide, head, bullets, color, font_size=10):
    """Right-side annotation panel with colored heading and bullet notes."""
    # light tinted background strip
    bg = slide.shapes.add_shape(
        1, _px(TEXT_X - 0.1), _px(0.65), _px(TEXT_W + 0.1), _px(H_IN - 0.75))
    bg.fill.solid()
    # very light tint: mix color toward white at 95%
    r = int(0.95 * 255 + 0.05 * color[0])
    g = int(0.95 * 255 + 0.05 * color[1])
    b = int(0.95 * 255 + 0.05 * color[2])
    bg.fill.fore_color.rgb = RGBColor(r, g, b)
    bg.line.fill.background()

    _add_bullets(slide, bullets, TEXT_X, 0.72, TEXT_W, H_IN - 0.85,
                 font_size=font_size, color=C_TITLE,
                 head=head, head_color=color)


# ── image loaders ─────────────────────────────────────────────────────────────

def _load_visual_frames(model_dir: Path, n_frames: int = 2) -> list:
    vis_path = model_dir / "recon" / "visual.tif"
    if not vis_path.exists():
        return []
    stack = tifffile.imread(str(vis_path))
    if stack.ndim == 3:
        stack = stack[np.newaxis]
    total = len(stack)
    idxs = np.linspace(5, total - 5, n_frames, dtype=int)
    frames = []
    for i in idxs:
        arr = stack[i]
        if arr.ndim == 2 or (arr.ndim == 3 and arr.shape[-1] == 1):
            if arr.ndim == 3:
                arr = arr[..., 0]
            lo, hi = arr.min(), arr.max()
            arr8 = np.clip((arr - lo) / (hi - lo + 1e-8) * 255, 0, 255).astype(np.uint8)
            arr = np.stack([arr8] * 3, axis=-1)
        frames.append(arr)
    return frames


def _load_cluster_grid(model_dir: Path) -> np.ndarray | None:
    clust_path = model_dir / "eval" / "cluster_panels" / "all_clusters.tif"
    if not clust_path.exists():
        return None
    stack = tifffile.imread(str(clust_path))
    n, ph, pw = stack.shape
    rows, cols = 2, 5
    gap = 4
    grid = np.zeros((rows * ph + (rows-1)*gap, cols * pw + (cols-1)*gap),
                    dtype=np.float32)
    for i in range(n):
        r, c = i // cols, i % cols
        y0, x0 = r * (ph + gap), c * (pw + gap)
        grid[y0:y0+ph, x0:x0+pw] = stack[i]
    return grid


# ── composite figure makers ───────────────────────────────────────────────────

def _make_recon_row(model_dirs_labels, n_frames=2) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(n_frames, n, figsize=(n * 4.0, n_frames * 2.2),
                             facecolor="white")
    if n == 1: axes = axes[:, np.newaxis]
    if n_frames == 1: axes = axes[np.newaxis, :]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        frames = _load_visual_frames(d, n_frames)
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=8,
                               fontweight="bold", pad=3)
        for row in range(n_frames):
            ax = axes[row, col]
            if row < len(frames):
                ax.imshow(frames[row])
            ax.axis("off")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_umap_grid(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(2, n, figsize=(n * 3.2, 5.6), facecolor="white")
    if n == 1: axes = axes[:, np.newaxis]
    row_labels = ["FA annotation", "KMeans k=10"]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=8,
                               fontweight="bold")
        for row, (key, rlbl) in enumerate(
                zip(["umap_annotation", "umap_kmeans_k10"], row_labels)):
            ax = axes[row, col]
            png_path = (d / "eval" / "cluster_panels" / "umap_kmeans_k10.png"
                        if key == "umap_kmeans_k10"
                        else d / "eval" / f"{key}.png")
            if png_path.exists():
                ax.imshow(np.array(Image.open(str(png_path))))
            ax.axis("off")
            if col == 0:
                ax.set_ylabel(rlbl, fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_umap_row(model_dirs_labels, key: str) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 3.2, 3.0), facecolor="white")
    if n == 1: axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        png_path = (d / "eval" / "cluster_panels" / "umap_kmeans_k10.png"
                    if key == "umap_kmeans_k10"
                    else d / "eval" / f"{key}.png")
        if png_path.exists():
            ax.imshow(np.array(Image.open(str(png_path))))
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_cluster_grid_fig(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 5.0, 3.0), facecolor="white")
    if n == 1: axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        grid = _load_cluster_grid(d)
        if grid is not None:
            ax.imshow(grid, cmap="gray", vmin=0, vmax=1)
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.suptitle("16 patches closest to each cluster centroid  (k=10, 2×5 grid)",
                 fontsize=8, y=0.02)
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_violin_row(model_dirs_labels, metric="recon_nl1") -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 4.8, 3.2), facecolor="white")
    if n == 1: axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        png_path = d / f"._cross_dataset_{metric}.png"
        if png_path.exists():
            ax.imshow(np.array(Image.open(str(png_path))))
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_violin_2ch_perchannel(model_dirs_labels,
                                row_metrics=None,
                                row_labels=None) -> bytes:
    """Grid: one row per metric, one column per model."""
    if row_metrics is None:
        row_metrics = ["recon_nl1_ch0", "recon_nl1_ch1"]
    if row_labels is None:
        row_labels = ["paxillin (ch0)", "actin (ch1)"]
    n   = len(model_dirs_labels)
    nr  = len(row_metrics)
    fig, axes = plt.subplots(nr, n, figsize=(n * 4.8, nr * 3.2), facecolor="white")
    if n == 1:  axes = axes[:, np.newaxis]
    if nr == 1: axes = axes[np.newaxis, :]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
        for row, (metric, rlbl) in enumerate(zip(row_metrics, row_labels)):
            ax = axes[row, col]
            png_path = d / f"._cross_dataset_{metric}.png"
            if png_path.exists():
                ax.imshow(np.array(Image.open(str(png_path))))
            ax.axis("off")
            if col == 0:
                ax.set_ylabel(rlbl, fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


# ── slide builders ────────────────────────────────────────────────────────────

def _content_slide(prs, title, color, img_bytes, note_head, note_bullets,
                   note_font=12):
    slide = _add_slide(prs)
    _slide_header(slide, title, color=color)
    _note_panel(slide, note_head, note_bullets, color, font_size=note_font)
    _add_image_bytes(slide, img_bytes, MARGIN, 0.68, width=IMG_W)
    return slide


def _build_title_slide(prs):
    slide = _add_slide(prs)

    # white background accent bar (left edge, colored)
    bar = slide.shapes.add_shape(1, 0, 0, _px(0.18), _px(H_IN))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()

    # blue/orange accent dots
    for y, c in [(1.8, C_ACT), (5.5, C_PAX)]:
        dot = slide.shapes.add_shape(9, _px(0.5), _px(y), _px(0.4), _px(0.4))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    _add_textbox(slide, "Contrastive Autoencoder",
                 0.8, 1.5, W_IN - 1.2, 1.0,
                 font_size=32, bold=True, color=C_TITLE)
    _add_textbox(slide, "Actin (ch3) vs Paxillin (ch1) vs Paxillin+Actin (2ch): Latent Space & Reconstruction Results",
                 0.8, 2.5, W_IN - 1.2, 0.7,
                 font_size=17, color=C_ACT)

    _add_textbox(slide,
                 "What we evaluate:",
                 0.8, 3.6, 5.5, 0.35,
                 font_size=12, bold=True, color=C_TITLE)
    items = [
        "Reconstruction fidelity  (raw vs reconstructed patches)",
        "Latent space organisation  (UMAP colored by FA type)",
        "Unsupervised structure  (KMeans k=10 cluster panels)",
        "Generalisation to unseen datasets  (cross-dataset L1 violin plots)",
    ]
    _add_bullets(slide, items, 0.8, 4.0, 7.5, 2.5, font_size=12, color=C_TITLE)

    _add_textbox(slide,
                 "Models: baseline AE · ConAE nL1 · ConAE nL1 λ=¼  |  Single-ch: EnlargedJitterCrop  ·  2ch: sc2 augmentation\n"
                 "Latent dim 12  |  Proj dim 8  |  500 epochs  |  Training data: vinc (ds1)",
                 0.8, 6.3, W_IN - 1.2, 0.8,
                 font_size=10, color=C_LGRAY, italic=True)


def _build_model_table(prs):
    slide = _add_slide(prs)
    _slide_header(slide, "Models Compared", color=C_TITLE)

    # Table data
    rows_data = [
        ["Channel",       "Model",           "Recon loss", "λ_contrast", "Aug",              "Notes"],
        ["actin (ch3)",   "AE Baseline",      "nL1",        "—",          "enlcrop sc2",      "No contrastive loss"],
        ["actin (ch3)",   "ConAE nL1",        "nL1",        "½",          "enlcrop sc2",      "EnlargedJitterCrop views"],
        ["actin (ch3)",   "ConAE nL1 λ=¼",   "nL1",        "¼",          "enlcrop sc2",      "Weaker contrastive pull"],
        ["pax (ch1)",     "ConAE nL1",        "nL1",        "½",          "enlcrop sc2",      "EnlargedJitterCrop views"],
        ["pax (ch1)",     "ConAE nL1 λ=¼",   "nL1",        "¼",          "enlcrop sc2",      "Weaker contrastive pull"],
        ["pax+act (2ch)", "AE Baseline",      "nL1",        "—",          "sc2",              "ch0=paxillin, ch1=actin"],
        ["pax+act (2ch)", "ConAE nL1",        "nL1",        "½",          "sc2",              "Both channels jointly encoded"],
        ["pax+act (2ch)", "ConAE nL1 λ=¼",   "nL1",        "¼",          "sc2",              "Weaker contrastive pull"],
    ]
    col_w = [1.5, 1.9, 1.1, 1.2, 1.2, 2.5]
    row_h = 0.45
    x0, y0 = MARGIN, 0.75

    for ri, row in enumerate(rows_data):
        x = x0
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            is_hdr = ri == 0
            is_act = ri in (1, 2, 3)
            is_pax = ri in (4, 5)
            is_2ch = ri in (6, 7, 8)
            if is_hdr:
                bg = slide.shapes.add_shape(
                    1, _px(x), _px(y0 + ri*row_h), _px(cw), _px(row_h))
                bg.fill.solid(); bg.fill.fore_color.rgb = C_TITLE
                bg.line.fill.background()
            txt_color = (C_WHITE if is_hdr
                         else C_ACT if is_act
                         else C_2CH if is_2ch
                         else C_PAX)
            _add_textbox(slide, cell, x+0.06, y0+ri*row_h+0.07,
                         cw-0.08, row_h-0.08,
                         font_size=9, bold=is_hdr, color=txt_color)
            x += cw

    # Right-side explanation
    _note_panel(slide, "Design choices", [
        "nL1 = normalised L1 loss: per-patch L1 divided by mean intensity, "
        "reducing bias toward bright patches",
        "λ_contrast ½ (enlcrop): EnlargedJitterCrop halves intensity scale, "
        "so ½× prevents contrastive term dominating reconstruction",
        "λ_contrast ¼: further reduces contrastive pull to emphasise reconstruction",
        "2ch models encode ch0 (paxillin) + ch1 (actin) jointly in 12 latent dims "
        "— both channels share the same latent code",
        "2ch uses standard sc2 augmentation (not enlcrop) so λ_contrast uses full scale",
        "All models: 500 epochs, Adam, cosine LR decay, latent dim 12, proj dim 8; "
        "evaluated on ds1 (vinc, train) + ds2 (pfak) + ds3 (ppax) + ds4 (nih3t3)",
    ], C_TITLE, font_size=11)


def _build_actin_recon(prs, act):
    img = _make_recon_row(act, n_frames=2)
    _content_slide(prs,
        "Actin (ch3)  —  Reconstruction  (raw left · recon right of each panel)",
        C_ACT, img,
        "What to look for",
        [
            "Baseline: encoder optimises reconstruction only — expect best fidelity",
            "ConAE: contrastive objective consumes some encoder capacity; "
            "expect slightly softer reconstruction vs baseline",
            "λ=¼ recovers some reconstruction detail relative to λ=½",
        ])


def _build_actin_umap(prs, act):
    img = _make_umap_grid(act)
    _content_slide(prs,
        "Actin (ch3)  —  UMAP  (top: FA annotation · bottom: KMeans k=10)",
        C_ACT, img,
        "Interpreting UMAP",
        [
            "Top: coloured by FA annotation; Bottom: KMeans k=10",
            "Cluster–annotation alignment (no labels used) = model captured FA structure",
            "ConAE contrastive loss should produce tighter, more separable clusters than baseline",
            "Actin signal is diffuse; expect softer FA-type boundaries than paxillin",
        ])


def _build_actin_clusters(prs, act):
    img = _make_cluster_grid_fig(act)
    _content_slide(prs,
        "Actin (ch3)  —  KMeans k=10 cluster panels  (2×5 grid, 16 patches per cluster)",
        C_ACT, img,
        "Cluster interpretation",
        [
            "16 patches nearest each cluster centroid — the model's representative examples",
            "Visual coherence within a panel = latent code captured that morphology",
            "N = cluster size. Small N = rare / transitional",
            "ConAE clusters tend to be more visually homogeneous than baseline",
        ])


def _build_actin_violin(prs, act):
    img = _make_violin_row(act)
    _content_slide(prs,
        "Actin (ch3)  —  Cross-dataset normalised L1  (ds1=vinc · ds2=pfak · ds3=ppax · ds4=nih3t3)",
        C_ACT, img,
        "Generalisation quality",
        [
            "ds1 (vinc) = train; ds2–4 = unseen cell lines",
            "Small train→test gap = model learned cell-line-agnostic features",
            "ds4 (nih3t3) = hardest domain shift",
            "λ=¼ may generalise better: less contrastive pressure, "
            "more focus on reconstruction",
        ])


def _build_pax_recon(prs, pax):
    img = _make_recon_row(pax, n_frames=2)
    _content_slide(prs,
        "Paxillin (ch1)  —  Reconstruction  (raw left · recon right of each panel)",
        C_PAX, img,
        "What to look for",
        [
            "Paxillin patches: high-contrast puncta on dark background — "
            "harder to reconstruct than diffuse actin signal",
            "Baseline: best fidelity; ConAE trades some detail for latent structure",
            "Spot blurring = latent code averaging sub-spot details; "
            "acceptable for representation, not for segmentation",
        ])


def _build_pax_umap(prs, pax):
    img = _make_umap_grid(pax)
    _content_slide(prs,
        "Paxillin (ch1)  —  UMAP  (top: FA annotation · bottom: KMeans k=10)",
        C_PAX, img,
        "Interpreting UMAP",
        [
            "Paxillin is a direct FA marker — expect cleaner class separation "
            "than actin",
            "Cluster–annotation alignment (unsupervised) = contrastive objective worked",
            "λ=¼: check whether lower contrastive pressure reduces cluster separation",
        ])


def _build_pax_clusters(prs, pax):
    img = _make_cluster_grid_fig(pax)
    _content_slide(prs,
        "Paxillin (ch1)  —  KMeans k=10 cluster panels  (2×5 grid, 16 patches per cluster)",
        C_PAX, img,
        "Cluster interpretation",
        [
            "16 patches nearest each centroid — the model's representative examples",
            "Spot-size / density coherence within a panel = model encoded FA maturation",
            "Mixed-appearance clusters = ambiguous or underrepresented morphologies",
            "Paxillin clusters should be more homogeneous than actin (higher specificity)",
        ])


def _build_pax_violin(prs, pax):
    img = _make_violin_row(pax)
    _content_slide(prs,
        "Paxillin (ch1)  —  Cross-dataset normalised L1  (ds1=vinc · ds2=pfak · ds3=ppax · ds4=nih3t3)",
        C_PAX, img,
        "Generalisation quality",
        [
            "ds1 (vinc) = train; ds2–4 = unseen cell lines",
            "Paxillin signal is more cell-line-specific than actin — "
            "expect a larger train→test gap than actin models",
            "ds4 (nih3t3) = hardest domain shift",
        ])


def _build_2ch_recon(prs, ch2):
    img = _make_recon_row(ch2, n_frames=2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  Reconstruction  (raw left · recon right of each panel)",
        C_2CH, img,
        "What to look for",
        [
            "Both channels encoded into one 12-dim latent; decoder reconstructs both",
            "If one channel dominates, reconstruction of the other visibly degrades",
            "Compare per-channel quality with single-channel models",
        ])


def _build_2ch_umap(prs, ch2):
    img = _make_umap_grid(ch2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  UMAP  (top: FA annotation · bottom: KMeans k=10)",
        C_2CH, img,
        "Interpreting UMAP",
        [
            "Does joint encoding sharpen FA-type separation vs single-channel?",
            "UMAP similarity to pax-only → pax dominates; to actin-only → actin dominates",
            "KMeans clusters may reflect joint morphotypes (spot size + fiber angle)",
        ])


def _build_2ch_clusters(prs, ch2):
    img = _make_cluster_grid_fig(ch2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  KMeans k=10 cluster panels  (pax top · actin bottom per cluster)",
        C_2CH, img,
        "Cluster interpretation",
        [
            "Each panel: paxillin row (top) + actin row (bottom)",
            "Good cluster = coherent in BOTH channels simultaneously",
            "Pax-only variation → pax drives the encoding; "
            "actin-only variation → actin contributes independently",
            "Compare morphotype count with single-channel panels",
        ])


def _build_2ch_violin(prs, ch2):
    img = _make_violin_row(ch2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  Cross-dataset normalised L1  (ds1=vinc · ds2=pfak · ds3=ppax · ds4=nih3t3)",
        C_2CH, img,
        "Generalisation quality",
        [
            "L1 averaged across both channels — harder generalisation target than single-ch",
            "Actin may regularise the encoder: lower ds2–4 error than pax-only?",
            "Compare train→test gap directly with single-channel violins",
        ])


def _build_2ch_violin_perchannel_nl1(prs, ch2):
    img = _make_violin_2ch_perchannel(
        ch2,
        row_metrics=["recon_nl1_ch0", "recon_nl1_ch1"],
        row_labels=["paxillin (ch0)  nL1", "actin (ch1)  nL1"],
    )
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  Per-channel normalised L1  (top: pax ch0 · bottom: actin ch1)",
        C_2CH, img,
        "Per-channel nL1",
        [
            "nL1 = L1 / mean|raw|: scale-invariant reconstruction error per channel",
            "Does one channel dominate the shared latent? "
            "Higher nL1 for that channel = latent code biased toward the other",
            "ConAE λ=¼: reduced contrastive pull may balance channel contribution",
            "Compare row-by-row with single-channel actin and paxillin violins",
        ])


def _build_2ch_violin_perchannel_l1(prs, ch2):
    img = _make_violin_2ch_perchannel(
        ch2,
        row_metrics=["recon_l1_ch0", "recon_l1_ch1"],
        row_labels=["paxillin (ch0)  L1", "actin (ch1)  L1"],
    )
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  Per-channel L1 (MAE)  (top: pax ch0 · bottom: actin ch1)",
        C_2CH, img,
        "Per-channel L1 (MAE)",
        [
            "Raw L1 (MAE) per channel at sc2-corrected intensity scale",
            "Compare absolute error magnitude between paxillin and actin channels",
            "Use together with nL1 slide to separate scale effects from reconstruction quality",
        ])


def _build_comparison(prs, act, pax, ch2):
    slide = _add_slide(prs)
    _slide_header(slide,
        "Actin vs Paxillin vs Pax+Actin  —  ConAE nL1: UMAP annotation  ·  Cross-dataset nL1",
        color=C_TITLE)

    compare = [act[1], pax[0], ch2[1]]  # ConAE nL1 for actin, pax, 2ch
    img_umap = _make_umap_row(compare, "umap_annotation")
    img_viol = _make_violin_row(compare, metric="recon_nl1")

    top_h   = 3.05
    bot_top = 0.68 + top_h + 0.05
    _add_image_bytes(slide, img_umap, MARGIN, 0.68, width=IMG_W)
    _add_image_bytes(slide, img_viol, MARGIN, bot_top, width=IMG_W)

    _note_panel(slide, "Key comparisons", [
        "All three: ConAE nL1 (λ=½) trained on vinc, same latent dim 12",
        "Paxillin: tighter FA-type UMAP separation; actin signal is diffuse",
        "2ch: does joint encoding sharpen separation over either single channel?",
        "Wider train→test nL1 gap = features more cell-line-specific",
        "If 2ch generalises better than pax-only: actin acts as a regulariser",
    ], C_TITLE, font_size=12)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    out_path = Path("slides_actin_vs_pax.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    act = [(RUNS / k, lbl) for k, lbl in ACTIN_MODELS]
    pax = [(RUNS / k, lbl) for k, lbl in PAX_MODELS]
    ch2 = [(RUNS / k, lbl) for k, lbl in CH2_MODELS]

    _SECTION_HEADERS = {
        "Actin section hdr":  ("Results: Actin (ch3)",         C_ACT,
                               "Baseline AE  ·  ConAE nL1  ·  ConAE nL1 λ=¼"),
        "Pax section hdr":    ("Results: Paxillin (ch1)",      C_PAX,
                               "ConAE nL1  ·  ConAE nL1 λ=¼"),
        "2ch section hdr":    ("Results: Paxillin + Actin (2ch)", C_2CH,
                               "AE Baseline  ·  ConAE nL1  ·  ConAE nL1 λ=¼"),
    }

    steps = [
        ("Title",              _build_title_slide,    (prs,)),
        ("Model table",        _build_model_table,    (prs,)),
        ("Actin section hdr",  None,                  None),
        ("Actin recon",        _build_actin_recon,    (prs, act)),
        ("Actin UMAP",         _build_actin_umap,     (prs, act)),
        ("Actin clusters",     _build_actin_clusters, (prs, act)),
        ("Actin violin",       _build_actin_violin,   (prs, act)),
        ("Pax section hdr",    None,                  None),
        ("Pax recon",          _build_pax_recon,      (prs, pax)),
        ("Pax UMAP",           _build_pax_umap,       (prs, pax)),
        ("Pax clusters",       _build_pax_clusters,   (prs, pax)),
        ("Pax violin",         _build_pax_violin,     (prs, pax)),
        ("2ch section hdr",    None,                  None),
        ("2ch recon",          _build_2ch_recon,      (prs, ch2)),
        ("2ch UMAP",           _build_2ch_umap,       (prs, ch2)),
        ("2ch clusters",       _build_2ch_clusters,   (prs, ch2)),
        ("2ch violin",         _build_2ch_violin,              (prs, ch2)),
        ("2ch violin pch nL1", _build_2ch_violin_perchannel_nl1, (prs, ch2)),
        ("2ch violin pch L1",  _build_2ch_violin_perchannel_l1,  (prs, ch2)),
        ("Comparison",         _build_comparison,              (prs, act, pax, ch2)),
    ]

    for label, fn, args in steps:
        print(f"  {label} …", flush=True)
        if fn is None:
            title_text, color, subtitle = _SECTION_HEADERS[label]
            slide = _add_slide(prs)
            _slide_header(slide, title_text, color=color, font_size=28)
            bar = slide.shapes.add_shape(
                1, _px(0), _px(H_IN * 0.82), _px(W_IN), _px(H_IN * 0.18))
            bar.fill.solid(); bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            _add_textbox(slide, subtitle, MARGIN, H_IN * 0.83,
                         W_IN - 2*MARGIN, 0.5,
                         font_size=14, color=C_WHITE, bold=False)
        else:
            fn(*args)

    prs.save(str(out_path))
    print(f"\nSaved → {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
