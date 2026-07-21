#!/usr/bin/env python3
"""
make_intensity_norm_pptx.py
===========================
Before/after CIO-RB normalization histograms, restricted to cell regions.

- After : all pixels from tiff_patches32_label (non-overlapping grid patches,
          in cell/FA region), aggregated per source image frame
- Before: raw uint16 PAX channel from CZI at the same patch (x, y) locations,
          aggregated per source image frame

Each curve = one source image (all its label patches aggregated).
Colors are drawn from a per-dataset colormap for visibility.

Run with PYTHONPATH for matplotlib + pptx:
  PYTHONPATH=".../core_env/lib/python3.11/site-packages" python3 scripts/make_intensity_norm_pptx.py
"""

from __future__ import annotations

import collections
import io
import re
import sys
from pathlib import Path

import czifile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT       = Path("/net/projects/CLS/lding/data/fa_data_analysis")
PATCH_BASE = ROOT / "ae_results/patches/cio_rb"
OUT_PPTX   = Path("intensity_normalization.pptx")

PAX_CH = 1  # CZI channel index for PAX in all 4 datasets
N_BINS = 300

# Shifted-log mapping parameters (must match intensity_transform.py defaults)
LOG_MAP_X_MIN  = -0.03   # lower reference → 0
LOG_MAP_X_REF  = 10.0    # upper reference → 1
LOG_MAP_DELTA  = 0.5     # compression strength

DATASETS = [
    {
        "key":   "vinc",
        "label": "ds1 - vinc",
        "cmap":  "Blues",
        "color": "#2E86C1",
        "conditions": [
            {
                "cond":      "control",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Control",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/vinc/control/tiff_patches32_label",
            },
            {
                "cond":      "ycomp",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Ycomp",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/vinc/ycomp/tiff_patches32_label",
            },
        ],
    },
    {
        "key":   "pfak",
        "label": "ds2 - pfak",
        "cmap":  "Reds",
        "color": "#C44E52",
        "conditions": [
            {
                "cond":      "control",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Control",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/pfak/control/tiff_patches32_label",
            },
            {
                "cond":      "ycomp",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Ycomp",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/pfak/ycomp/tiff_patches32_label",
            },
        ],
    },
    {
        "key":   "ppax",
        "label": "ds3 - ppax",
        "cmap":  "Greens",
        "color": "#55A868",
        "conditions": [
            {
                "cond":      "control",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Control",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/ppax/control/tiff_patches32_label",
            },
            {
                "cond":      "ycomp",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Y-comp",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/ppax/ycomp/tiff_patches32_label",
            },
        ],
    },
    {
        "key":   "nih3t3",
        "label": "ds4 - nih3t3",
        "cmap":  "Oranges",
        "color": "#DD8452",
        "conditions": [
            {
                "cond":      "control",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/Control",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/nih3t3/control/tiff_patches32_label",
            },
            {
                "cond":      "ycomp",
                "czi_dir":   ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/YCompound",
                "patch_dir": ROOT / "ae_results/patches/cio_rb/nih3t3/ycomp/tiff_patches32_label",
            },
        ],
    },
]

PATCH_RE = re.compile(r"_f(\d{4})x(\d{4})y(\d{4})ps(\d+)\.tif$")


def _parse_patch(name: str):
    m = PATCH_RE.search(name)
    return (int(m.group(1)), int(m.group(2)), int(m.group(3)), int(m.group(4))) if m else None


def _load_czi_pax(czi_path: Path) -> np.ndarray:
    """Return float32 (H, W) array: raw uint16 / 65535 for PAX channel."""
    arr = czifile.imread(str(czi_path))
    # shape: (1, 1, C, 1, H, W, 1)
    ch = arr[0, 0, PAX_CH, 0, :, :, 0]
    return ch.astype(np.float32) / 65535.0


def gather_frame_pixels(cond: dict) -> list[tuple[np.ndarray, np.ndarray, str]]:
    """
    Returns list of (before_pixels, after_pixels, label) per source frame.
    before_pixels: raw uint16/65535 at patch locations in CZI
    after_pixels : CIO-RB normalized values from tiff_patches32_label
    """
    patch_dir = Path(cond["patch_dir"])
    czi_dir   = Path(cond["czi_dir"])

    if not patch_dir.exists():
        print(f"    SKIP patch_dir not found: {patch_dir}")
        return []
    if not czi_dir.exists():
        print(f"    SKIP czi_dir not found: {czi_dir}")
        return []

    czi_files = sorted(czi_dir.glob("*.czi"))
    if not czi_files:
        print(f"    SKIP no CZI in {czi_dir}")
        return []

    # Group patch paths by frame index
    frame_patches: dict[int, list] = collections.defaultdict(list)
    for pf in sorted(patch_dir.glob("*.tif")):
        parsed = _parse_patch(pf.name)
        if parsed:
            fidx, x, y, ps = parsed
            frame_patches[fidx].append((pf, x, y, ps))

    results = []
    for fidx in sorted(frame_patches.keys()):
        if fidx >= len(czi_files):
            print(f"    WARN frame {fidx} >= {len(czi_files)} CZI files, skipping")
            continue

        patches = frame_patches[fidx]
        czi_path = czi_files[fidx]

        # After: load all patches -> float32 pixels
        after_parts = []
        for pf, x, y, ps in patches:
            img = tifffile.imread(str(pf)).astype(np.float32).ravel()
            after_parts.append(img)
        after_pixels = np.concatenate(after_parts)

        # Before: extract same regions from raw CZI
        try:
            raw = _load_czi_pax(czi_path)
            H, W = raw.shape
            before_parts = []
            for pf, x, y, ps in patches:
                crop = raw[y: min(y + ps, H), x: min(x + ps, W)]
                before_parts.append(crop.ravel())
            before_pixels = np.concatenate(before_parts)
        except Exception as e:
            print(f"    WARN CZI load failed {czi_path.name}: {e}")
            continue

        results.append((before_pixels, after_pixels, f"{cond['cond']}_f{fidx:04d}"))

    return results


def _plot_curves(ax, arrays: list[np.ndarray], alpha: float, lw: float,
                 bold_color: str, bins=None, n_bins=None):
    """
    Plot one histogram curve per array, each in a unique rainbow color.
    bold_color is used for the bold mean line.
    Returns (bins, mean_hist) so the caller can reuse them for summary plots.
    """
    n = len(arrays)
    if n == 0:
        return None, None
    if n_bins is None:
        n_bins = N_BINS
    cmap = plt.get_cmap("gist_rainbow")
    all_vals = np.concatenate(arrays)
    if bins is None:
        xmin, xmax = float(all_vals.min()), float(all_vals.max())
        bins = np.linspace(xmin, xmax, n_bins + 1)
    n_bins = len(bins) - 1
    mean_hist = np.zeros(n_bins, dtype=np.float64)

    for i, arr in enumerate(arrays):
        color = cmap(i / max(n - 1, 1))
        counts, _ = np.histogram(arr, bins=bins, density=True)
        ax.plot(bins[:-1], counts, color=color, alpha=alpha, lw=lw)
        mean_hist += counts

    mean_hist /= n
    ax.plot(bins[:-1], mean_hist, color=bold_color, alpha=0.9, lw=2.5, label="mean")
    return bins, mean_hist


def make_hist_figure(ds: dict, frames_data: list) -> plt.Figure:
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    fig.suptitle(ds["label"], fontsize=15, fontweight="bold")

    before_arrays = [b for b, a, _ in frames_data]
    after_arrays  = [a for b, a, _ in frames_data]

    _plot_curves(axes[0], before_arrays, alpha=0.35, lw=0.9, bold_color=ds["color"])
    axes[0].set_title("Before normalization  (raw PAX, uint16 / 65535)")
    axes[0].set_xlabel("Pixel intensity (0-1 scale)")
    axes[0].set_ylabel("Density")
    axes[0].axvline(1.0, color="red", ls="--", lw=1.5, label="sigmoid ceiling (1.0)")
    axes[0].legend(fontsize=8)

    _plot_curves(axes[1], after_arrays, alpha=0.35, lw=0.9, bold_color=ds["color"])
    axes[1].set_title("After CIO-RB normalization  (tiff_patches32_label)")
    axes[1].set_xlabel("Normalized pixel value")
    axes[1].set_ylabel("Density")
    axes[1].set_xlim(-0.03, 2.0)
    axes[1].axvline(1.0, color="red", ls="--", lw=1.5, label="sigmoid ceiling (1.0)")
    axes[1].legend(fontsize=8)

    n_frames = len(frames_data)
    fig.text(
        0.5, 0.01,
        f"{n_frames} source images  |  each curve = one image (label patches aggregated)  |  bold = mean",
        ha="center", fontsize=9, color="#555555",
    )
    plt.tight_layout(rect=[0, 0.04, 1, 1])
    return fig


def _draw_summary_axes(axes, summary: list[dict], xlim_before=None, xlim_after=None):
    """Fill two axes with dataset mean curves; optionally restrict x-axis."""
    for s in summary:
        axes[0].plot(s["bins_before"][:-1], s["before_mean"],
                     color=s["color"], lw=1.25, label=s["label"])
        axes[1].plot(s["bins_after"][:-1], s["after_mean"],
                     color=s["color"], lw=1.25, label=s["label"])

    panels = [
        (axes[0], "Before normalization  (raw PAX, uint16 / 65535)",
         "Pixel intensity (0-1 scale)", xlim_before),
        (axes[1], "After CIO-RB normalization  (tiff_patches32_label)",
         "Normalized pixel value", xlim_after),
    ]
    for ax, title, xlabel, xlim in panels:
        ax.set_title(title)
        ax.set_xlabel(xlabel)
        ax.set_ylabel("Density")
        ax.axvline(1.0, color="red", ls="--", lw=1.5, label="sigmoid ceiling (1.0)")
        ax.legend(fontsize=9)
        if xlim is not None:
            ax.set_xlim(xlim)


def make_summary_figure(summary: list[dict]) -> plt.Figure:
    """Full-range summary: before (left) and after (right) with all 4 dataset means."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    fig.suptitle("All datasets — mean PAX intensity distributions (full range)",
                 fontsize=14, fontweight="bold")
    _draw_summary_axes(axes, summary)
    fig.text(0.5, 0.01, "Mean over all source images per dataset",
             ha="center", fontsize=9, color="#555555")
    plt.tight_layout(rect=[0, 0.04, 1, 1])
    return fig


def make_summary_zoom_figure(summary: list[dict],
                              xlim_before=None,
                              xlim_after=(-0.03, 2.0)) -> plt.Figure:
    """
    Zoomed summary: 2-row x 2-col layout.
    Row 1 = full range (same as make_summary_figure).
    Row 2 = zoomed to the body of each distribution (2nd–98th percentile of means).
    """
    fig, axes = plt.subplots(2, 2, figsize=(14, 9))
    fig.suptitle("All datasets — mean PAX intensity distributions (full + zoomed)",
                 fontsize=14, fontweight="bold")

    # Full range (top row)
    _draw_summary_axes(axes[0], summary)
    axes[0, 0].set_title("Before norm  (full range)")
    axes[0, 1].set_title("After norm   (full range)")

    # Compute zoom range from the means (2nd-98th percentile of binned density)
    def _percentile_xlim(summary_key, bins_key, lo=1, hi=99):
        all_bins, all_means = [], []
        for s in summary:
            bins = s[bins_key][:-1]
            mean = s[summary_key]
            if mean is None:
                continue
            all_bins.append(bins)
            all_means.append(mean)
        if not all_bins:
            return None
        # Weight bins by mean density to find the occupied range
        bins_cat  = np.concatenate(all_bins)
        means_cat = np.concatenate(all_means)
        mask = means_cat > 0
        if mask.sum() == 0:
            return None
        xlo = np.percentile(bins_cat[mask], lo)
        xhi = np.percentile(bins_cat[mask], hi)
        pad = (xhi - xlo) * 0.05
        return (xlo - pad, xhi + pad)

    xlim_b = xlim_before if xlim_before is not None else _percentile_xlim("before_mean", "bins_before", lo=2, hi=98)
    xlim_a = xlim_after if xlim_after is not None else _percentile_xlim("after_mean", "bins_after", lo=2, hi=98)

    _draw_summary_axes(axes[1], summary, xlim_before=xlim_b, xlim_after=xlim_a)
    axes[1, 0].set_title(f"Before norm  (zoomed x: [{xlim_b[0]:.3f}, {xlim_b[1]:.3f}])" if xlim_b else "Before norm  (zoomed)")
    after_lim_str = f"[{xlim_a[0]:.2f}, {xlim_a[1]:.2f}]" if xlim_a else "auto"
    axes[1, 1].set_title(f"After norm   (zoomed x: {after_lim_str})")

    fig.text(0.5, 0.01, "Mean over all source images per dataset",
             ha="center", fontsize=9, color="#555555")
    plt.tight_layout(rect=[0, 0.04, 1, 1])
    return fig


def _apply_transform(arrays: list[np.ndarray], mode: str) -> list[np.ndarray]:
    """mode: 'clip' or 'clip_sqrt'."""
    if mode == "clip":
        return [np.clip(a, 0, None) for a in arrays]
    elif mode == "clip_sqrt":
        return [np.sqrt(np.clip(a, 0, None)) for a in arrays]
    return arrays


# ── shifted-log mapping → [0, 1] ─────────────────────────────────────────────

def _log_map_forward_np(a: np.ndarray) -> np.ndarray:
    """log_map_forward applied per-array using module constants."""
    x = np.maximum(a.astype(np.float64), LOG_MAP_X_MIN)
    norm = np.log1p((LOG_MAP_X_REF - LOG_MAP_X_MIN) / LOG_MAP_DELTA)
    return (np.log1p((x - LOG_MAP_X_MIN) / LOG_MAP_DELTA) / norm).astype(np.float32)


def _apply_log_map(arrays: list[np.ndarray]) -> list[np.ndarray]:
    return [_log_map_forward_np(a) for a in arrays]


def _log_map_key_y(x_val: float) -> float:
    """Return the mapped y value for a given x."""
    x = max(x_val, LOG_MAP_X_MIN)
    norm = np.log1p((LOG_MAP_X_REF - LOG_MAP_X_MIN) / LOG_MAP_DELTA)
    return float(np.log1p((x - LOG_MAP_X_MIN) / LOG_MAP_DELTA) / norm)


def make_hist_figure_log_map(ds: dict, frames_data: list) -> plt.Figure:
    """Before panel unchanged; after panel shows shifted-log → [0, 1] mapping."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    y0 = _log_map_key_y(0.0)
    y1 = _log_map_key_y(1.0)
    fig.suptitle(
        f"{ds['label']}  [after: log_map(x_min={LOG_MAP_X_MIN}, x_ref={LOG_MAP_X_REF}, δ={LOG_MAP_DELTA})]",
        fontsize=14, fontweight="bold",
    )

    before_arrays = [b for b, a, _ in frames_data]
    after_t = _apply_log_map([a for b, a, _ in frames_data])

    _plot_curves(axes[0], before_arrays, alpha=0.35, lw=0.9, bold_color=ds["color"], n_bins=100)
    axes[0].set_title("Before normalization  (raw PAX, uint16 / 65535)")
    axes[0].set_xlabel("Pixel intensity (0-1 scale)")
    axes[0].set_ylabel("Density")
    axes[0].legend(fontsize=8)

    _plot_curves(axes[1], after_t, alpha=0.35, lw=0.9, bold_color=ds["color"], n_bins=100)
    axes[1].set_title("After CIO-RB norm  →  shifted-log  →  [0, 1]")
    axes[1].set_xlabel("Log-mapped value")
    axes[1].set_ylabel("Density")
    axes[1].set_xlim(-0.02, 1.05)
    axes[1].axvline(y0, color="blue", ls=":", lw=1.5, label=f"x=0 (bg) → {y0:.3f}")
    axes[1].axvline(y1, color="red",  ls="--", lw=1.5, label=f"x=1 (sigmoid) → {y1:.3f}")
    axes[1].legend(fontsize=8)

    n_frames = len(frames_data)
    fig.text(
        0.5, 0.01,
        (f"x_min={LOG_MAP_X_MIN}, x_ref={LOG_MAP_X_REF}, δ={LOG_MAP_DELTA}  "
         f"|  x=0→{y0:.3f}, x=1→{y1:.3f}, x=10→1.000  "
         f"|  {n_frames} images  |  bold=mean"),
        ha="center", fontsize=8, color="#555555",
    )
    plt.tight_layout(rect=[0, 0.05, 1, 1])
    return fig


def _add_log_norm_section(prs, ds_list, ds_frames_list):
    """Add shifted-log → [0,1] section: per-dataset slides + summary slides."""
    y0 = _log_map_key_y(0.0)
    y1 = _log_map_key_y(1.0)

    _section_slide(
        prs,
        f"After norm: shifted-log → [0,1]",
        (f"log1p((x − {LOG_MAP_X_MIN}) / {LOG_MAP_DELTA}) / log1p(({LOG_MAP_X_REF} − {LOG_MAP_X_MIN}) / {LOG_MAP_DELTA})"
         f"  |  x=0→{y0:.3f}, x=1→{y1:.3f}  |  Before panel unchanged"),
    )

    summary_t = []
    for ds, all_frames in zip(ds_list, ds_frames_list):
        fig = make_hist_figure_log_map(ds, all_frames)
        _image_slide(prs, _fig_to_buf(fig))

        before_arrays = [b for b, a, _ in all_frames]
        after_t = _apply_log_map([a for b, a, _ in all_frames])

        fig_tmp, axes_tmp = plt.subplots(1, 2)
        bins_b, mean_b = _plot_curves(axes_tmp[0], before_arrays, alpha=0, lw=0, bold_color=ds["color"], n_bins=100)
        bins_a, mean_a = _plot_curves(axes_tmp[1], after_t,       alpha=0, lw=0, bold_color=ds["color"], n_bins=100)
        plt.close(fig_tmp)
        summary_t.append({
            "label": ds["label"], "color": ds["color"],
            "before_mean": mean_b, "after_mean": mean_a,
            "bins_before": bins_b, "bins_after": bins_a,
        })

    fig = make_summary_figure(summary_t)
    fig.suptitle(
        f"All datasets — mean distributions  [shifted-log→[0,1]]  (full range)",
        fontsize=14, fontweight="bold",
    )
    _image_slide(prs, _fig_to_buf(fig))

    fig = make_summary_zoom_figure(summary_t,
                                   xlim_before=(-0.05, 0.2),
                                   xlim_after=(-0.02, 1.05))
    fig.suptitle(
        f"All datasets — mean distributions  [shifted-log→[0,1]]  (full + zoomed)",
        fontsize=14, fontweight="bold",
    )
    _image_slide(prs, _fig_to_buf(fig))

    print("  Adding pax variance table slide (log-mapped, per-patch)...")
    fig = make_stats_table_figure(ds_list, apply_log=True)
    _image_slide(prs, _fig_to_buf(fig))
    print("  Adding pixel-pool table slide (log-mapped)...")
    fig = make_pixel_pool_table_figure(ds_list, ds_frames_list, mode="log")
    _image_slide(prs, _fig_to_buf(fig))


def make_hist_figure_transformed(ds: dict, frames_data: list, mode: str) -> plt.Figure:
    """Before panel is unchanged; transform applied only to the after (CIO-RB) panel."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    mode_label = {"clip": "clip(neg→0)", "clip_sqrt": "clip + sqrt"}[mode]
    fig.suptitle(f"{ds['label']}  [after panel: {mode_label}]", fontsize=15, fontweight="bold")

    before_arrays = [b for b, a, _ in frames_data]
    after_arrays  = _apply_transform([a for b, a, _ in frames_data], mode)

    _plot_curves(axes[0], before_arrays, alpha=0.35, lw=0.9, bold_color=ds["color"])
    axes[0].set_title("Before normalization  (raw PAX, uint16 / 65535)")
    axes[0].set_xlabel("Pixel intensity (0-1 scale)")
    axes[0].set_ylabel("Density")
    axes[0].legend(fontsize=8)

    _plot_curves(axes[1], after_arrays, alpha=0.35, lw=0.9, bold_color=ds["color"])
    axes[1].set_title(f"After CIO-RB normalization  [{mode_label}]")
    axes[1].set_xlabel("Transformed normalized value")
    axes[1].set_ylabel("Density")
    # sqrt(1.0)=1.0, so the sigmoid ceiling landmark stays at x=1.0 for both modes
    axes[1].axvline(1.0, color="red", ls="--", lw=1.5,
                    label="sqrt(sigmoid ceiling)=1.0" if mode == "clip_sqrt" else "sigmoid ceiling (1.0)")
    if mode == "clip":
        axes[1].set_xlim(-0.05, 2.0)
    axes[1].legend(fontsize=8)

    n_frames = len(frames_data)
    fig.text(
        0.5, 0.01,
        f"{n_frames} source images  |  each curve = one image  |  bold = mean",
        ha="center", fontsize=9, color="#555555",
    )
    plt.tight_layout(rect=[0, 0.04, 1, 1])
    return fig


def make_pixel_pool_table_figure(ds_list: list[dict],
                                  ds_frames_list: list[list],
                                  mode: str = "after") -> plt.Figure:
    """
    Pool ALL pixels for each dataset and compute percentile statistics.

    mode : "before"  — raw uint16/65535 pixels from CZI
           "after"   — CIO-RB normalized pixels (tiff_patches32_label)
           "log"     — CIO-RB pixels after shifted-log map
    """
    mode_labels = {
        "before": "Before norm  (raw PAX, uint16 / 65535)",
        "after":  "After CIO-RB normalization",
        "log":    f"After CIO-RB → shifted-log  (x_ref={LOG_MAP_X_REF}, δ={LOG_MAP_DELTA})",
    }
    pcts = [0, 1, 5, 25, 50, 75, 95, 99, 100]
    pct_labels = ["min", "p1", "p5", "p25", "p50\n(median)", "p75", "p95", "p99", "max"]

    rows, colors_used = [], []
    for ds, frames in zip(ds_list, ds_frames_list):
        if mode == "before":
            pixels = np.concatenate([b for b, a, _ in frames])
        else:
            pixels = np.concatenate([a for b, a, _ in frames])
            if mode == "log":
                pixels = _log_map_forward_np(pixels)

        n_px = len(pixels)
        vals = np.percentile(pixels, pcts)
        mean = float(pixels.mean())
        std  = float(pixels.std())
        row  = [ds["label"], f"{n_px/1e6:.2f} M", f"{mean:.4f}", f"{std:.4f}"]
        row += [f"{v:.4f}" for v in vals]
        rows.append(row)
        colors_used.append(ds["color"])
        print(f"    {ds['key']:8s}  n={n_px:,}  mean={mean:.4f}  std={std:.4f}  "
              f"p50={vals[4]:.4f}  p95={vals[6]:.4f}  max={vals[8]:.4f}")

    col_labels = ["Dataset", "n pixels", "mean", "std"] + pct_labels

    fig, ax = plt.subplots(figsize=(16, 3.5))
    ax.axis("off")
    tbl = ax.table(cellText=rows, colLabels=col_labels,
                   loc="center", cellLoc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9.5)
    tbl.scale(1, 2.5)

    import matplotlib.colors as mcolors
    for (r, c), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#2E4057")
            cell.set_text_props(color="white", fontweight="bold")
        elif r > 0:
            rgba = list(mcolors.to_rgba(colors_used[r - 1]))
            rgba[3] = 0.18
            cell.set_facecolor(rgba)
        # highlight p95 and max columns
        if c in (len(col_labels) - 3, len(col_labels) - 2, len(col_labels) - 1):
            cell.set_edgecolor("#CC4444")
            cell.set_linewidth(1.5)

    fig.suptitle(
        f"Pax pixel-level statistics — pooled per dataset\n{mode_labels[mode]}",
        fontsize=12, y=0.98,
    )
    fig.text(0.5, 0.01,
             "All pixels from all patch locations (tiff_patches32_label) pooled per dataset  |  "
             "red borders = p95 / p99 / max",
             ha="center", fontsize=8, color="#555555")
    plt.tight_layout(rect=[0, 0.06, 1, 0.92])
    return fig


def _compute_mr10_stats(ds_key: str, apply_log: bool = False) -> dict:
    """Per-patch statistics from tiff_patches32_mr10 (both conditions, max 3000/cond)."""
    import glob, random
    patch_means, patch_stds, patch_maxs, patch_p95s = [], [], [], []
    random.seed(42)
    for cond in ["control", "ycomp"]:
        pat = str(PATCH_BASE / ds_key / cond / "tiff_patches32_mr10")
        files = glob.glob(f"{pat}/*.tif")
        if len(files) > 3000:
            files = random.sample(files, 3000)
        for fp in files:
            im = tifffile.imread(fp).astype(np.float32)
            if apply_log:
                im = _log_map_forward_np(im)
            patch_means.append(float(im.mean()))
            patch_stds.append(float(im.std()))
            patch_maxs.append(float(im.max()))
            patch_p95s.append(float(np.percentile(im, 95)))
    n = len(patch_means)
    return dict(
        n=n,
        mean=np.mean(patch_means),
        std=np.mean(patch_stds),
        p95=np.mean(patch_p95s),
        max=np.mean(patch_maxs),
    )


def make_stats_table_figure(ds_list: list[dict], apply_log: bool = False) -> plt.Figure:
    """Violin plots + summary table of pax patch statistics from mr10 patches."""
    print(f"  Computing mr10 patch stats (log={apply_log})...")
    ds_stats = []
    for ds in ds_list:
        s = _compute_mr10_stats(ds["key"], apply_log=apply_log)
        s["label"] = ds["label"]
        s["color"] = ds["color"]
        ds_stats.append(s)
        print(f"    {ds['key']:8s}  n={s['n']:5d}  mean={s['mean']:.3f}  "
              f"std={s['std']:.3f}  p95={s['p95']:.3f}  max={s['max']:.3f}")

    # Reload full arrays for violins (sample)
    import glob, random
    random.seed(42)
    metrics_data = {m: [] for m in ["mean", "std", "p95", "max"]}
    ds_labels = [s["label"] for s in ds_stats]
    # build a ds_key lookup from label
    _label_to_key = {ds["label"]: ds["key"] for ds in ds_list}
    all_vals = {m: [] for m in metrics_data}
    for s in ds_stats:
        import glob
        vals_m, vals_s, vals_p, vals_x = [], [], [], []
        ds_key = _label_to_key[s["label"]]
        for cond in ["control", "ycomp"]:
            pat = str(PATCH_BASE / ds_key / cond / "tiff_patches32_mr10")
            files = glob.glob(f"{pat}/*.tif")
            if len(files) > 1500:
                files = random.sample(files, 1500)
            for fp in files:
                im = tifffile.imread(fp).astype(np.float32)
                if apply_log:
                    im = _log_map_forward_np(im)
                vals_m.append(im.mean())
                vals_s.append(im.std())
                vals_p.append(np.percentile(im, 95))
                vals_x.append(im.max())
        all_vals["mean"].append(np.array(vals_m))
        all_vals["std"].append(np.array(vals_s))
        all_vals["p95"].append(np.array(vals_p))
        all_vals["max"].append(np.array(vals_x))

    log_tag = " [log-mapped]" if apply_log else " [CIO-RB]"
    fig = plt.figure(figsize=(14, 8))
    gs = fig.add_gridspec(2, 4, hspace=0.45, wspace=0.35)

    metric_titles = [("mean", "Patch mean"), ("std", "Patch std\n(local contrast)"),
                     ("p95",  "Patch p95"),  ("max",  "Patch max")]
    colors = [s["color"] for s in ds_stats]

    for col, (metric, title) in enumerate(metric_titles):
        ax = fig.add_subplot(gs[0, col])
        data = all_vals[metric]
        vp = ax.violinplot(data, positions=range(len(ds_stats)),
                           showmedians=True, showextrema=False)
        for body, c in zip(vp["bodies"], colors):
            body.set_facecolor(c); body.set_alpha(0.7)
        vp["cmedians"].set_color("black")
        ax.set_xticks(range(len(ds_stats)))
        ax.set_xticklabels([s["label"].split(" - ")[0] for s in ds_stats], fontsize=8)
        ax.set_title(title, fontsize=9)

    # Summary table (bottom row spanning all columns)
    ax_tbl = fig.add_subplot(gs[1, :])
    ax_tbl.axis("off")
    col_labels = ["Dataset", "n patches", "mean", "std", "p95", "max"]
    rows = []
    for s in ds_stats:
        rows.append([s["label"], f"{s['n']:,}", f"{s['mean']:.3f}",
                     f"{s['std']:.3f}", f"{s['p95']:.3f}", f"{s['max']:.3f}"])
    tbl = ax_tbl.table(cellText=rows, colLabels=col_labels,
                       loc="center", cellLoc="center")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(11)
    tbl.scale(1, 2.2)
    for (r, c), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#2E4057")
            cell.set_text_props(color="white", fontweight="bold")
        elif r > 0:
            import matplotlib.colors as mcolors
            rgba = list(mcolors.to_rgba(colors[r - 1]))
            rgba[3] = 0.18
            cell.set_facecolor(rgba)

    fig.suptitle(f"Pax patch statistics across 4 datasets — mr10{log_tag}", fontsize=13)
    return fig


def _fig_to_buf(fig: plt.Figure) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


# ── PPT helpers ──────────────────────────────────────────────────────────────
W_SLIDE = Inches(13.33)
H_SLIDE = Inches(7.5)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W_SLIDE
    prs.slide_height = H_SLIDE
    return prs


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.33), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Intensity Normalization: Before / After CIO-RB"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.33), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = ("PAX channel pixels at label-patch locations  |  "
               "each curve = one image frame  |  "
               "bold = dataset mean")
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC)
    p2.alignment = PP_ALIGN.CENTER


def _image_slide(prs: Presentation, buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(buf, Inches(0.1), Inches(0.1),
                             width=Inches(13.13), height=Inches(7.3))


def _section_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x1A)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.33), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xAA)
        p2.alignment = PP_ALIGN.CENTER


def _add_transform_section(prs, ds_list, ds_frames_list, mode):
    """Add per-dataset slides + summary slides; transform applied to after panel only."""
    mode_label = {"clip": "clip(neg→0)", "clip_sqrt": "clip + sqrt"}[mode]
    _section_slide(prs, f"After norm: {mode_label}",
                   "Before panel unchanged  |  transform applied to CIO-RB normalized panel only")

    summary_t = []
    for ds, all_frames in zip(ds_list, ds_frames_list):
        fig = make_hist_figure_transformed(ds, all_frames, mode)
        _image_slide(prs, _fig_to_buf(fig))

        before_arrays = [b for b, a, _ in all_frames]
        after_t       = _apply_transform([a for b, a, _ in all_frames], mode)

        fig_tmp, axes_tmp = plt.subplots(1, 2)
        bins_b, mean_b = _plot_curves(axes_tmp[0], before_arrays, alpha=0, lw=0, bold_color=ds["color"])
        bins_a, mean_a = _plot_curves(axes_tmp[1], after_t,       alpha=0, lw=0, bold_color=ds["color"])
        plt.close(fig_tmp)
        summary_t.append({
            "label": ds["label"], "color": ds["color"],
            "before_mean": mean_b, "after_mean": mean_a,
            "bins_before": bins_b, "bins_after": bins_a,
        })

    fig = make_summary_figure(summary_t)
    fig.suptitle(f"All datasets — mean distributions  [after: {mode_label}]  (full range)",
                 fontsize=14, fontweight="bold")
    _image_slide(prs, _fig_to_buf(fig))

    # for zoom: auto-compute after xlim from data (not hardcoded -0.5 to 2.0)
    fig = make_summary_zoom_figure(summary_t, xlim_after=None)
    fig.suptitle(f"All datasets — mean distributions  [after: {mode_label}]  (full + zoomed)",
                 fontsize=14, fontweight="bold")
    _image_slide(prs, _fig_to_buf(fig))


# ── main ────────────────────────────────────────────────────────────────────
def main() -> None:
    prs = _new_prs()
    _title_slide(prs)

    summary = []
    ds_list, ds_frames_list = [], []

    for ds in DATASETS:
        print(f"\n=== {ds['label']} ===")
        all_frames: list = []
        for cond in ds["conditions"]:
            print(f"  [{cond['cond']}] loading...")
            frames = gather_frame_pixels(cond)
            print(f"    {len(frames)} frames")
            all_frames.extend(frames)

        if not all_frames:
            print(f"  No data — skipping slide")
            continue

        before_arrays = [b for b, a, _ in all_frames]
        after_arrays  = [a for b, a, _ in all_frames]

        print(f"  Plotting {len(all_frames)} curves...")
        fig = make_hist_figure(ds, all_frames)
        _image_slide(prs, _fig_to_buf(fig))
        print(f"  Slide added")

        fig_tmp, axes_tmp = plt.subplots(1, 2)
        bins_b, mean_b = _plot_curves(axes_tmp[0], before_arrays,
                                      alpha=0, lw=0, bold_color=ds["color"])
        bins_a, mean_a = _plot_curves(axes_tmp[1], after_arrays,
                                      alpha=0, lw=0, bold_color=ds["color"])
        plt.close(fig_tmp)

        summary.append({
            "label": ds["label"], "color": ds["color"],
            "before_mean": mean_b, "after_mean": mean_a,
            "bins_before": bins_b, "bins_after": bins_a,
        })
        ds_list.append(ds)
        ds_frames_list.append(all_frames)

    if summary:
        print("\n  Adding summary slides (original)...")
        fig = make_summary_figure(summary)
        _image_slide(prs, _fig_to_buf(fig))
        fig = make_summary_zoom_figure(summary,
                                       xlim_before=(-0.05, 0.2),
                                       xlim_after=(-0.05, 1.5))
        _image_slide(prs, _fig_to_buf(fig))
        print("\n  Adding pax variance table slide (CIO-RB, per-patch)...")
        fig = make_stats_table_figure(ds_list, apply_log=False)
        _image_slide(prs, _fig_to_buf(fig))
        print("\n  Adding pixel-pool table slides (before / after / log)...")
        for mode in ("before", "after", "log"):
            print(f"    mode={mode}")
            fig = make_pixel_pool_table_figure(ds_list, ds_frames_list, mode=mode)
            _image_slide(prs, _fig_to_buf(fig))

    for mode in ("clip", "clip_sqrt"):
        print(f"\n  Adding transform section: {mode}...")
        _add_transform_section(prs, ds_list, ds_frames_list, mode)

    print("\n  Adding transform section: log_norm...")
    _add_log_norm_section(prs, ds_list, ds_frames_list)

    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
