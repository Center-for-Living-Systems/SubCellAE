#!/usr/bin/env python3
"""
make_actin_vs_pax_ppt_zproj.py

Same as make_actin_vs_pax_ppt.py but UMAP and KMeans use z_proj (projector
head output) instead of z_recon.  Reconstruction images and nL1 violin plots
are identical.

Requires regen_zproj_analysis.py to have been run first to generate:
  eval/umap_proj_annotation.png
  eval/umap_proj_condition.png
  eval/umap_proj_split.png
  eval/cluster_panels_proj/all_clusters.tif
  eval/cluster_panels_proj/umap_proj_kmeans_k10.png

Output: slides_actin_vs_pax_zproj.pptx
"""
import io
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# ── import shared helpers from the base PPT script ────────────────────────────
sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_actin_vs_pax_ppt as _base

# re-export everything we don't override
from make_actin_vs_pax_ppt import (
    RUNS, ACTIN_MODELS, PAX_MODELS, CH2_MODELS,
    W_IN, H_IN, MARGIN, IMG_W, TEXT_X, TEXT_W,
    C_WHITE, C_TITLE, C_ACT, C_PAX, C_2CH, C_GRAY, C_LGRAY,
    _px, _add_textbox, _add_bullets, _add_image_bytes, _fig_to_bytes,
    _add_slide, _slide_header, _note_panel,
    _load_visual_frames,
    _make_recon_row, _make_violin_row, _make_violin_2ch_perchannel,
    _content_slide,
    _build_title_slide, _build_model_table,
    _build_actin_recon, _build_actin_violin,
    _build_pax_recon,   _build_pax_violin,
    _build_2ch_recon,   _build_2ch_violin,
    _build_2ch_violin_perchannel_nl1, _build_2ch_violin_perchannel_l1,
)


# ── z_proj path helpers ───────────────────────────────────────────────────────

def _load_cluster_grid_proj(model_dir: Path) -> np.ndarray | None:
    clust_path = model_dir / "eval" / "cluster_panels_proj" / "all_clusters.tif"
    if not clust_path.exists():
        return None
    stack = tifffile.imread(str(clust_path))
    n, ph, pw = stack.shape
    rows, cols = 2, 5
    gap = 4
    grid = np.zeros((rows * ph + (rows-1)*gap, cols * pw + (cols-1)*gap),
                    dtype=np.float32)
    for i in range(n):
        r, c = i // cols, i % cols
        y0, x0 = r * (ph + gap), c * (pw + gap)
        grid[y0:y0+ph, x0:x0+pw] = stack[i]
    return grid


def _make_umap_grid_proj(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(2, n, figsize=(n * 3.2, 5.6), facecolor="white")
    if n == 1: axes = axes[:, np.newaxis]
    row_labels = ["FA annotation (z_proj)", "KMeans k=10 (z_proj)"]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
        for row, (png_path, rlbl) in enumerate(zip([
            d / "eval" / "umap_proj_annotation.png",
            d / "eval" / "cluster_panels_proj" / f"umap_proj_kmeans_k10.png",
        ], row_labels)):
            ax = axes[row, col]
            if png_path.exists():
                ax.imshow(np.array(Image.open(str(png_path))))
            ax.axis("off")
            if col == 0:
                ax.set_ylabel(rlbl, fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_umap_row_proj(model_dirs_labels, key: str) -> bytes:
    """key: 'umap_proj_annotation' or 'umap_proj_kmeans_k10'."""
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 3.2, 3.0), facecolor="white")
    if n == 1: axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        if key == "umap_proj_kmeans_k10":
            png_path = d / "eval" / "cluster_panels_proj" / "umap_proj_kmeans_k10.png"
        else:
            png_path = d / "eval" / f"{key}.png"
        if png_path.exists():
            ax.imshow(np.array(Image.open(str(png_path))))
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_cluster_grid_fig_proj(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 5.0, 3.0), facecolor="white")
    if n == 1: axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        grid = _load_cluster_grid_proj(d)
        if grid is not None:
            ax.imshow(grid, cmap="gray", vmin=0, vmax=1)
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.suptitle("16 patches closest to each cluster centroid  (k=10, 2×5 grid)  [z_proj]",
                 fontsize=8, y=0.02)
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


# ── z_proj slide builders ─────────────────────────────────────────────────────

def _build_actin_umap_proj(prs, act):
    img = _make_umap_grid_proj(act)
    _content_slide(prs,
        "Actin (ch3)  —  UMAP z_proj  (top: FA annotation · bottom: KMeans k=10)",
        C_ACT, img,
        "z_proj UMAP",
        [
            "UMAP and KMeans fitted on projector head output (8-dim), not encoder z",
            "z_proj is optimised directly by NT-Xent — clusters should reflect "
            "contrastive structure more sharply than z_recon",
            "Baseline (no projector): z_recon used as fallback",
        ])


def _build_actin_clusters_proj(prs, act):
    img = _make_cluster_grid_fig_proj(act)
    _content_slide(prs,
        "Actin (ch3)  —  KMeans k=10 cluster panels  [z_proj]",
        C_ACT, img,
        "z_proj clusters",
        [
            "KMeans centroids in z_proj space; patches shown are raw image crops",
            "Compare panel homogeneity with z_recon clusters",
            "If z_proj clusters are visually cleaner: contrastive objective "
            "organised the projector space around image morphology",
        ])


def _build_pax_umap_proj(prs, pax):
    img = _make_umap_grid_proj(pax)
    _content_slide(prs,
        "Paxillin (ch1)  —  UMAP z_proj  (top: FA annotation · bottom: KMeans k=10)",
        C_PAX, img,
        "z_proj UMAP",
        [
            "UMAP and KMeans on 8-dim projector output",
            "Paxillin: expect sharper FA-type boundaries in z_proj vs z_recon",
            "λ=¼: weaker contrastive pull — check if z_proj clusters are less tight",
        ])


def _build_pax_clusters_proj(prs, pax):
    img = _make_cluster_grid_fig_proj(pax)
    _content_slide(prs,
        "Paxillin (ch1)  —  KMeans k=10 cluster panels  [z_proj]",
        C_PAX, img,
        "z_proj clusters",
        [
            "KMeans in z_proj space; visual content = raw paxillin patches",
            "Tighter spot-size grouping than z_recon = contrastive loss captured FA size",
        ])


def _build_2ch_umap_proj(prs, ch2):
    img = _make_umap_grid_proj(ch2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  UMAP z_proj  (top: FA annotation · bottom: KMeans k=10)",
        C_2CH, img,
        "z_proj UMAP",
        [
            "Joint 2-ch projector output (8-dim) used for UMAP and KMeans",
            "Does z_proj separate FA types better than z_recon for 2ch?",
            "Baseline (no projector): z_recon fallback",
        ])


def _build_2ch_clusters_proj(prs, ch2):
    img = _make_cluster_grid_fig_proj(ch2)
    _content_slide(prs,
        "Paxillin+Actin (2ch)  —  KMeans k=10 cluster panels  [z_proj]",
        C_2CH, img,
        "z_proj clusters",
        [
            "Each panel: paxillin row (top) + actin row (bottom)",
            "Clusters from z_proj — contrastive objective directly shaped this space",
        ])


def _build_comparison_proj(prs, act, pax, ch2):
    slide = _add_slide(prs)
    _slide_header(slide,
        "Actin vs Paxillin vs Pax+Actin  —  ConAE nL1: UMAP z_proj annotation  ·  Cross-dataset nL1",
        color=C_TITLE)

    compare = [act[1], pax[0], ch2[1]]
    img_umap = _make_umap_row_proj(compare, "umap_proj_annotation")
    img_viol = _make_violin_row(compare, metric="recon_nl1")

    top_h   = 3.05
    bot_top = 0.68 + top_h + 0.05
    _add_image_bytes(slide, img_umap, MARGIN, 0.68, width=IMG_W)
    _add_image_bytes(slide, img_viol, MARGIN, bot_top, width=IMG_W)

    _note_panel(slide, "Key comparisons", [
        "All three: ConAE nL1 (λ=½), z_proj UMAP  vs  nL1 generalisation",
        "Paxillin: tighter FA-type separation in z_proj expected",
        "2ch: does joint encoding in z_proj sharpen separation over single channel?",
        "nL1 violin unchanged — reconstruction quality is z_recon-driven",
    ], C_TITLE, font_size=12)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    out_path = Path("slides_actin_vs_pax_zproj.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    act = [(RUNS / k, lbl) for k, lbl in ACTIN_MODELS]
    pax = [(RUNS / k, lbl) for k, lbl in PAX_MODELS]
    ch2 = [(RUNS / k, lbl) for k, lbl in CH2_MODELS]

    _SECTION_HEADERS = {
        "Actin section hdr": ("Results: Actin (ch3)",            C_ACT,
                               "Baseline AE  ·  ConAE nL1  ·  ConAE nL1 λ=¼"),
        "Pax section hdr":   ("Results: Paxillin (ch1)",         C_PAX,
                               "ConAE nL1  ·  ConAE nL1 λ=¼"),
        "2ch section hdr":   ("Results: Paxillin + Actin (2ch)", C_2CH,
                               "AE Baseline  ·  ConAE nL1  ·  ConAE nL1 λ=¼"),
    }

    steps = [
        ("Title",              _build_title_slide,             (prs,)),
        ("Model table",        _build_model_table,             (prs,)),
        ("Actin section hdr",  None,                           None),
        ("Actin recon",        _build_actin_recon,             (prs, act)),
        ("Actin UMAP",         _build_actin_umap_proj,         (prs, act)),
        ("Actin clusters",     _build_actin_clusters_proj,     (prs, act)),
        ("Actin violin",       _build_actin_violin,            (prs, act)),
        ("Pax section hdr",    None,                           None),
        ("Pax recon",          _build_pax_recon,               (prs, pax)),
        ("Pax UMAP",           _build_pax_umap_proj,           (prs, pax)),
        ("Pax clusters",       _build_pax_clusters_proj,       (prs, pax)),
        ("Pax violin",         _build_pax_violin,              (prs, pax)),
        ("2ch section hdr",    None,                           None),
        ("2ch recon",          _build_2ch_recon,               (prs, ch2)),
        ("2ch UMAP",           _build_2ch_umap_proj,           (prs, ch2)),
        ("2ch clusters",       _build_2ch_clusters_proj,       (prs, ch2)),
        ("2ch violin",         _build_2ch_violin,              (prs, ch2)),
        ("2ch violin pch nL1", _build_2ch_violin_perchannel_nl1, (prs, ch2)),
        ("2ch violin pch L1",  _build_2ch_violin_perchannel_l1,  (prs, ch2)),
        ("Comparison",         _build_comparison_proj,         (prs, act, pax, ch2)),
    ]

    for label, fn, args in steps:
        print(f"  {label} …", flush=True)
        if fn is None:
            title_text, color, subtitle = _SECTION_HEADERS[label]
            slide = _add_slide(prs)
            _slide_header(slide, title_text, color=color, font_size=28)
            bar = slide.shapes.add_shape(
                1, _px(0), _px(H_IN * 0.82), _px(W_IN), _px(H_IN * 0.18))
            bar.fill.solid(); bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            _add_textbox(slide, subtitle, MARGIN, H_IN * 0.83,
                         W_IN - 2*MARGIN, 0.5,
                         font_size=14, color=C_WHITE, bold=False)
        else:
            fn(*args)

    prs.save(str(out_path))
    print(f"\nSaved → {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
