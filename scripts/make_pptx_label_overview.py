#!/usr/bin/env python3
"""
make_pptx_label_overview.py
============================
PPT summarising all FA annotation labels across two batches:
  Batch 1 – LabelStudio sessions (Nov 2025 – Feb 2026)
  Batch 2 – Prototype interface sessions (Apr – Aug 2026)

Slides
------
  1. Title
  2. Terminology & dataset definitions
  3. Annotation taxonomy (5 FA classes + colour key)
  4. Batch 1 – Session overview table
  5. Batch 1 – Dataset 1 class distribution (control vs ycomp)
  6. Batch 1 – Dataset 2 & 3 class distribution
  7. Batch 2 – Session overview table
  8. Batch 2 – Dataset 1 class distribution (control vs ycomp)
  9. Batch 2 – Dataset 2 & 3 class distribution
 10. Grand summary – both batches side-by-side per dataset

Output: results/label_overview.pptx
"""
from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
_REPO = Path(__file__).resolve().parents[1]
LAB   = Path("/net/projects/CLS/lding/data/fa_data_analysis/labelling")
RES   = _REPO / "results"
OUT   = RES / "label_overview.pptx"

# ---------------------------------------------------------------------------
# Dataset naming  (official terms — stick to these throughout)
DS_KEYS = ["ds1", "ds2", "ds3"]
DS_SHORT = {           # what to show in slide titles / axis labels
    "ds1": "Dataset 1",
    "ds2": "Dataset 2",
    "ds3": "Dataset 3",
}
DS_LONG = {            # full name with channel protein
    "ds1": "Dataset 1  (Vinculin – VINC)",
    "ds2": "Dataset 2  (Phospho-FAK – pFAK)",
    "ds3": "Dataset 3  (Phospho-Paxillin – pPAX)",
}
DS_CODENAME = {        # internal code name — shown only on the terminology slide
    "ds1": "vinc",
    "ds2": "pfak",
    "ds3": "ppax",
}
DS_COLORS = {
    "ds1": "#1f77b4",
    "ds2": "#ff7f0e",
    "ds3": "#2ca02c",
}

# FA class taxonomy
FA5_ORDER = [
    "No adhesion", "Nascent Adhesion", "focal complex",
    "focal adhesion", "fibrillar adhesion",
]
FA5_COLORS = {
    "No adhesion":        "#9467bd",
    "Nascent Adhesion":   "#1f77b4",
    "focal complex":      "#ff7f0e",
    "focal adhesion":     "#2ca02c",
    "fibrillar adhesion": "#d62728",
}
FA5_SHORT = {
    "No adhesion":        "No Adh",
    "Nascent Adhesion":   "NA",
    "focal complex":      "FC",
    "focal adhesion":     "FA",
    "fibrillar adhesion": "Fib",
}
UNCERTAIN_COLOR = "#aaaaaa"
CONDITIONS = ["control", "ycomp"]
COND_LABELS = {"control": "Control", "ycomp": "Y-compound"}
COND_HATCH  = {"control": "", "ycomp": ""}   # no hatch — condition shown on x-axis

SLIDE_W, SLIDE_H = 13.33, 7.5

# ---------------------------------------------------------------------------
# Batch metadata
# ---------------------------------------------------------------------------
BATCH1_SESSIONS = [
    dict(session="Project 1",    date="Nov 26–28 2025", annotator="Annabel + Margaret", ds="ds3", condition="control",  n=60,  note="pilot, cross-annotator"),
    dict(session="Project 13a",  date="Dec 17–18 2025", annotator="Margaret",           ds="ds1", condition="control",  n=53,  note=""),
    dict(session="Project 13b",  date="Dec 19 2025",    annotator="Margaret",           ds="ds1", condition="control",  n=176, note=""),
    dict(session="Project 15",   date="Dec 18–22 2025", annotator="Margaret",           ds="ds1", condition="ycomp",   n=194, note=""),
    dict(session="Project 17",   date="Feb 6–10 2026",  annotator="Annabel",            ds="ds1", condition="control",  n=42,  note=""),
    dict(session="Project 18",   date="Feb 10 2026",    annotator="Annabel",            ds="ds1", condition="ycomp",   n=53,  note=""),
    dict(session="Project 19",   date="Feb 5 2026",     annotator="Margaret",           ds="ds1", condition="control",  n=200, note=""),
    dict(session="Project 20",   date="Feb 5–6 2026",   annotator="Margaret",           ds="ds1", condition="ycomp",   n=759, note=""),
]
BATCH1_FILES = {
    "ds1": LAB / "labels_vinc_20260521.csv",
    "ds2": LAB / "labels_pfak_20260521.csv",
    "ds3": LAB / "labels_ppax_20260521.csv",
}
B1_LABEL_COL = "classification"
B1_COND_COL  = "condition"

BATCH2_SESSIONS = [
    dict(session="Dataset 2 – Annabel v1", date="Apr 27 2026",  annotator="Annabel", ds="ds2", condition="control",  n=54,  file="pfak_labels_Annabel_20260427_1035.csv"),
    dict(session="Dataset 3 – Ernest ctrl",date="Jul 27 2026",  annotator="Ernest",  ds="ds3", condition="control",  n=111, file="ppax_control_label_Ernest_20260727_1142.csv"),
    dict(session="Dataset 1 – Annabel",    date="Aug 16 2026",  annotator="Annabel", ds="ds1", condition="ctrl+ycomp",n=1224,file="vinc_combined_label_Annabel_20260816.csv"),
    dict(session="Dataset 2 – Annabel v2", date="Aug 2026",     annotator="Annabel", ds="ds2", condition="control",  n=211, file="pfak_combined_label_Annabel_aug2026.csv"),
    dict(session="Dataset 3 – Ernest all", date="Jul 2026",     annotator="Ernest",  ds="ds3", condition="control",  n=261, file="ppax_combined_label_Ernest_latest.csv"),
]
BATCH2_FILES = {
    "ds1": LAB / "vinc_combined_label_Annabel_20260816.csv",
    "ds2": LAB / "pfak_combined_label_Annabel_aug2026.csv",
    "ds3": LAB / "ppax_combined_label_Ernest_latest.csv",
}
B2_LABEL_COL = "label"

# ---------------------------------------------------------------------------
# PPT helpers
# ---------------------------------------------------------------------------

def _hex2rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _paste_pil(slide, img: Image.Image, left, top, width, height):
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(left), Inches(top), Inches(width), Inches(height))


def _fig_to_pil(fig) -> Image.Image:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return Image.open(buf)


def _txt(slide, text, left, top, width, height,
         size=11, bold=False, color="#333333", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = _hex2rgb(color)


def _rect(slide, left, top, width, height, fill_hex):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex2rgb(fill_hex)
    shape.line.fill.background()


def _header(slide, title_text, subtitle=None):
    _rect(slide, 0, 0, SLIDE_W, 0.12, "#2c3e50")
    _txt(slide, title_text, 0.4, 0.18, SLIDE_W - 0.8, 0.55,
         size=20, bold=True, color="#2c3e50")
    if subtitle:
        _txt(slide, subtitle, 0.4, 0.72, SLIDE_W - 0.8, 0.38,
             size=10, color="#666666")


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def _czi_to_fframe(czi_series: pd.Series) -> pd.Series:
    """Extract 0-indexed frame ID (f0000 style) from czi_filename like '..._control-01.czi'."""
    return czi_series.str.extract(r'-(\d+)\.czi', expand=False).apply(
        lambda x: f'f{int(x)-1:04d}' if pd.notna(x) else None)


def _n_images(df: pd.DataFrame, cond_col: str, ds_key: str = None,
              id_col: str = None) -> dict:
    """Return {COND_LABEL: n_images, 'Total': n_total} by counting unique image IDs."""
    df = df.copy()
    if id_col and id_col in df.columns:
        img_col = id_col
    elif "czi_filename" in df.columns and df["czi_filename"].notna().any():
        # Use condition + frame-number as image ID (czi_filename alone is not condition-specific)
        df["_img"] = df[cond_col] + "_" + _czi_to_fframe(df["czi_filename"])
        img_col = "_img"
    elif "filename" in df.columns:
        # Batch-2 style: extract condition_frame as image id
        df["_img"] = (df["filename"].str.extract(r'^(control|ycomp)', expand=False) + "_" +
                      df["filename"].str.extract(r'^(?:control|ycomp)_(f\d+)', expand=False))
        img_col = "_img"
    elif "unique_ID" in df.columns:
        # Batch-1 ds2 style: condition-fXXXX...
        df["_img"] = (df[cond_col] + "_" +
                      df["unique_ID"].str.extract(r'-(f\d+)x', expand=False))
        img_col = "_img"
    else:
        return {}

    conds = [c for c in CONDITIONS if c in df[cond_col].unique()]
    result = {}
    for c in conds:
        result[COND_LABELS[c]] = df[df[cond_col] == c][img_col].nunique()
    result["Total"] = df[img_col].nunique()
    return result


def _n_images_combined(b1_df: pd.DataFrame, b2_df: pd.DataFrame,
                        ds_key: str) -> dict:
    """Count unique images across B1+B2 combined per condition (deduplicates overlapping frames)."""
    # Build B1 image sets per condition (in f-format)
    b1 = b1_df.copy()
    b1_imgs: dict[str, set] = {}
    if "czi_filename" in b1.columns and b1["czi_filename"].notna().any():
        b1["_f"] = _czi_to_fframe(b1["czi_filename"])
        for cond, grp in b1.groupby(B1_COND_COL):
            b1_imgs[cond] = set((grp[B1_COND_COL] + "_" + grp["_f"]).dropna())
    elif "unique_ID" in b1.columns:
        b1["_f"] = b1["unique_ID"].str.extract(r'-(f\d+)x', expand=False)
        for cond, grp in b1.groupby(B1_COND_COL):
            b1_imgs[cond] = set((grp[B1_COND_COL] + "_" + grp["_f"]).dropna())

    # Build B2 image sets per condition
    b2 = b2_df.copy()
    if "condition" not in b2.columns and "filename" in b2.columns:
        b2["condition"] = b2["filename"].str.extract(r'^(control|ycomp)', expand=False)
    b2["_f"] = b2["filename"].str.extract(r'^(?:control|ycomp)_(f\d+)', expand=False)
    b2_imgs: dict[str, set] = {}
    for cond, grp in b2.groupby("condition"):
        b2_imgs[cond] = set((grp["condition"] + "_" + grp["_f"]).dropna())

    all_conds = set(b1_imgs) | set(b2_imgs)
    result = {}
    for cond in [c for c in CONDITIONS if c in all_conds]:
        combined = b1_imgs.get(cond, set()) | b2_imgs.get(cond, set())
        result[COND_LABELS[cond]] = len(combined)
    result["Total"] = sum(v for v in result.values())
    return result


def _count_table(df: pd.DataFrame, label_col: str, cond_col: str,
                 img_counts: dict = None) -> pd.DataFrame:
    """Return class × condition count DataFrame. First data row = # Images if provided."""
    all_cls = FA5_ORDER + ["Uncertain"]
    conds   = [c for c in CONDITIONS if c in df[cond_col].unique()]
    rows = []

    # # Images row first
    if img_counts:
        img_row = {"Class": "# Images"}
        for c in conds:
            img_row[COND_LABELS[c]] = img_counts.get(COND_LABELS[c], 0)
        rows.append(img_row)

    for cls in all_cls:
        sub = df[df[label_col] == cls]
        row = {"Class": cls}
        for c in conds:
            row[COND_LABELS[c]] = int((sub[cond_col] == c).sum())
        rows.append(row)
    tbl = pd.DataFrame(rows)
    tbl["Total"] = tbl.apply(
        lambda r: img_counts.get("Total", 0) if r["Class"] == "# Images"
        else sum(r[COND_LABELS[c]] for c in conds), axis=1)
    # drop class rows with zero total (keep # Images and Total always)
    tbl = tbl[(tbl["Total"] > 0) | tbl["Class"].isin(["# Images"])].copy()
    # Total row
    cls_rows = tbl[~tbl["Class"].isin(["# Images"])]
    totals = {"Class": "Total"}
    for col in tbl.columns[1:]:
        totals[col] = cls_rows[col].sum()
    tbl = pd.concat([tbl, pd.DataFrame([totals])], ignore_index=True)
    return tbl


def _add_class_table(slide, tbl: pd.DataFrame, ds_key: str,
                     left: float, top: float, width: float, height: float):
    """Add a pptx table to slide. tbl from _count_table."""
    from pptx.util import Inches as I, Pt
    from pptx.dml.color import RGBColor

    n_rows = len(tbl) + 1   # +1 header
    n_cols = len(tbl.columns)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, I(left), I(top), I(width), I(height)
    ).table

    # Column widths: class col wider, rest equal
    col_w_cls  = width * 0.42
    col_w_rest = (width - col_w_cls) / (n_cols - 1)
    tbl_shape.columns[0].width = I(col_w_cls)
    for ci in range(1, n_cols):
        tbl_shape.columns[ci].width = I(col_w_rest)

    def _cell(r, c, text, bold=False, bg=None, fg="#222222", align=PP_ALIGN.CENTER):
        cell = tbl_shape.cell(r, c)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.clear()
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(10)
        run.font.bold = bold
        run.font.color.rgb = _hex2rgb(fg)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex2rgb(bg)
        else:
            cell.fill.background()

    # Header row
    for ci, col in enumerate(tbl.columns):
        _cell(0, ci, col, bold=True, bg="#2c3e50", fg="#ffffff",
              align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    # Data rows
    for ri, (_, row) in enumerate(tbl.iterrows()):
        is_total  = row["Class"] == "Total"
        is_images = row["Class"] == "# Images"
        for ci, col in enumerate(tbl.columns):
            val = row[col]
            if is_images:
                bg = "#dde8f0"
                _cell(ri + 1, ci, str(int(val)) if ci > 0 else val,
                      bold=True, bg=bg, fg="#2c3e50",
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            elif is_total:
                bg = "#e8f0fe"
                fg = DS_COLORS.get(ds_key, "#222") if ci > 0 else DS_COLORS.get(ds_key, "#222")
                _cell(ri + 1, ci, str(int(val)) if ci > 0 else val,
                      bold=True, bg=bg, fg=fg,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            else:
                bg = "#f0f4f8" if ri % 2 == 0 else "#ffffff"
                if ci == 0:
                    fg = FA5_COLORS.get(val, UNCERTAIN_COLOR if val == "Uncertain"
                                        else "#222222")
                    _cell(ri + 1, ci, val, bold=True, bg=bg, fg=fg,
                          align=PP_ALIGN.LEFT)
                else:
                    _cell(ri + 1, ci, "" if val == 0 else str(int(val)),
                          bold=False, bg=bg, fg="#333333",
                          align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, 0.12, "#2c3e50")
    _txt(slide, "Focal Adhesion Annotation Dataset Overview",
         0.4, 1.6, SLIDE_W - 0.8, 1.2, size=32, bold=True, color="#2c3e50")
    lines = [
        "Two annotation batches: LabelStudio sessions (Batch 1) & Prototype interface sessions (Batch 2)",
        "Annotators: Margaret · Annabel · Ernest",
        "Datasets: Dataset 1 · Dataset 2 · Dataset 3",
    ]
    for i, line in enumerate(lines):
        _txt(slide, line, 0.4, 3.1 + i * 0.48, SLIDE_W - 0.8, 0.4,
             size=13, color="#555555")

    for i, ds in enumerate(DS_KEYS):
        x = 0.4 + i * 3.1
        _rect(slide, x, 5.0, 0.22, 0.22, DS_COLORS[ds])
        _txt(slide, DS_SHORT[ds], x + 0.3, 4.96, 2.8, 0.32,
             size=11, color="#444444")


def _slide_terminology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Terminology & Dataset Definitions",
            subtitle="Use Dataset 1 / 2 / 3 consistently. Internal code names are listed for reference only.")

    # Table: Dataset | Code name | Protein marker | Cell line | Conditions | Notes
    cols   = ["Name",       "Code",  "Protein marker",         "Cell line",          "Conditions",      "Patches (labeled)"]
    col_x  = [0.25,          1.6,     2.6,                      5.3,                  7.3,               9.3]
    col_w  = [1.3,           0.95,    2.65,                     1.95,                 1.95,              3.8]
    hdr_y  = 1.28

    _rect(slide, 0.2, hdr_y - 0.04, SLIDE_W - 0.4, 0.35, "#2c3e50")
    for cx, cw, ch in zip(col_x, col_w, cols):
        _txt(slide, ch, cx, hdr_y, cw, 0.3, size=9, bold=True, color="#ffffff")

    rows = [
        ("ds1", "Vinculin (VINC)",        "U2OS",      "control, ycomp", "B1: 1,340  ·  B2: 1,224"),
        ("ds2", "Phospho-FAK (pFAK)",     "U2OS",      "control",        "B1: 54     ·  B2: 211"),
        ("ds3", "Phospho-Paxillin (pPAX)","U2OS",      "control",        "B1: 60     ·  B2: 261"),
        ("ds4", "Paxillin (NIH3T3-Zyxin)","NIH3T3",    "control, ycomp", "no labels yet"),
    ]

    for ri, (ds, protein, cell, cond, patches) in enumerate(rows):
        row_y = hdr_y + 0.35 + ri * 0.78
        bg = "#f0f4f8" if ri % 2 == 0 else "#ffffff"
        _rect(slide, 0.2, row_y - 0.05, SLIDE_W - 0.4, 0.72, bg)

        dc = DS_COLORS.get(ds, "#aaaaaa")
        _rect(slide, 0.2, row_y - 0.05, 0.06, 0.72, dc)

        vals = [DS_SHORT.get(ds, "Dataset 4"), DS_CODENAME.get(ds, "nih3t3"),
                protein, cell, cond, patches]
        bolds = [True, False, False, False, False, False]
        colors = [DS_COLORS.get(ds, "#888888")] + ["#333333"] * 5
        for cx, cw, v, bd, co in zip(col_x, col_w, vals, bolds, colors):
            _txt(slide, v, cx + 0.06, row_y, cw, 0.65,
                 size=9.5, bold=bd, color=co)

    # Bottom note
    _txt(slide,
         "▸ 'ycomp' = Y-27632 ROCK inhibitor compound treatment  ·  "
         "B1 = LabelStudio batch  ·  B2 = Prototype interface batch",
         0.25, 6.85, SLIDE_W - 0.5, 0.4, size=8.5, color="#888888", italic=True)


def _slide_taxonomy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Annotation Taxonomy: 5 FA Sub-classes")

    descriptions = {
        "No adhesion":        "Background / no detectable focal adhesion structure",
        "Nascent Adhesion":   "Small dot-like early adhesions near the cell edge",
        "focal complex":      "Intermediate adhesions at the lamellipodia front",
        "focal adhesion":     "Mature, elongated adhesions in the cell interior",
        "fibrillar adhesion": "Fibrillar, actin-aligned structures, central region",
    }
    for i, cls in enumerate(FA5_ORDER):
        y = 1.1 + i * 1.1
        _rect(slide, 0.4, y, 0.38, 0.38, FA5_COLORS[cls])
        _txt(slide, f"{FA5_SHORT[cls]}  —  {cls}",
             0.9, y - 0.02, 4.0, 0.42, size=13, bold=True, color="#222222")
        _txt(slide, descriptions[cls],
             0.9, y + 0.34, 7.5, 0.38, size=10, color="#555555")

    _txt(slide,
         "Note: 'Uncertain' labels appear in some sessions and are carried in the CSVs "
         "but excluded from model training.",
         0.4, 6.85, SLIDE_W - 0.8, 0.38, size=9, color="#999999", italic=True)


def _slide_batch1_sessions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Batch 1: LabelStudio Sessions  (Nov 2025 – Feb 2026)",
            subtitle="8 LabelStudio projects. Consolidated into 3 dataset CSVs on 2026-05-21. "
                     "Later annotations overwrite earlier ones for duplicate patches.")

    cols  = ["Session",     "Date",           "Annotator",        "Dataset",    "Condition",  "# Patches", "Note"]
    col_x = [0.25,           1.75,             3.45,               5.2,          6.25,          7.5,         8.65]
    col_w = [1.45,           1.65,             1.7,                1.0,          0.95,          1.1,         4.6]
    hdr_y = 1.25

    _rect(slide, 0.2, hdr_y - 0.04, SLIDE_W - 0.4, 0.35, "#2c3e50")
    for cx, cw, ch in zip(col_x, col_w, cols):
        _txt(slide, ch, cx, hdr_y, cw, 0.3, size=9, bold=True, color="#ffffff")

    for ri, s in enumerate(BATCH1_SESSIONS):
        row_y = hdr_y + 0.35 + ri * 0.52
        bg = "#f5f5f5" if ri % 2 == 0 else "#ffffff"
        _rect(slide, 0.2, row_y - 0.04, SLIDE_W - 0.4, 0.48, bg)
        dc = DS_COLORS.get(s["ds"], "#888")
        _rect(slide, 0.2, row_y - 0.04, 0.06, 0.48, dc)
        vals = [s["session"], s["date"], s["annotator"],
                DS_SHORT[s["ds"]], s["condition"], str(s["n"]), s.get("note", "")]
        for cx, cw, v in zip(col_x, col_w, vals):
            _txt(slide, v, cx + 0.06, row_y, cw, 0.42, size=8.5, color="#222222")

    _txt(slide,
         "▶  Consolidated files:  labels_vinc_20260521.csv (1,340)  ·  "
         "labels_pfak_20260521.csv (54)  ·  labels_ppax_20260521.csv (60)",
         0.25, 6.88, SLIDE_W - 0.5, 0.38, size=9, bold=True, color="#2c3e50")


def _slide_batch1_ds1(prs, b1_dfs):
    df = b1_dfs["ds1"]
    n_ctrl  = (df[B1_COND_COL] == "control").sum()
    n_ycomp = (df[B1_COND_COL] == "ycomp").sum()
    imgs = _n_images(df, B1_COND_COL)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide,
            f"Batch 1 — {DS_SHORT['ds1']}: Class Counts by Condition",
            subtitle=f"Total: {len(df):,} patches  ·  Control: {n_ctrl:,}  ·  Y-compound: {n_ycomp:,}  "
                     f"·  Annotator: Margaret")
    tbl = _count_table(df, B1_LABEL_COL, B1_COND_COL, img_counts=imgs)
    _add_class_table(slide, tbl, "ds1", left=2.5, top=1.3, width=8.3, height=5.9)


def _slide_batch1_ds23(prs, b1_dfs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Batch 1 — Dataset 2 & Dataset 3: Class Counts by Condition",
            subtitle="Dataset 2 annotated by Annabel · Dataset 3 (pilot) annotated by Annabel + Margaret · control only")
    for i, ds_key in enumerate(["ds2", "ds3"]):
        df = b1_dfs[ds_key]
        imgs = _n_images(df, B1_COND_COL, ds_key=ds_key)
        tbl = _count_table(df, B1_LABEL_COL, B1_COND_COL, img_counts=imgs)
        left = 0.3 + i * 6.5
        _txt(slide, DS_SHORT[ds_key], left, 1.28, 6.0, 0.38,
             size=12, bold=True, color=DS_COLORS[ds_key])
        _add_class_table(slide, tbl, ds_key, left=left, top=1.65, width=6.2, height=5.3)


def _slide_batch2_sessions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Batch 2: Prototype Interface Sessions  (Apr – Aug 2026)",
            subtitle="Custom annotation interface. Annabel labeled Datasets 1 & 2. Ernest labeled Dataset 3.")

    cols  = ["Session",                  "Date",          "Annotator",   "Dataset",   "Condition",    "# Patches", "Output file"]
    col_x = [0.25,                        2.0,             3.5,           5.1,          6.1,            7.35,        8.55]
    col_w = [1.7,                         1.45,            1.55,          0.95,         1.2,            1.15,        4.7]
    hdr_y = 1.25

    _rect(slide, 0.2, hdr_y - 0.04, SLIDE_W - 0.4, 0.35, "#2c3e50")
    for cx, cw, ch in zip(col_x, col_w, cols):
        _txt(slide, ch, cx, hdr_y, cw, 0.3, size=9, bold=True, color="#ffffff")

    for ri, s in enumerate(BATCH2_SESSIONS):
        row_y = hdr_y + 0.35 + ri * 0.88
        bg = "#f5f5f5" if ri % 2 == 0 else "#ffffff"
        _rect(slide, 0.2, row_y - 0.04, SLIDE_W - 0.4, 0.82, bg)
        dc = DS_COLORS.get(s["ds"], "#888")
        _rect(slide, 0.2, row_y - 0.04, 0.06, 0.82, dc)
        vals = [s["session"], s.get("date", "").strip(), s["annotator"],
                DS_SHORT[s["ds"]], s["condition"], str(s["n"]), s.get("file", "")]
        for cx, cw, v in zip(col_x, col_w, vals):
            _txt(slide, v, cx + 0.06, row_y, cw, 0.76, size=8.5, color="#222222")


def _slide_batch2_ds1(prs, b2_dfs):
    df = b2_dfs["ds1"].copy()
    df["condition"] = df["filename"].str.extract(r'^(control|ycomp)', expand=False)
    n_ctrl  = (df["condition"] == "control").sum()
    n_ycomp = (df["condition"] == "ycomp").sum()
    imgs = _n_images(df, "condition")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide,
            f"Batch 2 — {DS_SHORT['ds1']}: Class Counts by Condition",
            subtitle=f"Total: {len(df):,} patches  ·  Control: {n_ctrl:,}  ·  Y-compound: {n_ycomp:,}  "
                     f"·  Annotator: Annabel")
    tbl = _count_table(df, B2_LABEL_COL, "condition", img_counts=imgs)
    _add_class_table(slide, tbl, "ds1", left=2.5, top=1.3, width=8.3, height=5.7)


def _slide_batch2_ds23(prs, b2_dfs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Batch 2 — Dataset 2 & Dataset 3: Class Counts by Condition",
            subtitle="Both datasets: control condition only.")
    for i, ds_key in enumerate(["ds2", "ds3"]):
        df = b2_dfs[ds_key].copy()
        df["condition"] = df["filename"].str.extract(r'^(control|ycomp)', expand=False)
        imgs = _n_images(df, "condition")
        tbl = _count_table(df, B2_LABEL_COL, "condition", img_counts=imgs)
        left = 0.3 + i * 6.5
        _txt(slide, DS_SHORT[ds_key], left, 1.28, 6.0, 0.38,
             size=12, bold=True, color=DS_COLORS[ds_key])
        _add_class_table(slide, tbl, ds_key, left=left, top=1.65, width=6.2, height=5.3)


def _fig_summary_by_condition(b1_dfs: dict, b2_dfs: dict) -> Image.Image:
    """For each dataset × condition, show B1 vs B2 stacked bars."""
    # Build a 2-row layout: top = control, bottom = ycomp
    # Columns = Dataset 1, Dataset 2, Dataset 3
    fig, axes = plt.subplots(2, 3, figsize=(11, 7.0), sharey=False)

    cond_rows = {"control": 0, "ycomp": 1}

    for col_i, ds_key in enumerate(DS_KEYS):
        b1 = b1_dfs[ds_key].copy()
        b2 = b2_dfs[ds_key].copy()

        # Add condition column
        if "condition" not in b1.columns:
            b1["condition"] = b1.get("filename", b1.get("crop_img_filename", pd.Series())).str.extract(r'^(control|ycomp)', expand=False)
        if "condition" not in b2.columns:
            b2["condition"] = b2["filename"].str.extract(r'^(control|ycomp)', expand=False)

        for cond, row_i in cond_rows.items():
            ax = axes[row_i][col_i]
            sub1 = b1[b1["condition"] == cond]
            sub2 = b2[b2["condition"] == cond]

            x = [0, 0.5]
            bottoms = [0, 0]

            for cls in FA5_ORDER:
                vals = [(sub1[B1_LABEL_COL] == cls).sum(),
                        (sub2[B2_LABEL_COL] == cls).sum()]
                for xi, v in enumerate(vals):
                    ax.bar(x[xi], v, bottom=bottoms[xi], width=0.38,
                           color=FA5_COLORS[cls], linewidth=0)
                    if v >= 8:
                        ax.text(x[xi], bottoms[xi] + v / 2, str(v),
                                ha="center", va="center", fontsize=7,
                                color="white", fontweight="bold")
                    bottoms[xi] += v

            # Uncertain
            for xi, (sub, lc) in enumerate([(sub1, B1_LABEL_COL), (sub2, B2_LABEL_COL)]):
                unc = (sub[lc] == "Uncertain").sum()
                if unc:
                    ax.bar(x[xi], unc, bottom=bottoms[xi], width=0.38,
                           color=UNCERTAIN_COLOR, linewidth=0)
                    bottoms[xi] += unc

            # Totals
            for xi, bot in enumerate(bottoms):
                ax.text(x[xi], bot + max(bottoms) * 0.02 if max(bottoms) else 0.5,
                        str(bot), ha="center", va="bottom",
                        fontsize=8, fontweight="bold", color=DS_COLORS[ds_key])

            ax.set_xticks(x)
            ax.set_xticklabels(["Batch 1", "Batch 2"], fontsize=8)
            ax.spines[["top", "right"]].set_visible(False)

            if col_i == 0:
                ax.set_ylabel("# patches", fontsize=9)
            if row_i == 0:
                ax.set_title(DS_SHORT[ds_key], fontsize=10, fontweight="bold",
                             color=DS_COLORS[ds_key], pad=6)

            # Condition label on right side of rightmost column
            if col_i == 2:
                ax.text(1.05, 0.5, COND_LABELS[cond],
                        transform=ax.transAxes,
                        fontsize=10, fontweight="bold", va="center",
                        rotation=-90, color="#444444")

            # If no data at all, show placeholder
            if max(bottoms) == 0:
                ax.text(0.25, 0.5, "no data", transform=ax.transAxes,
                        ha="center", va="center", fontsize=9, color="#aaa")
                ax.set_ylim(0, 10)

    # Row labels (left margin)
    for cond, row_i in cond_rows.items():
        fig.text(0.01, 0.74 - row_i * 0.48, COND_LABELS[cond],
                 fontsize=11, fontweight="bold", va="center",
                 rotation=90, color="#333333")

    fig.suptitle("Summary by Condition: Control vs Y-compound", fontsize=12,
                 fontweight="bold", y=1.01)
    fig.tight_layout(rect=[0.04, 0.06, 1, 1])
    return _fig_to_pil(fig)


def _slide_summary_by_condition(prs, b1_dfs, b2_dfs):
    """Table: rows=classes, cols=DS1-ctrl, DS1-ycomp, DS2-ctrl, DS3-ctrl for each batch."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    b1_total = sum(len(df) for df in b1_dfs.values())
    b2_total = sum(len(df) for df in b2_dfs.values())
    _header(slide, "Grand Summary: Control vs Y-compound Totals",
            subtitle=f"Grand total: {b1_total + b2_total:,} patches  ·  "
                     f"Batch 1: {b1_total:,}  ·  Batch 2: {b2_total:,}")

    # Build one combined table per batch: rows=classes, cols=DS×condition
    for bi, (batch_label, dfs, label_col) in enumerate([
        ("Batch 1 (LabelStudio)", b1_dfs, B1_LABEL_COL),
        ("Batch 2 (Prototype)",   b2_dfs, B2_LABEL_COL),
    ]):
        top = 1.3 + bi * 3.0
        _txt(slide, batch_label, 0.25, top - 0.28, 12.0, 0.3,
             size=11, bold=True, color="#2c3e50")

        # Columns: Class | DS1-Control | DS1-Ycomp | DS2-Control | DS3-Control | Total
        col_specs = []
        for ds_key in DS_KEYS:
            df = dfs[ds_key].copy()
            if "condition" not in df.columns:
                df["condition"] = df["filename"].str.extract(r'^(control|ycomp)', expand=False)
            lc = label_col
            for cond in [c for c in CONDITIONS if c in df["condition"].unique()]:
                col_specs.append((ds_key, cond, df[df["condition"] == cond], lc))

        # Image counts per condition: per-batch (sub is already filtered by condition)
        img_counts_per_col = [
            _n_images(sub, "condition").get("Total", 0)
            for ds_key, cond, sub, lc in col_specs
        ]
        img_row = ["# Images"] + img_counts_per_col + [sum(img_counts_per_col)]

        all_cls = FA5_ORDER + ["Uncertain"]
        rows_data = []
        for cls in all_cls:
            row = [cls]
            total = 0
            for ds_key, cond, sub, lc in col_specs:
                cnt = int((sub[lc] == cls).sum())
                row.append(cnt)
                total += cnt
            row.append(total)
            rows_data.append(row)

        # Total row
        total_row = ["Total"]
        for j in range(len(col_specs)):
            total_row.append(sum(r[j + 1] for r in rows_data))
        total_row.append(sum(total_row[1:]))
        rows_data.append(total_row)

        # Drop all-zero rows (except Total), then prepend # Images
        rows_data = [r for r in rows_data if r[-1] > 0 or r[0] == "Total"]
        rows_data.insert(0, img_row)

        col_headers = ["Class"] + [
            f"{DS_SHORT[ds]}\n{COND_LABELS[c]}" for ds, c, _, _ in col_specs
        ] + ["Total"]

        n_rows = len(rows_data) + 1
        n_cols = len(col_headers)
        tbl_h  = 2.55
        tbl_w  = SLIDE_W - 0.5
        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.25), Inches(top), Inches(tbl_w), Inches(tbl_h)
        ).table

        col_w_cls  = tbl_w * 0.20
        col_w_tot  = tbl_w * 0.10
        col_w_rest = (tbl_w - col_w_cls - col_w_tot) / (n_cols - 2)
        tbl_shape.columns[0].width = Inches(col_w_cls)
        for ci in range(1, n_cols - 1):
            tbl_shape.columns[ci].width = Inches(col_w_rest)
        tbl_shape.columns[n_cols - 1].width = Inches(col_w_tot)

        def _c(r, c, text, bold=False, bg=None, fg="#222222", align=PP_ALIGN.CENTER):
            cell = tbl_shape.cell(r, c)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.clear()
            p.alignment = align
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(8.5)
            run.font.bold = bold
            run.font.color.rgb = _hex2rgb(fg)
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex2rgb(bg)
            else:
                cell.fill.background()

        for ci, ch in enumerate(col_headers):
            ds_col = col_specs[ci - 1][0] if 1 <= ci <= len(col_specs) else None
            hdr_bg = DS_COLORS.get(ds_col, "#2c3e50") if ds_col else "#2c3e50"
            _c(0, ci, ch, bold=True, bg=hdr_bg, fg="#ffffff",
               align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

        for ri, row in enumerate(rows_data):
            is_total  = row[0] == "Total"
            is_images = row[0] == "# Images"
            bg = ("#dde8f0" if is_images else
                  "#e8f0fe" if is_total else
                  "#f5f5f5" if ri % 2 == 0 else "#ffffff")
            for ci, val in enumerate(row):
                if is_images:
                    _c(ri + 1, ci, str(int(val)) if ci > 0 else val,
                       bold=True, bg=bg, fg="#2c3e50",
                       align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
                elif ci == 0:
                    fg = FA5_COLORS.get(val, "#333333") if not is_total else "#2c3e50"
                    _c(ri + 1, ci, val, bold=True, bg=bg, fg=fg,
                       align=PP_ALIGN.LEFT)
                else:
                    _c(ri + 1, ci, "" if (val == 0 and not is_total) else str(int(val)),
                       bold=is_total, bg=bg, fg="#2c3e50" if is_total else "#333333",
                       align=PP_ALIGN.CENTER)


def _slide_summary(prs, b1_dfs, b2_dfs):
    """Table: rows=classes, cols=Batch1/Batch2 per dataset."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    b1_total = sum(len(df) for df in b1_dfs.values())
    b2_total = sum(len(df) for df in b2_dfs.values())
    _header(slide, "Grand Summary: All Labels — Both Batches",
            subtitle=f"Batch 1 (LabelStudio): {b1_total:,} patches  ·  "
                     f"Batch 2 (Prototype): {b2_total:,} patches  ·  "
                     f"Grand total: {b1_total + b2_total:,} patches")

    # Columns: Class | DS1-B1 | DS1-B2 | DS2-B1 | DS2-B2 | DS3-B1 | DS3-B2 | Total
    col_specs = []
    for ds_key in DS_KEYS:
        col_specs.append((ds_key, "B1", b1_dfs[ds_key], B1_LABEL_COL))
        col_specs.append((ds_key, "B2", b2_dfs[ds_key], B2_LABEL_COL))

    # Image counts per column (per-batch)
    img_counts_per_col = []
    for ds_key, batch, df, lc in col_specs:
        dc = df.copy()
        if "condition" not in dc.columns and "filename" in dc.columns:
            dc["condition"] = dc["filename"].str.extract(r'^(control|ycomp)', expand=False)
        cond_col = "condition" if "condition" in dc.columns else B1_COND_COL
        img_counts_per_col.append(_n_images(dc, cond_col, ds_key=ds_key).get("Total", 0))
    # Total = unique images across B1+B2 per dataset (deduplicated), summed across datasets
    combined_total = sum(
        _n_images_combined(b1_dfs[ds_key], b2_dfs[ds_key], ds_key).get("Total", 0)
        for ds_key in DS_KEYS
    )
    img_row = ["# Images"] + img_counts_per_col + [combined_total]

    all_cls = FA5_ORDER + ["Uncertain"]
    rows_data = []
    for cls in all_cls:
        row = [cls]
        total = 0
        for ds_key, batch, df, lc in col_specs:
            cnt = int((df[lc] == cls).sum())
            row.append(cnt)
            total += cnt
        row.append(total)
        rows_data.append(row)

    total_row = ["Total"] + [sum(r[j + 1] for r in rows_data) for j in range(len(col_specs))]
    total_row.append(sum(total_row[1:]))
    rows_data.append(total_row)
    rows_data = [r for r in rows_data if r[-1] > 0 or r[0] == "Total"]
    rows_data.insert(0, img_row)

    col_headers = ["Class"] + [
        f"{DS_SHORT[ds]}\n{b}" for ds, b, _, _ in col_specs
    ] + ["Total"]

    n_rows = len(rows_data) + 1
    n_cols = len(col_headers)
    tbl_w  = SLIDE_W - 0.5
    tbl_h  = 5.5

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.25), Inches(1.3), Inches(tbl_w), Inches(tbl_h)
    ).table

    col_w_cls  = tbl_w * 0.18
    col_w_tot  = tbl_w * 0.08
    col_w_rest = (tbl_w - col_w_cls - col_w_tot) / (n_cols - 2)
    tbl_shape.columns[0].width = Inches(col_w_cls)
    for ci in range(1, n_cols - 1):
        tbl_shape.columns[ci].width = Inches(col_w_rest)
    tbl_shape.columns[n_cols - 1].width = Inches(col_w_tot)

    def _c(r, c, text, bold=False, bg=None, fg="#222222", align=PP_ALIGN.CENTER):
        cell = tbl_shape.cell(r, c)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.clear()
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(8.5)
        run.font.bold = bold
        run.font.color.rgb = _hex2rgb(fg)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex2rgb(bg)
        else:
            cell.fill.background()

    for ci, ch in enumerate(col_headers):
        ds_col = col_specs[ci - 1][0] if 1 <= ci <= len(col_specs) else None
        hdr_bg = DS_COLORS.get(ds_col, "#2c3e50") if ds_col else "#2c3e50"
        _c(0, ci, ch, bold=True, bg=hdr_bg, fg="#ffffff",
           align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    for ri, row in enumerate(rows_data):
        is_total  = row[0] == "Total"
        is_images = row[0] == "# Images"
        bg = ("#dde8f0" if is_images else
              "#e8f0fe" if is_total else
              "#f5f5f5" if ri % 2 == 0 else "#ffffff")
        for ci, val in enumerate(row):
            if is_images:
                _c(ri + 1, ci, str(int(val)) if ci > 0 else val,
                   bold=True, bg=bg, fg="#2c3e50",
                   align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            elif ci == 0:
                fg = FA5_COLORS.get(val, "#333333") if not is_total else "#2c3e50"
                _c(ri + 1, ci, val, bold=True, bg=bg, fg=fg, align=PP_ALIGN.LEFT)
            else:
                _c(ri + 1, ci, "" if (val == 0 and not is_total) else str(int(val)),
                   bold=is_total, bg=bg,
                   fg="#2c3e50" if is_total else "#333333", align=PP_ALIGN.CENTER)

    _txt(slide,
         "Cross-batch overlap: 56 Dataset 1, 21 Dataset 2, 26 Dataset 3 patches in both batches "
         "(28 label conflicts). Batch 2 labels take precedence.",
         0.3, 7.1, SLIDE_W - 0.6, 0.35, size=8, color="#888888", italic=True)


# ---------------------------------------------------------------------------
# Image-level helpers
# ---------------------------------------------------------------------------

def _cond_abbr(c: str) -> str:
    return "ctrl" if c == "control" else c


def _image_df_b1(df: pd.DataFrame, ds_key: str) -> pd.DataFrame:
    """Return DataFrame with [image_id, short_id, condition, n_patches] for Batch 1."""
    df = df.copy()
    if ds_key == "ds2":
        df["frame_id"] = df["unique_ID"].str.extract(r'-(f\d+)x', expand=False)
        df["image_id"] = df["condition"] + "_" + df["frame_id"]
        df["short_id"] = ds_key + "_" + df["condition"].map(_cond_abbr) + "_" + df["frame_id"]
    else:
        # czi_filename always says "control-XX.czi" even for ycomp patches — use condition
        # column from CSV and 0-indexed frame number for correct per-condition image IDs
        frame_f = _czi_to_fframe(df["czi_filename"]).fillna("?")
        df["image_id"] = df["condition"] + "_" + frame_f
        df["short_id"] = ds_key + "_" + df["condition"].map(_cond_abbr) + "_" + frame_f
    grp = (df.groupby(["image_id", "short_id", "condition"])
             .size().reset_index(name="n_patches"))
    return grp.sort_values(["condition", "n_patches"], ascending=[True, False])


def _image_df_b2(df: pd.DataFrame, ds_key: str) -> pd.DataFrame:
    """Return DataFrame with [image_id, short_id, condition, n_patches] for Batch 2."""
    df = df.copy()
    df["condition"] = df["filename"].str.extract(r'^(control|ycomp)', expand=False)
    df["frame_id"]  = df["filename"].str.extract(r'^(?:control|ycomp)_(f\d+)', expand=False)
    df["image_id"]  = df["condition"] + "_" + df["frame_id"]
    df["short_id"]  = ds_key + "_" + df["condition"].map(_cond_abbr) + "_" + df["frame_id"]
    grp = (df.groupby(["image_id", "short_id", "condition"])
             .size().reset_index(name="n_patches"))
    return grp.sort_values(["condition", "n_patches"], ascending=[True, False])


def _fig_patches_per_image(img_dfs: list[tuple], title: str,
                           ref_n_bars: int | None = None) -> Image.Image:
    """Bar chart: one bar per source image, grouped by dataset panel, coloured by condition.
    img_dfs: list of (ds_key, label, img_df) where img_df has [image_id, condition, n_patches].
    """
    COND_COLORS = {"control": "#4e79a7", "ycomp": "#f28e2b"}
    BAR_W_IN  = 0.22   # inches per bar slot — same across all figures
    FIG_H_IN  = 4.2
    FS        = 8      # unified font size for every text element

    n_panels = len(img_dfs)
    n_bars   = [max(1, len(d[2])) for d in img_dfs]

    if ref_n_bars is not None:
        # Uniform-width mode: every figure uses the same x-span so bars have
        # identical physical width; small panels get whitespace padding.
        fig_w = ref_n_bars * BAR_W_IN
        width_ratios = [ref_n_bars] * n_panels
    else:
        fig_w = max(3.0 * n_panels, sum(n_bars) * BAR_W_IN)
        width_ratios = n_bars

    fig, axes = plt.subplots(1, n_panels, figsize=(fig_w, FIG_H_IN),
                             gridspec_kw={"width_ratios": width_ratios})
    if n_panels == 1:
        axes = [axes]

    for ax, (ds_key, batch_label, idf) in zip(axes, img_dfs):
        if idf.empty:
            ax.text(0.5, 0.5, "no data", transform=ax.transAxes,
                    ha="center", va="center", fontsize=FS, color="#aaa")
            ax.set_title(f"{DS_SHORT[ds_key]}\n{batch_label}", fontsize=FS,
                         fontweight="bold", color=DS_COLORS[ds_key])
            ax.axis("off")
            continue

        conds_present = [c for c in CONDITIONS if c in idf["condition"].values]
        bars_df = idf.copy().reset_index(drop=True)
        bars_df = pd.concat([
            bars_df[bars_df["condition"] == c].sort_values("n_patches", ascending=False)
            for c in conds_present
        ]).reset_index(drop=True)

        x = range(len(bars_df))
        colors = [COND_COLORS.get(c, "#888") for c in bars_df["condition"]]
        ax.bar(x, bars_df["n_patches"], color=colors, linewidth=0, width=0.75)

        for xi, (_, row) in enumerate(bars_df.iterrows()):
            ax.text(xi, row["n_patches"] + bars_df["n_patches"].max() * 0.01,
                    str(row["n_patches"]), ha="center", va="bottom", fontsize=FS)

        if len(conds_present) == 2:
            n_ctrl = (bars_df["condition"] == "control").sum()
            ax.axvline(n_ctrl - 0.5, color="#aaa", linewidth=0.8, linestyle="--")

        n_imgs   = len(bars_df)
        mean_ppi = bars_df["n_patches"].mean()
        ax.set_title(
            f"{DS_SHORT[ds_key]}  ·  {batch_label}\n"
            f"{n_imgs} images  ·  mean {mean_ppi:.0f} patches/image",
            fontsize=FS, fontweight="bold", color=DS_COLORS[ds_key], pad=5
        )
        ax.set_xticks(list(x))
        ax.set_xticklabels(bars_df["short_id"].tolist(), rotation=90,
                           ha="right", fontsize=FS)
        ax.set_ylabel("# patches", fontsize=FS)
        ax.tick_params(axis="y", labelsize=FS)
        ax.spines[["top", "right"]].set_visible(False)
        # Expand x-axis so bars occupy the same physical width as the reference
        if ref_n_bars is not None and ref_n_bars > len(bars_df):
            pad = (ref_n_bars - len(bars_df)) / 2
            ax.set_xlim(-0.5 - pad, len(bars_df) - 0.5 + pad)

    fig.suptitle(title, fontsize=FS, fontweight="bold", y=1.02)
    fig.tight_layout()
    return _fig_to_pil(fig)


def _slide_images_per_batch(prs, b1_dfs, b2_dfs, batch: int):
    """One slide showing patches-per-image for Batch 1 or Batch 2."""
    if batch == 1:
        img_dfs = [
            ("ds1", "Batch 1", _image_df_b1(b1_dfs["ds1"], "ds1")),
            ("ds2", "Batch 1", _image_df_b1(b1_dfs["ds2"], "ds2")),
            ("ds3", "Batch 1", _image_df_b1(b1_dfs["ds3"], "ds3")),
        ]
        title_txt = "Batch 1 (LabelStudio): Images Labeled & Patches per Image"
        n_imgs = {k: len(d) for k, _, d in img_dfs}
        subtitle = (
            f"Dataset 1: {n_imgs['ds1']} images  ·  "
            f"Dataset 2: {n_imgs['ds2']} images  ·  "
            f"Dataset 3: {n_imgs['ds3']} images  ·  "
            "Each bar = one source image. Bars grouped by condition."
        )
    else:
        img_dfs = [
            ("ds1", "Batch 2", _image_df_b2(b2_dfs["ds1"], "ds1")),
            ("ds2", "Batch 2", _image_df_b2(b2_dfs["ds2"], "ds2")),
            ("ds3", "Batch 2", _image_df_b2(b2_dfs["ds3"], "ds3")),
        ]
        title_txt = "Batch 2 (Prototype): Images Labeled & Patches per Image"
        n_imgs = {k: len(d) for k, _, d in img_dfs}
        subtitle = (
            f"Dataset 1: {n_imgs['ds1']} images  ·  "
            f"Dataset 2: {n_imgs['ds2']} images  ·  "
            f"Dataset 3: {n_imgs['ds3']} images  ·  "
            "Each bar = one source image. Bars grouped by condition."
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title_txt, subtitle=subtitle)
    img = _fig_patches_per_image(img_dfs, "")
    _paste_pil(slide, img, 0.3, 1.1, SLIDE_W - 0.6, 6.15)


# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print("Loading label CSVs...")
    b1_dfs = {ds: pd.read_csv(p) for ds, p in BATCH1_FILES.items()}
    b2_dfs = {ds: pd.read_csv(p) for ds, p in BATCH2_FILES.items()}

    # Batch 1 ppax has annotators 3/4 (integers), cast to string
    b1_dfs["ds3"]["annotator"] = b1_dfs["ds3"]["annotator"].astype(str)

    print("Building slides...")
    _slide_title(prs)
    _slide_terminology(prs)
    _slide_taxonomy(prs)
    _slide_batch1_sessions(prs)
    _slide_batch1_ds1(prs, b1_dfs)
    _slide_batch1_ds23(prs, b1_dfs)
    _slide_images_per_batch(prs, b1_dfs, b2_dfs, batch=1)
    _slide_batch2_sessions(prs)
    _slide_batch2_ds1(prs, b2_dfs)
    _slide_batch2_ds23(prs, b2_dfs)
    _slide_images_per_batch(prs, b1_dfs, b2_dfs, batch=2)
    _slide_summary(prs, b1_dfs, b2_dfs)
    _slide_summary_by_condition(prs, b1_dfs, b2_dfs)

    RES.mkdir(exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}  ({OUT.stat().st_size // 1024} KB)")

    # Save per-dataset patch-per-image plots as individual PNGs
    print("Saving individual per-dataset bar charts...")
    img_sources = {
        1: {
            "ds1": _image_df_b1(b1_dfs["ds1"], "ds1"),
            "ds2": _image_df_b1(b1_dfs["ds2"], "ds2"),
            "ds3": _image_df_b1(b1_dfs["ds3"], "ds3"),
        },
        2: {
            "ds1": _image_df_b2(b2_dfs["ds1"], "ds1"),
            "ds2": _image_df_b2(b2_dfs["ds2"], "ds2"),
            "ds3": _image_df_b2(b2_dfs["ds3"], "ds3"),
        },
    }
    # Use the largest bar count as reference so all figures share the same bar width
    ref_n = max(len(idf) for ds_imgs in img_sources.values() for idf in ds_imgs.values())
    for batch, ds_imgs in img_sources.items():
        for ds_key, idf in ds_imgs.items():
            fig_img = _fig_patches_per_image(
                [(ds_key, f"Batch {batch}", idf)],
                f"{DS_SHORT[ds_key]} — Batch {batch}: Patches per Image",
                ref_n_bars=ref_n,
            )
            out_path = RES / f"label_overview_b{batch}_{ds_key}_patches_per_image.png"
            fig_img.save(str(out_path))
            print(f"  {out_path.name}")


if __name__ == "__main__":
    main()
