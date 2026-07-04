#!/usr/bin/env python3
"""
Generate a PowerPoint explaining the reconstruction quality panel and violin plots.

Usage:
  python scripts/make_quality_pptx.py <variant_dir> [--out <output.pptx>]
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── colour palette ────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_MID    = RGBColor(0xF2, 0xF5, 0xFA)   # very light blue-grey (alt row)
C_ACCENT = RGBColor(0x1F, 0x5C, 0x99)   # medium blue (box headers)
C_GOLD   = RGBColor(0x1A, 0x3A, 0x6B)   # dark navy (section headers)
C_WHITE  = RGBColor(0x11, 0x11, 0x11)   # near-black (body text)
C_LIGHT  = RGBColor(0x33, 0x33, 0x33)   # dark grey (secondary text)
C_GREY   = RGBColor(0x66, 0x66, 0x66)   # mid grey (captions)


def _bg(slide, prs):
    """Fill slide background with white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK


def _add_title(slide, text, top=Inches(0.25), fontsize=32, color=C_WHITE):
    txb = slide.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.7))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(fontsize)
    run.font.bold   = True
    run.font.color.rgb = color
    return txb


def _add_subtitle(slide, text, top=Inches(0.9), fontsize=16, color=C_GOLD):
    txb = slide.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.4))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fontsize)
    run.font.italic = True
    run.font.color.rgb = color
    return txb


def _add_body(slide, lines, left=Inches(0.5), top=Inches(1.35),
              width=Inches(9.0), height=Inches(5.5),
              fontsize=14, color=C_WHITE, bullet=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        if isinstance(line, tuple):
            text, size, bold, col, indent = line
        else:
            text  = line
            size  = fontsize
            bold  = False
            col   = color
            indent = bullet
        run = p.add_run()
        run.text = ("• " if indent else "") + text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = col
    return txb


def _add_image(slide, img_path, left, top, width=None, height=None):
    if not Path(img_path).exists():
        return None
    try:
        if width and height:
            return slide.shapes.add_picture(str(img_path), left, top, width, height)
        elif width:
            return slide.shapes.add_picture(str(img_path), left, top, width=width)
        elif height:
            return slide.shapes.add_picture(str(img_path), left, top, height=height)
        else:
            return slide.shapes.add_picture(str(img_path), left, top)
    except Exception as e:
        print(f"  [warn] could not embed {img_path}: {e}")
        return None


def _add_box(slide, text, left, top, width, height,
             bg=C_ACCENT, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=12):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xBB, 0xCC, 0xDD); shape.line.width = Pt(0.75)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top  = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(fontsize); run.font.color.rgb = fg
    return shape


def make_pptx(variant_dir: Path, out_path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # blank layout

    qp_dir   = variant_dir / "quality_panels"
    qpb_dir  = variant_dir / "quality_panels_bulk"
    viol_dir = variant_dir / "violin_plots"
    model_name = variant_dir.name

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Reconstruction Quality Analysis", top=Inches(2.2), fontsize=36)
    _add_subtitle(sl, "Linking metric values to visual patch appearance", top=Inches(3.0), fontsize=18)
    _add_body(sl, [f"Model: {model_name}"], top=Inches(3.8),
              fontsize=13, color=C_GREY, bullet=False)
    _add_body(sl, [
        "Panel plots  →  what do patches at each quality level look like?",
        "Violin plots →  how is quality distributed across FA types and datasets?",
    ], top=Inches(4.4), fontsize=13, color=C_LIGHT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Motivation
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Why These Plots?")
    _add_subtitle(sl, "Connecting numbers to biology")
    _add_body(sl, [
        ("The problem", 16, True, C_GOLD, False),
        "Reconstruction metrics (L1, MSE, Hessian L1) are aggregate numbers.",
        "A value of 0.05 vs 0.12 is meaningless without visual reference.",
        "We cannot tell if high error = blurry recon, wrong structure, or just noisy patch.",
        "",
        ("The approach", 16, True, C_GOLD, False),
        "Divide the metric distribution into 9 decile bands (10th – 90th percentile).",
        "For each band: show example raw patches alongside their reconstructions.",
        "This lets reviewers directly judge reconstruction fidelity at each quality level.",
        "",
        ("Key design choice", 16, True, C_GOLD, False),
        "Percentile boundaries are computed from the same pooled reference (e.g. all vinc",
        "labelled patches) so panels across FA types are on the same scale.",
    ], fontsize=13)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Three metrics explained
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Three Reconstruction Metrics")
    _add_subtitle(sl, "Each captures a different aspect of reconstruction error")

    boxes = [
        ("L1  (MAE)", "Mean absolute pixel error.\nLow = globally close in intensity.\nSensitive to brightness offset."),
        ("MSE", "Mean squared pixel error.\nPenalises large deviations more.\nAffected by outlier pixels."),
        ("Hessian L1", "Mean Frobenius norm of the\nHessian of the residual.\nMeasures curvature-level error —\nhigh = structure/texture not captured."),
    ]
    for i, (title, body) in enumerate(boxes):
        x = Inches(0.4 + i * 3.1)
        _add_box(sl, title, x, Inches(1.5), Inches(2.8), Inches(0.5),
                 bg=C_ACCENT, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=14)
        _add_box(sl, body,  x, Inches(2.1), Inches(2.8), Inches(2.5),
                 bg=RGBColor(0xE8,0xF2,0xFF), fg=RGBColor(0x11,0x11,0x11), fontsize=12)

    _add_body(sl, [
        "Note: Hessian L1 is computed on the residual (raw − recon), not on the raw image.",
        "This means high Hessian L1 = the error image has curvature = missed fine structure.",
        "A smoothed reconstruction will have high Hessian L1 even if global L1 is low.",
    ], top=Inches(5.0), fontsize=12, color=C_LIGHT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Panel plot design
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Panel Plot Design")
    _add_subtitle(sl, "make_recon_quality_panels.py  |  make_recon_quality_panels_bulk.py")
    _add_body(sl, [
        ("Layout", 14, True, C_GOLD, False),
        "Each panel: 1 row raw + 1 row recon × 6 columns  (= 6 patch pairs)",
        "Raw/recon share the same intensity scale per column for fair comparison.",
        "Per-patch title shows train or val split.",
        "Empty slots left blank if fewer than 6 patches fall in that band.",
        "",
        ("Percentile windows", 14, True, C_GOLD, False),
        "9 panels per metric per group: centred at 10th, 20th, …, 90th percentile.",
        "Window = [pct(P−2.5), pct(P+2.5)] of the pooled reference distribution.",
        "Vinc labelled panels: reference = all labelled patches pooled (FA types + splits).",
        "Vinc unlabelled panels: reference = all unlabelled vinc patches.",
        "ppax panels: reference = all ppax patches (control + ycomp).",
    ], fontsize=12, top=Inches(1.3))

    # embed one example panel
    ex = qp_dir / "focalad-l1-50p.png"
    _add_image(sl, ex, Inches(6.3), Inches(2.8), width=Inches(3.5))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Panel groups (labelled vinc)
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Panel Groups — Labelled Vinc")
    _add_subtitle(sl, "FA type × metric × decile  |  quality_panels/")

    _add_body(sl, [
        "4 FA types (No adhesion excluded):  Nascent Adhesion · Focal Complex · Focal Adhesion · Fibrillar",
        "3 metrics × 9 deciles = 27 panels per FA type  (up to 108 panels total).",
        "Filename: {fatype}-{metric}-{P}p.png   e.g.  focalad-l1-50p.png",
        "",
        "What to look for:",
        "  Low decile panels (10–30p): model reconstructs well — sharp, correct morphology.",
        "  High decile panels (70–90p): model struggles — blurry, missing structure.",
        "  Compare FA types at the same decile: is one class systematically harder?",
        "  Fibrillar has very few patches (~16) — expect sparse coverage at some deciles.",
    ], fontsize=12, top=Inches(1.3))

    # embed two side-by-side examples
    for i, (fname, label) in enumerate([
        ("focalad-l1-10p.png", "Focal adhesion  L1 pct10"),
        ("focalad-l1-90p.png", "Focal adhesion  L1 pct90"),
    ]):
        path = qp_dir / fname
        _add_image(sl, path, Inches(0.3 + i * 4.8), Inches(4.5), width=Inches(4.5))
        _add_body(sl, [label], top=Inches(7.0), left=Inches(0.3 + i * 4.8),
                  width=Inches(4.5), height=Inches(0.4), fontsize=10,
                  color=C_GREY, bullet=False)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Panel groups (unlabelled vinc + ppax)
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Panel Groups — Unlabelled Vinc & ppax")
    _add_subtitle(sl, "30 patches per panel  |  quality_panels_bulk/")

    _add_body(sl, [
        ("Unlabelled vinc  (subset: unlabelled)", 13, True, C_GOLD, False),
        "~23 000 patches with no FA-type annotation (majority of training data).",
        "Percentile reference: all unlabelled vinc patches pooled.",
        "Train/val label shown per patch.  Filename: unlabelled-{metric}-{P}p.png",
        "",
        ("ppax  (subset: ppax)", 13, True, C_GOLD, False),
        "~5 800 patches from a different cell line — never seen during training.",
        "Inference run on-the-fly from model_final.pt.",
        "Condition (control / ycomp) shown as patch label.",
        "Percentile reference: all ppax patches pooled (control + ycomp).",
        "Filename: ppax-{metric}-{P}p.png",
        "",
        ("Key comparison", 13, True, C_GOLD, False),
        "Compare the L1 decile ranges: if ppax pct50 ≈ vinc pct80, the model",
        "generalises poorly — typical patches from ppax look like hard cases in vinc.",
    ], fontsize=11, top=Inches(1.3))

    ex = qpb_dir / "ppax-l1-50p.png"
    _add_image(sl, ex, Inches(7.0), Inches(3.5), width=Inches(2.8))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Violin plot design
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Violin Plot Design")
    _add_subtitle(sl, "plot_recon_metric_violins.py  |  violin_plots/")

    _add_body(sl, [
        ("Vinc plot", 13, True, C_GOLD, False),
        "x-axis: 4 FA types (No adhesion excluded).",
        "For each FA type: train violin (blue) and val violin (orange) side by side.",
        "9 horizontal lines = decile boundaries from the FA-type's own pooled distribution.",
        "  → train and val violins for the same FA type share identical lines.",
        "",
        ("ppax plot", 13, True, C_GOLD, False),
        "x-axis: control and ycomp conditions.",
        "9 lines from the pooled ppax distribution (control + ycomp).",
        "",
        ("What to look for", 13, True, C_GOLD, False),
        "Train vs val gap: large gap suggests overfitting in reconstruction.",
        "FA type differences: Fibrillar often harder (higher L1) — less training data.",
        "Hessian L1 > L1 pattern: model may smooth fine structure while getting mean right.",
        "ppax vs vinc scale: if ppax violins shift upward, model generalises poorly.",
    ], fontsize=11, top=Inches(1.3))

    ex = viol_dir / "vinc_recon_l1.png"
    _add_image(sl, ex, Inches(5.8), Inches(2.0), width=Inches(4.0))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Violin examples (vinc + ppax side by side)
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Violin Plots — Vinc vs ppax")
    _add_subtitle(sl, f"Model: {model_name}")

    for i, (fname, label, metric) in enumerate([
        ("vinc_recon_l1.png",  "Vinc  —  L1 (MAE)", "recon_l1"),
        ("ppax_recon_l1.png",  "ppax  —  L1 (MAE)", "recon_l1"),
    ]):
        path = viol_dir / fname
        _add_image(sl, path, Inches(0.2 + i * 5.0), Inches(1.4), width=Inches(4.7))
        _add_body(sl, [label], top=Inches(5.85), left=Inches(0.2 + i * 5.0),
                  width=Inches(4.7), height=Inches(0.3), fontsize=10,
                  color=C_GOLD, bullet=False)

    _add_body(sl, [
        "Lines inside each violin = same 9 decile boundaries used in the panel plots.",
        "A patch shown in panel pct50 corresponds to the median line in the violin.",
    ], top=Inches(6.5), fontsize=11, color=C_LIGHT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Connecting panels to violins
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Connecting Violin Lines to Patch Panels")
    _add_subtitle(sl, "Same percentile reference — direct lookup")

    _add_body(sl, [
        "The 9 lines in each violin mark exactly the boundaries used to select patches",
        "for the corresponding panel plots.",
        "",
        ("Example workflow", 14, True, C_GOLD, False),
        "1.  Look at the vinc L1 violin for Focal Adhesion.",
        "2.  The 7th line (70th pct) separates moderate from high-error patches.",
        "3.  Open  focalad-l1-70p.png  to see what those patches look like.",
        "4.  If the recon row looks blurry / misses puncta → model loses structure at this level.",
        "",
        ("Cross-dataset use", 14, True, C_GOLD, False),
        "The ppax violin uses its own pooled reference (ppax-only percentiles).",
        "To compare ppax and vinc on the same absolute scale, check where the",
        "ppax median falls relative to the vinc violin body.",
    ], fontsize=12, top=Inches(1.3))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Summary and file map
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl, prs)
    _add_title(sl, "Summary — Output File Map")

    rows = [
        ("quality_panels/",          "{fatype}-{metric}-{P}p.png",
         "Labelled vinc, per FA type, global vinc percentiles"),
        ("quality_panels_bulk/",     "unlabelled-{metric}-{P}p.png",
         "Unlabelled vinc, own percentile reference"),
        ("quality_panels_bulk/",     "ppax-{metric}-{P}p.png",
         "ppax inference, own percentile reference"),
        ("violin_plots/",            "vinc_{metric}.png",
         "Vinc FA type × train/val, same lines as panels"),
        ("violin_plots/",            "ppax_{metric}.png",
         "ppax control/ycomp, same lines as panels"),
    ]
    col_widths = [Inches(2.5), Inches(3.2), Inches(3.8)]
    headers    = ["Directory", "Filename pattern", "Description"]
    top_start  = Inches(1.3)
    row_h      = Inches(0.55)

    for ci, (hdr, w) in enumerate(zip(headers, col_widths)):
        x = Inches(0.3) + sum(col_widths[:ci])
        _add_box(sl, hdr, x, top_start, w, row_h, bg=C_ACCENT, fontsize=11)

    for ri, row in enumerate(rows):
        for ci, (cell, w) in enumerate(zip(row, col_widths)):
            x   = Inches(0.3) + sum(col_widths[:ci])
            y   = top_start + row_h * (ri + 1)
            bg  = RGBColor(0xF2,0xF5,0xFA) if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
            _add_box(sl, cell, x, y, w, row_h, bg=bg,
                     fg=RGBColor(0x11,0x11,0x11), fontsize=10)

    _add_body(sl, [
        "Metrics:  recon_l1 · recon_mse · recon_hessian_l1     |     P: 10, 20, …, 90",
    ], top=Inches(6.8), fontsize=11, color=C_GREY, bullet=False)

    # ── save ─────────────────────────────────────────────────────────────────
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("variant_dir", type=Path)
    parser.add_argument("--out", type=Path, default=None)
    args = parser.parse_args()

    vdir = args.variant_dir
    if not vdir.is_dir():
        sys.exit(f"Not a directory: {vdir}")

    out = args.out or vdir / "reconstruction_quality_overview.pptx"
    make_pptx(vdir, out)


if __name__ == "__main__":
    main()
