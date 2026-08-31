#!/usr/bin/env python3
"""
make_pptx_perclass_analysis.py
==============================
Comprehensive PPT on per-class F1 analysis for FA4 cross-dataset classification.

Slides
------
 1. Cover
 2. 4-class setup: FA hierarchy + class imbalance per dataset
 3. Baseline (A_zrecon) — full 5-scenario per-class F1 lines
 4. Hard-class deep dive — FC and Fib across scenarios and fractions
 5. Effect of projection features (zproj vs zrecon) — side-by-side
 6. Effect of SMOTE — side-by-side on hard classes
 7. Summary heatmap — per-class F1 at 75% × all scenarios × 4 variants
 8. Takeaways

Usage
-----
  python scripts/make_pptx_perclass_analysis.py
  python scripts/make_pptx_perclass_analysis.py --out custom_name.pptx
"""

from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
EVAL_DIR  = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run/fa4_xds_eval")
LABEL_DIR = Path("/net/projects/CLS/lding/data/fa_data_analysis/labelling")
OUT_DEFAULT = Path("results/fa4_perclass_analysis.pptx")

FA_LABELS   = ["Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
LABEL_SHORT = {"Nascent Adhesion": "NA", "focal complex": "FC",
               "focal adhesion": "FA", "fibrillar adhesion": "Fib"}
CLASS_COLORS = {"NA": "#4393c3", "FC": "#f4a582", "FA": "#2ca02c", "Fib": "#9467bd"}

SCENARIOS = ["vinc_only", "pfak_only", "vinc->pfak", "pfak->vinc", "combined"]
SCENARIO_LABELS = {
    "vinc_only":  "Vinc only",
    "pfak_only":  "pFAK only",
    "vinc->pfak": "Vinc→pFAK",
    "pfak->vinc": "pFAK→Vinc",
    "combined":   "Combined",
}

VARIANTS = ["A_zrecon", "A_zproj", "A_zrecon_smote", "A_zproj_smote"]
VARIANT_LABELS = {
    "A_zrecon":       "zrecon (baseline)",
    "A_zproj":        "zproj (proj head)",
    "A_zrecon_smote": "zrecon + SMOTE",
    "A_zproj_smote":  "zproj + SMOTE",
}

LABEL_FILES = {
    "vinc_ctrl":  LABEL_DIR / "vinc_control_label_Annabel_20260715_1554.csv",
    "vinc_ycomp": LABEL_DIR / "vinc_combined_label_Annabel_20260816.csv",
    "pfak_ctrl":  LABEL_DIR / "pfak_combined_label_Annabel_aug2026.csv",
}
DS_LABELS = {"vinc_ctrl": "Vinc/Ctrl", "vinc_ycomp": "Vinc/Ycomp", "pfak_ctrl": "pFAK/Ctrl"}

# ── slide geometry ─────────────────────────────────────────────────────────────
SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.15)
BODY_T  = TITLE_H + Inches(0.08)
BODY_H  = SH - BODY_T - PAD

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_GREEN = RGBColor(0x37, 0x8B, 0x4A)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)

# ── helpers ────────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2*PAD, Inches(0.25),
             subtitle, size=9, color=C_GREY)

def _fig_to_pil(fig) -> Image.Image:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    return Image.open(buf).convert("RGB")

def _place_pil(slide, pil, l, t, max_w, max_h):
    iw, ih = pil.size
    scale  = min(max_w / Inches(iw / 150), max_h / Inches(ih / 150))
    rw, rh = Inches(iw / 150) * scale, Inches(ih / 150) * scale
    buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                             t + (max_h - rh) / 2, rw, rh)

def _place_png(slide, path, l, t, max_w, max_h, ph="[pending]"):
    p = Path(path)
    if not p.exists():
        _txt(slide, l, t + max_h/2, max_w, Inches(0.3),
             ph, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
        return False
    pil = Image.open(str(p)).convert("RGB")
    _place_pil(slide, pil, l, t, max_w, max_h)
    return True

def _bullet_box(slide, l, t, w, h, bullets: list[tuple[str, RGBColor]], title=""):
    _rect(slide, l, t, w, h)
    y = t
    if title:
        _txt(slide, l + Inches(0.08), y, w - Inches(0.08), Inches(0.3),
             title, size=10, bold=True, color=C_DARK)
        y += Inches(0.3)
    line_h = Inches(0.29)
    for text, color in bullets:
        _txt(slide, l + Inches(0.08), y, w - Inches(0.1), line_h,
             f"• {text}", size=9, color=color, wrap=True)
        y += line_h

def _load_results_at_frac(suffix: str, frac: float) -> dict[str, dict]:
    """Return {scenario: {cls_short: mean_f1}} at a given fraction."""
    out = {}
    for sc in SCENARIOS:
        p = EVAL_DIR / f"results_{sc}_{suffix}.csv"
        if not p.exists():
            out[sc] = None
            continue
        df = pd.read_csv(p)
        df_f = df[df["frac"] == frac]
        if df_f.empty:
            out[sc] = None
            continue
        row = {}
        for cls in FA_LABELS:
            short = LABEL_SHORT[cls]
            col = f"f1_{short}"
            if col in df_f.columns:
                row[short] = df_f[col].mean()
        out[sc] = row
    return out

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(1), Inches(1.8), SW - Inches(2), Inches(1.2),
         "FA4 Per-Class F1 Analysis",
         size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.1), SW - Inches(2), Inches(0.6),
         "Where does the classifier succeed and fail across FA subtypes?",
         size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.8), SW - Inches(2), Inches(0.45),
         "Option A · Stage-2 SupCon AE (s3v1) · LightGBM · 5 eval scenarios",
         size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    lines = [
         "Variants:  z_recon (12-d)  ·  z_proj (8-d)  ·  ×  SMOTE oversampling",
    ]
    _txt(sl, Inches(1), Inches(4.3), SW - Inches(2), Inches(0.4),
         lines[0], size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), SH - Inches(0.55), SW - Inches(2), Inches(0.4),
         "2026-08-21", size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_class_setup(prs):
    """FA hierarchy + class distribution per dataset."""
    sl = _blank(prs)
    _title_bar(sl, "The 4-Class FA Subtype Problem",
               "Classes ordered by adhesion maturity  ·  severe class imbalance across all datasets")

    # ── left panel: class hierarchy description ────────────────────────────
    LW = Inches(3.8)
    lt = BODY_T + Inches(0.1)

    _txt(sl, PAD, lt, LW, Inches(0.3), "FA Maturation Hierarchy", size=11, bold=True, color=C_DARK)

    entries = [
        ("NA",  "Nascent Adhesion",   "Earliest, small, near cell edge",         "#4393c3"),
        ("FC",  "Focal Complex",      "Intermediate, dot-shaped",                 "#f4a582"),
        ("FA",  "Focal Adhesion",     "Mature, elongated, most common in labels", "#2ca02c"),
        ("Fib", "Fibrillar Adhesion", "Late-stage, fibrillar, rarest class",      "#9467bd"),
    ]
    ey = lt + Inches(0.38)
    for short, full, desc, color in entries:
        rgb = RGBColor(*[int(color.lstrip("#")[i:i+2], 16) for i in (0, 2, 4)])
        _rect(sl, PAD, ey, Inches(0.42), Inches(0.28), fill=rgb)
        _txt(sl, PAD + Inches(0.48), ey, LW - Inches(0.55), Inches(0.17),
             f"{short}  —  {full}", size=9, bold=True, color=C_BLACK)
        _txt(sl, PAD + Inches(0.48), ey + Inches(0.17), LW - Inches(0.55), Inches(0.17),
             desc, size=8, color=C_GREY)
        ey += Inches(0.62)

    _txt(sl, PAD, ey + Inches(0.05), LW, Inches(0.28),
         "Challenge: FC and Fib are rare and intermediate in appearance",
         size=9, bold=True, color=C_AMBER)

    # ── right panel: class distribution per dataset (inline bar chart) ─────
    ds_counts = {}
    for ds, path in LABEL_FILES.items():
        df = pd.read_csv(path)
        df = df[df["label"].isin(FA_LABELS)]
        ds_counts[ds] = {LABEL_SHORT[cls]: int((df["label"] == cls).sum()) for cls in FA_LABELS}

    fig, ax = plt.subplots(figsize=(8.5, 4.5), facecolor="white")
    x = np.arange(len(FA_LABELS))
    width = 0.25
    ds_color = {"vinc_ctrl": "#1f77b4", "vinc_ycomp": "#ff7f0e", "pfak_ctrl": "#d62728"}
    for i, (ds, counts) in enumerate(ds_counts.items()):
        vals = [counts[LABEL_SHORT[cls]] for cls in FA_LABELS]
        bars = ax.bar(x + (i - 1) * width, vals, width=width*0.9,
                      color=ds_color[ds], label=DS_LABELS[ds], edgecolor="white")
        for b, v in zip(bars, vals):
            if v > 0:
                ax.text(b.get_x() + b.get_width()/2, b.get_height() + 0.5, str(v),
                        ha="center", va="bottom", fontsize=8, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels([LABEL_SHORT[c] for c in FA_LABELS], fontsize=12)
    ax.set_ylabel("Labeled patch count", fontsize=11)
    ax.set_title("Class distribution — labeled FA patches per dataset", fontsize=12, fontweight="bold")
    ax.legend(fontsize=10)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    rw = SW - LW - PAD * 3
    _place_pil(sl, _fig_to_pil(fig), PAD + LW + PAD, BODY_T + Inches(0.05), rw, BODY_H - Inches(0.1))
    plt.close(fig)


def slide_baseline_overview(prs):
    """Full 5-scenario per-class F1 lines — Option A zrecon."""
    sl = _blank(prs)
    _title_bar(sl, "Per-Class F1 vs Label Fraction — Option A (z_recon baseline)",
               "5 eval scenarios  ·  4 FA classes  ·  error bars = ±1 SD  ·  "
               "NA=Nascent  FC=Focal Complex  FA=Focal Adhesion  Fib=Fibrillar")

    png = EVAL_DIR / "perclass_lines_A.png"
    img_h = BODY_H - Inches(1.5)
    _place_png(sl, png, PAD, BODY_T + Inches(0.05), SW - 2*PAD, img_h)

    # Annotation bullets below the plot
    by = BODY_T + img_h + Inches(0.12)
    _bullet_box(sl, PAD, by, SW * 0.5 - PAD, Inches(1.3), [
        ("FA (focal adhesion) is learned well across all scenarios — reaches F1 ≥ 0.7 at 10% labels", C_GREEN),
        ("NA (nascent) is moderate — learned faster within-dataset than cross-dataset", C_DARK),
        ("FC (focal complex) is the hardest class — nearly zero in most scenarios", C_RED),
        ("Fib (fibrillar) is sparse but sometimes reaches ≥ 0.5 within vinc_only", C_AMBER),
    ])
    _bullet_box(sl, SW * 0.5, by, SW * 0.5 - PAD, Inches(1.3), [
        ("pfak_only has near-zero FC and Fib F1 — not enough minority examples", C_RED),
        ("vinc→pfak: model trained on vinc generalizes better to pfak than reverse", C_DARK),
        ("pfak→vinc: FC completely unlearned (F1≈0) — domain gap dominates", C_RED),
        ("combined: modest improvement in FC across all fractions", C_AMBER),
    ], title="Cross-dataset patterns")


def slide_hardclass_deepdive(prs):
    """FC and Fib deep dive — inline comparison across scenarios."""
    sl = _blank(prs)
    _title_bar(sl, "Deep Dive: Hard Classes  —  FC (Focal Complex) and Fib (Fibrillar)",
               "Option A  z_recon  ·  error bars = ±1 SD  ·  showing within-dataset vs cross-dataset")

    # Load all scenario results for FC and Fib at all fracs
    fracs = [0.10, 0.25, 0.50, 0.75]

    fig, axes = plt.subplots(2, 5, figsize=(16, 6.5), sharey="row", facecolor="white")
    fig.suptitle("Hard class F1 (FC and Fib) across scenarios and label fractions",
                 fontsize=11, fontweight="bold")

    sc_colors = {
        "vinc_only": "#1f77b4", "pfak_only": "#d62728",
        "vinc->pfak": "#ff7f0e", "pfak->vinc": "#9467bd", "combined": "#2ca02c",
    }

    for col, sc in enumerate(SCENARIOS):
        p = EVAL_DIR / f"results_{sc}_A.csv"
        for row, (cls, short, color) in enumerate([
            ("focal complex",      "FC",  "#f4a582"),
            ("fibrillar adhesion", "Fib", "#9467bd"),
        ]):
            ax = axes[row][col]
            ax.spines[["top", "right"]].set_visible(False)
            ax.set_facecolor("white")
            ax.set_ylim(-0.05, 1.05)
            ax.axhline(0, color="lightgray", lw=0.8, zorder=0)
            if row == 0:
                ax.set_title(SCENARIO_LABELS[sc], fontsize=9, fontweight="bold",
                             color=sc_colors[sc])
            if col == 0:
                ax.set_ylabel(f"F1 ({short})", fontsize=9)
            ax.set_xticks(range(len(fracs)))
            ax.set_xticklabels([f"{int(f*100)}%" for f in fracs], fontsize=7)

            if not p.exists():
                ax.text(0.5, 0.5, "no data", ha="center", va="center",
                        transform=ax.transAxes, fontsize=8, color="gray")
                continue

            df = pd.read_csv(p)
            col_key = f"f1_{short}"
            if col_key not in df.columns:
                continue

            means = [df[df["frac"] == f][col_key].mean() for f in fracs]
            stds  = [df[df["frac"] == f][col_key].std()  for f in fracs]
            x = np.arange(len(fracs))
            ax.errorbar(x, means, yerr=stds, fmt="o-", color=color,
                        capsize=3, linewidth=2, markersize=6)

            # Annotate 75% value
            ax.text(len(fracs) - 1, means[-1] + 0.05,
                    f"{means[-1]*100:.0f}%", fontsize=7, ha="center", color=color,
                    fontweight="bold")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05),
               SW - 2*PAD, BODY_H - Inches(0.9))
    plt.close(fig)

    # Summary annotation
    by = SH - Inches(0.95)
    _txt(sl, PAD, by, SW * 0.5, Inches(0.85),
         "Key finding:  FC F1 is near zero in pfak_only and pfak→vinc — pFAK labels have too few FC examples "
         "(5 FC out of 151 labeled).  Within vinc_only, FC reaches ~25% at 75% labels.",
         size=9, color=C_DARK, wrap=True)
    _txt(sl, SW * 0.5, by, SW * 0.5 - PAD, Inches(0.85),
         "Fib is more volatile due to extreme sparsity (≤18 patches per dataset).  "
         "Best Fib F1 is within vinc_only at 75% labels (~50%).  Completely absent from pfak cross-dataset.",
         size=9, color=C_DARK, wrap=True)


def slide_zproj_comparison(prs):
    """zproj vs zrecon side-by-side on per-class lines."""
    sl = _blank(prs)
    _title_bar(sl, "Feature Variant: z_proj vs z_recon — Per-Class F1",
               "z_recon = 12-d reconstruction latent  ·  z_proj = 8-d contrastive projection head")

    half_w = (SW - 3 * PAD) / 2

    _txt(sl, PAD, BODY_T + Inches(0.05), half_w, Inches(0.28),
         "z_recon (baseline, 12-d)", size=11, bold=True, color=C_DARK)
    _place_png(sl, EVAL_DIR / "perclass_lines_A.png",
               PAD, BODY_T + Inches(0.3), half_w, BODY_H - Inches(1.3))

    _txt(sl, PAD * 2 + half_w, BODY_T + Inches(0.05), half_w, Inches(0.28),
         "z_proj (8-d projection head)", size=11, bold=True, color=RGBColor(0xFF, 0x7F, 0x0E))
    _place_png(sl, EVAL_DIR / "perclass_lines_A_zproj.png",
               PAD * 2 + half_w, BODY_T + Inches(0.3), half_w, BODY_H - Inches(1.3))

    # Annotation row
    by = SH - Inches(1.1)
    _bullet_box(sl, PAD, by, SW - 2*PAD, Inches(1.0), [
        ("FA class is similar between z_recon and z_proj — both reach F1 ≥ 0.7 at 10% labels", C_DARK),
        ("FC improves with z_proj in vinc_only and combined scenarios — projection head sharpens "
         "intra-class boundaries learned from contrastive training", C_GREEN),
        ("NA (nascent) shows moderate improvement with z_proj in pfak_only and combined", C_DARK),
        ("Fib is still volatile — both variants limited by extreme sparsity (4–18 samples per dataset)", C_AMBER),
    ])


def slide_smote_comparison(prs):
    """SMOTE effect on hard classes FC and Fib."""
    sl = _blank(prs)
    _title_bar(sl, "SMOTE Oversampling — Effect on Hard Classes (FC and Fib)",
               "Comparing zrecon vs zrecon+SMOTE  ·  SMOTE oversamples minority to match majority count")

    half_w = (SW - 3 * PAD) / 2

    _txt(sl, PAD, BODY_T + Inches(0.05), half_w, Inches(0.28),
         "Without SMOTE (z_recon baseline)", size=11, bold=True, color=C_DARK)
    _place_png(sl, EVAL_DIR / "perclass_lines_A.png",
               PAD, BODY_T + Inches(0.3), half_w, BODY_H - Inches(1.3))

    _txt(sl, PAD * 2 + half_w, BODY_T + Inches(0.05), half_w, Inches(0.28),
         "With SMOTE (z_recon + oversampling)", size=11, bold=True,
         color=RGBColor(0xD6, 0x22, 0x28))
    _place_png(sl, EVAL_DIR / "perclass_lines_A_zrecon_smote.png",
               PAD * 2 + half_w, BODY_T + Inches(0.3), half_w, BODY_H - Inches(1.3))

    by = SH - Inches(1.1)
    _bullet_box(sl, PAD, by, SW - 2*PAD, Inches(1.0), [
        ("SMOTE provides modest FC improvement in vinc_only — the class is still hard to separate", C_AMBER),
        ("Fib F1 increases slightly within vinc_only but remains noisy due to very few real samples", C_AMBER),
        ("FA class is largely unaffected by SMOTE — already well-learned without oversampling", C_DARK),
        ("SMOTE + zproj (not shown) gives the best FC performance — see heatmap slide for comparison", C_GREEN),
    ])


def slide_heatmap_75(prs):
    """Heatmap: per-class F1 at 75% labels × scenarios × variants."""
    sl = _blank(prs)
    _title_bar(sl, "Summary Heatmap — Per-Class F1 at 75% Label Fraction",
               "All 4 FA classes  ·  5 eval scenarios  ·  4 feature variants  ·  Option A")

    classes  = ["NA", "FC", "FA", "Fib"]
    n_sc     = len(SCENARIOS)
    n_cls    = len(classes)
    n_var    = len(VARIANTS)

    fig, axes = plt.subplots(1, n_var, figsize=(15, 4.5), facecolor="white")
    fig.suptitle("Per-class F1 at 75% labels  (mean over 4 repeats)",
                 fontsize=12, fontweight="bold")

    for ax, variant in zip(axes, VARIANTS):
        data = _load_results_at_frac(variant, 0.75)
        mat = np.full((n_cls, n_sc), np.nan)
        for j, sc in enumerate(SCENARIOS):
            if data[sc] is None:
                continue
            for i, cls_s in enumerate(classes):
                mat[i, j] = data[sc].get(cls_s, np.nan)

        im = ax.imshow(mat, vmin=0, vmax=1, cmap="RdYlGn", aspect="auto")
        ax.set_xticks(range(n_sc))
        ax.set_xticklabels([SCENARIO_LABELS[s].replace("→", "→\n") for s in SCENARIOS],
                           fontsize=7, ha="center")
        ax.set_yticks(range(n_cls))
        if ax == axes[0]:
            ax.set_yticklabels(classes, fontsize=9)
        else:
            ax.set_yticklabels([])
        ax.set_title(VARIANT_LABELS[variant], fontsize=9, fontweight="bold")

        for i in range(n_cls):
            for j in range(n_sc):
                v = mat[i, j]
                if np.isnan(v):
                    ax.text(j, i, "—", ha="center", va="center", fontsize=9, color="gray")
                else:
                    txt_color = "white" if v < 0.35 or v > 0.7 else "black"
                    ax.text(j, i, f"{v:.2f}", ha="center", va="center",
                            fontsize=8, color=txt_color, fontweight="bold")

    fig.colorbar(im, ax=axes[-1], fraction=0.04, pad=0.02, label="F1 score")
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05),
               SW - 2*PAD, BODY_H - Inches(1.05))
    plt.close(fig)

    by = SH - Inches(1.1)
    _bullet_box(sl, PAD, by, SW - 2*PAD, Inches(1.0), [
        ("FA is the only class learned reliably across all scenarios and variants (F1 ≥ 0.65)", C_GREEN),
        ("FC is near zero except in vinc_only; zproj+SMOTE gives the highest FC F1 within vinc_only", C_AMBER),
        ("pfak→vinc NA is surprisingly poor — pFAK training data biases toward FA", C_RED),
        ("Fib: vinc_only sometimes reaches 0.4–0.5; all cross-dataset and pfak_only are near zero", C_DARK),
    ])


def slide_takeaways(prs):
    """Key findings and next steps."""
    sl = _blank(prs)
    _title_bar(sl, "Key Findings and Next Steps",
               "Per-class F1 analysis  ·  Option A Stage-2 SupCon AE  ·  5 scenarios  ·  4 variants")

    col_w = (SW - 3 * PAD) / 2
    col2_x = PAD * 2 + col_w

    # Left: findings
    _txt(sl, PAD, BODY_T + Inches(0.1), col_w, Inches(0.3),
         "Key Findings", size=13, bold=True, color=C_DARK)

    findings = [
        (C_GREEN, "FA (focal adhesion) is always well-learned — strong signal in the 12-d/8-d latent space"),
        (C_AMBER, "NA is learned within-dataset but degrades significantly in cross-dataset transfer"),
        (C_RED,   "FC (focal complex) is the dominant problem — near-zero F1 in most cross-dataset scenarios"),
        (C_RED,   "FC has severe label imbalance: 5/151 in pFAK, 5–85/454 in vinc depending on condition"),
        (C_DARK,  "Fib is extremely sparse (4–18 samples per dataset) — any F1 estimate is high-variance"),
        (C_AMBER, "z_proj marginally helps FC in within-dataset; combined with SMOTE gives best hard-class F1"),
        (C_DARK,  "pfak→vinc is catastrophic for FC (0%) — pFAK FC appearance differs from vinc FC"),
    ]
    fy = BODY_T + Inches(0.5)
    for color, text in findings:
        _txt(sl, PAD, fy, col_w - Inches(0.1), Inches(0.48),
             f"• {text}", size=9, color=color, wrap=True)
        fy += Inches(0.48)

    # Right: next steps
    _txt(sl, col2_x, BODY_T + Inches(0.1), col_w, Inches(0.3),
         "Recommended Next Steps", size=13, bold=True, color=C_DARK)

    steps = [
        (C_DARK,  "Verify label sampling is stratified by class (not random) — important for FC/Fib"),
        (C_DARK,  "Run cumulative label test: 10% labels are a subset of 25% set (isolate sampling vs data effect)"),
        (C_AMBER, "Add class-balanced SupCon batching in AE retraining (WeightedRandomSampler) to learn better "
                  "FC/Fib representations at the feature level"),
        (C_GREEN, "Acquire more FC and Fib annotations — especially in pFAK dataset (currently 5 FC / 4 Fib)"),
        (C_DARK,  "Margaret to label 1 ycomp image from pFAK — expand ycomp representation in cross-dataset"),
        (C_AMBER, "Consider 2-stage approach: train FA-vs-other first, then refine FA subtypes separately"),
        (C_DARK,  "2D sweep (n_images × labels/image) once sampling logic is confirmed to be stratified"),
    ]
    sy = BODY_T + Inches(0.5)
    for color, text in steps:
        _txt(sl, col2_x, sy, col_w - Inches(0.1), Inches(0.48),
             f"• {text}", size=9, color=color, wrap=True)
        sy += Inches(0.48)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default=str(OUT_DEFAULT))
    args = parser.parse_args()

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("[1/8] Cover ...")
    slide_cover(prs)
    print("[2/8] Class setup ...")
    slide_class_setup(prs)
    print("[3/8] Baseline overview ...")
    slide_baseline_overview(prs)
    print("[4/8] Hard-class deep dive ...")
    slide_hardclass_deepdive(prs)
    print("[5/8] z_proj comparison ...")
    slide_zproj_comparison(prs)
    print("[6/8] SMOTE comparison ...")
    slide_smote_comparison(prs)
    print("[7/8] Summary heatmap ...")
    slide_heatmap_75(prs)
    print("[8/8] Takeaways ...")
    slide_takeaways(prs)

    prs.save(str(out_path))
    print(f"\n[done] {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
