#!/usr/bin/env python3
"""
make_pptx_lambda_sweep.py
=========================
Generate a focused PPT comparing three ConAE lat12 lambda_contrast values:
  λ=0.0001  contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc1e4
  λ=0.5     contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1
  λ=100     contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc100

Run from the SubCellAE repo root:
  python scripts/make_pptx_lambda_sweep.py
"""
from __future__ import annotations

import io
from pathlib import Path

import numpy as np
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
OUT  = Path("lambda_sweep_results.pptx")

# ── models (ordered: small → large lambda) ─────────────────────────────────────

MODELS = [
    # (label, dirname, lambda_contrast, note)
    ("ConAE lat12  λ=0.0001  (recon-dominated)",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc1e4",
     0.0001, "Near-zero contrastive weight — essentially a plain AE"),
    ("ConAE lat12  λ=0.5  (balanced)",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1",
     0.5, "Balanced recon + contrastive loss"),
    ("ConAE lat12  λ=100  (contrast-dominated)",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc100",
     100.0, "Contrastive loss dominates — recon nearly suppressed"),
]

# ── slide geometry (16:9) ──────────────────────────────────────────────────────

SW = Inches(13.33)
SH = Inches(7.5)
TITLE_H = Inches(0.65)
PAD     = Inches(0.12)

# ── colours ────────────────────────────────────────────────────────────────────

C_DARK   = RGBColor(0x1F, 0x4E, 0x79)
C_MID    = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT  = RGBColor(0xBD, 0xD7, 0xEE)
C_ALT    = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_GREY   = RGBColor(0x88, 0x88, 0x88)

# ── low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh


def _txt(slide, l, t, w, h, text, size=14, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def _img(slide, src, l, t, max_w, max_h):
    """Add image (path or bytes), fitting within max_w x max_h, centred."""
    if src is None:
        return
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return
        data = p.read_bytes()
    else:
        data = bytes(src)
    try:
        im = Image.open(io.BytesIO(data))
        iw, ih = im.size
    except Exception:
        return
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    lo = l + (max_w - w) // 2
    to = t + (max_h - h) // 2
    slide.shapes.add_picture(io.BytesIO(data), lo, to, w, h)


def _title_bar(slide, text, size=18, bg=C_DARK):
    _rect(slide, 0, 0, SW, TITLE_H, fill=bg)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         text, size=size, bold=True, color=C_WHITE)


def _tif_to_png(path: Path) -> bytes | None:
    try:
        arr = tifffile.imread(str(path)).astype(np.float32)
    except Exception:
        return None

    if arr.ndim == 2:
        img2d = arr
    elif arr.ndim == 3 and arr.shape[2] in (3, 4):
        img2d = arr
    elif arr.ndim == 3:
        K, H, W = arr.shape
        ncols = min(K, 5)
        nrows = (K + ncols - 1) // ncols
        canvas = np.zeros((nrows * H, ncols * W), dtype=np.float32)
        for i in range(K):
            r, c = divmod(i, ncols)
            canvas[r*H:(r+1)*H, c*W:(c+1)*W] = arr[i]
        img2d = canvas
    else:
        img2d = arr[0]

    vmin, vmax = img2d.min(), img2d.max()
    if vmax > vmin:
        img2d = ((img2d - vmin) / (vmax - vmin) * 255).clip(0, 255).astype(np.uint8)
    else:
        img2d = np.zeros_like(img2d, dtype=np.uint8)

    pil = Image.fromarray(img2d)
    buf = io.BytesIO()
    pil.save(buf, format='PNG')
    return buf.getvalue()


# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_DARK)
    _txt(slide, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.8),
         title, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(2),
             subtitle, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_section(prs, title, detail=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_MID)
    _txt(slide, Inches(1), Inches(2.4), SW - Inches(2), Inches(1.6),
         title, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if detail:
        _txt(slide, Inches(1.5), Inches(4.2), SW - Inches(3), Inches(2),
             detail, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_1img(prs, title, img, caption=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top = TITLE_H + PAD
    avail_h = SH - top - (Inches(0.32) if caption else PAD)
    _img(slide, img, PAD, top, SW - 2*PAD, avail_h)
    if caption:
        _txt(slide, PAD, SH - Inches(0.32), SW - 2*PAD, Inches(0.28),
             caption, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_2img(prs, title, img1, img2, cap1="", cap2=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top = TITLE_H + PAD
    cap_h = Inches(0.28) if (cap1 or cap2) else PAD
    h = SH - top - cap_h - PAD
    w = (SW - 3*PAD) / 2
    _img(slide, img1, PAD,           top, w, h)
    _img(slide, img2, PAD + w + PAD, top, w, h)
    if cap1:
        _txt(slide, PAD, SH - cap_h, w, cap_h,
             cap1, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
    if cap2:
        _txt(slide, PAD + w + PAD, SH - cap_h, w, cap_h,
             cap2, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_3img(prs, title, imgs, caps=None):
    slide = _blank(prs)
    _title_bar(slide, title)
    caps = (list(caps) + [""]*3)[:3] if caps else [""]*3
    imgs = (list(imgs) + [None]*3)[:3]
    top  = TITLE_H + PAD
    cap_h = Inches(0.28) if any(caps) else PAD
    h = SH - top - cap_h - PAD
    w = (SW - 4*PAD) / 3
    for i, (im, cap) in enumerate(zip(imgs, caps)):
        l = PAD + i * (w + PAD)
        if im:
            _img(slide, im, l, top, w, h)
        else:
            _txt(slide, l, top + h/2 - Inches(0.25), w, Inches(0.5),
                 "[not available]", size=10, color=C_GREY, align=PP_ALIGN.CENTER)
        if cap:
            _txt(slide, l, SH - cap_h, w, cap_h,
                 cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def _style_cell(cell, text, size=8, bold=False, align=PP_ALIGN.CENTER,
                bg=None, fg=C_BLACK):
    cell.text = str(text)
    tf = cell.text_frame
    tf.paragraphs[0].alignment = align
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def slide_table(prs, title, headers, rows, font_size=8, title_bg=C_DARK):
    slide = _blank(prs)
    _title_bar(slide, title, bg=title_bg)
    if not rows:
        _txt(slide, PAD, TITLE_H + PAD, SW - 2*PAD, Inches(1),
             "(no data)", size=12, color=C_GREY, align=PP_ALIGN.CENTER)
        return
    nr = len(rows) + 1
    nc = len(headers)
    row_h = min(Inches(0.33), (SH - TITLE_H - 2*PAD) / nr)
    tbl_h = row_h * nr
    tbl_w = SW - 2*PAD
    top   = TITLE_H + PAD + max(0, (SH - TITLE_H - 2*PAD - tbl_h) / 2)
    tbl = slide.shapes.add_table(nr, nc, PAD, top, tbl_w, tbl_h).table
    for j, h in enumerate(headers):
        _style_cell(tbl.cell(0, j), h, size=font_size, bold=True, bg=C_DARK, fg=C_WHITE)
    for i, row in enumerate(rows):
        bg = C_ALT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            _style_cell(tbl.cell(i+1, j), str(val) if val is not None else "",
                        size=font_size, bg=bg)


# ── per-model slides ───────────────────────────────────────────────────────────

def _model_slides(prs, label: str, d: Path):
    if not d.exists():
        slide_section(prs, f"{label} -- NOT FOUND", f"Directory missing: {d.name}")
        return

    ev = d / "eval"
    prefix = "contrastive"

    # 1. Training loss + component losses
    loss = d / f"{prefix}_train_val_loss.png"
    comp = d / f"{prefix}_component_losses.png"
    if loss.exists() and comp.exists():
        slide_2img(prs, f"{label} -- Training Loss",
                   loss, comp, "Train/Val Loss", "Component Losses")
    elif loss.exists():
        slide_1img(prs, f"{label} -- Training Loss", loss)

    # 2. ep500 views + reconstruction
    views = d / f"{prefix}_views_ep500.png"
    recon = d / f"{prefix}_recon_ep500.png"
    if views.exists() and recon.exists():
        slide_2img(prs, f"{label} -- Epoch 500",
                   views, recon, "Contrastive views", "Reconstruction")
    elif recon.exists():
        slide_1img(prs, f"{label} -- Epoch 500 Recon", recon)

    # 3. UMAP z_proj 4ds
    p4_ann  = ev / "umap_proj_4ds_annotation.png"
    p4_cond = ev / "umap_proj_4ds_condition.png"
    p4_ds   = ev / "umap_proj_4ds_dataset.png"
    p4_km   = ev / "umap_proj_4ds_kmeans.png"
    if p4_ann.exists():
        slide_3img(prs, f"{label} -- UMAP z_proj (cross-dataset)",
                   [p4_ann,
                    p4_cond if p4_cond.exists() else None,
                    p4_ds   if p4_ds.exists()   else None],
                   ["Annotation (labels)", "Condition", "Dataset"])
    if p4_km.exists():
        slide_2img(prs, f"{label} -- UMAP z_proj: K-means & training labels",
                   p4_km,
                   p4_cond if p4_cond.exists() else None,
                   "K-means cluster ID", "Condition (training labels)")

    # 4. PHATE z_proj 4ds
    pp_ann  = ev / "phate_proj_4ds_annotation.png"
    pp_cond = ev / "phate_proj_4ds_condition.png"
    if pp_ann.exists():
        slide_2img(prs, f"{label} -- PHATE z_proj (cross-dataset)",
                   pp_ann,
                   pp_cond if pp_cond.exists() else None,
                   "Annotation (labels)", "Condition")

    # 5. Cluster panels (all_clusters.tif → tiled PNG)
    cp_dir  = ev / "cluster_panels"
    all_tif = cp_dir / "all_clusters.tif"
    if not all_tif.exists():
        all_tif = next(iter(sorted(cp_dir.glob("all*.tif"))), None) if cp_dir.is_dir() else None
    if all_tif:
        png = _tif_to_png(all_tif)
        if png:
            slide_1img(prs, f"{label} -- K-means Cluster Panels", png, all_tif.name)


# ── comparison slides (one per figure type, all 3 models side-by-side) ──────────

def _comparison_slides(prs):
    """Side-by-side comparison slides: one figure type, three lambda values."""
    labels = [lbl for lbl, _, _, _ in MODELS]
    dirs   = [RUNS / dn for _, dn, _, _ in MODELS]
    lambdas = [lc for _, _, lc, _ in MODELS]
    short_caps = [f"λ={lc}" for lc in lambdas]

    def _get(d, *rel):
        p = d.joinpath(*rel)
        return p if p.exists() else None

    # Training loss
    slide_3img(prs, "Comparison — Train/Val Loss",
               [_get(d, "contrastive_train_val_loss.png") for d in dirs],
               short_caps)

    # Component losses
    slide_3img(prs, "Comparison — Component Losses",
               [_get(d, "contrastive_component_losses.png") for d in dirs],
               short_caps)

    # Epoch 500 recon
    slide_3img(prs, "Comparison — Epoch 500 Reconstruction",
               [_get(d, "contrastive_recon_ep500.png") for d in dirs],
               short_caps)

    # UMAP z_proj annotation
    slide_3img(prs, "Comparison — UMAP z_proj (FA annotation)",
               [_get(d, "eval", "umap_proj_4ds_annotation.png") for d in dirs],
               short_caps)

    # UMAP z_proj condition
    slide_3img(prs, "Comparison — UMAP z_proj (condition)",
               [_get(d, "eval", "umap_proj_4ds_condition.png") for d in dirs],
               short_caps)

    # UMAP z_proj K-means
    slide_3img(prs, "Comparison — UMAP z_proj (K-means clusters)",
               [_get(d, "eval", "umap_proj_4ds_kmeans.png") for d in dirs],
               short_caps)

    # PHATE z_proj annotation
    slide_3img(prs, "Comparison — PHATE z_proj (FA annotation)",
               [_get(d, "eval", "phate_proj_4ds_annotation.png") for d in dirs],
               short_caps)

    # PHATE z_proj condition
    slide_3img(prs, "Comparison — PHATE z_proj (condition)",
               [_get(d, "eval", "phate_proj_4ds_condition.png") for d in dirs],
               short_caps)

    # Cluster panels
    cluster_pngs = []
    for d in dirs:
        cp  = d / "eval" / "cluster_panels" / "all_clusters.tif"
        if not cp.exists():
            cp = next(iter(sorted((d / "eval" / "cluster_panels").glob("all*.tif"))), None) \
                 if (d / "eval" / "cluster_panels").is_dir() else None
        cluster_pngs.append(_tif_to_png(cp) if cp else None)
    slide_3img(prs, "Comparison — K-means Cluster Panels",
               cluster_pngs, short_caps)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs,
        "ConAE λ_contrast Sweep",
        "ConAE lat12  CIO normalisation  vinc training set  nL1 recon loss  500 epochs\n"
        "λ=0.0001  vs  λ=0.5  vs  λ=100\n\n"
        "2026-07-27")

    # Config overview table
    hdrs = ["Model", "λ_contrast", "Description", "Result dir"]
    rows = [[lbl, str(lc), note, dn]
            for lbl, dn, lc, note in MODELS]
    slide_table(prs, "Model Configurations", hdrs, rows, font_size=10)

    # Side-by-side comparison slides first
    print("Building comparison slides ...")
    slide_section(prs, "Direct Comparison",
                  "All three λ values shown side-by-side for each figure type")
    _comparison_slides(prs)

    # Then per-model deep-dive
    for lbl, dname, lc, note in MODELS:
        d = RUNS / dname
        exists_note = "[OK]" if d.exists() else "[MISSING]"
        print(f"  {lbl} ({exists_note}) ...")
        slide_section(prs, lbl, f"{dname}  {exists_note}\n{note}")
        _model_slides(prs, lbl, d)

    prs.save(str(OUT))
    n = len(prs.slides)
    print(f"\nSaved: {OUT}  ({n} slides)")
    print(f"Rsync to laptop:  rsync -avh --progress {OUT} lding@lding-Precision-3680:/home/lding/Desktop/")


if __name__ == "__main__":
    main()
