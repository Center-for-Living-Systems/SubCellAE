#!/usr/bin/env python3
"""
make_ds_combo_sweep_ppt.py

PPT comparing three training-loss variants of the v2 ConAE dataset-combo sweep:
  nL1  (ds_combo_enlcrop_sc2_lc010_bal)
  MSE  (ds_combo_enlcrop_sc2_lc010_bal_mse)
  L1   (ds_combo_enlcrop_sc2_lc010_bal_l1)

Slides
------
  1.  Title
  2.  Experiment overview
  3.  Mean nL1 table — nL1 training
  4.  Mean nL1 table — MSE training
  5.  Mean nL1 table — L1  training
  [Section] nL1 training loss
  6–10.   Violin plots (Singles / Pairs A / Pairs B / Triples / All-4)
           each slide: 3 rows (nL1 eval / MSE eval / L1 eval) × N combos
  11–15.  UMAP — condition (nL1 training only)
  [Section] MSE training loss
  16–20.  Violin plots
  [Section] L1 training loss
  21–25.  Violin plots

Usage
-----
  python scripts/make_ds_combo_sweep_ppt.py [--out slides_ds_combo_sweep.pptx]
"""

from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths & sweep definitions ─────────────────────────────────────────────────

_RUNS_BASE = Path("/net/projects/CLS/lding/data/fa_data_analysis"
                  "/ae_results/contrastive_run")

COMBO_LIST_TXT = Path("config/contrastive_config/ds_combo_v2/combo_list.txt")

# (key, display_label, runs_dir, accent_color)
SWEEPS = [
    ("nl1", "nL1",
     _RUNS_BASE / "ds_combo_enlcrop_sc2_lc010_bal",
     RGBColor(0x2E, 0x86, 0xC1)),   # blue
    ("mse", "MSE",
     _RUNS_BASE / "ds_combo_enlcrop_sc2_lc010_bal_mse",
     RGBColor(0xC0, 0x50, 0x20)),   # orange
    ("l1",  "L1",
     _RUNS_BASE / "ds_combo_enlcrop_sc2_lc010_bal_l1",
     RGBColor(0x1A, 0x7A, 0x40)),   # green
]

# ── combo groups ──────────────────────────────────────────────────────────────

SINGLES  = ["vinc", "nih3t3", "ppax", "pfak"]
PAIRS_A  = ["vinc_nih3t3", "vinc_ppax", "vinc_pfak"]
PAIRS_B  = ["nih3t3_ppax", "nih3t3_pfak", "ppax_pfak"]
TRIPLES  = ["vinc_nih3t3_ppax", "vinc_nih3t3_pfak",
            "vinc_ppax_pfak",   "nih3t3_ppax_pfak"]
ALL_DS   = ["vinc_nih3t3_ppax_pfak"]

GROUPS_VIOLIN = [SINGLES, PAIRS_A, PAIRS_B, TRIPLES, ALL_DS]

DATASETS_ORDERED = ["vinc", "pfak", "ppax", "nih3t3"]
DS_NAME  = {"vinc": "ds1", "pfak": "ds2", "ppax": "ds3", "nih3t3": "ds4"}
DS_LABEL = {k: v for k, v in DS_NAME.items()}

REPEAT = {"vinc": 1, "nih3t3": 2, "ppax": 2, "pfak": 3}

VIOLIN_METRICS = [
    ("recon_nl1", "nL1 eval"),
    ("recon_mse", "MSE eval"),
    ("recon_l1",  "L1 eval"),
]


def _combo_label(combo: str) -> str:
    return "+".join(DS_NAME.get(p, p) for p in combo.split("_"))


# ── slide geometry & colours ──────────────────────────────────────────────────

W_IN, H_IN = 13.33, 7.5
MARGIN   = 0.25
HEADER_H = 0.68

C_TITLE  = RGBColor(0x1F, 0x2D, 0x3D)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_LGRAY  = RGBColor(0xAA, 0xAA, 0xAA)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ── pptx helpers ──────────────────────────────────────────────────────────────

def _px(v): return Inches(v)

def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _textbox(slide, text, left, top, width, height,
             font_size=11, bold=False, color=C_TITLE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def _slide_header(slide, title, subtitle=None, color=C_TITLE):
    bar = slide.shapes.add_shape(1, _px(0), _px(0), _px(W_IN), _px(0.62))
    bar.fill.background(); bar.line.fill.background()
    _textbox(slide, title, 0.25, 0.05, W_IN - 0.5, 0.45,
             font_size=22, bold=True, color=color)
    if subtitle:
        _textbox(slide, subtitle, 0.25, 0.47, W_IN - 0.5, 0.22,
                 font_size=11, color=C_GRAY)

def _add_image_bytes(slide, img_bytes, left, top, width=None, height=None):
    buf = io.BytesIO(img_bytes)
    return slide.shapes.add_picture(
        buf, _px(left), _px(top),
        width=_px(width)  if width  is not None else None,
        height=_px(height) if height is not None else None,
    )

def _fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()

def _load_png(path: Path) -> np.ndarray | None:
    if path.exists():
        return np.array(Image.open(str(path)).convert("RGB"))
    return None


# ── composite figure builders ─────────────────────────────────────────────────

def _violin_3metrics_grid(combos: list[str], runs_dir: Path,
                           n_cols: int | None = None) -> bytes:
    """3-row grid: nL1 eval / MSE eval / L1 eval, one column per combo."""
    n = len(combos)
    if n_cols is None:
        n_cols = min(n, 4)
    combo_rows = (n + n_cols - 1) // n_cols
    n_mrows = len(VIOLIN_METRICS)

    fig, axes = plt.subplots(
        n_mrows * combo_rows, n_cols,
        figsize=(n_cols * 5.5, n_mrows * combo_rows * 3.8),
        facecolor="white",
    )
    axes = np.array(axes).reshape(n_mrows * combo_rows, n_cols)

    for mi, (metric, metric_label) in enumerate(VIOLIN_METRICS):
        for ci, combo in enumerate(combos):
            crow = ci // n_cols; ccol = ci % n_cols
            ax = axes[mi * combo_rows + crow, ccol]
            img = _load_png(runs_dir / combo / f"cross_dataset_{metric}.png")
            if img is not None:
                ax.imshow(img)
            else:
                ax.text(0.5, 0.5, "not yet generated",
                        ha="center", va="center", transform=ax.transAxes,
                        fontsize=10, color="gray")
            title = _combo_label(combo)
            if ccol == 0:
                title = f"[{metric_label}]  {title}"
            ax.set_title(title, fontsize=10, fontweight="bold", pad=3)
            ax.axis("off")
        for ci in range(len(combos), combo_rows * n_cols):
            axes[mi * combo_rows + ci // n_cols, ci % n_cols].axis("off")

    fig.tight_layout(pad=0.6)
    return _fig_to_bytes(fig)


def _umap_grid(combos: list[str], runs_dir: Path,
               n_cols: int | None = None, with_test: bool = False) -> bytes:
    n = len(combos)
    if n_cols is None:
        n_cols = min(n, 4)
    n_rows = (n + n_cols - 1) // n_cols
    fname = "umap_combo_condition_with_test.png" if with_test else "umap_combo_condition.png"

    fig, axes = plt.subplots(n_rows, n_cols,
                             figsize=(n_cols * 4.0, n_rows * 3.8),
                             facecolor="white")
    axes = np.array(axes).reshape(n_rows, n_cols)

    for i, combo in enumerate(combos):
        r, c = i // n_cols, i % n_cols
        img = _load_png(runs_dir / combo / "eval" / fname)
        ax = axes[r, c]
        if img is not None:
            ax.imshow(img)
        else:
            ax.text(0.5, 0.5, "not yet generated",
                    ha="center", va="center", transform=ax.transAxes,
                    fontsize=10, color="gray")
        ax.set_title(_combo_label(combo), fontsize=11, fontweight="bold", pad=4)
        ax.axis("off")

    for i in range(len(combos), n_rows * n_cols):
        axes[i // n_cols, i % n_cols].axis("off")

    fig.tight_layout(pad=0.8)
    return _fig_to_bytes(fig)


# ── slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _textbox(slide, "Dataset Combination ConAE Sweep  —  v2",
             1.0, 1.3, W_IN - 2.0, 1.1,
             font_size=38, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    _textbox(slide,
             "Three training losses compared:  nL1  ·  MSE  ·  L1 (MAE)",
             1.0, 2.55, W_IN - 2.0, 0.5,
             font_size=20, color=C_GRAY, align=PP_ALIGN.CENTER)
    _textbox(slide,
             "enlcrop · sc2 · λ_contrast = 0.10  |  15 combinations of {ds1, ds2, ds3, ds4}",
             1.0, 3.12, W_IN - 2.0, 0.45,
             font_size=15, color=C_LGRAY, align=PP_ALIGN.CENTER)
    items = [
        "ds1 = vinc  ·  ds2 = pfak  ·  ds3 = ppax  ·  ds4 = nih3t3",
        "Balanced sampling: ds1×1 (40% train), ds4×2, ds3×2, ds2×3  (~8–12k patches/dataset)",
        "Each violin slide: 3 eval metrics (nL1 / MSE / L1) per training-loss variant",
    ]
    y = 3.75
    for item in items:
        _textbox(slide, f"• {item}", 1.5, y, W_IN - 3.0, 0.42,
                 font_size=14, color=C_GRAY)
        y += 0.45
    _textbox(slide, "Latent 12  ·  Proj 8  ·  500 epochs  ·  Adam cosine LR",
             1.0, 6.5, W_IN - 2.0, 0.4,
             font_size=12, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)


def _slide_overview(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _slide_header(slide,
                  "Experiment Overview — 15 Dataset Combinations  (v2)",
                  subtitle="Balanced split · enlcrop/sc2 augmentation · "
                           "Three training losses (nL1 / MSE / L1) · λ_contrast = 0.10")

    combos = SINGLES + PAIRS_A + PAIRS_B + TRIPLES + ALL_DS
    col_x = [MARGIN, 1.6, 3.5, 5.4, 7.3]
    col_w = [1.3,    1.8, 1.8, 1.8, 1.8]
    headers = ["#", "Combo", "ds1", "ds3/ds2", "ds4"]
    row_h = 0.38; y0 = HEADER_H

    for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        box = slide.shapes.add_shape(1, _px(cx), _px(y0), _px(cw), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, hdr, cx + 0.04, y0 + 0.06, cw - 0.08, row_h - 0.08,
                 font_size=9, bold=True, color=C_TITLE)

    for i, combo in enumerate(combos):
        y = y0 + row_h * (i + 1)
        parts = set(combo.split("_"))
        row_vals = [
            str(i + 1),
            _combo_label(combo),
            f"×{REPEAT['vinc']}" if "vinc" in parts else "—",
            "+".join(f"×{REPEAT[d]}" for d in ["ppax", "pfak"] if d in parts) or "—",
            f"×{REPEAT['nih3t3']}" if "nih3t3" in parts else "—",
        ]
        for ci, (val, cw, cx) in enumerate(zip(row_vals, col_w, col_x)):
            box = slide.shapes.add_shape(1, _px(cx), _px(y), _px(cw), _px(row_h))
            box.fill.background(); box.line.color.rgb = C_LGRAY
            color = SWEEPS[0][3] if ci == 1 else C_GRAY
            _textbox(slide, val, cx + 0.04, y + 0.05, cw - 0.08, row_h - 0.08,
                     font_size=8, color=color)

    rx = 9.3
    _textbox(slide, "Training Settings", rx, HEADER_H, 3.8, 0.38,
             font_size=13, bold=True, color=C_TITLE)
    params = [
        ("Architecture", "ConAE — contrastive autoencoder"),
        ("Recon loss",   "nL1  |  MSE  |  L1 (three parallel sweeps)"),
        ("Contrastive",  "SC2  ·  λ = 0.10"),
        ("Augmentation", "EnlargedJitterCrop (±15°, ±4px) + sc2 jitter"),
        ("Latent / Proj","12 / 8"),
        ("Epochs",       "500"),
        ("Optimizer",    "Adam, cosine LR decay"),
        ("Patch size",   "32×32 px  (context 58×58)"),
        ("Balanced",     "ds1×1 (40%tr), ds4×2, ds3×2, ds2×3"),
    ]
    y = HEADER_H + 0.42
    for label, val in params:
        _textbox(slide, label + ":", rx, y, 1.7, 0.38,
                 font_size=9, bold=True, color=C_TITLE)
        _textbox(slide, val, rx + 1.75, y, 2.1, 0.42,
                 font_size=9, color=C_GRAY)
        y += 0.42
    _textbox(slide,
             "Eval metrics: nL1 = L1/mean|raw|, MSE, L1 = MAE  (all three shown)",
             rx, y + 0.1, 3.8, 0.55,
             font_size=9, italic=True, color=SWEEPS[0][3])


def _slide_mean_table(prs: Presentation, sweep_key: str, sweep_label: str,
                      runs_dir: Path, accent: RGBColor) -> None:
    csv_path = runs_dir / "cross_dataset_recon_metrics.csv"
    if not csv_path.exists():
        print(f"  [SKIP] mean table ({sweep_label}) — CSV not found")
        return

    df = pd.read_csv(csv_path)
    if "variant" not in df.columns or "recon_nl1" not in df.columns:
        print(f"  [SKIP] mean table ({sweep_label}) — unexpected columns")
        return

    combos = SINGLES + PAIRS_A + PAIRS_B + TRIPLES + ALL_DS
    slide = _add_slide(prs)
    _slide_header(slide,
                  f"Mean Metrics — {sweep_label} Training Loss",
                  subtitle="Lower = better  ·  bold = training dataset  ·  "
                            "non-training = unseen test",
                  color=accent)

    # show nL1 eval (most interpretable across all training losses)
    mean_df = (df.groupby(["variant", "dataset"])["recon_nl1"]
               .mean().reset_index())
    table = {}
    for _, row in mean_df.iterrows():
        table.setdefault(row["variant"], {})[row["dataset"]] = row["recon_nl1"]

    datasets = DATASETS_ORDERED
    col_labels = ["Combo"] + [DS_LABEL[d] for d in datasets]
    col_w = [3.2] + [2.1] * len(datasets)
    col_x: list[float] = []
    x = MARGIN
    for cw in col_w:
        col_x.append(x); x += cw

    row_h = 0.35; y0 = HEADER_H

    for ci, (hdr, cw, cx) in enumerate(zip(col_labels, col_w, col_x)):
        box = slide.shapes.add_shape(1, _px(cx), _px(y0), _px(cw), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, hdr, cx + 0.04, y0 + 0.04, cw - 0.08, row_h - 0.08,
                 font_size=9, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)

    all_vals = [v for c in combos for v in table.get(c, {}).values()]
    vmin = min(all_vals) if all_vals else 0.0
    vmax = max(all_vals) if all_vals else 1.0

    def _heat(val):
        t = (val - vmin) / max(vmax - vmin, 1e-6)
        return RGBColor(
            max(0, min(255, int(255 * (0.2 + 0.6 * t)))),
            max(0, min(255, int(255 * (0.9 - 0.5 * t)))),
            max(0, min(255, int(255 * (0.9 - 0.7 * t)))),
        )

    for i, combo in enumerate(combos):
        y = y0 + row_h * (i + 1)
        row_data = table.get(combo, {})
        parts = set(combo.split("_"))

        box = slide.shapes.add_shape(1, _px(col_x[0]), _px(y), _px(col_w[0]), _px(row_h))
        box.fill.background(); box.line.color.rgb = C_LGRAY
        _textbox(slide, _combo_label(combo),
                 col_x[0] + 0.04, y + 0.04, col_w[0] - 0.08, row_h - 0.06,
                 font_size=9, color=C_TITLE)

        for ci, ds in enumerate(datasets, start=1):
            val = row_data.get(ds)
            cx = col_x[ci]; cw = col_w[ci]
            box = slide.shapes.add_shape(1, _px(cx), _px(y), _px(cw), _px(row_h))
            box.fill.background(); box.line.color.rgb = C_LGRAY
            txt = f"{val:.3f}" if val is not None else "—"
            is_train = ds in parts
            _textbox(slide, txt, cx + 0.04, y + 0.04, cw - 0.08, row_h - 0.06,
                     font_size=9, bold=is_train, color=C_TITLE,
                     align=PP_ALIGN.CENTER)

    y_note = y0 + row_h * (len(combos) + 1) + 0.08
    _textbox(slide,
             "Metric shown: mean normalised L1 (nL1) eval  ·  bold = dataset was in training set",
             MARGIN, y_note, W_IN - 2 * MARGIN, 0.3,
             font_size=9, italic=True, color=C_GRAY)


def _slide_section(prs: Presentation, sweep_label: str,
                   accent: RGBColor) -> None:
    slide = _add_slide(prs)
    bar = slide.shapes.add_shape(1, _px(0), _px(2.8), _px(W_IN), _px(1.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    _textbox(slide, f"Training Loss:  {sweep_label}",
             0.5, 3.05, W_IN - 1.0, 1.1,
             font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide,
             "Violin plots show all three eval metrics (nL1 / MSE / L1) per combo group",
             0.5, 4.25, W_IN - 1.0, 0.5,
             font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER)


def _slide_violin_group(prs: Presentation, combos: list[str],
                        group_label: str, sweep_label: str,
                        runs_dir: Path, accent: RGBColor) -> None:
    slide = _add_slide(prs)
    n_cols = min(len(combos), 4)
    _slide_header(slide,
                  f"Violin — {sweep_label} training  ·  {group_label}",
                  subtitle="Rows: nL1 eval / MSE eval / L1 eval  ·  "
                            "training datasets: left of divider  ·  test: right",
                  color=accent)
    img = _violin_3metrics_grid(combos, runs_dir, n_cols=n_cols)
    _add_image_bytes(slide, img, MARGIN, HEADER_H + 0.05,
                     width=W_IN - 2 * MARGIN)


def _slide_umap_group(prs: Presentation, combos: list[str],
                      group_label: str, runs_dir: Path,
                      with_test: bool = False) -> None:
    slide = _add_slide(prs)
    n_cols = min(len(combos), 4)
    title = f"UMAP + Test — {group_label}" if with_test else f"UMAP — {group_label}"
    subtitle = ("z_proj  ·  filled = train  ·  hollow = test"
                if with_test else
                "z_proj latent space  ·  training conditions  ·  tab20 colors")
    _slide_header(slide, title, subtitle=subtitle, color=C_GREEN)
    img = _umap_grid(combos, runs_dir, n_cols=n_cols, with_test=with_test)
    _add_image_bytes(slide, img, MARGIN, HEADER_H + 0.05,
                     width=W_IN - 2 * MARGIN)


# ── main ──────────────────────────────────────────────────────────────────────

GROUP_LABELS = [
    "Singles  (ds1 · ds4 · ds3 · ds2)",
    "Pairs A  (ds1+ds4 · ds1+ds3 · ds1+ds2)",
    "Pairs B  (ds4+ds3 · ds4+ds2 · ds3+ds2)",
    "Triples  (all 4 triples)",
    "All 4 datasets",
]


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="slides_ds_combo_sweep.pptx")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    # ── title + overview ──────────────────────────────────────────────────────
    print("  Title …"); _slide_title(prs)
    print("  Overview …"); _slide_overview(prs)

    # ── mean tables (one per sweep) ───────────────────────────────────────────
    for key, label, runs_dir, accent in SWEEPS:
        print(f"  Mean table [{label}] …")
        _slide_mean_table(prs, key, label, runs_dir, accent)

    # ── per-sweep violin slides + UMAP for nL1 ───────────────────────────────
    nl1_runs = SWEEPS[0][2]
    for key, label, runs_dir, accent in SWEEPS:
        print(f"  [Section] {label} …")
        _slide_section(prs, label, accent)
        for combos, glbl in zip(GROUPS_VIOLIN, GROUP_LABELS):
            print(f"    Violin [{label}] {glbl} …")
            _slide_violin_group(prs, combos, glbl, label, runs_dir, accent)

        # UMAP only for nL1 (latent structure; all three have same architecture)
        if key == "nl1":
            for combos, glbl in zip(GROUPS_VIOLIN, GROUP_LABELS):
                print(f"    UMAP {glbl} …")
                _slide_umap_group(prs, combos, glbl, nl1_runs)

    out = Path(args.out)
    prs.save(str(out))
    print(f"\nSaved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
