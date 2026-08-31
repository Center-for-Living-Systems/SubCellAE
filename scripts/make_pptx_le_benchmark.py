#!/usr/bin/env python3
"""
make_pptx_le_benchmark.py

Slide deck for the label-efficiency benchmark comparing SupCon-AE vs
CellProfiler vs ilastik on DS1 B2 and DS1 B12 combined labels.

Each slide has:
  - Left/centre: balanced accuracy vs label budget line plot
  - Right panel: dataset / label / training / testing / model info

Usage:
  python scripts/make_pptx_le_benchmark.py
  python scripts/make_pptx_le_benchmark.py --out results/le_benchmark.pptx
"""
from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parents[1]
EVAL_DIR  = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/features/eval_results")
OUT_DIR   = REPO_ROOT / "results"

DATA_ROOT  = Path("/net/projects/CLS/lding/data/fa_data_analysis")
CP_CSV     = DATA_ROOT / "ae_results/features/cellprofiler/ds1.csv"
IL_CSV     = DATA_ROOT / "ae_results/features/ilastik/ds1.csv"
ANN_B2_CSV = DATA_ROOT / "labelling/vinc_combined_label_Annabel_20260816.csv"
FOLD_SPLITS_B2 = DATA_ROOT / "labelling/le_b2_supcon/fold_splits_ds1.csv"
PATCH_DIR  = DATA_ROOT / "ae_results/patches/cio/vinc"
LE_B2_DIR  = DATA_ROOT / "ae_results/contrastive_run/le_b2_lat12p8"

# ---------------------------------------------------------------------------
# Slide geometry
# ---------------------------------------------------------------------------
SW = Inches(13.33)
SH = Inches(7.5)

PLOT_LEFT   = Inches(0.35)
PLOT_TOP    = Inches(1.05)
PLOT_W      = Inches(8.8)
PLOT_H      = Inches(6.1)

PANEL_LEFT  = Inches(9.3)
PANEL_TOP   = Inches(1.05)
PANEL_W     = Inches(3.8)
PANEL_H     = Inches(6.1)

# ---------------------------------------------------------------------------
# Colours
# ---------------------------------------------------------------------------
C_HEAD   = RGBColor(0x16, 0x21, 0x3E)
C_BODY   = RGBColor(0x1A, 0x1A, 0x1A)
C_GREY   = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL  = RGBColor(0xF0, 0xF4, 0xFF)
C_ACCENT = RGBColor(0x0F, 0x3D, 0x79)
C_CP     = RGBColor(0x1F, 0x77, 0xB4)
C_IL     = RGBColor(0xD6, 0x5F, 0x0E)
C_ROW_A  = RGBColor(0xF5, 0xF7, 0xFF)
C_ROW_B  = RGBColor(0xE8, 0xED, 0xFF)

COL_SUPCON  = "#e377c2"   # pink
COL_CP      = "#1f77b4"   # blue
COL_ILASTIK = "#ff7f0e"   # orange

# Histogram class colors (feature distributions)
HIST_COL_AD   = "#555555"   # adhesion: dark grey
HIST_COL_NOAD = "#7b52ab"   # no adhesion: purple

CP_HIST  = OUT_DIR / "cp_feature_histograms.png"
IL_HIST  = OUT_DIR / "ilastik_feature_histograms.png"

# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _txt(slide, text, left, top, width, height, *,
         bold=False, italic=False, size_pt=12, color=C_BODY,
         align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.size      = Pt(size_pt)
    run.font.color.rgb = color
    return txb


def _rule(slide, top, width=None, left=None, thickness_pt=0.75):
    w = width or SW - Inches(1.0)
    l = left  or Inches(0.5)
    ln = slide.shapes.add_connector(1, l, top, l + w, top)
    ln.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.width = Pt(thickness_pt)


def _slide_header(slide, title, subtitle=""):
    _txt(slide, title,
         Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.55),
         bold=True, size_pt=20, color=C_HEAD)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.35),
             size_pt=11, color=C_GREY)
    _rule(slide, Inches(0.97))


def _fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_fig(slide, fig, left, top, width, height):
    buf = _fig_to_buf(fig)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)


def _img(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return
    kw = {}
    if width:  kw["width"]  = width
    if height: kw["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kw)


def _row_rect(slide, left, top, width, height, even):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ROW_A if even else C_ROW_B
    shape.line.fill.background()


def _bullet_box(slide, left, top, width, height, header, bullets,
                header_color=None, bullet_size=11):
    header_color = header_color or C_ACCENT
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PANEL
    shape.line.color.rgb = RGBColor(0xBB, 0xC8, 0xEE)
    shape.line.width = Pt(0.75)
    _txt(slide, header,
         left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.3),
         bold=True, size_pt=11, color=header_color)
    y = top + Inches(0.42)
    per = (height - Inches(0.5)) / max(len(bullets), 1)
    for b in bullets:
        _txt(slide, f"• {b}",
             left + Inches(0.18), y, width - Inches(0.3), per + Inches(0.05),
             size_pt=bullet_size, color=C_BODY)
        y += per


def _panel_bg(slide, left, top, width, height):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PANEL
    shape.line.color.rgb = RGBColor(0xCC, 0xD4, 0xEE)
    shape.line.width = Pt(0.75)


def _panel_section(slide, label, lines, left, top, width, y_offset):
    """Draw a labelled section in the right panel. Returns new y_offset."""
    _txt(slide, label,
         left + Inches(0.12), top + y_offset, width - Inches(0.24), Inches(0.28),
         bold=True, size_pt=10, color=C_HEAD)
    y_offset += Inches(0.27)
    for line in lines:
        _txt(slide, f"  {line}",
             left + Inches(0.12), top + y_offset, width - Inches(0.24), Inches(0.22),
             size_pt=9, color=C_BODY)
        y_offset += Inches(0.21)
    y_offset += Inches(0.12)
    _rule(slide, top + y_offset,
          width=width - Inches(0.24), left=left + Inches(0.12), thickness_pt=0.5)
    y_offset += Inches(0.10)
    return y_offset


def _add_info_panel(slide, cfg: dict):
    """Draw the right-side info panel from a config dict."""
    _panel_bg(slide, PANEL_LEFT, PANEL_TOP, PANEL_W, PANEL_H)

    _txt(slide, "Experiment Details",
         PANEL_LEFT + Inches(0.12), PANEL_TOP + Inches(0.1),
         PANEL_W - Inches(0.24), Inches(0.3),
         bold=True, size_pt=12, color=C_HEAD, align=PP_ALIGN.CENTER)

    y = Inches(0.48)
    for section_label, lines in cfg.items():
        y = _panel_section(slide, section_label, lines,
                           PANEL_LEFT, PANEL_TOP, PANEL_W, y)


# ---------------------------------------------------------------------------
# Design slide helper
# ---------------------------------------------------------------------------

def _design_box(slide, label, lines, left, top, width, height):
    """Labelled info box with a light background."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xFF)
    shape.line.color.rgb      = RGBColor(0xBB, 0xC8, 0xEE)
    shape.line.width = Pt(0.75)

    _txt(slide, label,
         left + Inches(0.1), top + Inches(0.08),
         width - Inches(0.2), Inches(0.28),
         bold=True, size_pt=10, color=C_HEAD)

    y = Inches(0.34)
    for line in lines:
        _txt(slide, f"• {line}",
             left + Inches(0.12), top + y,
             width - Inches(0.24), Inches(0.22),
             size_pt=9, color=C_BODY)
        y += Inches(0.21)


def _slide_design(prs, title: str, subtitle: str,
                  sections: dict,
                  goal: str,
                  hypothesis: str = ""):
    """
    Full-slide experiment design slide.

    Layout:
      Left  (0.35 → 8.3 in): grid of info boxes (2 cols × rows)
      Right (8.6 → 13.0 in): Goal + Hypothesis highlighted panel
    """
    slide = _blank(prs)
    _slide_header(slide, title, subtitle)

    # ---- Left grid of design boxes ----
    BOX_LEFT   = Inches(0.35)
    BOX_TOP    = Inches(1.1)
    GRID_W     = Inches(8.0)
    BOX_GAP    = Inches(0.12)

    items   = list(sections.items())
    n_cols  = 2
    n_rows  = (len(items) + 1) // n_cols
    col_w   = (GRID_W - BOX_GAP * (n_cols - 1)) / n_cols

    avail_h = Inches(6.1)
    # estimate box height based on max lines in any box
    max_lines = max(len(v) for v in sections.values())
    box_h   = min(
        (avail_h - BOX_GAP * (n_rows - 1)) / n_rows,
        Inches(0.38) + Inches(0.21) * max_lines + Inches(0.08),
    )

    for i, (sec_label, sec_lines) in enumerate(items):
        col = i % n_cols
        row = i // n_cols
        lft = BOX_LEFT + col * (col_w + BOX_GAP)
        top = BOX_TOP  + row * (box_h + BOX_GAP)
        _design_box(slide, sec_label, sec_lines, lft, top, col_w, box_h)

    # ---- Right panel: Goal + Hypothesis ----
    R_LEFT = Inches(8.6)
    R_TOP  = Inches(1.1)
    R_W    = Inches(4.5)
    R_H    = Inches(6.1)

    _panel_bg(slide, R_LEFT, R_TOP, R_W, R_H)

    _txt(slide, "Goal",
         R_LEFT + Inches(0.15), R_TOP + Inches(0.15),
         R_W - Inches(0.3), Inches(0.32),
         bold=True, size_pt=13, color=C_HEAD)
    _rule(slide, R_TOP + Inches(0.5),
          width=R_W - Inches(0.3), left=R_LEFT + Inches(0.15), thickness_pt=0.5)

    _txt(slide, goal,
         R_LEFT + Inches(0.15), R_TOP + Inches(0.58),
         R_W - Inches(0.3), Inches(2.4),
         size_pt=10.5, color=C_BODY, wrap=True)

    if hypothesis:
        _txt(slide, "Hypothesis",
             R_LEFT + Inches(0.15), R_TOP + Inches(3.1),
             R_W - Inches(0.3), Inches(0.32),
             bold=True, size_pt=13, color=C_HEAD)
        _rule(slide, R_TOP + Inches(3.45),
              width=R_W - Inches(0.3), left=R_LEFT + Inches(0.15), thickness_pt=0.5)
        _txt(slide, hypothesis,
             R_LEFT + Inches(0.15), R_TOP + Inches(3.53),
             R_W - Inches(0.3), Inches(2.3),
             size_pt=10.5, color=C_BODY, wrap=True)


# ---------------------------------------------------------------------------
# Plot helper
# ---------------------------------------------------------------------------

BUDGETS_B2      = [10, 20, 25, 50, 75, 100, 150, 200, 300, 500, 750]
BUDGETS_B12     = [10, 25, 50, 75, 150, 400, 750]
BUDGETS_B12_DS2 = [10, 20, 25, 50, 75, 100, 150]


def _summarize(path, budgets):
    df  = pd.read_csv(path)
    num = df[df["budget"] != "all"].copy()
    num["budget_int"] = num["budget"].astype(int)
    s   = num.groupby("budget_int")["bal_acc"].agg(["mean", "std"])
    s   = s.reindex(budgets)
    return s


def _fig_le_curves(datasets: list[dict], budgets: list[int],
                   title: str) -> plt.Figure:
    """
    datasets: list of {label, color, mean, std}
    """
    fig, ax = plt.subplots(figsize=(8.5, 5.8), facecolor="white")

    x = np.arange(len(budgets))

    for ds in datasets:
        mean = ds["mean"]
        std  = ds["std"]
        valid = ~np.isnan(mean)
        ax.plot(x[valid], mean[valid] * 100,
                color=ds["color"], linewidth=2.2,
                marker="o", markersize=5,
                label=ds["label"])
        ax.fill_between(x[valid],
                        (mean[valid] - std[valid]) * 100,
                        (mean[valid] + std[valid]) * 100,
                        color=ds["color"], alpha=0.12)

    ax.axhline(95, color="#888888", linestyle="--", linewidth=1.2,
               alpha=0.75, zorder=0)
    ax.text(len(budgets) - 0.08, 95.5, "95 %",
            color="#888888", fontsize=8.5, ha="right", va="bottom")

    ax.set_xticks(x)
    ax.set_xticklabels([str(b) for b in budgets], fontsize=9)
    ax.set_xlabel("Label budget (n patches)", fontsize=11)
    ax.set_ylabel("Balanced accuracy (%)", fontsize=11)
    ax.set_ylim(42, 102)
    ax.legend(fontsize=10, framealpha=0.9, loc="lower right")
    ax.set_facecolor("white")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title(title, fontsize=11, fontweight="bold", pad=8)
    ax.grid(axis="y", color="#EEEEEE", linewidth=0.7)
    fig.tight_layout()
    return fig


# ---------------------------------------------------------------------------
# Baseline method slides — CellProfiler and ilastik
# ---------------------------------------------------------------------------

def _slide_cp_overview(prs):
    slide = _blank(prs)
    _slide_header(slide, "CellProfiler-Style Features — Overview",
                  "50 features per 32×32 FA patch · 11 intensity + 39 Haralick/GLCM texture")
    _bullet_box(slide,
        Inches(0.35), Inches(1.1), Inches(6.3), Inches(2.7),
        "What is CellProfiler?",
        ["Open-source image analysis platform (Broad Institute)",
         "Designed for high-throughput microscopy quantification",
         "Measures morphology, intensity, and texture per object",
         "Standard baseline tool in cell biology phenotype profiling"],
        header_color=C_CP)
    _bullet_box(slide,
        Inches(6.85), Inches(1.1), Inches(6.3), Inches(2.7),
        "How features are used here",
        ["Extracted per segmented FA patch (32×32 px, PAX channel)",
         "CIO-normalized (background subtracted, ÷ 5× cell contrast)",
         "Fed into LGBM or logistic regression classifier",
         "Benchmark baseline vs. SupCon-AE learned representation"],
        header_color=C_CP)
    _bullet_box(slide,
        Inches(0.35), Inches(4.0), Inches(12.8), Inches(3.1),
        "Feature Set — 50 features total",
        ["Intensity features (11): mean, std, median, min, max, integrated sum, MAD, p10, p25, p75, p90",
         "Haralick / GLCM texture (39): 13 features × 3 distances (d=1, 2, 4 px), averaged over 4 angles (0°, 45°, 90°, 135°)",
         "GLCM computed with 16 gray levels after per-patch min–max quantization → texture is intensity-scale invariant",
         "Key finding: texture features (39) carry virtually all discriminative power; intensity alone is much weaker"],
        header_color=C_ACCENT, bullet_size=11)


def _slide_cp_intensity(prs):
    slide = _blank(prs)
    _slide_header(slide, "CellProfiler Features — Intensity (11)",
                  "Direct pixel statistics over the 32×32 FA patch")
    rows = [
        ("intensity_mean",       "Mean pixel value across the patch"),
        ("intensity_std",        "Standard deviation of pixel values"),
        ("intensity_median",     "Median pixel value"),
        ("intensity_min",        "Minimum pixel value"),
        ("intensity_max",        "Maximum pixel value (bright FA peak)"),
        ("intensity_integrated", "Sum of all pixel values (total signal)"),
        ("intensity_mad",        "Median Absolute Deviation — robust spread measure"),
        ("intensity_p10",        "10th percentile — lower tail of signal"),
        ("intensity_p25",        "25th percentile — lower quartile"),
        ("intensity_p75",        "75th percentile — upper quartile"),
        ("intensity_p90",        "90th percentile — upper tail / bright spots"),
    ]
    col_x = [Inches(0.4), Inches(4.8)]
    row_h = Inches(0.50)
    top0  = Inches(1.1)
    _txt(slide, "Feature name", col_x[0], top0, Inches(4.2), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _txt(slide, "What it captures", col_x[1], top0, Inches(8.2), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _rule(slide, top0 + row_h, thickness_pt=0.5)
    for i, (name, desc) in enumerate(rows):
        top = top0 + row_h + Inches(0.04) + i * row_h
        _row_rect(slide, Inches(0.35), top, Inches(12.6), row_h, i % 2 == 0)
        _txt(slide, name, col_x[0], top + Inches(0.07), Inches(4.2), row_h,
             size_pt=11, color=C_CP, bold=True)
        _txt(slide, desc, col_x[1], top + Inches(0.07), Inches(8.2), row_h,
             size_pt=11, color=C_BODY)


def _slide_cp_haralick(prs):
    slide = _blank(prs)
    _slide_header(slide, "CellProfiler Features — Haralick / GLCM Texture (39)",
                  "13 features × 3 distances (d=1, 2, 4 px) · averaged over 4 angles · 16 gray levels")
    rows = [
        ("angular_second_moment",     "Energy / uniformity — high for smooth/uniform regions"),
        ("contrast",                  "Local intensity variation between neighbouring pixels"),
        ("correlation",               "Linear dependency of grey levels across neighbours"),
        ("variance",                  "Variance of marginal intensity distribution"),
        ("inverse_difference_moment", "Homogeneity — high when neighbouring values are similar"),
        ("sum_average",               "Mean of the sum distribution p_{x+y}"),
        ("sum_variance",              "Variance of the sum distribution"),
        ("sum_entropy",               "Entropy of the sum distribution"),
        ("entropy",                   "Randomness of grey-level pairs — high for complex textures"),
        ("difference_variance",       "Variance of the difference distribution p_{x−y}"),
        ("difference_entropy",        "Entropy of the difference distribution"),
        ("info_meas1",                "Information measure of correlation 1  (HXY1−HXY) / max(HX,HY)"),
        ("info_meas2",                "Information measure of correlation 2  √(1 − exp(−2(HXY2−HXY)))"),
    ]
    col_x = [Inches(0.4), Inches(4.8)]
    row_h = Inches(0.432)
    top0  = Inches(1.1)
    _txt(slide, "Feature (suffix _d1 / _d2 / _d4)", col_x[0], top0, Inches(4.2), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _txt(slide, "What it captures", col_x[1], top0, Inches(8.2), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _rule(slide, top0 + row_h, thickness_pt=0.5)
    for i, (name, desc) in enumerate(rows):
        top = top0 + row_h + Inches(0.02) + i * row_h
        _row_rect(slide, Inches(0.35), top, Inches(12.6), row_h, i % 2 == 0)
        _txt(slide, name, col_x[0], top + Inches(0.06), Inches(4.2), row_h,
             size_pt=10.5, color=C_CP, bold=True)
        _txt(slide, desc, col_x[1], top + Inches(0.06), Inches(8.2), row_h,
             size_pt=10.5, color=C_BODY)


def _slide_cp_histograms(prs):
    slide = _blank(prs)
    _slide_header(slide, "CellProfiler Features — Distributions (DS1 / DS2 / DS3)",
                  "50 features · 1–99th percentile clipped · blue=DS1(vinc)  red=DS2(pfak)  green=DS3(ppax)")
    _img(slide, CP_HIST, Inches(0.1), Inches(1.1), width=Inches(13.1))


def _slide_ilastik_overview(prs):
    slide = _blank(prs)
    _slide_header(slide, "ilastik-Style Features — Overview",
                  "80 features per 32×32 FA patch · 8 filter types × 5 scales, summarized as mean + std")
    _bullet_box(slide,
        Inches(0.35), Inches(1.1), Inches(6.3), Inches(2.7),
        "What is ilastik?",
        ["Interactive machine learning tool for image segmentation",
         "Uses multiscale filter banks as pixel-level features",
         "User labels pixels → Random Forest classifies the rest",
         "Widely used in bioimage analysis for semantic segmentation"],
        header_color=C_IL)
    _bullet_box(slide,
        Inches(6.85), Inches(1.1), Inches(6.3), Inches(2.7),
        "How features are used here",
        ["8 filter types × 5 scales (σ = 0.3, 0.7, 1.0, 1.6, 3.5 px) = 40 maps",
         "Each map summarized: mean + std over 32×32 patch → 80 features",
         "Fed into LGBM or logistic regression classifier",
         "Captures multi-scale texture, edges, blobs, ridges"],
        header_color=C_IL)
    _bullet_box(slide,
        Inches(0.35), Inches(4.0), Inches(12.8), Inches(3.1),
        "Filter types (8 total)",
        ["Gaussian smoothing — local average intensity at each scale",
         "Laplacian of Gaussian (LoG) — blob detector (bright spots, holes)",
         "Gaussian gradient magnitude — edge strength at FA boundaries",
         "Difference of Gaussians (DoG) — band-pass, approximates LoG",
         "Structure tensor eigenvalues (×2: λ₁ large, λ₂ small) — orientation coherence, edges vs corners",
         "Hessian eigenvalues (×2: λ₁ large, λ₂ small) — ridges, elongated FA structures, focal complexes"],
        header_color=C_ACCENT, bullet_size=11)


def _slide_ilastik_filters(prs):
    slide = _blank(prs)
    _slide_header(slide, "ilastik Features — Filter Definitions",
                  "Each applied at σ = 0.3, 0.7, 1.0, 1.6, 3.5 px · summarized as mean + std per patch → 80 features")
    rows = [
        ("Gaussian smoothing",
         "Low-pass filter; captures local average intensity at each scale. "
         "Large σ → coarse structure; small σ → fine detail."),
        ("Laplacian of Gaussian (LoG)",
         "Second-derivative blob detector. Highlights bright spots and dark holes. "
         "Response peaks where image curvature matches scale σ."),
        ("Gaussian gradient magnitude",
         "First-derivative edge detector: √(Gx² + Gy²) after Gaussian smoothing. "
         "High at FA boundaries and intensity transitions."),
        ("Difference of Gaussians (DoG)",
         "Band-pass filter: Gaussian(σ) − Gaussian(σ√2). "
         "Approximates LoG; highlights structure between two spatial scales."),
        ("Structure tensor — large eigenvalue",
         "Describes local orientation coherence. Large λ₁ indicates strong directional gradient "
         "(edges, oriented fibres). Built from outer product of smoothed gradients."),
        ("Structure tensor — small eigenvalue",
         "Small λ₂ ≈ 0 at edges (one direction), large at corners/junctions. "
         "Ratio λ₁/λ₂ distinguishes edges from isotropic texture."),
        ("Hessian — large eigenvalue",
         "Second-order shape: large |λ₁| at ridges or valleys. "
         "Detects elongated FA structures and fibrillar adhesion ridges."),
        ("Hessian — small eigenvalue",
         "Near-zero at ridges (one principal curvature). Both eigenvalues large at "
         "blob-like focal complexes. Combination distinguishes FA sub-types."),
    ]
    col_x = [Inches(0.4), Inches(3.9)]
    row_h = Inches(0.72)
    top0  = Inches(1.1)
    _txt(slide, "Filter", col_x[0], top0, Inches(3.3), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _txt(slide, "What it captures", col_x[1], top0, Inches(9.2), row_h,
         bold=True, size_pt=11, color=C_HEAD)
    _rule(slide, top0 + row_h, thickness_pt=0.5)
    colors = [C_IL, C_IL, C_IL, C_IL, C_CP, C_CP, C_ACCENT, C_ACCENT]
    for i, (name, desc) in enumerate(rows):
        top = top0 + row_h + Inches(0.02) + i * row_h
        _row_rect(slide, Inches(0.35), top, Inches(12.6), row_h, i % 2 == 0)
        _txt(slide, name, col_x[0], top + Inches(0.1), Inches(3.3), row_h,
             size_pt=10.5, color=colors[i], bold=True)
        _txt(slide, desc, col_x[1], top + Inches(0.1), Inches(9.2), row_h,
             size_pt=10.5, color=C_BODY)


def _slide_ilastik_histograms(prs):
    slide = _blank(prs)
    _slide_header(slide, "ilastik Features — Distributions (DS1 / DS2 / DS3)",
                  "80 features · 1–99th percentile clipped · blue=DS1(vinc)  red=DS2(pfak)  green=DS3(ppax)")
    _img(slide, IL_HIST, Inches(0.1), Inches(1.1), width=Inches(13.1))


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def _slide_title(prs):
    slide = _blank(prs)
    _txt(slide,
         "Label Efficiency Benchmark",
         Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0),
         bold=True, size_pt=36, color=C_HEAD, align=PP_ALIGN.CENTER)
    _txt(slide,
         "SupCon-AE  vs  CellProfiler  vs  ilastik",
         Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.6),
         size_pt=20, color=C_GREY, align=PP_ALIGN.CENTER)
    _rule(slide, Inches(3.85), width=Inches(9), left=Inches(2.17))
    _txt(slide,
         "DS1 vinc ctrl+ycomp  ·  B2 (Annabel) labels  ·  LGBM & logreg classifiers  ·  5-fold CV × 5 repeats",
         Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.5),
         size_pt=12, color=C_GREY, align=PP_ALIGN.CENTER)
    _txt(slide,
         "Cold-start protocol: SupCon-AE (lat=12, proj=8) retrained from scratch at each label budget  ·  n_labeled_per_class bug fixed",
         Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.4),
         size_pt=11, color=RGBColor(0x8B, 0x00, 0x00), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slides 2a/2b — DS1 B2 design + results
# ---------------------------------------------------------------------------

_DS1_B2_SECTIONS = {
    "Dataset": [
        "DS1: vinc (vincristine)",
        "Conditions: ctrl + ycomp (mixed)",
        "1,224 labeled patches total",
        "ctrl: 539  ycomp: 685",
        "32×32 px, PAX channel",
    ],
    "Labels": [
        "B2: Annabel labels only",
        "Binary: adhesion (454) / No adhesion (770)",
        "5-fold stratified CV",
        "~245 test / ~979 train per fold",
        "5 repeats per budget",
    ],
    "Training (SupCon-AE)": [
        "Cold-start: retrained from scratch",
        "per budget subset (ctrl+ycomp labels)",
        "lat=12, proj=8, 500 epochs",
        "λ_supcon=5, λ_recon=1, λ_contrast=0.5",
        "No n_labeled_per_class cap (fixed)",
        "Budgets: 10,20,25,50,75,100,150,200,300,500,750",
    ],
    "Evaluation": [
        "Held-out fold (~245 patches, ctrl+ycomp)",
        "Classifiers: LGBM  and  logistic regression",
        "Frozen z-latents (SupCon-AE, dim=12)",
        "Raw features (CP / ilastik)",
        "Metric: balanced accuracy",
    ],
    "Models compared": [
        "SupCon-AE (this work, lat=12)",
        "CellProfiler (2D feature set)",
        "ilastik pixel features",
    ],
}

def _slide_ds1_b2_design(prs):
    _slide_design(
        prs,
        title    = "Experiment Design — DS1 B2 Label Efficiency",
        subtitle = "Single-annotator baseline: how does SupCon-AE scale with Annabel's labels?",
        sections = _DS1_B2_SECTIONS,
        goal     = (
            "Measure how balanced accuracy scales with the labeled-patch budget "
            "for SupCon-AE vs. CellProfiler vs. ilastik on DS1 (vincristine) "
            "using annotations from a single annotator (Annabel, B2).\n\n"
            "This establishes the single-annotator baseline before testing "
            "whether pooling two annotators (B12) improves label efficiency."
        ),
        hypothesis = (
            "SupCon-AE should approach CP/ilastik at moderate budgets "
            "(~150–300 patches) because its learned latent representation "
            "should capture morphological variation better than hand-crafted "
            "pixel statistics — especially in the low-label regime where "
            "unsupervised pretraining provides the most benefit."
        ),
    )


def _slide_ds1_b2_lgbm(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B2 — Label Efficiency: SupCon-AE vs CP vs ilastik  [LGBM]",
                  "vinc ctrl+ycomp · 1,224 patches · lat=12 fixed · LGBM classifier · 5-fold CV × 5 repeats")

    sc = _summarize(EVAL_DIR / "supcon_le_b2_lat12p8_ds1.csv",  BUDGETS_B2)
    cp = _summarize(EVAL_DIR / "cp_ds1.csv",                    BUDGETS_B2)
    il = _summarize(EVAL_DIR / "ilastik_ds1.csv",               BUDGETS_B2)

    datasets = [
        dict(label="SupCon-AE (lat=12, fixed)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(
        datasets, BUDGETS_B2,
        "DS1 B2 · Balanced Accuracy vs Label Budget (LGBM)",
    )
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS1_B2_SECTIONS)


def _slide_ds1_b2_logreg(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B2 — Label Efficiency: SupCon-AE vs CP vs ilastik  [Logreg]",
                  "vinc ctrl+ycomp · 1,224 patches · lat=12 fixed · logistic regression · 5-fold CV × 5 repeats")

    sc = _summarize(EVAL_DIR / "supcon_le_b2_lat12p8_logreg_ds1.csv", BUDGETS_B2)
    cp = _summarize(EVAL_DIR / "cp_b2_logreg_ds1.csv",                BUDGETS_B2)
    il = _summarize(EVAL_DIR / "ilastik_b2_logreg_ds1.csv",           BUDGETS_B2)

    datasets = [
        dict(label="SupCon-AE (lat=12, fixed)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(
        datasets, BUDGETS_B2,
        "DS1 B2 · Balanced Accuracy vs Label Budget (logreg)",
    )
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS1_B2_SECTIONS)


# ---------------------------------------------------------------------------
# Slides 3a/3b — DS1 B12 design + results
# ---------------------------------------------------------------------------

_DS1_B12_SECTIONS = {
    "Dataset": [
        "DS1: vinc (vincristine)",
        "Conditions: ctrl + ycomp",
        "~2,428 labeled patches",
        "32×32 px, PAX channel",
    ],
    "Labels": [
        "B12: B1 (Margaret) + B2 (Annabel)",
        "53 conflicts dropped, B2 priority",
        "Binary: adhesion / No adhesion",
        "adh=1,014 / noad=1,414",
        "5-fold stratified CV × 5 repeats",
    ],
    "Training (SupCon-AE)": [
        "Cold-start: retrained from scratch",
        "per budget subset",
        "lat=32, proj=8, 500 epochs",
        "λ_supcon=5, λ_recon=1",
        "Budgets: 10,25,50,75,150,400,750",
    ],
    "Evaluation": [
        "Held-out fold (~486 patches)",
        "Classifier: logistic regression",
        "Frozen z-latents (SupCon-AE)",
        "Raw features (CP / ilastik)",
        "Metric: balanced accuracy",
    ],
    "Models compared": [
        "SupCon-AE (this work)",
        "CellProfiler (2D feature set)",
        "ilastik pixel features",
    ],
}

def _slide_ds1_b12_design(prs):
    _slide_design(
        prs,
        title    = "Experiment Design — DS1 B12 Label Efficiency",
        subtitle = "Dual-annotator pool: does combining B1+B2 labels close the gap to handcrafted features?",
        sections = _DS1_B12_SECTIONS,
        goal     = (
            "Test whether merging two annotators' labels (B1: Margaret + B2: Annabel, "
            "~2,428 patches total) improves label efficiency for SupCon-AE relative "
            "to CellProfiler and ilastik, compared to the single-annotator (B2) baseline.\n\n"
            "If SupCon-AE benefits more from annotation diversity than handcrafted "
            "features, the gap should narrow at equivalent budgets."
        ),
        hypothesis = (
            "More diverse labels (two annotators) should provide SupCon-AE with "
            "richer supervisory signal, helping the contrastive loss learn more "
            "robust boundaries. CP/ilastik are feature-fixed, so they benefit "
            "less from annotator diversity — expect the gap to narrow vs. B2-only."
        ),
    )


def _slide_ds1_b12(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik",
                  "vinc ctrl+ycomp · B1 (Margaret) + B2 (Annabel) combined (~2,428 patches) · logreg · 5-fold CV × 5 repeats")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_supcon_logreg_ds1.csv", BUDGETS_B12)
    cp = _summarize(EVAL_DIR / "cp_b12_logreg_ds1.csv",               BUDGETS_B12)
    il = _summarize(EVAL_DIR / "ilastik_b12_logreg_ds1.csv",          BUDGETS_B12)

    datasets = [
        dict(label="SupCon-AE (lat=32, proj=8)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(
        datasets, BUDGETS_B12,
        "DS1 B12 · Balanced Accuracy vs Label Budget (logreg)",
    )
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS1_B12_SECTIONS)


# ---------------------------------------------------------------------------
# Slides 4a/4b — B2 vs B12 design + comparison
# ---------------------------------------------------------------------------

_B2_VS_B12_SECTIONS = {
    "Comparison": [
        "Same dataset, same protocol",
        "Only the label pool differs",
        "Shared budgets: 10,25,50,75,150,750",
        "Classifier: logistic regression",
    ],
    "B2 (Annabel only)": [
        "~1,225 patches",
        "1 annotator",
        "ctrl + ycomp frames",
        "adh / noad balanced",
    ],
    "B12 (B1+B2 merged)": [
        "~2,428 patches (2× B2)",
        "2 annotators merged",
        "53 conflicts dropped",
        "adh=1,014 / noad=1,414",
    ],
    "Protocol": [
        "SupCon-AE: cold-start per budget",
        "lat=32, proj=8, 500 epochs",
        "CP / ilastik: fixed features",
        "5-fold CV × 5 repeats",
    ],
}

def _slide_b2_vs_b12_design(prs):
    _slide_design(
        prs,
        title    = "Experiment Design — B2 vs B12 Comparison",
        subtitle = "Does pooling two annotators' labels help SupCon-AE close the gap?",
        sections = _B2_VS_B12_SECTIONS,
        goal     = (
            "Directly compare label efficiency under the same experimental "
            "conditions (DS1, cold-start SupCon-AE, logreg) but with two "
            "different annotation pools: B2-only (~1,225 patches, 1 annotator) "
            "vs. B12 combined (~2,428 patches, 2 annotators).\n\n"
            "This isolates the effect of annotation pool size and diversity "
            "on the SupCon-AE vs. CP gap."
        ),
        hypothesis = (
            "If the gap between SupCon-AE and CP narrows in B12 relative to B2 "
            "at the same label budget, it suggests SupCon-AE benefits more from "
            "annotator diversity. If the gap stays constant, annotation volume "
            "alone cannot explain any advantage of learned representations."
        ),
    )


def _slide_b2_vs_b12(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "B2 vs B12 Label Set Comparison — DS1",
                  "Same protocol, more labels: does SupCon-AE close the gap to CP/ilastik?")

    budgets_shared = [10, 25, 50, 75, 150, 750]

    def _get(path, budgets):
        df  = pd.read_csv(path)
        num = df[df["budget"] != "all"].copy()
        num["budget_int"] = num["budget"].astype(int)
        s   = num.groupby("budget_int")["bal_acc"].agg(["mean", "std"])
        return s.reindex(budgets)

    sc_b2  = _get(EVAL_DIR / "supcon_le_b2_lat12p8_logreg_ds1.csv", budgets_shared)
    cp_b2  = _get(EVAL_DIR / "cp_b2_logreg_ds1.csv",                budgets_shared)
    sc_b12 = _get(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_logreg_ds1.csv", budgets_shared)
    cp_b12 = _get(EVAL_DIR / "cp_b12_logreg_ds1.csv",               budgets_shared)

    fig, axes = plt.subplots(1, 2, figsize=(11, 5.5), facecolor="white")

    x = np.arange(len(budgets_shared))
    xlabels = [str(b) for b in budgets_shared]

    for ax, (sc, cp, title_tag) in zip(axes, [
        (sc_b2,  cp_b2,  "B2 (~1,225 patches)"),
        (sc_b12, cp_b12, "B12 (~2,428 patches)"),
    ]):
        for vals, color, label in [
            (sc, COL_SUPCON,  "SupCon-AE"),
            (cp, COL_CP,      "CellProfiler"),
        ]:
            m = vals["mean"].values
            s = vals["std"].values
            valid = ~np.isnan(m)
            ax.plot(x[valid], m[valid]*100, color=color, linewidth=2.2,
                    marker="o", markersize=5, label=label)
            ax.fill_between(x[valid], (m[valid]-s[valid])*100, (m[valid]+s[valid])*100,
                            color=color, alpha=0.12)

        ax.set_xticks(x)
        ax.set_xticklabels(xlabels, fontsize=9)
        ax.set_xlabel("Label budget", fontsize=10)
        ax.set_ylabel("Balanced accuracy (%)", fontsize=10)
        ax.set_ylim(42, 102)
        ax.axhline(95, color="#888888", linestyle="--", linewidth=1.2, alpha=0.75, zorder=0)
        ax.text(len(budgets_shared) - 0.08, 95.5, "95 %",
                color="#888888", fontsize=7.5, ha="right", va="bottom")
        ax.legend(fontsize=9, framealpha=0.9, loc="lower right")
        ax.set_facecolor("white")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_title(f"DS1 {title_tag}\nSupCon-AE vs CP (logreg)", fontsize=10, fontweight="bold")
        ax.grid(axis="y", color="#EEEEEE", linewidth=0.7)

    fig.tight_layout(pad=1.5)
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)

    _add_info_panel(slide, _B2_VS_B12_SECTIONS)


# ---------------------------------------------------------------------------
# Slides — DS2 B12 LGBM
# ---------------------------------------------------------------------------

_DS2_B12_SECTIONS = {
    "Dataset": [
        "DS2: pfak (pFAK antibody)",
        "244 labeled patches total",
        "32×32 px, PAX channel",
        "adh=175 / noad=69",
    ],
    "Labels": [
        "B12: B1 (Margaret) + B2 (Annabel)",
        "Binary: adhesion / No adhesion",
        "5-fold stratified CV × 5 repeats",
        "~49 patches per test fold",
    ],
    "Training (SupCon-AE)": [
        "Cold-start per budget subset",
        "lat=12, proj=8, 500 epochs",
        "λ_supcon=5, λ_recon=1",
        "Budgets: 10,20,25,50,75,100,150",
    ],
    "Evaluation": [
        "Held-out fold (~49 patches)",
        "Classifier: LightGBM",
        "Frozen z-latents (SupCon-AE)",
        "Raw features (CP / ilastik)",
        "Metric: balanced accuracy",
    ],
    "Models compared": [
        "SupCon-AE (lat=12, this work)",
        "CellProfiler (50 features)",
        "ilastik pixel features (80)",
    ],
}


def _slide_ds2_b12_lgbm(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS2 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik (LGBM)",
                  "pfak · B1+B2 combined (244 patches) · LightGBM · 5-fold CV × 5 repeats · lat=12/proj=8")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_ds2_lat12p8_ds2.csv", BUDGETS_B12_DS2)
    cp = _summarize(EVAL_DIR / "cp_b12_ds2.csv",                    BUDGETS_B12_DS2)
    il = _summarize(EVAL_DIR / "ilastik_b12_ds2.csv",               BUDGETS_B12_DS2)

    datasets = [
        dict(label="SupCon-AE (lat=12, proj=8)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(
        datasets, BUDGETS_B12_DS2,
        "DS2 B12 · Balanced Accuracy vs Label Budget (LGBM)",
    )
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS2_B12_SECTIONS)


# ---------------------------------------------------------------------------
# Feature descriptions
# ---------------------------------------------------------------------------

_FEAT_DESC = {
    "intensity_mean":             "Mean brightness — FAs are brighter protein clusters",
    "intensity_std":              "Brightness spread — FAs have high local variance",
    "intensity_median":           "Median brightness — robust FA vs. background contrast",
    "intensity_min":              "Min pixel — background level in patch",
    "intensity_max":              "Peak intensity — highlights bright FA core",
    "intensity_integrated":       "Total signal — larger/brighter FAs have higher sum",
    "intensity_mad":              "MAD — robust measure of heterogeneous intensity",
    "intensity_p10":              "10th percentile — captures dark patch background",
    "intensity_p25":              "Lower quartile — separates FA from cytoplasm",
    "intensity_p75":              "Upper quartile — bright FA signal above background",
    "intensity_p90":              "90th percentile — highlights brightest FA pixels",
    "angular_second_moment":      "GLCM energy — high for smooth/uniform patches (No adhesion)",
    "contrast":                   "GLCM contrast — local intensity variation; high in FA edges",
    "correlation":                "GLCM correlation — grey-level linear dependency across neighbours",
    "variance":                   "GLCM variance — spread of intensity distribution",
    "inverse_difference_moment":  "GLCM homogeneity — high for uniform patches (No adhesion)",
    "sum_average":                "GLCM sum average — captures mean grey-level pair intensity",
    "sum_variance":               "GLCM sum variance — spread of paired-pixel sums",
    "sum_entropy":                "GLCM sum entropy — texture complexity of pair sums",
    "entropy":                    "GLCM entropy — randomness of grey-level pairs; high = complex texture",
    "difference_variance":        "GLCM diff variance — variance of pixel-pair differences",
    "difference_entropy":         "GLCM diff entropy — entropy of pixel-pair differences",
    "info_meas1":                 "GLCM info measure 1 — correlation of marginal distributions",
    "info_meas2":                 "GLCM info measure 2 — sqrt(1−exp(−2×ΔHXY))",
}

def _feat_desc(name: str) -> str:
    base = name
    for suf in ("_d1", "_d2", "_d4"):
        if name.endswith(suf):
            base = name[: -len(suf)]
            dist = suf[1:]
            break
    else:
        dist = None
    desc = _FEAT_DESC.get(base, "")
    if dist:
        desc += f"  [d={dist}px]"
    return desc


# ---------------------------------------------------------------------------
# Patch loading
# ---------------------------------------------------------------------------

def _load_tif(filename: str, patch_dir: Path) -> np.ndarray:
    cond = filename.split("_f")[0]
    p = patch_dir / cond / "tiff_patches32_mr10" / filename
    if not p.exists():
        return np.zeros((32, 32), dtype=np.float32)
    from PIL import Image
    img = np.array(Image.open(str(p))).astype(np.float32)
    if img.ndim == 3:
        img = img[..., 0]
    vmin, vmax = np.percentile(img, 2), np.percentile(img, 98)
    return np.clip((img - vmin) / max(float(vmax - vmin), 1e-6), 0.0, 1.0)


# ---------------------------------------------------------------------------
# Feature importance — LGBM voting across folds
# ---------------------------------------------------------------------------

def _anova_f(X: np.ndarray, y: np.ndarray) -> np.ndarray:
    """Pure-numpy per-column ANOVA F-statistic (binary y)."""
    classes = np.unique(y)
    n = len(y)
    grand_mean = np.nanmean(X, axis=0)
    ss_between = np.zeros(X.shape[1])
    ss_within  = np.zeros(X.shape[1])
    for c in classes:
        mask = y == c
        nc = mask.sum()
        if nc == 0:
            continue
        cm = np.nanmean(X[mask], axis=0)
        ss_between += nc * (cm - grand_mean) ** 2
        ss_within  += np.nansum((X[mask] - cm) ** 2, axis=0)
    k = len(classes)
    ms_between = ss_between / max(k - 1, 1)
    ms_within  = ss_within  / max(n - k, 1)
    with np.errstate(all="ignore"):
        f = np.where(ms_within > 0, ms_between / ms_within, 0.0)
    return np.nan_to_num(f, nan=0.0)


def _top_features_voted(feat_df: pd.DataFrame, ann_df: pd.DataFrame,
                        fold_splits: pd.DataFrame, n: int = 3):
    """
    Compute univariate ANOVA F-statistics per feature on each fold's train
    set (binary: adhesion vs No adhesion), average across folds, return
    top-n (name, avg_F).  Pure numpy — no sklearn or lightgbm.
    """
    feat_cols = [c for c in feat_df.columns if c != "filename"]

    merged = ann_df[["filename", "label"]].merge(feat_df, on="filename")
    merged["y"] = (merged["label"] != "No adhesion").astype(int)

    fs = fold_splits[["unique_ID", "fold"]].rename(columns={"unique_ID": "filename"})
    merged = merged.merge(fs, on="filename", how="left")
    merged["fold"] = merged["fold"].fillna(-1).astype(int)

    f_scores_list = []
    for fold in range(5):
        train = merged[merged["fold"] != fold]
        if len(train) < 10 or train["y"].nunique() < 2:
            continue
        X = train[feat_cols].values.astype(float)
        y = train["y"].values
        f_scores_list.append(_anova_f(X, y))

    if not f_scores_list:
        X = merged[feat_cols].values.astype(float)
        y = merged["y"].values
        f_scores_list = [_anova_f(X, y)]

    avg = np.mean(f_scores_list, axis=0)
    top_idx = np.argsort(avg)[::-1][:n]
    return [(feat_cols[i], float(avg[i])) for i in top_idx]


# ---------------------------------------------------------------------------
# Feature visualization figure — histogram + patches at percentiles
# ---------------------------------------------------------------------------

FA_COLORS_MAP = {
    "focal adhesion":    "#e05c5c",
    "Nascent Adhesion":  "#e09a3c",
    "focal complex":     "#5c7de0",
    "fibrillar adhesion":"#5cb85c",
    "No adhesion":       "#888888",
}


def _fig_feature_histo_patches(feat_df: pd.DataFrame,
                                ann_full: pd.DataFrame,
                                top_features: list,
                                patch_dir: Path,
                                n_pct: int = 5) -> plt.Figure:
    """
    3-row figure. Each row = one top feature.
    Left: histogram AD vs No-AD.
    Right: n_pct patches for each class at spread percentiles.
    """
    feat_cols = [c for c in feat_df.columns if c != "filename"]

    merged = ann_full[["filename", "label"]].merge(feat_df, on="filename")
    merged["binary"] = merged["label"].map(
        lambda x: "adhesion" if x != "No adhesion" else "No adhesion"
    )

    pcts = np.linspace(5, 95, n_pct).astype(int)

    n_rows = len(top_features)
    fig = plt.figure(figsize=(13.5, 2.5 * n_rows), facecolor="white")
    gs_outer = gridspec.GridSpec(n_rows, 1, figure=fig,
                                 hspace=0.55)

    for row_i, (feat_name, feat_imp) in enumerate(top_features):
        gs_row = gridspec.GridSpecFromSubplotSpec(
            1, 2, subplot_spec=gs_outer[row_i],
            width_ratios=[3.5, 9.5], wspace=0.08,
        )

        # ---- Left: histogram ----
        ax_hist = fig.add_subplot(gs_row[0])

        ad_vals   = merged.loc[merged["binary"] == "adhesion",    feat_name].dropna().values
        noad_vals = merged.loc[merged["binary"] == "No adhesion", feat_name].dropna().values

        all_vals = np.concatenate([ad_vals, noad_vals])
        lo, hi = np.percentile(all_vals, 1), np.percentile(all_vals, 99)

        for vals, color, label in [
            (noad_vals, HIST_COL_NOAD, "No adhesion"),
            (ad_vals,   HIST_COL_AD,   "adhesion"),
        ]:
            vals_c = np.clip(vals, lo, hi)
            if len(vals_c) < 5:
                continue
            counts, edges = np.histogram(vals_c, bins=40, density=True)
            centers = (edges[:-1] + edges[1:]) / 2
            ax_hist.fill_between(centers, counts, alpha=0.35, color=color)
            ax_hist.plot(centers, counts, color=color, linewidth=1.5, label=label)

        ax_hist.set_xlim(lo, hi)
        ax_hist.set_yticks([])
        ax_hist.set_xlabel("Feature value", fontsize=7)
        desc = _feat_desc(feat_name)
        short_name = feat_name.replace("_d1", "\n(d=1)").replace("_d2", "\n(d=2)").replace("_d4", "\n(d=4)")
        ax_hist.set_title(f"#{row_i+1}  {short_name}\nimp={feat_imp:.0f}",
                          fontsize=7.5, fontweight="bold", loc="left", pad=3)
        ax_hist.spines[["top", "right", "left"]].set_visible(False)
        ax_hist.tick_params(labelsize=6)
        ax_hist.legend(fontsize=6, framealpha=0.8, loc="upper right")
        ax_hist.grid(axis="x", color="#EEEEEE", linewidth=0.5)

        # ---- Right: patches at percentiles ----
        gs_patches = gridspec.GridSpecFromSubplotSpec(
            2, n_pct, subplot_spec=gs_row[1],
            hspace=0.08, wspace=0.05,
        )

        for cls_i, (cls_label, cls_vals, cls_df) in enumerate([
            ("No adhes.", noad_vals, merged[merged["binary"] == "No adhesion"]),
            ("adhesion",  ad_vals,   merged[merged["binary"] == "adhesion"]),
        ]):
            if len(cls_vals) == 0:
                continue
            cutoffs = np.percentile(cls_vals, pcts)
            cls_sorted = cls_df.sort_values(feat_name).reset_index(drop=True)

            for pct_i, (pct, cutoff) in enumerate(zip(pcts, cutoffs)):
                ax_p = fig.add_subplot(gs_patches[cls_i, pct_i])

                # find patch closest to this percentile
                dists = np.abs(cls_sorted[feat_name].values - cutoff)
                best_idx = int(np.argmin(dists))
                best_row = cls_sorted.iloc[best_idx]
                img = _load_tif(best_row["filename"], patch_dir)

                ax_p.imshow(img, cmap="gray", vmin=0, vmax=1, interpolation="nearest")
                ax_p.set_xticks([])
                ax_p.set_yticks([])
                border_color = FA_COLORS_MAP.get(best_row["label"], "#888888")
                for spine in ax_p.spines.values():
                    spine.set_edgecolor(border_color)
                    spine.set_linewidth(2.0)

                if cls_i == 0:  # top row: add percentile label above
                    ax_p.set_title(f"p{pct}", fontsize=5.5, pad=1.5)
                if pct_i == 0:  # leftmost patch: add class label
                    ax_p.set_ylabel(cls_label, fontsize=5.5, rotation=0,
                                    ha="right", va="center", labelpad=22)

    fig.tight_layout(pad=0.5)
    return fig


# ---------------------------------------------------------------------------
# New slides: CP / IL top features
# ---------------------------------------------------------------------------

def _slide_cp_top_features(prs):
    fold_splits = pd.read_csv(FOLD_SPLITS_B2)
    ann = pd.read_csv(ANN_B2_CSV)
    cp  = pd.read_csv(CP_CSV)

    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)
    fold_splits["unique_ID"] = fold_splits["unique_ID"].str.replace("-f", "_f", n=1)

    print("  Computing CP top features by LGBM voting …", flush=True)
    top3 = _top_features_voted(cp, ann, fold_splits, n=3)
    print(f"    Top features: {[t[0] for t in top3]}", flush=True)

    fig = _fig_feature_histo_patches(cp, ann, top3, PATCH_DIR)

    slide = _blank(prs)
    _slide_header(slide,
                  "CellProfiler — Top 3 Discriminative Features (DS1 B2)",
                  "LGBM importance voted across 5 folds · grey border=No adhesion · coloured border=FA subtype")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))

    # explanation text strip
    y = Inches(1.07)
    dy = Inches(2.08)
    for feat_name, _ in top3:
        desc = _feat_desc(feat_name)
        if desc:
            _txt(slide, f"  → {desc}",
                 Inches(0.22), y + dy * 0.82,
                 Inches(3.1), Inches(0.32),
                 size_pt=7, color=C_GREY)
        y += dy


def _slide_il_top_features(prs):
    fold_splits = pd.read_csv(FOLD_SPLITS_B2)
    ann = pd.read_csv(ANN_B2_CSV)
    il  = pd.read_csv(IL_CSV)

    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)
    fold_splits["unique_ID"] = fold_splits["unique_ID"].str.replace("-f", "_f", n=1)

    # ensure filename key alignment
    il_fn = set(il["filename"])
    cp_fn = set(pd.read_csv(CP_CSV, usecols=["filename"])["filename"])
    if not (il_fn & cp_fn):
        # il may use different key; try matching by index order
        pass

    print("  Computing ilastik top features by LGBM voting …", flush=True)
    top3 = _top_features_voted(il, ann, fold_splits, n=3)
    print(f"    Top features: {[t[0] for t in top3]}", flush=True)

    fig = _fig_feature_histo_patches(il, ann, top3, PATCH_DIR)

    slide = _blank(prs)
    _slide_header(slide,
                  "ilastik — Top 3 Discriminative Features (DS1 B2)",
                  "LGBM importance voted across 5 folds · grey border=No adhesion · coloured border=FA subtype")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))

    y = Inches(1.07)
    dy = Inches(2.08)
    for feat_name, _ in top3:
        _txt(slide, f"  → {feat_name}: multi-scale filter response (σ varies)",
             Inches(0.22), y + dy * 0.82,
             Inches(3.1), Inches(0.32),
             size_pt=7, color=C_GREY)
        y += dy


# ---------------------------------------------------------------------------
# SupCon-AE top latent dims
# ---------------------------------------------------------------------------

def _top_supcon_dims(latent_path: Path, ann_full: pd.DataFrame, n: int = 3):
    """
    Cohen's d per latent dim between adhesion and No adhesion classes.
    Returns top-n (dim_name, cohens_d).
    """
    lat = pd.read_csv(latent_path)
    z_cols = [c for c in lat.columns if c.startswith("z_")]
    lat = lat[["filename"] + z_cols]
    lat["filename"] = lat["filename"].apply(lambda x: x.replace("-f", "_f", 1)
                                            if isinstance(x, str) else x)

    merged = ann_full[["filename", "label"]].merge(lat, on="filename")
    merged["binary"] = (merged["label"] != "No adhesion").astype(int)

    ad   = merged[merged["binary"] == 1]
    noad = merged[merged["binary"] == 0]

    scores = []
    for z in z_cols:
        a = ad[z].dropna().values
        b = noad[z].dropna().values
        if len(a) < 2 or len(b) < 2:
            scores.append(0.0)
            continue
        pooled_std = np.sqrt((np.var(a) + np.var(b)) / 2 + 1e-9)
        d = abs(np.mean(a) - np.mean(b)) / pooled_std
        scores.append(d)

    scores = np.array(scores)
    top_idx = np.argsort(scores)[::-1][:n]
    return [(z_cols[i], float(scores[i])) for i in top_idx]


def _slide_supcon_top_features(prs, budget: str = "150", fold: int = 0, repeat: int = 0):
    ann = pd.read_csv(ANN_B2_CSV)
    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)

    run_name = f"le_b2_lat12p8_ds1_fv{fold}_nb{budget}_r{repeat}"
    latent_path = LE_B2_DIR / run_name / "latents.csv"
    if not latent_path.exists():
        print(f"  SupCon-AE latents not found: {latent_path}", flush=True)
        return

    print(f"  Computing SupCon-AE top latent dims (budget={budget}) …", flush=True)
    top3 = _top_supcon_dims(latent_path, ann, n=3)
    print(f"    Top dims: {[t[0] for t in top3]}", flush=True)

    lat_df = pd.read_csv(latent_path)
    lat_df["filename"] = lat_df["filename"].apply(
        lambda x: x.replace("-f", "_f", 1) if isinstance(x, str) else x
    )
    z_cols = [c for c in lat_df.columns if c.startswith("z_")]
    feat_df = lat_df[["filename"] + z_cols]

    fig = _fig_feature_histo_patches(feat_df, ann, top3, PATCH_DIR)

    slide = _blank(prs)
    _slide_header(slide,
                  f"SupCon-AE — Top 3 Discriminative Latent Dims (DS1 B2, nb={budget})",
                  f"Cohen's d between adhesion/No-adhesion · fold={fold}, repeat={repeat} · lat=12, proj=8")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))


# ---------------------------------------------------------------------------
# UMAP comparison slides
# ---------------------------------------------------------------------------

def _compute_umap_embedding(X: np.ndarray, seed: int = 42) -> np.ndarray:
    """PCA-based 2D embedding via covariance eigendecomposition (fast, avoids large-matrix SVD)."""
    X_c = X - X.mean(axis=0)
    col_std = X_c.std(axis=0)
    col_std[col_std == 0] = 1.0
    X_c = X_c / col_std
    # Covariance matrix is (n_features × n_features) — far smaller than data matrix
    cov = X_c.T @ X_c / max(X_c.shape[0] - 1, 1)
    eigenvalues, eigenvectors = np.linalg.eigh(cov)
    # eigh returns ascending order; take top 2
    top2 = eigenvectors[:, -2:][:, ::-1]
    return X_c @ top2


def _ax_umap(ax, emb, labels, title, bg_idx=None):
    colors = {"adhesion": HIST_COL_AD, "No adhesion": HIST_COL_NOAD}
    if bg_idx is not None:
        ax.scatter(emb[bg_idx, 0], emb[bg_idx, 1],
                   c="#DDDDDD", s=2, alpha=0.3, linewidths=0, rasterized=True)
        fg_idx = np.array([i for i in range(len(emb)) if i not in set(bg_idx)])
    else:
        fg_idx = np.arange(len(emb))

    for cls, col in colors.items():
        idx = fg_idx[[labels[i] == cls for i in fg_idx]]
        if len(idx):
            ax.scatter(emb[idx, 0], emb[idx, 1],
                       c=col, s=10, alpha=0.8, linewidths=0, label=cls, rasterized=True)

    ax.set_title(title, fontsize=9, fontweight="bold")
    ax.set_xticks([]); ax.set_yticks([])
    ax.spines[["top", "right", "left", "bottom"]].set_visible(False)


def _slide_umap_comparison(prs, budget: str = "150", fold: int = 0, repeat: int = 0):
    ann = pd.read_csv(ANN_B2_CSV)
    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)
    ann_map = dict(zip(ann["filename"], ann["label"].map(
        lambda x: "adhesion" if x != "No adhesion" else "No adhesion"
    )))

    cp_df = pd.read_csv(CP_CSV)
    il_df = pd.read_csv(IL_CSV)

    run_name = f"le_b2_lat12p8_ds1_fv{fold}_nb{budget}_r{repeat}"
    latent_path = LE_B2_DIR / run_name / "latents.csv"
    if not latent_path.exists():
        print(f"  UMAP: SupCon latents not found for nb={budget}, skipping", flush=True)
        return

    lat_df = pd.read_csv(latent_path)
    lat_df["filename"] = lat_df["filename"].apply(
        lambda x: x.replace("-f", "_f", 1) if isinstance(x, str) else x
    )
    z_cols = [c for c in lat_df.columns if c.startswith("z_")]

    common = set(cp_df["filename"]) & set(il_df["filename"]) & set(lat_df["filename"])
    common = sorted(common)
    print(f"  UMAP nb={budget}: {len(common)} common patches", flush=True)

    labels = np.array([ann_map.get(f, "unknown") for f in common])
    bg_idx = np.where(labels == "unknown")[0]

    cp_X  = cp_df.set_index("filename").reindex(common).values.astype(float)
    il_X  = il_df.set_index("filename").reindex(common).values.astype(float)
    lat_X = lat_df.set_index("filename").reindex(common)[z_cols].values.astype(float)

    # replace NaNs with column means
    for X in [cp_X, il_X, lat_X]:
        col_means = np.nanmean(X, axis=0)
        nan_mask = np.isnan(X)
        X[nan_mask] = np.take(col_means, np.where(nan_mask)[1])

    print("  Computing dimensionality reductions …", flush=True)
    cp_emb  = _compute_umap_embedding(cp_X)
    il_emb  = _compute_umap_embedding(il_X)
    lat_emb = _compute_umap_embedding(lat_X)

    fig, axes = plt.subplots(1, 3, figsize=(13, 4.8), facecolor="white")
    _ax_umap(axes[0], cp_emb,  labels, f"CellProfiler (nb={budget})", bg_idx)
    _ax_umap(axes[1], il_emb,  labels, f"ilastik (nb={budget})",      bg_idx)
    _ax_umap(axes[2], lat_emb, labels, f"SupCon-AE lat=12 (nb={budget})", bg_idx)

    handles = [mpatches.Patch(color=HIST_COL_AD,   label="adhesion"),
               mpatches.Patch(color=HIST_COL_NOAD, label="No adhesion"),
               mpatches.Patch(color="#DDDDDD", label="unlabeled")]
    axes[2].legend(handles=handles, fontsize=8, loc="lower right", framealpha=0.9)
    fig.tight_layout()

    slide = _blank(prs)
    _slide_header(slide,
                  f"Embedding: CellProfiler vs ilastik vs SupCon-AE  [nb={budget}]",
                  f"PCA (2D) · DS1 B2 · fold={fold}, repeat={repeat} · pink=adhesion, grey=No adhesion")
    _add_fig(slide, fig, Inches(0.2), Inches(1.05), Inches(13.0), Inches(6.2))


# ---------------------------------------------------------------------------
# Split CP / ilastik histogram slides
# ---------------------------------------------------------------------------

def _fig_cp_histograms_group(feat_cols_grp: list, title: str) -> plt.Figure:
    """Generate histogram figure for a subset of CP features, AD vs No-AD overlaid."""
    ann = pd.read_csv(ANN_B2_CSV)
    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)
    cp  = pd.read_csv(CP_CSV)
    merged = ann[["filename", "label"]].merge(cp, on="filename")
    merged["binary"] = merged["label"].map(
        lambda x: "adhesion" if x != "No adhesion" else "No adhesion"
    )
    ad   = merged[merged["binary"] == "adhesion"]
    noad = merged[merged["binary"] == "No adhesion"]

    n = len(feat_cols_grp)
    n_cols = min(n, 7)
    n_rows = (n + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                             figsize=(n_cols * 1.9, n_rows * 1.7),
                             facecolor="white")
    axes = np.array(axes).ravel()

    for fi, col in enumerate(feat_cols_grp):
        ax = axes[fi]
        for df, color, label in [(noad, HIST_COL_NOAD, "No ad"), (ad, HIST_COL_AD, "adh")]:
            vals = df[col].dropna().values
            lo, hi = np.percentile(vals, 1), np.percentile(vals, 99)
            ax.hist(np.clip(vals, lo, hi), bins=40, color=color, alpha=0.5,
                    density=True, linewidth=0, label=label)
        short = col.replace("intensity_", "int_")
        for suf in ("_d1", "_d2", "_d4"):
            short = short.replace(suf, suf)
        ax.set_title(short, fontsize=6, pad=2)
        ax.set_yticks([])
        ax.tick_params(axis="x", labelsize=5)
        ax.spines[["top", "right", "left"]].set_visible(False)

    for fi in range(n, len(axes)):
        axes[fi].set_visible(False)

    handles = [mpatches.Patch(color=HIST_COL_AD,   label="adhesion"),
               mpatches.Patch(color=HIST_COL_NOAD, label="No adhesion")]
    fig.legend(handles=handles, fontsize=7, loc="upper right",
               bbox_to_anchor=(0.99, 0.99))
    fig.suptitle(title, fontsize=9, fontweight="bold", y=1.01)
    fig.tight_layout(pad=0.4)
    return fig


def _slide_cp_histograms_intensity(prs):
    cp = pd.read_csv(CP_CSV, nrows=0)
    cols = [c for c in cp.columns if c.startswith("intensity_")]
    fig = _fig_cp_histograms_group(cols, "CellProfiler Intensity Features — adhesion vs No adhesion (DS1 B2)")
    slide = _blank(prs)
    _slide_header(slide, "CellProfiler Intensity Features — Distributions",
                  "11 features · adhesion (dark grey) vs No adhesion (purple) · DS1 vinc B2 annotations")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))


def _slide_cp_histograms_haralick(prs, distance: str = "d1"):
    cp = pd.read_csv(CP_CSV, nrows=0)
    cols = [c for c in cp.columns if c.endswith(f"_{distance}")]
    fig = _fig_cp_histograms_group(
        cols,
        f"CellProfiler Haralick/{distance.upper()} Texture Features — adhesion vs No adhesion",
    )
    slide = _blank(prs)
    dist_px = {"d1": "1 px", "d2": "2 px", "d4": "4 px"}[distance]
    _slide_header(slide,
                  f"CellProfiler Haralick Features (distance={dist_px}) — Distributions",
                  f"13 features · adhesion (dark grey) vs No adhesion (purple) · DS1 vinc B2 annotations")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))


def _fig_il_histograms_group(feat_cols_grp: list, title: str) -> plt.Figure:
    ann = pd.read_csv(ANN_B2_CSV)
    ann["filename"] = ann["filename"].str.replace("-f", "_f", n=1)
    il  = pd.read_csv(IL_CSV)
    merged = ann[["filename", "label"]].merge(il, on="filename")
    merged["binary"] = merged["label"].map(
        lambda x: "adhesion" if x != "No adhesion" else "No adhesion"
    )
    ad   = merged[merged["binary"] == "adhesion"]
    noad = merged[merged["binary"] == "No adhesion"]

    n = len(feat_cols_grp)
    n_cols = min(n, 8)
    n_rows = (n + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                             figsize=(n_cols * 1.7, n_rows * 1.6),
                             facecolor="white")
    axes = np.array(axes).ravel()

    for fi, col in enumerate(feat_cols_grp):
        ax = axes[fi]
        for df, color in [(noad, HIST_COL_NOAD), (ad, HIST_COL_AD)]:
            vals = df[col].dropna().values
            lo, hi = np.percentile(vals, 1), np.percentile(vals, 99)
            ax.hist(np.clip(vals, lo, hi), bins=40, color=color, alpha=0.5,
                    density=True, linewidth=0)
        short = col.replace("gaussian_", "G_").replace("laplacian_of_gaussian_", "LoG_")
        short = short.replace("gaussian_gradient_magnitude_", "Gmag_")
        short = short.replace("difference_of_gaussians_", "DoG_")
        short = short.replace("structure_tensor_eigenvalue_", "ST_")
        short = short.replace("hessian_eigenvalue_", "Hess_")
        ax.set_title(short, fontsize=5.5, pad=2)
        ax.set_yticks([])
        ax.tick_params(axis="x", labelsize=5)
        ax.spines[["top", "right", "left"]].set_visible(False)

    for fi in range(n, len(axes)):
        axes[fi].set_visible(False)

    handles = [mpatches.Patch(color=HIST_COL_AD,   label="adhesion"),
               mpatches.Patch(color=HIST_COL_NOAD, label="No adhesion")]
    fig.legend(handles=handles, fontsize=7, loc="upper right",
               bbox_to_anchor=(0.99, 0.99))
    fig.suptitle(title, fontsize=8, fontweight="bold", y=1.01)
    fig.tight_layout(pad=0.3)
    return fig


def _slide_il_histograms_by_scale(prs, scales: list, label: str):
    il = pd.read_csv(IL_CSV, nrows=0)
    all_il_cols = [c for c in il.columns if c != "filename"]

    # keep cols whose name contains any of the scale strings
    scale_strs = [f"s{s.replace('.', 'p')}" for s in scales]
    cols = [c for c in all_il_cols
            if any(ss in c for ss in scale_strs)]

    fig = _fig_il_histograms_group(cols,
        f"ilastik Features — scales σ={', '.join(scales)} — adhesion vs No adhesion")
    slide = _blank(prs)
    _slide_header(slide,
                  f"ilastik Features (σ={', '.join(scales)} px) — Distributions",
                  f"adhesion (pink) vs No adhesion (grey) · DS1 vinc B2 annotations · {label}")
    _add_fig(slide, fig, Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.3))


# ---------------------------------------------------------------------------
# 2-panel LGBM + logreg figure
# ---------------------------------------------------------------------------

def _fig_le_2panel(lgbm_sets, logreg_sets, budgets,
                   title_lgbm="LGBM", title_logreg="Logreg") -> plt.Figure:
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(9.5, 5.5), facecolor="white")
    x = np.arange(len(budgets))

    for ax, sets, panel_title in [
        (ax1, lgbm_sets,  title_lgbm),
        (ax2, logreg_sets, title_logreg),
    ]:
        for ds in sets:
            if ds is None:
                continue
            mean, std = np.array(ds["mean"], dtype=float), np.array(ds["std"], dtype=float)
            valid = ~np.isnan(mean)
            ls = ds.get("linestyle", "-")
            ax.plot(x[valid], mean[valid] * 100,
                    color=ds["color"], linewidth=2.2,
                    marker="o", markersize=5,
                    label=ds["label"], linestyle=ls)
            ax.fill_between(x[valid],
                            (mean[valid] - std[valid]) * 100,
                            (mean[valid] + std[valid]) * 100,
                            color=ds["color"], alpha=0.12)
        ax.axhline(95, color="#888888", linestyle="--", linewidth=1.2,
                   alpha=0.75, zorder=0)
        ax.text(len(budgets) - 0.08, 95.5, "95 %",
                color="#888888", fontsize=7.5, ha="right", va="bottom")
        ax.set_xticks(x)
        ax.set_xticklabels([str(b) for b in budgets], fontsize=8)
        ax.set_xlabel("Label budget (n patches)", fontsize=10)
        ax.set_ylabel("Balanced accuracy (%)", fontsize=10)
        ax.set_ylim(42, 102)
        ax.legend(fontsize=9, framealpha=0.9, loc="lower right")
        ax.set_facecolor("white")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_title(panel_title, fontsize=10, fontweight="bold", pad=6)
        ax.grid(axis="y", color="#EEEEEE", linewidth=0.7)

    fig.tight_layout(pad=1.5)
    return fig


# ---------------------------------------------------------------------------
# Updated DS1 B2 slide — LGBM + logreg combined
# ---------------------------------------------------------------------------

def _slide_ds1_b2_results(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B2 — Label Efficiency: SupCon-AE vs CP vs ilastik",
                  "vinc ctrl+ycomp · 1,224 patches · lat=12 fixed · 5-fold CV × 5 repeats · left=LGBM, right=logreg")

    sc_lgbm  = _summarize(EVAL_DIR / "supcon_le_b2_lat12p8_ds1.csv",       BUDGETS_B2)
    cp_lgbm  = _summarize(EVAL_DIR / "cp_ds1.csv",                          BUDGETS_B2)
    il_lgbm  = _summarize(EVAL_DIR / "ilastik_ds1.csv",                     BUDGETS_B2)
    sc_lr    = _summarize(EVAL_DIR / "supcon_le_b2_lat12p8_logreg_ds1.csv", BUDGETS_B2)
    cp_lr    = _summarize(EVAL_DIR / "cp_b2_logreg_ds1.csv",                BUDGETS_B2)
    il_lr    = _summarize(EVAL_DIR / "ilastik_b2_logreg_ds1.csv",           BUDGETS_B2)

    def _ds(label, color, s):
        return dict(label=label, color=color,
                    mean=s["mean"].values, std=s["std"].values)

    lgbm_sets  = [_ds("SupCon-AE (lat=12)", COL_SUPCON, sc_lgbm),
                  _ds("CellProfiler",       COL_CP,      cp_lgbm),
                  _ds("ilastik",            COL_ILASTIK, il_lgbm)]
    logreg_sets = [_ds("SupCon-AE (lat=12)", COL_SUPCON, sc_lr),
                   _ds("CellProfiler",       COL_CP,      cp_lr),
                   _ds("ilastik",            COL_ILASTIK, il_lr)]

    fig = _fig_le_2panel(lgbm_sets, logreg_sets, BUDGETS_B2,
                         "DS1 B2 · LGBM", "DS1 B2 · Logreg")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, Inches(9.5), PLOT_H)
    _add_info_panel(slide, _DS1_B2_SECTIONS)


# ---------------------------------------------------------------------------
# Updated DS1 B12 slide — LGBM + logreg combined
# ---------------------------------------------------------------------------

def _slide_ds1_b12_results(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik",
                  "vinc ctrl+ycomp · B1+B2 (~2,428 patches) · 5-fold CV × 5 repeats · SupCon-AE lat=12/proj=8 · left=LGBM, right=logreg")

    sc_lgbm = _summarize(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_ds1.csv",        BUDGETS_B12)
    cp_lgbm = _summarize(EVAL_DIR / "cp_b12_ds1.csv",                           BUDGETS_B12)
    il_lgbm = _summarize(EVAL_DIR / "ilastik_b12_ds1.csv",                      BUDGETS_B12)
    sc_lr   = _summarize(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_logreg_ds1.csv", BUDGETS_B12)
    cp_lr   = _summarize(EVAL_DIR / "cp_b12_logreg_ds1.csv",                    BUDGETS_B12)
    il_lr   = _summarize(EVAL_DIR / "ilastik_b12_logreg_ds1.csv",               BUDGETS_B12)

    def _ds(label, color, s):
        return dict(label=label, color=color,
                    mean=s["mean"].values, std=s["std"].values)

    lgbm_sets  = [_ds("SupCon-AE (lat=12)", COL_SUPCON,  sc_lgbm),
                  _ds("CellProfiler",        COL_CP,      cp_lgbm),
                  _ds("ilastik",             COL_ILASTIK, il_lgbm)]
    logreg_sets = [_ds("SupCon-AE (lat=12)", COL_SUPCON,  sc_lr),
                   _ds("CellProfiler",        COL_CP,      cp_lr),
                   _ds("ilastik",             COL_ILASTIK, il_lr)]

    fig = _fig_le_2panel(lgbm_sets, logreg_sets, BUDGETS_B12,
                         "DS1 B12 · LGBM (lat=12, proj=8)", "DS1 B12 · Logreg (lat=12, proj=8)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, Inches(9.5), PLOT_H)
    _add_info_panel(slide, _DS1_B12_SECTIONS)


# ---------------------------------------------------------------------------
# Updated DS2 B12 slide — LGBM + logreg combined
# ---------------------------------------------------------------------------

def _slide_ds2_b12_results(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS2 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik",
                  "pfak · B1+B2 (244 patches) · 5-fold CV × 5 repeats · SupCon-AE lat=64/proj=32 · left=LGBM, right=logreg")

    sc_lgbm = _summarize(EVAL_DIR / "supcon_le_b12_ds2_lat64p32_ds2.csv",         BUDGETS_B12_DS2)
    cp_lgbm = _summarize(EVAL_DIR / "cp_b12_ds2.csv",                             BUDGETS_B12_DS2)
    il_lgbm = _summarize(EVAL_DIR / "ilastik_b12_ds2.csv",                        BUDGETS_B12_DS2)
    sc_lr   = _summarize(EVAL_DIR / "supcon_le_b12_ds2_lat64p32_logreg_ds2.csv",  BUDGETS_B12_DS2)
    cp_lr   = _summarize(EVAL_DIR / "cp_b12_logreg_ds2.csv",                      BUDGETS_B12_DS2)
    il_lr   = _summarize(EVAL_DIR / "ilastik_b12_logreg_ds2.csv",                 BUDGETS_B12_DS2)

    def _ds(label, color, s):
        return dict(label=label, color=color,
                    mean=s["mean"].values, std=s["std"].values)

    lgbm_sets  = [_ds("SupCon-AE (lat=64)", COL_SUPCON,  sc_lgbm),
                  _ds("CellProfiler",        COL_CP,       cp_lgbm),
                  _ds("ilastik",             COL_ILASTIK,  il_lgbm)]
    logreg_sets = [_ds("SupCon-AE (lat=64)", COL_SUPCON,  sc_lr),
                   _ds("CellProfiler",        COL_CP,       cp_lr),
                   _ds("ilastik",             COL_ILASTIK,  il_lr)]

    fig = _fig_le_2panel(lgbm_sets, logreg_sets, BUDGETS_B12_DS2,
                         "DS2 B12 · LGBM (lat=64, proj=32)", "DS2 B12 · Logreg (lat=64, proj=32)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, Inches(9.5), PLOT_H)
    _add_info_panel(slide, _DS2_B12_SECTIONS)


# ---------------------------------------------------------------------------
# LGBM-only and logreg-only single-panel slides (for main deck / appendix)
# ---------------------------------------------------------------------------

def _slide_ds1_b12_lgbm(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik  [LGBM]",
                  "vinc ctrl+ycomp · B1+B2 (~2,428 patches) · lat=12/proj=8 · LGBM · 5-fold CV × 5 repeats")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_ds1.csv", BUDGETS_B12)
    cp = _summarize(EVAL_DIR / "cp_b12_ds1.csv",                    BUDGETS_B12)
    il = _summarize(EVAL_DIR / "ilastik_b12_ds1.csv",               BUDGETS_B12)

    datasets = [
        dict(label="SupCon-AE (lat=12, proj=8)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(datasets, BUDGETS_B12,
                         "DS1 B12 · Balanced Accuracy vs Label Budget (LGBM)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS1_B12_SECTIONS)


def _slide_ds1_b12_logreg(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS1 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik  [Logreg]",
                  "vinc ctrl+ycomp · B1+B2 (~2,428 patches) · lat=12/proj=8 · logistic regression · 5-fold CV × 5 repeats")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_logreg_ds1.csv", BUDGETS_B12)
    cp = _summarize(EVAL_DIR / "cp_b12_logreg_ds1.csv",                    BUDGETS_B12)
    il = _summarize(EVAL_DIR / "ilastik_b12_logreg_ds1.csv",               BUDGETS_B12)

    datasets = [
        dict(label="SupCon-AE (lat=12, proj=8)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(datasets, BUDGETS_B12,
                         "DS1 B12 · Balanced Accuracy vs Label Budget (Logreg)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS1_B12_SECTIONS)


def _slide_b2_vs_b12_lgbm(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "B2 vs B12 Label Set Comparison — DS1  [LGBM]",
                  "Same protocol, more labels: does SupCon-AE close the gap to CP/ilastik?")

    budgets_shared = [10, 25, 50, 75, 150, 750]

    def _get(path, budgets):
        df  = pd.read_csv(path)
        num = df[df["budget"] != "all"].copy()
        num["budget_int"] = num["budget"].astype(int)
        s   = num.groupby("budget_int")["bal_acc"].agg(["mean", "std"])
        return s.reindex(budgets)

    sc_b2  = _get(EVAL_DIR / "supcon_le_b2_lat12p8_ds1.csv",          budgets_shared)
    cp_b2  = _get(EVAL_DIR / "cp_ds1.csv",                            budgets_shared)
    sc_b12 = _get(EVAL_DIR / "supcon_le_b12_ds1_lat12p8_ds1.csv",     budgets_shared)
    cp_b12 = _get(EVAL_DIR / "cp_b12_ds1.csv",                        budgets_shared)

    fig, axes = plt.subplots(1, 2, figsize=(11, 5.5), facecolor="white")
    x = np.arange(len(budgets_shared))
    xlabels = [str(b) for b in budgets_shared]

    for ax, (sc, cp, title_tag) in zip(axes, [
        (sc_b2,  cp_b2,  "B2 (~1,225 patches)"),
        (sc_b12, cp_b12, "B12 (~2,428 patches)"),
    ]):
        for vals, color, label in [
            (sc, COL_SUPCON, "SupCon-AE"),
            (cp, COL_CP,     "CellProfiler"),
        ]:
            m = vals["mean"].values
            s = vals["std"].values
            valid = ~np.isnan(m)
            ax.plot(x[valid], m[valid]*100, color=color, linewidth=2.2,
                    marker="o", markersize=5, label=label)
            ax.fill_between(x[valid], (m[valid]-s[valid])*100, (m[valid]+s[valid])*100,
                            color=color, alpha=0.12)
        ax.set_xticks(x)
        ax.set_xticklabels(xlabels, fontsize=9)
        ax.set_xlabel("Label budget", fontsize=10)
        ax.set_ylabel("Balanced accuracy (%)", fontsize=10)
        ax.set_ylim(42, 102)
        ax.axhline(95, color="#888888", linestyle="--", linewidth=1.2, alpha=0.75, zorder=0)
        ax.text(len(budgets_shared) - 0.08, 95.5, "95 %",
                color="#888888", fontsize=7.5, ha="right", va="bottom")
        ax.legend(fontsize=9, framealpha=0.9, loc="lower right")
        ax.set_facecolor("white")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_title(f"DS1 {title_tag}\nSupCon-AE vs CP (LGBM)", fontsize=10, fontweight="bold")
        ax.grid(axis="y", color="#EEEEEE", linewidth=0.7)

    fig.tight_layout(pad=1.5)
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _B2_VS_B12_SECTIONS)


def _slide_ds2_b12_lgbm_lat64(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS2 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik  [LGBM]",
                  "pfak · B1+B2 combined (244 patches) · LightGBM · 5-fold CV × 5 repeats · lat=64/proj=32")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_ds2_lat64p32_ds2.csv", BUDGETS_B12_DS2)
    cp = _summarize(EVAL_DIR / "cp_b12_ds2.csv",                     BUDGETS_B12_DS2)
    il = _summarize(EVAL_DIR / "ilastik_b12_ds2.csv",                BUDGETS_B12_DS2)

    datasets = [
        dict(label="SupCon-AE (lat=64, proj=32)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(datasets, BUDGETS_B12_DS2,
                         "DS2 B12 · Balanced Accuracy vs Label Budget (LGBM)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS2_B12_SECTIONS)


def _slide_ds2_b12_logreg(prs):
    slide = _blank(prs)
    _slide_header(slide,
                  "DS2 B12 — Label Efficiency: SupCon-AE vs CP vs ilastik  [Logreg]",
                  "pfak · B1+B2 combined (244 patches) · logistic regression · 5-fold CV × 5 repeats · lat=64/proj=32")

    sc = _summarize(EVAL_DIR / "supcon_le_b12_ds2_lat64p32_logreg_ds2.csv", BUDGETS_B12_DS2)
    cp = _summarize(EVAL_DIR / "cp_b12_logreg_ds2.csv",                      BUDGETS_B12_DS2)
    il = _summarize(EVAL_DIR / "ilastik_b12_logreg_ds2.csv",                 BUDGETS_B12_DS2)

    datasets = [
        dict(label="SupCon-AE (lat=64, proj=32)",
             color=COL_SUPCON,  mean=sc["mean"].values, std=sc["std"].values),
        dict(label="CellProfiler",
             color=COL_CP,      mean=cp["mean"].values, std=cp["std"].values),
        dict(label="ilastik",
             color=COL_ILASTIK, mean=il["mean"].values, std=il["std"].values),
    ]

    fig = _fig_le_curves(datasets, BUDGETS_B12_DS2,
                         "DS2 B12 · Balanced Accuracy vs Label Budget (Logreg)")
    _add_fig(slide, fig, PLOT_LEFT, PLOT_TOP, PLOT_W, PLOT_H)
    _add_info_panel(slide, _DS2_B12_SECTIONS)


def _slide_appendix_divider(prs):
    slide = _blank(prs)
    _slide_header(slide, "Appendix — Logistic Regression Results",
                  "Supplementary: same experiments repeated with logistic regression classifier")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_pptx(out_path: Path):
    print("Building LE benchmark PPT …")
    prs = _prs()

    print("  Slide 1 — Title")
    _slide_title(prs)

    # ---- CellProfiler feature slides ----
    print("  Slide 2 — CP overview")
    _slide_cp_overview(prs)
    print("  Slide 3 — CP intensity features table")
    _slide_cp_intensity(prs)
    print("  Slide 4 — CP Haralick/GLCM features table")
    _slide_cp_haralick(prs)
    print("  Slide 5 — CP top-3 discriminative features (histogram + patches)")
    _slide_cp_top_features(prs)
    print("  Slide 6 — CP histograms: intensity")
    _slide_cp_histograms_intensity(prs)
    print("  Slide 7 — CP histograms: Haralick d1")
    _slide_cp_histograms_haralick(prs, "d1")
    print("  Slide 8 — CP histograms: Haralick d2")
    _slide_cp_histograms_haralick(prs, "d2")
    print("  Slide 9 — CP histograms: Haralick d4")
    _slide_cp_histograms_haralick(prs, "d4")

    # ---- ilastik feature slides ----
    print("  Slide 10 — ilastik overview")
    _slide_ilastik_overview(prs)
    print("  Slide 11 — ilastik filter definitions")
    _slide_ilastik_filters(prs)
    print("  Slide 12 — IL top-3 discriminative features (histogram + patches)")
    _slide_il_top_features(prs)
    print("  Slide 13 — ilastik histograms: small scales (0.3, 0.7)")
    _slide_il_histograms_by_scale(prs, ["0.3", "0.7"], "small scales")
    print("  Slide 14 — ilastik histograms: medium scales (1.0, 1.6)")
    _slide_il_histograms_by_scale(prs, ["1.0", "1.6"], "medium scales")
    print("  Slide 15 — ilastik histograms: large scale (3.5)")
    _slide_il_histograms_by_scale(prs, ["3.5"], "large scale")

    # ---- SupCon-AE top features ----
    print("  Slide 16 — SupCon-AE top latent dims: nb=150")
    _slide_supcon_top_features(prs, budget="150")
    print("  Slide 17 — SupCon-AE top latent dims: nb=750")
    _slide_supcon_top_features(prs, budget="750")

    # ---- UMAP comparison ----
    print("  Slide 18 — UMAP comparison: nb=150")
    _slide_umap_comparison(prs, budget="150")
    print("  Slide 19 — UMAP comparison: nb=750")
    _slide_umap_comparison(prs, budget="750")

    # ---- Benchmark results (LGBM) ----
    print("  Slide 20a — DS1 B2 design")
    _slide_ds1_b2_design(prs)
    print("  Slide 20b — DS1 B2 LGBM")
    _slide_ds1_b2_lgbm(prs)

    print("  Slide 21a — DS1 B12 design")
    _slide_ds1_b12_design(prs)
    print("  Slide 21b — DS1 B12 LGBM")
    _slide_ds1_b12_lgbm(prs)

    print("  Slide 22a — B2 vs B12 design")
    _slide_b2_vs_b12_design(prs)
    print("  Slide 22b — B2 vs B12 LGBM")
    _slide_b2_vs_b12_lgbm(prs)

    print("  Slide 23 — DS2 B12 LGBM")
    _slide_ds2_b12_lgbm_lat64(prs)

    # ---- Appendix: logistic regression results ----
    print("  Appendix divider")
    _slide_appendix_divider(prs)
    print("  Appendix A1 — DS1 B2 Logreg")
    _slide_ds1_b2_logreg(prs)
    print("  Appendix A2 — DS1 B12 Logreg")
    _slide_ds1_b12_logreg(prs)
    print("  Appendix A3 — B2 vs B12 Logreg")
    _slide_b2_vs_b12(prs)
    print("  Appendix A4 — DS2 B12 Logreg")
    _slide_ds2_b12_logreg(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"\nSaved → {out_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(OUT_DIR / "le_benchmark.pptx"))
    args = ap.parse_args()
    build_pptx(Path(args.out))


if __name__ == "__main__":
    main()
