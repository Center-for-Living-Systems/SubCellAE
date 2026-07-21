#!/usr/bin/env python3
"""
make_vinc_ablation_views_pptx.py

One slide per run for every folder starting with 'contrastive_cio_rb_vinc_'.
Each slide shows:
  - Last contrastive_views_ep*.png (left 2/3)
  - Model + training config extracted from the YAML (right 1/3)

Usage:
  python scripts/make_vinc_ablation_views_pptx.py [--out vinc_ablation_views.pptx]
"""

from __future__ import annotations

import argparse
import io
import re
from pathlib import Path

import yaml

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path("/net/projects/CLS/lding/data/fa_data_analysis"
            "/ae_results/contrastive_run")

W_IN, H_IN = 13.33, 7.5

C_DARK   = RGBColor(0x1F, 0x2D, 0x3D)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_LGRAY  = RGBColor(0xAA, 0xAA, 0xAA)
C_LGRAY2 = RGBColor(0xDD, 0xDD, 0xDD)
C_ACCENT = RGBColor(0x2E, 0x86, 0xC1)
C_GREEN  = RGBColor(0x1A, 0x7A, 0x40)
C_ORANGE = RGBColor(0xC0, 0x55, 0x10)


def _px(v): return Inches(v)

def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _tb(slide, text, left, top, w, h,
        size=10, bold=False, italic=False,
        color=C_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(w), _px(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


# ── config parsing ────────────────────────────────────────────────────────────

def _load_yaml(folder: Path) -> dict:
    yamls = sorted(folder.glob("*.yaml"))
    if not yamls:
        return {}
    # prefer shortest name (most likely the primary config)
    yamls.sort(key=lambda p: len(p.name))
    txt = yamls[0].read_text()
    # strip root_folder concatenations so yaml.safe_load doesn't choke
    txt = re.sub(r'root_folder \+ "([^"]*)"', r'"\1"', txt)
    txt = re.sub(r"root_folder \+ '([^']*)'", r"'\1'", txt)
    try:
        return yaml.safe_load(txt) or {}
    except Exception:
        return {}


def _cfg_rows(cfg: dict) -> list[tuple[str, str]]:
    """Extract key model + training fields as (label, value) pairs."""
    rows = []
    m = cfg.get("model", {})
    t = cfg.get("training", {})
    ec = cfg.get("enlarged_crop", {})

    def _v(val, fmt=None):
        if val is None:
            return "—"
        return fmt % val if fmt else str(val)

    # model
    rows.append(("— Model —", ""))
    rows.append(("Type",          _v(m.get("model_type"))))
    rows.append(("Latent / Proj", f"{_v(m.get('latent_dim'))} / {_v(m.get('proj_dim'))}"))
    rows.append(("Channels",      _v(m.get("no_ch"))))
    rows.append(("Recon loss",    _v(m.get("recon_loss_type", "mse")).upper()))
    rows.append(("λ_contrast",    _v(m.get("lambda_contrast"))))
    rows.append(("λ_recon",       _v(m.get("lambda_recon", 1.0))))
    rows.append(("Noise prob",    _v(m.get("noise_prob"))))
    rows.append(("Temp",          _v(m.get("temperature"))))
    rows.append(("BN / Dropout",  f"{m.get('BN_flag',False)} / {m.get('dropout_flag',False)}"))

    # augmentation
    rows.append(("— Augmentation —", ""))
    enlcrop_on = ec.get("enabled", False)
    rows.append(("Enlarged crop", "yes" if enlcrop_on else "no"))
    if enlcrop_on:
        rows.append(("Context size",  _v(ec.get("context_size"))))
        rows.append(("Max shift",     f"±{_v(ec.get('max_shift_px'))} px"))
        rows.append(("Max rotation",  f"±{_v(ec.get('max_angle_deg'))}°"))
        rows.append(("Input divisor", _v(ec.get("input_divisor", 1.0))))

    # training
    rows.append(("— Training —", ""))
    rows.append(("Epochs",        _v(t.get("epochs"))))
    rows.append(("LR",            _v(t.get("lr"))))
    rows.append(("Batch size",    _v(t.get("batch_size"))))
    rows.append(("LR scheduler",  _v(t.get("lr_scheduler"))))
    rows.append(("Warmup epochs", _v(t.get("warmup_epochs", 0))))
    rows.append(("Weight decay",  _v(t.get("weight_decay"))))
    rows.append(("Val split",     _v(t.get("val_split"))))

    return rows


# ── slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs, folders):
    sl = _add_slide(prs)
    _tb(sl, "Vinc ConAE Ablation — Final Epoch Views",
        0.8, 1.8, W_IN - 1.6, 1.0,
        size=32, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    _tb(sl, f"contrastive_cio_rb_vinc_*  ·  {len(folders)} runs",
        0.8, 2.95, W_IN - 1.6, 0.45,
        size=16, color=C_GRAY, align=PP_ALIGN.CENTER)
    _tb(sl, "Last contrastive_views_ep*.png + model/training config per slide",
        0.8, 3.5, W_IN - 1.6, 0.4,
        size=13, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)


def _slide_model(prs, folder: Path, idx: int, total: int):
    sl = _add_slide(prs)
    name = folder.name

    # ── header bar ────────────────────────────────────────────────────────────
    bar = sl.shapes.add_shape(1, _px(0), _px(0), _px(W_IN), _px(0.58))
    bar.fill.background(); bar.line.fill.background()
    _tb(sl, f"[{idx}/{total}]  {name}",
        0.2, 0.05, W_IN - 0.4, 0.36,
        size=14, bold=True, color=C_DARK)

    # ── image (left) ──────────────────────────────────────────────────────────
    img_w = 9.0
    img_top = 0.62
    img_h = H_IN - img_top - 0.05

    views = sorted(folder.glob("contrastive_views_ep*.png"),
                   key=lambda p: int(re.search(r"ep(\d+)", p.name).group(1)))
    if views:
        ep = re.search(r"ep(\d+)", views[-1].name).group(1)
        _tb(sl, f"ep {ep}", 0.2, img_top - 0.01, 1.0, 0.2,
            size=8, italic=True, color=C_LGRAY)
        with open(views[-1], "rb") as f:
            buf = io.BytesIO(f.read())
        sl.shapes.add_picture(buf, _px(0.0), _px(img_top),
                               height=_px(img_h))
    else:
        _tb(sl, "no views PNG found", 0.2, 3.5, img_w, 0.4,
            size=11, color=C_LGRAY, align=PP_ALIGN.CENTER)

    # ── config panel (right) ──────────────────────────────────────────────────
    rx = img_w + 0.15
    rw = W_IN - rx - 0.1
    cfg = _load_yaml(folder)
    rows = _cfg_rows(cfg)

    if not cfg:
        _tb(sl, "no YAML config found", rx, 1.0, rw, 0.4,
            size=9, italic=True, color=C_LGRAY)
        return

    row_h = 0.245
    y = img_top + 0.02
    for label, val in rows:
        if val == "":  # section header
            _tb(sl, label, rx, y, rw, row_h,
                size=8, bold=True, color=C_ACCENT)
        else:
            # label left, value right
            _tb(sl, label, rx, y, rw * 0.55, row_h,
                size=8, color=C_GRAY)
            _tb(sl, val, rx + rw * 0.55, y, rw * 0.45, row_h,
                size=8, bold=True, color=C_DARK)
        y += row_h
        if y > H_IN - 0.15:
            break  # overflow guard


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="vinc_ablation_views.pptx")
    args = parser.parse_args()

    folders = sorted(
        [d for d in BASE.iterdir()
         if d.is_dir() and d.name.startswith("contrastive_cio_rb_vinc_")],
        key=lambda p: p.name,
    )
    print(f"Found {len(folders)} folders")

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    _slide_title(prs, folders)
    for i, folder in enumerate(folders, start=1):
        print(f"  [{i}/{len(folders)}] {folder.name}")
        _slide_model(prs, folder, i, len(folders))

    out = Path(args.out)
    prs.save(str(out))
    print(f"\nSaved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
