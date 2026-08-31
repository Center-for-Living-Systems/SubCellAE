#!/usr/bin/env python3
"""
make_pptx_le_clean_perclass.py
================================
Short 3-slide PPT:
  1. Definitions — precision, recall, balanced accuracy with confusion matrix
  2. Recall curves — balanced acc + adhesion recall + no-adhesion recall
  3. Precision curves — balanced acc + adhesion precision + no-adhesion precision

Usage
-----
  python scripts/make_pptx_le_clean_perclass.py
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
REPO    = Path(__file__).resolve().parents[1]
RESULTS = REPO / "results"
OUT     = RESULTS / "le_clean_perclass_metrics.pptx"

SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.15)
BODY_T  = TITLE_H + Inches(0.08)
BODY_H  = SH - BODY_T - PAD

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_GREEN = RGBColor(0x2C, 0xA0, 0x2C)
C_RED   = RGBColor(0xD6, 0x27, 0x28)
C_BLUE  = RGBColor(0x1F, 0x77, 0xB4)
C_PURP  = RGBColor(0x94, 0x67, 0xBD)

# ── helpers ────────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None, line_color=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2*PAD, Inches(0.25),
             subtitle, size=9, color=C_GREY)

def _place_png(slide, path, l, t, max_w, max_h):
    p = Path(path)
    if not p.exists():
        _txt(slide, l, t, max_w, Inches(0.3), f"[missing: {p.name}]",
             size=9, color=C_GREY, align=PP_ALIGN.CENTER)
        return
    pil = Image.open(str(p)).convert("RGB")
    iw, ih = pil.size
    scale  = min(max_w / Inches(iw / 150), max_h / Inches(ih / 150))
    rw, rh = Inches(iw / 150) * scale, Inches(ih / 150) * scale
    buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                             t + (max_h - rh) / 2, rw, rh)

def _fig_to_slide(slide, fig, l, t, max_w, max_h):
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    pil = Image.open(buf).convert("RGB")
    iw, ih = pil.size
    scale  = min(max_w / Inches(iw / 150), max_h / Inches(ih / 150))
    rw, rh = Inches(iw / 150) * scale, Inches(ih / 150) * scale
    buf2 = io.BytesIO(); pil.save(buf2, format="PNG"); buf2.seek(0)
    slide.shapes.add_picture(buf2, l + (max_w - rw) / 2,
                             t + (max_h - rh) / 2, rw, rh)


# ── helpers for editable formula cards ────────────────────────────────────────

def _hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _formula_card(slide, l, t, w, h, title, formula, desc, color_hex):
    """Editable formula card: coloured header bar + native text boxes."""
    hdr_h  = Inches(0.30)
    pad_i  = Inches(0.10)
    color  = _hex_to_rgb(color_hex)
    light  = RGBColor(0xF5, 0xF5, 0xF5)

    # Header rectangle + title
    _rect(slide, l, t, w, hdr_h, fill=color)
    _txt(slide, l + pad_i, t + Inches(0.03), w - 2*pad_i, hdr_h - Inches(0.03),
         title, size=10, bold=True, color=C_WHITE)

    # Body background
    body_h = h - hdr_h - Inches(0.04)
    _rect(slide, l, t + hdr_h, w, body_h, fill=light)

    # Formula line (bold, coloured)
    _txt(slide, l + pad_i, t + hdr_h + Inches(0.04),
         w - 2*pad_i, Inches(0.26),
         formula, size=10, bold=True, color=color)

    # Description
    _txt(slide, l + pad_i, t + hdr_h + Inches(0.30),
         w - 2*pad_i, body_h - Inches(0.32),
         desc, size=9, color=C_BLACK, wrap=True)


# ── slide 1: definitions ───────────────────────────────────────────────────────

def slide_definitions(prs):
    sl = _blank(prs)
    _title_bar(sl, "Metric Definitions — Precision, Recall, Balanced Accuracy",
               "Binary classifier: adhesion (positive) vs no-adhesion (negative)")

    LEFT_W  = Inches(7.2)
    RIGHT_X = LEFT_W + Inches(0.20)
    RIGHT_W = SW - RIGHT_X - PAD

    # ── LEFT: confusion matrix (matplotlib figure) ─────────────────────────
    fig, ax = plt.subplots(figsize=(7, 6), facecolor="white")
    ax.set_xlim(0, 4)
    ax.set_ylim(0, 4)
    ax.set_aspect("equal")
    ax.axis("off")

    CELL_COLORS = {"TP": "#c6efce", "FN": "#ffc7ce",
                   "FP": "#ffc7ce", "TN": "#c6efce"}
    cells = [
        (1, 2, "TP\n(True Positive)",  CELL_COLORS["TP"], "e.g. 281"),
        (2, 2, "FN\n(False Negative)", CELL_COLORS["FN"], "e.g. 58"),
        (1, 1, "FP\n(False Positive)", CELL_COLORS["FP"], "e.g. 254"),
        (2, 1, "TN\n(True Negative)",  CELL_COLORS["TN"], "e.g. 339"),
    ]
    for col, row, lbl, color, ex in cells:
        ax.add_patch(plt.Rectangle((col, row), 1, 1, facecolor=color,
                                   edgecolor="gray", linewidth=1.5))
        ax.text(col + 0.5, row + 0.68, lbl, ha="center", va="center",
                fontsize=9, fontweight="bold", linespacing=1.5)
        ax.text(col + 0.5, row + 0.22, ex, ha="center", va="center",
                fontsize=8, color="#555555")

    ax.text(1.5, 3.65, "Predicted: Adhesion", ha="center", fontsize=9,
            fontweight="bold", color="#2ca02c")
    ax.text(2.5, 3.65, "Predicted: No-adh",   ha="center", fontsize=9,
            fontweight="bold", color="#1f77b4")
    ax.text(0.48, 2.5, "True:\nAdhesion", ha="center", va="center",
            fontsize=9, fontweight="bold", color="#2ca02c")
    ax.text(0.48, 1.5, "True:\nNo-adh",   ha="center", va="center",
            fontsize=9, fontweight="bold", color="#1f77b4")
    ax.text(2.0, 0.68, "Confusion matrix (ycomp zero-shot, example values)",
            ha="center", fontsize=7, color="#888888", style="italic")

    ax.annotate("", xy=(1.0, 2.0), xytext=(1.0, 3.0),
                arrowprops=dict(arrowstyle="<->", color="#d62728", lw=1.8))
    ax.text(0.82, 2.5, "Recall\ncolumn", ha="center", fontsize=7,
            color="#d62728", rotation=90)
    ax.annotate("", xy=(3.0, 2.0), xytext=(3.0, 1.0),
                arrowprops=dict(arrowstyle="<->", color="#9467bd", lw=1.8))
    ax.text(3.22, 1.5, "Precision\nrow", ha="center", fontsize=7,
            color="#9467bd", rotation=90)

    ax.set_title("Confusion Matrix", fontsize=11, fontweight="bold", pad=8)
    fig.tight_layout(pad=0.5)
    _fig_to_slide(sl, fig, PAD, BODY_T + Inches(0.05),
                  LEFT_W - PAD, BODY_H - Inches(0.1))
    plt.close(fig)

    # ── RIGHT: editable formula cards (native pptx) ────────────────────────
    defs = [
        ("#2ca02c", "Adhesion Recall  (Sensitivity)",
         "= TP / (TP + FN)",
         "Of all true adhesion patches, what fraction did the model correctly "
         "identify?  →  Low recall = model misses real adhesions"),
        ("#1f77b4", "No-adh Recall  (Specificity)",
         "= TN / (TN + FP)",
         "Of all true no-adhesion patches, what fraction did the model correctly "
         "reject?  →  Low specificity = too many false alarms"),
        ("#d62728", "Adhesion Precision  (PPV)",
         "= TP / (TP + FP)",
         "Of all patches predicted as adhesion, what fraction are truly adhesion?"
         "  →  Low precision = many false positives called as adhesion"),
        ("#9467bd", "No-adh Precision  (NPV)",
         "= TN / (TN + FN)",
         "Of all patches predicted as no-adhesion, what fraction are truly "
         "no-adhesion?  →  Low NPV = missed adhesions inflate the negative bucket"),
        ("#333333", "Balanced Accuracy  (BAcc)",
         "= (Adhesion Recall + No-adh Recall) / 2",
         "Average recall across both classes.  Robust to class imbalance — "
         "chance level = 50%."),
    ]

    card_h   = BODY_H / len(defs) - Inches(0.04)
    card_gap = Inches(0.05)
    for i, (color, title, formula, desc) in enumerate(defs):
        cy = BODY_T + Inches(0.05) + i * (card_h + card_gap)
        _formula_card(sl, RIGHT_X, cy, RIGHT_W, card_h,
                      title, formula, desc, color)


# ── slide 2: recall curves ─────────────────────────────────────────────────────

def slide_recall(prs):
    sl = _blank(prs)
    _title_bar(sl, "Label Efficiency — Recall Curves  (adhesion vs no-adhesion)",
               "SupCon AE s3v1 + LGBM  ·  vinc/control  ·  image-held-out  ·  3 repeats  ·  "
               "Adhesion recall = sensitivity  ·  No-adhesion recall = specificity")
    _place_png(sl, RESULTS / "le_clean_curve.png",
               PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.05))

    by = SH - Inches(1.05)
    bullets = [
        ("cfg0 npi=10: adhesion recall only 67% — model defaults to 'no adhesion' when data-starved; "
         "recovers to 98% at npi=75", C_GREEN),
        ("No-adhesion recall stays high (82–96%) across all label counts — the 'no adhesion' class is "
         "easier to learn (homogeneous background appearance)", C_BLUE),
        ("cfg2 (3 train images): both curves track closely from npi=10 — more images stabilises the "
         "decision boundary without needing many labels per image", C_BLACK),
    ]
    lh = Inches(0.28)
    y = by
    for text, color in bullets:
        _txt(sl, PAD, y, SW - 2*PAD, lh,
             f"• {text}", size=9, color=color, wrap=True)
        y += lh


# ── slide 3: precision curves ──────────────────────────────────────────────────

def slide_precision(prs):
    sl = _blank(prs)
    _title_bar(sl, "Label Efficiency — Precision Curves  (adhesion vs no-adhesion)",
               "SupCon AE s3v1 + LGBM  ·  vinc/control  ·  image-held-out  ·  3 repeats  ·  "
               "Adhesion precision = PPV  ·  No-adhesion precision = NPV")
    _place_png(sl, RESULTS / "le_clean_curve_precision.png",
               PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.05))

    by = SH - Inches(1.05)
    bullets = [
        ("Adhesion precision is consistently 10–20 pp below no-adhesion precision at low label counts — "
         "the model produces many false positive adhesions even when recall is already high", C_RED),
        ("The gap narrows with more labels: at npi=all (full annotation), adhesion precision reaches "
         "93–95%, close to no-adhesion precision of 97–99%", C_BLACK),
        ("Implication: the decision boundary is biased toward 'adhesion' — diverse adhesion patch "
         "appearance makes the boundary fuzzy, pulling uncertain patches to the positive side", C_PURP),
    ]
    lh = Inches(0.28)
    y = by
    for text, color in bullets:
        _txt(sl, PAD, y, SW - 2*PAD, lh,
             f"• {text}", size=9, color=color, wrap=True)
        y += lh


# ── main ───────────────────────────────────────────────────────────────────────

def slide_trainvaltest(prs):
    sl = _blank(prs)
    _title_bar(sl, "Label Efficiency — Train / Val / Test Balanced Accuracy",
               "Train = the selected npi labels  ·  Val = remaining labeled patches from "
               "training frame(s) not used for training  ·  Test = held-out frames")
    _place_png(sl, RESULTS / "le_clean_curve_trainvaltest.png",
               PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.15))

    by = SH - Inches(1.15)
    bullets = [
        ("Train BAcc = 100% at all npi — the GBM perfectly memorises the small training set "
         "(10–all patches), so the training curve is uninformative about generalisation", C_RED),
        ("Val BAcc (remaining labels from the same training frame) tracks test BAcc within "
         "~5–10 pp and rises consistently with npi — a reliable in-image sanity check", C_BLACK),
        ("Val and test converge at high npi (≥75): with enough labels the model generalises "
         "well both within and across frames, confirming the latent space is consistent", C_GREEN),
        ("Val is unavailable for npi=all because all frame labels are used for training — "
         "points are omitted from the val curve at that setting", C_GREY),
    ]
    lh = Inches(0.26)
    y = by
    for text, color in bullets:
        _txt(sl, PAD, y, SW - 2*PAD, lh, f"• {text}", size=9, color=color, wrap=True)
        y += lh


def slide_val_vs_test(prs):
    sl = _blank(prs)
    _title_bar(sl, "Label Efficiency — Val vs Test: Per-class Recall & Precision Breakdown",
               "Solid = test frames (held-out)  ·  Dashed = val (remaining labeled patches "
               "from training frame, same color)  ·  Val omitted at npi=all")
    _place_png(sl, RESULTS / "le_clean_curve_val_vs_test.png",
               PAD, BODY_T + Inches(0.05), SW - 2*PAD, BODY_H - Inches(1.15))

    by = SH - Inches(1.15)
    bullets = [
        ("Val BAcc < Test BAcc across most settings because frame 0 contains harder adhesion "
         "patches — the selected npi labels are a random draw, so remaining val patches include "
         "more borderline cases", C_BLACK),
        ("Adhesion recall is the main gap: val adh recall (dashed green) lags test adh recall "
         "(solid green), especially at low npi — the classifier under-calls adhesion on the "
         "ambiguous frame-0 patches", C_GREEN),
        ("No-adhesion recall (blue) and both precision curves (red/purple) align much more "
         "closely between val and test, confirming that no-adhesion patches are consistent "
         "across frames", C_MID),
    ]
    lh = Inches(0.28)
    y = by
    for text, color in bullets:
        _txt(sl, PAD, y, SW - 2*PAD, lh, f"• {text}", size=9, color=color, wrap=True)
        y += lh


def main():
    RESULTS.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("[1/5] Definitions slide ...")
    slide_definitions(prs)
    print("[2/5] Recall curves slide ...")
    slide_recall(prs)
    print("[3/5] Precision curves slide ...")
    slide_precision(prs)
    print("[4/5] Train / Val / Test slide ...")
    slide_trainvaltest(prs)
    print("[5/5] Val vs Test breakdown slide ...")
    slide_val_vs_test(prs)

    prs.save(str(OUT))
    print(f"\n[done] {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
