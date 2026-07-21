#!/usr/bin/env python3
"""
make_recon_quality_4ds_pptx.py

Reconstruction-quality panel PPT for a model trained on all 4 datasets.

Reads the training recon TIFs from <model_dir>/recon/, computes normalised L1
(nL1 = L1 / mean|raw|) for every training patch globally, then for each
10th-percentile band (10 – 90 %) generates a 6-column × 5-pair panel showing
patches sampled near that quality level.  Each patch is labelled:
    dsN_cond_split   e.g.  ds1_ctrl_tr  or  ds4_yc_val

Output
------
  <model_dir>/quality_panels_4ds_nl1/nl1_NNp.png   (9 panel images)
  reconstruction_quality_overview_4ds.pptx           (11 slides)

Usage
-----
  python scripts/make_recon_quality_4ds_pptx.py <model_dir>
  python scripts/make_recon_quality_4ds_pptx.py <model_dir> --out foo.pptx
"""
from __future__ import annotations

import argparse
import io
import sys
from collections import Counter
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import tifffile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── constants ──────────────────────────────────────────────────────────────────

PERCENTILES = list(range(10, 100, 10))   # 10, 20, ..., 90
N_PATCHES   = 30
COLS        = 6
PCT_WINDOW  = 2    # ± 2 pct points around each target

_DS_MAP    = {"vinc": "ds1", "pfak": "ds2", "ppax": "ds3", "nih3t3": "ds4"}
_COND_MAP  = {"control": "ctrl", "ycomp": "yc"}
_SPLIT_MAP = {"train": "tr", "val": "val"}


def _patch_label(condition_name: str, split: str) -> str:
    """'vinc_control', 'train'  →  'ds1_ctrl_tr'"""
    parts = str(condition_name).split("_")
    ds   = _DS_MAP.get(parts[0], parts[0])
    cond = _COND_MAP.get(parts[1], parts[1]) if len(parts) > 1 else ""
    spl  = _SPLIT_MAP.get(str(split), str(split))
    return f"{ds}_{cond}_{spl}" if cond else f"{ds}_{spl}"


# ── panel figure ──────────────────────────────────────────────────────────────

def _panel_png(raw_patches: list, recon_patches: list,
               labels: list[str], title: str) -> bytes:
    """6-col panel of raw/recon pairs; return PNG bytes."""
    n_show      = min(N_PATCHES, len(raw_patches))
    n_row_pairs = (n_show + COLS - 1) // COLS
    fig_rows    = n_row_pairs * 2

    fig, axes = plt.subplots(fig_rows, COLS,
                             figsize=(COLS * 1.25, fig_rows * 1.25),
                             facecolor="white")
    axes = np.array(axes).reshape(fig_rows, COLS)
    for ax in axes.flat:
        ax.axis("off")

    for idx in range(n_show):
        pr = idx // COLS
        pc = idx  % COLS
        raw  = raw_patches[idx]
        recon = recon_patches[idx]
        if raw.ndim == 3 and raw.shape[0] == 1:
            raw = raw[0]; recon = recon[0]
        vmin = float(min(raw.min(), recon.min()))
        vmax = float(max(raw.max(), recon.max()))
        if vmax <= vmin:
            vmax = vmin + 1e-6
        axes[pr * 2,     pc].imshow(raw,   cmap="gray", vmin=vmin, vmax=vmax)
        axes[pr * 2,     pc].set_title(labels[idx], fontsize=5, pad=1.5)
        axes[pr * 2 + 1, pc].imshow(recon, cmap="gray", vmin=vmin, vmax=vmax)

    # row labels: "raw" / "recon" on first column
    for pr in range(n_row_pairs):
        for row_off, tag in ((0, "raw"), (1, "recon")):
            ax = axes[pr * 2 + row_off, 0]
            ax.axis("on")
            ax.set_xticks([]); ax.set_yticks([])
            for sp in ax.spines.values():
                sp.set_visible(False)
            ax.set_ylabel(tag, fontsize=7, fontweight="bold", labelpad=2)

    fig.suptitle(title, fontsize=8, y=0.998)
    fig.tight_layout(rect=[0, 0, 1, 0.988], pad=0.3)
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


# ── distribution overview figure ─────────────────────────────────────────────

def _dist_png(nl1: np.ndarray, idx_df: pd.DataFrame,
              pct_vals: np.ndarray) -> bytes:
    """Violin of nL1 per dataset × split, with global percentile lines."""
    datasets = ["vinc", "pfak", "ppax", "nih3t3"]
    ds_ids   = ["ds1",  "ds2",  "ds3",  "ds4"]

    fig, axes = plt.subplots(1, 4, figsize=(14, 4.5), sharey=True, facecolor="white")
    y_max = float(np.percentile(nl1, 99.5)) * 1.1

    for ax, ds, ds_id in zip(axes, datasets, ds_ids):
        mask_ds = idx_df["condition_name"].str.startswith(ds)
        positions, colors, split_labels = [], [], []
        for pos, (split, color, lbl) in enumerate(
                [("train", "#4C72B0", "tr"), ("val", "#DD8452", "val")]):
            vals = nl1[mask_ds & (idx_df["split"] == split)]
            if len(vals) == 0:
                continue
            parts = ax.violinplot([vals], positions=[pos], widths=0.7,
                                  showmedians=True, showextrema=False)
            for pc_body in parts["bodies"]:
                pc_body.set_facecolor(color); pc_body.set_alpha(0.65)
            parts["cmedians"].set_color("black"); parts["cmedians"].set_linewidth(1.5)
            positions.append(pos); colors.append(color); split_labels.append(lbl)

        for pv in pct_vals:
            ax.axhline(pv, color="gray", linewidth=0.4, alpha=0.6, linestyle="--")

        ax.set_title(f"{ds_id}  ({ds})\nn={int(mask_ds.sum())}", fontsize=10)
        ax.set_xticks([0, 1]); ax.set_xticklabels(["tr", "val"], fontsize=9)
        ax.set_ylim(0, y_max)

    axes[0].set_ylabel("nL1", fontsize=10)
    fig.suptitle(
        "Normalised L1 by dataset & split  —  dashed = global 10th–90th pct thresholds",
        fontsize=11)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


# ── PPT helpers ───────────────────────────────────────────────────────────────

W_IN, H_IN = 13.33, 7.5
MARGIN     = 0.2
HEADER_H   = 0.68

C_NAVY  = RGBColor(0x1F, 0x2D, 0x3D)
C_BLUE  = RGBColor(0x2E, 0x86, 0xC1)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x66, 0x66, 0x66)
C_LTBLUE = RGBColor(0xBB, 0xCC, 0xEE)
C_PALE   = RGBColor(0xCC, 0xDD, 0xFF)


def _px(in_: float):
    return Inches(in_)


def _add_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, text, left, top, width, height,
             font_size=11, bold=False, color=C_NAVY,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb


def _header(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(1, _px(0), _px(0), _px(W_IN), _px(0.62))
    bar.fill.background()
    bar.line.fill.background()
    _textbox(slide, title, 0.25, 0.04, W_IN - 0.5, 0.46,
             font_size=22, bold=True, color=C_NAVY)
    if subtitle:
        _textbox(slide, subtitle, 0.25, 0.47, W_IN - 0.5, 0.22,
                 font_size=10, color=C_GRAY)


def _embed(slide, png_bytes: bytes, left, top, width=None, height=None):
    buf = io.BytesIO(png_bytes)
    return slide.shapes.add_picture(
        buf, _px(left), _px(top),
        width=_px(width)  if width  is not None else None,
        height=_px(height) if height is not None else None,
    )


# ── main logic ────────────────────────────────────────────────────────────────

def run(model_dir: Path, out_path: Path, seed: int = 42) -> None:
    recon_dir = model_dir / "recon"
    for f in ("patches_raw.tif", "patches_recon.tif", "patches_index.csv"):
        if not (recon_dir / f).exists():
            sys.exit(f"Required file missing: {recon_dir / f}")

    out_panels = model_dir / "quality_panels_4ds_nl1"
    out_panels.mkdir(exist_ok=True)
    rng = np.random.default_rng(seed)

    # ── load TIFs ─────────────────────────────────────────────────────────────
    print("Loading patches_raw.tif …", flush=True)
    raw_all   = tifffile.imread(str(recon_dir / "patches_raw.tif"))
    print("Loading patches_recon.tif …", flush=True)
    recon_all = tifffile.imread(str(recon_dir / "patches_recon.tif"))
    idx_df    = pd.read_csv(recon_dir / "patches_index.csv")

    assert len(raw_all) == len(idx_df), \
        f"TIF rows {len(raw_all)} ≠ index rows {len(idx_df)}"
    print(f"  {len(raw_all)} patches  |  "
          f"datasets: {sorted(idx_df['condition_name'].str.split('_').str[0].unique())}",
          flush=True)

    # ── compute nL1 ───────────────────────────────────────────────────────────
    print("Computing nL1 …", flush=True)
    r = raw_all.astype(np.float32)
    p = recon_all.astype(np.float32)
    if r.ndim == 4 and r.shape[1] == 1:
        r = r[:, 0]; p = p[:, 0]
    nl1      = np.abs(r - p).mean(axis=(1, 2)) / (np.abs(r).mean(axis=(1, 2)) + 1e-8)
    pct_vals = np.percentile(nl1, PERCENTILES)
    print(f"  nL1: min={nl1.min():.4f}  med={np.median(nl1):.4f}  "
          f"max={nl1.max():.4f}", flush=True)
    print(f"  Pct values: " +
          "  ".join(f"{P}%={v:.3f}" for P, v in zip(PERCENTILES, pct_vals)),
          flush=True)

    # per-patch labels: "ds1_ctrl_tr", "ds4_yc_val" …
    patch_labels = [
        _patch_label(row["condition_name"], row["split"])
        for _, row in idx_df.iterrows()
    ]

    # ── generate panels ───────────────────────────────────────────────────────
    panel_info: list[tuple] = []   # (P, pct_v, lo, hi, n_draw, n_win, png_bytes)

    for P, pct_v in zip(PERCENTILES, pct_vals):
        lo  = float(np.percentile(nl1, max(P - PCT_WINDOW, 0)))
        hi  = float(np.percentile(nl1, min(P + PCT_WINDOW, 100)))
        win = np.where((nl1 >= lo) & (nl1 <= hi))[0]
        if len(win) == 0:
            print(f"  {P}p: no patches in window — skip", flush=True)
            continue

        n_draw = min(N_PATCHES, len(win))
        chosen = rng.choice(win, size=n_draw, replace=False)
        chosen = chosen[np.argsort(nl1[chosen])]   # sorted ascending by nL1

        chosen_labels = [patch_labels[i] for i in chosen]

        # breakdown for title (ds+cond part only, strip split)
        cnt = Counter("_".join(lbl.split("_")[:2]) for lbl in chosen_labels)
        cnt_str = "  ".join(f"{k}:{v}" for k, v in sorted(cnt.items()))

        title = (f"nL1  {P:2d}th pct = {pct_v:.3f}"
                 f"  [{lo:.3f}–{hi:.3f}]"
                 f"  n={n_draw}/{len(win)}\n{cnt_str}")

        png = _panel_png(list(r[chosen]), list(p[chosen]), chosen_labels, title)
        (out_panels / f"nl1_{P:02d}p.png").write_bytes(png)
        panel_info.append((P, pct_v, lo, hi, n_draw, len(win), png))
        print(f"  {P:2d}p: nL1={pct_v:.3f}  n={n_draw}/{len(win):5d}  {cnt_str}",
              flush=True)

    # ── distribution overview figure ─────────────────────────────────────────
    print("Generating distribution figure …", flush=True)
    dist_bytes = _dist_png(nl1, idx_df, pct_vals)

    # ── build PPT ─────────────────────────────────────────────────────────────
    print("Building PPT …", flush=True)
    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    # Slide 1 — title
    sl = _add_slide(prs)
    _textbox(sl, "Reconstruction Quality Overview",
             1.0, 1.6, W_IN - 2.0, 1.1,
             font_size=38, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    _textbox(sl,
             "ds1+ds2+ds3+ds4  ·  ConAE enlcrop / sc2 / nL1 / λ=0.25  ·  "
             "Normalised L1 global percentiles",
             1.0, 2.9, W_IN - 2.0, 0.55,
             font_size=17, color=C_GRAY,
             align=PP_ALIGN.CENTER)
    bullets = [
        f"Model dir: {model_dir.name}",
        f"Training patches: {len(raw_all):,}  "
        f"(ds1 {int((idx_df['condition_name'].str.startswith('vinc')).sum()):,}  "
        f"ds2 {int((idx_df['condition_name'].str.startswith('pfak')).sum()):,}  "
        f"ds3 {int((idx_df['condition_name'].str.startswith('ppax')).sum()):,}  "
        f"ds4 {int((idx_df['condition_name'].str.startswith('nih3t3')).sum()):,})",
        "Metric: nL1 = L1 / mean|raw|  (lower = better reconstruction)",
        "9 panels: 10th – 90th pct  ·  30 patches each  ·  sorted low→high nL1",
        "Patch label: dsN_cond_split   e.g.  ds1_ctrl_tr  ·  ds4_yc_val",
    ]
    y = 3.65
    for b in bullets:
        _textbox(sl, f"• {b}", 1.6, y, W_IN - 3.2, 0.42,
                 font_size=13, color=C_GRAY)
        y += 0.44

    # Slide 2 — distribution
    sl = _add_slide(prs)
    _header(sl,
            "nL1 Distribution — all training patches by dataset & split",
            subtitle="blue = train  ·  orange = val  ·  "
                     "dashed lines = global 10th–90th percentile thresholds  ·  "
                     "labels: ds1=vinc  ds2=pfak  ds3=ppax  ds4=nih3t3")
    _embed(sl, dist_bytes, MARGIN, HEADER_H + 0.05, width=W_IN - 2 * MARGIN)

    # Slides 3–11 — one per percentile
    for P, pct_v, lo, hi, n_draw, n_win, png in panel_info:
        sl = _add_slide(prs)
        _header(sl,
                f"nL1 — {P:2d}th percentile  =  {pct_v:.3f}",
                subtitle=(f"window [{lo:.3f} – {hi:.3f}]  ·  "
                          f"n = {n_draw} sampled / {n_win} in window  ·  "
                          "raw (top) / recon (bottom)  ·  "
                          "label: dsN_cond_split"))
        _embed(sl, png, MARGIN, HEADER_H + 0.03, width=W_IN - 2 * MARGIN)

    prs.save(str(out_path))
    print(f"\nSaved → {out_path}  ({len(prs.slides)} slides)")
    print(f"Panels → {out_panels}/")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("model_dir", type=Path)
    parser.add_argument("--out",  type=Path, default=None)
    parser.add_argument("--seed", type=int,  default=42)
    args = parser.parse_args()

    if not args.model_dir.is_dir():
        sys.exit(f"Not a directory: {args.model_dir}")

    out = args.out or Path("reconstruction_quality_overview_4ds.pptx")
    run(args.model_dir, out, seed=args.seed)


if __name__ == "__main__":
    main()
