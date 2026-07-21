#!/usr/bin/env python3
"""
make_multichannel_ppt.py

PPT comparing multi-channel ConAE results:
  - 3ch pza: paxillin + zyxin + actin (conae_3ch_pza_v_nl1)
  - 4ch vinc: vinculin + paxillin + zyxin + actin (conae_4ch_vinc_v_nl1)

Shows: reconstruction samples, UMAP (FA annotation + KMeans k=10),
       cluster panels (all channels), cross-dataset violin plots.

Output: slides_multichannel.pptx  (in the current working directory)

Usage:
    python scripts/make_multichannel_ppt.py [--out slides_multichannel.pptx]
"""
import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── constants ──────────────────────────────────────────────────────────────────
SWEEP = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/protein_sweep")

MODEL_3CH = SWEEP / "conae_3ch_pza_v_nl1"
MODEL_4CH = SWEEP / "conae_4ch_vinc_v_nl1"

LABEL_3CH = "ConAE 3ch\n(pax+zyx+act)"
LABEL_4CH = "ConAE 4ch\n(vinc+pax+zyx+act)"

W_IN, H_IN = 13.33, 7.5
MARGIN     = 0.25
IMG_W      = 9.6
TEXT_X     = IMG_W + MARGIN + 0.1
TEXT_W     = W_IN - TEXT_X - MARGIN

C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE = RGBColor(0x1F, 0x2D, 0x3D)
C_3CH   = RGBColor(0x62, 0x39, 0xA8)   # purple for 3ch
C_4CH   = RGBColor(0xD4, 0x6C, 0x08)   # orange for 4ch
C_GRAY  = RGBColor(0x66, 0x66, 0x66)
C_LGRAY = RGBColor(0xAA, 0xAA, 0xAA)


# ── pptx helpers ──────────────────────────────────────────────────────────────

def _px(inches):
    return Inches(inches)


def _add_textbox(slide, text, left, top, width, height,
                 font_size=11, bold=False, color=C_TITLE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, bullets, left, top, width, height,
                 font_size=10, color=C_TITLE, head=None, head_color=None):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if head:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = head
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = head_color or color
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if (head or i > 0) else tf.paragraphs[0]
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def _add_image_bytes(slide, img_bytes, left, top, width=None, height=None):
    return slide.shapes.add_picture(
        io.BytesIO(img_bytes), _px(left), _px(top),
        width=_px(width) if width else None,
        height=_px(height) if height else None,
    )


def _fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_header(slide, title, color=C_TITLE, font_size=18):
    _add_textbox(slide, title, MARGIN, 0.08, W_IN - 2*MARGIN, 0.5,
                 font_size=font_size, bold=True, color=color)
    rule = slide.shapes.add_shape(
        1, _px(MARGIN), _px(0.60), _px(W_IN - 2*MARGIN), _px(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def _note_panel(slide, head, bullets, color, font_size=10):
    bg = slide.shapes.add_shape(
        1, _px(TEXT_X - 0.1), _px(0.65), _px(TEXT_W + 0.1), _px(H_IN - 0.75))
    bg.fill.solid()
    r = int(0.95 * 255 + 0.05 * color[0])
    g = int(0.95 * 255 + 0.05 * color[1])
    b = int(0.95 * 255 + 0.05 * color[2])
    bg.fill.fore_color.rgb = RGBColor(r, g, b)
    bg.line.fill.background()
    _add_bullets(slide, bullets, TEXT_X, 0.72, TEXT_W, H_IN - 0.85,
                 font_size=font_size, color=C_TITLE,
                 head=head, head_color=color)


def _content_slide(prs, title, color, img_bytes, note_head, note_bullets,
                   note_font=12):
    slide = _add_slide(prs)
    _slide_header(slide, title, color=color)
    _note_panel(slide, note_head, note_bullets, color, font_size=note_font)
    _add_image_bytes(slide, img_bytes, MARGIN, 0.68, width=IMG_W)
    return slide


# ── image loaders ─────────────────────────────────────────────────────────────

def _load_visual_frames(model_dir: Path, n_frames: int = 2) -> list:
    """Load frames from visual.tif; supports (N,H,W,3) RGB and grayscale."""
    vis_path = model_dir / "recon" / "visual.tif"
    if not vis_path.exists():
        return []
    stack = tifffile.imread(str(vis_path))
    if stack.ndim == 3:
        stack = stack[np.newaxis]
    total = len(stack)
    idxs = np.linspace(5, total - 5, n_frames, dtype=int)
    frames = []
    for i in idxs:
        arr = stack[i]
        if arr.ndim == 3 and arr.shape[-1] == 3:
            frames.append(arr.astype(np.uint8))
        elif arr.ndim == 2 or (arr.ndim == 3 and arr.shape[-1] == 1):
            if arr.ndim == 3:
                arr = arr[..., 0]
            lo, hi = arr.min(), arr.max()
            arr8 = np.clip((arr - lo) / (hi - lo + 1e-8) * 255, 0, 255).astype(np.uint8)
            frames.append(np.stack([arr8] * 3, axis=-1))
        else:
            frames.append(arr.astype(np.uint8))
    return frames


def _load_cluster_grid(model_dir: Path) -> np.ndarray | None:
    """Arrange all_clusters.tif (N, H, W) as a 2×5 grid."""
    clust_path = model_dir / "eval" / "cluster_panels" / "all_clusters.tif"
    if not clust_path.exists():
        return None
    stack = tifffile.imread(str(clust_path))
    n, ph, pw = stack.shape
    rows, cols = 2, 5
    gap = 4
    grid = np.zeros((rows * ph + (rows - 1) * gap, cols * pw + (cols - 1) * gap),
                    dtype=np.float32)
    for i in range(n):
        r, c = i // cols, i % cols
        y0, x0 = r * (ph + gap), c * (pw + gap)
        grid[y0:y0 + ph, x0:x0 + pw] = stack[i]
    return grid


# ── figure makers ─────────────────────────────────────────────────────────────

def _make_recon_row(model_dirs_labels, n_frames: int = 2) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(n_frames, n, figsize=(n * 5.0, n_frames * 2.8),
                             facecolor="white")
    if n == 1:
        axes = axes[:, np.newaxis]
    if n_frames == 1:
        axes = axes[np.newaxis, :]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        frames = _load_visual_frames(d, n_frames)
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=9,
                               fontweight="bold", pad=3)
        for row in range(n_frames):
            ax = axes[row, col]
            if row < len(frames):
                ax.imshow(frames[row])
            ax.axis("off")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_umap_grid(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(2, n, figsize=(n * 3.5, 5.6), facecolor="white")
    if n == 1:
        axes = axes[:, np.newaxis]
    row_labels = ["FA annotation", "KMeans k=10"]
    for col, (d, lbl) in enumerate(model_dirs_labels):
        axes[0, col].set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
        for row, (key, rlbl) in enumerate(
                zip(["umap_annotation", "umap_kmeans_k10"], row_labels)):
            ax = axes[row, col]
            png_path = (d / "eval" / "cluster_panels" / "umap_kmeans_k10.png"
                        if key == "umap_kmeans_k10"
                        else d / "eval" / f"{key}.png")
            if png_path.exists():
                ax.imshow(np.array(Image.open(str(png_path))))
            else:
                ax.text(0.5, 0.5, "pending", ha="center", va="center",
                        fontsize=8, color="gray", transform=ax.transAxes)
            ax.axis("off")
            if col == 0:
                ax.set_ylabel(rlbl, fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_umap_row(model_dirs_labels, key: str) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 3.5, 3.2), facecolor="white")
    if n == 1:
        axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        png_path = (d / "eval" / "cluster_panels" / "umap_kmeans_k10.png"
                    if key == "umap_kmeans_k10"
                    else d / "eval" / f"{key}.png")
        if png_path.exists():
            ax.imshow(np.array(Image.open(str(png_path))))
        else:
            ax.text(0.5, 0.5, "pending", ha="center", va="center",
                    fontsize=8, color="gray", transform=ax.transAxes)
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_cluster_grid_fig(model_dirs_labels) -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 5.5, 3.5), facecolor="white")
    if n == 1:
        axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        grid = _load_cluster_grid(d)
        if grid is not None:
            ax.imshow(grid, cmap="gray", vmin=0, vmax=1, aspect="auto")
        else:
            ax.text(0.5, 0.5, "pending", ha="center", va="center",
                    fontsize=8, color="gray", transform=ax.transAxes)
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.suptitle("16 patches closest to each cluster centroid  (k=10, 2×5 grid, all channels shown)",
                 fontsize=8, y=0.02)
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


def _make_violin_row(model_dirs_labels, metric="recon_nl1") -> bytes:
    n = len(model_dirs_labels)
    fig, axes = plt.subplots(1, n, figsize=(n * 4.8, 3.2), facecolor="white")
    if n == 1:
        axes = [axes]
    for ax, (d, lbl) in zip(axes, model_dirs_labels):
        # current code saves without ._  prefix
        png_path = d / f"cross_dataset_{metric}.png"
        if png_path.exists():
            ax.imshow(np.array(Image.open(str(png_path))))
        else:
            ax.text(0.5, 0.5, "pending\n(eval running)", ha="center", va="center",
                    fontsize=8, color="gray", transform=ax.transAxes)
        ax.axis("off")
        ax.set_title(lbl.replace("\n", " "), fontsize=8, fontweight="bold")
    fig.tight_layout(pad=0.3)
    return _fig_to_bytes(fig)


# ── slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs):
    slide = _add_slide(prs)

    bar = slide.shapes.add_shape(1, 0, 0, _px(0.18), _px(H_IN))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()

    for y, c in [(1.8, C_3CH), (5.5, C_4CH)]:
        dot = slide.shapes.add_shape(9, _px(0.5), _px(y), _px(0.4), _px(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    _add_textbox(slide, "Contrastive Autoencoder",
                 0.8, 1.5, W_IN - 1.2, 1.0,
                 font_size=32, bold=True, color=C_TITLE)
    _add_textbox(slide,
                 "Multi-Channel Models: 3ch (pax+zyx+act) and 4ch (vinc+pax+zyx+act)",
                 0.8, 2.5, W_IN - 1.2, 0.7,
                 font_size=17, color=C_3CH)

    _add_textbox(slide, "What we evaluate:",
                 0.8, 3.6, 5.5, 0.35,
                 font_size=12, bold=True, color=C_TITLE)
    items = [
        "Reconstruction fidelity (each channel shown as separate row in panels)",
        "Latent space organisation  (UMAP colored by FA type)",
        "Unsupervised structure  (KMeans k=10 cluster panels, all channels)",
        "Generalisation to unseen datasets  (cross-dataset nL1 violin plots)",
    ]
    _add_bullets(slide, items, 0.8, 4.0, 7.5, 2.5, font_size=12, color=C_TITLE)

    _add_textbox(slide,
                 "Models: ConAE nL1  |  EnlargedJitterCrop sc2 augmentation\n"
                 "3ch: latent dim 12, proj 8  |  4ch: latent dim 12, proj 8  "
                 "|  500 epochs  |  Training: vinc (ds1)  |  Eval: ds1–ds4",
                 0.8, 6.3, W_IN - 1.2, 0.8,
                 font_size=10, color=C_LGRAY, italic=True)


def _build_model_table(prs):
    slide = _add_slide(prs)
    _slide_header(slide, "Models Compared", color=C_TITLE)

    rows_data = [
        ["Channels",               "Model",     "Recon loss", "λ_c",    "Aug",         "Latent", "Notes"],
        ["pax + zyx + act (3ch)",  "ConAE nL1", "nL1",        "0.03",   "enlcrop sc2", "12/8",   "ch0=pax, ch1=zyx, ch2=act"],
        ["vinc+pax+zyx+act (4ch)", "ConAE nL1", "nL1",        "0.03",   "enlcrop sc2", "12/8",   "all 4 channels jointly"],
    ]
    col_w = [2.5, 1.5, 1.1, 0.8, 1.4, 1.0, 2.5]
    row_h = 0.50
    x0, y0 = MARGIN, 0.80

    for ri, row in enumerate(rows_data):
        x = x0
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            is_hdr = ri == 0
            is_3ch = ri == 1
            if is_hdr:
                bg = slide.shapes.add_shape(
                    1, _px(x), _px(y0 + ri * row_h), _px(cw), _px(row_h))
                bg.fill.solid()
                bg.fill.fore_color.rgb = C_TITLE
                bg.line.fill.background()
            txt_color = (C_WHITE if is_hdr else C_3CH if is_3ch else C_4CH)
            _add_textbox(slide, cell,
                         x + 0.06, y0 + ri * row_h + 0.08,
                         cw - 0.08, row_h - 0.08,
                         font_size=9, bold=is_hdr, color=txt_color)
            x += cw

    _note_panel(slide, "Design choices", [
        "Both models share the same architecture: ConAE with normalised-L1 (nL1) recon loss",
        "λ_c = 0.03: contrastive weight scaled down to balance the multi-channel recon loss",
        "EnlargedJitterCrop (enlcrop sc2): 58×58 context → 32×32 input, intensity ÷ 2",
        "3ch encodes paxillin + zyxin + actin jointly in 12 latent dims",
        "4ch adds vinculin (ch0) to the 3ch set — all four channels share one latent code",
        "Evaluated on ds1 (vinc, train) + ds2 (pfak) + ds3 (ppax) + ds4 (nih3t3)",
        "Cluster panels show all N channels per cluster (one row per channel)",
    ], C_TITLE, font_size=11)


def _build_recon(prs, model_dirs_labels, section_label, color):
    img = _make_recon_row(model_dirs_labels, n_frames=2)
    _content_slide(prs,
        f"{section_label}  —  Reconstruction  (channels shown left→right in each frame)",
        color, img,
        "What to look for",
        [
            "Each frame: raw | reconstructed side-by-side, all channels composited as RGB",
            "Channel order in visual.tif: left stripe = ch0, right stripe = chN−1",
            "Blurring / smoothing reflects latent bottleneck averaging channel details",
            "Compare reconstruction quality channel-by-channel across models",
        ])


def _build_umap(prs, model_dirs_labels, section_label, color):
    img = _make_umap_grid(model_dirs_labels)
    _content_slide(prs,
        f"{section_label}  —  UMAP  (top: FA annotation  ·  bottom: KMeans k=10)",
        color, img,
        "Interpreting UMAP",
        [
            "Top row: coloured by FA annotation (nascent/growing/mature/retracting)",
            "Bottom row: KMeans k=10 cluster labels",
            "Cluster–annotation alignment (unsupervised) = model captured FA maturation",
            "Multi-channel latent encodes joint morphotype — expect richer structure "
            "than single-channel if channels carry complementary information",
        ])


def _build_clusters(prs, model_dirs_labels, section_label, color, ch_names: str):
    img = _make_cluster_grid_fig(model_dirs_labels)
    _content_slide(prs,
        f"{section_label}  —  KMeans k=10 cluster panels  (2×5 grid · {ch_names})",
        color, img,
        "Cluster interpretation",
        [
            "Each panel: one row per channel (label bar + 4×4 patch grid)",
            "16 patches nearest each cluster centroid — the model's representative examples",
            "Good cluster = coherent in ALL channels simultaneously",
            "Channel-specific variation within a cluster reveals which channel drives the split",
            "N = cluster size; small N = rare or transitional morphotype",
        ])


def _build_violin(prs, model_dirs_labels, section_label, color):
    img = _make_violin_row(model_dirs_labels, metric="recon_nl1")
    _content_slide(prs,
        f"{section_label}  —  Cross-dataset normalised L1  (ds1=vinc · ds2=pfak · ds3=ppax · ds4=nih3t3)",
        color, img,
        "Generalisation quality",
        [
            "nL1 = L1 / mean|raw|: scale-invariant error averaged across all channels",
            "ds1 (vinc) = train; ds2–4 = unseen cell lines",
            "Small train→test gap = model learned cell-line-agnostic features",
            "Multi-channel models may generalise better if channels regularise each other",
            "ds4 (nih3t3) = hardest domain shift: different cell type",
        ])


def _build_comparison(prs, m3ch, m4ch):
    slide = _add_slide(prs)
    _slide_header(slide,
        "3ch (pax+zyx+act) vs 4ch (vinc+pax+zyx+act)  —  UMAP annotation  ·  Cross-dataset nL1",
        color=C_TITLE)

    both = [m3ch, m4ch]
    img_umap = _make_umap_row(both, "umap_annotation")
    img_viol = _make_violin_row(both, metric="recon_nl1")

    top_h   = 3.05
    bot_top = 0.68 + top_h + 0.05
    _add_image_bytes(slide, img_umap, MARGIN, 0.68, width=IMG_W)
    _add_image_bytes(slide, img_viol, MARGIN, bot_top, width=IMG_W)

    _note_panel(slide, "Key comparisons", [
        "Both ConAE nL1, λ_c=0.03, trained on vinc (ds1), lat dim 12",
        "Does adding vinculin (ch0) change UMAP structure vs 3ch?",
        "4ch has 4× the reconstruction channels in the same 12-dim latent — "
        "expect higher nL1 but richer cluster morphotypes",
        "Compare train→test gap: does extra channel diversity regularise better?",
        "Cluster panels: look for clusters that differ primarily in the vinculin channel",
    ], C_TITLE, font_size=12)


def _build_section_header(prs, title, subtitle, color):
    slide = _add_slide(prs)
    _slide_header(slide, title, color=color, font_size=28)
    bar = slide.shapes.add_shape(
        1, _px(0), _px(H_IN * 0.82), _px(W_IN), _px(H_IN * 0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    _add_textbox(slide, subtitle, MARGIN, H_IN * 0.83,
                 W_IN - 2 * MARGIN, 0.5,
                 font_size=14, color=C_WHITE, bold=False)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path, default=Path("slides_multichannel.pptx"))
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = Inches(W_IN)
    prs.slide_height = Inches(H_IN)

    m3ch = (MODEL_3CH, LABEL_3CH)
    m4ch = (MODEL_4CH, LABEL_4CH)

    steps = [
        ("Title",          _build_title_slide,  (prs,)),
        ("Model table",    _build_model_table,  (prs,)),
        # ── 3ch section ──
        ("3ch hdr",        None,
         ("Results: 3ch — paxillin + zyxin + actin", "ConAE nL1  ·  λ_c=0.03  ·  latent 12/8", C_3CH)),
        ("3ch recon",      _build_recon,
         (prs, [m3ch], "3ch (pax+zyx+act)", C_3CH)),
        ("3ch UMAP",       _build_umap,
         (prs, [m3ch], "3ch (pax+zyx+act)", C_3CH)),
        ("3ch clusters",   _build_clusters,
         (prs, [m3ch], "3ch (pax+zyx+act)", C_3CH, "pax · zyx · act rows")),
        ("3ch violin",     _build_violin,
         (prs, [m3ch], "3ch (pax+zyx+act)", C_3CH)),
        # ── 4ch section ──
        ("4ch hdr",        None,
         ("Results: 4ch — vinculin + paxillin + zyxin + actin", "ConAE nL1  ·  λ_c=0.03  ·  latent 12/8", C_4CH)),
        ("4ch recon",      _build_recon,
         (prs, [m4ch], "4ch (vinc+pax+zyx+act)", C_4CH)),
        ("4ch UMAP",       _build_umap,
         (prs, [m4ch], "4ch (vinc+pax+zyx+act)", C_4CH)),
        ("4ch clusters",   _build_clusters,
         (prs, [m4ch], "4ch (vinc+pax+zyx+act)", C_4CH, "vinc · pax · zyx · act rows")),
        ("4ch violin",     _build_violin,
         (prs, [m4ch], "4ch (vinc+pax+zyx+act)", C_4CH)),
        # ── side-by-side ──
        ("Comparison",     _build_comparison,  (prs, m3ch, m4ch)),
    ]

    for label, fn, fn_args in steps:
        print(f"  {label} …", flush=True)
        if fn is None:
            title_text, subtitle, color = fn_args
            _build_section_header(prs, title_text, subtitle, color)
        else:
            fn(*fn_args)

    prs.save(str(args.out))
    print(f"\nSaved → {args.out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
