#!/usr/bin/env python3
"""
Generate a PowerPoint explaining Contrastive AE and Supervised Contrastive AE.

Usage:
  python scripts/make_contrastive_pptx.py [--out output.pptx]
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── palette (white theme) ─────────────────────────────────────────────────────
C_BG     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x11, 0x11, 0x11)
C_DARK   = RGBColor(0x1A, 0x3A, 0x6B)   # dark navy — section headers
C_BLUE   = RGBColor(0x1F, 0x5C, 0x99)   # medium blue — box headers
C_LBLUE  = RGBColor(0xD6, 0xE8, 0xF8)   # light blue — box bodies
C_GREEN  = RGBColor(0x1A, 0x7A, 0x4A)   # green — positive pairs
C_RED    = RGBColor(0xB0, 0x20, 0x20)   # red — negative pairs
C_ORANGE = RGBColor(0xC0, 0x60, 0x10)   # orange — z_proj
C_PURPLE = RGBColor(0x5A, 0x1F, 0x8A)   # purple — z_recon
C_GREY   = RGBColor(0x66, 0x66, 0x66)
C_LGREY  = RGBColor(0xF2, 0xF5, 0xFA)


def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def _title(slide, text, top=Inches(0.2), fontsize=30, color=C_DARK):
    txb = slide.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.65))
    p   = txb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(fontsize)
    run.font.bold  = True
    run.font.color.rgb = color


def _sub(slide, text, top=Inches(0.85), fontsize=13, color=C_GREY):
    txb = slide.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.35))
    p   = txb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text   = text
    run.font.size   = Pt(fontsize)
    run.font.italic = True
    run.font.color.rgb = color


def _body(slide, lines, left=Inches(0.5), top=Inches(1.25),
          width=Inches(9.0), height=Inches(5.8),
          fontsize=13, color=C_BLACK):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(line, tuple):
            text, sz, bold, col, bullet = line
        else:
            text, sz, bold, col, bullet = line, fontsize, False, color, True
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if bullet else "") + text
        run.font.size  = Pt(sz)
        run.font.bold  = bold
        run.font.color.rgb = col


def _box(slide, text, left, top, width, height,
         bg=C_BLUE, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=11, bold=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xBB,0xCC,0xDD); shape.line.width = Pt(0.75)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.07)
    tf.margin_top  = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(fontsize); run.font.color.rgb = fg
    run.font.bold = bold
    return shape


def _arrow(slide, x1, y, x2, color=C_BLACK, width=Pt(1.5)):
    """Horizontal arrow from (x1,y) to (x2,y)."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y, x2, y)
    connector.line.color.rgb = color
    connector.line.width = width


def _hline(slide, left, top, width, color=C_GREY):
    line = slide.shapes.add_shape(1, left, top, width, Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.color.rgb = color


def make_pptx(out_path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Contrastive Autoencoder for Focal Adhesion Morphology",
           top=Inches(2.0), fontsize=32)
    _sub(sl, "Self-supervised and supervised contrastive learning on paxillin patches",
         top=Inches(2.85), fontsize=16, color=C_BLUE)
    _body(sl, [
        ("Contrastive AE (ConAE) — NT-Xent self-supervised loss", 14, False, C_GREY, False),
        ("Supervised Contrastive AE (SupCon AE) — label-guided positive pairs", 14, False, C_GREY, False),
    ], top=Inches(3.6), fontsize=14, color=C_GREY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — ConAE architecture overview
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Contrastive AE — Architecture")
    _sub(sl,   "Encoder feeds both a reconstruction head and a contrastive projection head")

    # architecture flow: patch → encoder → z_recon → decoder → recon
    #                                    ↓
    #                               projector → z_proj → contrastive loss
    boxes = [
        (Inches(0.3),  Inches(2.0), Inches(1.3), Inches(0.9), "Patch\n(32×32)", C_LGREY, C_BLACK),
        (Inches(2.0),  Inches(2.0), Inches(1.5), Inches(0.9), "Encoder\n(CNN)", C_BLUE,  RGBColor(0xFF,0xFF,0xFF)),
        (Inches(4.1),  Inches(1.5), Inches(1.4), Inches(0.9), "z_recon\n(latent)", C_PURPLE, RGBColor(0xFF,0xFF,0xFF)),
        (Inches(6.2),  Inches(1.5), Inches(1.5), Inches(0.9), "Decoder\n(CNN)", C_BLUE,  RGBColor(0xFF,0xFF,0xFF)),
        (Inches(8.3),  Inches(1.5), Inches(1.4), Inches(0.9), "Recon\npatch", C_LGREY, C_BLACK),
        (Inches(4.1),  Inches(3.2), Inches(1.4), Inches(0.9), "Projector\n(MLP)", C_ORANGE, RGBColor(0xFF,0xFF,0xFF)),
        (Inches(6.2),  Inches(3.2), Inches(1.4), Inches(0.9), "z_proj\n(smaller)", C_ORANGE, RGBColor(0xFF,0xFF,0xFF)),
        (Inches(8.3),  Inches(3.2), Inches(1.4), Inches(0.9), "Contrastive\nloss", RGBColor(0x8B,0x00,0x00), RGBColor(0xFF,0xFF,0xFF)),
    ]
    for x, y, w, h, txt, bg, fg in boxes:
        _box(sl, txt, x, y, w, h, bg=bg, fg=fg, fontsize=11, bold=True)

    # arrows (approximate, using thin rectangles as lines)
    for x1, x2, y in [
        (Inches(1.65), Inches(2.0),  Inches(2.45)),
        (Inches(3.55), Inches(4.1),  Inches(1.95)),
        (Inches(5.55), Inches(6.2),  Inches(1.95)),
        (Inches(7.75), Inches(8.3),  Inches(1.95)),
        (Inches(5.55), Inches(6.2),  Inches(3.65)),
        (Inches(7.65), Inches(8.3),  Inches(3.65)),
    ]:
        line = sl.shapes.add_shape(1, x1, y, x2-x1, Pt(1.5))
        line.fill.solid(); line.fill.fore_color.rgb = C_BLACK
        line.line.color.rgb = C_BLACK

    # vertical arrow from z_recon down to projector
    vline = sl.shapes.add_shape(1, Inches(4.77), Inches(2.45), Pt(1.5), Inches(0.75))
    vline.fill.solid(); vline.fill.fore_color.rgb = C_BLACK
    vline.line.color.rgb = C_BLACK

    # labels
    _body(sl, [
        ("z_recon  — encodes all patch information; used for downstream classification & analysis", 11, False, C_PURPLE, False),
        ("z_proj   — compact projection for contrastive loss only; discarded at inference", 11, False, C_ORANGE, False),
        ("Recon loss (MSE) on z_recon path    +    Contrastive loss (NT-Xent) on z_proj path", 11, False, C_BLACK, False),
    ], top=Inches(4.6), fontsize=11)

    _hline(sl, Inches(0.4), Inches(4.5), Inches(9.2))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Augmentation
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Augmentation Strategy")
    _sub(sl,   "Two views per patch created on-the-fly each training step")

    _body(sl, [
        ("For each patch in the batch, two views are generated:", 14, True, C_DARK, False),
        "",
        ("View 1  —  Clean", 13, True, C_BLUE, False),
        "Original patch, no modification.",
        "Used for reconstruction loss (MSE between clean and decoded z_recon).",
        "",
        ("View 2  —  Augmented", 13, True, C_ORANGE, False),
        "Random 90° rotation  (rot90 mode — preserves FA morphology orientation structure)",
        "Optional random horizontal / vertical flip",
        "Optional soft salt-and-pepper noise  (intensity = mean ± std/3, NOT 0/1)",
        "  → avoids bright spots that mimic nascent adhesion puncta",
        "Optional per-patch intensity scaling  (default ×[0.8, 1.2])",
        "",
        ("Key design choice", 13, True, C_DARK, False),
        "Augmentations are morphology-preserving — the FA structure is unchanged,",
        "only orientation, brightness, and minor noise differ.",
        "The encoder must learn to produce similar z_proj for both views → orientation-invariant features.",
    ], fontsize=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Positive and negative pairs
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Positive and Negative Pairs (NT-Xent)")
    _sub(sl,   "Self-supervised contrastive learning — no labels required")

    _body(sl, [
        ("Batch construction: 2N embeddings", 14, True, C_DARK, False),
        "Each batch of N patches produces 2N projection vectors:",
        "  [ z_proj(clean_1), …, z_proj(clean_N),  z_proj(aug_1), …, z_proj(aug_N) ]",
        "",
        ("Positive pair  (pull together  ↑)", 13, True, C_GREEN, False),
        "Each patch i is paired with its own augmented version i+N.",
        "These two views of the same patch should have similar z_proj.",
        "Only 1 positive per anchor in standard ConAE.",
        "",
        ("Negative pairs  (push apart  ↓)", 13, True, C_RED, False),
        "All other 2N−2 embeddings in the batch are negatives.",
        "Different patches should have dissimilar z_proj.",
        "Batch size matters: more negatives = stronger contrastive signal.",
        "",
        ("NT-Xent loss per anchor i:", 13, True, C_DARK, False),
        "L_i  =  −log [ exp(sim(i, i+N) / τ)  /  Σ_{k≠i} exp(sim(i, k) / τ) ]",
        "τ = temperature (default 0.5) — lower = sharper discrimination",
    ], fontsize=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — z_recon vs z_proj
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "z_recon vs z_proj — Two Latent Spaces")
    _sub(sl,   "Why have two separate representations?")

    # two columns
    for ci, (title, color, items) in enumerate([
        ("z_recon  (latent_dim = 12)", C_PURPLE, [
            "Output of the encoder.",
            "Encodes all patch information needed to reconstruct the image.",
            "Higher dimensionality — richer representation.",
            "Used for: downstream FA classification, UMAP, clustering, latent distance metrics.",
            "Trained by: reconstruction loss (MSE).",
            "Not directly exposed to contrastive gradient — protected from representation collapse.",
            "",
            "Think of it as:  'What does this patch look like?'",
        ]),
        ("z_proj  (proj_dim = 8)", C_ORANGE, [
            "Output of a small MLP projector on top of z_recon.",
            "Lower dimensionality — compact invariant summary.",
            "Used for: NT-Xent / SupCon contrastive loss only.",
            "Discarded at inference — never used for downstream tasks.",
            "Trained by: contrastive loss — must be similar for same-structure patches.",
            "The projector absorbs the contrastive invariance pressure,",
            "  protecting z_recon from collapsing.",
            "",
            "Think of it as:  'Is this patch structurally equivalent to another?'",
        ]),
    ]):
        x = Inches(0.3 + ci * 4.9)
        _box(sl, title, x, Inches(1.2), Inches(4.5), Inches(0.45),
             bg=color, fontsize=12, bold=True)
        _box(sl, "\n".join(("  " if i.startswith(" ") else "• ") + i if i else ""
                            for i in items),
             x, Inches(1.7), Inches(4.5), Inches(4.5),
             bg=C_LGREY, fg=C_BLACK, fontsize=11)

    _body(sl, [
        "Key insight (SimCLR / MoCo finding): using the projector output for contrastive loss",
        "instead of the raw encoder output consistently improves downstream task performance.",
    ], top=Inches(6.5), fontsize=11, color=C_GREY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Warmup phase
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Training Strategy — Warmup Phase")
    _sub(sl,   "Reconstruction-only warmup before contrastive loss activates")

    _body(sl, [
        ("Problem without warmup", 14, True, C_RED, False),
        "If contrastive loss starts from a random encoder, z_proj is random noise.",
        "NT-Xent loss on random embeddings provides misleading gradient — encoder may diverge.",
        "The decoder has no useful z_recon to work with in early epochs.",
        "",
        ("Warmup solution (warmup_epochs = 100)", 14, True, C_GREEN, False),
        "Phase 1  (epochs 0–100):   λ_contrast = 0   →   reconstruction-only training.",
        "  Encoder learns to compress meaningful patch structure into z_recon.",
        "  Decoder learns to reconstruct from this structure.",
        "Phase 2  (epochs 100–500):  λ_contrast = 0.5  →  full contrastive + recon training.",
        "  LR reset to original value at transition, scheduler restarted.",
        "  Contrastive gradients now act on a well-formed representation.",
        "",
        ("Effect on z_proj", 14, True, C_DARK, False),
        "By epoch 100 the projector maps from a meaningful z_recon,",
        "so positive pairs are already somewhat similar → stable NT-Xent gradient from the start.",
    ], fontsize=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Supervised ConAE overview
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Supervised Contrastive AE (SupCon AE)")
    _sub(sl,   "Label-guided positive pairs — more positives for labeled patches")

    _body(sl, [
        ("Key difference from ConAE", 14, True, C_DARK, False),
        "In standard ConAE: each patch has exactly 1 positive (its own augmented view).",
        "In SupCon AE: labeled patches have all same-class patches in the batch as positives.",
        "",
        ("Hybrid loss (labeled + unlabeled in the same batch)", 14, True, C_DARK, False),
        ("Labeled anchor  (label ≥ 0)", 13, True, C_GREEN, False),
        "Positives = all other patches in the 2N batch sharing the same FA-type label.",
        "  (includes own aug view, plus all same-class clean and aug patches)",
        "The more labeled patches in the batch, the stronger the supervisory signal.",
        ("Unlabeled anchor  (label = −1)", 13, True, C_BLUE, False),
        "Falls back to standard NT-Xent: only own augmented view as positive.",
        "Unlabeled patches still learn general morphological structure.",
        "",
        ("SupCon loss (labeled anchors)", 13, True, C_DARK, False),
        "L_i  =  −(1/|P_i|) · Σ_{p∈P_i} sim(i,p)/τ  +  log Σ_{k≠i} exp(sim(i,k)/τ)",
        "  P_i = set of all positives for anchor i  (can be > 1)",
        "  Larger |P_i| → more stable gradient, encourages tight class clusters.",
    ], fontsize=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — SupCon for FA data
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "SupCon AE Applied to Focal Adhesion Data")
    _sub(sl,   "Sparse labels, 5 FA types, strong class imbalance")

    _body(sl, [
        ("Dataset composition (vinc, control + ycomp)", 14, True, C_DARK, False),
        "Total patches:        ~24 400   (training + validation)",
        "Labeled patches:       ~1 295   (5.3% of total)",
        "  No adhesion:            692   — dominant class, excluded from metrics",
        "  Focal complex:          257",
        "  Focal adhesion:         218",
        "  Nascent Adhesion:       112",
        "  Fibrillar adhesion:      16   — very rare",
        "Unlabeled patches:    ~23 100   (95% — use NT-Xent fallback)",
        "",
        ("What SupCon learns for FA types", 14, True, C_DARK, False),
        "Nascent adhesions:  small puncta at cell edge → cluster together in z_proj.",
        "Focal complexes:    slightly larger, still peripheral → adjacent to nascent.",
        "Focal adhesions:    elongated, mature → clearly separated from nascent.",
        "Fibrillar:          long, linear, cell centre → most distinct cluster.",
        "",
        ("Challenge", 13, True, C_RED, False),
        "With only 16 fibrillar patches and batch_size=128, many batches have 0 fibrillar.",
        "→ Fibrillar cluster is shaped primarily by unlabeled self-supervised signal.",
    ], fontsize=11)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — FA maturation continuum + pair weights
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "FA Maturation Continuum — Weighted SupCon")
    _sub(sl,   "Encoding biological knowledge into the contrastive loss (exp/nonad-vs-ad branch)")

    _body(sl, [
        ("FA adhesion maturation pathway", 14, True, C_DARK, False),
        "Nascent Adhesion  →  Focal Complex  →  Focal Adhesion  →  Fibrillar Adhesion",
        "Adjacent stages share morphological features; distant stages are structurally distinct.",
        "",
        ("Standard SupCon limitation", 13, True, C_RED, False),
        "Binary positive/negative: same class = pull together, different class = push apart equally.",
        "But Nascent and Focal Complex are more similar than Nascent and Fibrillar.",
        "Treating Nascent vs Focal Complex as a hard negative discards this gradient.",
        "",
        ("Weighted SupCon (pair weight matrix)", 13, True, C_GREEN, False),
        "Each class pair gets a continuous weight  w[y_i, y_j]  from a K×K matrix:",
    ], fontsize=12)

    # pair weight matrix visual
    labels = ["Nascent", "F.Complex", "F.Adhesion", "Fibrillar"]
    matrix = [[1.0, 0.3, -0.5, -1.0],
              [0.3, 1.0,  0.3, -0.5],
              [-0.5, 0.3, 1.0,  0.3],
              [-1.0, -0.5, 0.3,  1.0]]
    cell_w = Inches(0.95); cell_h = Inches(0.42)
    x0 = Inches(1.8); y0 = Inches(4.4)
    # header row
    for ci, lbl in enumerate(labels):
        _box(sl, lbl, x0 + (ci+1)*cell_w, y0, cell_w, cell_h,
             bg=C_BLUE, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=9)
    for ri, row_lbl in enumerate(labels):
        _box(sl, row_lbl, x0, y0 + (ri+1)*cell_h, cell_w, cell_h,
             bg=C_BLUE, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=9)
        for ci, val in enumerate(matrix[ri]):
            if val > 0.5:
                bg = RGBColor(0xC8, 0xE6, 0xC9)  # green
            elif val > 0:
                bg = RGBColor(0xE8, 0xF5, 0xE9)  # light green
            elif val > -0.6:
                bg = RGBColor(0xFF, 0xEB, 0xEE)  # light red
            else:
                bg = RGBColor(0xFF, 0xCC, 0xCC)  # red
            _box(sl, f"{val:+.1f}",
                 x0 + (ci+1)*cell_w, y0 + (ri+1)*cell_h, cell_w, cell_h,
                 bg=bg, fg=C_BLACK, fontsize=10)

    _body(sl, [
        "Positive weight → attract (pull together).   Negative weight → repel (push apart).",
        "The loss becomes:  L_i = −Σ_j w[y_i,y_j]·sim(i,j) / Σ_j |w[y_i,y_j]| + log_denom",
    ], top=Inches(6.65), fontsize=10, color=C_GREY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Summary comparison
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Summary — ConAE vs SupCon AE")

    rows = [
        ("",             "ConAE",                    "SupCon AE",             "Weighted SupCon"),
        ("Positives",    "Own aug view only (1/batch)","All same-class in batch","All same-class, weighted"),
        ("Labels needed","None",                      "Optional (fallback NT-Xent if absent)", "Required for labeled patches"),
        ("Loss",         "NT-Xent",                   "SupCon + NT-Xent hybrid", "Weighted SupCon + NT-Xent"),
        ("z_proj shape", "Morphology-invariant",       "Class-cluster-aware",   "Continuum-aware"),
        ("z_recon use",  "Reconstruction + downstream","Reconstruction + downstream","Reconstruction + downstream"),
        ("FA data fit",  "All 24k patches used equally","5% labeled guides clusters","Maturation continuum encoded"),
        ("Status",       "Trained (multiple strategies)","Trained (baseline)",    "Exp branch — in development"),
    ]
    col_w = [Inches(1.8), Inches(2.4), Inches(2.6), Inches(2.6)]
    row_h = Inches(0.52)
    x0    = Inches(0.3)
    y0    = Inches(1.1)
    col_bgs = [C_BLUE, C_LGREY, RGBColor(0xE8,0xF5,0xE9), RGBColor(0xFF,0xF3,0xE0)]
    col_fgs = [RGBColor(0xFF,0xFF,0xFF), C_BLACK, C_BLACK, C_BLACK]

    for ri, row in enumerate(rows):
        for ci, (cell, w) in enumerate(zip(row, col_w)):
            x  = x0 + sum(col_w[:ci])
            y  = y0 + ri * row_h
            if ri == 0:
                bg = C_DARK; fg = RGBColor(0xFF,0xFF,0xFF)
            elif ri % 2 == 1:
                bg = col_bgs[ci]; fg = col_fgs[ci]
            else:
                bg = RGBColor(0xF8,0xFA,0xFF) if ci > 0 else C_BLUE
                fg = RGBColor(0xFF,0xFF,0xFF) if ci == 0 else C_BLACK
            _box(sl, cell, x, y, w, row_h, bg=bg, fg=fg, fontsize=9)

    # ── save ─────────────────────────────────────────────────────────────────
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path,
                        default=Path("/net/projects/CLS/lding/data/fa_data_analysis"
                                     "/ae_results/contrastive_ae_overview.pptx"))
    args = parser.parse_args()
    make_pptx(args.out)


if __name__ == "__main__":
    main()
