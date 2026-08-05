#!/usr/bin/env python3
"""
make_pptx_blind_test_crossds.py

Generate PowerPoint summarising the blind cross-dataset evaluation results
from run_blind_test_crossds.py.

For each of the 9 Annabel-vinc models, for each of 4 (ds, cond) pairs,
for each of 2 feature spaces (zrecon, zproj):
  - confusion matrix (normalised)
  - per-class F1
  - overall accuracy / macro-F1

Slide layout:
  1. Title slide
  2. Approach overview
  3. Label statistics (Margaret vs Annabel)
  Per model (9 × 3 slides = 27):
    4–6.   conae  s1v3
    7–9.   conae  s2v2
   10–12.  conae  s3v1
   13–15.  supcon2 s1v3
   16–18.  supcon2 s2v2
   19–21.  supcon2 s3v1
   22–24.  supcon5 s1v3
   25–27.  supcon5 s2v2
   28–30.  supcon5 s3v1
  Comparison:
   31.  Macro-F1 summary table (all 9 models × 4 ds/cond × 2 feat)

Usage:
  python scripts/make_pptx_blind_test_crossds.py
  python scripts/make_pptx_blind_test_crossds.py --out results/blind_test.pptx
"""
from __future__ import annotations

import argparse
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
REPO_ROOT  = Path(__file__).resolve().parents[1]
DATA_ROOT  = Path("/net/projects/CLS/lding/data/fa_data_analysis")
RUN_DIR    = DATA_ROOT / "ae_results" / "contrastive_run"

MODELS    = ["conae",   "supcon2",  "supcon5"]
SPLITS    = ["s1v3",    "s2v2",     "s3v1"]
FEATS     = ["zrecon",  "zproj"]
DS_CONDS  = [
    ("vinc",  "control"),
    ("vinc",  "ycomp"),
    ("ppax",  "control"),
    ("pfak",  "control"),
]

# slide dimensions (widescreen 13.33 × 7.5 in)
W = Inches(13.33)
H = Inches(7.5)

TITLE_COLOR  = RGBColor(0x1F, 0x49, 0x7D)   # dark blue
BODY_COLOR   = RGBColor(0x26, 0x26, 0x26)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs: Presentation):
    blank = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank)


def _add_text(slide, text: str, left, top, width, height,
              bold=False, size_pt=14, color=BODY_COLOR, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(size_pt)
    run.font.color.rgb = color
    return txb


def _img_or_ph(slide, path: Path, left, top, width, height, label="[pending]"):
    if path and path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}\n{path}" if path else label
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)


def _result_dir(model: str, split: str) -> Path:
    return RUN_DIR / f"annabel_vinc_{model}_{split}"


def _blind_dir(model: str, split: str) -> Path:
    return _result_dir(model, split) / "blind_test"


def _cm_path(model: str, split: str, ds: str, cond: str, feat: str) -> Path:
    return _blind_dir(model, split) / f"{ds}_{cond}_{feat}" / "confusion_matrix_norm.png"


def _metrics_path(model: str, split: str, ds: str, cond: str, feat: str) -> Path:
    return _blind_dir(model, split) / f"{ds}_{cond}_{feat}" / "metrics.csv"


def _load_metric(model, split, ds, cond, feat, col="macro_f1") -> float | None:
    p = _metrics_path(model, split, ds, cond, feat)
    if not p.exists():
        return None
    try:
        df = pd.read_csv(p)
        return float(df[col].iloc[0])
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------

def slide_title(prs: Presentation):
    sl = _blank(prs)
    _add_text(sl, "Blind Cross-Dataset Evaluation",
              Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
              bold=True, size_pt=32, color=TITLE_COLOR)
    _add_text(sl,
              "9 Annabel-vinc AE+LightGBM models → vinc (ctrl+ycomp), ppax (ctrl), pfak (ctrl)\n"
              "Evaluated against Margaret's independent label CSVs (labels_*_20260521.csv)\n"
              "Features: z_recon (12-d) | z_proj (8-d)  •  Models: ConAE, SupCon-2cls, SupCon-5cls",
              Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.5),
              size_pt=14)


# ---------------------------------------------------------------------------
# Slide 2: Approach
# ---------------------------------------------------------------------------

def slide_approach(prs: Presentation):
    sl = _blank(prs)
    _add_text(sl, "Approach", Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5),
              bold=True, size_pt=22, color=TITLE_COLOR)
    body = (
        "Training domain: Annabel-labeled vinc/control patches (539 patches, 4 frames, 5 FA classes)\n\n"
        "Inference on:\n"
        "  • vinc control   (14879 patches, pax channel)\n"
        "  • vinc ycomp     (12758 patches, pax channel)\n"
        "  • ppax control   (pax channel, ch1)\n"
        "  • pfak control   (pax channel, ch1)\n\n"
        "Evaluation: match patches by unique_ID to Margaret labels (labels_*_20260521.csv)\n"
        "  Exclude 'Uncertain' labels.  For SupCon-2cls: remap adhesion subtypes → 'adhesion'\n\n"
        "Output per model × dataset × feature:\n"
        "  blind_test/{ds}_{cond}_{feat}/  →  confusion_matrix_norm.png, metrics.csv, predictions.csv"
    )
    _add_text(sl, body, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.0), size_pt=13)


# ---------------------------------------------------------------------------
# Slide 3: Label statistics
# ---------------------------------------------------------------------------

def slide_label_stats(prs: Presentation):
    sl = _blank(prs)
    _add_text(sl, "Label Statistics: Margaret (blind) vs Annabel (training)",
              Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5),
              bold=True, size_pt=18, color=TITLE_COLOR)

    stats = {
        "vinc/control (Margaret)": {"No adhesion": 107, "Nascent Adhesion": 14,
                                     "focal complex": 52, "focal adhesion": 192,
                                     "fibrillar adhesion": 12, "Uncertain": 14},
        "vinc/ycomp  (Margaret)":  {"No adhesion": 594, "Nascent Adhesion": 99,
                                     "focal complex": 208, "focal adhesion": 28,
                                     "fibrillar adhesion": 4, "Uncertain": 16},
        "ppax/control (Margaret)": {"No adhesion": 15, "Nascent Adhesion": 9,
                                     "focal complex": 5, "focal adhesion": 22,
                                     "fibrillar adhesion": 0, "Uncertain": 9},
        "pfak/control (Margaret)": {"No adhesion": 14, "Nascent Adhesion": 4,
                                     "focal complex": 7, "focal adhesion": 29,
                                     "fibrillar adhesion": 0, "Uncertain": 0},
        "vinc/control (Annabel, train)": {"No adhesion": "~260", "Nascent Adhesion": "~65",
                                           "focal complex": "~5", "focal adhesion": "~190",
                                           "fibrillar adhesion": "~6", "Uncertain": "-"},
    }

    y = Inches(0.9)
    for title, counts in stats.items():
        row = "  ".join(f"{k}: {v}" for k, v in counts.items())
        _add_text(sl, f"{title}\n  {row}", Inches(0.5), y, Inches(12.3), Inches(0.75),
                  size_pt=11)
        y += Inches(0.85)

    _add_text(sl,
              "Note: Models trained on Annabel vinc/control only. "
              "ppax and pfak are fully out-of-distribution (different cell line markers).",
              Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
              size_pt=11, color=RGBColor(0x80, 0x40, 0x00))


# ---------------------------------------------------------------------------
# Per-model slides (3 per model×split): one slide per feature
# ---------------------------------------------------------------------------

def slide_model_blind(prs: Presentation, model: str, split: str, feat: str):
    sl = _blank(prs)
    title = f"Blind Test — {model}  {split}  |  {feat}"
    _add_text(sl, title, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45),
              bold=True, size_pt=16, color=TITLE_COLOR)

    # 2×2 grid of confusion matrices: row=feat, col=ds/cond
    ds_labels = ["vinc/ctrl", "vinc/ycomp", "ppax/ctrl", "pfak/ctrl"]
    ds_conds  = [("vinc", "control"), ("vinc", "ycomp"), ("ppax", "control"), ("pfak", "control")]

    col_w = Inches(3.1)
    col_gap = Inches(0.1)
    row_h = Inches(3.1)
    x0 = Inches(0.2)
    y0 = Inches(0.65)

    for ci, ((ds, cond), label) in enumerate(zip(ds_conds, ds_labels)):
        x = x0 + ci * (col_w + col_gap)
        # Label above
        _add_text(sl, label, x, y0, col_w, Inches(0.25), bold=True, size_pt=11)
        # Confusion matrix
        cm_p = _cm_path(model, split, ds, cond, feat)
        _img_or_ph(sl, cm_p, x, y0 + Inches(0.27), col_w, row_h,
                   label=f"[{ds}/{cond}/{feat}]")

    # Metrics table row at bottom
    _add_text(sl, "Macro-F1 (accuracy)",
              Inches(0.2), Inches(4.1), Inches(2.5), Inches(0.3),
              bold=True, size_pt=11)
    col_w2 = Inches(3.1)
    for ci, (ds, cond) in enumerate(ds_conds):
        x = x0 + ci * (col_w2 + col_gap)
        mf1 = _load_metric(model, split, ds, cond, feat, "macro_f1")
        acc = _load_metric(model, split, ds, cond, feat, "accuracy")
        val = f"F1={mf1:.3f}  acc={acc:.3f}" if mf1 is not None else "[pending]"
        _add_text(sl, val, x, Inches(4.1), col_w2, Inches(0.35), size_pt=10)


# ---------------------------------------------------------------------------
# Summary slide: macro-F1 table
# ---------------------------------------------------------------------------

def slide_summary(prs: Presentation):
    sl = _blank(prs)
    _add_text(sl, "Summary: Macro-F1 across all models (blind test)",
              Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45),
              bold=True, size_pt=18, color=TITLE_COLOR)

    rows = []
    for model in MODELS:
        for split in SPLITS:
            for feat in FEATS:
                row = {"model": model, "split": split, "feat": feat}
                for ds, cond in DS_CONDS:
                    key = f"{ds}_{cond}"
                    mf1 = _load_metric(model, split, ds, cond, feat, "macro_f1")
                    row[key] = f"{mf1:.3f}" if mf1 is not None else "—"
                rows.append(row)

    df = pd.DataFrame(rows)

    # Save as temp png table
    fig, ax = plt.subplots(figsize=(13, 9))
    ax.axis("off")
    col_labels = ["model", "split", "feat"] + [f"{ds}/{cond}" for ds, cond in DS_CONDS]
    tbl = ax.table(
        cellText=df.values,
        colLabels=col_labels,
        loc="center",
        cellLoc="center",
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)
    tbl.scale(1, 1.5)
    fig.tight_layout()
    tmp_path = Path("/tmp/blind_test_summary_table.png")
    fig.savefig(str(tmp_path), dpi=150, bbox_inches="tight")
    plt.close(fig)

    _img_or_ph(sl, tmp_path, Inches(0.2), Inches(0.65), Inches(12.9), Inches(6.5),
               label="[summary table pending]")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path,
                    default=REPO_ROOT / "results" / "blind_test_crossds.pptx")
    args = ap.parse_args()
    args.out.parent.mkdir(parents=True, exist_ok=True)

    prs = _prs()
    slide_title(prs)
    slide_approach(prs)
    slide_label_stats(prs)

    for model in MODELS:
        for split in SPLITS:
            for feat in FEATS:
                slide_model_blind(prs, model, split, feat)

    slide_summary(prs)

    prs.save(str(args.out))
    print(f"Saved: {args.out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
