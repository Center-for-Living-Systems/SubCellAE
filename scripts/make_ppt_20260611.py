"""
make_ppt_20260611.py
Generate slides_20260611.pptx — EnlargedCropDataset & Contrastive AE development log.

Run with:
    /net/projects/CLS/lding/conda_env/core_env/bin/python3 scripts/make_ppt_20260611.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml as oxml
from lxml import etree

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
TEAL   = RGBColor(0x2E, 0x86, 0xAB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

SW = Inches(13.33)
SH = Inches(7.50)
TITLE_H = Inches(1.1)
MARGIN  = Inches(0.35)

# Image directories
BASE_IMG = "/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run"
CONAE_DIR  = os.path.join(BASE_IMG, "contrastive_cio_rb_vinc_lat12proj8_enlcrop")
SUPCON_DIR = os.path.join(BASE_IMG, "supcon_cio_rb_vinc_lat12proj8_enlcrop")
SC2_DIR    = os.path.join(BASE_IMG, "supcon_cio_rb_vinc_lat12proj8_enlcrop_sc2")

OUTPUT = "/net/projects/CLS/lding/gitcode/SubCellAE/slides_20260611.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # layout 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add navy title bar with white bold title, optional lighter subtitle row."""
    # Navy bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        0, 0, SW, TITLE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    bar.left = 0
    bar.top  = 0
    # Add left padding via XML margin
    txBody = bar.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Inches(0.25))))
        bodyPr.set('tIns', str(int(Inches(0.25))))

    if subtitle_text:
        sub_top = TITLE_H
        sub_h = Inches(0.42)
        sub_bar = slide.shapes.add_shape(1, 0, sub_top, SW, sub_h)
        sub_bar.fill.solid()
        sub_bar.fill.fore_color.rgb = RGBColor(0x2A, 0x5A, 0x9F)
        sub_bar.line.fill.background()
        stf = sub_bar.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run()
        sr.text = subtitle_text
        sr.font.size = Pt(16)
        sr.font.color.rgb = WHITE
        sr.font.name = "Calibri"
        sr.font.italic = True
        stxBody = sub_bar.text_frame._txBody
        sbodyPr = stxBody.find(qn('a:bodyPr'))
        if sbodyPr is not None:
            sbodyPr.set('lIns', str(int(Inches(0.25))))
            sbodyPr.set('tIns', str(int(Inches(0.10))))
        return sub_top + sub_h
    return TITLE_H


def content_top(has_subtitle=False):
    return TITLE_H + (Inches(0.42) if has_subtitle else 0) + Inches(0.15)


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=DARK_GRAY,
                font_name="Calibri", align=PP_ALIGN.LEFT,
                word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return txb


def add_bullet_para(tf, text, level=0, font_size=14, bold=False,
                    color=DARK_GRAY, font_name="Calibri"):
    p = tf.add_paragraph()
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return p


def add_bullet_textbox(slide, bullets, left, top, width, height,
                       font_size=14, color=DARK_GRAY, font_name="Calibri",
                       bullet_char="•"):
    """Add a textbox with bullet points. bullets is list of (text, level) tuples or plain strings."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, str):
            text, level = item, 0
        else:
            text, level = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        indent = "    " * level
        r = p.add_run()
        r.text = f"{bullet_char} {text}" if level == 0 else f"  {'–'} {text}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = font_name
    return txb


def add_numbered_textbox(slide, items, left, top, width, height,
                         font_size=14, color=DARK_GRAY, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items, 1):
        if i == 1:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{i}. {text}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = font_name
        # spacing after
        p.space_after = Pt(4)
    return txb


def add_section_header(slide, text, left, top, width, height=Inches(0.35),
                       color=TEAL):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    txBody = box.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Inches(0.10))))
        bodyPr.set('tIns', str(int(Inches(0.05))))
    return top + height


def safe_add_picture(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    else:
        # Placeholder box
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[Image not found]\n{os.path.basename(path)}"
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide01_title(prs):
    slide = blank_slide(prs)
    # Full navy background
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # White accent rectangle
    accent = slide.shapes.add_shape(1, Inches(0.5), Inches(2.8), Inches(12.33), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Title
    tb = add_textbox(slide, "EnlargedCropDataset & Contrastive AE",
                     Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.1),
                     font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, "Development Log — 2026-06-11",
                Inches(0.5), Inches(2.95), Inches(10.0), Inches(0.6),
                font_size=22, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.LEFT)

    # Author
    add_textbox(slide, "Liyading",
                Inches(0.5), Inches(3.7), Inches(5.0), Inches(0.5),
                font_size=18, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)


def slide02_motivation(prs):
    slide = blank_slide(prs)
    ct = add_title_bar(slide, "Why EnlargedCropDataset?")
    top = content_top()

    bullets = [
        "JitterCropDataset: scipy rotation per patch inside __getitem__ → slow DataLoader",
        "Two interpolation rounds: load-time rotation + augment_contrastive_view → degraded quality",
        "Goal: single GPU affine pass, independent per-view, zero numpy allocation per sample",
    ]
    add_bullet_textbox(slide, bullets,
                       MARGIN, top, SW - 2*MARGIN, SH - top - MARGIN,
                       font_size=17)


def slide03_design(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "EnlargedCropDataset Design")
    top = content_top()

    col_w = Inches(5.8)
    col_gap = Inches(0.3)
    right_left = MARGIN + col_w + col_gap
    right_w = SW - right_left - MARGIN

    # Left column bullets
    bullets = [
        "Context size: 2 × ⌈√2 × (ps/2 + max_shift)⌉",
        "For ps=32, shift=4 → context = 58 px",
        "Init: pad all frames with np.pad(mode='reflect')",
        "__getitem__: pure numpy slice → (1, 58, 58) tensor",
        "No rotation/jitter at load time",
    ]
    add_bullet_textbox(slide, bullets,
                       MARGIN, top, col_w, SH - top - MARGIN,
                       font_size=15)

    # Right column — diagram
    diagram_top = top + Inches(0.1)
    # Outer rectangle: 58×58 context
    outer_w = Inches(3.2)
    outer_h = Inches(3.2)
    outer_left = right_left + (right_w - outer_w) / 2
    outer = slide.shapes.add_shape(1, outer_left, diagram_top, outer_w, outer_h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = RGBColor(0xD6, 0xEA, 0xF8)
    outer.line.color.rgb = NAVY
    outer.line.width = Pt(1.5)
    add_textbox(slide, "58×58 context",
                outer_left, diagram_top + Inches(0.08),
                outer_w, Inches(0.35),
                font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Inner rectangle: 32×32 crop
    inner_w = Inches(1.8)
    inner_h = Inches(1.8)
    inner_left = outer_left + (outer_w - inner_w) / 2
    inner_top  = diagram_top + (outer_h - inner_h) / 2
    inner = slide.shapes.add_shape(1, inner_left, inner_top, inner_w, inner_h)
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor(0xA9, 0xCC, 0xE3)
    inner.line.color.rgb = TEAL
    inner.line.width = Pt(1.5)
    # dashed line via XML
    ln = inner.line._ln
    prstDash = etree.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    add_textbox(slide, "32×32 crop",
                inner_left, inner_top + Inches(0.65),
                inner_w, Inches(0.35),
                font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Labels below
    arrow_top = diagram_top + outer_h + Inches(0.15)
    add_textbox(slide, "→ view1 (random jitter+rot)",
                outer_left, arrow_top, outer_w / 2, Inches(0.5),
                font_size=11, color=TEAL, align=PP_ALIGN.LEFT)
    add_textbox(slide, "→ view2 (random jitter+rot)",
                outer_left + outer_w / 2, arrow_top, outer_w / 2, Inches(0.5),
                font_size=11, color=TEAL, align=PP_ALIGN.RIGHT)


def slide04_gpu_aug(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "_jitter_rot_crop: Batched GPU Affine")
    top = content_top()

    code = (
        "angles ~ Uniform(-15°, +15°)   dx,dy ~ randint(-4, 4)\n"
        "cos_a, sin_a = cos(angles), sin(angles)\n"
        "s = out_size / H                tx = 2·dx/W,  ty = 2·dy/H\n\n"
        "theta = [[cos_a·s,  sin_a·s,  tx],\n"
        "         [-sin_a·s, cos_a·s,  ty]]   # (B, 2, 3)\n\n"
        "grid  = F.affine_grid(theta, (B,C,32,32))\n"
        "out   = F.grid_sample(x, grid, mode='bilinear')"
    )

    code_h = Inches(2.5)
    code_box = slide.shapes.add_textbox(MARGIN, top, SW - 2*MARGIN, code_h)
    # Light gray fill
    spPr = code_box._element.find(qn('p:spPr'))
    # Add solid fill to code box shape
    fill_xml = (
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="F0F0F0"/>'
        '</a:solidFill>'
    )
    if spPr is not None:
        spPr.insert(0, etree.fromstring(fill_xml))

    tf = code_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Courier New"
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x10, 0x10, 0x60)

    bullet_top = top + code_h + Inches(0.2)
    bullets = [
        "One GPU kernel call for entire batch — no Python loop",
        "view1 and view2 get independent random angles/shifts",
        "padding_mode='border' handles boundary pixels",
    ]
    add_bullet_textbox(slide, bullets,
                       MARGIN, bullet_top, SW - 2*MARGIN, SH - bullet_top - MARGIN,
                       font_size=15)


def slide05_perf(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Performance & Logging Fixes")
    top = content_top()

    half_w = (SW - 3*MARGIN) / 2
    right_left = MARGIN + half_w + MARGIN

    # Left column header
    sec_top = add_section_header(slide, "np.pad Bottleneck",
                                 MARGIN, top, half_w)
    left_bullets = [
        "np.pad(full_frame, pad, 'reflect') called per patch per epoch",
        "First submit (934185/934186): 16+ min, zero epoch prints",
        "Fix: pre-pad all frames once at __init__",
        "→ __getitem__ becomes a zero-allocation slice",
    ]
    add_bullet_textbox(slide, left_bullets,
                       MARGIN, sec_top + Inches(0.1), half_w,
                       SH - sec_top - MARGIN - Inches(0.1),
                       font_size=14)

    # Right column header
    sec_top2 = add_section_header(slide, "SLURM stdout Buffering",
                                  right_left, top, half_w)
    right_bullets = [
        "Python print() fully buffered in non-TTY (~8 KB)",
        "Epoch loss prints never appeared in log file",
        "Fix: flush=True on all epoch print() calls",
        "logging.INFO does flush after each call (already correct)",
    ]
    add_bullet_textbox(slide, right_bullets,
                       right_left, sec_top2 + Inches(0.1), half_w,
                       SH - sec_top2 - MARGIN - Inches(0.1),
                       font_size=14)


def slide06_experiments(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Experiments — 2026-06-11")
    top = content_top()

    headers = ["Job", "Model", "Input", "Sigmoid", "Warmup", "Status", "Key Finding"]
    rows_data = [
        ["934303", "ConAE (enlcrop)", "raw", "yes", "200 ep", "Done",
         "Recon stuck at mean; warmup no benefit"],
        ["934304", "SupCon (enlcrop)", "raw", "yes", "0", "Done",
         "Best recon (max 0.67); class labels help"],
        ["934341", "SupCon (enlcrop)", "raw", "yes", "100 ep", "Done",
         "Similar to 934304; best ckpt ep118"],
        ["934342", "ConAE (sc2)", "÷2", "no", "0", "Done",
         "Mean collapse; recon constant"],
        ["934343", "SupCon (sc2)", "÷2", "no", "100 ep", "Done",
         "Full collapse; contrast stuck at max"],
        ["934376/77", "Both (sc2+Sigmoid)", "÷2", "yes", "0", "Running", "—"],
    ]

    n_rows = len(rows_data) + 1  # +1 for header
    n_cols = len(headers)
    tbl_left = MARGIN
    tbl_top  = top + Inches(0.05)
    tbl_w    = SW - 2*MARGIN
    tbl_h    = SH - tbl_top - MARGIN

    tbl = slide.shapes.add_table(n_rows, n_cols,
                                 tbl_left, tbl_top, tbl_w, tbl_h).table

    # Column widths
    col_widths = [Inches(1.1), Inches(2.0), Inches(0.7), Inches(0.8),
                  Inches(0.9), Inches(0.8), Inches(6.53)]
    # adjust last col to fill
    used = sum(col_widths[:-1])
    col_widths[-1] = tbl_w - used
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = int(w)

    def set_cell(cell, text, bold=False, bg=None, fg=WHITE, size=12, align=PP_ALIGN.LEFT):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = fg
        r.font.name = "Calibri"
        if bg is not None:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', f'{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}')

    # Header row
    for j, h in enumerate(headers):
        set_cell(tbl.cell(0, j), h, bold=True, bg=(0x1B, 0x3A, 0x6B),
                 fg=WHITE, size=12)

    # Data rows
    for i, row in enumerate(rows_data):
        bg = None if i % 2 == 0 else (0xE8, 0xF4, 0xFD)
        fg_color = DARK_GRAY if bg is None else DARK_GRAY
        for j, val in enumerate(row):
            set_cell(tbl.cell(i+1, j), val, bg=bg,
                     fg=DARK_GRAY, size=11)


def slide07_conae_loss(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "ConAE enlcrop (934303) — Loss Curve")
    top = content_top()

    img_path = os.path.join(CONAE_DIR, "contrastive_train_val_loss.png")
    img_h = Inches(4.3)
    img_w = Inches(7.5)
    img_left = (SW - img_w) / 2
    safe_add_picture(slide, img_path, img_left, top, img_w, img_h)

    note_top = top + img_h + Inches(0.12)
    notes = [
        "Warmup (ep 0–200): total loss ≈ recon only (0.007) — flat throughout",
        "Post-warmup (ep 200+): contrast kicks in, val diverges from train",
        "Best checkpoint saved at ep 200 (last warmup epoch) — broken by design",
    ]
    add_bullet_textbox(slide, notes,
                       MARGIN, note_top, SW - 2*MARGIN, SH - note_top - MARGIN,
                       font_size=13)


def slide08_conae_recon(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "ConAE enlcrop — Reconstruction @ ep300")
    top = content_top()

    half_w = (SW - 3*MARGIN) / 2
    right_left = MARGIN + half_w + MARGIN

    img_h = Inches(3.6)
    img_path_recon = os.path.join(CONAE_DIR, "contrastive_recon_ep300.png")
    img_path_views = os.path.join(CONAE_DIR, "contrastive_views_ep300.png")

    safe_add_picture(slide, img_path_recon, MARGIN, top, half_w, img_h)
    safe_add_picture(slide, img_path_views, right_left, top, half_w, img_h)

    cap_top = top + img_h + Inches(0.10)
    add_textbox(slide,
                "Recon — Row 0: input center-crop  |  Row 1: recon — near-blank, max ~0.33 vs input max ~1.2",
                MARGIN, cap_top, half_w, Inches(0.7),
                font_size=12, color=DARK_GRAY)
    add_textbox(slide,
                "Views — Row 0: recon of view1  |  Row 1: view1 (jitter+rot crop)  |  Row 2: view2",
                right_left, cap_top, half_w, Inches(0.7),
                font_size=12, color=DARK_GRAY)


def slide09_supcon_results(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "SupCon enlcrop (934304) — Best Reconstruction")
    top = content_top()

    half_w = (SW - 3*MARGIN) / 2
    right_left = MARGIN + half_w + MARGIN
    img_h = Inches(4.0)

    img_recon = os.path.join(SUPCON_DIR, "supcon_recon_ep300.png")
    img_loss  = os.path.join(SUPCON_DIR, "supcon_train_val_loss.png")

    safe_add_picture(slide, img_recon, MARGIN, top, half_w, img_h)
    safe_add_picture(slide, img_loss, right_left, top, half_w, img_h)

    cap_top = top + img_h + Inches(0.12)
    add_textbox(slide,
                "SupCon (no warmup): recon max 0.67, visible FA structure — "
                "class labels drive better encoder",
                MARGIN, cap_top, SW - 2*MARGIN, Inches(0.6),
                font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)


def slide10_bugs(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Bugs Found & Fixed")
    top = content_top()

    half_w = (SW - 3*MARGIN) / 2
    right_left = MARGIN + half_w + MARGIN

    # Section A
    sec_a_top = add_section_header(slide, "Best-Checkpoint Selection Bug",
                                   MARGIN, top, half_w)
    bugs_a = [
        "min_epochs_for_best=200 → tracking started at ep200 (last warmup epoch)",
        "Warmup val loss (recon only, ~0.008) always beat post-warmup (recon+contrast, ~2.1)",
        "model_best.pt = last warmup epoch, before any contrastive training",
        "Fix: past_warmup = (epoch+1) > warmup_epochs required before tracking",
    ]
    add_bullet_textbox(slide, bugs_a,
                       MARGIN, sec_a_top + Inches(0.1), half_w,
                       SH - sec_a_top - MARGIN - Inches(0.1),
                       font_size=13)

    # Section B
    sec_b_top = add_section_header(slide, "SupCon Warmup Missing",
                                   right_left, top, half_w)
    bugs_b = [
        "train_supervised_contrastive_ae had no warmup support at all",
        "Added: warmup_epochs, eff_lambda_contrast, LR reset at transition, model_best.pt",
    ]
    add_bullet_textbox(slide, bugs_b,
                       right_left, sec_b_top + Inches(0.1), half_w,
                       SH - sec_b_top - MARGIN - Inches(0.1),
                       font_size=13)


def slide11_sc2(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "sc2: Input ÷2 + No Sigmoid (934342/343)")
    top = content_top()

    col_w = Inches(4.8)
    right_left = MARGIN + col_w + MARGIN
    right_w = SW - right_left - MARGIN

    # Left — Motivation
    sec_top = add_section_header(slide, "Motivation", MARGIN, top, col_w)
    left_bullets = [
        "FA pixels reach ~1.2, Sigmoid clips at 1.0 → 20% of bright pixels clipped",
        "Remove Sigmoid → decoder unconstrained, free to match any amplitude",
        "Divide input by 2 → range [0, 0.6] fits within (0,1)",
    ]
    add_bullet_textbox(slide, left_bullets,
                       MARGIN, sec_top + Inches(0.1), col_w,
                       Inches(2.0),
                       font_size=13)

    # Right — What happened
    sec_top2 = add_section_header(slide, "What happened", right_left, top, right_w)
    right_bullets = [
        "Without Sigmoid, decoder immediately finds MSE minimum = dataset mean constant",
        "SupCon: 100-ep warmup collapsed encoder before contrast could help",
        "Contrast stuck at ln(256)=5.53 — random-init ceiling for batch=128",
    ]
    add_bullet_textbox(slide, right_bullets,
                       right_left, sec_top2 + Inches(0.1), right_w,
                       Inches(2.0),
                       font_size=13)

    # Image bottom right
    img_path = os.path.join(SC2_DIR, "supcon_recon_ep300.png")
    img_top  = top + Inches(2.5)
    img_h    = Inches(2.2)
    safe_add_picture(slide, img_path, right_left, img_top, right_w, img_h)
    add_textbox(slide, "SupCon sc2: constant output (min=max=mean=0.05)",
                right_left, img_top + img_h + Inches(0.05), right_w, Inches(0.4),
                font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


def slide12_collapse(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Why Removing Sigmoid Caused Collapse")
    top = content_top()

    bullets = [
        "MSE optimal predictor for any distribution = conditional mean E[y|z]",
        "For sparse FA patches (~85% background), mean ≈ 0.05 → constant output minimizes MSE",
        "Without Sigmoid: no nonlinear friction; gradient drives decoder to constant in < 10 epochs",
        "With Sigmoid: slope near equilibrium is small but non-zero → slows collapse, contrastive has time to act",
        "Class labels (SupCon) help decoder learn structure — but only if encoder isn't already dead",
    ]
    add_bullet_textbox(slide, bullets,
                       MARGIN, top, SW - 2*MARGIN, Inches(3.2),
                       font_size=15)

    diag_top = top + Inches(3.35)
    diag_box = slide.shapes.add_textbox(MARGIN, diag_top,
                                        SW - 2*MARGIN, Inches(1.8))
    diag_box_sp = diag_box._element.find(qn('p:spPr'))
    fill_xml = (
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="EAF4FB"/>'
        '</a:solidFill>'
    )
    if diag_box_sp is not None:
        diag_box_sp.insert(0, etree.fromstring(fill_xml))
    tf = diag_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Sigmoid output landscape:   flat near constant → slight gradient → learning possible"
    r1.font.size = Pt(13)
    r1.font.name = "Courier New"
    r1.font.color.rgb = RGBColor(0x10, 0x50, 0x10)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Linear output landscape:    constant minimized instantly → encoder receives no input-dependent gradient"
    r2.font.size = Pt(13)
    r2.font.name = "Courier New"
    r2.font.color.rgb = RGBColor(0x80, 0x10, 0x10)


def slide13_fix(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Fix: sc2 + Sigmoid (934376 / 934377)")
    top = content_top()

    bullets = [
        "Restore output_sigmoid: true — Sigmoid prevents collapse, provides gradient friction",
        "input_divisor: 2.0 retained — input [0, 0.6] fits inside (0, 1) cleanly, no clipping at 1.2",
        "logit(0.6) ≈ 0.43 — decoder pre-sigmoid only needs small positive activation to reach FA peaks",
        "SupCon warmup_epochs: 0 — class labels drive encoder from ep1 (matches best 934304 result)",
    ]
    add_bullet_textbox(slide, bullets,
                       MARGIN, top, SW - 2*MARGIN, Inches(3.2),
                       font_size=16)

    # Status box
    status_top = top + Inches(3.4)
    status_box = slide.shapes.add_shape(1, MARGIN, status_top,
                                        SW - 2*MARGIN, Inches(0.85))
    status_box.fill.solid()
    status_box.fill.fore_color.rgb = TEAL
    status_box.line.fill.background()
    stf = status_box.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "Jobs 934376 (ConAE) and 934377 (SupCon) running — results tomorrow"
    sr.font.size = Pt(18)
    sr.font.bold = True
    sr.font.color.rgb = WHITE
    sr.font.name = "Calibri"
    stxBody = stf._txBody
    sbodyPr = stxBody.find(qn('a:bodyPr'))
    if sbodyPr is not None:
        sbodyPr.set('tIns', str(int(Inches(0.18))))


def slide14_next(prs):
    slide = blank_slide(prs)
    add_title_bar(slide, "Next Steps")
    top = content_top()

    items = [
        "Check 934376/934377: expect recon amplitude up to ~0.6, visible FA structure in SupCon",
        "If sc2+Sigmoid improves recon → run full analysis pipeline (latent CSV, UMAP, cls eval)",
        "Contrastive overfitting: all runs show val > train on contrast loss — may need larger batch or more augmentation",
        "nonad-vs-ad branch: port JitterCropDataset pre-padding fix; fix num_workers=0 in ordered DataLoaders",
        "Consider SSIM or normalized-MSE loss if sharper reconstruction is needed",
    ]
    add_numbered_textbox(slide, items,
                         MARGIN, top, SW - 2*MARGIN, SH - top - MARGIN,
                         font_size=16)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    print("Building slide 1 — Title ...")
    slide01_title(prs)
    print("Building slide 2 — Motivation ...")
    slide02_motivation(prs)
    print("Building slide 3 — Design ...")
    slide03_design(prs)
    print("Building slide 4 — GPU Augmentation ...")
    slide04_gpu_aug(prs)
    print("Building slide 5 — Performance Fixes ...")
    slide05_perf(prs)
    print("Building slide 6 — Experiments Overview ...")
    slide06_experiments(prs)
    print("Building slide 7 — ConAE loss curve ...")
    slide07_conae_loss(prs)
    print("Building slide 8 — ConAE reconstruction ...")
    slide08_conae_recon(prs)
    print("Building slide 9 — SupCon results ...")
    slide09_supcon_results(prs)
    print("Building slide 10 — Bugs ...")
    slide10_bugs(prs)
    print("Building slide 11 — sc2 experiment ...")
    slide11_sc2(prs)
    print("Building slide 12 — Mean collapse explanation ...")
    slide12_collapse(prs)
    print("Building slide 13 — Fix & status ...")
    slide13_fix(prs)
    print("Building slide 14 — Next steps ...")
    slide14_next(prs)

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f"File size: {size_kb:.1f} KB")


if __name__ == "__main__":
    main()
