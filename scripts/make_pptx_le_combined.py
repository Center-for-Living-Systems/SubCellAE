#!/usr/bin/env python3
"""
make_pptx_le_combined.py

Comprehensive slide deck for label × image efficiency experiments:
  - 3 train/test splits (s1v3, s2v2, s3v1)
  - SupCon AE + LGBM classifier
  - 5 npi levels × N_images × 3 series = 225 runs per split

Usage (from repo root):
  python scripts/make_pptx_le_combined.py
  python scripts/make_pptx_le_combined.py --out results/my_out.pptx
"""
from __future__ import annotations

import argparse
import io
import re
import glob
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.gridspec as gridspec
import numpy as np
import pandas as pd
import tifffile
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT  = Path(__file__).resolve().parents[1]
DATA_ROOT  = Path("/net/projects/CLS/lding/data/fa_data_analysis")
PATCH_DIR  = DATA_ROOT / "ae_results/patches/cio/vinc/control/tiff_patches32_mr10"
FRAME_DIR  = DATA_ROOT / "ae_results/source_frames/cio_mode_prt/vinc/control"
ANN_DIR    = DATA_ROOT / "labelling/le_combined_s2v2"
RESULTS    = REPO_ROOT / "results"

# Pre-generated PNGs
PNG = {
    "s1v3_heatmap": RESULTS / "le_combined_heatmap.png",
    "s1v3_curves":  RESULTS / "le_combined_curves.png",
    "s2v2_heatmap": RESULTS / "le_combined_s2v2_heatmap.png",
    "s2v2_curves":  RESULTS / "le_combined_s2v2_curves.png",
    "s3v1_heatmap": RESULTS / "le_combined_s3v1_heatmap.png",
    "s3v1_curves":  RESULTS / "le_combined_s3v1_curves.png",
}

# ---------------------------------------------------------------------------
# Slide dimensions
# ---------------------------------------------------------------------------
SW = Inches(13.33)
SH = Inches(7.5)

# ---------------------------------------------------------------------------
# Colours
# ---------------------------------------------------------------------------
C_TITLE  = RGBColor(0x1A, 0x1A, 0x2E)
C_HEAD   = RGBColor(0x16, 0x21, 0x3E)
C_BODY   = RGBColor(0x1A, 0x1A, 0x1A)
C_GREY   = RGBColor(0x66, 0x66, 0x66)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_ACCENT = RGBColor(0x0F, 0x3D, 0x79)
C_GOOD   = RGBColor(0x1A, 0x6B, 0x30)
C_WARN   = RGBColor(0x8B, 0x45, 0x00)

# Patch group colours (matplotlib hex)
COL_NPI10  = "#1f77b4"   # blue
COL_NPI25  = "#2ca02c"   # green
COL_NPI50  = "#ff7f0e"   # orange
COL_NPI75  = "#d62728"   # red

# RGB tuples for drawing on numpy arrays
COL_NPI10_RGB  = (31,  119, 180)
COL_NPI25_RGB  = (44,  160, 44)
COL_NPI50_RGB  = (255, 127, 14)
COL_NPI75_RGB  = (214, 39,  40)

# ---------------------------------------------------------------------------
# Core slide helpers  (same pattern as make_pptx_noad_vs_ad_story.py)
# ---------------------------------------------------------------------------

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _txt(slide, text: str, left, top, width, height, *,
         bold=False, italic=False, size_pt=13, color=C_BODY,
         align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.size      = Pt(size_pt)
    run.font.color.rgb = color
    return txb


def _rule(slide, top, width=None, left=None, thickness_pt=0.75):
    w = width or SW - Inches(1.0)
    l = left  or Inches(0.5)
    ln = slide.shapes.add_connector(1, l, top, l + w, top)
    ln.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.width = Pt(thickness_pt)


def _slide_header(slide, title: str, subtitle: str = ""):
    _txt(slide, title,
         Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.55),
         bold=True, size_pt=20, color=C_HEAD)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.35),
             size_pt=11, color=C_GREY)
    _rule(slide, Inches(0.97))


def _img_or_ph(slide, path, left, top, width, height, label="[pending]"):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf  = box.text_frame
        tf.paragraphs[0].add_run().text = label
        tf.paragraphs[0].runs[0].font.size = Pt(9)
        tf.paragraphs[0].runs[0].font.color.rgb = C_GREY


def _fig_to_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_fig(slide, fig, left, top, width, height):
    buf = _fig_to_buf(fig)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def _norm_frame(arr: np.ndarray) -> np.ndarray:
    """Normalize float32 frame: clip 1st-99th percentile → [0, 1]."""
    lo, hi = np.percentile(arr, 1), np.percentile(arr, 99)
    if hi <= lo:
        return np.zeros_like(arr, dtype=np.float32)
    return np.clip((arr - lo) / (hi - lo), 0.0, 1.0).astype(np.float32)


def _norm_patch(p: np.ndarray) -> np.ndarray:
    """Normalize float32 patch to [0, 1]."""
    return np.clip((p - p.min()) / (p.max() - p.min() + 1e-8), 0, 1).astype(np.float32)


def _load_frame(frame_idx: int) -> np.ndarray:
    path = FRAME_DIR / f"control_f{frame_idx:04d}_vinc.tif"
    arr  = tifffile.imread(str(path)).astype(np.float32)
    return _norm_frame(arr)


def _load_patch(filename: str) -> np.ndarray:
    path = PATCH_DIR / filename
    arr  = tifffile.imread(str(path)).astype(np.float32)
    return _norm_patch(arr)


def _parse_patch_coords(filename: str):
    """Return (frame, x_col, y_row) from patch filename."""
    m = re.search(r"_f(\d+)x(\d+)y(\d+)ps32", filename)
    return int(m.group(1)), int(m.group(2)), int(m.group(3))


def _draw_marker_on_frame(frame_img: np.ndarray, x_col: int, y_row: int,
                           rgb: tuple, thumb_size: int = 256,
                           marker_px: int = 6) -> None:
    """Draw a small filled square on `frame_img` (modified in-place).
    frame_img is a (H, W, 3) uint8 array at thumb_size resolution.
    Coordinates are in original 1024×1024 space.
    """
    scale = thumb_size / 1024.0
    cx = int(x_col * scale) + int(16 * scale)   # centre of 32-px patch
    cy = int(y_row * scale) + int(16 * scale)
    half = marker_px // 2
    r0, r1 = max(0, cy - half), min(thumb_size, cy + half)
    c0, c1 = max(0, cx - half), min(thumb_size, cx + half)
    frame_img[r0:r1, c0:c1] = rgb


# ---------------------------------------------------------------------------
# Load annotation data (cached at module level on first use)
# ---------------------------------------------------------------------------

def _load_ann(npi: int, series: int = 0) -> pd.DataFrame:
    path = ANN_DIR / f"le_comb_s2_npi{npi}_r{series}.csv"
    return pd.read_csv(str(path))


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def _slide_title(prs: Presentation):
    slide = _blank(prs)
    _txt(slide,
         "Label × Image Efficiency for Focal Adhesion Detection",
         Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2),
         bold=True, size_pt=32, color=C_HEAD, align=PP_ALIGN.CENTER)
    _txt(slide,
         "SupCon AE + LGBM  ·  3 train-test splits  ·  225 runs per split",
         Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.6),
         size_pt=16, color=C_GREY, align=PP_ALIGN.CENTER)
    _rule(slide, Inches(4.55), width=Inches(9), left=Inches(2.17))
    _txt(slide,
         "vinc control dataset  ·  32×32 patches  ·  npi ∈ {10, 25, 50, 75, 100}  ·  N_images ∈ 1-49",
         Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
         size_pt=12, color=C_GREY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — Experiment Design Overview
# ---------------------------------------------------------------------------

def _slide_design(prs: Presentation):
    slide = _blank(prs)
    _slide_header(slide, "Experiment Design Overview",
                  "3 train-test splits · SupCon AE trains on N images · LGBM trains on npi labeled patches")

    # Pipeline description
    pipeline_txt = (
        "Pipeline: Patches  →  SupCon AE (trained on N total images)  "
        "→  latent features  →  LGBM classifier (trained on npi labeled patches/frame)  "
        "→  Balanced Accuracy on held-out test frame(s)"
    )
    _txt(slide, pipeline_txt,
         Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.5),
         size_pt=12, color=C_BODY)

    # Split table via matplotlib
    fig, ax = plt.subplots(figsize=(12, 3.8), facecolor="white")
    ax.axis("off")

    headers = ["Split", "AE Training Frames", "LGBM Labels From", "Test Frames", "N_images range", "Runs"]
    rows = [
        ["s1v3", "frame 0 only", "frame 0", "frames 1, 2, 3", "1 – 47", "225"],
        ["s2v2", "frames 0 + 1", "frames 0, 1", "frames 2, 3", "2 – 48", "225"],
        ["s3v1", "frames 0 + 1 + 2", "frames 0, 1, 2", "frame 3", "3 – 49", "225"],
    ]
    col_widths = [0.10, 0.20, 0.20, 0.18, 0.18, 0.10]

    # Draw header
    x = 0.01
    y = 0.88
    for w, h in zip(col_widths, headers):
        ax.text(x + w/2, y, h, ha="center", va="center", fontsize=11,
                fontweight="bold", color="#16213E",
                transform=ax.transAxes)
        x += w

    ax.axhline(y - 0.06, xmin=0.01, xmax=0.99, color="#CCCCCC", lw=1.2)

    row_colors = ["#EEF2FF", "#FFFFFF", "#EEF2FF"]
    for ri, (row, bg) in enumerate(zip(rows, row_colors)):
        y -= 0.26
        x = 0.01
        ax.axhspan(y - 0.10, y + 0.10, xmin=0.005, xmax=0.995, color=bg, alpha=0.7,
                   transform=ax.transAxes)
        for w, cell in zip(col_widths, row):
            fw = "bold" if w == col_widths[0] else "normal"
            fc = "#16213E" if w == col_widths[0] else "#1A1A1A"
            ax.text(x + w/2, y, cell, ha="center", va="center", fontsize=10,
                    fontweight=fw, color=fc, transform=ax.transAxes)
            x += w

    # npi legend row
    ax.axhline(0.05, xmin=0.01, xmax=0.99, color="#CCCCCC", lw=0.8)
    ax.text(0.01, 0.02,
            "npi levels: 10 · 25 · 50 · 75 · 100  (labels per frame, cumulative; 3 random series each)",
            ha="left", va="bottom", fontsize=9, color="#666666", transform=ax.transAxes)

    fig.tight_layout(pad=0.2)
    _add_fig(slide, fig, Inches(0.4), Inches(1.7), Inches(12.5), Inches(3.6))

    # Bottom text
    details = (
        "Each run: (n_images, npi, series) triple.  "
        "Series 0/1/2 = 3 independent random label samples.  "
        "Metric: Balanced Accuracy (adh recall + no-adh recall) / 2."
    )
    _txt(slide, details,
         Inches(0.55), Inches(5.5), Inches(12.2), Inches(0.6),
         size_pt=11, color=C_GREY)

    extra_info = (
        "Extra unlabeled images (beyond the 2 labeled frames in s2v2) come from frames 4, 5, 6, ...  "
        "— they provide more AE training diversity without additional annotation cost."
    )
    _txt(slide, extra_info,
         Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.6),
         size_pt=11, color=C_GREY)


# ---------------------------------------------------------------------------
# Slides 3–5 — Label Sampling Visualization
# ---------------------------------------------------------------------------

# s2v2 frame lists for selected N values
N_FRAMES = {
    3:  [0, 1, 4],
    6:  [0, 1, 4, 5, 6, 7],
    11: [0, 1, 4, 5, 6, 7, 8, 9, 10, 11, 12],
}

LABELED_FRAMES = [0, 1]   # always for s2v2


def _build_label_vis_fig(n_val: int) -> plt.Figure:
    """Build the combined top+bottom matplotlib figure for a label vis slide."""
    frames_list = N_FRAMES[n_val]
    n_frames    = len(frames_list)

    # Load annotations (series 0)
    df10  = _load_ann(10,  series=0)
    df25  = _load_ann(25,  series=0)
    df50  = _load_ann(50,  series=0)
    df75  = _load_ann(75,  series=0)

    uid2label75 = dict(zip(df75["unique_ID"], df75["label"]))  # superset for label lookup

    def _frame_ids(df, frame):
        return set(df[df["frame"] == frame]["unique_ID"])

    # Cumulative extras for frame 0
    f0_ids10  = _frame_ids(df10, 0)
    f0_ids25  = _frame_ids(df25, 0)
    f0_ids50  = _frame_ids(df50, 0)
    f0_ids75  = _frame_ids(df75, 0)
    f0_extra25 = f0_ids25 - f0_ids10
    f0_extra50 = f0_ids50 - f0_ids25
    f0_extra75 = f0_ids75 - f0_ids50

    # Cumulative extras for frame 1
    f1_ids10  = _frame_ids(df10, 1)
    f1_ids25  = _frame_ids(df25, 1)
    f1_ids50  = _frame_ids(df50, 1)
    f1_ids75  = _frame_ids(df75, 1)
    f1_extra25 = f1_ids25 - f1_ids10
    f1_extra50 = f1_ids50 - f1_ids25
    f1_extra75 = f1_ids75 - f1_ids50

    # -----------------------------------------------------------------------
    # Figure layout: top = frame thumbnails, middle = frame 0 patches, bottom = frame 1 patches
    # -----------------------------------------------------------------------
    THUMB = 160  # pixels for square frame thumbnail (displayed)
    MAX_SHOW_PATCH = 15

    def _make_groups(ids10, extra25, extra50, extra75):
        return [
            ("npi=10\n(10 patches)",         sorted(ids10),                          10,             COL_NPI10),
            ("extra npi=25\n(+15 patches)",  sorted(extra25),                        15,             COL_NPI25),
            ("extra npi=50\n(+25, show 15)", sorted(list(extra50))[:MAX_SHOW_PATCH], MAX_SHOW_PATCH, COL_NPI50),
            ("extra npi=75\n(+25, show 15)", sorted(list(extra75))[:MAX_SHOW_PATCH], MAX_SHOW_PATCH, COL_NPI75),
        ]

    groups_f0 = _make_groups(f0_ids10, f0_extra25, f0_extra50, f0_extra75)
    groups_f1 = _make_groups(f1_ids10, f1_extra25, f1_extra50, f1_extra75)

    # Build figure (taller to accommodate two patch rows)
    fig_w = 15.0
    fig_h = 11.0
    fig = plt.figure(figsize=(fig_w, fig_h), facecolor="white")

    gs_top  = gridspec.GridSpec(1, n_frames,
                                top=0.97, bottom=0.70,
                                left=0.02, right=0.98, wspace=0.08)
    gs_f0   = gridspec.GridSpec(1, 4,
                                top=0.63, bottom=0.37,
                                left=0.02, right=0.98, wspace=0.05)
    gs_f1   = gridspec.GridSpec(1, 4,
                                top=0.30, bottom=0.02,
                                left=0.02, right=0.98, wspace=0.05)

    # ---- TOP: frame thumbnails ----
    # Markers for frame 0: use frame-0 id sets
    F0_MARKER_LAYERS = [
        (f0_ids10,   COL_NPI10_RGB),
        (f0_extra25, COL_NPI25_RGB),
        (f0_extra50, COL_NPI50_RGB),
        (f0_extra75, COL_NPI75_RGB),
    ]
    # Markers for frame 1: use frame-1 id sets
    F1_MARKER_LAYERS = [
        (f1_ids10,   COL_NPI10_RGB),
        (f1_extra25, COL_NPI25_RGB),
        (f1_extra50, COL_NPI50_RGB),
        (f1_extra75, COL_NPI75_RGB),
    ]
    MARKER_PX_THUMB = max(4, THUMB // 50)

    for fi, frame_idx in enumerate(frames_list):
        ax_t = fig.add_subplot(gs_top[0, fi])

        frame_arr = _load_frame(frame_idx)
        # Convert to uint8 RGB for marker drawing
        frame_rgb = (frame_arr * 255).astype(np.uint8)
        frame_rgb = np.stack([frame_rgb, frame_rgb, frame_rgb], axis=-1)
        # Resize to THUMB×THUMB
        frame_pil  = Image.fromarray(frame_rgb).resize((THUMB, THUMB), Image.BILINEAR)
        frame_thumb = np.array(frame_pil)

        if frame_idx in LABELED_FRAMES:
            # Use per-frame marker layers
            marker_layers = F0_MARKER_LAYERS if frame_idx == 0 else F1_MARKER_LAYERS
            for npi_ids, rgb in marker_layers:
                layer_df = df75[
                    (df75["frame"] == frame_idx) &
                    (df75["unique_ID"].isin(npi_ids))
                ]
                for _, row in layer_df.iterrows():
                    _, x_col, y_row = _parse_patch_coords(row["filename"])
                    _draw_marker_on_frame(frame_thumb, x_col, y_row,
                                          rgb, thumb_size=THUMB,
                                          marker_px=MARKER_PX_THUMB)

        ax_t.imshow(frame_thumb, cmap=None, aspect="equal")
        ax_t.axis("off")

        if frame_idx in LABELED_FRAMES:
            label_str = f"f{frame_idx}\n(labeled)"
            fc = "#EEF6EE"
            ec = "#2ca02c"
        else:
            label_str = f"f{frame_idx}\n(unlabeled)"
            fc = "#F5F5F5"
            ec = "#AAAAAA"

        ax_t.set_title(label_str, fontsize=9, color="#1A1A1A",
                       pad=3, fontweight="bold" if frame_idx in LABELED_FRAMES else "normal",
                       bbox=dict(boxstyle="round,pad=0.2", fc=fc, ec=ec, lw=0.8))

    # Legend for markers (top section)
    legend_elements = [
        mpatches.Patch(facecolor=COL_NPI10,  edgecolor="grey", lw=0.5, label="npi=10 patches"),
        mpatches.Patch(facecolor=COL_NPI25,  edgecolor="grey", lw=0.5, label="extra npi=25"),
        mpatches.Patch(facecolor=COL_NPI50,  edgecolor="grey", lw=0.5, label="extra npi=50"),
        mpatches.Patch(facecolor=COL_NPI75,  edgecolor="grey", lw=0.5, label="extra npi=75"),
    ]
    # Place legend in top-right of figure
    fig.legend(handles=legend_elements, loc="upper right",
               bbox_to_anchor=(0.99, 0.99),
               fontsize=8, framealpha=0.9,
               title="Labeled patch layer", title_fontsize=8)

    # ---- BOTTOM: patch grids per npi group (frame 0 then frame 1) ----
    PATCH_SZ = 32
    GRID_COLS = 15

    COLOR_BORDER_MAP = {
        COL_NPI10: COL_NPI10_RGB, COL_NPI25: COL_NPI25_RGB,
        COL_NPI50: COL_NPI50_RGB, COL_NPI75: COL_NPI75_RGB,
    }

    def _render_patch_group(gs, row_idx, groups):
        for gi, (group_title, uid_list, n_show, color) in enumerate(groups):
            ax_b = fig.add_subplot(gs[row_idx, gi])
            ax_b.axis("off")

            patches_to_show = uid_list[:n_show]
            n_patches = len(patches_to_show)

            if n_patches == 0:
                ax_b.text(0.5, 0.5, "(none)", ha="center", va="center",
                          transform=ax_b.transAxes, fontsize=9, color="#888888")
                ax_b.set_title(group_title, fontsize=9, color=color, fontweight="bold")
                continue

            n_cols = min(n_patches, GRID_COLS)
            n_rows = int(np.ceil(n_patches / n_cols))
            canvas_h = n_rows * (PATCH_SZ + 2) + 2
            canvas_w = n_cols * (PATCH_SZ + 2) + 2
            canvas = np.ones((canvas_h, canvas_w, 3), dtype=np.float32)

            adh_count, noadh_count = 0, 0
            for idx, uid in enumerate(patches_to_show):
                ri, ci = idx // n_cols, idx % n_cols
                row_off = 1 + ri * (PATCH_SZ + 2)
                col_off = 1 + ci * (PATCH_SZ + 2)
                fname = uid.replace("control-f", "control_f")
                p = _load_patch(fname)
                p_rgb = np.stack([p, p, p], axis=-1)
                canvas[row_off:row_off+PATCH_SZ, col_off:col_off+PATCH_SZ] = p_rgb
                r0, r1 = row_off - 1, row_off + PATCH_SZ
                c0, c1 = col_off - 1, col_off + PATCH_SZ
                br = tuple(v / 255.0 for v in COLOR_BORDER_MAP[color])
                canvas[r0:r1+1, c0, :] = br
                canvas[r0:r1+1, c1, :] = br
                canvas[r0, c0:c1+1, :] = br
                canvas[r1, c0:c1+1, :] = br
                lbl = uid2label75.get(uid, "")
                if lbl == "adhesion":
                    adh_count += 1
                elif lbl == "No adhesion":
                    noadh_count += 1

            ax_b.imshow(canvas, aspect="equal", interpolation="nearest")
            ax_b.set_title(group_title, fontsize=9, color=color, fontweight="bold", pad=4)
            ax_b.text(0.5, -0.05, f"Adh: {adh_count}  No-adh: {noadh_count}",
                      ha="center", va="top", fontsize=8, color="#1A1A1A",
                      transform=ax_b.transAxes)
            for spine in ax_b.spines.values():
                spine.set_visible(True)
                spine.set_edgecolor(color)
                spine.set_linewidth(2.0)

    _render_patch_group(gs_f0, 0, groups_f0)
    _render_patch_group(gs_f1, 0, groups_f1)

    # Section labels
    fig.text(0.5, 0.99, f"Frame thumbnails (N={n_val} images, s2v2 series 0)  ·  colored squares = labeled patches",
             ha="center", va="top", fontsize=10, color="#16213E", fontweight="bold")
    fig.text(0.5, 0.645, "Frame 0 — labeled patch grid  ·  cumulative npi levels",
             ha="center", va="top", fontsize=10, color="#16213E", fontweight="bold")
    fig.text(0.5, 0.315, "Frame 1 — labeled patch grid  ·  cumulative npi levels",
             ha="center", va="top", fontsize=10, color="#16213E", fontweight="bold")

    fig.patch.set_facecolor("white")
    return fig


def _slides_label_vis(prs: Presentation):
    for n_val in [3, 6, 11]:
        slide = _blank(prs)
        _slide_header(
            slide,
            f"s2v2 — Labeling Strategy · N={n_val} images (series 0)",
            "Top: frame thumbnails with labeled patch locations (colored squares) · "
            "Bottom: patch grid per npi level from frame 0"
        )
        fig = _build_label_vis_fig(n_val)
        _add_fig(slide, fig, Inches(0.15), Inches(1.05), Inches(13.0), Inches(6.3))


# ---------------------------------------------------------------------------
# Slides 6–8 — Results (one per split)
# ---------------------------------------------------------------------------

SPLIT_INFO = {
    "s1v3": {
        "title":    "s1v3 Results — Train: frame 0 only · Test: frames 1+2+3",
        "subtitle": "225 runs (1–47 images × 5 npi × 3 series) · metric = Balanced Accuracy",
        "csv":      RESULTS / "le_combined_results.csv",
        "heatmap":  PNG["s1v3_heatmap"],
    },
    "s2v2": {
        "title":    "s2v2 Results — Train: frames 0+1 · Test: frames 2+3",
        "subtitle": "225 runs (2–48 images × 5 npi × 3 series) · metric = Balanced Accuracy",
        "csv":      RESULTS / "le_combined_s2v2_results.csv",
        "heatmap":  PNG["s2v2_heatmap"],
    },
    "s3v1": {
        "title":    "s3v1 Results — Train: frames 0+1+2 · Test: frame 3",
        "subtitle": "225 runs (3–49 images × 5 npi × 3 series) · metric = Balanced Accuracy",
        "csv":      RESULTS / "le_combined_s3v1_results.csv",
        "heatmap":  PNG["s3v1_heatmap"],
    },
}

NPI_COLORS = {10: "#d62728", 25: "#ff7f0e", 50: "#2ca02c", 75: "#1f77b4", 100: "#9467bd"}
SERIES_STYLES = {0: "-", 1: "--", 2: ":"}
SERIES_MARKERS = {0: "o", 1: "s", 2: "^"}


def _fig_curves_per_series(df: pd.DataFrame, split: str) -> plt.Figure:
    """BAcc vs N_images: one line per (npi, series) pair, colored by npi."""
    npi_levels = sorted(df["npi"].unique())
    n_values   = sorted(df["n_images"].unique())
    x = np.log2(n_values)

    fig, ax = plt.subplots(figsize=(7, 5.5), facecolor="white")

    for npi in npi_levels:
        color = NPI_COLORS.get(npi, "#333333")
        series_list = sorted(df["series"].unique())
        for si, series in enumerate(series_list):
            sub = (df[(df["npi"] == npi) & (df["series"] == series)]
                   .set_index("n_images")
                   .reindex(n_values))
            y = sub["balanced_acc"].values * 100
            ls = SERIES_STYLES.get(si, "-")
            mk = SERIES_MARKERS.get(si, "o")
            label = f"npi={npi} r{series}" if si == 0 else None
            ax.plot(x, y, color=color, linestyle=ls, marker=mk,
                    markersize=4, linewidth=1.4, alpha=0.85, label=label)

    # Legend: one entry per npi (color) + one entry per series (linestyle)
    npi_handles = [
        plt.Line2D([0], [0], color=NPI_COLORS.get(n, "#333"), linewidth=2, label=f"npi={n}")
        for n in npi_levels
    ]
    series_handles = [
        plt.Line2D([0], [0], color="#555", linestyle=SERIES_STYLES[s],
                   marker=SERIES_MARKERS[s], markersize=4, linewidth=1.4, label=f"series {s}")
        for s in sorted(df["series"].unique())
    ]
    leg1 = ax.legend(handles=npi_handles, loc="lower right", fontsize=8,
                     title="npi level", title_fontsize=8, framealpha=0.9)
    ax.add_artist(leg1)
    ax.legend(handles=series_handles, loc="upper left", fontsize=8,
              title="series", title_fontsize=8, framealpha=0.9)

    ax.set_xticks(x)
    ax.set_xticklabels([str(n) for n in n_values], fontsize=7, rotation=45, ha="right")
    ax.set_xlabel("N images in AE training (log₂ scale)", fontsize=10)
    ax.set_ylabel("Balanced accuracy (%)", fontsize=10)
    ax.set_ylim(20, 108)
    ax.axhline(90, color="#CCCCCC", linestyle="--", linewidth=0.8, zorder=0)
    ax.set_facecolor("white")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title(f"{split} — BAcc per series\nSupCon AE + LGBM · each line = one (npi, series) run",
                 fontsize=9, fontweight="bold")
    fig.tight_layout()
    return fig


def _slides_results(prs: Presentation):
    for split in ["s1v3", "s2v2", "s3v1"]:
        info  = SPLIT_INFO[split]
        slide = _blank(prs)
        _slide_header(slide, info["title"], info["subtitle"])

        top   = Inches(1.1)
        h     = Inches(6.2)
        w     = Inches(6.4)
        gap   = Inches(0.3)
        left1 = Inches(0.3)
        left2 = left1 + w + gap

        # Left: pre-generated heatmap
        _img_or_ph(slide, info["heatmap"], left1, top, w, h,
                   label=f"[{split} heatmap not found]")

        # Right: per-series curves generated inline
        df = pd.read_csv(info["csv"])
        fig = _fig_curves_per_series(df, split)
        _add_fig(slide, fig, left2, top, w, h)

        _txt(slide, "Heatmap: BAcc % · rows=npi · cols=N_images",
             left1, Inches(7.2), w, Inches(0.25),
             size_pt=9, color=C_GREY, align=PP_ALIGN.CENTER)
        _txt(slide, "Curves: each line = one series · colored by npi level",
             left2, Inches(7.2), w, Inches(0.25),
             size_pt=9, color=C_GREY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 9 — Max-N: BAcc vs npi, all splits compared
# ---------------------------------------------------------------------------

SPLIT_COLORS  = {"s1v3": "#1f77b4", "s2v2": "#2ca02c", "s3v1": "#d62728"}
SPLIT_LABELS  = {
    "s1v3": "s1v3 (N=47, train f0)",
    "s2v2": "s2v2 (N=48, train f0+1)",
    "s3v1": "s3v1 (N=49, train f0+1+2)",
}
NPI_LEVELS_ALL = [10, 25, 50, 75, 100]


def _fig_max_n_npi_comparison() -> plt.Figure:
    """BAcc vs npi at max N_images for all 3 splits, individual series lines."""
    dfs = {
        "s1v3": pd.read_csv(RESULTS / "le_combined_results.csv"),
        "s2v2": pd.read_csv(RESULTS / "le_combined_s2v2_results.csv"),
        "s3v1": pd.read_csv(RESULTS / "le_combined_s3v1_results.csv"),
    }

    fig, ax = plt.subplots(figsize=(9, 5.5), facecolor="white")

    for split, df in dfs.items():
        max_n = df["n_images"].max()
        sub   = df[df["n_images"] == max_n]
        color = SPLIT_COLORS[split]

        for si, series in enumerate(sorted(sub["series"].unique())):
            row = (sub[sub["series"] == series]
                   .set_index("npi")
                   .reindex(NPI_LEVELS_ALL))
            y = row["balanced_acc"].values * 100
            ls = SERIES_STYLES.get(si, "-")
            mk = SERIES_MARKERS.get(si, "o")
            # Only label the first series per split (for legend)
            label = SPLIT_LABELS[split] if si == 0 else None
            ax.plot(NPI_LEVELS_ALL, y, color=color, linestyle=ls, marker=mk,
                    markersize=5, linewidth=1.6, alpha=0.85, label=label)

    # Split legend (color)
    split_handles = [
        plt.Line2D([0], [0], color=SPLIT_COLORS[s], linewidth=2.5,
                   label=SPLIT_LABELS[s])
        for s in ["s1v3", "s2v2", "s3v1"]
    ]
    # Series legend (line style)
    series_handles = [
        plt.Line2D([0], [0], color="#555", linestyle=SERIES_STYLES[s],
                   marker=SERIES_MARKERS[s], markersize=4, linewidth=1.4,
                   label=f"series {s}")
        for s in [0, 1, 2]
    ]
    leg1 = ax.legend(handles=split_handles, loc="lower right", fontsize=9,
                     title="Split", title_fontsize=9, framealpha=0.92)
    ax.add_artist(leg1)
    ax.legend(handles=series_handles, loc="upper left", fontsize=9,
              title="Series", title_fontsize=9, framealpha=0.92)

    ax.set_xticks(NPI_LEVELS_ALL)
    ax.set_xticklabels([str(n) for n in NPI_LEVELS_ALL], fontsize=10)
    ax.set_xlabel("Labels per frame (npi)", fontsize=11)
    ax.set_ylabel("Balanced accuracy (%)", fontsize=11)
    ax.set_ylim(20, 108)
    ax.axhline(90, color="#CCCCCC", linestyle="--", linewidth=0.8, zorder=0)
    ax.set_facecolor("white")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_title(
        "BAcc vs npi at maximum N_images per split\n"
        "s1v3: N=47 · s2v2: N=48 · s3v1: N=49 · each line = one series",
        fontsize=10, fontweight="bold"
    )
    fig.tight_layout()
    return fig


def _slide_max_n_npi(prs: Presentation):
    slide = _blank(prs)
    _slide_header(
        slide,
        "BAcc vs Labels (npi) at Max N_images — Split Comparison",
        "All 3 splits at their largest N · x-axis = labels per frame · "
        "each line = one series · colored by split"
    )
    fig = _fig_max_n_npi_comparison()
    _add_fig(slide, fig, Inches(1.5), Inches(1.05), Inches(10.3), Inches(6.3))


# ---------------------------------------------------------------------------
# Slide 10 — Cross-Split Comparison
# ---------------------------------------------------------------------------

def _slide_comparison(prs: Presentation):
    slide = _blank(prs)
    _slide_header(slide,
                  "Train/Test Split Comparison · BAcc Heatmaps",
                  "Same colour scale across splits · each cell = mean BAcc over 3 series")

    top  = Inches(1.1)
    h    = Inches(6.1)
    w    = Inches(4.15)
    gap  = Inches(0.1)

    for i, split in enumerate(["s1v3", "s2v2", "s3v1"]):
        left = Inches(0.2) + i * (w + gap)
        _img_or_ph(slide, PNG[f"{split}_heatmap"], left, top, w, h,
                   label=f"[{split} heatmap]")
        _txt(slide, split,
             left, top + h + Inches(0.05), w, Inches(0.3),
             size_pt=11, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 10 — Summary
# ---------------------------------------------------------------------------

def _compute_summary() -> dict:
    """Compute key statistics from results CSVs."""
    dfs = {
        "s1v3": pd.read_csv(RESULTS / "le_combined_results.csv"),
        "s2v2": pd.read_csv(RESULTS / "le_combined_s2v2_results.csv"),
        "s3v1": pd.read_csv(RESULTS / "le_combined_s3v1_results.csv"),
    }
    out = {}
    for name, df in dfs.items():
        best_idx = df["balanced_acc"].idxmax()
        best     = df.loc[best_idx]
        mean_by_n = df.groupby("n_images")["balanced_acc"].mean()
        global_max = mean_by_n.max()
        sat_thresh = 0.9 * global_max
        sat_n = mean_by_n[mean_by_n >= sat_thresh].index.min()
        mean_by_npi = df.groupby("npi")["balanced_acc"].mean()
        best_npi    = mean_by_npi.idxmax()
        worst_npi   = mean_by_npi.idxmin()
        out[name] = {
            "best_bacc":  best["balanced_acc"],
            "best_n":     int(best["n_images"]),
            "best_npi":   int(best["npi"]),
            "sat_n":      int(sat_n),
            "best_npi_mean":  mean_by_npi[best_npi],
            "worst_npi_mean": mean_by_npi[worst_npi],
            "best_npi_key":   int(best_npi),
            "worst_npi_key":  int(worst_npi),
        }
    return out


def _slide_summary(prs: Presentation):
    slide = _blank(prs)
    _slide_header(slide, "Summary — Key Findings",
                  "Label × Image Efficiency · SupCon AE + LGBM · vinc control · 3 splits")

    stats = _compute_summary()

    # Build a matplotlib figure with bullet text
    fig, ax = plt.subplots(figsize=(12.5, 5.5), facecolor="white")
    ax.axis("off")

    lines = []

    # Best BAcc per split
    lines.append(("BEST BALANCED ACCURACY PER SPLIT", 14, "#16213E", True))
    for split in ["s1v3", "s2v2", "s3v1"]:
        s = stats[split]
        lines.append((
            f"  {split}:  {s['best_bacc']:.1%}   "
            f"(N_images={s['best_n']}, npi={s['best_npi']})",
            12, "#1A1A1A", False
        ))
    lines.append(("", 6, "white", False))

    # Saturation
    lines.append(("SATURATION WITH N_IMAGES", 14, "#16213E", True))
    lines.append((
        "  BAcc mean across npi levels reaches ~90% of its plateau quickly:",
        12, "#1A1A1A", False
    ))
    for split in ["s1v3", "s2v2", "s3v1"]:
        s = stats[split]
        lines.append((
            f"  {split}:  plateau reached around N ≈ {s['sat_n']} images",
            12, "#1A1A1A", False
        ))
    lines.append(("", 6, "white", False))

    # npi effect
    lines.append(("EFFECT OF LABELED PATCHES PER FRAME (npi)", 14, "#16213E", True))
    lines.append((
        "  Higher npi generally yields higher BAcc, but gains diminish beyond npi=50.",
        12, "#1A1A1A", False
    ))
    for split in ["s1v3", "s2v2", "s3v1"]:
        s = stats[split]
        lines.append((
            f"  {split}:  best npi={s['best_npi_key']} (mean {s['best_npi_mean']:.1%})  "
            f"vs worst npi={s['worst_npi_key']} (mean {s['worst_npi_mean']:.1%})",
            11, "#444444", False
        ))
    lines.append(("", 6, "white", False))

    # General takeaways
    lines.append(("GENERAL TAKEAWAYS", 14, "#16213E", True))
    takeaways = [
        "  AE trained on even a few unlabeled frames provides useful latent features for LGBM.",
        "  Adding more unlabeled AE frames (higher N) helps more than adding labeled patches (npi).",
        "  s3v1 (3 training frames) achieves the highest peak BAcc but tests on only 1 frame.",
        "  s2v2 and s1v3 are more conservative test scenarios; results are noisier.",
        "  With npi ≥ 50 and N ≥ 10, all splits reach stable, high BAcc (≥ 85%).",
    ]
    for t in takeaways:
        lines.append((t, 11, "#1A1A2E", False))

    y = 0.97
    for text, fs, col, bold in lines:
        if text == "":
            y -= (fs / 600.0)
            continue
        ax.text(0.01, y, text,
                ha="left", va="top",
                fontsize=fs, color=col,
                fontweight="bold" if bold else "normal",
                transform=ax.transAxes)
        y -= max(0.055, fs / 180.0)

    fig.tight_layout(pad=0.3)
    _add_fig(slide, fig, Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.1))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_pptx(out_path: Path):
    print("Building label-efficiency PPT …")
    prs = _prs()

    print("  Slide 1 — Title")
    _slide_title(prs)

    print("  Slide 2 — Experiment Design")
    _slide_design(prs)

    print("  Slides 3–5 — Label Sampling Visualization (N=3, 6, 11)")
    _slides_label_vis(prs)

    print("  Slides 6–8 — Results per split")
    _slides_results(prs)

    print("  Slide 9 — Max-N npi comparison (all splits)")
    _slide_max_n_npi(prs)

    print("  Slide 10 — Cross-Split Comparison")
    _slide_comparison(prs)

    print("  Slide 11 — Summary")
    _slide_summary(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"\nSaved → {out_path}")


def main():
    parser = argparse.ArgumentParser(description="Generate label-efficiency PPT")
    parser.add_argument("--out", default=str(RESULTS / "le_combined_label_efficiency.pptx"),
                        help="Output .pptx path")
    args = parser.parse_args()
    build_pptx(Path(args.out))


if __name__ == "__main__":
    main()
