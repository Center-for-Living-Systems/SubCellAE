#!/usr/bin/env python3
"""
make_pptx_annabel_vinc_sweep.py
================================
Build a PowerPoint summarising the Annabel vinc control label sweep.

Three models  ×  three train/val splits  ×  two classifiers (z_recon / z_proj)

Run before results are ready → placeholders.
Re-run after jobs finish     → results fill in automatically.

Usage:
  python scripts/make_pptx_annabel_vinc_sweep.py
  python scripts/make_pptx_annabel_vinc_sweep.py --out my_results.pptx
"""
from __future__ import annotations

import argparse
import io
from pathlib import Path

import numpy as np
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
OUT  = Path("annabel_vinc_sweep.pptx")

# ── sweep definition ───────────────────────────────────────────────────────────

MODEL_ORDER = ["conae", "supcon2", "supcon5"]
SPLIT_ORDER = ["s1v3", "s2v2", "s3v1"]
SPLIT_LABEL = {"s1v3": "1 train / 3 val", "s2v2": "2 / 2", "s3v1": "3 train / 1 val"}
MODEL_LABEL = {
    "conae":   "ConAE (unsupervised)",
    "supcon2": "SupCon 2-class",
    "supcon5": "SupCon 5-class",
}
MODEL_DETAIL = {
    "conae":   "NT-Xent contrastive loss  |  5-class labels in latents only",
    "supcon2": "Supervised contrastive  |  No adhesion vs Adhesion",
    "supcon5": "Supervised contrastive  |  5 FA subtypes",
}

def result_dir(model: str, split: str) -> Path:
    return RUNS / f"annabel_vinc_{model}_{split}"

# ── slide dimensions ───────────────────────────────────────────────────────────

SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.12)

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_ALT   = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_GREEN = RGBColor(0x37, 0x8B, 0x4A)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)

# ── low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _img(slide, src, l, t, max_w, max_h):
    if src is None:
        return False
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return False
        data = p.read_bytes()
    else:
        data = bytes(src)
    try:
        # Handle TIFF → PNG conversion
        if isinstance(src, (str, Path)) and str(src).lower().endswith(".tif"):
            arr = tifffile.imread(str(src)).astype(np.float32)
            if arr.ndim == 3 and arr.shape[0] > 4:   # (K, H, W) stacked
                # show first cluster only for single-panel display
                arr = arr[0]
            if arr.ndim == 3 and arr.shape[2] in (3, 4):
                pass
            elif arr.ndim == 3:
                arr = arr[0]
            vmin, vmax = arr.min(), arr.max()
            if vmax > vmin:
                arr = ((arr - vmin) / (vmax - vmin) * 255).clip(0, 255).astype(np.uint8)
            else:
                arr = np.zeros_like(arr, dtype=np.uint8)
            buf = io.BytesIO()
            Image.fromarray(arr).save(buf, format="PNG")
            data = buf.getvalue()
        im = Image.open(io.BytesIO(data))
        iw, ih = im.size
    except Exception:
        return False
    scale = min(max_w / iw, max_h / ih)
    w, h  = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(io.BytesIO(data),
                             l + (max_w - w) // 2,
                             t + (max_h - h) // 2, w, h)
    return True

def _placeholder(slide, l, t, w, h, text="[pending]"):
    _rect(slide, l, t, w, h, fill=C_ALT)
    _txt(slide, l, t + h/2 - Inches(0.2), w, Inches(0.4),
         text, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def _img_or_ph(slide, src, l, t, w, h, ph_text="[pending]"):
    if not _img(slide, src, l, t, w, h):
        _placeholder(slide, l, t, w, h, ph_text)

def _title_bar(slide, text, size=15, bg=C_DARK):
    _rect(slide, 0, 0, SW, TITLE_H, fill=bg)
    _txt(slide, PAD, Inches(0.07), SW - 2*PAD, TITLE_H - Inches(0.07),
         text, size=size, bold=True, color=C_WHITE)

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_DARK)
    _txt(slide, Inches(1), Inches(2.0), SW - Inches(2), Inches(2.0),
         title, size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(1.8),
             subtitle, size=15, color=C_LIGHT, align=PP_ALIGN.CENTER)

def slide_section(prs, title, detail=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_MID)
    _txt(slide, Inches(1), Inches(2.5), SW - Inches(2), Inches(1.5),
         title, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if detail:
        _txt(slide, Inches(1.5), Inches(4.2), SW - Inches(3), Inches(2),
             detail, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

def slide_design(prs):
    slide = _blank(prs)
    _title_bar(slide, "Experiment Design", bg=C_DARK)
    top = TITLE_H + PAD

    col_w = (SW - 2*PAD) / 3
    headers = ["Model", "Label scheme", "Training signal"]
    rows = [
        ["ConAE", "5-class (latents only)", "NT-Xent contrastive (unsupervised)"],
        ["SupCon 2-class", "No adhesion vs Adhesion", "Supervised contrastive"],
        ["SupCon 5-class", "No adhes / Nascent / FC / FA / Fibrillar",
         "Supervised contrastive"],
    ]

    # header row
    for ci, h in enumerate(headers):
        l = PAD + ci * col_w
        _rect(slide, l, top, col_w - PAD/2, Inches(0.38), fill=C_MID)
        _txt(slide, l + PAD/2, top + Inches(0.06), col_w - PAD,
             Inches(0.3), h, size=11, bold=True, color=C_WHITE)
    top += Inches(0.42)

    for ri, row in enumerate(rows):
        bg = C_ALT if ri % 2 == 0 else C_WHITE
        for ci, cell in enumerate(row):
            l = PAD + ci * col_w
            _rect(slide, l, top, col_w - PAD/2, Inches(0.42), fill=bg)
            _txt(slide, l + PAD/2, top + Inches(0.07), col_w - PAD,
                 Inches(0.35), cell, size=10)
        top += Inches(0.44)

    top += PAD
    _txt(slide, PAD, top, SW - 2*PAD, Inches(0.32),
         "Train / val splits (per image, 4 frames total):", size=11, bold=True)
    top += Inches(0.34)

    split_rows = [
        ["s1v3 (1/3)", "val_split = 0.75", "1 frame train,  3 frames val"],
        ["s2v2 (2/2)", "val_split = 0.50", "2 frames train, 2 frames val"],
        ["s3v1 (3/1)", "val_split = 0.25", "3 frames train, 1 frame val"],
    ]
    for ri, row in enumerate(split_rows):
        bg = C_ALT if ri % 2 == 0 else C_WHITE
        for ci, cell in enumerate(row):
            l = PAD + ci * col_w
            _rect(slide, l, top, col_w - PAD/2, Inches(0.38), fill=bg)
            _txt(slide, l + PAD/2, top + Inches(0.06), col_w - PAD,
                 Inches(0.3), cell, size=10)
        top += Inches(0.40)

    top += PAD * 2
    arch = ("Architecture: 32×32 px patch  |  58×58 px context (EnlargedCrop)  |  "
            "latent dim=12  |  proj dim=8  |  cio_mode_prt normalisation  |  "
            "output_sigmoid=false  |  nl1 loss  |  λ_recon=1.0  λ_contrast=0.5")
    _txt(slide, PAD, top, SW - 2*PAD, Inches(0.6), arch, size=9, color=C_GREY)

def slide_label_stats(prs):
    slide = _blank(prs)
    _title_bar(slide, "Label Distribution — Annabel vinc control (Annabel, 2026-07-15)")
    top = TITLE_H + PAD

    # 5-class breakdown
    _txt(slide, PAD, top, SW/2, Inches(0.3),
         "5-class labels (539 patches, 4 frames):", size=12, bold=True)
    top5 = top + Inches(0.34)
    rows5 = [
        ("No adhesion",     342, "63%"),
        ("Focal adhesion",  168, "31%"),
        ("Nascent Adhesion", 18,  "3%"),
        ("Fibrillar adhesion", 6, "1%"),
        ("Focal complex",     5,  "1%"),
    ]
    for ri, (label, n, pct) in enumerate(rows5):
        bg = C_ALT if ri % 2 == 0 else C_WHITE
        lx = PAD
        _rect(slide, lx, top5, Inches(2.4), Inches(0.36), fill=bg)
        _txt(slide, lx + PAD/2, top5 + Inches(0.07), Inches(2.3),
             Inches(0.28), label, size=10)
        _rect(slide, lx + Inches(2.45), top5, Inches(0.6), Inches(0.36), fill=bg)
        _txt(slide, lx + Inches(2.5), top5 + Inches(0.07), Inches(0.55),
             Inches(0.28), str(n), size=10, align=PP_ALIGN.RIGHT)
        _rect(slide, lx + Inches(3.1), top5, Inches(0.6), Inches(0.36), fill=bg)
        _txt(slide, lx + Inches(3.15), top5 + Inches(0.07), Inches(0.55),
             Inches(0.28), pct, size=10, color=C_GREY)
        top5 += Inches(0.38)

    # 2-class breakdown
    rx = SW / 2 + PAD
    _txt(slide, rx, top, SW/2 - 2*PAD, Inches(0.3),
         "2-class labels (merged for SupCon 2-class):", size=12, bold=True)
    top2 = top + Inches(0.34)
    rows2 = [
        ("No adhesion", 342, "63%"),
        ("Adhesion (all subtypes merged)", 197, "37%"),
    ]
    for ri, (label, n, pct) in enumerate(rows2):
        bg = C_ALT if ri % 2 == 0 else C_WHITE
        _rect(slide, rx, top2, Inches(3.2), Inches(0.36), fill=bg)
        _txt(slide, rx + PAD/2, top2 + Inches(0.07), Inches(3.1),
             Inches(0.28), label, size=10)
        _rect(slide, rx + Inches(3.25), top2, Inches(0.6), Inches(0.36), fill=bg)
        _txt(slide, rx + Inches(3.3), top2 + Inches(0.07), Inches(0.55),
             Inches(0.28), str(n), size=10, align=PP_ALIGN.RIGHT)
        _rect(slide, rx + Inches(3.9), top2, Inches(0.5), Inches(0.36), fill=bg)
        _txt(slide, rx + Inches(3.95), top2 + Inches(0.07), Inches(0.45),
             Inches(0.28), pct, size=10, color=C_GREY)
        top2 += Inches(0.38)

    # frame breakdown
    bot = max(top5, top2) + PAD * 2
    _txt(slide, PAD, bot, SW - 2*PAD, Inches(0.3),
         "Frames:  f0000, f0001, f0002, f0003  (4 total — split is per-image)", size=10, color=C_GREY)
    bot += Inches(0.34)
    frame_rows = [
        ("f0000", "~143 patches"),
        ("f0001", "~132 patches"),
        ("f0002", "~138 patches"),
        ("f0003", "~126 patches"),
    ]
    for fi, (fname, desc) in enumerate(frame_rows):
        _txt(slide, PAD + fi * Inches(2.5), bot, Inches(2.4), Inches(0.28),
             f"{fname}  —  {desc}", size=9, color=C_GREY)

    bot += Inches(0.32) + PAD
    _txt(slide, PAD, bot, SW - 2*PAD, Inches(0.3),
         "Classification: LightGBM on z_recon (12-dim) and z_proj (8-dim) independently. "
         "Train/val split matches the AE training split (from_csv).", size=9, color=C_GREY)

def slide_3img_split(prs, title, model, img_fn, cap_fn=None, ph="[pending]"):
    """3-panel slide: one column per split (s1v3 | s2v2 | s3v1)."""
    slide = _blank(prs)
    _title_bar(slide, title)
    top   = TITLE_H + PAD
    cap_h = Inches(0.28)
    h     = SH - top - cap_h - PAD
    w     = (SW - 4*PAD) / 3

    for i, split in enumerate(SPLIT_ORDER):
        l   = PAD + i * (w + PAD)
        rd  = result_dir(model, split)
        src = img_fn(rd, split)
        _img_or_ph(slide, src, l, top, w, h, ph_text=ph)
        cap = cap_fn(split) if cap_fn else SPLIT_LABEL[split]
        _txt(slide, l, SH - cap_h, w, cap_h,
             cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def slide_3img_model(prs, title, split, img_fn, cap_fn=None, ph="[pending]"):
    """3-panel slide: one column per model (conae | supcon2 | supcon5)."""
    slide = _blank(prs)
    _title_bar(slide, title, bg=C_GREEN)
    top   = TITLE_H + PAD
    cap_h = Inches(0.28)
    h     = SH - top - cap_h - PAD
    w     = (SW - 4*PAD) / 3

    for i, model in enumerate(MODEL_ORDER):
        l   = PAD + i * (w + PAD)
        rd  = result_dir(model, split)
        src = img_fn(rd, model)
        _img_or_ph(slide, src, l, top, w, h, ph_text=ph)
        cap = cap_fn(model) if cap_fn else MODEL_LABEL[model]
        _txt(slide, l, SH - cap_h, w, cap_h,
             cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def slide_6img(prs, title, items, ph="[pending]"):
    """2-row × 3-col layout. items = list of (img_path, caption)."""
    slide = _blank(prs)
    _title_bar(slide, title)
    top   = TITLE_H + PAD
    cap_h = Inches(0.24)
    h     = (SH - top - 3*PAD - 2*cap_h) / 2
    w     = (SW - 4*PAD) / 3

    for idx, (src, cap) in enumerate(items[:6]):
        row, col = divmod(idx, 3)
        l = PAD + col * (w + PAD)
        t = top + row * (h + cap_h + PAD)
        _img_or_ph(slide, src, l, t, w, h, ph_text=ph)
        _txt(slide, l, t + h, w, cap_h,
             cap, size=8, color=C_GREY, align=PP_ALIGN.CENTER)

# ── image path helpers ─────────────────────────────────────────────────────────

def _umap_annotation(rd, _):
    return rd / "analysis" / "UMAP" / "by_annotation.png"

def _umap_split(rd, _):
    return rd / "analysis" / "UMAP" / "by_split.png"

def _umap_kmeans(rd, _, k):
    return rd / f"eval/cluster_panels_k{k}" / f"umap_kmeans_k{k}.png"

def _confusion(rd, feat, split_type="val"):
    return rd / f"fa_cls_{feat}" / f"confusion_matrix_norm_{split_type}.png"

def _violin(rd, metric):
    return rd / "violin_plots" / f"vinc_{metric}.png"

def _cluster_all(rd, k):
    # all_clusters.tif — stacked TIFF (one frame per cluster)
    return rd / f"eval/cluster_panels_k{k}" / "all_clusters.tif"

# ── per-model section ──────────────────────────────────────────────────────────

def model_section(prs, model):
    label  = MODEL_LABEL[model]
    detail = MODEL_DETAIL[model]
    slide_section(prs, label, detail)

    # UMAP by annotation (3 splits side-by-side)
    slide_3img_split(prs,
        f"{label} — UMAP coloured by FA type  (3 splits)",
        model,
        img_fn=_umap_annotation,
        cap_fn=lambda sp: SPLIT_LABEL[sp])

    # UMAP by split (train/val)
    slide_3img_split(prs,
        f"{label} — UMAP coloured by train/val split",
        model,
        img_fn=_umap_split,
        cap_fn=lambda sp: SPLIT_LABEL[sp])

    # Confusion z_recon (3 splits)
    slide_3img_split(prs,
        f"{label} — Confusion matrix (z_recon, normalised)  [val set]",
        model,
        img_fn=lambda rd, sp: _confusion(rd, "zrecon"),
        cap_fn=lambda sp: SPLIT_LABEL[sp])

    # Confusion z_proj (3 splits)
    slide_3img_split(prs,
        f"{label} — Confusion matrix (z_proj, normalised)  [val set]",
        model,
        img_fn=lambda rd, sp: _confusion(rd, "zproj"),
        cap_fn=lambda sp: SPLIT_LABEL[sp])

    # Violin: L1, MSE, Hessian for best split (s2v2)
    slide = _blank(prs)
    _title_bar(slide, f"{label} — Reconstruction quality violins  (s2v2: 2/2 split)")
    top = TITLE_H + PAD
    h   = SH - top - PAD
    w   = (SW - 4*PAD) / 3
    rd  = result_dir(model, "s2v2")
    for i, metric in enumerate(["recon_l1", "recon_mse", "recon_hessian_l1"]):
        l = PAD + i * (w + PAD)
        _img_or_ph(slide, _violin(rd, metric), l, top, w, h)
    _txt(slide, PAD, SH - Inches(0.24), w, Inches(0.22),
         "Recon L1", size=8, color=C_GREY, align=PP_ALIGN.CENTER)
    _txt(slide, PAD + w + PAD, SH - Inches(0.24), w, Inches(0.22),
         "Recon MSE", size=8, color=C_GREY, align=PP_ALIGN.CENTER)
    _txt(slide, PAD + 2*(w + PAD), SH - Inches(0.24), w, Inches(0.22),
         "Hessian L1", size=8, color=C_GREY, align=PP_ALIGN.CENTER)

    # KMeans cluster scatter (UMAP coloured by cluster) for each k
    for k in [3, 6, 10]:
        slide_3img_split(prs,
            f"{label} — KMeans k={k} scatter (UMAP coloured by cluster)",
            model,
            img_fn=lambda rd, sp, _k=k: _umap_kmeans(rd, sp, _k),
            cap_fn=lambda sp: SPLIT_LABEL[sp])

    # KMeans cluster panels (all_clusters.tif) for s2v2, k=3/6/10
    rd = result_dir(model, "s2v2")
    items = []
    for k in [3, 6, 10]:
        src = _cluster_all(rd, k)
        items.append((src, f"k={k}  all clusters"))
    # pad to 6
    while len(items) < 6:
        items.append((None, ""))
    slide_6img(prs,
        f"{label} — KMeans cluster panels (s2v2: 2/2 split)",
        items[:3] + [(None, "")] * 3)   # top row only, 3 panels

# ── model comparison section ───────────────────────────────────────────────────

def comparison_section(prs, split="s2v2"):
    label = f"Model Comparison — {SPLIT_LABEL[split]} split"
    slide_section(prs, label,
                  "ConAE  vs  SupCon 2-class  vs  SupCon 5-class")

    # Confusion z_recon: 3 models
    slide_3img_model(prs,
        f"Confusion matrix — z_recon (normalised)  [{SPLIT_LABEL[split]}]",
        split,
        img_fn=lambda rd, m: _confusion(rd, "zrecon"),
        cap_fn=lambda m: MODEL_LABEL[m])

    # Confusion z_proj: 3 models
    slide_3img_model(prs,
        f"Confusion matrix — z_proj (normalised)  [{SPLIT_LABEL[split]}]",
        split,
        img_fn=lambda rd, m: _confusion(rd, "zproj"),
        cap_fn=lambda m: MODEL_LABEL[m])

    # UMAP by annotation: 3 models
    slide_3img_model(prs,
        f"UMAP coloured by FA type  [{SPLIT_LABEL[split]}]",
        split,
        img_fn=_umap_annotation,
        cap_fn=lambda m: MODEL_LABEL[m])

    # UMAP by split: 3 models
    slide_3img_model(prs,
        f"UMAP coloured by train/val  [{SPLIT_LABEL[split]}]",
        split,
        img_fn=_umap_split,
        cap_fn=lambda m: MODEL_LABEL[m])

# ── main ───────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", type=Path, default=OUT)
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    # ── cover ──────────────────────────────────────────────────────────────────
    slide_title(prs,
        "Supervised Contrastive AE\nClassification — Annabel vinc Control",
        "ConAE  ·  SupCon 2-class  ·  SupCon 5-class\n"
        "3 train/val splits  ·  LightGBM on z_recon & z_proj\n"
        "cio_mode_prt normalisation  ·  latent=12  proj=8\n"
        "2026-08-02")

    # ── design & labels ────────────────────────────────────────────────────────
    slide_design(prs)
    slide_label_stats(prs)

    # ── per-model sections ─────────────────────────────────────────────────────
    for model in MODEL_ORDER:
        model_section(prs, model)

    # ── cross-model comparison ─────────────────────────────────────────────────
    comparison_section(prs, split="s2v2")

    # ── save ───────────────────────────────────────────────────────────────────
    prs.save(str(args.out))
    n = len(prs.slides)
    print(f"Saved {n} slides → {args.out}")


if __name__ == "__main__":
    main()
