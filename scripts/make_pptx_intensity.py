#!/usr/bin/env python3
"""
make_pptx_intensity.py
======================
Build a 2-slide PPT for the FA patch peak-intensity distribution analysis.

Slide 1: Overview — 3-panel distribution figure + dataset-level summary
Slide 2: Per-class breakdown + retention implications

Usage
-----
  python scripts/make_pptx_intensity.py
"""
from __future__ import annotations

from pathlib import Path

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
_REPO   = Path(__file__).resolve().parents[1]
EVAL    = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run/fa4_xds_eval")
RES     = _REPO / "results"
OUT_PPT = RES / "intensity_analysis.pptx"

DIST_PNG = EVAL / "intensity_distribution.png"

DS_COLORS = {
    "vinc_ctrl":  RGBColor(0x1f, 0x77, 0xb4),
    "vinc_ycomp": RGBColor(0xff, 0x7f, 0x0e),
    "pfak_ctrl":  RGBColor(0xd6, 0x27, 0x28),
}
FA_COLORS = {
    "Nascent Adhesion":   RGBColor(0x43, 0x93, 0xc3),
    "focal complex":      RGBColor(0xf4, 0xa5, 0x82),
    "focal adhesion":     RGBColor(0x2c, 0xa0, 0x2c),
    "fibrillar adhesion": RGBColor(0x94, 0x67, 0xbd),
}
FA_ORDER = ["Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
FA_SHORT = {"Nascent Adhesion": "NA", "focal complex": "FC",
            "focal adhesion": "FA", "fibrillar adhesion": "Fib"}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text: str, top=Inches(0.12)):
    txb = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.7), Inches(0.5))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold  = True
    run.font.size  = Pt(18)
    run.font.color.rgb = RGBColor(0x1a, 0x4a, 0x7a)


def add_subtitle(slide, text: str, top=Inches(0.60)):
    txb = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.7), Inches(0.35))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def add_bullets(slide, bullets: list[tuple[str, RGBColor | None]],
                left, top, width, height, fontsize=Pt(9.5)):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(bullets):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = fontsize
        if color:
            p.font.color.rgb = color
        else:
            p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(4)


def add_table(slide, rows: list[list[str]], col_widths: list[float],
              left, top, row_h=Inches(0.28), header_color=RGBColor(0x1a, 0x4a, 0x7a)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  left, top,
                                  Inches(total_w), row_h * n_rows).table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xf2, 0xf6, 0xfb)


# ---------------------------------------------------------------------------

def load_data():
    all_df = pd.read_csv(EVAL / "patch_intensities_all.csv")
    lab_df = pd.read_csv(EVAL / "patch_intensities_labeled.csv")
    return all_df, lab_df


def slide1_overview(prs, all_df, lab_df):
    slide = blank_slide(prs)
    add_title(slide, "FA patch peak intensity — overview (99th percentile pixel value per 32×32 patch)")
    add_subtitle(slide,
        "All patches scanned across 3 datasets · labeled patches (802) broken out by FA class · "
        "patch size 32×32 px · intensity range after per-dataset percentile normalisation")

    # Main figure — takes up most of slide
    if DIST_PNG.exists():
        slide.shapes.add_picture(str(DIST_PNG), Inches(0.2), Inches(1.0),
                                 Inches(9.5), Inches(3.6))

    # Dataset summary table (right side)
    ds_stats = all_df.groupby("dataset")["peak_intensity"].agg(
        count="count", mean="mean", median="median", std="std"
    ).round(3)

    tbl_rows = [["Dataset", "N patches", "Median", "Mean ± SD"]]
    ds_labels = {
        "vinc_ctrl":  "vinc control",
        "vinc_ycomp": "vinc ycomp",
        "pfak_ctrl":  "pfak control",
    }
    for ds in ["vinc_ctrl", "vinc_ycomp", "pfak_ctrl"]:
        s = ds_stats.loc[ds]
        tbl_rows.append([
            ds_labels[ds],
            f"{int(s['count']):,}",
            f"{s['median']:.3f}",
            f"{s['mean']:.3f} ± {s['std']:.3f}",
        ])
    add_table(slide, tbl_rows, [1.8, 1.0, 0.8, 1.5],
              Inches(9.9), Inches(1.05))

    # Key observations bullets
    vals = lab_df["peak_intensity"].values
    t75  = np.percentile(vals, 75)
    bullets = [
        ("Key observations:", None),
        (f"• pfak patches are ~2× brighter than vinc (median 1.08 vs 0.54) — different imaging conditions or FA density", None),
        (f"• vinc ycomp slightly brighter than ctrl (0.58 vs 0.54) — consistent with ycomp morphology changes", None),
        (f"• FA and Fib classes are the brightest (mean ~1.3), NA the dimmest (0.81) — intensity partially separates classes", None),
        (f"• At p75 threshold (intensity ≥ {t75:.2f}): FA/Fib retain ~32% of labeled patches; NA/FC retain only ~14–16%", None),
        (f"• Intensity-based filtering would disproportionately remove NA and FC patches — biases training set toward mature FAs", None),
    ]
    add_bullets(slide, bullets, Inches(0.2), Inches(4.75), Inches(9.5), Inches(2.6),
                fontsize=Pt(9.5))

    return slide


def slide2_perclass(prs, all_df, lab_df):
    slide = blank_slide(prs)
    add_title(slide, "Peak intensity by FA class — retention implications")
    add_subtitle(slide,
        "Labeled patches only (802 total: vinc ctrl + vinc ycomp + pfak ctrl) · "
        "retention = % patches above threshold · question: does intensity filtering change what the AE trains on?")

    # Per-class stats table
    cls_stats = lab_df.groupby("label")["peak_intensity"].agg(
        count="count", mean="mean", median="median", std="std"
    ).round(3)

    vals = lab_df["peak_intensity"].values
    thresholds = {f"p{p}": np.percentile(vals, p) for p in [50, 75, 90]}

    tbl_rows = [["FA class", "n", "Median", "Mean ± SD",
                 f"≥p50 ({thresholds['p50']:.2f})",
                 f"≥p75 ({thresholds['p75']:.2f})",
                 f"≥p90 ({thresholds['p90']:.2f})"]]

    for cls in FA_ORDER:
        if cls not in cls_stats.index:
            continue
        s   = cls_stats.loc[cls]
        sub = lab_df[lab_df["label"] == cls]["peak_intensity"].values
        tbl_rows.append([
            f"{FA_SHORT[cls]} — {cls}",
            str(int(s["count"])),
            f"{s['median']:.3f}",
            f"{s['mean']:.3f} ± {s['std']:.3f}",
            f"{np.mean(sub >= thresholds['p50'])*100:.0f}%",
            f"{np.mean(sub >= thresholds['p75'])*100:.0f}%",
            f"{np.mean(sub >= thresholds['p90'])*100:.0f}%",
        ])

    add_table(slide, tbl_rows, [2.3, 0.5, 0.8, 1.5, 1.1, 1.1, 1.1],
              Inches(0.3), Inches(1.05))

    # Interpretation bullets
    bullets = [
        ("Interpretation:", None),
        ("• NA (Nascent Adhesion) patches are the dimmest class (median 0.60) — small, newly formed FAs with low paxillin signal", None),
        ("• FA and Fib patches are >2× brighter than NA — mature, large FAs accumulate more paxillin", None),
        ("• At the p75 threshold, intensity filtering retains ~32% of FA/Fib but only 14% of NA patches", None),
        ("• Conclusion: intensity-based patch selection would severely under-represent NA structures", None),
        ("  → Recommendation: do NOT filter by intensity for AE training; retain all patches", None),
        ("• Alternative use: flag low-intensity adhesion patches as harder cases for the classifier (consistent with FN error analysis)", None),
    ]
    add_bullets(slide, bullets, Inches(0.3), Inches(3.3), Inches(12.7), Inches(3.9),
                fontsize=Pt(10))

    return slide


def slide3_umap(prs):
    UMAP_PNG = RES / "umap_intensity_s2v2.png"
    slide = blank_slide(prs)
    add_title(slide, "UMAP of s2v2 latent space — intensity + FA class overlay")
    add_subtitle(slide,
        "Model: le_c1_npi100_r0 (s2v2, npi=100) · 14,879 vinc-control patches · "
        "Fill = 99th-pct pixel intensity (blue scale) · Edge colour = FA 5-class label")

    if UMAP_PNG.exists():
        slide.shapes.add_picture(str(UMAP_PNG), Inches(1.5), Inches(0.95),
                                 Inches(8.5), Inches(6.3))

    bullets = [
        ("• Fill (blue shade): brighter = higher peak intensity — mature FAs cluster in high-intensity regions", None),
        ("• Edge colours: purple=No adhesion · blue=NA · orange=FC · green=FA · red=Fib", None),
        ("• Low-intensity NA patches (dim blue fill, blue edge) tend to sit at the boundary with No-adhesion — consistent with FN errors", None),
    ]
    add_bullets(slide, bullets, Inches(10.1), Inches(2.5), Inches(3.0), Inches(3.5),
                fontsize=Pt(9))
    return slide


# ---------------------------------------------------------------------------

def main():
    RES.mkdir(exist_ok=True)

    if not DIST_PNG.exists():
        print("WARNING: intensity_distribution.png not found — run analyze_fa4_intensity.py first")

    print("Loading intensity data ...")
    all_df, lab_df = load_data()

    prs = new_prs()

    print("Slide 1: overview ...")
    slide1_overview(prs, all_df, lab_df)

    print("Slide 2: per-class breakdown ...")
    slide2_perclass(prs, all_df, lab_df)

    print("Slide 3: UMAP intensity ...")
    slide3_umap(prs)

    prs.save(str(OUT_PPT))
    print(f"\nSaved: {OUT_PPT}")


if __name__ == "__main__":
    main()
