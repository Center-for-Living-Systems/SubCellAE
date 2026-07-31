#!/usr/bin/env python3
"""
make_pptx_three_groups.py
=========================
Generate a comprehensive PPT for three ConAE/SupConAE training groups:
  Group 1: CIO-vinc flat models (lat12 lc-sweep + lat18/lat24 multichannel)
  Group 2: ds_combo_enlcrop_clip01_l1  (15 dataset combos)
  Group 3: ds_combo_enlcrop_sc2_clip02_l1  (15 dataset combos)

Run from the SubCellAE repo root:
  python scripts/make_pptx_three_groups.py
"""
from __future__ import annotations

import io
from pathlib import Path

import numpy as np
import pandas as pd
import tifffile
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
OUT  = Path("three_groups_results.pptx")

# ── slide geometry (16:9) ──────────────────────────────────────────────────────

SW = Inches(13.33)
SH = Inches(7.5)
TITLE_H = Inches(0.65)
PAD     = Inches(0.12)

# ── colours ────────────────────────────────────────────────────────────────────

C_DARK   = RGBColor(0x1F, 0x4E, 0x79)
C_MID    = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT  = RGBColor(0xBD, 0xD7, 0xEE)
C_ALT    = RGBColor(0xF2, 0xF2, 0xF2)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_GREY   = RGBColor(0x88, 0x88, 0x88)

# ── low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh


def _txt(slide, l, t, w, h, text, size=14, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def _img(slide, src, l, t, max_w, max_h):
    """Add image (path or bytes), fitting within max_w x max_h, centred."""
    if src is None:
        return
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return
        data = p.read_bytes()
    else:
        data = bytes(src)
    try:
        im = Image.open(io.BytesIO(data))
        iw, ih = im.size
    except Exception:
        return
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    lo = l + (max_w - w) // 2
    to = t + (max_h - h) // 2
    slide.shapes.add_picture(io.BytesIO(data), lo, to, w, h)


def _title_bar(slide, text, size=18, bg=C_DARK):
    _rect(slide, 0, 0, SW, TITLE_H, fill=bg)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         text, size=size, bold=True, color=C_WHITE)


def _tif_to_png(path: Path) -> bytes | None:
    """Convert a cluster-panel TIF (K, H, W) to a tiled PNG (bytes)."""
    try:
        arr = tifffile.imread(str(path)).astype(np.float32)
    except Exception:
        return None

    if arr.ndim == 2:
        img2d = arr
    elif arr.ndim == 3 and arr.shape[2] in (3, 4):
        img2d = arr
    elif arr.ndim == 3:
        K, H, W = arr.shape
        ncols = min(K, 5)
        nrows = (K + ncols - 1) // ncols
        canvas = np.zeros((nrows * H, ncols * W), dtype=np.float32)
        for i in range(K):
            r, c = divmod(i, ncols)
            canvas[r*H:(r+1)*H, c*W:(c+1)*W] = arr[i]
        img2d = canvas
    else:
        img2d = arr[0]

    vmin, vmax = img2d.min(), img2d.max()
    if vmax > vmin:
        img2d = ((img2d - vmin) / (vmax - vmin) * 255).clip(0, 255).astype(np.uint8)
    else:
        img2d = np.zeros_like(img2d, dtype=np.uint8)

    pil = Image.fromarray(img2d)
    buf = io.BytesIO()
    pil.save(buf, format='PNG')
    return buf.getvalue()


# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_DARK)
    _txt(slide, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.8),
         title, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _txt(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(2),
             subtitle, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_section(prs, title, detail=""):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, fill=C_MID)
    _txt(slide, Inches(1), Inches(2.4), SW - Inches(2), Inches(1.6),
         title, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if detail:
        _txt(slide, Inches(1.5), Inches(4.2), SW - Inches(3), Inches(2),
             detail, size=13, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_1img(prs, title, img, caption=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top = TITLE_H + PAD
    avail_h = SH - top - (Inches(0.32) if caption else PAD)
    _img(slide, img, PAD, top, SW - 2*PAD, avail_h)
    if caption:
        _txt(slide, PAD, SH - Inches(0.32), SW - 2*PAD, Inches(0.28),
             caption, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_2img(prs, title, img1, img2, cap1="", cap2=""):
    slide = _blank(prs)
    _title_bar(slide, title)
    top = TITLE_H + PAD
    cap_h = Inches(0.28) if (cap1 or cap2) else PAD
    h = SH - top - cap_h - PAD
    w = (SW - 3*PAD) / 2
    _img(slide, img1, PAD,             top, w, h)
    _img(slide, img2, PAD + w + PAD,   top, w, h)
    if cap1:
        _txt(slide, PAD, SH - cap_h, w, cap_h,
             cap1, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
    if cap2:
        _txt(slide, PAD + w + PAD, SH - cap_h, w, cap_h,
             cap2, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_3img(prs, title, imgs, caps=None):
    slide = _blank(prs)
    _title_bar(slide, title)
    caps = (list(caps) + [""]*3)[:3] if caps else [""]*3
    imgs = (list(imgs) + [None]*3)[:3]
    top  = TITLE_H + PAD
    cap_h = Inches(0.28) if any(caps) else PAD
    h = SH - top - cap_h - PAD
    w = (SW - 4*PAD) / 3
    for i, (im, cap) in enumerate(zip(imgs, caps)):
        l = PAD + i * (w + PAD)
        if im:
            _img(slide, im, l, top, w, h)
        else:
            _txt(slide, l, top + h/2 - Inches(0.25), w, Inches(0.5),
                 "[not available]", size=10, color=C_GREY, align=PP_ALIGN.CENTER)
        if cap:
            _txt(slide, l, SH - cap_h, w, cap_h,
                 cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def _style_cell(cell, text, size=8, bold=False, align=PP_ALIGN.CENTER,
                bg=None, fg=C_BLACK):
    cell.text = str(text)
    tf = cell.text_frame
    tf.paragraphs[0].alignment = align
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def slide_table(prs, title, headers, rows, font_size=8, title_bg=C_DARK):
    slide = _blank(prs)
    _title_bar(slide, title, bg=title_bg)
    if not rows:
        _txt(slide, PAD, TITLE_H + PAD, SW - 2*PAD, Inches(1),
             "(no data)", size=12, color=C_GREY, align=PP_ALIGN.CENTER)
        return
    nr = len(rows) + 1
    nc = len(headers)
    row_h = min(Inches(0.33), (SH - TITLE_H - 2*PAD) / nr)
    tbl_h = row_h * nr
    tbl_w = SW - 2*PAD
    top   = TITLE_H + PAD + max(0, (SH - TITLE_H - 2*PAD - tbl_h) / 2)
    tbl = slide.shapes.add_table(nr, nc, PAD, top, tbl_w, tbl_h).table
    for j, h in enumerate(headers):
        _style_cell(tbl.cell(0, j), h, size=font_size, bold=True, bg=C_DARK, fg=C_WHITE)
    for i, row in enumerate(rows):
        bg = C_ALT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            _style_cell(tbl.cell(i+1, j), str(val) if val is not None else "",
                        size=font_size, bg=bg)


# ── Group 1 definitions ────────────────────────────────────────────────────────

G1_MODELS = [
    # (label, dirname, prefix, lat, proj, nch, channels_str, lc)
    ("ConAE lat12 lc=0.5",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1",
     "contrastive", 12, 8, 1, "pax", 0.5),
    ("ConAE lat12 lc=100",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc100",
     "contrastive", 12, 8, 1, "pax", 100),
    ("ConAE lat12 lc=1e-4",
     "contrastive_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc1e4",
     "contrastive", 12, 8, 1, "pax", 0.0001),
    ("SupCon lat12 lc=0.5",
     "supcon_cio_vinc_lat12proj8_enlcrop_sc2_nl1",
     "supcon", 12, 8, 1, "pax", 0.5),
    ("SupCon lat12 lc=100",
     "supcon_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc100",
     "supcon", 12, 8, 1, "pax", 100),
    ("SupCon lat12 lc=1e-4",
     "supcon_cio_vinc_lat12proj8_enlcrop_sc2_nl1_lc1e4",
     "supcon", 12, 8, 1, "pax", 0.0001),
    ("ConAE lat18 3ch (pax+zyx+act) lc=0.5",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza",
     "contrastive", 18, 12, 3, "pax,zyx,act", 0.5),
    ("ConAE lat18 3ch (pax+zyx+act) lc=100",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza_lc100",
     "contrastive", 18, 12, 3, "pax,zyx,act", 100),
    ("ConAE lat18 3ch (pax+zyx+act) lc=1e-4",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_3ch_pza_lc1e4",
     "contrastive", 18, 12, 3, "pax,zyx,act", 0.0001),
    ("ConAE lat18 4ch (vinc+pax+zyx+act) lc=0.5",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc",
     "contrastive", 18, 12, 4, "vinc,pax,zyx,act", 0.5),
    ("ConAE lat18 4ch (vinc+pax+zyx+act) lc=100",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc_lc100",
     "contrastive", 18, 12, 4, "vinc,pax,zyx,act", 100),
    ("ConAE lat18 4ch (vinc+pax+zyx+act) lc=1e-4",
     "contrastive_cio_vinc_lat18proj12_enlcrop_sc2_nl1_4ch_vinc_lc1e4",
     "contrastive", 18, 12, 4, "vinc,pax,zyx,act", 0.0001),
    ("ConAE lat24 3ch (pax+zyx+act) lc=0.5",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza",
     "contrastive", 24, 16, 3, "pax,zyx,act", 0.5),
    ("ConAE lat24 3ch (pax+zyx+act) lc=100",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza_lc100",
     "contrastive", 24, 16, 3, "pax,zyx,act", 100),
    ("ConAE lat24 3ch (pax+zyx+act) lc=1e-4",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_3ch_pza_lc1e4",
     "contrastive", 24, 16, 3, "pax,zyx,act", 0.0001),
    ("ConAE lat24 4ch (vinc+pax+zyx+act) lc=0.5",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc",
     "contrastive", 24, 16, 4, "vinc,pax,zyx,act", 0.5),
    ("ConAE lat24 4ch (vinc+pax+zyx+act) lc=100",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc_lc100",
     "contrastive", 24, 16, 4, "vinc,pax,zyx,act", 100),
    ("ConAE lat24 4ch (vinc+pax+zyx+act) lc=1e-4",
     "contrastive_cio_vinc_lat24proj16_enlcrop_sc2_nl1_4ch_vinc_lc1e4",
     "contrastive", 24, 16, 4, "vinc,pax,zyx,act", 0.0001),
]

COMBO_LIST = [
    "vinc", "nih3t3", "ppax", "pfak",
    "vinc_nih3t3", "vinc_ppax", "vinc_pfak",
    "nih3t3_ppax", "nih3t3_pfak", "ppax_pfak",
    "vinc_nih3t3_ppax", "vinc_nih3t3_pfak", "vinc_ppax_pfak",
    "nih3t3_ppax_pfak", "vinc_nih3t3_ppax_pfak",
]

DS_GROUPS = [
    ("Group 2", "ds_combo_enlcrop_clip01_l1",
     "ConAE  clip[0,1] input  L1 recon loss  lambda_contrast=0.10\n"
     "15 dataset combinations, 500 epochs each"),
    ("Group 3", "ds_combo_enlcrop_sc2_clip02_l1",
     "ConAE  clip[0,2] input  L1 recon loss  scale_contrast=2  lambda_contrast=0.10\n"
     "15 dataset combinations, 500 epochs each"),
]


# ── per-model slide block ──────────────────────────────────────────────────────

def _cluster_panels(prs, label: str, ev: Path):
    """Add slides for cluster_panels and cluster_panels_combo."""
    for cp_name, tag in [("cluster_panels", ""), ("cluster_panels_combo", " (training combo)")]:
        cp_dir = ev / cp_name
        if not cp_dir.is_dir():
            continue
        all_tif = cp_dir / "all_clusters.tif"
        if not all_tif.exists():
            all_tif = next(iter(sorted(cp_dir.glob("all*.tif"))), None)
        if all_tif:
            png = _tif_to_png(all_tif)
            if png:
                slide_1img(prs, f"{label} -- K-means Cluster Panels{tag}",
                           png, f"{all_tif.name}")


def _model_slides(prs, label: str, d: Path, prefix: str):
    """Add all result slides for one model directory."""
    if not d.exists():
        slide_section(prs, f"{label} -- NOT FOUND",
                      f"Directory missing: {d.name}")
        return

    ev = d / "eval"

    # 1. Training loss + component losses
    loss = d / f"{prefix}_train_val_loss.png"
    comp = d / f"{prefix}_component_losses.png"
    if loss.exists() and comp.exists():
        slide_2img(prs, f"{label} -- Training Loss",
                   loss, comp, "Train/Val Loss", "Component Losses")
    elif loss.exists():
        slide_1img(prs, f"{label} -- Training Loss", loss)

    # 2. ep500 views + reconstruction
    views = d / f"{prefix}_views_ep500.png"
    recon = d / f"{prefix}_recon_ep500.png"
    if views.exists() and recon.exists():
        slide_2img(prs, f"{label} -- Epoch 500",
                   views, recon, "Contrastive views", "Reconstruction")
    elif views.exists():
        slide_1img(prs, f"{label} -- Epoch 500 Views", views)
    elif recon.exists():
        slide_1img(prs, f"{label} -- Epoch 500 Recon", recon)

    # 3. Cross-dataset violin/box plots (recon metrics)
    cross_nl1  = d / "cross_dataset_recon_nl1.png"
    cross_l1   = d / "cross_dataset_recon_l1.png"
    cross_mse  = d / "cross_dataset_recon_mse.png"
    cross_hess = d / "cross_dataset_recon_hessian_l1.png"
    has_cross = cross_nl1.exists() or cross_l1.exists() or cross_mse.exists()
    if has_cross:
        slide_3img(prs, f"{label} -- Cross-Dataset Recon Metrics (violin plots)",
                   [cross_nl1 if cross_nl1.exists() else None,
                    cross_l1  if cross_l1.exists()  else None,
                    cross_mse if cross_mse.exists() else None],
                   ["nL1 (normalised)", "L1 (absolute)", "MSE"])
        if cross_hess.exists():
            slide_1img(prs, f"{label} -- Cross-Dataset Hessian L1", cross_hess)

    # Per-channel recon metrics (multichannel models)
    for ch in range(4):
        nl1c = d / f"cross_dataset_recon_nl1_ch{ch}.png"
        l1c  = d / f"cross_dataset_recon_l1_ch{ch}.png"
        if nl1c.exists() and l1c.exists():
            slide_2img(prs, f"{label} -- Cross-Dataset Recon ch{ch}",
                       nl1c, l1c, f"nL1 ch{ch}", f"L1 ch{ch}")

    # 4. Metrics table from cross_dataset_recon_metrics.csv
    mcsv = d / "cross_dataset_recon_metrics.csv"
    if mcsv.exists():
        try:
            df = pd.read_csv(mcsv)
            grp_cols = [c for c in ("dataset", "split") if c in df.columns]
            mc = [c for c in ("recon_nl1", "recon_l1", "recon_mse") if c in df.columns]
            if mc and grp_cols:
                agg = (df.groupby(grp_cols)[mc]
                         .mean().reset_index().round(4))
                hdrs = list(agg.columns)
                rows = [[str(v) for v in r] for _, r in agg.iterrows()]
                slide_table(prs,
                            f"{label} -- Reconstruction Metrics Table",
                            hdrs, rows, font_size=9)
        except Exception as e:
            print(f"  WARNING metrics table: {e}")

    # 5. UMAP z_recon 4ds (lat18/lat24 + ds_combo models)
    u4_ann  = ev / "umap_4ds_annotation.png"
    u4_cond = ev / "umap_4ds_condition.png"
    u4_ds   = ev / "umap_4ds_dataset.png"
    if u4_ann.exists():
        slide_3img(prs, f"{label} -- UMAP z_recon (cross-dataset)",
                   [u4_ann,
                    u4_cond if u4_cond.exists() else None,
                    u4_ds   if u4_ds.exists()   else None],
                   ["Annotation (labels)", "Condition", "Dataset"])

    # 6. UMAP z_proj 4ds (all models)
    p4_ann  = ev / "umap_proj_4ds_annotation.png"
    p4_cond = ev / "umap_proj_4ds_condition.png"
    p4_ds   = ev / "umap_proj_4ds_dataset.png"
    p4_km   = ev / "umap_proj_4ds_kmeans.png"
    if p4_ann.exists():
        slide_3img(prs, f"{label} -- UMAP z_proj (cross-dataset)",
                   [p4_ann,
                    p4_cond if p4_cond.exists() else None,
                    p4_ds   if p4_ds.exists()   else None],
                   ["Annotation (labels)", "Condition (training labels)", "Dataset"])
    # K-means cluster ID + training condition side by side
    if p4_km.exists():
        slide_2img(prs, f"{label} -- UMAP z_proj: K-means clusters & training labels",
                   p4_km,
                   p4_cond if p4_cond.exists() else None,
                   "K-means cluster ID", "Condition (training labels)")

    # 7. PHATE z_recon 4ds
    ph_ann  = ev / "phate_4ds_annotation.png"
    ph_cond = ev / "phate_4ds_condition.png"
    if ph_ann.exists():
        slide_2img(prs, f"{label} -- PHATE z_recon (cross-dataset)",
                   ph_ann,
                   ph_cond if ph_cond.exists() else None,
                   "Annotation (labels)", "Condition")

    # 8. PHATE z_proj 4ds
    pp_ann  = ev / "phate_proj_4ds_annotation.png"
    pp_cond = ev / "phate_proj_4ds_condition.png"
    if pp_ann.exists():
        slide_2img(prs, f"{label} -- PHATE z_proj (cross-dataset)",
                   pp_ann,
                   pp_cond if pp_cond.exists() else None,
                   "Annotation (labels)", "Condition")

    # 9. umap_combo (multi-dataset training models)
    uc_cond  = ev / "umap_combo_condition.png"
    uc_condt = ev / "umap_combo_condition_with_test.png"
    uc_split = ev / "umap_combo_split.png"
    if uc_cond.exists():
        slide_3img(prs, f"{label} -- UMAP combo (training datasets)",
                   [uc_cond,
                    uc_condt if uc_condt.exists() else None,
                    uc_split if uc_split.exists() else None],
                   ["Condition", "Condition + test", "Train/Val/Test split"])

    # 10. Local UMAP (training data only): annotation / condition / kmeans
    u_ann  = ev / "umap_annotation.png"
    u_cond = ev / "umap_condition.png"
    u_km   = ev / "umap_kmeans.png"
    if u_ann.exists():
        slide_3img(prs, f"{label} -- UMAP (training data)",
                   [u_ann,
                    u_cond if u_cond.exists() else None,
                    u_km   if u_km.exists()   else None],
                   ["FA annotation", "Condition", "K-means clusters"])

    # 11. PHATE (local training data)
    p_ann  = ev / "phate_annotation.png"
    p_cond = ev / "phate_condition.png"
    p_km   = ev / "phate_kmeans.png"
    if p_ann.exists():
        slide_3img(prs, f"{label} -- PHATE (training data)",
                   [p_ann,
                    p_cond if p_cond.exists() else None,
                    p_km   if p_km.exists()   else None],
                   ["FA annotation", "Condition", "K-means clusters"])

    # 12. Cluster panels (tiled TIF)
    _cluster_panels(prs, label, ev)

    # 13. Classification confusion matrices
    for cls_sub in ("fa_cls_zrecon", "fa_cls_zproj", "pos_cls_zrecon"):
        cls_dir = d / cls_sub
        if not cls_dir.is_dir():
            continue
        conf = next(iter(sorted(cls_dir.glob("confusion_*_norm.png"))), None)
        if conf:
            slide_1img(prs, f"{label} -- FA Classification ({cls_sub})",
                       conf, conf.stem)
        rep_csv = cls_dir / "metrics.csv"
        if rep_csv.exists():
            try:
                rep = pd.read_csv(rep_csv)
                hdrs = list(rep.columns)
                rows = [[str(v) for v in r] for _, r in rep.iterrows()]
                slide_table(prs,
                            f"{label} -- Classification Metrics ({cls_sub})",
                            hdrs, rows, font_size=8)
            except Exception:
                pass


# ── build group 1 ─────────────────────────────────────────────────────────────

def build_group1(prs):
    slide_section(prs,
        "Group 1: CIO-vinc ConAE / SupConAE",
        "18 models  CIO normalisation  vinc training dataset  nL1 recon loss  500 epochs\n"
        "lat12 (1-channel) lc sweep  +  lat18 / lat24 multichannel (3ch/4ch) lc sweep")

    hdrs = ["Model", "Type", "Lat", "Proj", "Ch", "Channels", "lambda_contrast",
            "Recon loss", "Epochs", "Norm"]
    rows = [[lbl, pfx.capitalize(), str(lat), str(proj), str(ch), chstr,
             str(lc), "nL1", "500", "CIO"]
            for lbl, _, pfx, lat, proj, ch, chstr, lc in G1_MODELS]
    slide_table(prs, "Group 1 -- Model Configurations", hdrs, rows)

    for lbl, dname, pfx, *_ in G1_MODELS:
        d = RUNS / dname
        exists_note = "[OK]" if d.exists() else "[MISSING]"
        slide_section(prs, lbl, f"{dname}  {exists_note}")
        _model_slides(prs, lbl, d, pfx)


# ── build ds_combo group ───────────────────────────────────────────────────────

def _parse_combo(combo: str) -> list[str]:
    parts = combo.split("_")
    ds = []
    i = 0
    while i < len(parts):
        if i + 1 < len(parts) and parts[i] == "nih" and parts[i+1] == "3t3":
            ds.append("nih3t3"); i += 2
        else:
            ds.append(parts[i]); i += 1
    return ds


def build_ds_combo_group(prs, group_label: str, parent_name: str, description: str):
    parent_dir = RUNS / parent_name

    slide_section(prs, f"{group_label}: {parent_name}", description)

    # Overview table
    hdrs = ["Combo", "Datasets", "# DS", "Dir exists"]
    rows = []
    for combo in COMBO_LIST:
        ds_list = _parse_combo(combo)
        exists = "Yes" if (parent_dir / combo).exists() else "No"
        rows.append([combo, " + ".join(ds_list), str(len(ds_list)), exists])
    slide_table(prs, f"{group_label} -- Dataset Combinations", hdrs, rows, font_size=9)

    # Group-level summary metrics (parent-level CSV)
    parent_mcsv = parent_dir / "cross_dataset_recon_metrics.csv"
    if parent_mcsv.exists():
        try:
            df = pd.read_csv(parent_mcsv)
            mc = [c for c in ("recon_nl1", "recon_l1", "recon_mse") if c in df.columns]
            if mc and "variant" in df.columns and "split" in df.columns:
                agg = (df.groupby(["variant", "split"])[mc]
                         .mean().reset_index().round(4))
                splits = ["train", "val", "test"]
                col_order = ["variant"] + [f"{s}_{m}" for s in splits for m in mc]
                wide_rows = {}
                for var in sorted(agg["variant"].unique()):
                    row = [var]
                    sub = agg[agg["variant"] == var]
                    for s in splits:
                        sv = sub[sub["split"] == s]
                        for m in mc:
                            row.append(f"{sv[m].values[0]:.4f}" if len(sv) else "--")
                    wide_rows[var] = row
                hdrs2 = col_order
                rows2 = [wide_rows[var] for var in sorted(wide_rows)]
                slide_table(prs, f"{group_label} -- Summary Metrics (all combos)",
                            hdrs2, rows2, font_size=7)
        except Exception as e:
            print(f"  WARNING group metrics: {e}")

    # Per-combo results
    for combo in COMBO_LIST:
        model_dir = parent_dir / combo
        label = f"{group_label} -- {combo}"
        ds_list = _parse_combo(combo)
        detail = " + ".join(ds_list) + f"  ({len(ds_list)} dataset{'s' if len(ds_list)>1 else ''})"
        slide_section(prs, combo, f"{group_label}  --  {detail}")
        _model_slides(prs, label, model_dir, "contrastive")


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs,
        "ConAE / SupConAE Training Results",
        "Group 1: CIO-vinc Flat Models  (lat12 + lat18/lat24 multichannel, lambda_contrast sweep)\n"
        "Group 2: ds_combo_enlcrop_clip01_l1  (15 dataset combinations)\n"
        "Group 3: ds_combo_enlcrop_sc2_clip02_l1  (15 dataset combinations)\n\n"
        "2026-07-26")

    print("Building Group 1 ...")
    build_group1(prs)

    for g_label, p_name, desc in DS_GROUPS:
        print(f"Building {g_label} ({p_name}) ...")
        build_ds_combo_group(prs, g_label, p_name, desc)

    prs.save(str(OUT))
    n = len(prs.slides)
    print(f"\nSaved: {OUT}  ({n} slides)")
    print(f"Rsync to laptop:  rsync -avh --progress {OUT} lding@lding-Precision-3680:/home/lding/Desktop/")


if __name__ == "__main__":
    main()
