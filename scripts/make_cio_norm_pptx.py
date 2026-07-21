#!/usr/bin/env python3
"""
make_cio_norm_pptx.py
=====================
Before / after CIO normalization (no rolling ball, scale=1) histograms.

- Before : raw uint16 / 65535 CZI pixels, all 4 channels, subsampled per frame
- After  : CIO-normalized source_frames/cio pixels, all 4 channels, same frames

Slides
------
  1. Title
  2-5. Per-dataset: 4×2 grid (ch0/pax/zyx/act × before/after), each curve = one frame
  6-9. Per-channel summary: all 4 dataset mean curves (before vs after)
  10.  Patch stats table (pax patches from patches/cio/{ds}/{cond}/tiff_patches32_mr10)

Usage:
  cd /net/projects/CLS/lding/gitcode/SubCellAE
  PYTHONPATH=... python3 scripts/make_cio_norm_pptx.py
"""

from __future__ import annotations

import io
import random
from pathlib import Path

import czifile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import tifffile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT        = Path("/net/projects/CLS/lding/data/fa_data_analysis")
FRAMES_BASE = ROOT / "ae_results/source_frames/cio"
PATCH_BASE  = ROOT / "ae_results/patches/cio"
OUT_PPTX    = Path("cio_norm_before_after.pptx")

N_PIXELS_PER_FRAME = 10_000   # random pixel subsample per frame per channel
N_FRAMES_MAX       = 15       # max frames loaded per condition
N_BINS             = 200
random.seed(42)
rng = np.random.default_rng(42)

# ch0 name varies by dataset
DATASETS = [
    {
        "key":       "vinc",
        "label":     "ds1 - vinc",
        "color":     "#2E86C1",
        "ch_names":  ["vinc", "pax", "zyx", "act"],
        "conditions": [
            {
                "cond":    "control",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Control",
                "frame_dir": FRAMES_BASE / "vinc/control",
            },
            {
                "cond":    "ycomp",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250311_eGFPZyxin488_Phalloidin405_Vinculin(rb)647_paxillin(m)568/Ycomp",
                "frame_dir": FRAMES_BASE / "vinc/ycomp",
            },
        ],
    },
    {
        "key":       "pfak",
        "label":     "ds2 - pfak",
        "color":     "#C44E52",
        "ch_names":  ["pfak", "pax", "zyx", "act"],
        "conditions": [
            {
                "cond":    "control",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Control",
                "frame_dir": FRAMES_BASE / "pfak/control",
            },
            {
                "cond":    "ycomp",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250720_eGFP-Zyxin 488, Phalloidin 405, pFAK (rb) 647, paxillin(m)568/072025/Ycomp",
                "frame_dir": FRAMES_BASE / "pfak/ycomp",
            },
        ],
    },
    {
        "key":       "ppax",
        "label":     "ds3 - ppax",
        "color":     "#55A868",
        "ch_names":  ["ppax", "pax", "zyx", "act"],
        "conditions": [
            {
                "cond":    "control",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Control",
                "frame_dir": FRAMES_BASE / "ppax/control",
            },
            {
                "cond":    "ycomp",
                "czi_dir": ROOT / "fa_data/other_paxillin/20250721_eGFP-Zyxin 488_Phalloidin405_pPaxy118(rb) 647_Pax(m)568/Y-comp",
                "frame_dir": FRAMES_BASE / "ppax/ycomp",
            },
        ],
    },
    {
        "key":       "nih3t3",
        "label":     "ds4 - nih3t3",
        "color":     "#DD8452",
        "ch_names":  ["vinc", "pax", "zyx", "act"],
        "conditions": [
            {
                "cond":    "control",
                "czi_dir": ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/Control",
                "frame_dir": FRAMES_BASE / "nih3t3/control",
            },
            {
                "cond":    "ycomp",
                "czi_dir": ROOT / "fa_data/other_paxillin/20260227_NIH3T3_ZyxinGFP,Phalloidin405,Vinc_rb647,Pax_m555_reduced_size_AH/YCompound",
                "frame_dir": FRAMES_BASE / "nih3t3/ycomp",
            },
        ],
    },
]


def _subsample(arr: np.ndarray, n: int) -> np.ndarray:
    if len(arr) <= n:
        return arr.ravel().astype(np.float32)
    idx = rng.choice(len(arr.ravel()), size=n, replace=False)
    return arr.ravel()[idx].astype(np.float32)


def _load_czi_ch(czi_path: Path, ch_idx: int) -> np.ndarray:
    """Return float32 (H, W): raw uint16 / 65535 for given channel."""
    arr = czifile.imread(str(czi_path))
    # shape typically (1, 1, C, 1, H, W, 1) or (1, C, H, W)
    arr = arr.squeeze()
    if arr.ndim == 3:
        ch = arr[ch_idx]
    elif arr.ndim == 2:
        ch = arr
    else:
        ch = arr[ch_idx]
    return ch.astype(np.float32) / 65535.0


def _load_frame(frame_dir: Path, cond: str, fidx: int, ch_name: str) -> np.ndarray | None:
    """Load CIO-normalized frame; return None if file missing."""
    p = frame_dir / f"{cond}_f{fidx:04d}_{ch_name}.tif"
    if not p.exists():
        return None
    return tifffile.imread(str(p)).astype(np.float32)


def gather_dataset(ds: dict) -> dict:
    """
    Returns {
      "before": {ch_name: [arr_frame0, arr_frame1, ...]},
      "after":  {ch_name: [arr_frame0, arr_frame1, ...]},
      "n_frames": int,
    }
    Each arr is a 1-D float32 subsample of N_PIXELS_PER_FRAME pixels.
    """
    ch_names = ds["ch_names"]
    before = {c: [] for c in ch_names}
    after  = {c: [] for c in ch_names}
    n_frames = 0

    for cond_info in ds["conditions"]:
        cond      = cond_info["cond"]
        czi_dir   = Path(cond_info["czi_dir"])
        frame_dir = Path(cond_info["frame_dir"])

        czi_files = sorted(czi_dir.glob("*.czi")) if czi_dir.exists() else []
        if not czi_files:
            print(f"    SKIP — no CZI in {czi_dir}")
            continue
        if not frame_dir.exists():
            print(f"    SKIP — frame_dir missing: {frame_dir}")
            continue

        take = min(len(czi_files), N_FRAMES_MAX)
        indices = list(range(take))
        print(f"    [{cond}] {take}/{len(czi_files)} frames × {len(ch_names)} channels")

        for fidx in indices:
            czi_path = czi_files[fidx]
            frame_ok = True

            for ci, ch_name in enumerate(ch_names):
                # After CIO
                after_frame = _load_frame(frame_dir, cond, fidx, ch_name)
                if after_frame is None:
                    print(f"      WARN missing: {frame_dir}/{cond}_f{fidx:04d}_{ch_name}.tif")
                    frame_ok = False
                    break
                after[ch_name].append(_subsample(after_frame, N_PIXELS_PER_FRAME))

            if not frame_ok:
                continue

            # Before: load CZI once per frame, extract each channel
            try:
                for ci, ch_name in enumerate(ch_names):
                    raw = _load_czi_ch(czi_path, ci)
                    before[ch_name].append(_subsample(raw, N_PIXELS_PER_FRAME))
            except Exception as e:
                print(f"      WARN CZI load failed {czi_path.name}: {e}")
                for ch_name in ch_names:
                    if before[ch_name] and len(before[ch_name]) > len(after[after_ch_name := ch_names[0]]) - 1:
                        before[ch_name].pop()
                continue

            n_frames += 1

    return {"before": before, "after": after, "n_frames": n_frames}


def _plot_curves_ch(ax, arrays: list[np.ndarray], bold_color: str,
                    alpha: float = 0.35, lw: float = 0.9) -> tuple:
    n = len(arrays)
    if n == 0:
        return None, None
    cmap = plt.get_cmap("gist_rainbow")
    all_vals = np.concatenate(arrays)
    xmin, xmax = float(np.percentile(all_vals, 0.5)), float(np.percentile(all_vals, 99.5))
    bins = np.linspace(xmin, xmax, N_BINS + 1)
    mean_hist = np.zeros(N_BINS, dtype=np.float64)
    for i, arr in enumerate(arrays):
        color = cmap(i / max(n - 1, 1))
        counts, _ = np.histogram(arr, bins=bins, density=True)
        ax.plot(bins[:-1], counts, color=color, alpha=alpha, lw=lw)
        mean_hist += counts
    mean_hist /= n
    ax.plot(bins[:-1], mean_hist, color=bold_color, lw=2.5, alpha=0.9, label="mean")
    return bins, mean_hist


def make_dataset_figure(ds: dict, data: dict) -> plt.Figure:
    """4-row × 2-col figure: row=channel, col=before/after."""
    ch_names = ds["ch_names"]
    n_ch = len(ch_names)
    fig, axes = plt.subplots(n_ch, 2, figsize=(14, 4 * n_ch))
    fig.suptitle(f"{ds['label']}  —  Before / After CIO  (no rolling ball, scale=1)",
                 fontsize=14, fontweight="bold")

    for row, ch_name in enumerate(ch_names):
        before_arrs = data["before"].get(ch_name, [])
        after_arrs  = data["after"].get(ch_name, [])

        ax_b = axes[row, 0]
        ax_a = axes[row, 1]

        _plot_curves_ch(ax_b, before_arrs, bold_color=ds["color"])
        ax_b.set_title(f"{ch_name}  — Before  (raw, uint16/65535)", fontsize=9)
        ax_b.set_xlabel("Intensity (0–1 scale)")
        ax_b.set_ylabel("Density")
        ax_b.legend(fontsize=7)

        _plot_curves_ch(ax_a, after_arrs, bold_color=ds["color"])
        ax_a.set_title(f"{ch_name}  — After CIO  (scale=1, no rolling ball)", fontsize=9)
        ax_a.set_xlabel("CIO normalized value")
        ax_a.set_ylabel("Density")
        ax_a.axvline(1.0, color="red", ls="--", lw=1.2, label="cell mean (=1)")
        ax_a.axvline(0.0, color="gray", ls=":", lw=1.0, label="background (≈0)")
        ax_a.legend(fontsize=7)

    fig.text(0.5, 0.005,
             f"{data['n_frames']} frames  |  {N_PIXELS_PER_FRAME:,} pixels/frame  |  bold=mean  |  each color=one frame",
             ha="center", fontsize=8, color="#555555")
    plt.tight_layout(rect=[0, 0.02, 1, 0.97])
    return fig


def make_channel_summary_figure(ch_name: str, ds_list: list[dict],
                                 ds_data: list[dict]) -> plt.Figure:
    """All 4 datasets' before vs after for one channel."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    fig.suptitle(f"Channel: {ch_name}  —  Before / After CIO  (all datasets)",
                 fontsize=13, fontweight="bold")

    for ds, data in zip(ds_list, ds_data):
        before_arrs = data["before"].get(ch_name, [])
        after_arrs  = data["after"].get(ch_name, [])
        if not before_arrs or not after_arrs:
            continue
        # Just plot the mean curve per dataset
        fig_tmp, ax_tmp = plt.subplots(1, 2)
        bins_b, mean_b = _plot_curves_ch(ax_tmp[0], before_arrs, bold_color="k", alpha=0, lw=0)
        bins_a, mean_a = _plot_curves_ch(ax_tmp[1], after_arrs,  bold_color="k", alpha=0, lw=0)
        plt.close(fig_tmp)

        if mean_b is not None:
            axes[0].plot(bins_b[:-1], mean_b, color=ds["color"], lw=1.5, label=ds["label"])
        if mean_a is not None:
            axes[1].plot(bins_a[:-1], mean_a, color=ds["color"], lw=1.5, label=ds["label"])

    for ax, title, xlabel in [
        (axes[0], f"{ch_name} — Before  (raw, uint16/65535)",  "Intensity (0–1)"),
        (axes[1], f"{ch_name} — After CIO  (scale=1)",          "CIO normalized value"),
    ]:
        ax.set_title(title, fontsize=10)
        ax.set_xlabel(xlabel)
        ax.set_ylabel("Density")
        ax.legend(fontsize=9)

    axes[1].axvline(1.0, color="red", ls="--", lw=1.2, label="cell mean (=1)")
    axes[1].legend(fontsize=9)

    fig.text(0.5, 0.01, "Mean over all frames per dataset",
             ha="center", fontsize=9, color="#555555")
    plt.tight_layout(rect=[0, 0.04, 1, 1])
    return fig


def make_patch_stats_figure(ds_list: list[dict]) -> plt.Figure:
    """Violin + summary table of pax patch stats from patches/cio/."""
    import glob

    print("  Computing mr10 patch stats from patches/cio ...")
    ds_stats, all_vals = [], {"mean": [], "std": [], "p95": [], "max": []}
    for ds in ds_list:
        vals_m, vals_s, vals_p, vals_x = [], [], [], []
        for cond in ["control", "ycomp"]:
            pat = str(PATCH_BASE / ds["key"] / cond / "tiff_patches32_mr10")
            files = glob.glob(f"{pat}/*.tif")
            if not files:
                print(f"    WARN no patches in {pat}")
                continue
            if len(files) > 2000:
                files = random.sample(files, 2000)
            for fp in files:
                im = tifffile.imread(fp).astype(np.float32)
                vals_m.append(float(im.mean()))
                vals_s.append(float(im.std()))
                vals_p.append(float(np.percentile(im, 95)))
                vals_x.append(float(im.max()))
        n = len(vals_m)
        s = dict(label=ds["label"], color=ds["color"], n=n,
                 mean=np.mean(vals_m) if vals_m else 0,
                 std=np.mean(vals_s) if vals_s else 0,
                 p95=np.mean(vals_p) if vals_p else 0,
                 max=np.mean(vals_x) if vals_x else 0)
        ds_stats.append(s)
        all_vals["mean"].append(np.array(vals_m))
        all_vals["std"].append(np.array(vals_s))
        all_vals["p95"].append(np.array(vals_p))
        all_vals["max"].append(np.array(vals_x))
        print(f"    {ds['key']:8s}  n={n:5d}  mean={s['mean']:.3f}  "
              f"std={s['std']:.3f}  p95={s['p95']:.3f}  max={s['max']:.3f}")

    if not any(s["n"] > 0 for s in ds_stats):
        fig, ax = plt.subplots(figsize=(14, 3))
        ax.text(0.5, 0.5, "patches/cio not yet available — run sbatch_patchprep_cio_mr10.sh first",
                ha="center", va="center", fontsize=14, transform=ax.transAxes, color="red")
        ax.axis("off")
        return fig

    fig = plt.figure(figsize=(14, 8))
    gs = fig.add_gridspec(2, 4, hspace=0.45, wspace=0.35)
    metric_titles = [("mean","Patch mean"), ("std","Patch std\n(local contrast)"),
                     ("p95","Patch p95"),   ("max","Patch max")]
    colors = [s["color"] for s in ds_stats]
    ds_with_data = [i for i, s in enumerate(ds_stats) if s["n"] > 0]

    for col, (metric, title) in enumerate(metric_titles):
        ax = fig.add_subplot(gs[0, col])
        data = [all_vals[metric][i] for i in ds_with_data]
        if not any(len(d) > 0 for d in data):
            continue
        positions = list(range(len(data)))
        vp = ax.violinplot(data, positions=positions, showmedians=True, showextrema=False)
        for body, i in zip(vp["bodies"], ds_with_data):
            body.set_facecolor(colors[i]); body.set_alpha(0.7)
        vp["cmedians"].set_color("black")
        ax.set_xticks(positions)
        ax.set_xticklabels([ds_stats[i]["label"].split(" - ")[0] for i in ds_with_data], fontsize=8)
        ax.set_title(title, fontsize=9)
        ax.axhline(1.0, color="red", ls="--", lw=0.8, alpha=0.6)

    ax_tbl = fig.add_subplot(gs[1, :])
    ax_tbl.axis("off")
    rows = []
    for s in ds_stats:
        rows.append([s["label"], f"{s['n']:,}",
                     f"{s['mean']:.3f}", f"{s['std']:.3f}",
                     f"{s['p95']:.3f}",  f"{s['max']:.3f}"])
    tbl = ax_tbl.table(cellText=rows,
                       colLabels=["Dataset", "n patches", "mean", "std", "p95", "max"],
                       loc="center", cellLoc="center")
    tbl.auto_set_font_size(False); tbl.set_fontsize(11); tbl.scale(1, 2.2)
    import matplotlib.colors as mcolors
    for (r, c), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#2E4057")
            cell.set_text_props(color="white", fontweight="bold")
        elif r > 0:
            rgba = list(mcolors.to_rgba(colors[r - 1])); rgba[3] = 0.18
            cell.set_facecolor(rgba)

    fig.suptitle("PAX patch statistics — patches/cio mr10  (CIO, no rolling ball, scale=1)",
                 fontsize=13)
    return fig


# ── PPT helpers ──────────────────────────────────────────────────────────────
W_SLIDE = Inches(13.33)
H_SLIDE = Inches(7.5)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W_SLIDE
    prs.slide_height = H_SLIDE
    return prs


def _fig_to_buf(fig: plt.Figure) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.33), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Intensity Normalization: Before / After CIO"
    p.font.size = Pt(34); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.33), Inches(1.2))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = ("No rolling ball  |  scale=1  |  all 4 channels  |  all 4 datasets\n"
               "Before: raw uint16/65535  |  After: cell_insideoutside / cell_mean")
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC)
    p2.alignment = PP_ALIGN.CENTER


def _image_slide(prs: Presentation, buf: io.BytesIO) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(buf, Inches(0.1), Inches(0.1),
                             width=Inches(13.13), height=Inches(7.3))


def main() -> None:
    prs = _new_prs()
    _title_slide(prs)

    ds_list, ds_data = [], []
    for ds in DATASETS:
        print(f"\n=== {ds['label']} ===")
        data = gather_dataset(ds)
        print(f"  {data['n_frames']} frames loaded")
        ds_list.append(ds)
        ds_data.append(data)

        if data["n_frames"] == 0:
            print("  No data — skipping slide")
            continue

        print("  Making dataset figure...")
        fig = make_dataset_figure(ds, data)
        _image_slide(prs, _fig_to_buf(fig))

    print("\n=== Channel summary slides ===")
    all_ch_names = ["pax", "zyx", "act"]
    ch0_names    = list({ds["ch_names"][0] for ds in DATASETS})
    for ch_name in (ch0_names + all_ch_names):
        print(f"  Channel: {ch_name}")
        fig = make_channel_summary_figure(ch_name, ds_list, ds_data)
        _image_slide(prs, _fig_to_buf(fig))

    print("\n=== Patch stats slide ===")
    fig = make_patch_stats_figure(ds_list)
    _image_slide(prs, _fig_to_buf(fig))

    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
