#!/usr/bin/env python3
"""
make_pptx_supcon_analysis.py

Slide deck: SupCon label-efficiency — guaranteed labeled pairs analysis.

Slides
------
1. Title
2. Problem: sparse labels → SupCon degenerates to NT-Xent
3. Solution 1: LabeledAwareBatchSampler
4. Solution 2: split loss weights (lambda_supcon)
5. Experiment design
6. Result — sampler only (λ_supcon = -1, i.e. same as λ_contrast)
7. Result — lambda_supcon sweep (1.0, 1.5, 2.0)
8. Conclusion

Usage
-----
  python scripts/make_pptx_supcon_analysis.py
  python scripts/make_pptx_supcon_analysis.py --out path/to/out.pptx
"""
from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from sklearn.ensemble import GradientBoostingClassifier
from sklearn.metrics import balanced_accuracy_score
from sklearn.preprocessing import LabelEncoder
from sklearn.utils.class_weight import compute_sample_weight
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parents[1]
DATA_ROOT = Path("/net/projects/CLS/lding/data/fa_data_analysis")
RUN_DIR   = DATA_ROOT / "ae_results" / "contrastive_run"
LAB_DIR   = DATA_ROOT / "labelling"
RESULTS   = REPO_ROOT / "results"

SW = Inches(13.33)
SH = Inches(7.5)

C_TITLE  = RGBColor(0x1A, 0x1A, 0x2E)
C_HEAD   = RGBColor(0x16, 0x21, 0x3E)
C_BODY   = RGBColor(0x1A, 0x1A, 0x1A)
C_ACCENT = RGBColor(0x0F, 0x3D, 0x79)
C_GOOD   = RGBColor(0x1A, 0x6B, 0x30)
C_WARN   = RGBColor(0x8B, 0x45, 0x00)
C_BAD    = RGBColor(0xAA, 0x22, 0x22)
C_GREY   = RGBColor(0x66, 0x66, 0x66)

NPI_ORDER  = ["10", "25", "50", "75", "100", "all"]
Z_COLS     = [f"z_{i}" for i in range(12)]
FULL_ANN   = LAB_DIR / "vinc_control_label_Annabel_20260715_1554_2cls.csv"
FRAMES     = {"train": [0], "test": [1, 2, 3]}

# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _txt(slide, text, left, top, width, height, *,
         bold=False, italic=False, size_pt=13, color=C_BODY,
         align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.size      = Pt(size_pt)
    run.font.color.rgb = color
    return txb

def _rule(slide, top, width=None, left=None, thickness_pt=0.75):
    w = width or SW - Inches(1.0)
    l = left  or Inches(0.5)
    ln = slide.shapes.add_connector(1, l, top, l + w, top)
    ln.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.width = Pt(thickness_pt)

def _header(slide, title, subtitle=""):
    _txt(slide, title,
         Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.55),
         bold=True, size_pt=22, color=C_HEAD)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.35),
             size_pt=12, color=C_GREY)
    _rule(slide, Inches(0.97))

def _add_fig(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)

def _img(slide, path, left, top, width, height):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)

def _bullet(slide, items, left, top, width, height, size_pt=12, indent="  • "):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.size      = Pt(size_pt)
        run.font.color.rgb = C_BODY

# ---------------------------------------------------------------------------
# Evaluation helpers
# ---------------------------------------------------------------------------

def _extract_frame(fn):
    m = re.search(r"_f(\d+)", fn)
    return int(m.group(1)) if m else -1

def _eval_dir(base_dir, full_ann):
    runs = sorted(base_dir.glob("le_c0_npi*_r*"))
    records = []
    for run_dir in runs:
        lat_csv = run_dir / "latents.csv"
        name    = run_dir.name
        ann_csv = LAB_DIR / "le_clean" / f"{name}.csv"
        if not lat_csv.exists() or not ann_csv.exists():
            continue
        latents = pd.read_csv(lat_csv)
        latents["frame"] = latents["filename"].apply(_extract_frame)
        ann_train = pd.read_csv(ann_csv)
        train_l = latents[latents["frame"].isin(FRAMES["train"])].merge(
            ann_train[["filename", "label"]], on="filename", how="inner")
        if len(train_l) == 0:
            continue
        le = LabelEncoder()
        y_tr = le.fit_transform(train_l["label"])
        clf  = GradientBoostingClassifier(n_estimators=200, max_depth=4,
                                          learning_rate=0.05, random_state=42)
        clf.fit(train_l[Z_COLS].values, y_tr,
                sample_weight=compute_sample_weight("balanced", y_tr))
        test_l = latents[latents["frame"].isin(FRAMES["test"])].merge(
            full_ann[["filename", "label"]], on="filename", how="inner")
        if len(test_l) == 0:
            continue
        m2 = re.match(r"le_c\d+_npi(\w+)_r(\d+)$", name)
        records.append({"npi": m2.group(1), "repeat": int(m2.group(2)),
                        "balanced_acc": balanced_accuracy_score(
                            le.transform(test_l["label"]),
                            clf.predict(test_l[Z_COLS].values))})
    df = pd.DataFrame(records)
    s  = df.groupby("npi")["balanced_acc"].agg(mean="mean", std="std").reset_index()
    s["mean_pct"] = (s["mean"] * 100).round(1)
    s["std_pct"]  = (s["std"]  * 100).round(1)
    s["order"]    = s["npi"].apply(lambda x: NPI_ORDER.index(x) if x in NPI_ORDER else 99)
    return s.sort_values("order")

# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = _blank(prs)
    _txt(slide,
         "SupCon Label Efficiency:\nGuaranteed Labeled Pairs",
         Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.8),
         bold=True, size_pt=36, color=C_TITLE, align=PP_ALIGN.CENTER)
    _txt(slide,
         "Can we force labeled patches into every batch to strengthen supervision?\n"
         "And does up-weighting the supervised loss term help?",
         Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0),
         size_pt=16, color=C_GREY, align=PP_ALIGN.CENTER)
    _rule(slide, Inches(5.3), width=Inches(10.0), left=Inches(1.65))


def slide_problem(prs):
    slide = _blank(prs)
    _header(slide, "Problem: Sparse Labels → SupCon Degenerates to NT-Xent",
            "With small K, labeled patches almost never appear in a random batch")

    # Left: explanation text
    _txt(slide, "Setup", Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "14,879 total patches per frame (vinc control)",
        "K labels spread across 1 training frame",
        "Batch size = 128",
        "SupCon loss uses class positives for labeled anchors,\n     augmentation pairs for unlabeled anchors",
    ], Inches(0.5), Inches(1.45), Inches(5.8), Inches(1.8), size_pt=12)

    _txt(slide, "Expected labeled patches per batch", Inches(0.5), Inches(3.3),
         Inches(5.8), Inches(0.35), bold=True, size_pt=14, color=C_ACCENT)
    _txt(slide, "E[labeled/batch]  =  K / N_patches  ×  batch_size",
         Inches(0.5), Inches(3.65), Inches(5.8), Inches(0.4),
         size_pt=13, italic=True, color=C_BODY)

    # Table: npi → expected labeled per batch
    rows = [
        ("npi=10",  "10",  "0.09",  "~91% of batches have zero labeled patches"),
        ("npi=25",  "25",  "0.21",  ""),
        ("npi=50",  "50",  "0.43",  ""),
        ("npi=75",  "75",  "0.64",  ""),
        ("npi=100", "100", "0.86",  ""),
        ("npi=all", "145", "1.25",  "Still < 2 on average"),
    ]
    headers = ["npi",  "K labels", "E[labeled/batch]", ""]
    col_x   = [Inches(0.5), Inches(1.5), Inches(2.8), Inches(4.1)]
    col_w   = [Inches(0.9), Inches(1.1), Inches(1.2), Inches(2.3)]
    row_h   = Inches(0.32)
    top0    = Inches(4.15)

    for ci, h in enumerate(headers):
        _txt(slide, h, col_x[ci], top0, col_w[ci], row_h,
             bold=True, size_pt=10, color=C_HEAD)
    _rule(slide, top0 + row_h, width=Inches(5.8), left=Inches(0.5), thickness_pt=0.5)
    for ri, row in enumerate(rows):
        top = top0 + row_h + Inches(0.02) + ri * row_h
        for ci, val in enumerate(row):
            color = C_WARN if ri == 0 and ci == 2 else C_BODY
            _txt(slide, val, col_x[ci], top, col_w[ci], row_h,
                 size_pt=10, color=color)

    _txt(slide, "Consequence: contrast loss is identical across K=10 to K=145\n→ SupCon behaves like NT-Xent (fully self-supervised)",
         Inches(0.5), Inches(6.7), Inches(5.8), Inches(0.6),
         size_pt=11, italic=True, color=C_WARN)

    # Right: figure showing E[labeled/batch] vs npi
    fig, ax = plt.subplots(figsize=(5.5, 4.2), facecolor="white")
    ax.set_facecolor("white")
    npi_k   = np.array([10, 25, 50, 75, 100, 145])
    N_patch = 14879
    batch   = 128
    e_lab   = npi_k / N_patch * batch
    ax.bar(range(len(npi_k)), e_lab, color="#4E79A7", alpha=0.8, edgecolor="white")
    ax.axhline(2, color="#E15759", linestyle="--", linewidth=1.5,
               label="Target: 2 labeled/batch\n(needed for class pairs)")
    ax.set_xticks(range(len(npi_k)))
    ax.set_xticklabels(["10","25","50","75","100","all"], fontsize=10)
    ax.set_xlabel("n_per_img (K labels)", fontsize=11)
    ax.set_ylabel("E[labeled patches / batch]", fontsize=11)
    ax.set_title("With standard shuffle:\nlabeled patches rarely appear", fontsize=11)
    ax.legend(fontsize=9)
    ax.grid(axis="y", alpha=0.3)
    _add_fig(slide, fig, Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.5))


def slide_solution_sampler(prs):
    slide = _blank(prs)
    _header(slide, "Solution 1: LabeledAwareBatchSampler",
            "Guarantee n_per_class labeled patches of each class in every batch")

    _txt(slide, "How it works", Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "Custom PyTorch Sampler replacing standard shuffle",
        "At each batch: reserve n_per_class slots per class from labeled pool",
        "Fill remaining slots with random unlabeled patches",
        "n_per_class = 0 → identical to standard shuffle (backward compat)",
        "Experiment: n_per_class = 2  →  4 labeled guaranteed / batch of 128",
    ], Inches(0.5), Inches(1.45), Inches(5.8), Inches(2.2), size_pt=12)

    _txt(slide, "Effect on batch composition", Inches(0.5), Inches(3.75),
         Inches(5.8), Inches(0.35), bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "npi=10:  labeled/batch: 0.09 → 4  (44× increase)",
        "npi=100: labeled/batch: 0.86 → 4  (4.7× increase)",
        "Labeled fraction still only 3.1%  (4/128)",
        "Log line confirms activation:",
    ], Inches(0.5), Inches(4.1), Inches(5.8), Inches(1.6), size_pt=12)

    _txt(slide,
         "LabeledAwareBatchSampler: n_per_class=2  n_classes=2\n"
         "  labeled=10  labeled_per_batch=4 / 128",
         Inches(0.7), Inches(5.75), Inches(5.4), Inches(0.7),
         size_pt=10, italic=True, color=RGBColor(0x33, 0x33, 0x33))

    # Right: batch composition diagram
    fig, axes = plt.subplots(1, 2, figsize=(5.8, 4.0), facecolor="white")
    for ax, mode, lab_n, title in zip(
            axes,
            ["Standard\nShuffle", "LabeledAware\nSampler"],
            [0.09, 4.0],
            ["Standard shuffle\n(npi=10)", "LabeledAwareBatchSampler\n(npi=10, n_per_class=2)"]):
        unl_n = 128 - lab_n
        ax.set_facecolor("white")
        ax.bar([0], [unl_n], color="#AAAAAA", label="Unlabeled", width=0.5)
        ax.bar([0], [lab_n], bottom=[unl_n], color="#E15759", label="Labeled", width=0.5)
        ax.set_xlim(-0.5, 0.5)
        ax.set_ylim(0, 135)
        ax.set_xticks([])
        ax.set_ylabel("Patches per batch", fontsize=10)
        ax.set_title(title, fontsize=9.5)
        ax.text(0, unl_n + lab_n / 2 + 1, f"{lab_n:.1f}", ha="center",
                va="bottom", fontsize=10, color="#AA2222", fontweight="bold")
        ax.legend(fontsize=8, loc="upper right")
        ax.grid(axis="y", alpha=0.3)
    fig.suptitle("Batch composition (batch_size=128)", fontsize=11, fontweight="bold")
    fig.tight_layout()
    _add_fig(slide, fig, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.2))


def slide_solution_lambda(prs):
    slide = _blank(prs)
    _header(slide, "Solution 2: Split Loss Weights (lambda_supcon)",
            "Apply higher weight to labeled-anchor losses, lower to unlabeled")

    _txt(slide, "Loss formulation", Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)

    _txt(slide,
         "Original:   total = λ_recon · L_recon  +  λ_contrast · L_contrast",
         Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.4),
         size_pt=13, italic=True, color=C_BODY)
    _txt(slide,
         "New:        total = λ_recon · L_recon  +  λ_supcon · L_labeled  +  λ_contrast · L_unlabeled",
         Inches(0.7), Inches(1.9), Inches(12.0), Inches(0.4),
         size_pt=13, italic=True, color=C_ACCENT)

    _txt(slide, "Key design decisions", Inches(0.5), Inches(2.5), Inches(5.8), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "L_labeled  = mean loss over labeled anchors in batch",
        "L_unlabeled = mean loss over unlabeled anchors in batch",
        "Shared denominator: all 2N samples in partition function",
        "λ_supcon = −1 (sentinel) → same as λ_contrast (backward compat)",
        "Existing experiments unaffected — only new YAMLs opt in",
        "Epoch log now reports sc= (supervised) and uc= (unsupervised) separately",
    ], Inches(0.5), Inches(2.85), Inches(5.8), Inches(2.4), size_pt=12)

    _txt(slide, "Expectation", Inches(0.5), Inches(5.4), Inches(5.8), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "λ_supcon > λ_contrast → labeled class pairs pulled together more strongly",
        "Combined with LabeledAwareBatchSampler: effect fires every batch",
        "Risk: too high → unlabeled structure degraded → GBC accuracy drops",
    ], Inches(0.5), Inches(5.75), Inches(5.8), Inches(1.5), size_pt=12)

    # Right: loss diagram
    fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor="white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    def box(x, y, w, h, label, sublabel="", fc="#EAF0FB", ec="#4E79A7"):
        rect = plt.Rectangle((x, y), w, h, fc=fc, ec=ec, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2 + (0.25 if sublabel else 0), label,
                ha="center", va="center", fontsize=10, fontweight="bold")
        if sublabel:
            ax.text(x + w/2, y + h/2 - 0.35, sublabel,
                    ha="center", va="center", fontsize=8.5, color="#555555")

    box(0.3, 7.5, 2.2, 1.2, "L_recon", "reconstruction", fc="#F5F5F5", ec="#999999")
    box(0.3, 5.2, 2.2, 1.2, "L_labeled", "labeled anchors", fc="#FDECEA", ec="#E15759")
    box(0.3, 2.9, 2.2, 1.2, "L_unlabeled", "unlabeled anchors", fc="#EAF0FB", ec="#4E79A7")

    ax.annotate("", xy=(4.0, 8.1), xytext=(2.5, 8.1),
                arrowprops=dict(arrowstyle="->", color="#999999", lw=1.5))
    ax.annotate("", xy=(4.0, 5.8), xytext=(2.5, 5.8),
                arrowprops=dict(arrowstyle="->", color="#E15759", lw=2.0))
    ax.annotate("", xy=(4.0, 3.5), xytext=(2.5, 3.5),
                arrowprops=dict(arrowstyle="->", color="#4E79A7", lw=1.5))

    ax.text(3.1, 8.3, "λ_recon = 1.0", fontsize=9, color="#555555")
    ax.text(3.1, 6.0, "λ_supcon = ?", fontsize=9.5, color="#AA2222", fontweight="bold")
    ax.text(3.1, 3.7, "λ_contrast = 0.5", fontsize=9, color="#4E79A7")

    ax.annotate("", xy=(6.8, 5.5), xytext=(5.5, 8.1),
                arrowprops=dict(arrowstyle="->", color="#333333", lw=1.2))
    ax.annotate("", xy=(6.8, 5.5), xytext=(5.5, 5.8),
                arrowprops=dict(arrowstyle="->", color="#333333", lw=1.2))
    ax.annotate("", xy=(6.8, 5.5), xytext=(5.5, 3.5),
                arrowprops=dict(arrowstyle="->", color="#333333", lw=1.2))

    box(6.8, 4.7, 2.8, 1.6, "Total Loss", "minimize", fc="#F0F7EE", ec="#1A6B30")
    ax.set_title("New loss decomposition", fontsize=11, fontweight="bold", pad=6)
    _add_fig(slide, fig, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.8))


def slide_experiment(prs):
    slide = _blank(prs)
    _header(slide, "Experiment Design",
            "cfg0 only — train frame 0, test frames 1/2/3 — reuses le_clean annotation CSVs")

    left_w = Inches(5.5)

    _txt(slide, "Conditions tested", Inches(0.5), Inches(1.1), left_w, Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    conds = [
        ("le_clean baseline",    "Standard random shuffle, λ_contrast=0.5 only"),
        ("le_supcon (sampler)",  "LabeledAwareBatchSampler n_per_class=2, λ_supcon=-1 (same as baseline)"),
        ("le_supcon λ=1.0",      "Sampler + labeled anchor weight = 1.0 (2× unlabeled)"),
        ("le_supcon λ=1.5",      "Sampler + labeled anchor weight = 1.5 (3× unlabeled)"),
        ("le_supcon λ=2.0",      "Sampler + labeled anchor weight = 2.0 (4× unlabeled)"),
    ]
    colors_cond = [C_BODY, C_ACCENT, C_BODY, C_BODY, C_BODY]
    top0 = Inches(1.5)
    for i, ((name, desc), col) in enumerate(zip(conds, colors_cond)):
        top = top0 + i * Inches(0.55)
        _txt(slide, f"• {name}:", Inches(0.5), top, Inches(2.5), Inches(0.5),
             bold=True, size_pt=11, color=col)
        _txt(slide, desc, Inches(3.1), top, Inches(2.9), Inches(0.5),
             size_pt=11, color=C_BODY)

    _txt(slide, "Fixed parameters", Inches(0.5), Inches(4.4), left_w, Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    _bullet(slide, [
        "Model: SupCon AE (ContrastiveAE + supervised NT-Xent head)",
        "Architecture: latent_dim=12, proj_dim=8, recon_loss=nl1",
        "Training: 500 epochs, lr=0.001, batch_size=128, weight_decay=1e-4",
        "Augmentation: enlarged crop (58px→32px), intensity scale [0.8, 1.2]",
        "Evaluation: GradientBoostingClassifier on 12-dim latents",
        "Test: all Annabel-labeled patches from frames 1, 2, 3 (539 patches)",
    ], Inches(0.5), Inches(4.75), left_w, Inches(2.4), size_pt=11)

    # Right: npi sweep table
    _txt(slide, "npi sweep", Inches(6.2), Inches(1.1), Inches(6.8), Inches(0.35),
         bold=True, size_pt=14, color=C_ACCENT)
    rows = [
        ("npi=10",  "10 labels (frame 0)",  "3 repeats"),
        ("npi=25",  "25 labels",             "3 repeats"),
        ("npi=50",  "50 labels",             "3 repeats"),
        ("npi=75",  "75 labels",             "3 repeats"),
        ("npi=100", "100 labels",            "3 repeats"),
        ("npi=all", "145 labels (all annotated)", "1 repeat"),
    ]
    col_x = [Inches(6.2), Inches(7.5), Inches(9.8)]
    col_w = [Inches(1.2), Inches(2.2), Inches(1.8)]
    row_h = Inches(0.36)
    top0  = Inches(1.5)
    for ci, h in enumerate(["Condition", "Labels", "Repeats"]):
        _txt(slide, h, col_x[ci], top0, col_w[ci], row_h,
             bold=True, size_pt=11, color=C_HEAD)
    _rule(slide, top0 + row_h, width=Inches(5.8), left=Inches(6.2), thickness_pt=0.5)
    for ri, row in enumerate(rows):
        top = top0 + row_h + ri * row_h
        for ci, val in enumerate(row):
            _txt(slide, val, col_x[ci], top, col_w[ci], row_h, size_pt=11)

    _txt(slide, "→  16 jobs per condition  ×  5 conditions  =  80 total runs",
         Inches(6.2), Inches(4.1), Inches(6.8), Inches(0.4),
         size_pt=11, color=C_GREY, italic=True)

    _txt(slide, "Metric: balanced accuracy on held-out frames\n(mean ± std over 3 repeats)",
         Inches(6.2), Inches(4.6), Inches(6.8), Inches(0.6),
         size_pt=11, color=C_BODY)


def slide_result_sampler(prs, full_ann):
    slide = _blank(prs)
    _header(slide, "Result 1: LabeledAwareBatchSampler (λ_supcon = baseline)",
            "Does guaranteeing 4 labeled patches per batch improve accuracy?")

    # Compute summaries
    s_cln = _eval_dir(RUN_DIR / "le_clean", full_ann)
    s_sup = _eval_dir(RUN_DIR / "le_supcon", full_ann)
    # filter le_clean to cfg0
    s_cln = s_cln  # already cfg0 only from _eval_dir

    # Left: table
    _txt(slide, "cfg0  |  balanced accuracy (mean ± std, %)",
         Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.35),
         bold=True, size_pt=13, color=C_ACCENT)

    col_x = [Inches(0.5), Inches(1.5), Inches(2.9), Inches(4.1)]
    col_w = [Inches(0.9), Inches(1.3), Inches(1.1), Inches(1.8)]
    row_h = Inches(0.34)
    top0  = Inches(1.5)
    for ci, h in enumerate(["npi", "le_clean", "le_supcon", "Δ"]):
        _txt(slide, h, col_x[ci], top0, col_w[ci], row_h,
             bold=True, size_pt=11, color=C_HEAD)
    _rule(slide, top0 + row_h, width=Inches(5.3), left=Inches(0.5), thickness_pt=0.5)

    npi_vals = s_cln["npi"].tolist()
    for ri, npi in enumerate(npi_vals):
        top = top0 + row_h + ri * row_h
        r_c = s_cln[s_cln["npi"] == npi].iloc[0]
        r_s = s_sup[s_sup["npi"] == npi].iloc[0] if npi in s_sup["npi"].values else None
        delta = (r_s["mean_pct"] - r_c["mean_pct"]) if r_s is not None else float("nan")
        delta_str = f"{delta:+.1f}%" if not np.isnan(delta) else "—"
        delta_col = C_GOOD if delta > 1 else (C_WARN if delta < -1 else C_BODY)

        _txt(slide, npi, col_x[0], top, col_w[0], row_h, size_pt=11)
        _txt(slide, f"{r_c['mean_pct']:.1f} ±{r_c['std_pct']:.1f}",
             col_x[1], top, col_w[1], row_h, size_pt=11)
        if r_s is not None:
            _txt(slide, f"{r_s['mean_pct']:.1f} ±{r_s['std_pct']:.1f}",
                 col_x[2], top, col_w[2], row_h, size_pt=11)
        _txt(slide, delta_str, col_x[3], top, col_w[3], row_h,
             size_pt=11, bold=True, color=delta_col)

    _txt(slide,
         "Key observation: npi=10 variance halved (±18.8 → ±8.9)\n"
         "but mean differences are within noise (3 repeats).\n"
         "sc loss does decrease; uc loss slightly increases.",
         Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.2),
         size_pt=11, color=C_BODY)

    # Right: curve
    _img(slide, RESULTS / "le_supcon_vs_clean_curve.png",
         Inches(6.2), Inches(1.05), Inches(6.9), Inches(5.8))


def slide_result_sweep(prs, full_ann):
    slide = _blank(prs)
    _header(slide, "Result 2: lambda_supcon Sweep (1.0, 1.5, 2.0)",
            "Does boosting labeled-anchor weight improve classification?")

    TAGS = {
        "le_clean":        (RUN_DIR / "le_clean",       "#4E79A7", "o", "-"),
        "ls1 (λ=1.0)":    (RUN_DIR / "le_supcon_ls1",  "#59A14F", "s", "--"),
        "ls15 (λ=1.5)":   (RUN_DIR / "le_supcon_ls15", "#F28E2B", "^", "--"),
        "ls2 (λ=2.0)":    (RUN_DIR / "le_supcon_ls2",  "#E15759", "D", "--"),
    }
    summaries = {}
    for lbl, (d, *_) in TAGS.items():
        summaries[lbl] = _eval_dir(d, full_ann)

    # Left: figure
    fig, ax = plt.subplots(figsize=(6.0, 4.8), facecolor="white")
    ax.set_facecolor("white")
    for lbl, (_, color, mk, ls) in TAGS.items():
        s = summaries[lbl]
        x = np.arange(len(s))
        ax.errorbar(x, s["mean_pct"], yerr=s["std_pct"],
                    fmt=f"{mk}{ls}", color=color, capsize=4,
                    linewidth=1.8, markersize=7, label=lbl)
    ax.set_xticks(x)
    ax.set_xticklabels(summaries["le_clean"]["npi"].tolist(), fontsize=10)
    ax.set_xlabel("n_per_img", fontsize=11)
    ax.set_ylabel("Balanced accuracy (%)", fontsize=11)
    ax.set_title("cfg0 — all λ_supcon values vs le_clean baseline", fontsize=10.5)
    ax.set_ylim(20, 105)
    ax.axhline(90, color="#AAAAAA", linestyle=":", linewidth=1.0, label="90% target")
    ax.legend(fontsize=9)
    ax.grid(axis="y", alpha=0.3)
    _add_fig(slide, fig, Inches(0.4), Inches(1.05), Inches(6.6), Inches(5.7))

    # Right: result table
    _txt(slide, "Balanced accuracy (mean%, std%)",
         Inches(7.2), Inches(1.1), Inches(5.9), Inches(0.35),
         bold=True, size_pt=12, color=C_ACCENT)

    col_x = [Inches(7.2), Inches(8.2), Inches(9.2), Inches(10.2), Inches(11.3)]
    col_w = [Inches(0.9)] * 5
    row_h = Inches(0.32)
    top0  = Inches(1.5)
    headers = ["npi", "clean", "λ=1.0", "λ=1.5", "λ=2.0"]
    tag_keys = ["le_clean", "ls1 (λ=1.0)", "ls15 (λ=1.5)", "ls2 (λ=2.0)"]
    for ci, h in enumerate(headers):
        _txt(slide, h, col_x[ci], top0, col_w[ci], row_h,
             bold=True, size_pt=10, color=C_HEAD)
    _rule(slide, top0 + row_h, width=Inches(5.0), left=Inches(7.2), thickness_pt=0.5)

    for ri, npi in enumerate(NPI_ORDER):
        top = top0 + row_h + ri * row_h
        _txt(slide, npi, col_x[0], top, col_w[0], row_h, size_pt=10)
        vals = []
        for tk in tag_keys:
            s = summaries[tk]
            row = s[s["npi"] == npi]
            vals.append(row.iloc[0]["mean_pct"] if len(row) else float("nan"))

        best = max(v for v in vals if not np.isnan(v))
        for ci, (val, tk) in enumerate(zip(vals, tag_keys)):
            if np.isnan(val):
                _txt(slide, "—", col_x[ci+1], top, col_w[ci+1], row_h, size_pt=10)
            else:
                is_best = abs(val - best) < 0.05
                is_clean = tk == "le_clean"
                col = C_GOOD if (is_best and is_clean) else (C_WARN if (is_best and not is_clean) else C_BODY)
                _txt(slide, f"{val:.1f}", col_x[ci+1], top, col_w[ci+1], row_h,
                     size_pt=10, bold=is_best, color=col)

    _txt(slide,
         "le_clean wins at every npi.\n"
         "Higher λ_supcon → worse (even at npi=all).\n"
         "Labeled-anchor boost degrades unlabeled structure.",
         Inches(7.2), Inches(5.0), Inches(5.9), Inches(1.0),
         size_pt=11, color=C_WARN)

    # sc vs uc loss annotation
    _txt(slide,
         "Loss trace (task 0, npi=10):\n"
         "  sc (labeled): 4.13 → 2.39  ✓ supervised signal works\n"
         "  uc (unlabeled): 4.71 → 5.51  ✗ self-supervised degraded",
         Inches(7.2), Inches(6.1), Inches(5.9), Inches(1.1),
         size_pt=10, italic=True, color=C_BODY)


def slide_conclusion(prs):
    slide = _blank(prs)
    _header(slide, "Conclusions",
            "Guaranteed labeled pairs + boosted loss weight — what worked and what didn't")

    findings = [
        ("LabeledAwareBatchSampler",
         "✓ Effective at forcing labeled patches into every batch (0.09 → 4 per batch)\n"
         "  Stabilises training: variance at npi=10 halved (±18.8% → ±8.9%)\n"
         "  No consistent mean accuracy improvement over 3 repeats"),
        ("lambda_supcon = 1.0–2.0",
         "✗ Accuracy drops at all npi values compared to le_clean baseline\n"
         "  sc loss decreases (good) but uc loss increases (bad)\n"
         "  The boosted term degrades unlabeled self-supervised geometry\n"
         "  GBC relies on the full latent space — not just the labeled cluster separation"),
        ("Root cause",
         "Even with guaranteed sampling, labeled patches are only 3.1% of the batch.\n"
         "Any boost to the supervised term must compete with 124 unlabeled anchors\n"
         "sharing the same partition function denominator.\n"
         "At small K the labeled class signal is too sparse to guide the whole space."),
        ("Take-away",
         "The unlabeled self-supervised contrast is the backbone of the embedding.\n"
         "Injecting more supervision via loss weighting is counterproductive here.\n"
         "The le_clean baseline (standard shuffle, no boost) remains the best strategy."),
    ]

    icons  = ["", "", "→", "★"]
    colors = [C_GOOD, C_WARN, C_ACCENT, C_HEAD]
    tops   = [Inches(1.1), Inches(2.55), Inches(4.0), Inches(5.45)]

    for (title, body), icon, color, top in zip(findings, icons, colors, tops):
        _txt(slide, f"{icon}  {title}", Inches(0.5), top, Inches(12.3), Inches(0.4),
             bold=True, size_pt=13, color=color)
        _txt(slide, body, Inches(0.8), top + Inches(0.4), Inches(12.0), Inches(1.0),
             size_pt=11, color=C_BODY)
        if top != tops[-1]:
            _rule(slide, top + Inches(1.25), thickness_pt=0.4)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="results/supcon_analysis.pptx")
    args = ap.parse_args()

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    print("Loading annotation data …")
    full_ann = pd.read_csv(FULL_ANN)

    print("Building slides …")
    prs = _prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution_sampler(prs)
    slide_solution_lambda(prs)
    slide_experiment(prs)
    slide_result_sampler(prs, full_ann)
    slide_result_sweep(prs, full_ann)
    slide_conclusion(prs)

    prs.save(str(out_path))
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
