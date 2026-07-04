#!/usr/bin/env python3
"""
Generate a PowerPoint on the ds1 vs ds1+ds3 generalizability experiment.

Usage:
  python scripts/make_generalizability_pptx.py [--out output.pptx]
"""
import argparse
from pathlib import Path
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ───────────────────────────────────────────────────────────────────
C_BG    = RGBColor(0xFF,0xFF,0xFF)
C_BLACK = RGBColor(0x11,0x11,0x11)
C_DARK  = RGBColor(0x1A,0x3A,0x6B)
C_BLUE  = RGBColor(0x1F,0x5C,0x99)
C_LBLUE = RGBColor(0xD6,0xE8,0xF8)
C_GREEN = RGBColor(0x1A,0x7A,0x4A)
C_RED   = RGBColor(0xB0,0x20,0x20)
C_GREY  = RGBColor(0x66,0x66,0x66)
C_LGREY = RGBColor(0xF2,0xF5,0xFA)

# variant colours matching the bar chart
V_COLORS = {
    "vinc only":             RGBColor(0x4C,0x72,0xB0),
    "vinc+ppax (unbalanced)":RGBColor(0xDD,0x84,0x52),
    "vinc+ppax (balanced)":  RGBColor(0x55,0xA8,0x68),
    "vinc only + histmatch": RGBColor(0x81,0x72,0xB2),
    "vinc+ppax bal+histmatch":RGBColor(0xC4,0x4E,0x52),
}

GEN = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/generalizability")

def _bg(sl):
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = C_BG

def _title(sl, text, top=Inches(0.2), fontsize=30, color=C_DARK):
    txb = sl.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.65))
    p = txb.text_frame.paragraphs[0]; run = p.add_run()
    run.text = text; run.font.size = Pt(fontsize)
    run.font.bold = True; run.font.color.rgb = color

def _sub(sl, text, top=Inches(0.85), fontsize=13, color=C_GREY):
    txb = sl.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.35))
    p = txb.text_frame.paragraphs[0]; run = p.add_run()
    run.text = text; run.font.size = Pt(fontsize)
    run.font.italic = True; run.font.color.rgb = color

def _body(sl, lines, left=Inches(0.5), top=Inches(1.25),
          width=Inches(9.0), height=Inches(5.8), fontsize=13, color=C_BLACK):
    txb = sl.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True; first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        if isinstance(line, tuple): text,sz,bold,col,bul = line
        else: text,sz,bold,col,bul = line,fontsize,False,color,True
        p.space_before = Pt(3); run = p.add_run()
        run.text = ("• " if bul else "") + text
        run.font.size = Pt(sz); run.font.bold = bold; run.font.color.rgb = col

def _box(sl, text, left, top, width, height,
         bg=C_BLUE, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=11, bold=False,
         align=PP_ALIGN.CENTER):
    shape = sl.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xBB,0xCC,0xDD); shape.line.width = Pt(0.75)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top  = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fontsize); run.font.color.rgb = fg; run.font.bold = bold
    return shape

def _img(sl, path, left, top, width=None, height=None):
    if not Path(path).exists(): return
    try:
        if width and height: sl.shapes.add_picture(str(path),left,top,width,height)
        elif width:          sl.shapes.add_picture(str(path),left,top,width=width)
        elif height:         sl.shapes.add_picture(str(path),left,top,height=height)
        else:                sl.shapes.add_picture(str(path),left,top)
    except Exception as e: print(f"  [warn] {path}: {e}")

def _hline(sl, left, top, width, color=C_GREY):
    line = sl.shapes.add_shape(1, left, top, width, Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.color.rgb = color


def make_pptx(out_path: Path):
    df = pd.read_csv(GEN / "cross_dataset_recon_metrics.csv", low_memory=False)

    VARIANTS = [
        ("baseline_vinc_only",                    "vinc only",              V_COLORS["vinc only"]),
        ("baseline_vinc_ppax",                    "vinc+ppax (unbalanced)", V_COLORS["vinc+ppax (unbalanced)"]),
        ("baseline_vinc_ppax_balanced",           "vinc+ppax (balanced)",   V_COLORS["vinc+ppax (balanced)"]),
        ("baseline_vinc_only_histmatch",          "vinc only + histmatch",  V_COLORS["vinc only + histmatch"]),
        ("baseline_vinc_ppax_balanced_histmatch", "vinc+ppax bal+histmatch",V_COLORS["vinc+ppax bal+histmatch"]),
    ]
    EXT = ["ppax_control","ppax_ycomp","pfak_control","pfak_ycomp","nih3t3_control","nih3t3_ycomp"]
    EXT_SHORT = ["ppax\nctrl","ppax\nycomp","pfak\nctrl","pfak\nycomp","nih3t3\nctrl","nih3t3\nycomp"]

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Generalisation: Training Data Diversity & Intensity Normalisation",
           top=Inches(1.9), fontsize=28)
    _sub(sl, "Can adding more datasets and histogram matching help the AE generalise to unseen cell lines?",
         top=Inches(2.7), fontsize=14, color=C_BLUE)
    _body(sl, [
        ("ds1 = vinc   (vinculin experiment, ~24 400 patches)", 13, False, C_GREY, False),
        ("ds2 = pfak   (paxillin, FAK cell line)",              13, False, C_GREY, False),
        ("ds3 = ppax   (paxillin, standard cell line)",         13, False, C_GREY, False),
        ("ds4 = nih3t3 (paxillin, NIH3T3 cell line)",           13, False, C_GREY, False),
    ], top=Inches(3.5), fontsize=13)
    _body(sl, [
        "Training: ds1 only  vs  ds1 + ds3  |  3 metrics: L1, MSE, Hessian L1",
        "Testing:  external datasets ds2, ds3, ds4 (never seen during training)",
    ], top=Inches(5.5), fontsize=12, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Experimental setup
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Experimental Setup")
    _sub(sl,   "5 variants tested across 3 strategies")

    strategies = [
        ("Strategy 1 — Training data",    C_BLUE,  [
            "Baseline: train on ds1 (vinc) only.",
            "Multi-dataset: train on ds1 + ds3 (ppax).",
        ]),
        ("Strategy 2 — Balanced sampling",C_GREEN, [
            "Problem: ds1 ~24k patches vs ds3 ~5.8k → 4:1 imbalance.",
            "Fix: repeat ds3 patch directories ×4 in config (~23k ppax).",
            "Result: model sees equal numbers from each dataset per epoch.",
        ]),
        ("Strategy 3 — Histogram matching",RGBColor(0x5A,0x1F,0x8A), [
            "Problem: ds1 pixel std ≈ 0.18 vs ds2/ds3/ds4 std ≈ 0.31–0.34.",
            "Same mean (~0.185) but very different contrast after CIO-RB.",
            "Fix: compute reference CDF from all 4 datasets pooled (~12M pixels).",
            "     Per-dataset forward map applied at patch load time.",
            "     Inverse map applied to reconstructed patches (original appearance).",
        ]),
    ]
    for i, (title, color, items) in enumerate(strategies):
        y = Inches(1.3 + i * 1.8)
        _box(sl, title, Inches(0.3), y, Inches(3.0), Inches(0.42),
             bg=color, fontsize=11, bold=True)
        _body(sl, items, left=Inches(3.5), top=y, width=Inches(6.3),
              height=Inches(1.4), fontsize=11, color=C_BLACK)
        _hline(sl, Inches(0.3), y + Inches(1.55), Inches(9.4))

    _body(sl, [
        ("5 variants", 13, True, C_DARK, False),
        "① vinc only       ② vinc+ppax (unbal.)  ③ vinc+ppax (balanced)",
        "④ vinc only + histmatch       ⑤ vinc+ppax balanced + histmatch",
    ], top=Inches(6.7), fontsize=11)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3a — Intensity histogram motivation (stats table + pooled plot)
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Why Histogram Matching? — Dataset Statistics")
    _sub(sl,   "CIO-RB normalisation aligns means but not variance")

    _body(sl, [
        ("Pixel intensity statistics after CIO-RB normalisation", 14, True, C_DARK, False),
    ], top=Inches(1.2), fontsize=13)

    stats = [
        ("Dataset", "Mean", "Std", "N images", "Note"),
        ("ds1 (vinc)",   "0.1885", "0.179", "50", "Training data — low contrast"),
        ("ds2 (pfak)",   "0.1822", "0.337", "10", "External — 1.9× wider std"),
        ("ds3 (ppax)",   "0.1865", "0.310", "11", "External/training — 1.7× wider std"),
        ("ds4 (nih3t3)", "0.1855", "0.342", "16", "External — 1.9× wider std"),
    ]
    col_w = [Inches(1.7), Inches(1.0), Inches(1.0), Inches(1.0), Inches(4.5)]
    row_h = Inches(0.45)
    x0, y0 = Inches(0.3), Inches(1.85)
    for ri, row in enumerate(stats):
        for ci, (cell, w) in enumerate(zip(row, col_w)):
            x = x0 + sum(col_w[:ci])
            y = y0 + ri * row_h
            if ri == 0:
                bg = C_DARK; fg = RGBColor(0xFF,0xFF,0xFF)
            elif ri == 1:
                bg = RGBColor(0xE8,0xF2,0xFF); fg = C_BLACK
            else:
                bg = C_LGREY if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF); fg = C_BLACK
            _box(sl, cell, x, y, w, row_h, bg=bg, fg=fg, fontsize=10)

    _img(sl, GEN.parent / "dataset_intensity_histograms.png",
         Inches(0.3), Inches(4.2), width=Inches(9.4))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3b — Per-image histograms original vs histmatch
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Per-Image Intensity Histograms — Original vs Histogram-Matched")
    _sub(sl,   "Thin lines = individual source images  |  Bold = pooled  |  Top: original  |  Bottom: after histmatch")

    _img(sl, GEN.parent / "per_image_intensity_histograms.png",
         Inches(0.2), Inches(1.15), width=Inches(9.6))

    _body(sl, [
        "vinc (50 imgs): low-contrast, narrow distribution — histmatch STRETCHES it to match the wider reference.",
        "ppax/pfak/nih3t3: already wider — histmatch COMPRESSES them toward the reference.",
        "After histmatch, all datasets converge to the same pooled distribution (bold lines overlap).",
        "Caveat: stretching vinc beyond [0,1] inflates reconstruction error — see next slides.",
    ], top=Inches(6.3), fontsize=10, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — L1 plots: vinc only vs vinc+ppax balanced
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "L1 Error Plots — Effect of Adding ppax to Training")
    _sub(sl,   "Left: vinc only   |   Right: vinc+ppax balanced")

    _img(sl, GEN / "baseline_vinc_only_cross_dataset_recon_l1.png",
         Inches(0.2), Inches(1.25), width=Inches(4.85))
    _img(sl, GEN / "baseline_vinc_ppax_balanced_cross_dataset_recon_l1.png",
         Inches(5.1), Inches(1.25), width=Inches(4.85))

    _body(sl, [
        "ppax/pfak/nih3t3 violins shift downward when ppax is added to training.",
        "vinc train/val bands remain similar — no catastrophic forgetting of the original domain.",
    ], top=Inches(6.6), fontsize=11, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — L1 plots: histmatch variants
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "L1 Error Plots — Effect of Histogram Matching")
    _sub(sl,   "Left: vinc only + histmatch   |   Right: vinc+ppax bal + histmatch")

    _img(sl, GEN / "baseline_vinc_only_histmatch_cross_dataset_recon_l1.png",
         Inches(0.2), Inches(1.25), width=Inches(4.85))
    _img(sl, GEN / "baseline_vinc_ppax_balanced_histmatch_cross_dataset_recon_l1.png",
         Inches(5.1), Inches(1.25), width=Inches(4.85))

    _body(sl, [
        "Histogram matching alone (left) gives similar improvement to adding ppax data.",
        "Combining both (right) further reduces pfak error but shows diminishing returns on nih3t3.",
    ], top=Inches(6.6), fontsize=11, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6a — L1 comparison (vinc + external)
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "L1 (MAE) — All Strategies (vinc training + external datasets)")
    _sub(sl,   "Left two bars = vinc train/val  |  dashed line separates training from unseen")
    _img(sl, GEN / "all_variants_comparison_recon_l1.png",
         Inches(0.2), Inches(1.2), width=Inches(9.6))
    _body(sl, [
        "Vinc train/val L1 increases slightly when ppax is added (model shares capacity across domains).",
        "All strategies substantially reduce error on external datasets vs vinc-only baseline.",
    ], top=Inches(6.7), fontsize=11, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6b — MSE comparison
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "MSE — All Strategies (vinc training + external datasets)")
    _sub(sl,   "Same pattern as L1 — histmatch and balanced sampling both help")
    _img(sl, GEN / "all_variants_comparison_recon_mse.png",
         Inches(0.2), Inches(1.2), width=Inches(9.6))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6c — Hessian L1 comparison
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Hessian L1 — All Strategies (vinc training + external datasets)")
    _sub(sl,   "Fine-texture error is essentially unchanged by any strategy")
    _img(sl, GEN / "all_variants_comparison_recon_hessian_l1.png",
         Inches(0.2), Inches(1.2), width=Inches(9.6))
    _body(sl, [
        "Hessian L1 measures curvature-level reconstruction error (missed fine structure).",
        "No strategy improves this — suggests a fundamental model capacity / domain limit.",
        "May require larger latent dim or multi-scale architecture to address.",
    ], top=Inches(6.4), fontsize=11, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Results table L1
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Results Table — L1 (MAE)")
    _sub(sl,   "Lower is better  |  Green = best per column  |  metrics in original intensity space")

    variant_keys = [v[0] for v in VARIANTS]

    def vinc_split_mean(key, split, metric="recon_l1"):
        vdf  = df[df["variant"] == key]
        mask = (vdf["group"].str.contains("vinc", na=False) &
                vdf["group"].str.contains(split,  na=False) &
                ~vdf["group"].str.contains("ppax", na=False))
        return round(vdf[mask][metric].mean(), 4) if mask.any() else float("nan")

    COL_GROUPS = (
        [("vinc\ntrain", None, "train"), ("vinc\nval",   None, "val")] +
        [(s, g, None) for s, g in zip(
            ["ppax\nctrl","ppax\nycomp","pfak\nctrl","pfak\nycomp","nih3t3\nctrl","nih3t3\nycomp"],
            EXT)]
    )  # (header, ext_group_or_None, split_or_None)

    n_cols  = len(COL_GROUPS)
    col_w   = [Inches(2.05)] + [Inches(0.95)] * n_cols
    row_h   = Inches(0.48)
    x0, y0  = Inches(0.2), Inches(1.15)

    # build data dict
    data = {}
    for key in variant_keys:
        vdf = df[df["variant"] == key]
        row_vals = []
        for hdr, ext_g, split in COL_GROUPS:
            if split:
                row_vals.append(vinc_split_mean(key, split, "recon_l1"))
            else:
                row_vals.append(round(vdf[vdf["group"]==ext_g]["recon_l1"].mean(), 4))
        data[key] = row_vals

    col_mins = [min(data[k][ci] for k in variant_keys) for ci in range(n_cols)]

    # header row — vinc cols in lighter blue, external cols in dark
    hdrs = [c[0] for c in COL_GROUPS]
    for ci, (hdr, w) in enumerate(zip(["Variant"]+hdrs, col_w)):
        x   = x0 + sum(col_w[:ci])
        is_vinc = ci in (1, 2)
        bg  = C_BLUE if is_vinc else C_DARK
        _box(sl, hdr, x, y0, w, row_h, bg=bg,
             fg=RGBColor(0xFF,0xFF,0xFF), fontsize=8, bold=True)

    # separator line between vinc and external cols
    sep_x = x0 + col_w[0] + col_w[1] + col_w[2]
    sep_line = sl.shapes.add_shape(1, sep_x, y0, Pt(2), row_h*(len(VARIANTS)+1))
    sep_line.fill.solid(); sep_line.fill.fore_color.rgb = RGBColor(0x88,0xAA,0xCC)
    sep_line.line.color.rgb = RGBColor(0x88,0xAA,0xCC)

    for ri, (key, label, color) in enumerate(VARIANTS):
        y      = y0 + (ri+1) * row_h
        bg_row = C_LGREY if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        _box(sl, label, x0, y, col_w[0], row_h,
             bg=bg_row, fg=color, fontsize=8, bold=True, align=PP_ALIGN.LEFT)
        for ci, val in enumerate(data[key]):
            x    = x0 + sum(col_w[:ci+1])
            best = abs(val - col_mins[ci]) < 0.0001
            bg   = RGBColor(0xC8,0xE6,0xC9) if best else bg_row
            fg   = RGBColor(0x0A,0x50,0x20) if best else C_BLACK
            txt  = f"{val:.4f}" if not np.isnan(val) else "—"
            _box(sl, txt, x, y, col_w[ci+1], row_h, bg=bg, fg=fg, fontsize=9, bold=best)

    _body(sl, [
        "Vertical separator divides vinc (training) from external (unseen) datasets.",
        "Green = best per column.  histmatch variants achieve lowest error on pfak across all strategies.",
    ], top=Inches(6.7), fontsize=11, color=C_DARK)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Results table MSE + Hessian
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Results Table — MSE and Hessian L1")
    _sub(sl,   "Hessian L1 barely changes; MSE follows same pattern as L1")

    for panel_i, (metric, mlabel) in enumerate([("recon_mse","MSE"),
                                                ("recon_hessian_l1","Hessian L1")]):
        y_offset = Inches(1.15 + panel_i * 2.85)
        _box(sl, mlabel, Inches(0.2), y_offset, Inches(9.6), Inches(0.38),
             bg=C_BLUE, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=11, bold=True)

        col_w2 = [Inches(2.05)] + [Inches(0.93)] * n_cols
        row_h2 = Inches(0.37)

        # build data
        data2 = {}
        for key in variant_keys:
            vdf = df[df["variant"] == key]
            row_vals = []
            for hdr, ext_g, split in COL_GROUPS:
                if split:
                    row_vals.append(vinc_split_mean(key, split, metric))
                else:
                    v = vdf[vdf["group"]==ext_g][metric].mean()
                    row_vals.append(round(v, 4) if not np.isnan(v) else float("nan"))
            data2[key] = row_vals
        col_mins2 = [min(data2[k][ci] for k in variant_keys) for ci in range(n_cols)]

        hdrs2 = [c[0] for c in COL_GROUPS]
        for ci, (hdr, w) in enumerate(zip(["Variant"]+hdrs2, col_w2)):
            x   = Inches(0.2) + sum(col_w2[:ci])
            is_vinc = ci in (1, 2)
            bg  = C_BLUE if is_vinc else C_DARK
            _box(sl, hdr, x, y_offset+Inches(0.4), w, row_h2,
                 bg=bg, fg=RGBColor(0xFF,0xFF,0xFF), fontsize=7, bold=True)

        for ri, (key, vlabel, color) in enumerate(VARIANTS):
            y      = y_offset + Inches(0.4) + (ri+1)*row_h2
            bg_row = C_LGREY if ri%2==0 else RGBColor(0xFF,0xFF,0xFF)
            _box(sl, vlabel, Inches(0.2), y, col_w2[0], row_h2,
                 bg=bg_row, fg=color, fontsize=7, bold=True, align=PP_ALIGN.LEFT)
            for ci, val in enumerate(data2[key]):
                x    = Inches(0.2) + sum(col_w2[:ci+1])
                best = abs(val - col_mins2[ci]) < 0.0001
                bg   = RGBColor(0xC8,0xE6,0xC9) if best else bg_row
                fg   = RGBColor(0x0A,0x50,0x20) if best else C_BLACK
                txt  = f"{val:.4f}" if not np.isnan(val) else "—"
                _box(sl, txt, x, y, col_w2[ci+1], row_h2, bg=bg, fg=fg,
                     fontsize=8, bold=best)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Interpretation
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Interpretation and Key Findings")

    _body(sl, [
        ("1.  Histogram matching is highly effective", 14, True, C_DARK, False),
        "vinc only + histmatch beats vinc+ppax (unbalanced) on all 6 external columns.",
        "This suggests intensity distribution mismatch is a primary cause of generalisation error.",
        "Remarkably, no extra data is needed — just preprocessing alignment.",
        "",
        ("2.  Balanced ppax sampling matters", 14, True, C_DARK, False),
        "Unbalanced ds1+ds3 training gives only modest improvement over ds1 alone.",
        "Balanced (4× ppax repeat) consistently reduces error, especially on ppax.",
        "The model cannot learn ppax morphology well when it sees it 4× less per epoch.",
        "",
        ("3.  Hessian L1 is unchanged", 14, True, C_DARK, False),
        "Fine-texture reconstruction error is not improved by any strategy tested.",
        "Likely reflects a fundamental limit: the model capacity / latent dim is insufficient",
        "to capture high-frequency structural details across domains.",
        "",
        ("4.  Best strategy: balanced + histmatch", 14, True, C_GREEN, False),
        "Combining both gives the lowest L1 and MSE on 4/6 external columns.",
        "Open question: does this also improve downstream FA-type classification?",
    ], fontsize=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Next steps
    # ══════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank); _bg(sl)
    _title(sl, "Next Steps")

    _body(sl, [
        ("Immediate", 14, True, C_DARK, False),
        "Apply histmatch preprocessing to the full semisup / contrastive AE training pipeline.",
        "Evaluate downstream FA classification with histmatch-trained models.",
        "",
        ("Histogram matching limitations to address", 14, True, C_DARK, False),
        "Current reference: pooled from 4 datasets — may be biased toward paxillin contrast.",
        "Deployment concern: histogram map must be precomputed for each new dataset.",
        "Alternative: per-image percentile clipping (no dataset-level reference required).",
        "",
        ("Broader generalisation", 14, True, C_DARK, False),
        "Add pfak or nih3t3 to training (currently only vinc + ppax tested).",
        "Cross-dataset FA classification: train on vinc+ppax, test on pfak/nih3t3.",
        "",
        ("Future: learned normalisation", 14, True, C_DARK, False),
        "Instance normalisation layer inside the encoder — no preprocessing needed.",
        "Style transfer approach: domain-invariant features via adversarial training.",
    ], fontsize=12)

    prs.save(str(out_path))
    print(f"Saved: {out_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path,
                        default=Path("/net/projects/CLS/lding/data/fa_data_analysis"
                                     "/ae_results/generalizability_overview.pptx"))
    args = parser.parse_args()
    make_pptx(args.out)

if __name__ == "__main__":
    main()
