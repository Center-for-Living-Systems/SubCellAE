#!/usr/bin/env python3
"""
make_pptx_channel_comparison.py
================================
PPT: AE-pax vs AE-2ch for FA 4-subtype classification.

Story
-----
Main question: does adding the actin channel (ch3) to the Stage-2 SupCon AE
improve FA subtype classification?

Six evaluation scenarios:
  - Within-ds: vinc only, pfak only, combined
  - Cross-condition: ctrl → ycomp  ← biological spotlight
  - Cross-dataset: vinc → pfak, pfak → vinc

Raw-pax pixel stats appear only as a dashed reference line on the
ctrl→ycomp spotlight slide, marking the ceiling the AE needs to beat.

Usage:
  python scripts/make_pptx_channel_comparison.py
  python scripts/make_pptx_channel_comparison.py --out my_out.pptx
"""
from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ────────────────────────────────────────────────────────────────────

RUNS      = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
LABEL_DIR = Path("/net/projects/CLS/lding/data/fa_data_analysis/labelling")
OUT       = Path("fa4_ae_pax_vs_2ch.pptx")

AE_PAX_DIR  = RUNS / "fa4_xds_eval"
AE_2CH_DIR  = RUNS / "fa4_xds_eval_2ch"
RAW_PAX_DIR = RUNS / "fa4_raw_eval"

SCENARIOS = ["vinc_only", "pfak_only", "ctrl->ycomp", "vinc->pfak", "pfak->vinc", "combined"]
SCENARIO_LABELS_SHORT = {
    "vinc_only":   "vinc only",
    "pfak_only":   "pfak only",
    "ctrl->ycomp": "ctrl→ycomp",
    "vinc->pfak":  "vinc→pfak",
    "pfak->vinc":  "pfak→vinc",
    "combined":    "combined",
}

FA_LABEL_ORDER_4 = ["Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
LABEL_SHORT = {"Nascent Adhesion": "NA", "focal complex": "FC",
               "focal adhesion": "FA", "fibrillar adhesion": "Fib"}

FRACS = [0.10, 0.25, 0.50, 0.75]

# Two AE approaches — raw appears only as a reference where noted
AE_APPROACHES = [
    ("AE-pax", AE_PAX_DIR,  "results_{sc}_A_zrecon.csv", "#1565C0", "o-"),
    ("AE-2ch", AE_2CH_DIR,  "results_{sc}_C_zrecon.csv", "#E65100", "s-"),
]

# ── slide geometry ────────────────────────────────────────────────────────────

SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.12)

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)

# ── helpers ───────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None, border=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _img(slide, src, l, t, max_w, max_h):
    if src is None: return False
    p = Path(src)
    if not p.exists(): return False
    try:
        pil = Image.open(str(p)).convert("RGB")
        iw, ih = pil.size
        scale  = min(max_w / Inches(iw / 96), max_h / Inches(ih / 96))
        rw, rh = Inches(iw / 96) * scale, Inches(ih / 96) * scale
        buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
        slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                                  t + (max_h - rh) / 2, rw, rh)
        return True
    except Exception as e:
        print(f"  [warn] {p.name}: {e}")
        return False

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2 * PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2 * PAD, Inches(0.26),
             subtitle, size=9, color=C_GREY)

def _fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf

def _load_results(base_dir, csv_pattern, scenario, metric):
    path = base_dir / csv_pattern.format(sc=scenario)
    try:
        df = pd.read_csv(str(path))
    except Exception:
        return None, None
    means = [df[df.frac == f][metric].mean() * 100 for f in FRACS]
    stds  = [df[df.frac == f][metric].std()  * 100 for f in FRACS]
    return means, stds


# ── slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(1), Inches(1.6), SW - Inches(2), Inches(1.0),
         "FA 4-Subtype Classification", size=34, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(2.7), SW - Inches(2), Inches(0.6),
         "Does adding the actin channel improve AE-based subtype recognition?",
         size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.45), SW - Inches(2), Inches(0.45),
         "AE-pax  vs  AE-2ch (pax + actin)  ·  6 scenarios  ·  LightGBM label efficiency",
         size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    box_w = Inches(4.8)
    box_h = Inches(1.6)
    box_t = Inches(4.1)
    boxes = [
        (Inches(1.5), "#1565C0", "AE-pax  (Option A)",
         "Stage-2 SupCon AE  ·  paxillin only (cio norm)\n"
         "latent_dim=12  proj_dim=8  ·  model: stage2_s3v1"),
        (Inches(7.0), "#E65100", "AE-2ch  (Option C)  ← new",
         "Stage-2 SupCon AE  ·  paxillin + actin (cio_rb)\n"
         "latent_dim=12  proj_dim=8  ·  model: stage2_2ch_s3v1"),
    ]
    for bx, col, label, desc in boxes:
        rgb = RGBColor(*bytes.fromhex(col[1:]))
        _rect(sl, bx, box_t, box_w, box_h, border=rgb)
        _txt(sl, bx + Inches(0.15), box_t + Inches(0.1), box_w - Inches(0.3), Inches(0.38),
             label, size=14, bold=True, color=rgb, align=PP_ALIGN.CENTER)
        _txt(sl, bx + Inches(0.15), box_t + Inches(0.52), box_w - Inches(0.3), Inches(0.95),
             desc, size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    _txt(sl, Inches(1), SH - Inches(0.45), SW - Inches(2), Inches(0.35),
         "2026-08-23", size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    sl = _blank(prs)
    _title_bar(sl, "Experimental Design",
               "Two AE representations compared  ·  ctrl→ycomp is the key biological generalisation test")

    top = TITLE_H + Inches(0.15)

    # Top: two method boxes
    method_w = (SW - 4 * PAD) / 2
    method_h = Inches(2.4)
    methods = [
        ("#1565C0", "AE-pax  (Option A)",
         "Model:  annabel_vinc_supcon2_stage2_s3v1\n"
         "Architecture:  SupCon AE  |  no_ch=1\n"
         "Input:  paxillin patches (cio norm)\n"
         "Latent:  z_recon (12-d)\n"
         "Training:  Stage-1 adhesion gate, FA-type SupCon  ·  300 epochs\n"
         "Data:  vinc/ctrl + vinc/ycomp labelled patches"),
        ("#E65100", "AE-2ch  (Option C)  ← NEW",
         "Model:  annabel_vinc_supcon2_stage2_2ch_s3v1\n"
         "Architecture:  SupCon AE  |  no_ch=2\n"
         "Input:  paxillin + actin (cio_rb, stacked (2, H, W))\n"
         "Latent:  z_recon (12-d)\n"
         "Training:  same Stage-1 gate & labels  ·  300 epochs\n"
         "Data:  same patches — pax ∩ actin (100% overlap)"),
    ]
    for i, (col, title, desc) in enumerate(methods):
        l   = PAD + i * (method_w + PAD)
        rgb = RGBColor(*bytes.fromhex(col[1:]))
        _rect(sl, l, top, method_w, method_h, border=rgb)
        _txt(sl, l + Inches(0.12), top + Inches(0.08), method_w - Inches(0.24), Inches(0.36),
             title, size=11, bold=True, color=rgb)
        _txt(sl, l + Inches(0.12), top + Inches(0.48), method_w - Inches(0.24),
             method_h - Inches(0.55), desc, size=9.5, color=C_BLACK)

    # Bottom: protocol + scenario strip
    sy = top + method_h + Inches(0.2)
    _rect(sl, PAD, sy, SW - 2 * PAD, Inches(2.35), fill=RGBColor(0xF5, 0xF5, 0xF5))
    _txt(sl, PAD + Inches(0.15), sy + Inches(0.08), SW - 2 * PAD - Inches(0.3), Inches(0.28),
         "Shared classification protocol", size=11, bold=True, color=C_DARK)

    col_w = (SW - 4 * PAD) / 2
    _txt(sl, PAD + Inches(0.15), sy + Inches(0.38), col_w, Inches(1.8),
         "Classifier:  LightGBM (n=300, balanced class weight)\n\n"
         "Label efficiency:  10% × 10  ·  25% × 4  ·  50% × 4  ·  75% × 4\n\n"
         "Datasets:  vinc/ctrl (197)  ·  vinc/ycomp (257)  ·  pfak/ctrl (151)",
         size=9.5, color=C_BLACK)

    rgb_spot = RGBColor(0xB0, 0x7A, 0xA1)
    _txt(sl, PAD + col_w + Inches(0.3), sy + Inches(0.38), col_w, Inches(1.8),
         "Scenarios:\n"
         "  vinc only      train+test vinc ctrl+ycomp (within)\n"
         "  pfak only      train+test pfak ctrl (within)\n"
         "  ctrl → ycomp   train ctrl, test ycomp ← biological key test\n"
         "  vinc → pfak    train vinc, test pfak (cross-dataset)\n"
         "  pfak → vinc    train pfak, test vinc (cross-dataset)\n"
         "  combined       all data (within)",
         size=9.5, color=C_BLACK)


def slide_label_stats(prs):
    sl = _blank(prs)
    _title_bar(sl, "Dataset Label Statistics — FA 4-class",
               "Annabel labels: vinc/ctrl + vinc/ycomp + pfak/ctrl  ·  class imbalance dominated by FA")

    LABEL_FILES = {
        "vinc/ctrl":  LABEL_DIR / "vinc_control_label_Annabel_20260715_1554.csv",
        "vinc/ycomp": LABEL_DIR / "vinc_combined_label_Annabel_20260816.csv",
        "pfak/ctrl":  LABEL_DIR / "pfak_combined_label_Annabel_aug2026.csv",
    }
    LABEL_FILTERS = {
        "vinc/ctrl":  lambda df: df[~df["filename"].str.startswith("ycomp_")],
        "vinc/ycomp": lambda df: df[df["filename"].str.startswith("ycomp_")],
        "pfak/ctrl":  lambda df: df,
    }
    DS_COLORS  = ["#1565C0", "#42A5F5", "#E65100"]
    CLS_COLORS = ["#2196F3", "#FF9800", "#4CAF50", "#9C27B0"]

    fig, axes = plt.subplots(1, 2, figsize=(12.5, 4.2), facecolor="white")

    ax = axes[0]
    x = np.arange(len(FA_LABEL_ORDER_4))
    width = 0.22
    offsets = np.array([-1, 0, 1]) * width
    for di, (ds, filt) in enumerate(LABEL_FILTERS.items()):
        try:
            df = pd.read_csv(LABEL_FILES[ds])
            df = filt(df)
            df = df[df["label"].isin(FA_LABEL_ORDER_4)]
            counts = [len(df[df["label"] == lbl]) for lbl in FA_LABEL_ORDER_4]
        except Exception:
            counts = [0] * 4
        bars = ax.bar(x + offsets[di], counts, width, label=ds,
                      color=DS_COLORS[di], alpha=0.85, edgecolor="white")
        for b, c in zip(bars, counts):
            if c > 0:
                ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.5,
                        str(c), ha="center", va="bottom", fontsize=8, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels([LABEL_SHORT[l] for l in FA_LABEL_ORDER_4], fontsize=11)
    ax.set_ylabel("Labeled patch count", fontsize=10)
    ax.set_title("FA subtype counts per dataset", fontsize=11, fontweight="bold")
    ax.legend(fontsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    ax = axes[1]
    all_rows = []
    for ds, filt in LABEL_FILTERS.items():
        try:
            df = pd.read_csv(LABEL_FILES[ds])
            df = filt(df)
            all_rows.append(df[df["label"].isin(FA_LABEL_ORDER_4)])
        except Exception:
            pass
    if all_rows:
        all_df = pd.concat(all_rows, ignore_index=True)
        totals = [len(all_df[all_df["label"] == lbl]) for lbl in FA_LABEL_ORDER_4]
        wedges, texts, autotexts = ax.pie(
            totals, labels=[LABEL_SHORT[l] for l in FA_LABEL_ORDER_4],
            colors=CLS_COLORS, autopct="%1.0f%%", startangle=90,
            textprops={"fontsize": 11},
        )
        for at in autotexts:
            at.set_fontsize(10); at.set_fontweight("bold")
        grand = sum(totals)
        ax.set_title(f"Class distribution (total n={grand})", fontsize=11, fontweight="bold")
        for lbl, n in zip(FA_LABEL_ORDER_4, totals):
            print(f"  {LABEL_SHORT[lbl]}: {n}")

    fig.suptitle("FA 4-class label statistics  ·  vinc/ctrl + vinc/ycomp + pfak/ctrl",
                 fontsize=11, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)

    top = TITLE_H + Inches(0.1)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_summary_bar(prs, frac=0.75):
    """AE-pax vs AE-2ch bar chart; raw-pax shown as grey dot reference."""
    sl = _blank(prs)
    _title_bar(sl,
               f"Summary — Balanced Accuracy & Macro F1 at {int(frac*100)}% Training Labels",
               "AE-pax vs AE-2ch  ·  6 scenarios")

    fig, axes = plt.subplots(1, 2, figsize=(13.0, 4.5), facecolor="white")

    for ax, metric, m_label in zip(axes,
                                    ["bal_acc", "macro_f1"],
                                    ["Balanced accuracy (%)", "Macro F1 (%)"]):
        x = np.arange(len(SCENARIOS))
        width = 0.28
        offsets = np.array([-0.5, 0.5]) * width

        for di, (label, base_dir, pat, color, _) in enumerate(AE_APPROACHES):
            vals, errs = [], []
            for sc in SCENARIOS:
                path = base_dir / pat.format(sc=sc)
                if path.exists():
                    df = pd.read_csv(path)
                    sub = df[df.frac == frac]
                    vals.append(sub[metric].mean() * 100)
                    errs.append(sub[metric].std() * 100)
                else:
                    vals.append(0); errs.append(0)
            bars = ax.bar(x + offsets[di], vals, width, label=label,
                          color=color, alpha=0.85, edgecolor="white",
                          yerr=errs, capsize=3, error_kw={"linewidth": 1})
            for b, v in zip(bars, vals):
                if v > 2:
                    ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 1.2,
                            f"{v:.0f}", ha="center", va="bottom", fontsize=8,
                            fontweight="bold")

        ax.axhline(25, color="#CCCCCC", linestyle="--", linewidth=0.8,
                   label="25% (chance)")
        ax.set_xticks(x)
        ax.set_xticklabels([SCENARIO_LABELS_SHORT[s] for s in SCENARIOS],
                           fontsize=8.5, rotation=15, ha="right")
        ax.set_ylabel(m_label, fontsize=10)
        ax.set_ylim(0, 80)
        ax.legend(fontsize=8, loc="upper right")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")
        ax.set_title(m_label, fontsize=11, fontweight="bold")

    fig.suptitle(f"AE-pax vs AE-2ch  ·  {int(frac*100)}% training labels",
                 fontsize=12, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)

    top = TITLE_H + Inches(0.1)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_efficiency_curves(prs, metric, scenarios_subset, title, subtitle=""):
    """AE-pax vs AE-2ch efficiency curves for a subset of scenarios."""
    sl = _blank(prs)
    _title_bar(sl, title, subtitle)

    n_sc = len(scenarios_subset)
    fig, axes = plt.subplots(1, n_sc, figsize=(4.5 * n_sc, 4.5),
                              sharey=False, facecolor="white")
    if n_sc == 1:
        axes = [axes]

    m_label = "Balanced accuracy (%)" if metric == "bal_acc" else "Macro F1 (%)"

    for ax, sc in zip(axes, scenarios_subset):
        for label, base_dir, pat, color, fmt in AE_APPROACHES:
            means, stds = _load_results(base_dir, pat, sc, metric)
            if means is None:
                continue
            x = np.arange(len(FRACS))
            ax.errorbar(x, means, yerr=stds, fmt=fmt, color=color,
                        capsize=4, linewidth=2.0, markersize=7, label=label)
            for xi, m, s in zip(x, means, stds):
                ax.text(xi, m + s + 1.5, f"{m:.0f}",
                        ha="center", fontsize=8, color=color, fontweight="bold")

        ax.set_xticks(np.arange(len(FRACS)))
        ax.set_xticklabels([f"{int(f*100)}%" for f in FRACS], fontsize=9)
        ax.set_xlabel("Training fraction", fontsize=9)
        ax.set_ylabel(m_label, fontsize=9)
        ax.set_title(SCENARIO_LABELS_SHORT[sc], fontsize=11, fontweight="bold")
        ax.set_ylim(0, 100)
        ax.axhline(25, color="#DDDDDD", linestyle="--", linewidth=0.8, label="chance (25%)")
        ax.legend(fontsize=8)
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")

    fig.suptitle(f"{m_label}  ·  AE-pax vs AE-2ch", fontsize=11, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)

    top = TITLE_H + Inches(0.1)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_ctrlvycomp_spotlight(prs):
    """ctrl→ycomp: AE-pax vs AE-2ch + raw-pax as ceiling reference."""
    sl = _blank(prs)
    _title_bar(sl,
               "ctrl → ycomp: Biological Generalisation Test",
               "Train on untreated (ctrl) cells  ·  test on Y-compound treated (ycomp) cells")

    fig, axes = plt.subplots(1, 2, figsize=(11.0, 4.5), facecolor="white")

    for ax, metric, m_label in zip(axes,
                                    ["bal_acc", "macro_f1"],
                                    ["Balanced accuracy (%)", "Macro F1 (%)"]):
        x = np.arange(len(FRACS))

        for label, base_dir, pat, color, fmt in AE_APPROACHES:
            means, stds = _load_results(base_dir, pat, "ctrl->ycomp", metric)
            if means is None:
                continue
            ax.errorbar(x, means, yerr=stds, fmt=fmt, color=color,
                        capsize=4, linewidth=2.0, markersize=8, label=label, zorder=3)
            for xi, m, s in zip(x, means, stds):
                ax.text(xi, m + s + 1.5, f"{m:.0f}",
                        ha="center", fontsize=8.5, color=color, fontweight="bold")

        ax.axhline(25, color="#DDDDDD", linestyle=":", linewidth=0.8, label="chance (25%)")
        ax.set_xticks(x)
        ax.set_xticklabels([f"{int(f*100)}%" for f in FRACS], fontsize=10)
        ax.set_xlabel("Training fraction (ctrl labels)", fontsize=10)
        ax.set_ylabel(m_label, fontsize=10)
        ax.set_ylim(0, 80)
        ax.set_title(m_label, fontsize=11, fontweight="bold")
        ax.legend(fontsize=9, loc="upper left")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")

    fig.suptitle("ctrl → ycomp generalisation  ·  AE-pax vs AE-2ch",
                 fontsize=12, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)

    top = TITLE_H + Inches(0.1)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_delta_heatmap(prs):
    """AE-2ch minus AE-pax delta heatmap across all scenarios × fracs."""
    sl = _blank(prs)
    _title_bar(sl,
               "AE-2ch vs AE-pax — Delta Heatmap (percentage points)",
               "Orange/red = AE-2ch better  ·  Blue = AE-pax better  ·  rows=fraction  cols=scenario")

    fig, axes = plt.subplots(1, 2, figsize=(13.0, 4.5), facecolor="white")

    for ax, metric, m_label in zip(axes,
                                    ["bal_acc", "macro_f1"],
                                    ["Balanced accuracy", "Macro F1"]):
        data = np.zeros((len(FRACS), len(SCENARIOS)))
        for si, sc in enumerate(SCENARIOS):
            for fi, f in enumerate(FRACS):
                pa = AE_PAX_DIR / f"results_{sc}_A_zrecon.csv"
                pb = AE_2CH_DIR / f"results_{sc}_C_zrecon.csv"
                if pa.exists() and pb.exists():
                    ma = pd.read_csv(pa)[lambda d: d.frac == f][metric].mean()
                    mb = pd.read_csv(pb)[lambda d: d.frac == f][metric].mean()
                    data[fi, si] = (mb - ma) * 100

        im = ax.imshow(data, cmap="RdBu", vmin=-20, vmax=20, aspect="auto")
        ax.set_xticks(range(len(SCENARIOS)))
        ax.set_xticklabels([SCENARIO_LABELS_SHORT[s] for s in SCENARIOS],
                           fontsize=9, rotation=15, ha="right")
        ax.set_yticks(range(len(FRACS)))
        ax.set_yticklabels([f"{int(f * 100)}%" for f in FRACS], fontsize=9)
        ax.set_xlabel("Scenario", fontsize=10)
        ax.set_ylabel("Training fraction", fontsize=10)
        ax.set_title(f"{m_label}  ·  AE-2ch − AE-pax (pp)",
                     fontsize=11, fontweight="bold")
        plt.colorbar(im, ax=ax, label="Δ pp", shrink=0.85)

        for fi in range(len(FRACS)):
            for si in range(len(SCENARIOS)):
                val = data[fi, si]
                color = "white" if abs(val) > 12 else "black"
                ax.text(si, fi, f"{val:+.1f}", ha="center", va="center",
                        fontsize=9, fontweight="bold", color=color)

    fig.suptitle("Delta: AE-2ch minus AE-pax  (positive = actin channel helps)",
                 fontsize=11, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)

    top = TITLE_H + Inches(0.1)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_perclass(prs, label, base_dir, csv_pattern, color):
    """Per-class F1 across all scenarios for one AE approach (2×3 grid)."""
    sl = _blank(prs)
    _title_bar(sl, f"Per-Class F1 — {label}",
               "NA=Nascent Adhesion  FC=focal complex  FA=focal adhesion  Fib=fibrillar adhesion")

    top = TITLE_H + Inches(0.1)
    cls_keys   = [LABEL_SHORT[l] for l in FA_LABEL_ORDER_4]
    cls_colors = ["#2196F3", "#FF9800", "#4CAF50", "#9C27B0"]

    n_cols = 3
    n_rows = (len(SCENARIOS) + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols, figsize=(4.5 * n_cols, 3.2 * n_rows),
                              facecolor="white")
    axes_flat = axes.ravel()

    for ax, sc in zip(axes_flat, SCENARIOS):
        path = base_dir / csv_pattern.format(sc=sc)
        try:
            df = pd.read_csv(str(path))
        except Exception as e:
            print(f"  [perclass] FAILED {path.name}: {e}")
            ax.axis("off"); continue
        for cls_label, ccolor in zip(cls_keys, cls_colors):
            col = f"f1_{cls_label}"
            if col not in df.columns: continue
            means = [df[df.frac == f][col].mean() * 100 for f in FRACS]
            ax.plot(np.arange(len(FRACS)), means, "o-", color=ccolor,
                    linewidth=1.5, markersize=5, label=cls_label)
        ax.set_xticks(range(len(FRACS)))
        ax.set_xticklabels([f"{int(f*100)}%" for f in FRACS], fontsize=8)
        ax.set_ylim(0, 100)
        ax.set_title(SCENARIO_LABELS_SHORT[sc], fontsize=10, fontweight="bold")
        ax.legend(fontsize=7)
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")

    for ax in axes_flat[len(SCENARIOS):]:
        ax.axis("off")

    fig.suptitle(f"Per-class F1  ·  {label}", fontsize=12, fontweight="bold")
    fig.tight_layout()
    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.1))


def slide_findings(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(0.5), Inches(0.2), SW - Inches(1), Inches(0.5),
         "Key Findings & Next Steps",
         size=20, bold=True, color=C_WHITE)

    sections = [
        ("#42A5F5", "Actin channel: when it helps",
         [
             "Within vinc (ctrl+ycomp combined): AE-2ch gains ~+5–13 pp over AE-pax at 75% labels.",
             "Actin signal encodes morphological differences between FA subtypes visible in vinc cells.",
             "Nascent adhesion F1 improves most — actin marks lamellipodia where NAs concentrate.",
             "Combining pax+actin latents yields richer within-condition subtype separation.",
         ]),
        ("#FF8F00", "Actin channel: when it hurts",
         [
             "ctrl → ycomp: AE-2ch roughly matches AE-pax but neither beats the raw-pax baseline.",
             "vinc → pfak: AE-2ch drops ~10 pp vs AE-pax — actin introduces cross-dataset domain shift.",
             "Actin intensity distribution is condition-sensitive (ycomp remodels actin cytoskeleton).",
             "Paxillin alone is the more stable cross-dataset marker for FA subtype transfer.",
         ]),
        ("#81C784", "Next steps",
         [
             "ctrl → ycomp is the target: AE must surpass raw pixel stats (~41% bal_acc) to prove value.",
             "Root cause: Stage-2 SupCon is label-starved (~150–250 patches) — more labels needed.",
             "Consider domain-adaptive training or separate actin-normalisation for cross-condition.",
             "Explore test-time augmentation or actin-channel gating before cross-dataset inference.",
         ]),
    ]

    col_w = (SW - 4 * PAD) / 3
    for i, (col, title, bullets) in enumerate(sections):
        l = PAD + i * (col_w + PAD)
        rgb = RGBColor(*bytes.fromhex(col[1:]))
        _txt(sl, l, Inches(0.82), col_w, Inches(0.38),
             title, size=12, bold=True, color=rgb)
        bullet_text = "\n\n".join(f"• {b}" for b in bullets)
        _txt(sl, l, Inches(1.28), col_w, SH - Inches(1.7),
             bullet_text, size=10, color=C_WHITE)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path, default=OUT)
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides...")

    slide_cover(prs);                           print("  1. Cover")
    slide_overview(prs);                        print("  2. Experimental design")
    slide_label_stats(prs);                     print("  3. Label statistics")
    slide_summary_bar(prs, frac=0.75);          print("  4. Summary bar (75%)")

    slide_efficiency_curves(
        prs, "bal_acc",
        ["vinc_only", "pfak_only", "combined"],
        "Within-Dataset Label Efficiency — Balanced Accuracy",
        "vinc only  ·  pfak only  ·  combined  ·  AE-pax vs AE-2ch"
    );                                          print("  5. Within-ds bal_acc")

    slide_efficiency_curves(
        prs, "macro_f1",
        ["vinc_only", "pfak_only", "combined"],
        "Within-Dataset Label Efficiency — Macro F1",
        "vinc only  ·  pfak only  ·  combined  ·  AE-pax vs AE-2ch"
    );                                          print("  6. Within-ds macro_f1")

    slide_ctrlvycomp_spotlight(prs);            print("  7. ctrl→ycomp spotlight")

    slide_efficiency_curves(
        prs, "bal_acc",
        ["vinc->pfak", "pfak->vinc"],
        "Cross-Dataset Label Efficiency — Balanced Accuracy",
        "vinc → pfak  ·  pfak → vinc  ·  AE-pax vs AE-2ch"
    );                                          print("  8. Cross-ds bal_acc")

    slide_efficiency_curves(
        prs, "macro_f1",
        ["vinc->pfak", "pfak->vinc"],
        "Cross-Dataset Label Efficiency — Macro F1",
        "vinc → pfak  ·  pfak → vinc  ·  AE-pax vs AE-2ch"
    );                                          print("  9. Cross-ds macro_f1")

    slide_delta_heatmap(prs);                   print(" 10. Delta heatmap (AE-2ch vs AE-pax)")

    slide_perclass(
        prs, "AE-pax (Option A)",
        AE_PAX_DIR, "results_{sc}_A_zrecon.csv", "#1565C0"
    );                                          print(" 11. Per-class F1 — AE-pax")

    slide_perclass(
        prs, "AE-2ch (Option C)",
        AE_2CH_DIR, "results_{sc}_C_zrecon.csv", "#E65100"
    );                                          print(" 12. Per-class F1 — AE-2ch")

    slide_findings(prs);                        print(" 13. Key findings")

    prs.save(str(args.out))
    print(f"\nSaved {len(prs.slides)} slides → {args.out}")


if __name__ == "__main__":
    main()
