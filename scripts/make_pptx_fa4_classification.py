#!/usr/bin/env python3
"""
make_pptx_fa4_classification.py
================================
PPT v2 — 4 FA subtype classification story using Annabel vinc/control labels.

Covers:
  1. Label stats & setup
  2. Stage 1 binary (no-adh vs adh) results — s1v3 / s2v2 / s3v1
  3. FA4-coloured UMAPs in Stage 1 latent space
  4. 4-class LightGBM on Stage 1 latents (generated inline)
  5. Overlays on microscopy frames
  6. Stage 2 plan — dedicated SupCon AE (currently training)

Usage:
  python scripts/make_pptx_fa4_classification.py
  python scripts/make_pptx_fa4_classification.py --out fa4_v2.pptx
"""
from __future__ import annotations

import argparse
import io
import tempfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────

RUNS      = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run")
LABEL_DIR = Path("/net/projects/CLS/lding/data/fa_data_analysis/labelling")
OUT       = Path("fa4_classification_v2.pptx")

SPLITS      = ["s1v3", "s2v2", "s3v1"]
SPLIT_LABEL = {"s1v3": "1 train / 3 val", "s2v2": "2 / 2", "s3v1": "3 train / 1 val"}
LABEL_ORDER_4 = ["Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
LABEL_SHORT   = {"Nascent Adhesion": "NA", "focal complex": "FC",
                 "focal adhesion": "FA", "fibrillar adhesion": "Fib"}

def supcon2_dir(split):   return RUNS / f"annabel_vinc_supcon2_{split}"
def cls_dir(split):       return supcon2_dir(split) / "fa_cls_zrecon"
def stage2_dir(split):    return RUNS / f"annabel_vinc_supcon2_stage2_{split}"
def stage2_cls_dir(split): return stage2_dir(split) / "stage2_cls"

# ── slide geometry ─────────────────────────────────────────────────────────────

SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.12)

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_GREEN = RGBColor(0x37, 0x8B, 0x4A)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)
C_RED   = RGBColor(0xC0, 0x00, 0x00)

# ── low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _img(slide, src, l, t, max_w, max_h):
    if src is None: return False
    p = Path(src)
    if not p.exists(): return False
    try:
        if p.suffix.lower() in {".tif", ".tiff"}:
            import tifffile
            arr = tifffile.imread(str(p))
            if arr.ndim == 4: arr = arr[0]
            if arr.ndim == 3 and arr.shape[0] < arr.shape[-1]: arr = arr[0]
            pil = Image.fromarray(arr.astype(np.uint8) if arr.max() > 1 else
                                  (arr * 255).astype(np.uint8)).convert("RGB")
        else:
            pil = Image.open(str(p)).convert("RGB")
        iw, ih = pil.size
        scale  = min(max_w / Inches(iw / 96), max_h / Inches(ih / 96))
        rw, rh = Inches(iw / 96) * scale, Inches(ih / 96) * scale
        buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
        slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                                  t + (max_h - rh) / 2, rw, rh)
        return True
    except Exception as e:
        print(f"  [warn] {p.name}: {e}")
        return False

def _img_ar(slide, src, l, t, max_w, max_h, ph="[pending]"):
    if not _img(slide, src, l, t, max_w, max_h):
        _txt(slide, l, t + max_h/2 - Inches(0.15), max_w, Inches(0.3),
             ph, size=9, color=C_GREY, align=PP_ALIGN.CENTER)

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2*PAD, Inches(0.26),
             subtitle, size=9, color=C_GREY)

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(1), Inches(2.2), SW - Inches(2), Inches(1.2),
         "FA Subtype Classification\n4-class · Annabel vinc/control",
         size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.6), SW - Inches(2), Inches(0.5),
         "Stage 1: binary SupCon AE  →  Stage 2: dedicated 4-class SupCon AE  ✓ COMPLETE",
         size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(4.2), SW - Inches(2), Inches(0.4),
         "vinc / control  ·  cio_mode_prt norm  ·  latent=12  proj=8  ·  LightGBM",
         size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), SH - Inches(0.5), SW - Inches(2), Inches(0.4),
         "2026-08-10", size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_label_stats(prs):
    """Label counts per FA type × train/val split."""
    sl = _blank(prs)
    _title_bar(sl, "Label Statistics — Annabel vinc/control  (4 FA subtypes)",
               "197 annotated adhesion patches  ·  s2v2 split shown")

    csv = cls_dir("s2v2") / "classification_results.csv"
    top = TITLE_H + Inches(0.35)
    if not csv.exists():
        _txt(sl, PAD, top, SW - 2*PAD, Inches(1), "classification_results.csv not found",
             size=11, color=C_RED)
        return

    df = pd.read_csv(csv)
    adh = df[df["annotation_label_name"].isin(LABEL_ORDER_4)].copy()

    # Bar chart of counts
    fig, axes = plt.subplots(1, 2, figsize=(11, 3.8), facecolor="white")
    for ax, split_name in zip(axes, ["train", "val"]):
        sub = adh[adh["split"] == split_name]
        counts = [len(sub[sub["annotation_label_name"] == lbl]) for lbl in LABEL_ORDER_4]
        bars = ax.bar([LABEL_SHORT[l] for l in LABEL_ORDER_4], counts,
                      color=["#2196F3","#FF9800","#4CAF50","#9C27B0"], edgecolor="white")
        for b, c in zip(bars, counts):
            ax.text(b.get_x() + b.get_width()/2, b.get_height() + 0.3, str(c),
                    ha="center", va="bottom", fontsize=10, fontweight="bold")
        ax.set_title(f"{split_name.capitalize()} set  (n={len(sub)})", fontsize=12)
        ax.set_ylabel("patch count"); ax.spines[["top","right"]].set_visible(False)
        ax.set_ylim(0, max(counts) * 1.25 + 1)
    fig.suptitle("s2v2 split (frames 0+1 train / frames 2+3 val)", fontsize=11)
    fig.tight_layout()
    buf = io.BytesIO(); fig.savefig(buf, format="png", dpi=130, bbox_inches="tight"); buf.seek(0)
    plt.close(fig)
    sl.shapes.add_picture(buf, Inches(1.2), top, Inches(10.8), Inches(4.5))

    # Overall totals text
    totals = {lbl: len(adh[adh["annotation_label_name"] == lbl]) for lbl in LABEL_ORDER_4}
    txt = "Total labeled adhesion patches:  " + "  ·  ".join(
        f"{LABEL_SHORT[lbl]}={n}" for lbl, n in totals.items())
    _txt(sl, PAD, SH - Inches(0.4), SW - 2*PAD, Inches(0.35),
         txt, size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_binary_results(prs):
    """Stage 1 binary confusion matrices, 3 splits side by side."""
    sl = _blank(prs)
    _title_bar(sl, "Stage 1 — Binary SupCon AE: No-adhesion vs Adhesion  (val set)",
               "SupCon 2-class  ·  LightGBM on z_recon  ·  3 train/val splits")
    top = TITLE_H + PAD
    cap_h = Inches(0.28)
    w = (SW - 4*PAD) / 3
    h = SH - top - PAD - cap_h

    for i, sp in enumerate(SPLITS):
        l   = PAD + i * (w + PAD)
        img = cls_dir(sp) / "confusion_matrix_norm_val.png"
        _img_ar(sl, img, l, top, w, h)

        # Read metrics
        m_path = cls_dir(sp) / "metrics.csv"
        cap = SPLIT_LABEL[sp]
        if m_path.exists():
            m = pd.read_csv(m_path)
            f1s = m["f1"].values
            cap += f"\nF1: {f1s[0]:.3f} / {f1s[1]:.3f}"
        _txt(sl, l, top + h, w, cap_h, cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_fa4_umap(prs, img_fn, title, subtitle=""):
    """3-column FA4 UMAP side by side."""
    sl = _blank(prs)
    _title_bar(sl, title, subtitle)
    top   = TITLE_H + PAD
    cap_h = Inches(0.26)
    w = (SW - 4*PAD) / 3
    h = SH - top - PAD - cap_h
    for i, sp in enumerate(SPLITS):
        l   = PAD + i * (w + PAD)
        img = cls_dir(sp) / img_fn
        _img_ar(sl, img, l, top, w, h)
        _txt(sl, l, top + h, w, cap_h, SPLIT_LABEL[sp],
             size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_4cls_confusion(prs, split="s2v2"):
    """Run 4-class LightGBM on Stage 1 latents; show confusion matrix."""
    sl = _blank(prs)
    _title_bar(sl,
        f"4-class FA subtype — LightGBM on Stage 1 z_recon  ({SPLIT_LABEL[split]})",
        "Train on adhesion patches in train split; evaluate on val split adhesion patches")
    top = TITLE_H + Inches(0.1)

    # ── load data ────────────────────────────────────────────────────────────
    csv = cls_dir(split) / "classification_results.csv"
    if not csv.exists():
        _txt(sl, PAD, top + Inches(0.5), SW - 2*PAD, Inches(0.4),
             "classification_results.csv not found", size=11, color=C_RED)
        return

    df = pd.read_csv(csv)
    # We need z_* features — they're not in classification_results; use blind_test latents
    lat_csv = supcon2_dir(split) / "blind_test" / "vinc_control_latents.csv"
    if not lat_csv.exists():
        _txt(sl, PAD, top + Inches(0.5), SW - 2*PAD, Inches(0.4),
             "latents CSV not found", size=11, color=C_RED)
        return

    # Load latents and merge with labels+split from classification_results
    lat = pd.read_csv(lat_csv)
    z_cols = [c for c in lat.columns if c.startswith("z_")]

    # classification_results has filename + split column (binary labels only)
    # Use Annabel 4-class label file for FA subtypes
    ann_csv = LABEL_DIR / "vinc_control_label_Annabel_20260715_1554.csv"
    ann = pd.read_csv(ann_csv)[["filename", "label"]].rename(columns={"label": "fa4_label"})
    meta = df[["filename", "split"]].copy()
    merged = lat.merge(meta, on="filename", how="left")
    merged = merged.merge(ann, on="filename", how="left")

    ADHESION = set(LABEL_ORDER_4)
    train_mask = (merged["split"] == "train") & (merged["fa4_label"].isin(ADHESION))
    val_mask   = (merged["split"] == "val")   & (merged["fa4_label"].isin(ADHESION))
    train_df = merged[train_mask]
    val_df   = merged[val_mask]

    lo4_present = [l for l in LABEL_ORDER_4 if l in set(train_df["fa4_label"])]
    lo4_int     = {l: i for i, l in enumerate(lo4_present)}
    X_tr = train_df[z_cols].values.astype(np.float32)
    y_tr = np.array([lo4_int[l] for l in train_df["fa4_label"]])

    # ── fit classifier ───────────────────────────────────────────────────────
    try:
        from lightgbm import LGBMClassifier
        import joblib
        clf = LGBMClassifier(n_estimators=500, learning_rate=0.05, num_leaves=31,
                             min_child_samples=3, class_weight="balanced",
                             random_state=42, verbose=-1, n_jobs=2)
    except ImportError:
        from sklearn.ensemble import GradientBoostingClassifier
        from sklearn.utils.class_weight import compute_sample_weight
        w = compute_sample_weight("balanced", y_tr)
        clf = GradientBoostingClassifier(n_estimators=300, max_depth=4,
                                         learning_rate=0.05, random_state=42)
    clf.fit(X_tr, y_tr)

    # ── evaluate ─────────────────────────────────────────────────────────────
    from sklearn.metrics import (confusion_matrix, ConfusionMatrixDisplay,
                                 balanced_accuracy_score)
    n_tr, n_val = len(train_df), len(val_df)
    txt_lines = [
        f"Train: {n_tr} adhesion patches  ({' · '.join(f'{LABEL_SHORT[l]}={sum(y_tr==i)}' for l,i in lo4_int.items())})",
    ]

    has_val = len(val_df) > 0
    if has_val:
        val_present = [l for l in LABEL_ORDER_4 if l in set(val_df["fa4_label"])]
        vp_int = {l: lo4_int.get(l, -1) for l in val_present}
        X_val  = val_df[z_cols].values.astype(np.float32)
        y_val_str = val_df["fa4_label"].tolist()
        y_pred_int = clf.predict(X_val)
        y_pred_str = [lo4_present[int(p)] if int(p) < len(lo4_present) else "?"
                      for p in y_pred_int]
        acc = sum(a == b for a, b in zip(y_val_str, y_pred_str)) / len(y_val_str)
        try:
            bal = balanced_accuracy_score(y_val_str, y_pred_str)
        except Exception:
            bal = float("nan")

        n_val_by_cls = {l: sum(1 for s in y_val_str if s == l) for l in LABEL_ORDER_4}
        txt_lines.append(
            f"Val:   {n_val} adhesion patches  ({' · '.join(f'{LABEL_SHORT[l]}={n_val_by_cls.get(l,0)}' for l in LABEL_ORDER_4 if l in n_val_by_cls)})"
        )
        txt_lines.append(f"Val accuracy: {acc*100:.1f}%   Balanced acc: {bal*100:.1f}%")

        cm = confusion_matrix(y_val_str, y_pred_str, labels=val_present)
        cm_norm = cm.astype(float) / cm.sum(axis=1, keepdims=True).clip(1)
        short_labels = [LABEL_SHORT[l] for l in val_present]

        fig, ax = plt.subplots(figsize=(5, 4.5), facecolor="white")
        disp = ConfusionMatrixDisplay(cm_norm, display_labels=short_labels)
        disp.plot(ax=ax, colorbar=True, values_format=".2f", cmap="Blues")
        ax.set_title(f"Stage 1 latents → 4-class  |  {SPLIT_LABEL[split]}  val\n"
                     f"acc={acc*100:.1f}%  bal={bal*100:.1f}%", fontsize=10)
        fig.tight_layout()
        buf = io.BytesIO(); fig.savefig(buf, format="png", dpi=130, bbox_inches="tight"); buf.seek(0)
        plt.close(fig)
        sl.shapes.add_picture(buf, Inches(0.5), top + Inches(0.2), Inches(5.5), Inches(5.0))
    else:
        _txt(sl, PAD, top + Inches(0.5), SW - 2*PAD, Inches(0.4),
             "No val adhesion patches for this split", size=11, color=C_AMBER)

    # Stats text on right
    txt_x = Inches(6.3)
    for k, line in enumerate(txt_lines):
        _txt(sl, txt_x, top + Inches(0.3) + k * Inches(0.45),
             SW - txt_x - PAD, Inches(0.4),
             line, size=10, color=C_BLACK)

    # Interpretation note
    note = ("NOTE: Stage 1 SupCon was trained for binary (no-adh vs adh).\n"
            "It was NOT optimised to separate FA subtypes — this is a probe\n"
            "of how much subtype structure is encoded incidentally.\n\n"
            "Stage 2 (completed) uses a dedicated SupCon AE trained ONLY on\n"
            "predicted-adhesion patches with 4-class labels → see next slides.")
    _txt(sl, txt_x, top + Inches(1.8), SW - txt_x - PAD, Inches(2.5),
         note, size=10, color=C_GREY)


def slide_overlays(prs, split="s2v2"):
    """Binary adhesion overlays on 4 training frames."""
    sl = _blank(prs)
    _title_bar(sl, f"Binary adhesion overlays on microscopy frames  ({SPLIT_LABEL[split]})",
               "Red=No adhesion  ·  Green=Adhesion  ·  frames 0-3 shown")
    top   = TITLE_H + PAD
    cap_h = Inches(0.26)
    w = (SW - 4*PAD) / 4
    h = SH - top - PAD - cap_h
    frames = [0, 1, 2, 3]
    for i, fr in enumerate(frames):
        l   = PAD + i * (w + PAD)
        img = cls_dir(split) / f"overlay_frame{fr:04d}.png"
        _img_ar(sl, img, l, top, w, h, f"[frame {fr}]")
        split_tag = "train" if (split == "s2v2" and fr < 2) or (split == "s1v3" and fr == 0) \
                    else "val"
        _txt(sl, l, top + h, w, cap_h, f"frame {fr}  ({split_tag})",
             size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_fa4_overlays(prs, split="s2v2"):
    """FA4-colored UMAPs (predicted_control_fa4, split_val_fa4)."""
    sl = _blank(prs)
    _title_bar(sl, f"Stage 1 latent space — FA4 subtype coloring  ({SPLIT_LABEL[split]})",
               "UMAP of all control patches; adhesion patches coloured by FA subtype")
    top   = TITLE_H + PAD
    cap_h = Inches(0.26)
    w = (SW - 4*PAD) / 3
    h = SH - top - PAD - cap_h

    items = [
        (cls_dir(split) / "umap_split_val_fa4.png",         "Val patches, FA4 coloured"),
        (cls_dir(split) / "umap_predicted_control_fa4.png", "All control, FA4 coloured"),
        (cls_dir(split) / "umap_predicted_all_fa4.png",     "All patches, FA4 coloured"),
    ]
    for i, (img, cap) in enumerate(items):
        l = PAD + i * (w + PAD)
        _img_ar(sl, img, l, top, w, h)
        _txt(sl, l, top + h, w, cap_h, cap, size=9, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_stage2_results(prs):
    """Stage 2 dedicated 4-class SupCon AE — confusion matrices for all 3 splits."""
    sl = _blank(prs)
    _title_bar(sl, "Stage 2 — Dedicated 4-class SupCon AE: Results  (300 epochs)",
               "New AE trained only on Stage-1 predicted-adhesion patches (5771/14879)  ·  "
               "LightGBM on Stage 2 latents  ·  all 3 splits")
    top   = TITLE_H + PAD
    cap_h = Inches(0.46)
    w = (SW - 4*PAD) / 3
    h = SH - top - PAD - cap_h

    notes = {
        "s1v3": "val: NA=18 FC=5 FA=168 Fib=6\n(frames 1-3)",
        "s2v2": "val: NA=9 FA=37 (frames 2-3)",
        "s3v1": "val: FA=46 only\n(frame 3 — single class)",
    }

    for i, sp in enumerate(SPLITS):
        l   = PAD + i * (w + PAD)
        img = stage2_cls_dir(sp) / "confusion_matrix_norm.png"
        _img_ar(sl, img, l, top, w, h - Inches(0.15), f"[{sp} — pending]")

        # metrics caption
        m_path = stage2_cls_dir(sp) / "metrics.csv"
        if m_path.exists():
            m = pd.read_csv(m_path).iloc[0]
            acc = m["accuracy"] * 100
            bal = m["balanced_acc"] * 100
            cap = (f"{SPLIT_LABEL[sp]}\nacc={acc:.1f}%  bal={bal:.1f}%\n"
                   f"train n={int(m['n_train'])}  val n={int(m['n_val'])}\n"
                   + notes.get(sp, ""))
        else:
            cap = f"{SPLIT_LABEL[sp]}\n[no metrics]"
        _txt(sl, l, top + h - Inches(0.15), w, cap_h, cap,
             size=8, color=C_GREY, align=PP_ALIGN.CENTER)


XDS_EVAL_DIR = RUNS / "fa4_xds_eval"

SCENARIO_LABELS = {
    "vinc_only":  "vinc only\n(within)",
    "pfak_only":  "pfak only\n(within)",
    "vinc->pfak": "vinc → pfak\n(cross)",
    "pfak->vinc": "pfak → vinc\n(cross)",
    "combined":   "combined\n(within)",
}
SCENARIO_COLORS = {
    "vinc_only":  "#1565C0",
    "pfak_only":  "#E65100",
    "vinc->pfak": "#6A1B9A",
    "pfak->vinc": "#00695C",
    "combined":   "#2E7D32",
}
SCENARIOS = [
    "vinc_only", "pfak_only", "vinc->pfak", "pfak->vinc", "combined",
]


def slide_xds_strategy(prs):
    """Strategy overview: Option A vs B, 5 eval scenarios, label efficiency design."""
    sl = _blank(prs)
    _title_bar(sl,
               "Cross-Dataset FA Subtype Classification — Strategy",
               "Option A (vinc-only AE)  ·  Option B (combined AE)  ·  "
               "5 eval scenarios  ·  10/25/50/75% label efficiency")
    top = TITLE_H + PAD
    col_w = (SW - 4 * PAD) / 2

    # ── Option A ──────────────────────────────────────────────────────────────
    _rect(sl, PAD, top, col_w, Inches(3.0), fill=RGBColor(0xBD, 0xD7, 0xEE))
    _txt(sl, PAD + Inches(0.1), top + Inches(0.08), col_w - Inches(0.2), Inches(0.36),
         "Option A — vinc-only Stage-2 AE  (zero-shot cross-ds)",
         size=12, bold=True, color=C_DARK)
    _txt(sl, PAD + Inches(0.1), top + Inches(0.46), col_w - Inches(0.2), Inches(2.4),
         "Model:  annabel_vinc_supcon2_stage2_s3v1\n"
         "  · Trained on vinc/ctrl predicted-adhesion patches only\n"
         "  · Encode vinc/ycomp + pfak/ctrl zero-shot\n\n"
         "Eval:  LightGBM on Stage-2 latents\n"
         "  · Repeated subsampling: 10%×10, 25%×4, 50%×4, 75%×4\n"
         "  · 5 scenarios (below)",
         size=10, color=C_BLACK)

    # ── Option B ──────────────────────────────────────────────────────────────
    bx = PAD + col_w + PAD
    _rect(sl, bx, top, col_w, Inches(3.0), fill=RGBColor(0xD5, 0xE8, 0xD4))
    _txt(sl, bx + Inches(0.1), top + Inches(0.08), col_w - Inches(0.2), Inches(0.36),
         "Option B — combined Stage-2 AE  (trained on all datasets)",
         size=12, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
    _txt(sl, bx + Inches(0.1), top + Inches(0.46), col_w - Inches(0.2), Inches(2.4),
         "Model:  annabel_vinc_supcon2_stage2_combined\n"
         "  · Trained on all 4 datasets simultaneously\n"
         "  · vinc/ctrl: Stage-1 gate + 4-class Annabel labels\n"
         "  · vinc/ycomp: all + Annabel labels\n"
         "  · pfak/ctrl: all + Annabel labels\n"
         "  · ppax/ctrl: all + Ernest labels  (4-class FA only)\n"
         "  · Pretrained from corrected_s3v1  →  300 epochs fine-tune\n\n"
         "Eval:  same as Option A",
         size=10, color=C_BLACK)

    # ── Scenarios ─────────────────────────────────────────────────────────────
    sy = top + Inches(3.1)
    _txt(sl, PAD, sy, SW - 2 * PAD, Inches(0.3),
         "Evaluation Scenarios  (5):",
         size=11, bold=True, color=C_DARK)
    sy += Inches(0.32)
    scenario_descs = [
        ("vinc_only",  "vinc/ctrl + vinc/ycomp  →  vinc  (within-dataset)"),
        ("pfak_only",  "pfak/ctrl  →  pfak  (within-dataset)"),
        ("vinc→pfak",  "train vinc, test pfak  (cross-dataset zero-shot)"),
        ("pfak→vinc",  "train pfak, test vinc  (cross-dataset zero-shot)"),
        ("combined",   "vinc + pfak combined  →  combined  (within)"),
    ]
    box_w = (SW - 6 * PAD) / 5
    for i, (name, desc) in enumerate(scenario_descs):
        bxi = PAD + i * (box_w + PAD)
        col = SCENARIO_COLORS.get(name.replace("→", "->"), C_GREY)
        _rect(sl, bxi, sy, box_w, Inches(1.0), fill=None)
        sl.shapes[-1].line.color.rgb = RGBColor(*bytes.fromhex(col[1:]))
        sl.shapes[-1].line.width = Pt(1.5)
        _txt(sl, bxi + Inches(0.05), sy + Inches(0.05), box_w - Inches(0.1), Inches(0.35),
             name, size=10, bold=True, color=RGBColor(*bytes.fromhex(col[1:])))
        _txt(sl, bxi + Inches(0.05), sy + Inches(0.42), box_w - Inches(0.1), Inches(0.55),
             desc, size=8, color=C_BLACK)

    # ── ppax zero-shot note ───────────────────────────────────────────────────
    _txt(sl, PAD, sy + Inches(1.1), SW - 2 * PAD, Inches(0.7),
         "ppax zero-shot  (both options):  apply Stage-1 binary GBM → Stage-2 4-class GBM "
         "to ppax/ctrl patches.\n"
         "Ernest labels are 4-class FA only (no 'No adhesion') — "
         "Stage-1 false negatives = FA patches predicted as no-adhesion.",
         size=10, color=C_GREY)


def _xds_suffix(option: str, variant: str = "zrecon", smote: bool = False) -> str:
    return f"{option}_{variant}" + ("_smote" if smote else "")


def slide_xds_results(prs, option: str, variant: str = "zrecon", smote: bool = False):
    """Cross-dataset label efficiency results — balanced accuracy + macro F1."""
    suffix = _xds_suffix(option, variant, smote)
    tag = f"Option {option}  {variant}" + ("  SMOTE" if smote else "")
    sl = _blank(prs)
    _title_bar(sl,
               f"Cross-Dataset Results — {tag}",
               "Stage-2 SupCon AE  ·  LightGBM  ·  5 scenarios  ·  repeated stratified subsampling")

    top   = TITLE_H + PAD
    img_h = (SH - top - 2 * PAD) / 2

    _img_ar(sl, XDS_EVAL_DIR / f"efficiency_bal_acc_{suffix}.png",
            PAD, top, SW - 2 * PAD, img_h,
            f"[{tag} — balanced accuracy — pending]")
    _img_ar(sl, XDS_EVAL_DIR / f"efficiency_macro_f1_{suffix}.png",
            PAD, top + img_h + PAD, SW - 2 * PAD, img_h,
            f"[{tag} — macro F1 — pending]")


def slide_xds_perclass(prs, option: str, variant: str = "zrecon", smote: bool = False):
    """Per-class F1 breakdown."""
    suffix = _xds_suffix(option, variant, smote)
    tag = f"Option {option}  {variant}" + ("  SMOTE" if smote else "")
    sl = _blank(prs)
    _title_bar(sl,
               f"Cross-Dataset Per-Class F1 — {tag}",
               "NA / FC / FA / Fibrillar  ·  rows = class, columns = scenario")

    top = TITLE_H + PAD
    _img_ar(sl, XDS_EVAL_DIR / f"efficiency_perclass_{suffix}.png",
            PAD, top, SW - 2 * PAD, SH - top - PAD,
            f"[{tag} — per-class F1 — pending]")


def slide_ppax_zeroshot(prs, option: str, variant: str = "zrecon", smote: bool = False):
    """ppax zero-shot evaluation slide."""
    suffix = _xds_suffix(option, variant, smote)
    tag = f"Option {option}  {variant}" + ("  SMOTE" if smote else "")
    sl = _blank(prs)
    _title_bar(sl,
               f"ppax Zero-Shot — {tag}",
               "Ernest labels (4-class FA only)  ·  Stage-1 binary gate → Stage-2 4-class")

    top = TITLE_H + PAD
    confusion_img = XDS_EVAL_DIR / f"ppax_zeroshot_confusion_{suffix}.png"
    summary_csv   = XDS_EVAL_DIR / f"ppax_zeroshot_summary_{suffix}.csv"

    _img_ar(sl, confusion_img,
            PAD, top, Inches(6.0), SH - top - Inches(1.0),
            f"[ppax zero-shot confusion — {tag} — pending]")

    rx = PAD + Inches(6.2)
    rw = SW - rx - PAD
    _txt(sl, rx, top, rw, Inches(0.3),
         "Summary statistics:", size=11, bold=True, color=C_DARK)

    if summary_csv.exists():
        try:
            s = pd.read_csv(summary_csv).iloc[0]
            lines = [
                f"Ernest-labeled patches:  {int(s.get('n_ernest', 0))}",
                f"Stage-1 recall:          {float(s.get('s1_recall', float('nan'))):.3f}",
                f"Stage-1 false negatives: {int(s.get('n_false_neg', 0))}",
                "",
                f"Stage-2 evaluated:       {int(s.get('n_eval', 0))}",
                f"Stage-2 balanced acc:    {float(s.get('bal_acc_s2', float('nan'))):.3f}",
            ]
            _txt(sl, rx, top + Inches(0.38), rw, SH - top - Inches(1.4),
                 "\n".join(lines), size=12, color=C_BLACK)
        except Exception as e:
            _txt(sl, rx, top + Inches(0.38), rw, Inches(1.0),
                 f"[error: {e}]", size=10, color=C_RED)
    else:
        _txt(sl, rx, top + Inches(0.38), rw, Inches(0.4),
             "[pending]", size=11, color=C_GREY)

    _txt(sl, PAD, SH - Inches(0.35), SW - 2 * PAD, Inches(0.3),
         "Stage-2 GBM trained on all vinc + pfak labeled patches (100%).  "
         "Stage-1 recall < 1.0 = Ernest-labeled FA patches filtered by binary gate.",
         size=8, color=C_GREY)


def slide_variant_comparison(prs, option: str):
    """Grouped bar chart: bal_acc at 75% training fraction across variants × scenarios."""
    VARIANTS = [
        ("zrecon",       False, "zrecon",       "#1565C0"),
        ("zproj",        False, "zproj",        "#E65100"),
        ("zrecon",       True,  "zrecon+SMOTE", "#6A1B9A"),
        ("zproj",        True,  "zproj+SMOTE",  "#00695C"),
    ]

    frac_show = 0.75

    # Build data: variant → scenario → bal_acc
    data = {}
    for variant, smote, label, _ in VARIANTS:
        suffix = _xds_suffix(option, variant, smote)
        csv = XDS_EVAL_DIR / f"summary_{suffix}.csv"
        if not csv.exists():
            continue
        df = pd.read_csv(csv)
        row = df[df["frac"] == frac_show]
        if len(row) == 0:
            continue
        data[label] = {r["scenario"]: r["bal_acc_mean"] * 100
                       for _, r in row.iterrows()}

    if not data:
        sl = _blank(prs)
        _title_bar(sl, f"Variant Comparison — Option {option}  (pending)", "")
        return

    sc_order = ["vinc_only", "pfak_only", "vinc->pfak", "pfak->vinc", "combined"]
    sc_labels = ["vinc only", "pfak only", "vinc→pfak", "pfak→vinc", "combined"]
    var_labels = list(data.keys())
    colors = {label: c for _, _, label, c in VARIANTS if label in data}

    x = np.arange(len(sc_order))
    n_vars = len(var_labels)
    width  = 0.18
    offsets = np.linspace(-(n_vars - 1) / 2, (n_vars - 1) / 2, n_vars) * width

    fig, ax = plt.subplots(figsize=(12, 4.5), facecolor="white")
    for i, (vlab, offset) in enumerate(zip(var_labels, offsets)):
        vals = [data[vlab].get(sc, 0) for sc in sc_order]
        bars = ax.bar(x + offset, vals, width, label=vlab,
                      color=colors.get(vlab, "#888888"), alpha=0.85)
        for bar, v in zip(bars, vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, v + 0.5,
                        f"{v:.0f}", ha="center", va="bottom", fontsize=7)

    ax.axhline(50, color="#AAAAAA", linestyle="--", linewidth=0.8, label="50% (chance)")
    ax.set_xticks(x)
    ax.set_xticklabels(sc_labels, fontsize=10)
    ax.set_ylabel("Balanced accuracy (%)", fontsize=10)
    ax.set_ylim(0, 80)
    ax.legend(fontsize=9, loc="upper right")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")
    ax.set_title(f"Option {option} — variant comparison at {int(frac_show*100)}% training labels",
                 fontsize=11, fontweight="bold")
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)

    sl = _blank(prs)
    _title_bar(sl,
               f"Variant Comparison — Option {option}  (balanced accuracy @ 75% labels)",
               "zrecon vs zproj vs +SMOTE  ·  5 cross-dataset scenarios")
    top = TITLE_H + PAD
    sl.shapes.add_picture(buf, PAD, top, SW - 2 * PAD, SH - top - Inches(0.3))


# ── main ───────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path, default=OUT)
    ap.add_argument("--split", default="s2v2", choices=SPLITS,
                    help="Primary split for 4-class analysis and overlays")
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides...")

    slide_cover(prs)
    print("  1. cover")

    slide_label_stats(prs)
    print("  2. label stats")

    slide_binary_results(prs)
    print("  3. Stage 1 binary confusion (3 splits)")

    # FA4 UMAP views
    slide_fa4_umap(prs, "umap_split_val_fa4.png",
                   "Stage 1 latent space — FA4 coloured (val split only)",
                   "Does the binary SupCon AE incidentally separate the 4 FA subtypes?")
    print("  4. FA4 UMAP (val, 3 splits)")

    slide_fa4_umap(prs, "umap_split_fa4.png",
                   "Stage 1 latent space — FA4 coloured (all annotated patches)",
                   "Train + val combined; circle = train, cross = val")
    print("  5. FA4 UMAP (all, 3 splits)")

    slide_fa4_overlays(prs, args.split)
    print(f"  6. FA4 UMAP panels ({args.split})")

    slide_4cls_confusion(prs, args.split)
    print(f"  7. 4-class LightGBM probe on Stage 1 latents ({args.split})")

    slide_overlays(prs, args.split)
    print(f"  8. Binary overlays on frames 0-3 ({args.split})")

    slide_stage2_results(prs)
    print("  9. Stage 2 results — 3-split confusion matrices")

    slide_xds_strategy(prs)
    print(" 10. Cross-dataset strategy (Option A / B)")

    # Option A — zrecon (baseline)
    slide_xds_results(prs, "A", "zrecon")
    slide_xds_perclass(prs, "A", "zrecon")
    slide_ppax_zeroshot(prs, "A", "zrecon")
    print(" 11-13. Option A zrecon (efficiency + per-class + ppax)")

    # Option A — zproj
    slide_xds_results(prs, "A", "zproj")
    slide_xds_perclass(prs, "A", "zproj")
    slide_ppax_zeroshot(prs, "A", "zproj")
    print(" 14-16. Option A zproj")

    # Option A — zrecon + SMOTE
    slide_xds_results(prs, "A", "zrecon", smote=True)
    slide_xds_perclass(prs, "A", "zrecon", smote=True)
    slide_ppax_zeroshot(prs, "A", "zrecon", smote=True)
    print(" 17-19. Option A zrecon+SMOTE")

    # Option A — zproj + SMOTE
    slide_xds_results(prs, "A", "zproj", smote=True)
    slide_xds_perclass(prs, "A", "zproj", smote=True)
    slide_ppax_zeroshot(prs, "A", "zproj", smote=True)
    print(" 20-22. Option A zproj+SMOTE")

    # Option A — variant comparison
    slide_variant_comparison(prs, "A")
    print(" 23. Option A variant comparison (bar chart)")

    # Option B — zrecon
    slide_xds_results(prs, "B", "zrecon")
    slide_xds_perclass(prs, "B", "zrecon")
    slide_ppax_zeroshot(prs, "B", "zrecon")
    print(" 24-26. Option B results")

    prs.save(str(args.out))
    n = len(prs.slides)
    print(f"\nSaved {n} slides → {args.out}")


if __name__ == "__main__":
    main()
