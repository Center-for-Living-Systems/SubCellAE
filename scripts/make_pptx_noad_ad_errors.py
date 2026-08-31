#!/usr/bin/env python3
"""
make_pptx_noad_ad_errors.py
============================
Build a PPT for the ad/no-ad binary classifier error analysis.

Slides
------
  1. Training label examples — sample patches per FA sub-label (5 classes),
     coloured bounding box, from the full 2cls annotation (all labeled frames)
  2. s1v3 (cfg 0) error grid + intensity
  3. s2v2 (cfg 1) error grid + intensity
  4. s3v1 (cfg 2) error grid + intensity

Usage
-----
  python scripts/make_pptx_noad_ad_errors.py
"""
from __future__ import annotations

import re
from pathlib import Path

import numpy as np
import pandas as pd
import tifffile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io

# ---------------------------------------------------------------------------
_REPO   = Path(__file__).resolve().parents[1]
DATA    = Path("/net/projects/CLS/lding/data/fa_data_analysis")
LAB     = DATA / "labelling"
PATCH   = DATA / "ae_results/patches/cio/vinc/control/tiff_patches32_mr10"
RES     = _REPO / "results"
OUT_PPT = RES / "noad_ad_error_analysis.pptx"

FULL_2CLS = LAB / "vinc_control_label_Annabel_20260715_1554_2cls.csv"
FULL_4CLS = LAB / "vinc_control_label_Annabel_20260715_1554.csv"

FA5_COLORS_HEX = {
    "No adhesion":        "#9467bd",   # purple
    "Nascent Adhesion":   "#1f77b4",   # blue
    "focal complex":      "#ff7f0e",   # orange
    "focal adhesion":     "#2ca02c",   # green
    "fibrillar adhesion": "#d62728",   # red
}
FA5_SHORT = {
    "No adhesion":        "No adhesion",
    "Nascent Adhesion":   "Nascent Adhesion (NA)",
    "focal complex":      "Focal Complex (FC)",
    "focal adhesion":     "Focal Adhesion (FA)",
    "fibrillar adhesion": "Fibrillar Adhesion (Fib)",
}
FA5_ORDER = ["No adhesion", "Nascent Adhesion", "focal complex",
             "focal adhesion", "fibrillar adhesion"]

BORDER_PX  = 4
PATCH_SIZE = 32
N_EXAMPLES = 10   # patches per class on training-label slide

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

CONFIGS = {
    0: "s1v3  (train: frame 0 · test: frames 1,2,3)",
    1: "s2v2  (train: frames 0,1 · test: frames 2,3)",
    2: "s3v1  (train: frames 0,1,2 · test: frame 3)",
}


# ---------------------------------------------------------------------------
# Image helpers

def _hex_to_rgb01(h: str) -> tuple[float, float, float]:
    return (int(h[1:3], 16) / 255, int(h[3:5], 16) / 255, int(h[5:7], 16) / 255)


def _norm(img: np.ndarray) -> np.ndarray:
    lo, hi = img.min(), img.max()
    return (img - lo) / (hi - lo + 1e-9)


def add_border(img_norm: np.ndarray, color_hex: str, px: int = BORDER_PX) -> np.ndarray:
    r, g, b = _hex_to_rgb01(color_hex)
    h, w    = img_norm.shape
    rgb     = np.stack([img_norm, img_norm, img_norm], axis=-1)
    out     = np.ones((h + 2 * px, w + 2 * px, 3), dtype=np.float32)
    out[:, :] = [r, g, b]
    out[px:px + h, px:px + w] = rgb
    return out


def load_patch(fname: str) -> np.ndarray | None:
    p = PATCH / fname
    if not p.exists():
        return None
    return tifffile.imread(str(p)).astype(np.float32)


def patches_to_pil(patch_list: list[np.ndarray]) -> Image.Image:
    """Concatenate bordered patches horizontally into a PIL image."""
    strip = np.concatenate(patch_list, axis=1)
    u8    = (np.clip(strip, 0, 1) * 255).astype(np.uint8)
    return Image.fromarray(u8, mode="RGB")


def pil_to_stream(img: Image.Image) -> io.BytesIO:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# PPT helpers

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def add_title(slide, text: str, top=Inches(0.15), fontsize=Pt(18)):
    txb = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.7), Inches(0.5))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = fontsize
    run.font.color.rgb = RGBColor(0x1a, 0x4a, 0x7a)


def add_subtitle(slide, text: str, top=Inches(0.62), fontsize=Pt(11)):
    txb = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.7), Inches(0.35))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = fontsize
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def add_label_text(slide, text: str, left, top, color_hex: str, fontsize=Pt(9)):
    r, g, b = (int(color_hex[i:i+2], 16) for i in (1, 3, 5))
    txb = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.3))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold  = True
    run.font.size  = fontsize
    run.font.color.rgb = RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Slide 1: Training label examples

def slide_training_examples(prs: Presentation):
    slide = blank_slide(prs)
    add_title(slide, "Training label examples — binary ad / no-ad classifier")
    add_subtitle(slide,
        "Border colour = FA sub-label (5 classes) · patches drawn from full labeled set (all frames) · npi=100")

    # Load 4cls labels and merge with 2cls
    df4 = pd.read_csv(FULL_4CLS)[["filename", "label"]].rename(columns={"label": "fa_label"})
    df2 = pd.read_csv(FULL_2CLS)[["filename", "label"]].rename(columns={"label": "binary_label"})
    df  = df2.merge(df4, on="filename", how="left")
    df["fa_label"] = df["fa_label"].fillna("No adhesion")

    rng = np.random.default_rng(7)

    cell_h = PATCH_SIZE + 2 * BORDER_PX
    cell_w = PATCH_SIZE + 2 * BORDER_PX

    top_start = Inches(1.05)
    row_h     = Inches(1.15)
    left_img  = Inches(2.3)
    scale     = 3.5  # display scale factor for 32px patches

    for ri, cls in enumerate(FA5_ORDER):
        color = FA5_COLORS_HEX[cls]
        top   = top_start + ri * row_h

        # Class label on left
        add_label_text(slide, FA5_SHORT[cls], Inches(0.2), top + Inches(0.2),
                       color, fontsize=Pt(10))

        # Sample patches
        pool = df[df["fa_label"] == cls]["filename"].tolist()
        if not pool:
            continue
        chosen = rng.choice(pool, size=min(N_EXAMPLES, len(pool)), replace=False)

        imgs = []
        for fn in chosen:
            img = load_patch(fn)
            if img is not None:
                imgs.append(add_border(_norm(img), color))

        if not imgs:
            continue

        pil_img  = patches_to_pil(imgs)
        buf      = pil_to_stream(pil_img)
        img_w_px = pil_img.width
        img_h_px = pil_img.height

        # Scale to fit row height (≈ 0.9 inches)
        disp_h = Inches(0.88)
        disp_w = int(disp_h * img_w_px / img_h_px)

        slide.shapes.add_picture(buf, left_img, top + Inches(0.05), disp_w, disp_h)

    # Count label text on right
    for ri, cls in enumerate(FA5_ORDER):
        n   = (df["fa_label"] == cls).sum()
        top = top_start + ri * row_h
        add_label_text(slide, f"n={n}", Inches(12.5), top + Inches(0.2),
                       FA5_COLORS_HEX[cls], fontsize=Pt(9))

    return slide


# ---------------------------------------------------------------------------
# Slides 2-4: per-config error analysis

def slide_cfg_errors(prs: Presentation, cfg_id: int):
    cfg_label = CONFIGS[cfg_id]
    grid_png  = RES / f"noad_ad_errors_grid_cfg{cfg_id}.png"
    int_png   = RES / f"noad_ad_errors_intensity_cfg{cfg_id}.png"
    pred_csv  = RES / f"noad_ad_error_predictions_cfg{cfg_id}.csv"

    slide = blank_slide(prs)
    add_title(slide, f"Ad/No-ad errors — {cfg_label}")

    # Load stats
    if pred_csv.exists():
        df      = pd.read_csv(pred_csv)
        n_total = len(df)
        n_err   = (~df["correct"]).sum()
        n_fn    = ((df["true_2cls"] == "adhesion") & ~df["correct"]).sum()
        n_fp    = ((df["true_2cls"] == "No adhesion") & ~df["correct"]).sum()
        # Compute BAcc from file
        from sklearn.metrics import balanced_accuracy_score
        bacc = balanced_accuracy_score(df["true_2cls"], df["pred_2cls"])
        sub  = (f"BAcc = {bacc:.3f}  ·  test patches = {n_total}  ·  "
                f"errors = {n_err} ({n_err/n_total*100:.1f}%)  ·  "
                f"FN = {n_fn}  FP = {n_fp}")
    else:
        sub = "(results not found)"
    add_subtitle(slide, sub, top=Inches(0.62), fontsize=Pt(10))

    # Error grid (left, large)
    if grid_png.exists():
        slide.shapes.add_picture(str(grid_png), Inches(0.2), Inches(1.0),
                                 Inches(8.6), Inches(3.5))

    # Intensity plot (right)
    if int_png.exists():
        slide.shapes.add_picture(str(int_png), Inches(6.5), Inches(1.0),
                                 Inches(6.6), Inches(6.2))

    # Bullet annotations (bottom left)
    bullets = [
        "FN = adhesion patches predicted as No-adhesion (missed FA structures)",
        "FP = No-adhesion patches predicted as adhesion (false alarms)",
        "Border colour = FA sub-label; triangle markers = errors in intensity plot",
    ]
    txb = slide.shapes.add_textbox(Inches(0.2), Inches(4.7), Inches(6.1), Inches(2.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return slide


# ---------------------------------------------------------------------------

def main():
    RES.mkdir(exist_ok=True)

    # Check all images exist
    missing = []
    for cfg_id in [0, 1, 2]:
        for suffix in ["grid", "intensity"]:
            p = RES / f"noad_ad_errors_{suffix}_cfg{cfg_id}.png"
            if not p.exists():
                missing.append(p.name)
    if missing:
        print(f"WARNING: missing images: {missing}")
        print("Run analyze_noad_ad_errors.py first.")

    prs = new_prs()

    print("Slide 1: training label examples ...")
    slide_training_examples(prs)

    for cfg_id in [0, 1, 2]:
        print(f"Slide {cfg_id + 2}: cfg {cfg_id} = {CONFIGS[cfg_id]} ...")
        slide_cfg_errors(prs, cfg_id)

    prs.save(str(OUT_PPT))
    print(f"\nSaved: {OUT_PPT}")


if __name__ == "__main__":
    main()
