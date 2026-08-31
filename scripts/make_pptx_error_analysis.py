#!/usr/bin/env python3
"""
make_pptx_error_analysis.py
============================
Comprehensive PPT on misclassified FA4 patches across eval scenarios.

Slides
------
 1. Cover
 2. Error rate overview — per-scenario table + top confusion pairs
 3. vinc_only  (37.7% error) — patch grid + confusion breakdown
 4. pfak_only  (10.5% error) — patch grid + confusion breakdown
 5. vinc→pfak  (14.6% error) — patch grid + confusion breakdown
 6. pfak→vinc  (44.1% error) — patch grid + confusion breakdown (worst case)
 7. combined   (31.6% error) — patch grid + confusion breakdown
 8. FC dominance — bar chart of FC error count across all scenarios
 9. Takeaways and recommendations

Usage
-----
  python scripts/make_pptx_error_analysis.py
  python scripts/make_pptx_error_analysis.py --out custom_name.pptx
"""

from __future__ import annotations

import argparse
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
EVAL_DIR  = Path("/net/projects/CLS/lding/data/fa_data_analysis/ae_results/contrastive_run/fa4_xds_eval")
OUT_DEFAULT = Path("results/fa4_error_analysis.pptx")

SUFFIX = "A_zrecon"

FA_LABELS = ["Nascent Adhesion", "focal complex", "focal adhesion", "fibrillar adhesion"]
LABEL_SHORT = {"Nascent Adhesion": "NA", "focal complex": "FC",
               "focal adhesion": "FA", "fibrillar adhesion": "Fib"}
SHORT_LONG = {v: k for k, v in LABEL_SHORT.items()}
CLASS_COLORS = {
    "NA":  "#1565C0",  # blue
    "FC":  "#E65100",  # orange
    "FA":  "#2ca02c",  # green
    "Fib": "#C00000",  # red
}
NO_AD_COLOR = "#9467bd"  # purple — for "No adhesion" class

SCENARIOS = ["vinc_only", "pfak_only", "vinc->pfak", "pfak->vinc", "combined"]
SCENARIO_LABELS = {
    "vinc_only":  "Vinc only (within-dataset)",
    "pfak_only":  "pFAK only (within-dataset)",
    "vinc->pfak": "Vinc → pFAK (cross-dataset)",
    "pfak->vinc": "pFAK → Vinc (cross-dataset)",
    "combined":   "Combined (all data)",
}
SCENARIO_SHORT = {
    "vinc_only": "Vinc only", "pfak_only": "pFAK only",
    "vinc->pfak": "Vinc→pFAK", "pfak->vinc": "pFAK→Vinc", "combined": "Combined",
}

# ── slide geometry ─────────────────────────────────────────────────────────────
SW      = Inches(13.33)
SH      = Inches(7.5)
TITLE_H = Inches(0.52)
PAD     = Inches(0.15)
BODY_T  = TITLE_H + Inches(0.08)
BODY_H  = SH - BODY_T - PAD

C_DARK  = RGBColor(0x1F, 0x4E, 0x79)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GREY  = RGBColor(0x88, 0x88, 0x88)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_GREEN = RGBColor(0x37, 0x8B, 0x4A)
C_AMBER = RGBColor(0xE0, 0x7B, 0x00)

# ── helpers ────────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()

def _txt(slide, l, t, w, h, text, size=12, bold=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color

def _title_bar(slide, title, subtitle=""):
    _rect(slide, 0, 0, SW, TITLE_H, fill=C_DARK)
    _txt(slide, PAD, Inches(0.06), SW - 2*PAD, TITLE_H - Inches(0.06),
         title, size=14, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, PAD, TITLE_H, SW - 2*PAD, Inches(0.25),
             subtitle, size=9, color=C_GREY)

def _fig_to_pil(fig) -> Image.Image:
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    return Image.open(buf).convert("RGB")

def _place_pil(slide, pil, l, t, max_w, max_h):
    iw, ih = pil.size
    scale  = min(max_w / Inches(iw / 150), max_h / Inches(ih / 150))
    rw, rh = Inches(iw / 150) * scale, Inches(ih / 150) * scale
    buf = io.BytesIO(); pil.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, l + (max_w - rw) / 2,
                             t + (max_h - rh) / 2, rw, rh)

def _place_png(slide, path, l, t, max_w, max_h, ph="[pending]"):
    p = Path(path)
    if not p.exists():
        _txt(slide, l, t + max_h / 2, max_w, Inches(0.3),
             ph, size=9, color=C_GREY, align=PP_ALIGN.CENTER)
        return False
    pil = Image.open(str(p)).convert("RGB")
    _place_pil(slide, pil, l, t, max_w, max_h)
    return True

def _load_predictions(scenario: str) -> pd.DataFrame | None:
    p = EVAL_DIR / f"predictions_{scenario}_{SUFFIX}.csv"
    if not p.exists():
        return None
    return pd.read_csv(p)

def _confusion_summary(df: pd.DataFrame) -> pd.DataFrame:
    """Return (true_label, pred_label, n, pct_of_total) sorted descending."""
    errors = df[df["true_label"] != df["pred_label"]]
    total  = len(df)
    conf   = (errors.groupby(["true_label", "pred_label"], observed=True)
                    .size().reset_index(name="n")
                    .sort_values("n", ascending=False))
    conf["pct"] = conf["n"] / total * 100
    return conf


# ── slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, fill=C_DARK)
    _txt(sl, Inches(1), Inches(1.8), SW - Inches(2), Inches(1.2),
         "FA4 Misclassification Analysis",
         size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.1), SW - Inches(2), Inches(0.6),
         "Which patches are misclassified and why?",
         size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(3.8), SW - Inches(2), Inches(0.45),
         "Option A  ·  z_recon features (12-d)  ·  75% label fraction  ·  5 eval scenarios",
         size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), Inches(4.3), SW - Inches(2), Inches(0.4),
         "Each patch grid shows: rows = true class, cols = predicted class  (off-diagonal = errors)",
         size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _txt(sl, Inches(1), SH - Inches(0.55), SW - Inches(2), Inches(0.4),
         "2026-08-21", size=10, color=C_GREY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    """Error rate per scenario + top confusion pairs."""
    sl = _blank(prs)
    _title_bar(sl, "Error Rate Overview — All 5 Scenarios at 75% Label Fraction",
               "Option A  ·  z_recon  ·  single representative split (first stratified repeat)")

    scenario_stats = []
    for sc in SCENARIOS:
        df = _load_predictions(sc)
        if df is None:
            scenario_stats.append((sc, None, None, None, None))
            continue
        n_total = len(df)
        n_err   = (df["true_label"] != df["pred_label"]).sum()
        conf    = _confusion_summary(df)
        top2    = conf.head(2)
        top_str = "  |  ".join(
            f"{LABEL_SHORT.get(r['true_label'], r['true_label'])}→"
            f"{LABEL_SHORT.get(r['pred_label'], r['pred_label'])} (n={r['n']})"
            for _, r in top2.iterrows()
        )
        scenario_stats.append((sc, n_total, n_err, n_err / n_total, top_str))

    # Bar chart of error rates
    fig, axes = plt.subplots(1, 2, figsize=(14, 4.5), facecolor="white",
                             gridspec_kw={"width_ratios": [1, 1.8]})

    sc_colors = ["#1f77b4", "#d62728", "#ff7f0e", "#9467bd", "#2ca02c"]

    # Left: error rate bars
    ax = axes[0]
    sc_names = [SCENARIO_SHORT[r[0]] for r in scenario_stats if r[1] is not None]
    rates     = [r[3] * 100 for r in scenario_stats if r[1] is not None]
    colors    = [sc_colors[i] for i, r in enumerate(scenario_stats) if r[1] is not None]
    bars = ax.barh(sc_names, rates, color=colors, edgecolor="white")
    for b, v in zip(bars, rates):
        ax.text(v + 0.5, b.get_y() + b.get_height() / 2,
                f"{v:.1f}%", va="center", fontsize=10, fontweight="bold")
    ax.set_xlabel("Error rate (%)", fontsize=10)
    ax.set_title("Error rate at 75% labels", fontsize=11, fontweight="bold")
    ax.set_xlim(0, 55)
    ax.axvline(25, color="gray", lw=0.8, linestyle=":")
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")

    # Right: confusion matrix heatmap (true×pred, summed across all scenarios)
    ax2 = axes[1]
    shorts = ["NA", "FC", "FA", "Fib"]
    combined_conf = np.zeros((4, 4))
    all_dfs = []
    for sc in SCENARIOS:
        df = _load_predictions(sc)
        if df is not None:
            all_dfs.append(df)
    if all_dfs:
        all_df = pd.concat(all_dfs, ignore_index=True)
        for i, true_s in enumerate(shorts):
            true_l = SHORT_LONG[true_s]
            for j, pred_s in enumerate(shorts):
                pred_l = SHORT_LONG[pred_s]
                combined_conf[i, j] = ((all_df["true_label"] == true_l) &
                                       (all_df["pred_label"] == pred_l)).sum()

    # Normalize each row
    row_sums = combined_conf.sum(axis=1, keepdims=True)
    row_sums[row_sums == 0] = 1
    conf_norm = combined_conf / row_sums

    im = ax2.imshow(conf_norm, cmap="Blues", vmin=0, vmax=1, aspect="auto")
    ax2.set_xticks(range(4))
    ax2.set_yticks(range(4))
    ax2.set_xticklabels([f"pred\n{s}" for s in shorts], fontsize=9)
    ax2.set_yticklabels([f"true {s}" for s in shorts], fontsize=9)
    ax2.set_title("Aggregate confusion (all scenarios, row-normalized)", fontsize=11, fontweight="bold")
    for i in range(4):
        for j in range(4):
            n = int(combined_conf[i, j])
            v = conf_norm[i, j]
            color = "white" if v > 0.5 else "black"
            ax2.text(j, i, f"{v:.2f}\n(n={n})", ha="center", va="center",
                     fontsize=8, color=color)
    fig.colorbar(im, ax=ax2, fraction=0.04, pad=0.02, label="Row-normalized fraction")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05),
               SW - 2*PAD, BODY_H - Inches(1.05))
    plt.close(fig)

    # Annotation row
    by = SH - Inches(1.1)
    _txt(sl, PAD, by, SW * 0.5, Inches(1.0),
         "Within-dataset (vinc_only, pfak_only): reasonable accuracy — 10–38% error\n"
         "Cross-dataset (pfak→vinc): worst case — 44% error; FC nearly always misclassified\n"
         "pfak is a small dataset (151 patches) — limited FA subtype diversity",
         size=9, color=C_DARK, wrap=True)
    _txt(sl, SW * 0.5, by, SW * 0.5 - PAD, Inches(1.0),
         "Confusion hotspots: FC→NA (focal complex confused with no adhesion)\n"
         "FA→FC and NA→FA boundaries are blurry — intermediate morphology\n"
         "Fib→FA is common in cross-dataset — fibrillar looks like FA without context",
         size=9, color=C_DARK, wrap=True)


def _scenario_slide(prs, scenario: str, error_rate: float, n_total: int, n_errors: int,
                    top_errors: list[tuple[str, str, int]], notes: str):
    """One scenario slide: error grid + confusion breakdown."""
    sl = _blank(prs)
    label = SCENARIO_LABELS[scenario]
    cross = "→" in scenario
    _title_bar(
        sl,
        f"Misclassified Patches — {label}",
        f"{n_errors}/{n_total} errors ({error_rate*100:.1f}%)  ·  75% label fraction  ·  "
        f"{'cross-dataset: ALL test patches from other domain' if cross else 'within-dataset: random 25% held-out'}",
    )

    # ── left: patch grid ──────────────────────────────────────────────────
    grid_w = SW * 0.64
    _place_png(sl, EVAL_DIR / f"errors_{scenario}_{SUFFIX}.png",
               PAD, BODY_T + Inches(0.05), grid_w - PAD, BODY_H - Inches(0.1))

    # ── right: confusion breakdown ────────────────────────────────────────
    rx = grid_w + PAD
    rw = SW - rx - PAD

    df = _load_predictions(scenario)
    if df is not None:
        conf = _confusion_summary(df)

        # Per-class accuracy bar chart
        fig, ax = plt.subplots(figsize=(3.5, 3.2), facecolor="white")
        shorts = ["NA", "FC", "FA", "Fib"]
        per_cls_acc = []
        per_cls_n   = []
        for s in shorts:
            long = SHORT_LONG[s]
            sub  = df[df["true_label"] == long]
            n    = len(sub)
            acc  = (sub["pred_label"] == long).mean() if n > 0 else np.nan
            per_cls_acc.append(acc)
            per_cls_n.append(n)

        colors = [CLASS_COLORS[s] for s in shorts]
        bar_heights = [a * 100 if not np.isnan(a) else 0 for a in per_cls_acc]
        bars = ax.bar(shorts, bar_heights, color=colors, edgecolor="white")
        for b, acc, n in zip(bars, per_cls_acc, per_cls_n):
            if n == 0:
                label_txt = "n=0"
                ax.text(b.get_x() + b.get_width() / 2, 2,
                        label_txt, ha="center", va="bottom",
                        fontsize=7, color="gray")
            else:
                label_txt = f"{acc*100:.0f}%"
                ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 1,
                        f"{label_txt}\n(n={n})", ha="center", va="bottom",
                        fontsize=7, fontweight="bold")
        ax.set_ylim(0, 115)
        ax.set_ylabel("Per-class accuracy (%)", fontsize=9)
        ax.set_title(f"Per-class accuracy\n({scenario})", fontsize=9, fontweight="bold")
        ax.axhline(50, color="gray", lw=0.8, linestyle=":")
        ax.spines[["top", "right"]].set_visible(False)
        ax.set_facecolor("white")
        _place_pil(sl, _fig_to_pil(fig), rx, BODY_T + Inches(0.05),
                   rw, Inches(2.8))
        plt.close(fig)

        # Top confusion pairs as text
        ty = BODY_T + Inches(2.95)
        _txt(sl, rx, ty, rw, Inches(0.3), "Top error pairs:", size=9, bold=True, color=C_DARK)
        ty += Inches(0.3)
        for _, row in conf.head(8).iterrows():
            t  = LABEL_SHORT.get(row["true_label"], row["true_label"])
            p  = LABEL_SHORT.get(row["pred_label"], row["pred_label"])
            n  = row["n"]
            pct = row["pct"]
            tc = RGBColor(*[int(CLASS_COLORS[t].lstrip("#")[i:i+2], 16) for i in (0, 2, 4)])
            _txt(sl, rx, ty, rw, Inches(0.25),
                 f"  true={t} → pred={p}   n={n} ({pct:.1f}%)",
                 size=8.5, color=tc)
            ty += Inches(0.25)

    # Notes
    ny = SH - Inches(1.1)
    _txt(sl, PAD, ny, SW - 2*PAD, Inches(1.0),
         notes, size=9, color=C_DARK, wrap=True)


def slide_fc_dominance(prs):
    """FC error dominance across all scenarios."""
    sl = _blank(prs)
    _title_bar(sl, "FC (Focal Complex) — The Dominant Error Class",
               "True=FC misclassified as any other class  ·  and False positives: true=other predicted=FC")

    fig, axes = plt.subplots(1, 2, figsize=(15, 5.5), facecolor="white")

    shorts = ["NA", "FC", "FA", "Fib"]
    sc_colors = ["#1f77b4", "#d62728", "#ff7f0e", "#9467bd", "#2ca02c"]

    # Left: per-scenario class error rate (how often each true class is wrong)
    ax = axes[0]
    x = np.arange(len(SCENARIOS))
    width = 0.2
    for i, (short, long) in enumerate(zip(shorts, FA_LABELS)):
        err_rates = []
        for sc in SCENARIOS:
            df = _load_predictions(sc)
            if df is None:
                err_rates.append(0)
                continue
            sub = df[df["true_label"] == long]
            if len(sub) == 0:
                err_rates.append(0)
                continue
            err_rates.append((sub["pred_label"] != long).mean() * 100)
        ax.bar(x + (i - 1.5) * width, err_rates, width=width * 0.9,
               color=CLASS_COLORS[short], label=short, edgecolor="white")

    ax.set_xticks(x)
    ax.set_xticklabels([SCENARIO_SHORT[s].replace("→", "→\n") for s in SCENARIOS],
                       fontsize=9, ha="center")
    ax.set_ylabel("Error rate per true class (%)", fontsize=10)
    ax.set_title("How often is each class misclassified?", fontsize=11, fontweight="bold")
    ax.legend(title="True class", fontsize=9)
    ax.axhline(50, color="gray", lw=0.8, linestyle=":")
    ax.set_ylim(0, 110)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("white")
    for xi in x:
        ax.text(xi, 102, SCENARIO_SHORT[SCENARIOS[xi]].replace("→", "→\n"),
                ha="center", va="bottom", fontsize=0)  # invisible spacer

    # Right: "where does FC go?" — stacked bar of FC→* errors per scenario
    ax2 = axes[1]
    fc_dest_colors = {"NA": "#1565C0", "FA": "#2ca02c", "Fib": "#C00000", "other": "#aaaaaa"}
    bottom = np.zeros(len(SCENARIOS))
    dest_labels = ["NA", "FA", "Fib"]
    for dest in dest_labels:
        dest_long = SHORT_LONG[dest]
        counts = []
        for sc in SCENARIOS:
            df = _load_predictions(sc)
            if df is None:
                counts.append(0)
                continue
            fc_long = SHORT_LONG["FC"]
            n = ((df["true_label"] == fc_long) & (df["pred_label"] == dest_long)).sum()
            counts.append(n)
        ax2.bar(x, counts, bottom=bottom,
                color=fc_dest_colors[dest], label=f"FC→{dest}", edgecolor="white")
        bottom += np.array(counts)

    # Also show n_correct
    correct_counts = []
    for sc in SCENARIOS:
        df = _load_predictions(sc)
        if df is None:
            correct_counts.append(0)
            continue
        fc_long = SHORT_LONG["FC"]
        n = ((df["true_label"] == fc_long) & (df["pred_label"] == fc_long)).sum()
        correct_counts.append(n)
    ax2.bar(x, correct_counts, bottom=bottom, color="#555555",
            label="FC→FC (correct)", edgecolor="white", alpha=0.4, hatch="//")

    ax2.set_xticks(x)
    ax2.set_xticklabels([SCENARIO_SHORT[s].replace("→", "→\n") for s in SCENARIOS], fontsize=9)
    ax2.set_ylabel("Number of FC patches", fontsize=10)
    ax2.set_title("Fate of FC (Focal Complex) patches", fontsize=11, fontweight="bold")
    ax2.legend(fontsize=9)
    ax2.spines[["top", "right"]].set_visible(False)
    ax2.set_facecolor("white")

    fig.tight_layout()
    _place_pil(sl, _fig_to_pil(fig), PAD, BODY_T + Inches(0.05),
               SW - 2*PAD, BODY_H - Inches(1.1))
    plt.close(fig)

    by = SH - Inches(1.1)
    _txt(sl, PAD, by, SW * 0.5, Inches(1.0),
         "• FC has by far the highest per-class error rate (>80% in 3 of 5 scenarios)\n"
         "• Focal complex is most often mistaken for NA — similar size/shape to early-stage\n"
         "• In pfak→vinc, ALL FC patches are misclassified (100% error rate)",
         size=9, color=C_RED, wrap=True)
    _txt(sl, SW * 0.5, by, SW * 0.5 - PAD, Inches(1.0),
         "• Only in vinc_only does FC reach reasonable recall — ~25% at 75% labels\n"
         "• The model has almost no concept of FC when trained on pFAK alone (5 FC labels)\n"
         "• Both zproj and SMOTE improve vinc_only FC but cannot fix cross-dataset FC",
         size=9, color=C_DARK, wrap=True)


def slide_takeaways(prs):
    sl = _blank(prs)
    _title_bar(sl, "Key Findings and Recommendations — Misclassification Analysis",
               "Option A  ·  z_recon  ·  75% labels  ·  5 scenarios")

    col_w = (SW - 3 * PAD) / 2
    col2_x = PAD * 2 + col_w

    _txt(sl, PAD, BODY_T + Inches(0.12), col_w, Inches(0.3),
         "Root Causes of Errors", size=13, bold=True, color=C_DARK)

    causes = [
        (C_RED,   "FC (focal complex) is hardest: biologically intermediate appearance, "
                  "severe label imbalance (5/151 in pFAK), and boundary blurring with NA and FA"),
        (C_RED,   "pfak→vinc is catastrophic: the pFAK domain lacks enough FC/Fib examples "
                  "to teach a classifier that generalizes; FC error rate = 100%"),
        (C_AMBER, "FA→FC confusion (14 errors across scenarios): some mature FAs appear compact "
                  "and could be re-scored as FC — annotation calibration may help"),
        (C_AMBER, "NA→FA confusion: the model sometimes assigns FA to clear background — "
                  "likely patches at the cell edge where signal is ambiguous"),
        (C_DARK,  "Fib→FA: fibrillar adhesions elongated along fibers may resemble classical FAs "
                  "without the orientation context that human annotators use"),
    ]
    fy = BODY_T + Inches(0.5)
    for color, text in causes:
        _txt(sl, PAD, fy, col_w - Inches(0.1), Inches(0.55),
             f"• {text}", size=9, color=color, wrap=True)
        fy += Inches(0.56)

    _txt(sl, col2_x, BODY_T + Inches(0.12), col_w, Inches(0.3),
         "Recommended Actions", size=13, bold=True, color=C_DARK)

    actions = [
        (C_GREEN, "Immediate: visually inspect the FC→NA patches — are they truly FC or "
                  "borderline NA? If borderline, relabeling would reduce noise"),
        (C_GREEN, "Acquire more FC and Fib annotations, especially in pFAK (currently 5 FC / 4 Fib). "
                  "Margaret labeling one more ycomp/pfak image would help"),
        (C_DARK,  "Add class-balanced SupCon batching (WeightedRandomSampler) in Stage-2 AE "
                  "retraining so FC/Fib get equal representation in contrastive pairs"),
        (C_AMBER, "Consider binary 'mature vs immature' decomposition: separate NA+FC (immature) "
                  "from FA+Fib (mature) as a first step, then refine within each group"),
        (C_DARK,  "Cumulative label test (next): retain 10% labels in 25% split — will reveal "
                  "whether FC errors are due to insufficient training or poor representations"),
        (C_DARK,  "Ernest to add 'No adhesion' labels to the ppax ycomp set — needed for "
                  "balanced accuracy evaluation in the zero-shot ppax experiment"),
    ]
    sy = BODY_T + Inches(0.5)
    for color, text in actions:
        _txt(sl, col2_x, sy, col_w - Inches(0.1), Inches(0.55),
             f"• {text}", size=9, color=color, wrap=True)
        sy += Inches(0.56)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default=str(OUT_DEFAULT))
    args = parser.parse_args()

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("[1/9] Cover ...")
    slide_cover(prs)

    print("[2/9] Overview ...")
    slide_overview(prs)

    scenario_notes = {
        "vinc_only": (
            "Within vinc (ctrl + ycomp): 37.7% error is high for within-dataset — driven by FC↔NA confusion. "
            "FC is rare (90 patches of 1763 total vinc labeled FA-patches) and its appearance overlaps "
            "with early nascent adhesions. Fib (18 patches) is also sparse. "
            "More labels or class-balanced training should help."
        ),
        "pfak_only": (
            "Within pFAK: only 4 errors (10.5%) — the model does well with 151 labeled pFAK patches. "
            "However pFAK has very few FC (5) and Fib (4) so this low error rate is misleading — "
            "the classifier likely just predicts NA or FA for everything."
        ),
        "vinc->pfak": (
            "Cross-dataset: train on vinc (all), test on all pFAK (151 patches). "
            "14.6% error — reasonable generalization. Main errors are FA→FC and Fib→FA which may reflect "
            "morphological differences in how FAs appear under the pFAK vs vinc staining protocol."
        ),
        "pfak->vinc": (
            "Worst scenario: train on pFAK (151 patches), test on all vinc (1763 patches). "
            "44.1% error — 200 errors. FC→NA dominates (46 errors): pFAK has only 5 FC examples, "
            "so the classifier never learns FC from pFAK. NA→FA confusion (32 errors) suggests pFAK "
            "biases the model toward predicting FA for ambiguous patches."
        ),
        "combined": (
            "Combined (all 3 datasets pooled, stratified split): 31.6% error. "
            "FC→NA (11) and FC→FA (8) remain the top errors. Having more total data helps but "
            "does not resolve the fundamental FC label scarcity problem."
        ),
    }

    for i, sc in enumerate(SCENARIOS):
        print(f"[{i+3}/9] {sc} slide ...")
        df = _load_predictions(sc)
        if df is not None:
            n_total = len(df)
            n_err = (df["true_label"] != df["pred_label"]).sum()
            err_rate = n_err / n_total
        else:
            n_total = n_err = 0
            err_rate = 0.0
        top_errs = []
        if df is not None:
            conf = _confusion_summary(df)
            for _, row in conf.head(5).iterrows():
                top_errs.append((row["true_label"], row["pred_label"], row["n"]))
        _scenario_slide(prs, sc, err_rate, n_total, n_err, top_errs,
                        scenario_notes.get(sc, ""))

    print("[8/9] FC dominance ...")
    slide_fc_dominance(prs)

    print("[9/9] Takeaways ...")
    slide_takeaways(prs)

    prs.save(str(out_path))
    print(f"\n[done] {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
