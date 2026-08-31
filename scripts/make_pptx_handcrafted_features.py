#!/usr/bin/env python3
"""
make_pptx_handcrafted_features.py

PPT: Handcrafted feature methods for the label-efficiency benchmark.
  Slide 1  — Title
  Slide 2  — CellProfiler: what, how, features list
  Slide 3  — CellProfiler: intensity feature definitions
  Slide 4  — CellProfiler: Haralick/GLCM feature definitions
  Slide 5  — CellProfiler: feature histograms (DS1/2/3)
  Slide 6  — Ilastik: what, how, filter bank overview
  Slide 7  — Ilastik: filter definitions (6 types × 5 scales)
  Slide 8  — Ilastik: feature histograms (DS1/2/3)

Output: results/handcrafted_features.pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parents[1]
OUT  = REPO / "results" / "handcrafted_features.pptx"

CP_HIST    = REPO / "results" / "cp_feature_histograms.png"
IL_HIST    = REPO / "results" / "ilastik_feature_histograms.png"

W = Inches(13.33)
H = Inches(7.5)

BG_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID   = RGBColor(0x16, 0x21, 0x3E)
BG_PANEL = RGBColor(0x0F, 0x34, 0x56)
ACCENT1  = RGBColor(0xE9, 0x45, 0x60)
ACCENT2  = RGBColor(0x53, 0xD8, 0xFB)
ACCENT3  = RGBColor(0xF5, 0xA6, 0x23)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)


# ---------------------------------------------------------------------------
# Helpers

def _bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height,
         text="", fontsize=14, bold=False, color=WHITE,
         bg=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(fontsize)
    run.font.bold   = bold
    run.font.color.rgb = color
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    return txb


def _rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _img(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 height=Inches(height))


def _title_bar(slide, title, subtitle="", title_color=ACCENT2):
    _rect(slide, 0, 0, 13.33, 1.1, BG_MID)
    _box(slide, 0.3, 0.1, 12.5, 0.65, title,
         fontsize=28, bold=True, color=title_color)
    if subtitle:
        _box(slide, 0.3, 0.72, 12.5, 0.36, subtitle,
             fontsize=14, color=LGRAY)


def _bullet_frame(slide, left, top, width, height, header, bullets,
                  header_color=ACCENT2, bullet_color=WHITE,
                  header_size=15, bullet_size=12, bg=BG_MID):
    _rect(slide, left, top, width, height, bg)
    y = top + 0.12
    _box(slide, left + 0.15, y, width - 0.2, 0.35,
         header, fontsize=header_size, bold=True, color=header_color)
    y += 0.38
    per_bullet = (height - 0.55) / max(len(bullets), 1)
    for b in bullets:
        _box(slide, left + 0.25, y, width - 0.35, per_bullet + 0.05,
             "• " + b, fontsize=bullet_size, color=bullet_color)
        y += per_bullet


# ---------------------------------------------------------------------------
# Slides

def slide_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _rect(sl, 0, 2.5, 13.33, 2.8, BG_MID)
    _box(sl, 0.5, 2.65, 12.3, 1.1,
         "Handcrafted Feature Methods",
         fontsize=40, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    _box(sl, 0.5, 3.7, 12.3, 0.6,
         "CellProfiler-style  ·  Ilastik-style  ·  Label Efficiency Benchmark Prep",
         fontsize=20, color=LGRAY, align=PP_ALIGN.CENTER)
    _box(sl, 0.5, 6.8, 12.3, 0.45,
         "Features extracted from 32-px FA patches  |  DS1 (vinc) · DS2 (pfak) · DS3 (ppax)",
         fontsize=12, color=LGRAY, align=PP_ALIGN.CENTER)


def slide_cp_overview(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "CellProfiler-Style Features — Overview", title_color=ACCENT2)

    # What is CellProfiler
    _bullet_frame(sl, 0.2, 1.2, 6.3, 2.8,
        "What is CellProfiler?",
        ["Open-source image analysis platform (Broad Institute)",
         "Designed for high-throughput microscopy quantification",
         "Used to measure morphology, intensity, and texture per object",
         "Standard tool in cell biology for phenotype profiling"],
        header_color=ACCENT2, bg=BG_MID)

    # How features are used
    _bullet_frame(sl, 6.7, 1.2, 6.3, 2.8,
        "How features are typically used",
        ["Extracted per segmented object (cell, organelle, focus)",
         "Fed into classical ML: SVM, Random Forest, LightGBM",
         "Used for clustering, dimensionality reduction (UMAP/PCA)",
         "Benchmark baseline vs. deep learning representations"],
        header_color=ACCENT3, bg=BG_MID)

    # Our implementation
    _bullet_frame(sl, 0.2, 4.15, 12.8, 2.8,
        "Our Implementation — 50 features per 32×32 FA patch",
        ["Input: 32×32 px FA patches, CIO-normalized (background subtracted, divided by 5 × cell signal contrast)",
         "Intensity features (11): mean, std, median, min, max, integrated sum, MAD, p10, p25, p75, p90",
         "Haralick / GLCM texture features (39): 13 features × 3 distances (d=1, 2, 4 px), averaged over 4 angles (0°, 45°, 90°, 135°)",
         "GLCM computed with 16 gray levels after per-patch min–max quantization → texture is intensity-scale invariant"],
        header_color=ACCENT1, bg=BG_PANEL, bullet_size=12)


def slide_cp_intensity(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "CellProfiler Features — Intensity (11)", title_color=ACCENT2)

    rows = [
        ("intensity_mean",       "Mean pixel value across the patch"),
        ("intensity_std",        "Standard deviation of pixel values"),
        ("intensity_median",     "Median pixel value"),
        ("intensity_min",        "Minimum pixel value"),
        ("intensity_max",        "Maximum pixel value (bright FA peak)"),
        ("intensity_integrated", "Sum of all pixel values (total signal)"),
        ("intensity_mad",        "Median Absolute Deviation — robust spread measure"),
        ("intensity_p10",        "10th percentile — lower tail of signal"),
        ("intensity_p25",        "25th percentile — lower quartile"),
        ("intensity_p75",        "75th percentile — upper quartile"),
        ("intensity_p90",        "90th percentile — upper tail / bright spots"),
    ]

    _rect(sl, 0.2, 1.2, 12.9, 5.9, BG_MID)
    _box(sl, 0.5, 1.25, 4.5, 0.4, "Feature name",
         fontsize=13, bold=True, color=ACCENT2)
    _box(sl, 5.2, 1.25, 7.8, 0.4, "What it captures",
         fontsize=13, bold=True, color=ACCENT2)

    _rect(sl, 0.2, 1.6, 12.9, 0.04, ACCENT2)

    row_h = 0.48
    for i, (name, desc) in enumerate(rows):
        y   = 1.68 + i * row_h
        bg  = BG_PANEL if i % 2 == 0 else BG_MID
        _rect(sl, 0.2, y, 12.9, row_h, bg)
        _box(sl, 0.4, y + 0.05, 4.6, row_h - 0.05,
             name, fontsize=11, color=ACCENT3)
        _box(sl, 5.2, y + 0.05, 7.7, row_h - 0.05,
             desc, fontsize=11, color=WHITE)


def slide_cp_haralick(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "CellProfiler Features — Haralick / GLCM Texture (39)",
               subtitle="13 features × 3 distances (d=1, 2, 4 px) · averaged over 4 angles · 16 gray levels",
               title_color=ACCENT2)

    rows = [
        ("angular_second_moment",    "Energy / uniformity — high for smooth/uniform regions"),
        ("contrast",                 "Local intensity variation between neighbouring pixels"),
        ("correlation",              "Linear dependency of grey levels across neighbours"),
        ("variance",                 "Variance of marginal intensity distribution"),
        ("inverse_difference_moment","Homogeneity — high when neighbouring values are similar"),
        ("sum_average",              "Mean of the sum distribution p_{x+y}"),
        ("sum_variance",             "Variance of the sum distribution"),
        ("sum_entropy",              "Entropy of the sum distribution"),
        ("entropy",                  "Randomness of grey-level pairs — high for complex textures"),
        ("difference_variance",      "Variance of the difference distribution p_{x−y}"),
        ("difference_entropy",       "Entropy of the difference distribution"),
        ("info_meas1",               "Information measure of correlation 1 (HXY1−HXY) / max(HX,HY)"),
        ("info_meas2",               "Information measure of correlation 2  √(1 − e^{−2(HXY2−HXY)})"),
    ]

    _rect(sl, 0.2, 1.25, 12.9, 5.9, BG_MID)
    _box(sl, 0.4, 1.28, 4.8, 0.38, "Feature (suffix _d1 / _d2 / _d4)",
         fontsize=12, bold=True, color=ACCENT2)
    _box(sl, 5.3, 1.28, 7.7, 0.38, "What it captures",
         fontsize=12, bold=True, color=ACCENT2)
    _rect(sl, 0.2, 1.62, 12.9, 0.04, ACCENT2)

    row_h = 0.42
    for i, (name, desc) in enumerate(rows):
        y  = 1.70 + i * row_h
        bg = BG_PANEL if i % 2 == 0 else BG_MID
        _rect(sl, 0.2, y, 12.9, row_h, bg)
        _box(sl, 0.4, y + 0.04, 4.7, row_h,
             name, fontsize=10.5, color=ACCENT3)
        _box(sl, 5.3, y + 0.04, 7.7, row_h,
             desc, fontsize=10.5, color=WHITE)


def slide_cp_histograms(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "CellProfiler Features — Distributions (DS1 / DS2 / DS3)",
               subtitle="50 features · 1–99th percentile clipped · blue=DS1(vinc)  red=DS2(pfak)  green=DS3(ppax)",
               title_color=ACCENT2)
    if CP_HIST.exists():
        _img(sl, CP_HIST, left=0.1, top=1.15, width=13.1)


def slide_ilastik_overview(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "Ilastik-Style Features — Overview", title_color=ACCENT3)

    _bullet_frame(sl, 0.2, 1.2, 6.3, 2.8,
        "What is ilastik?",
        ["Interactive machine learning tool for image segmentation",
         "Uses multiscale filter banks as pixel-level features",
         "User labels a few pixels → Random Forest classifies rest",
         "Widely used in bioimage analysis for semantic segmentation"],
        header_color=ACCENT3, bg=BG_MID)

    _bullet_frame(sl, 6.7, 1.2, 6.3, 2.8,
        "How features are typically used",
        ["Computed per pixel at multiple spatial scales",
         "Capture local structure: edges, blobs, ridges, texture",
         "Fed into pixel-wise Random Forest classifier in ilastik",
         "Here: aggregated (mean + std) per patch → patch-level features"],
        header_color=ACCENT2, bg=BG_MID)

    _bullet_frame(sl, 0.2, 4.15, 12.8, 2.8,
        "Our Implementation — 80 features per 32×32 FA patch",
        ["8 filter types × 5 scales (σ = 0.3, 0.7, 1.0, 1.6, 3.5 px) = 40 feature maps",
         "Each feature map summarized with mean + std over the 32×32 patch → 80 features total",
         "Filters: Gaussian, Laplacian of Gaussian (LoG), Gradient magnitude, Difference of Gaussians (DoG),",
         "  Structure tensor eigenvalues (×2: large λ, small λ),  Hessian eigenvalues (×2: large λ, small λ)"],
        header_color=ACCENT1, bg=BG_PANEL, bullet_size=12)


def slide_ilastik_filters(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "Ilastik Features — Filter Definitions",
               subtitle="Each filter applied at σ = 0.3, 0.7, 1.0, 1.6, 3.5 px · summarized as mean + std per patch",
               title_color=ACCENT3)

    rows = [
        ("Gaussian smoothing",
         "Low-pass filter; captures local average intensity at each scale. "
         "Large σ → coarse structure; small σ → fine detail."),
        ("Laplacian of Gaussian (LoG)",
         "Second-derivative blob detector. Highlights bright spots and dark holes. "
         "Response peaks where image curvature matches the scale σ."),
        ("Gaussian gradient magnitude",
         "First-derivative edge detector: √(Gx² + Gy²) after Gaussian smoothing. "
         "High at FA boundaries and intensity transitions."),
        ("Difference of Gaussians (DoG)",
         "Band-pass filter: Gaussian(σ) − Gaussian(σ√2). "
         "Approximates LoG; highlights structure between two spatial scales."),
        ("Structure tensor — large eigenvalue",
         "Describes local orientation coherence. Large λ₁ indicates strong "
         "directional gradient (edges, oriented fibres). Built from outer product of smoothed gradients."),
        ("Structure tensor — small eigenvalue",
         "Small λ₂ ≈ 0 at edges (one direction), large at corners/junctions. "
         "Ratio λ₁/λ₂ distinguishes edges from isotropic texture."),
        ("Hessian — large eigenvalue",
         "Second-order shape descriptor. Large |λ₁| at ridges or valleys. "
         "Used to detect elongated FA structures and fibrillar adhesion ridges."),
        ("Hessian — small eigenvalue",
         "Near-zero at ridges (one principal curvature). Both eigenvalues large "
         "at blob-like focal complexes. Combination distinguishes FA sub-types."),
    ]

    colors = [ACCENT2, ACCENT3, ACCENT1,
              RGBColor(0x53, 0xD8, 0xFB),
              RGBColor(0xF5, 0xA6, 0x23),
              RGBColor(0xE9, 0x45, 0x60),
              RGBColor(0x2E, 0xCC, 0x71),
              RGBColor(0xAF, 0x7A, 0xC5)]

    _rect(sl, 0.2, 1.25, 12.9, 5.9, BG_MID)
    row_h = 0.69
    for i, (name, desc) in enumerate(rows):
        y  = 1.28 + i * row_h
        bg = BG_PANEL if i % 2 == 0 else BG_MID
        _rect(sl, 0.2, y, 12.9, row_h, bg)
        _box(sl, 0.4, y + 0.06, 3.2, row_h - 0.08,
             name, fontsize=11, bold=True, color=colors[i])
        _box(sl, 3.7, y + 0.06, 9.3, row_h - 0.08,
             desc, fontsize=10.5, color=WHITE)


def slide_ilastik_histograms(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl)
    _title_bar(sl, "Ilastik Features — Distributions (DS1 / DS2 / DS3)",
               subtitle="80 features · 1–99th percentile clipped · blue=DS1(vinc)  red=DS2(pfak)  green=DS3(ppax)",
               title_color=ACCENT3)
    if IL_HIST.exists():
        _img(sl, IL_HIST, left=0.1, top=1.15, width=13.1)


# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_cp_overview(prs)
    slide_cp_intensity(prs)
    slide_cp_haralick(prs)
    slide_cp_histograms(prs)
    slide_ilastik_overview(prs)
    slide_ilastik_filters(prs)
    slide_ilastik_histograms(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
